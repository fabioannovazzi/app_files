"""Mechanical helpers for the Clara Codex plugin."""

from __future__ import annotations

import html
import json
import logging
import os
import re
import shutil
import subprocess
import tempfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from datetime import datetime, timezone
from hashlib import sha256
from pathlib import Path
from typing import Any, Iterable, Mapping, Sequence
from zipfile import ZIP_DEFLATED, BadZipFile, ZipFile

from html_deck_runtime import (
    apply_fixed_16_9_deck_runtime,
    assert_fixed_16_9_deck_runtime,
)

LOGGER = logging.getLogger(__name__)

__all__ = [
    "CASE_BRIEF_FILENAME",
    "CLARA_KICKOFF_PREPARATION_FILENAME",
    "CLARA_KICKOFF_DECK_FILENAME",
    "CLARA_MANDATE_FILENAME",
    "CLARA_PARTNER_BRIEF_FILENAME",
    "CASE_FILES",
    "INCLUSION_BUNDLES_FILENAME",
    "CaseWorkspaceError",
    "CaseExchangeExportResult",
    "CaseExchangeImportResult",
    "EditablePptxMergeInputResult",
    "CaseFileCopyResult",
    "LegacyPptxNormalizationResult",
    "CaseSupportPackageResult",
    "CaseWorkspaceArchiveResult",
    "MaterialDeleteResult",
    "CaseBriefResult",
    "ClaraKickoffPreparationResult",
    "ClaraKickoffDeckResult",
    "ClaraPartnerBriefResult",
    "DecisionPackResult",
    "InclusionBundleApplyResult",
    "InclusionReviewResult",
    "ISSUE_STATUSES",
    "JUDGEMENT_KINDS",
    "JUDGEMENT_STATUSES",
    "MATERIAL_STATUSES",
    "MATERIAL_TYPES",
    "OPEN_QUESTION_STATUSES",
    "SUPPORTED_LANGUAGES",
    "add_judgement_entries",
    "add_open_question",
    "apply_inclusion_bundles",
    "audit_human_visible_document_text",
    "build_decision_pack",
    "build_inclusion_review",
    "copy_case_file",
    "delete_materials",
    "discover_material_paths",
    "export_case_update",
    "export_case_workspace_archive",
    "index_materials",
    "ingest_note_file",
    "ingest_note_text",
    "inspect_pptx_legacy_media",
    "import_case_update",
    "initialize_case",
    "load_case_file",
    "normalize_legacy_pptx_for_editable_merge",
    "prepare_editable_pptx_merge_input",
    "prepare_clara_kickoff",
    "prepare_support_package",
    "register_material",
    "refresh_case_brief",
    "render_clara_kickoff_deck",
    "render_clara_partner_brief",
    "resolve_soffice_binary",
    "set_judgement_status",
    "set_judgement_statuses",
    "update_clara_mandate_from_kickoff",
    "upsert_case_issues",
    "validate_case_workspace",
]

SCHEMA_VERSION = 1
EXCHANGE_SCHEMA_VERSION = 1
EXCHANGE_SOURCE = "case_notes_case_update"
SUPPORTED_LANGUAGES = {"it", "en", "fr", "de", "es"}
CASE_STATUSES = {"active", "paused", "complete", "archived"}
MATERIAL_TYPES = {
    "source",
    "note",
    "transcript",
    "memo",
    "proposal",
    "questionnaire",
    "other",
}
MATERIAL_STATUSES = {"indexed", "reviewed", "superseded", "excluded"}
JUDGEMENT_KINDS = {
    "fact",
    "advisor_judgement",
    "codex_inference",
    "open_question",
    "decision_implication",
}
JUDGEMENT_STATUSES = {"pending", "approved", "rejected"}
OPEN_QUESTION_STATUSES = {"open", "answered", "dismissed"}
ISSUE_STATUSES = {"active", "resolved", "parked"}
SUPPORTED_FILE_SUFFIXES = {".md", ".txt", ".docx", ".pdf", ".pptx"}
TEXT_SUFFIXES = {".md", ".txt"}
CASE_FILE_KINDS = {"auto", "source", "note", "audio", "deck", "presentation"}
AUDIO_FILE_SUFFIXES = {
    ".aac",
    ".aiff",
    ".flac",
    ".m4a",
    ".mp3",
    ".mp4",
    ".ogg",
    ".wav",
    ".wma",
}
PRESENTATION_FILE_SUFFIXES = {".key", ".ppt", ".pptx"}
NOTE_FILE_SUFFIXES = {".md", ".txt"}
PRESENTATION_NAME_TOKENS = {
    "bozza",
    "deck",
    "draft",
    "incontro",
    "presentazione",
    "presentation",
    "slide",
    "slides",
}
LEGACY_PPTX_MEDIA_SUFFIXES = {".emf", ".wmf"}
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
CUSTOM_PROPS_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.custom-properties+xml"
)
DEFAULT_SOFFICE_PATHS = (
    Path("/opt/homebrew/bin/soffice"),
    Path("/usr/local/bin/soffice"),
    Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"),
    Path("/mnt/c/Program Files/LibreOffice/program/soffice.exe"),
    Path("/mnt/c/Program Files (x86)/LibreOffice/program/soffice.exe"),
)
NOTE_NAME_TOKENS = {
    "appunti",
    "intervista",
    "interview",
    "nota",
    "note",
    "notes",
    "sintesi",
    "transcript",
    "trascrizione",
}

CASE_FILES = {
    "manifest": "case_manifest.json",
    "materials": "material_registry.json",
    "judgement": "judgement_log.json",
    "open_questions": "open_questions.json",
    "issues": "case_issues.json",
    "clara_mandate": "clara_mandate.json",
}
CASE_BRIEF_FILENAME = "case_brief.md"
CLARA_MANDATE_FILENAME = "clara_mandate.json"
CLARA_KICKOFF_PREPARATION_FILENAME = "clara_kickoff_preparation.md"
CLARA_KICKOFF_DECK_FILENAME = "clara_kickoff_deck.html"
CLARA_PARTNER_BRIEF_FILENAME = "clara_partner_brief.html"
INCLUSION_REVIEW_FILENAME = "inclusion_review.md"
INCLUSION_BUNDLES_FILENAME = "inclusion_bundles.json"
SUPPORT_REQUEST_FILENAME = "support_request.md"
SHARE_EXCLUDED_DIR_NAMES = {
    ".cache",
    ".git",
    ".mypy_cache",
    ".pytest_cache",
    ".ruff_cache",
    ".venv",
    "__pycache__",
    "case_support_exports",
    "case_share_exports",
    "env",
    "exchange_exports",
    "node_modules",
    "venv",
}
SHARE_EXCLUDED_FILE_NAMES = {
    ".DS_Store",
}


# Deterministic because these are fixed scaffolding, metadata, placeholder, and
# filler markers that should never be visible in ordinary human-readable Clara
# documents. Provenance workpapers are the explicit exception.
HUMAN_VISIBLE_DOCUMENT_BANNED_PATTERNS: tuple[tuple[str, str], ...] = (
    ("visible judgement/source id", r"\bjud-\d{4,}\b"),
    ("page-number scaffolding", r"\b(?:pagina|pagine|page|pages)\s+\d+"),
    ("page-count scaffolding", r"\b\d+\s*[-/]\s*\d+\s+pagine\b"),
    ("placeholder label", r"\betichetta da confermare\b"),
    ("working-pack jargon", r"\bworking\s+(?:review\s+)?pack\b"),
    ("working-draft jargon", r"\bworking\s+draft\b"),
    ("opaque dossier title", r"\bdossier di revisione operativo\b"),
    ("pseudo-editorial filler", r"\bdocumento piu decidibile\b"),
    ("pseudo-editorial filler", r"\baiutare a pensare\b"),
    ("pseudo-editorial filler", r"\blettura non lineare\b"),
    ("generic value language", r"\bcreare valore\b"),
    ("generic value language", r"\bgive concrete levers\b"),
    ("generic process language", r"\bcommenti generici\b"),
    ("generic process language", r"\bapprovazione del testo\b"),
    ("unhelpful Clara self-reference", r"\bclara ha preparato\b"),
    ("unhelpful Clara self-reference", r"\bcosa clara (?:pensa|deve|serve)\b"),
    ("unhelpful Clara self-reference", r"\bwhat clara (?:thinks|should|needs)\b"),
    ("unhelpful instruction copy", r"\busa questa scheda\b"),
    ("unhelpful instruction copy", r"\buse this brief\b"),
    ("decorative warning styling", r"\bclass=[\"'][^\"']*warning[^\"']*[\"']"),
    ("decorative shadow styling", r"\bbox-shadow\s*:"),
)


class CaseWorkspaceError(ValueError):
    """Raised when a case workspace file violates the plugin contract."""


def _document_quality_scan_text(value: str) -> str:
    normalized = value.casefold()
    normalized = normalized.replace("’", "'")
    normalized = normalized.replace("è", "e'")
    normalized = normalized.replace("é", "e'")
    normalized = normalized.replace("È", "e'")
    normalized = normalized.replace("É", "e'")
    normalized = normalized.replace("ù", "u")
    normalized = normalized.replace("Ù", "u")
    normalized = normalized.replace("&amp;", "&")
    return normalized


def audit_human_visible_document_text(text: str) -> list[str]:
    """Return mechanical BS markers in ordinary human-visible Clara output."""

    scan_text = _document_quality_scan_text(text)
    violations: list[str] = []
    for label, pattern in HUMAN_VISIBLE_DOCUMENT_BANNED_PATTERNS:
        if re.search(pattern, scan_text):
            violations.append(label)
    return _dedupe_preserve_order(violations)


def _assert_human_visible_document_quality(text: str, *, label: str) -> None:
    violations = audit_human_visible_document_text(text)
    if violations:
        joined = "; ".join(violations)
        raise CaseWorkspaceError(f"{label} failed Clara quality gate: {joined}")


@dataclass(frozen=True)
class CaseBriefResult:
    """Path and counts produced by derived case-brief rendering."""

    brief_path: Path
    approved_count: int
    pending_count: int
    rejected_count: int
    open_question_count: int
    issue_count: int


@dataclass(frozen=True)
class DecisionPackResult:
    """Paths and counts produced by decision-pack rendering."""

    markdown_path: Path
    docx_path: Path
    workpaper_markdown_path: Path
    workpaper_docx_path: Path
    approved_count: int
    pending_count: int
    rejected_count: int


@dataclass(frozen=True)
class InclusionReviewResult:
    """Path and counts produced by deterministic inclusion-review rendering."""

    review_path: Path
    approved_count: int
    pending_count: int
    rejected_count: int


@dataclass(frozen=True)
class InclusionBundleApplyResult:
    """Path and counts produced by deterministic inclusion-bundle persistence."""

    bundles_path: Path
    bundle_count: int
    bundled_entry_count: int


@dataclass(frozen=True)
class CaseExchangeExportResult:
    """Package metadata produced by deterministic case-update export."""

    package_path: Path
    exchange_id: str
    material_count: int
    judgement_count: int
    open_question_count: int
    included_file_count: int


@dataclass(frozen=True)
class CaseWorkspaceArchiveResult:
    """Metadata produced by clean case-workspace archive export."""

    package_path: Path
    included_file_count: int
    excluded_file_count: int
    excluded_bytes: int


@dataclass(frozen=True)
class CaseSupportPackageResult:
    """Metadata produced by a clean support package for a delivery backstop."""

    package_path: Path
    support_request_archive_path: str
    included_file_count: int
    excluded_file_count: int
    excluded_bytes: int


@dataclass(frozen=True)
class ClaraKickoffPreparationResult:
    """Paths and counts produced by Clara kickoff preparation."""

    mandate_path: Path
    preparation_path: Path
    material_count: int
    baseline_source_count: int


@dataclass(frozen=True)
class ClaraPartnerBriefResult:
    """Path produced by Clara partner brief rendering."""

    html_path: Path
    open_clarification_count: int
    next_step_count: int


@dataclass(frozen=True)
class ClaraKickoffDeckResult:
    """Path and counts produced by Clara's first partner kickoff deck."""

    html_path: Path
    hypothesis_count: int
    open_question_count: int


@dataclass(frozen=True)
class CaseExchangeImportResult:
    """Counts produced by deterministic append-only case-update import."""

    exchange_id: str
    imported_material_count: int
    imported_judgement_count: int
    imported_open_question_count: int
    skipped_count: int
    conflict_count: int


@dataclass(frozen=True)
class CaseFileCopyResult:
    """Destination and optional registry record for a copied case file."""

    source_path: Path
    destination_path: Path
    kind: str
    copied: bool
    registered_material: dict[str, Any] | None
    legacy_pptx_normalization: LegacyPptxNormalizationResult | None = None


@dataclass(frozen=True)
class MaterialDeleteResult:
    """Canonical references removed by a material deletion."""

    removed_material_ids: tuple[str, ...]
    missing_material_ids: tuple[str, ...]
    removed_material_paths: tuple[Path, ...]
    updated_judgement_ids: tuple[str, ...]
    unanchored_judgement_ids: tuple[str, ...]
    removed_mandate_source_material_ids: tuple[str, ...]
    removed_mandate_voice_session_paths: tuple[str, ...]
    removed_preparation_material_anchor_ids: tuple[str, ...]
    orphan_candidate_paths: tuple[Path, ...]
    removed_empty_orphan_dirs: tuple[Path, ...]
    brief_path: Path


@dataclass(frozen=True)
class EditablePptxMergeInputResult:
    """Audited PPTX base selection for an editable slide merge."""

    source_path: Path
    merge_base_path: Path
    report_path: Path
    status: str
    source_legacy_media: tuple[str, ...]
    merge_base_legacy_media: tuple[str, ...]
    normalized_path: Path | None
    skip_normalization_reason: str | None


@dataclass(frozen=True)
class LegacyPptxNormalizationResult:
    """Result of normalizing a legacy PowerPoint deck for editable slide merging."""

    source_path: Path
    output_path: Path
    report_path: Path
    normalized: bool
    legacy_media_before: tuple[str, ...]
    legacy_media_after: tuple[str, ...]
    soffice_binary: Path | None


def _now_iso(now: datetime | None = None) -> str:
    timestamp = now or datetime.now(timezone.utc)
    if timestamp.tzinfo is None:
        timestamp = timestamp.replace(tzinfo=timezone.utc)
    return timestamp.astimezone(timezone.utc).replace(microsecond=0).isoformat()


def _slugify(value: str, fallback: str = "note") -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower()).strip("-")
    return slug or fallback


def _compact_timestamp(value: str) -> str:
    return value.replace("-", "").replace(":", "")[:15] + "Z"


def _case_path(case_dir: Path, key: str) -> Path:
    return case_dir / CASE_FILES[key]


def _write_json(path: Path, payload: Mapping[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(payload, indent=2, sort_keys=True, ensure_ascii=True) + "\n",
        encoding="utf-8",
    )


def _read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _next_id(prefix: str, existing_ids: Iterable[str]) -> str:
    used_numbers: list[int] = []
    for item_id in existing_ids:
        match = re.fullmatch(rf"{re.escape(prefix)}-(\d+)", item_id)
        if match:
            used_numbers.append(int(match.group(1)))
    next_number = max(used_numbers, default=0) + 1
    return f"{prefix}-{next_number:04d}"


def _dedupe_preserve_order(values: Iterable[str]) -> list[str]:
    seen: set[str] = set()
    deduped: list[str] = []
    for value in values:
        clean_value = value.strip()
        if clean_value and clean_value not in seen:
            deduped.append(clean_value)
            seen.add(clean_value)
    return deduped


def _touch_manifest(case_dir: Path, timestamp: str) -> None:
    manifest_path = _case_path(case_dir, "manifest")
    manifest = _read_json(manifest_path)
    manifest["updated_at"] = timestamp
    _write_json(manifest_path, manifest)


def _validate_choice(value: str, allowed: set[str], field_name: str) -> None:
    if value not in allowed:
        allowed_values = ", ".join(sorted(allowed))
        raise CaseWorkspaceError(
            f"{field_name} must be one of: {allowed_values}; got {value!r}"
        )


def _empty_material_registry() -> dict[str, Any]:
    return {"schema_version": SCHEMA_VERSION, "materials": []}


def _empty_judgement_log() -> dict[str, Any]:
    return {"schema_version": SCHEMA_VERSION, "entries": []}


def _empty_open_questions() -> dict[str, Any]:
    return {"schema_version": SCHEMA_VERSION, "questions": []}


def _empty_case_issues() -> dict[str, Any]:
    return {"schema_version": SCHEMA_VERSION, "issues": []}


def _empty_clara_mandate() -> dict[str, Any]:
    return {
        "schema_version": SCHEMA_VERSION,
        "persona": "Clara",
        "partner_role": "senior partner",
        "status": "not_started",
        "prepared_at": None,
        "kickoff_imported_at": None,
        "updated_at": None,
        "preparation": {
            "case_snapshot": {},
            "material_anchors": [],
            "succession_lenses": [],
            "red_flags": [],
            "industry_context": [],
            "external_research": [],
            "preparation_note": "",
        },
        "mandate": {
            "engagement_objective": "",
            "client_decision": "",
            "clara_understanding": "",
            "partner_starting_orientation": "",
            "sensitive_points": [],
            "what_clara_should_investigate": [],
            "what_clara_should_not_waste_time_on": [],
            "essential_clarifications": [],
            "next_steps": [],
        },
        "source_material_ids": [],
        "voice_session_paths": [],
    }


def initialize_case(
    case_dir: Path,
    *,
    client: str,
    project: str,
    objective: str,
    audience: str,
    output_language: str = "it",
    status: str = "active",
    overwrite: bool = False,
    now: datetime | None = None,
) -> dict[str, Any]:
    """Create or load the four durable case workspace files."""

    _validate_choice(output_language, SUPPORTED_LANGUAGES, "output_language")
    _validate_choice(status, CASE_STATUSES, "status")

    case_dir.mkdir(parents=True, exist_ok=True)
    timestamp = _now_iso(now)
    manifest_path = _case_path(case_dir, "manifest")
    if overwrite or not manifest_path.exists():
        manifest = {
            "schema_version": SCHEMA_VERSION,
            "client": client,
            "project": project,
            "objective": objective,
            "audience": audience,
            "status": status,
            "output_language": output_language,
            "created_at": timestamp,
            "updated_at": timestamp,
        }
        _write_json(manifest_path, manifest)
    else:
        manifest = _read_json(manifest_path)

    defaults = {
        "materials": _empty_material_registry(),
        "judgement": _empty_judgement_log(),
        "open_questions": _empty_open_questions(),
        "issues": _empty_case_issues(),
        "clara_mandate": _empty_clara_mandate(),
    }
    for key, payload in defaults.items():
        path = _case_path(case_dir, key)
        if overwrite or not path.exists():
            _write_json(path, payload)

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))
    refresh_case_brief(case_dir, now=now)
    return manifest


def load_case_file(case_dir: Path, key: str) -> dict[str, Any]:
    """Load one canonical workspace JSON file by logical key."""

    if key not in CASE_FILES:
        raise CaseWorkspaceError(f"Unknown case file key: {key}")
    return _read_json(_case_path(case_dir, key))


def _relative_path_from_material(case_dir: Path, material: Mapping[str, Any]) -> str:
    path = _material_path(case_dir, material)
    try:
        return str(path.relative_to(case_dir.resolve()))
    except ValueError:
        return str(path)


def _validate_audio_pointer_links(
    *,
    case_dir: Path,
    materials: list[Any],
) -> list[str]:
    errors: list[str] = []
    material_by_id = {
        str(material.get("id")): material
        for material in materials
        if isinstance(material, dict) and material.get("id")
    }
    for material in materials:
        if not isinstance(material, dict):
            continue
        if material.get("material_type") != "transcript":
            continue
        metadata = material.get("source_metadata")
        if not isinstance(metadata, dict):
            continue
        pointer_id = str(metadata.get("raw_audio_pointer_material_id") or "").strip()
        if not pointer_id:
            continue
        transcript_id = str(material.get("id") or "").strip()
        pointer = material_by_id.get(pointer_id)
        if pointer is None:
            errors.append(
                "material_registry.json: transcript "
                f"{transcript_id} links missing raw audio pointer {pointer_id}"
            )
            continue
        pointer_metadata = pointer.get("source_metadata")
        if not isinstance(pointer_metadata, dict):
            errors.append(
                "material_registry.json: audio pointer "
                f"{pointer_id} linked to transcript {transcript_id} "
                "has no source_metadata"
            )
            pointer_metadata = {}
        if pointer_metadata.get("transcription_status") != "transcribed":
            errors.append(
                "material_registry.json: audio pointer "
                f"{pointer_id} linked to transcript {transcript_id} "
                "is not marked transcribed"
            )
        if pointer_metadata.get("linked_transcript_material_id") != transcript_id:
            errors.append(
                "material_registry.json: audio pointer "
                f"{pointer_id} does not link back to transcript {transcript_id}"
            )
        expected_path = _relative_path_from_material(case_dir, material)
        if pointer_metadata.get("linked_transcript_path") != expected_path:
            errors.append(
                "material_registry.json: audio pointer "
                f"{pointer_id} has stale linked transcript path for {transcript_id}"
            )
        pointer_path = _material_path(case_dir, pointer)
        if pointer_path.is_file():
            try:
                pointer_text = pointer_path.read_text(encoding="utf-8")
            except UnicodeDecodeError:
                pointer_text = ""
            stale_markers = (
                "non ancora trascritto",
                "not yet transcribed",
                "Quando sara' prodotta la trascrizione",
                "Quando sarà prodotta la trascrizione",
            )
            if any(marker in pointer_text for marker in stale_markers):
                errors.append(
                    "material_registry.json: audio pointer "
                    f"{pointer_id} linked to transcript {transcript_id} "
                    "still says not transcribed"
                )
    return errors


def validate_case_workspace(case_dir: Path) -> list[str]:
    """Validate mechanical JSON schemas without making semantic judgments."""

    errors: list[str] = []
    for key, filename in CASE_FILES.items():
        path = case_dir / filename
        if not path.exists():
            errors.append(f"missing {filename}")
            continue
        try:
            payload = _read_json(path)
        except json.JSONDecodeError as exc:
            errors.append(f"{filename}: invalid JSON: {exc}")
            continue
        if payload.get("schema_version") != SCHEMA_VERSION:
            errors.append(f"{filename}: unsupported schema_version")
        if key == "manifest":
            for field in (
                "client",
                "project",
                "objective",
                "audience",
                "status",
                "output_language",
            ):
                if not payload.get(field):
                    errors.append(f"{filename}: missing {field}")
            if payload.get("status") not in CASE_STATUSES:
                errors.append(f"{filename}: invalid status")
            if payload.get("output_language") not in SUPPORTED_LANGUAGES:
                errors.append(f"{filename}: invalid output_language")
        elif key == "materials":
            materials = payload.get("materials")
            if not isinstance(materials, list):
                errors.append(f"{filename}: materials must be a list")
            else:
                errors.extend(
                    _validate_audio_pointer_links(
                        case_dir=case_dir,
                        materials=materials,
                    )
                )
        elif key == "judgement":
            entries = payload.get("entries")
            if not isinstance(entries, list):
                errors.append(f"{filename}: entries must be a list")
            else:
                for entry in entries:
                    if entry.get("kind") not in JUDGEMENT_KINDS:
                        errors.append(f"{filename}: invalid judgement kind")
                    if entry.get("status") not in JUDGEMENT_STATUSES:
                        errors.append(f"{filename}: invalid judgement status")
        elif key == "open_questions":
            if not isinstance(payload.get("questions"), list):
                errors.append(f"{filename}: questions must be a list")
        elif key == "issues":
            issues = payload.get("issues")
            if not isinstance(issues, list):
                errors.append(f"{filename}: issues must be a list")
            else:
                for issue in issues:
                    if issue.get("status") not in ISSUE_STATUSES:
                        errors.append(f"{filename}: invalid issue status")
        elif key == "clara_mandate":
            if not isinstance(payload.get("preparation"), dict):
                errors.append(f"{filename}: preparation must be an object")
            if not isinstance(payload.get("mandate"), dict):
                errors.append(f"{filename}: mandate must be an object")
            if not isinstance(payload.get("source_material_ids"), list):
                errors.append(f"{filename}: source_material_ids must be a list")
            if not isinstance(payload.get("voice_session_paths"), list):
                errors.append(f"{filename}: voice_session_paths must be a list")
    return errors


def _mechanical_text_summary(text: str, *, limit: int = 420) -> str:
    clean = re.sub(r"\s+", " ", text).strip()
    if not clean:
        return "No readable preview text."
    if len(clean) <= limit:
        return clean
    return f"{clean[:limit].rstrip()}..."


def _summarize_file(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        try:
            return _mechanical_text_summary(path.read_text(encoding="utf-8"))
        except UnicodeDecodeError:
            return "Text file indexed; UTF-8 preview unavailable."
    if suffix == ".docx":
        try:
            from docx import Document
        except ImportError:
            return "DOCX indexed; install python-docx for mechanical preview."
        document = Document(path)
        text = " ".join(paragraph.text for paragraph in document.paragraphs[:12])
        return _mechanical_text_summary(text)
    if suffix == ".pdf":
        return "PDF indexed; semantic review remains in Codex."
    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            return "PPTX indexed; install python-pptx for mechanical preview."
        presentation = Presentation(path)
        slide_text: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            if index > 12:
                break
            texts = [
                str(getattr(shape, "text", "")).strip()
                for shape in slide.shapes
                if str(getattr(shape, "text", "")).strip()
            ]
            if texts:
                slide_text.append(f"Slide {index}: {' '.join(texts)}")
        return _mechanical_text_summary(" ".join(slide_text))
    return "File indexed without preview."


def discover_material_paths(paths: Sequence[Path]) -> list[Path]:
    """Return supported files from explicit files or folders, without copying."""

    discovered: list[Path] = []
    for source in paths:
        if not source.exists():
            raise FileNotFoundError(source)
        if source.is_file() and source.suffix.lower() in SUPPORTED_FILE_SUFFIXES:
            discovered.append(source)
        elif source.is_dir():
            discovered.extend(
                path
                for path in source.rglob("*")
                if path.is_file() and path.suffix.lower() in SUPPORTED_FILE_SUFFIXES
            )
    return sorted({path.resolve() for path in discovered})


def register_material(
    case_dir: Path,
    path: Path,
    *,
    material_type: str = "source",
    title: str | None = None,
    status: str = "indexed",
    summary: str | None = None,
    source_metadata: Mapping[str, Any] | None = None,
    now: datetime | None = None,
) -> dict[str, Any]:
    """Register one source path in place and return the material record."""

    _validate_choice(material_type, MATERIAL_TYPES, "material_type")
    _validate_choice(status, MATERIAL_STATUSES, "status")
    if not path.exists():
        raise FileNotFoundError(path)

    timestamp = _now_iso(now)
    registry_path = _case_path(case_dir, "materials")
    registry = _read_json(registry_path)
    materials = registry["materials"]
    resolved_path = str(path.resolve())
    existing_ids = [item["id"] for item in materials]
    material_summary = summary if summary is not None else _summarize_file(path)
    metadata = dict(source_metadata or {})

    for item in materials:
        if item["path"] == resolved_path:
            item.update(
                {
                    "title": title or item["title"],
                    "material_type": material_type,
                    "status": status,
                    "summary": material_summary,
                    "last_reviewed": item.get("last_reviewed"),
                    "updated_at": timestamp,
                }
            )
            if metadata:
                existing_metadata = item.get("source_metadata")
                if existing_metadata is not None and not isinstance(
                    existing_metadata, dict
                ):
                    raise CaseWorkspaceError("source_metadata must be an object")
                item["source_metadata"] = {
                    **dict(existing_metadata or {}),
                    **metadata,
                }
            _write_json(registry_path, registry)
            _touch_manifest(case_dir, timestamp)
            refresh_case_brief(case_dir, now=now)
            return item

    material = {
        "id": _next_id("mat", existing_ids),
        "path": resolved_path,
        "title": title or path.stem.replace("_", " ").replace("-", " ").title(),
        "material_type": material_type,
        "status": status,
        "summary": material_summary,
        "added_at": timestamp,
        "updated_at": timestamp,
        "last_reviewed": None,
    }
    if metadata:
        material["source_metadata"] = metadata
    materials.append(material)
    _write_json(registry_path, registry)
    _touch_manifest(case_dir, timestamp)
    refresh_case_brief(case_dir, now=now)
    return material


def index_materials(
    case_dir: Path,
    paths: Sequence[Path],
    *,
    material_type: str = "source",
    source_metadata: Mapping[str, Any] | None = None,
    now: datetime | None = None,
) -> list[dict[str, Any]]:
    """Index supported materials from files or folders without copying them."""

    return [
        register_material(
            case_dir,
            path,
            material_type=material_type,
            source_metadata=source_metadata,
            now=now,
        )
        for path in discover_material_paths(paths)
    ]


def _material_path(case_dir: Path, material: Mapping[str, Any]) -> Path:
    raw_path = Path(str(material.get("path", ""))).expanduser()
    if raw_path.is_absolute():
        return raw_path.resolve()
    return (case_dir / raw_path).resolve()


def _case_owned_path(case_dir: Path, path: Path) -> bool:
    try:
        path.resolve().relative_to(case_dir.resolve())
    except ValueError:
        return False
    return True


def _voice_session_path_for_path(case_dir: Path, path: Path) -> str | None:
    try:
        relative_path = path.resolve().relative_to(case_dir.resolve())
    except ValueError:
        return None
    if len(relative_path.parts) < 2 or relative_path.parts[0] != "voice_sessions":
        return None
    return str(Path(relative_path.parts[0]) / relative_path.parts[1])


def _dedupe_paths(paths: Iterable[Path]) -> tuple[Path, ...]:
    seen: set[str] = set()
    deduped: list[Path] = []
    for path in paths:
        key = str(path.resolve())
        if key not in seen:
            deduped.append(path)
            seen.add(key)
    return tuple(deduped)


def _filter_material_ids(
    values: Any,
    deleted_ids: set[str],
) -> tuple[list[str], list[str]]:
    if not isinstance(values, list):
        return [], []
    before = [str(item) for item in values if str(item).strip()]
    after = [item for item in before if item not in deleted_ids]
    removed = [item for item in before if item in deleted_ids]
    return _dedupe_preserve_order(after), _dedupe_preserve_order(removed)


def delete_materials(
    case_dir: Path,
    material_ids: Sequence[str],
    *,
    ignore_missing: bool = False,
    remove_empty_orphan_dirs: bool = False,
    now: datetime | None = None,
) -> MaterialDeleteResult:
    """Remove material records and scrub canonical source references.

    This is intentionally conservative about files. It removes JSON references,
    refreshes the derived brief, and reports case-owned orphan candidates. It
    deletes only empty orphan directories when explicitly requested.
    """

    requested_ids = _dedupe_preserve_order(str(item) for item in material_ids)
    if not requested_ids:
        raise CaseWorkspaceError("at least one material id is required")

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    timestamp = _now_iso(now)
    registry_path = _case_path(case_dir, "materials")
    registry = _read_json(registry_path)
    materials = registry.get("materials")
    if not isinstance(materials, list):
        raise CaseWorkspaceError("material_registry.json: materials must be a list")

    materials_by_id = {
        str(item.get("id")): item for item in materials if isinstance(item, dict)
    }
    missing_ids = tuple(
        item_id for item_id in requested_ids if item_id not in materials_by_id
    )
    if missing_ids and not ignore_missing:
        raise CaseWorkspaceError("unknown material id(s): " + ", ".join(missing_ids))

    removed_materials = [
        materials_by_id[item_id]
        for item_id in requested_ids
        if item_id in materials_by_id
    ]
    removed_ids = tuple(str(item["id"]) for item in removed_materials)
    removed_id_set = set(removed_ids)
    if not removed_materials:
        brief_result = refresh_case_brief(case_dir, now=now)
        return MaterialDeleteResult(
            removed_material_ids=(),
            missing_material_ids=missing_ids,
            removed_material_paths=(),
            updated_judgement_ids=(),
            unanchored_judgement_ids=(),
            removed_mandate_source_material_ids=(),
            removed_mandate_voice_session_paths=(),
            removed_preparation_material_anchor_ids=(),
            orphan_candidate_paths=(),
            removed_empty_orphan_dirs=(),
            brief_path=brief_result.brief_path,
        )

    removed_path_list = [_material_path(case_dir, item) for item in removed_materials]
    removed_paths = _dedupe_paths(removed_path_list)
    removed_session_paths = {
        session_path
        for path in removed_paths
        if (session_path := _voice_session_path_for_path(case_dir, path)) is not None
    }

    remaining_materials = [
        item for item in materials if str(item.get("id")) not in removed_id_set
    ]
    registry["materials"] = remaining_materials
    _write_json(registry_path, registry)

    judgement_path = _case_path(case_dir, "judgement")
    judgement = _read_json(judgement_path)
    updated_judgement_ids: list[str] = []
    unanchored_judgement_ids: list[str] = []
    for entry in judgement.get("entries", []):
        if not isinstance(entry, dict):
            continue
        source_ids = entry.get("source_material_ids", [])
        remaining_source_ids, removed_source_ids = _filter_material_ids(
            source_ids,
            removed_id_set,
        )
        if not removed_source_ids:
            continue
        entry["source_material_ids"] = remaining_source_ids
        entry_id = str(entry.get("id", "")).strip()
        if entry_id:
            updated_judgement_ids.append(entry_id)
            if not remaining_source_ids:
                unanchored_judgement_ids.append(entry_id)
    _write_json(judgement_path, judgement)

    mandate_path = _case_path(case_dir, "clara_mandate")
    mandate = _read_json(mandate_path)
    mandate_source_ids, removed_mandate_source_ids = _filter_material_ids(
        mandate.get("source_material_ids", []),
        removed_id_set,
    )
    mandate["source_material_ids"] = mandate_source_ids

    removed_anchor_ids: list[str] = []
    preparation = mandate.get("preparation")
    if isinstance(preparation, dict):
        anchors = preparation.get("material_anchors", [])
        if isinstance(anchors, list):
            kept_anchors = []
            for anchor in anchors:
                if not isinstance(anchor, dict):
                    kept_anchors.append(anchor)
                    continue
                anchor_id = str(anchor.get("id", "")).strip()
                if anchor_id in removed_id_set:
                    removed_anchor_ids.append(anchor_id)
                else:
                    kept_anchors.append(anchor)
            preparation["material_anchors"] = kept_anchors

    remaining_session_paths = {
        session_path
        for material in remaining_materials
        if isinstance(material, dict)
        for path in (_material_path(case_dir, material),)
        if (session_path := _voice_session_path_for_path(case_dir, path)) is not None
    }
    removable_session_paths = removed_session_paths - remaining_session_paths
    voice_session_paths = mandate.get("voice_session_paths", [])
    removed_voice_session_paths: list[str] = []
    if isinstance(voice_session_paths, list):
        kept_voice_session_paths: list[str] = []
        for session_path in voice_session_paths:
            session_path_str = str(session_path)
            if session_path_str in removable_session_paths:
                removed_voice_session_paths.append(session_path_str)
            else:
                kept_voice_session_paths.append(session_path_str)
        mandate["voice_session_paths"] = _dedupe_preserve_order(
            kept_voice_session_paths
        )

    if removed_mandate_source_ids or removed_anchor_ids or removed_voice_session_paths:
        mandate["updated_at"] = timestamp
    _write_json(mandate_path, mandate)

    orphan_candidates: list[Path] = []
    remaining_paths = {
        str(_material_path(case_dir, item))
        for item in remaining_materials
        if isinstance(item, dict)
    }
    for path in removed_paths:
        if (
            path.exists()
            and _case_owned_path(case_dir, path)
            and str(path) not in remaining_paths
        ):
            orphan_candidates.append(path)
    for session_path in removable_session_paths:
        session_dir = (case_dir / session_path).resolve()
        if session_dir.exists():
            orphan_candidates.append(session_dir)
    orphan_candidate_paths = _dedupe_paths(orphan_candidates)

    removed_empty_dirs: list[Path] = []
    if remove_empty_orphan_dirs:
        for candidate in orphan_candidate_paths:
            if candidate.is_dir():
                try:
                    candidate.rmdir()
                except OSError:
                    continue
                removed_empty_dirs.append(candidate)

    _touch_manifest(case_dir, timestamp)
    brief_result = refresh_case_brief(case_dir, now=now)
    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    return MaterialDeleteResult(
        removed_material_ids=removed_ids,
        missing_material_ids=missing_ids,
        removed_material_paths=removed_paths,
        updated_judgement_ids=tuple(_dedupe_preserve_order(updated_judgement_ids)),
        unanchored_judgement_ids=tuple(
            _dedupe_preserve_order(unanchored_judgement_ids)
        ),
        removed_mandate_source_material_ids=tuple(
            _dedupe_preserve_order(removed_mandate_source_ids)
        ),
        removed_mandate_voice_session_paths=tuple(
            _dedupe_preserve_order(removed_voice_session_paths)
        ),
        removed_preparation_material_anchor_ids=tuple(
            _dedupe_preserve_order(removed_anchor_ids)
        ),
        orphan_candidate_paths=orphan_candidate_paths,
        removed_empty_orphan_dirs=tuple(removed_empty_dirs),
        brief_path=brief_result.brief_path,
    )


def ingest_note_text(
    case_dir: Path,
    *,
    title: str,
    text: str,
    now: datetime | None = None,
) -> dict[str, Any]:
    """Persist pasted consultant notes as a source material inside the case."""

    if not text.strip():
        raise CaseWorkspaceError("note text cannot be empty")
    timestamp = _now_iso(now)
    compact_timestamp = timestamp.replace("-", "").replace(":", "")[:15] + "Z"
    notes_dir = case_dir / "notes"
    notes_dir.mkdir(parents=True, exist_ok=True)
    note_path = notes_dir / f"{compact_timestamp}-{_slugify(title)}.md"
    note_body = (
        f"# {title}\n\n"
        f"Captured: {timestamp}\n"
        "Source: pasted consultant note\n\n"
        f"{text.strip()}\n"
    )
    note_path.write_text(note_body, encoding="utf-8")
    return register_material(
        case_dir,
        note_path,
        material_type="note",
        title=title,
        summary=_mechanical_text_summary(text),
        now=now,
    )


def _owned_note_path(
    case_dir: Path,
    *,
    title: str,
    suffix: str,
    now: datetime | None = None,
) -> Path:
    timestamp = _now_iso(now)
    compact_timestamp = timestamp.replace("-", "").replace(":", "")[:15] + "Z"
    notes_dir = case_dir / "notes"
    notes_dir.mkdir(parents=True, exist_ok=True)
    clean_suffix = suffix if suffix.startswith(".") else f".{suffix}"
    return notes_dir / f"{compact_timestamp}-{_slugify(title)}{clean_suffix}"


def ingest_note_file(
    case_dir: Path,
    *,
    title: str,
    notes_file: Path,
    now: datetime | None = None,
) -> dict[str, Any]:
    """Copy an existing note file into the case workspace and register it."""

    if not title.strip():
        raise CaseWorkspaceError("note title cannot be empty")
    if not notes_file.exists():
        raise FileNotFoundError(notes_file)

    resolved_source = notes_file.resolve()
    notes_dir = (case_dir / "notes").resolve()
    if resolved_source.is_relative_to(notes_dir):
        note_path = resolved_source
    else:
        suffix = notes_file.suffix or ".md"
        note_path = _owned_note_path(case_dir, title=title, suffix=suffix, now=now)
        if note_path.resolve() != resolved_source:
            shutil.copy2(resolved_source, note_path)

    return register_material(
        case_dir,
        note_path,
        material_type="note",
        title=title,
        now=now,
    )


def _file_sha256(path: Path) -> str:
    digest = sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _filename_tokens(path: Path) -> set[str]:
    return {
        token for token in re.split(r"[^a-zA-Z0-9]+", path.stem.casefold()) if token
    }


def _infer_case_file_kind(path: Path) -> str:
    suffix = path.suffix.lower()
    tokens = _filename_tokens(path)
    if suffix in AUDIO_FILE_SUFFIXES:
        return "audio"
    if suffix in NOTE_FILE_SUFFIXES or tokens & NOTE_NAME_TOKENS:
        return "note"
    if suffix in PRESENTATION_FILE_SUFFIXES and tokens & PRESENTATION_NAME_TOKENS:
        return "presentation"
    return "source"


def _normalize_case_file_kind(kind: str, source_path: Path) -> str:
    _validate_choice(kind, CASE_FILE_KINDS, "kind")
    if kind == "auto":
        return _infer_case_file_kind(source_path)
    if kind == "deck":
        return "presentation"
    return kind


def _case_file_destination_dir(case_dir: Path, kind: str) -> Path:
    if kind == "presentation":
        return case_dir / "outputs" / "presentations" / "current"
    if kind == "note":
        return case_dir / "notes"
    if kind == "audio":
        return case_dir / "source_materials" / "interviews" / "audio"
    if kind == "source":
        return case_dir / "source_materials" / "project_docs"
    raise CaseWorkspaceError(f"Unsupported case file kind: {kind}")


def _copy_destination_for_source(
    source_path: Path,
    destination_dir: Path,
    *,
    overwrite: bool,
) -> tuple[Path, bool]:
    destination = destination_dir / source_path.name
    if not destination.exists():
        return destination, True
    if destination.resolve() == source_path.resolve():
        return destination, False
    if _file_sha256(destination) == _file_sha256(source_path):
        return destination, False
    if overwrite:
        return destination, True

    for index in range(2, 1000):
        candidate = destination.with_name(
            f"{destination.stem}-{index}{destination.suffix}"
        )
        if not candidate.exists():
            return candidate, True
        if _file_sha256(candidate) == _file_sha256(source_path):
            return candidate, False
    raise CaseWorkspaceError(
        f"Could not choose a non-conflicting destination for {source_path.name}"
    )


def _default_normalized_pptx_path(path: Path) -> Path:
    return path.with_name(f"{path.stem}_normalized_for_merge{path.suffix}")


def copy_case_file(
    case_dir: Path,
    source_file: Path,
    *,
    kind: str = "auto",
    register: bool = False,
    title: str | None = None,
    material_type: str | None = None,
    overwrite: bool = False,
    normalize_legacy_pptx: bool = True,
    soffice_binary: Path | None = None,
    now: datetime | None = None,
) -> CaseFileCopyResult:
    """Copy a downloaded/local file into the appropriate case folder."""

    if not source_file.exists():
        raise FileNotFoundError(source_file)
    if not source_file.is_file():
        raise CaseWorkspaceError(f"case file source must be a file: {source_file}")
    resolved_source = source_file.resolve()
    resolved_kind = _normalize_case_file_kind(kind, resolved_source)
    destination_dir = _case_file_destination_dir(case_dir, resolved_kind)
    destination_dir.mkdir(parents=True, exist_ok=True)
    destination_path, should_copy = _copy_destination_for_source(
        resolved_source,
        destination_dir,
        overwrite=overwrite,
    )
    if should_copy:
        shutil.copy2(resolved_source, destination_path)

    legacy_pptx_normalization: LegacyPptxNormalizationResult | None = None
    if normalize_legacy_pptx and resolved_kind == "presentation":
        legacy_pptx_normalization = _normalize_copied_presentation_for_merge(
            destination_path,
            soffice_binary=soffice_binary,
            now=now,
        )

    registered_material: dict[str, Any] | None = None
    if register:
        default_material_type = "note" if resolved_kind == "note" else "source"
        registered_material = register_material(
            case_dir,
            destination_path,
            material_type=material_type or default_material_type,
            title=title,
            now=now,
        )

    return CaseFileCopyResult(
        source_path=resolved_source,
        destination_path=destination_path.resolve(),
        kind=resolved_kind,
        copied=should_copy,
        registered_material=registered_material,
        legacy_pptx_normalization=legacy_pptx_normalization,
    )


def _normalize_copied_presentation_for_merge(
    deck_path: Path,
    *,
    soffice_binary: Path | None = None,
    now: datetime | None = None,
) -> LegacyPptxNormalizationResult | None:
    if deck_path.suffix.lower() != ".pptx":
        return None
    try:
        legacy_media = inspect_pptx_legacy_media(deck_path)
    except BadZipFile:
        LOGGER.warning(
            "Skipping legacy PPTX normalization because %s is not a readable PPTX package.",
            deck_path,
        )
        return None
    if not legacy_media:
        return None
    output_path = _default_normalized_pptx_path(deck_path)
    return normalize_legacy_pptx_for_editable_merge(
        deck_path,
        output_path=output_path,
        overwrite=True,
        soffice_binary=soffice_binary,
        now=now,
    )


def inspect_pptx_legacy_media(path: Path) -> tuple[str, ...]:
    """Return WMF/EMF media parts embedded in a PPTX package."""

    if not path.exists():
        raise FileNotFoundError(path)
    if path.suffix.lower() != ".pptx":
        raise CaseWorkspaceError(f"legacy deck normalization expects a .pptx: {path}")

    with ZipFile(path) as package:
        legacy_media = sorted(
            name
            for name in package.namelist()
            if name.startswith("ppt/media/")
            and Path(name).suffix.lower() in LEGACY_PPTX_MEDIA_SUFFIXES
        )
    return tuple(legacy_media)


def resolve_soffice_binary(configured: Path | None = None) -> Path:
    """Resolve the LibreOffice binary used for PPTX normalization."""

    candidates: list[Path] = []
    if configured is not None:
        candidates.append(configured)
    env_value = os.environ.get("SOFFICE_BINARY", "").strip()
    if env_value:
        candidates.append(Path(env_value))
    for binary_name in ("soffice", "libreoffice"):
        resolved = shutil.which(binary_name)
        if resolved:
            candidates.append(Path(resolved))
    candidates.extend(DEFAULT_SOFFICE_PATHS)

    for candidate in candidates:
        if candidate.exists():
            return candidate.resolve()
    raise CaseWorkspaceError(
        "No LibreOffice executable found. Install LibreOffice or set SOFFICE_BINARY."
    )


def _soffice_profile_uri(path: Path) -> str:
    return path.resolve().as_uri()


def _run_soffice_pptx_roundtrip(
    source_path: Path,
    *,
    output_dir: Path,
    soffice_binary: Path,
) -> Path:
    """Round-trip a presentation through LibreOffice and return its PPTX output."""

    output_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="clara_soffice_profile_") as profile_dir:
        command = [
            str(soffice_binary),
            f"-env:UserInstallation={_soffice_profile_uri(Path(profile_dir))}",
            "--headless",
            "--invisible",
            "--nodefault",
            "--nofirststartwizard",
            "--nolockcheck",
            "--norestore",
            "--convert-to",
            "pptx",
            "--outdir",
            str(output_dir),
            str(source_path.resolve()),
        ]
        completed = subprocess.run(
            command,
            check=False,
            capture_output=True,
            text=True,
        )
    output_path = output_dir / source_path.name
    if completed.returncode != 0 or not output_path.exists():
        message = (
            completed.stderr.strip() or completed.stdout.strip() or "unknown error"
        )
        raise CaseWorkspaceError(
            f"LibreOffice failed to normalize PPTX for editable merge: {message}"
        )
    return output_path


def _ensure_custom_props_content_type(content_types_xml: bytes) -> bytes:
    ET.register_namespace("", CONTENT_TYPES_NS)
    root = ET.fromstring(content_types_xml)
    override_tag = f"{{{CONTENT_TYPES_NS}}}Override"
    has_override = any(
        child.tag == override_tag
        and child.attrib.get("PartName") == "/docProps/custom.xml"
        for child in root
    )
    if not has_override:
        ET.SubElement(
            root,
            override_tag,
            {
                "PartName": "/docProps/custom.xml",
                "ContentType": CUSTOM_PROPS_CONTENT_TYPE,
            },
        )
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _replace_zip_parts(zip_path: Path, replacements: Mapping[str, bytes]) -> None:
    temporary_path = zip_path.with_name(f".{zip_path.name}.tmp")
    with (
        ZipFile(zip_path, "r") as source_zip,
        ZipFile(temporary_path, "w", compression=ZIP_DEFLATED) as target_zip,
    ):
        for item in source_zip.infolist():
            if item.filename in replacements:
                continue
            target_zip.writestr(item, source_zip.read(item.filename))
        for name, data in replacements.items():
            target_zip.writestr(name, data)
    temporary_path.replace(zip_path)


def _preserve_custom_properties(source_path: Path, output_path: Path) -> bool:
    with ZipFile(source_path, "r") as source_zip:
        source_names = set(source_zip.namelist())
        if "docProps/custom.xml" not in source_names:
            return False
        custom_xml = source_zip.read("docProps/custom.xml")

    with ZipFile(output_path, "r") as output_zip:
        output_names = set(output_zip.namelist())
        if "[Content_Types].xml" not in output_names:
            raise CaseWorkspaceError("Normalized PPTX is missing [Content_Types].xml")
        content_types_xml = _ensure_custom_props_content_type(
            output_zip.read("[Content_Types].xml")
        )

    _replace_zip_parts(
        output_path,
        {
            "[Content_Types].xml": content_types_xml,
            "docProps/custom.xml": custom_xml,
        },
    )
    return True


def _validate_pptx_package(path: Path) -> None:
    with ZipFile(path, "r") as package:
        corrupt_member = package.testzip()
    if corrupt_member is not None:
        raise CaseWorkspaceError(
            f"Normalized PPTX package contains a corrupt member: {corrupt_member}"
        )


def normalize_legacy_pptx_for_editable_merge(
    source_path: Path,
    *,
    output_path: Path | None = None,
    report_path: Path | None = None,
    force: bool = False,
    overwrite: bool = False,
    soffice_binary: Path | None = None,
    now: datetime | None = None,
) -> LegacyPptxNormalizationResult:
    """Normalize WMF/EMF-heavy PPTX decks before editable slide merging.

    The helper is deterministic around inspection, package validation, metadata
    preservation, and reporting. LibreOffice performs only the mechanical PPTX
    round-trip needed to rewrite legacy presentation internals.
    """

    resolved_source = source_path.resolve()
    before = inspect_pptx_legacy_media(resolved_source)
    resolved_output = (
        output_path.resolve()
        if output_path is not None
        else _default_normalized_pptx_path(resolved_source)
    )
    if resolved_output == resolved_source:
        raise CaseWorkspaceError("normalized PPTX output must not overwrite the source")
    if resolved_output.exists() and not overwrite:
        raise CaseWorkspaceError(
            f"normalized PPTX output already exists: {resolved_output}"
        )
    resolved_output.parent.mkdir(parents=True, exist_ok=True)

    used_soffice: Path | None = None
    normalized = bool(before) or force
    if normalized:
        used_soffice = resolve_soffice_binary(soffice_binary)
        with tempfile.TemporaryDirectory(prefix="clara_pptx_normalize_") as temp_dir:
            roundtrip_path = _run_soffice_pptx_roundtrip(
                resolved_source,
                output_dir=Path(temp_dir),
                soffice_binary=used_soffice,
            )
            shutil.copy2(roundtrip_path, resolved_output)
        custom_properties_preserved = _preserve_custom_properties(
            resolved_source, resolved_output
        )
    else:
        shutil.copy2(resolved_source, resolved_output)
        custom_properties_preserved = _preserve_custom_properties(
            resolved_source, resolved_output
        )

    _validate_pptx_package(resolved_output)
    after = inspect_pptx_legacy_media(resolved_output)
    resolved_report = (
        report_path.resolve()
        if report_path is not None
        else resolved_output.with_suffix(".normalization_report.json")
    )
    resolved_report.parent.mkdir(parents=True, exist_ok=True)
    report = {
        "schema_version": 1,
        "source_path": str(resolved_source),
        "output_path": str(resolved_output),
        "normalized": normalized,
        "legacy_media_before": list(before),
        "legacy_media_after": list(after),
        "legacy_media_before_count": len(before),
        "legacy_media_after_count": len(after),
        "soffice_binary": str(used_soffice) if used_soffice is not None else None,
        "custom_properties_preserved": custom_properties_preserved,
        "created_at": _now_iso(now),
        "merge_guidance": (
            "Use this normalized PPTX as the base for editable slide merging."
            if not after
            else "Legacy WMF/EMF media remain; prefer image fallback for affected slides."
        ),
    }
    _write_json(resolved_report, report)

    return LegacyPptxNormalizationResult(
        source_path=resolved_source,
        output_path=resolved_output,
        report_path=resolved_report,
        normalized=normalized,
        legacy_media_before=before,
        legacy_media_after=after,
        soffice_binary=used_soffice,
    )


def prepare_editable_pptx_merge_input(
    source_path: Path,
    *,
    normalized_path: Path | None = None,
    skip_normalization_reason: str | None = None,
    report_path: Path | None = None,
    now: datetime | None = None,
) -> EditablePptxMergeInputResult:
    """Select and audit the PPTX base allowed into editable slide merging.

    Legacy WMF/EMF-heavy decks must use a normalized merge base. A caller can
    bypass that only by recording an explicit skip reason in the generated
    report; this prevents silent fallback to fragile editable merging.
    """

    resolved_source = source_path.resolve()
    source_legacy_media = inspect_pptx_legacy_media(resolved_source)
    requested_normalized_path = (
        normalized_path.resolve() if normalized_path is not None else None
    )
    discovered_normalized_path = _default_normalized_pptx_path(resolved_source)
    candidate_normalized_path = requested_normalized_path
    if candidate_normalized_path is None and discovered_normalized_path.exists():
        candidate_normalized_path = discovered_normalized_path.resolve()

    reason = (skip_normalization_reason or "").strip()
    normalized_path_for_report: Path | None = None
    if source_legacy_media:
        if candidate_normalized_path is not None:
            if not candidate_normalized_path.exists():
                raise FileNotFoundError(candidate_normalized_path)
            merge_base_path = candidate_normalized_path
            normalized_path_for_report = candidate_normalized_path
            status = "normalized_merge_base_ready"
        elif reason:
            merge_base_path = resolved_source
            status = "normalization_skipped_with_reason"
        else:
            raise CaseWorkspaceError(
                "Editable PPTX merge requires a normalized source path for "
                f"legacy WMF/EMF media in {resolved_source}. Run "
                "scripts/normalize_legacy_pptx.py or pass a "
                "skip_normalization_reason to record why normalization is not used."
            )
    else:
        if candidate_normalized_path is not None:
            if not candidate_normalized_path.exists():
                raise FileNotFoundError(candidate_normalized_path)
            merge_base_path = candidate_normalized_path
            normalized_path_for_report = candidate_normalized_path
            status = "explicit_merge_base_ready"
        else:
            merge_base_path = resolved_source
            status = "source_clean_for_merge"

    merge_base_legacy_media = inspect_pptx_legacy_media(merge_base_path)
    if source_legacy_media and status == "normalized_merge_base_ready":
        if merge_base_path == resolved_source:
            raise CaseWorkspaceError(
                "normalized merge base must be a separate file from the legacy source"
            )
        if merge_base_legacy_media and not reason:
            raise CaseWorkspaceError(
                "Normalized merge base still contains WMF/EMF media; pass a "
                "skip_normalization_reason to record why editable merge is still "
                "being attempted, or use an image fallback for affected slides."
            )
        if merge_base_legacy_media:
            status = "normalized_merge_base_has_legacy_media_with_reason"

    resolved_report = (
        report_path.resolve()
        if report_path is not None
        else merge_base_path.with_suffix(".editable_merge_input_report.json")
    )
    resolved_report.parent.mkdir(parents=True, exist_ok=True)
    report = {
        "schema_version": 1,
        "source_path": str(resolved_source),
        "merge_base_path": str(merge_base_path),
        "normalized_path": (
            str(normalized_path_for_report)
            if normalized_path_for_report is not None
            else None
        ),
        "status": status,
        "source_legacy_media": list(source_legacy_media),
        "source_legacy_media_count": len(source_legacy_media),
        "merge_base_legacy_media": list(merge_base_legacy_media),
        "merge_base_legacy_media_count": len(merge_base_legacy_media),
        "skip_normalization_reason": reason or None,
        "editable_merge_guard": (
            "normalized_required"
            if source_legacy_media
            else "normalization_not_required"
        ),
        "created_at": _now_iso(now),
    }
    _write_json(resolved_report, report)

    return EditablePptxMergeInputResult(
        source_path=resolved_source,
        merge_base_path=merge_base_path,
        report_path=resolved_report,
        status=status,
        source_legacy_media=source_legacy_media,
        merge_base_legacy_media=merge_base_legacy_media,
        normalized_path=normalized_path_for_report,
        skip_normalization_reason=reason or None,
    )


def _known_material_ids(case_dir: Path) -> set[str]:
    registry = _read_json(_case_path(case_dir, "materials"))
    return {str(item["id"]) for item in registry["materials"]}


def _known_judgement_ids(case_dir: Path) -> set[str]:
    judgement = _read_json(_case_path(case_dir, "judgement"))
    return {str(item["id"]) for item in judgement["entries"]}


def _known_open_question_ids(case_dir: Path) -> set[str]:
    payload = _read_json(_case_path(case_dir, "open_questions"))
    return {str(item["id"]) for item in payload["questions"]}


def add_judgement_entries(
    case_dir: Path,
    entries: Sequence[Mapping[str, Any]],
    *,
    now: datetime | None = None,
) -> list[dict[str, Any]]:
    """Store Codex-structured judgement entries with explicit review status."""

    timestamp = _now_iso(now)
    known_material_ids = _known_material_ids(case_dir)
    judgement_path = _case_path(case_dir, "judgement")
    judgement = _read_json(judgement_path)
    existing_ids = [item["id"] for item in judgement["entries"]]
    added: list[dict[str, Any]] = []

    for raw_entry in entries:
        kind = str(raw_entry.get("kind", ""))
        status = str(raw_entry.get("status", "pending"))
        text = str(raw_entry.get("text", "")).strip()
        _validate_choice(kind, JUDGEMENT_KINDS, "kind")
        _validate_choice(status, JUDGEMENT_STATUSES, "status")
        if not text:
            raise CaseWorkspaceError("judgement text cannot be empty")
        source_ids = list(raw_entry.get("source_material_ids", []))
        unknown_sources = sorted(set(source_ids) - known_material_ids)
        if unknown_sources:
            raise CaseWorkspaceError(
                "unknown source_material_ids: " + ", ".join(unknown_sources)
            )

        entry = {
            "id": _next_id("jud", [*existing_ids, *(item["id"] for item in added)]),
            "kind": kind,
            "text": text,
            "status": status,
            "source_material_ids": source_ids,
            "rationale": str(raw_entry.get("rationale", "")).strip(),
            "captured_at": timestamp,
            "reviewed_at": timestamp if status != "pending" else None,
            "reviewer": str(raw_entry.get("reviewer", "")).strip() or None,
            "review_note": str(raw_entry.get("review_note", "")).strip(),
        }
        judgement["entries"].append(entry)
        added.append(entry)

    _write_json(judgement_path, judgement)
    _touch_manifest(case_dir, timestamp)
    refresh_case_brief(case_dir, now=now)
    return added


def set_judgement_status(
    case_dir: Path,
    entry_id: str,
    *,
    status: str,
    reviewer: str,
    review_note: str = "",
    now: datetime | None = None,
) -> dict[str, Any]:
    """Mark one judgement entry included, excluded, or pending."""

    updated_entries = set_judgement_statuses(
        case_dir,
        [entry_id],
        status=status,
        reviewer=reviewer,
        review_note=review_note,
        now=now,
    )
    return updated_entries[0]


def set_judgement_statuses(
    case_dir: Path,
    entry_ids: Sequence[str],
    *,
    status: str,
    reviewer: str,
    review_note: str = "",
    now: datetime | None = None,
) -> list[dict[str, Any]]:
    """Mark multiple judgement entries included, excluded, or pending."""

    _validate_choice(status, JUDGEMENT_STATUSES, "status")
    if status != "pending" and not reviewer.strip():
        raise CaseWorkspaceError(
            "recorded-by advisor name is required for include/exclude status"
        )
    requested_ids = _dedupe_preserve_order(entry_ids)
    if not requested_ids:
        raise CaseWorkspaceError("at least one judgement entry id is required")

    timestamp = _now_iso(now)
    judgement_path = _case_path(case_dir, "judgement")
    judgement = _read_json(judgement_path)
    entries_by_id = {entry["id"]: entry for entry in judgement["entries"]}
    missing_ids = [
        entry_id for entry_id in requested_ids if entry_id not in entries_by_id
    ]
    if missing_ids:
        raise CaseWorkspaceError(
            "judgement entries not found: " + ", ".join(missing_ids)
        )

    updated_entries: list[dict[str, Any]] = []
    selected_ids = set(requested_ids)
    for entry in judgement["entries"]:
        if entry["id"] in selected_ids:
            entry["status"] = status
            entry["reviewed_at"] = None if status == "pending" else timestamp
            entry["reviewer"] = reviewer.strip() or None
            entry["review_note"] = review_note.strip()
            updated_entries.append(entry)
    _write_json(judgement_path, judgement)
    _touch_manifest(case_dir, timestamp)
    refresh_case_brief(case_dir, now=now)
    return sorted(
        updated_entries,
        key=lambda entry: requested_ids.index(entry["id"]),
    )


def add_open_question(
    case_dir: Path,
    *,
    question: str,
    why_it_matters: str,
    source_entry_ids: Sequence[str] = (),
    status: str = "open",
    now: datetime | None = None,
) -> dict[str, Any]:
    """Store a targeted follow-up question for the advisor or client team."""

    _validate_choice(status, OPEN_QUESTION_STATUSES, "status")
    if not question.strip():
        raise CaseWorkspaceError("question cannot be empty")
    timestamp = _now_iso(now)
    questions_path = _case_path(case_dir, "open_questions")
    payload = _read_json(questions_path)
    question_item = {
        "id": _next_id("q", [item["id"] for item in payload["questions"]]),
        "question": question.strip(),
        "why_it_matters": why_it_matters.strip(),
        "status": status,
        "source_entry_ids": list(source_entry_ids),
        "created_at": timestamp,
        "updated_at": timestamp,
    }
    payload["questions"].append(question_item)
    _write_json(questions_path, payload)
    _touch_manifest(case_dir, timestamp)
    refresh_case_brief(case_dir, now=now)
    return question_item


def _clean_id_list(raw_values: Any, *, field_name: str) -> list[str]:
    if raw_values is None:
        return []
    if not isinstance(raw_values, list):
        raise CaseWorkspaceError(f"{field_name} must be a list")
    return _dedupe_preserve_order(str(value).strip() for value in raw_values)


def _clean_issue_id(issue_id: str) -> str:
    cleaned = issue_id.strip()
    if not cleaned:
        raise CaseWorkspaceError("issue id cannot be empty")
    if not re.fullmatch(r"[A-Za-z0-9][A-Za-z0-9_-]*", cleaned):
        raise CaseWorkspaceError(
            "issue id may contain only letters, numbers, hyphen, and underscore"
        )
    return cleaned


def upsert_case_issues(
    case_dir: Path,
    issues: Sequence[Mapping[str, Any]],
    *,
    now: datetime | None = None,
) -> list[dict[str, Any]]:
    """Create or update cross-interview case issues linked to evidence entries."""

    timestamp = _now_iso(now)
    known_judgement_ids = _known_judgement_ids(case_dir)
    known_question_ids = _known_open_question_ids(case_dir)
    issues_path = _case_path(case_dir, "issues")
    payload = _read_json(issues_path)
    existing_by_id = {str(item["id"]): item for item in payload["issues"]}
    existing_ids = [str(item["id"]) for item in payload["issues"]]
    updated: list[dict[str, Any]] = []

    for raw_issue in issues:
        title = str(raw_issue.get("title", "")).strip()
        if not title and not raw_issue.get("id"):
            raise CaseWorkspaceError("issue title is required for new issues")
        issue_id = (
            _clean_issue_id(str(raw_issue["id"]))
            if raw_issue.get("id")
            else _next_id("issue", [*existing_ids, *(item["id"] for item in updated)])
        )
        status = str(raw_issue.get("status", "active"))
        _validate_choice(status, ISSUE_STATUSES, "status")

        evidence_for = _clean_id_list(
            raw_issue.get("evidence_for", []),
            field_name="evidence_for",
        )
        evidence_against = _clean_id_list(
            raw_issue.get("evidence_against", []),
            field_name="evidence_against",
        )
        open_tests = _clean_id_list(
            raw_issue.get("open_tests", []),
            field_name="open_tests",
        )
        unknown_judgement_ids = sorted(
            (set(evidence_for) | set(evidence_against)) - known_judgement_ids
        )
        if unknown_judgement_ids:
            raise CaseWorkspaceError(
                "unknown judgement entry ids: " + ", ".join(unknown_judgement_ids)
            )
        unknown_question_ids = sorted(set(open_tests) - known_question_ids)
        if unknown_question_ids:
            raise CaseWorkspaceError(
                "unknown open question ids: " + ", ".join(unknown_question_ids)
            )

        existing = existing_by_id.get(issue_id)
        issue = {
            "id": issue_id,
            "title": title or str(existing.get("title", "") if existing else ""),
            "decision_area": str(
                raw_issue.get(
                    "decision_area",
                    existing.get("decision_area", "") if existing else "",
                )
            ).strip(),
            "current_synthesis": str(
                raw_issue.get(
                    "current_synthesis",
                    existing.get("current_synthesis", "") if existing else "",
                )
            ).strip(),
            "status": status,
            "evidence_for": evidence_for,
            "evidence_against": evidence_against,
            "open_tests": open_tests,
            "created_at": (
                existing.get("created_at", timestamp) if existing else timestamp
            ),
            "updated_at": timestamp,
        }
        if existing is None:
            payload["issues"].append(issue)
        else:
            existing.update(issue)
            issue = existing
        updated.append(issue)

    _write_json(issues_path, payload)
    _touch_manifest(case_dir, timestamp)
    refresh_case_brief(case_dir, now=now)
    return updated


def _plugin_root() -> Path:
    return Path(__file__).resolve().parents[1]


def _load_clara_playbook() -> dict[str, Any]:
    playbook_path = _plugin_root() / "playbook" / "clara_kickoff_playbook.json"
    if playbook_path.exists():
        payload = json.loads(playbook_path.read_text(encoding="utf-8"))
        return payload if isinstance(payload, dict) else {}
    return {}


def _fallback_clara_playbook() -> dict[str, Any]:
    return {
        "succession_lenses": [
            "Ownership transfer and economic rights",
            "Control, vetoes, and real decision rights",
            "Operating leadership and successor credibility",
            "Founder role after transition",
            "Family alignment and politically acceptable narrative",
            "Governance mechanisms that make agreement operational",
            "Tax, legal, and financial constraints",
            "Industry-specific capabilities needed by the next leader",
        ],
        "red_flags": [
            "A successor is named but authority is not actually transferred.",
            "Family agreement is stated but not translated into governance mechanisms.",
            "A tax or legal structure solves wealth transfer but not leadership continuity.",
            "The founder's future role is emotionally or operationally ambiguous.",
            "The case contains a private preference that is being disguised as analysis.",
        ],
        "baseline_sources": [
            {
                "title": "HBS Working Knowledge - Governing the Family-Run Business",
                "url": "https://www.library.hbs.edu/working-knowledge/governing-the-family-run-business",
                "takeaway": "Separate family, ownership, and business governance roles.",
            },
            {
                "title": "HBS Working Knowledge - Developing Your Next CEO for the Family Business",
                "url": "https://www.library.hbs.edu/working-knowledge/managing-the-family-business-developing-your-next-ceo",
                "takeaway": "CEO succession in family companies requires both capability and family-system fit.",
            },
            {
                "title": "IFC Family Business Governance Handbook",
                "url": "https://www.ifc.org/en/insights-reports/2011/ifc-family-business-governance-handbook",
                "takeaway": "Family-business governance needs explicit structures, roles, and policies.",
            },
        ],
    }


def _clara_playbook_value(key: str) -> list[Any]:
    playbook = _load_clara_playbook() or _fallback_clara_playbook()
    value = playbook.get(key)
    fallback = _fallback_clara_playbook()[key]
    return value if isinstance(value, list) else fallback


def _clean_string_list(values: Any) -> list[str]:
    if isinstance(values, str):
        values = [values]
    if not isinstance(values, Sequence):
        return []
    cleaned: list[str] = []
    for value in values:
        clean = " ".join(str(value).split())
        if clean:
            cleaned.append(clean)
    return _dedupe_preserve_order(cleaned)


def _clean_research_items(values: Sequence[Mapping[str, Any]]) -> list[dict[str, str]]:
    cleaned: list[dict[str, str]] = []
    for value in values:
        if not isinstance(value, Mapping):
            continue
        title = str(value.get("title", "")).strip()
        url = str(value.get("url", "")).strip()
        takeaway = str(value.get("takeaway", "")).strip()
        if title or url or takeaway:
            cleaned.append({"title": title, "url": url, "takeaway": takeaway})
    return cleaned


def _material_anchors(materials: Sequence[Mapping[str, Any]]) -> list[dict[str, str]]:
    anchors: list[dict[str, str]] = []
    for material in materials:
        anchors.append(
            {
                "id": str(material["id"]),
                "title": str(material["title"]),
                "material_type": str(material["material_type"]),
                "summary": str(material.get("summary", "")).strip(),
            }
        )
    return anchors


def _render_clara_kickoff_preparation_markdown(
    *,
    manifest: Mapping[str, Any],
    material_anchors: Sequence[Mapping[str, str]],
    succession_lenses: Sequence[str],
    red_flags: Sequence[str],
    industry_context: Sequence[str],
    external_research: Sequence[Mapping[str, str]],
    generated_at: str,
) -> str:
    language = _ui_language(manifest)
    if language == "es":
        copy = {
            "title": "Preparación del kickoff de Clara",
            "intro": (
                "Preparación interna para el kickoff con el socio. Clara ha leído "
                "los resúmenes de los materiales del caso y el playbook controlado. "
                "Toda investigación web debe registrarse mediante enlaces a las "
                "fuentes y conclusiones, sin copiar el texto de las fuentes."
            ),
            "generated": "Generado",
            "snapshot": "Resumen del caso",
            "project": "Proyecto",
            "objective": "Objetivo",
            "audience": "Destinatario",
            "output_language": "Idioma de salida",
            "materials": "Materiales de referencia",
            "no_materials": "Aún no hay materiales del caso indexados.",
            "lenses": "Perspectivas para la sucesión",
            "red_flags": "Señales de alerta que vigilar",
            "industry": "Contexto sectorial",
            "no_industry": (
                "Aún no se ha registrado investigación sectorial específica. "
                "Clara debe añadir implicaciones sectoriales concisas y respaldadas "
                "por fuentes antes de basarse en afirmaciones sobre el sector."
            ),
            "research": "Notas de investigación externa",
            "source": "Fuente",
            "no_takeaway": "No se ha registrado ninguna conclusión.",
            "no_research": "No se han registrado notas de investigación externa.",
            "posture": "Enfoque para el kickoff",
            "posture_items": (
                "El socio informa a Clara; Clara no dirige un interrogatorio.",
                "Clara formula únicamente las preguntas de aclaración imprescindibles "
                "cuando falta un dato que impide comprender el caso.",
                "Clara debe inferir de la explicación del socio si el caso es "
                "analítico, político, orientado primero a la presentación o mixto, "
                "sin preguntar directamente por esa clasificación.",
                "Clara no debe presentar el contenido del playbook como una "
                "conclusión. El playbook es una perspectiva para preparar y dar "
                "seguimiento al trabajo.",
            ),
        }
    else:
        copy = {
            "title": "Clara Kickoff Preparation",
            "intro": (
                "Internal preparation for the partner kickoff. Clara has read the "
                "case material summaries and the controlled playbook. Any web "
                "research should be recorded as source links and takeaways, not "
                "copied source text."
            ),
            "generated": "Generated",
            "snapshot": "Case Snapshot",
            "project": "Project",
            "objective": "Objective",
            "audience": "Audience",
            "output_language": "Output language",
            "materials": "Material Anchors",
            "no_materials": "No case materials indexed yet.",
            "lenses": "Succession Lenses",
            "red_flags": "Red Flags To Watch",
            "industry": "Industry Context",
            "no_industry": (
                "No industry-specific research has been recorded yet. Clara should "
                "add concise, source-backed industry implications before relying on "
                "industry claims."
            ),
            "research": "External Research Notes",
            "source": "Source",
            "no_takeaway": "No takeaway recorded.",
            "no_research": "No external research notes recorded.",
            "posture": "Kickoff Posture",
            "posture_items": (
                "The partner briefs Clara; Clara does not run an interrogation.",
                "Clara asks only essential clarification questions when a missing "
                "point blocks understanding.",
                "Clara should infer whether the case is analytical, political, "
                "presentation-first, or mixed from the partner's explanation rather "
                "than asking that classification directly.",
                "Clara should not present playbook material as a conclusion. The "
                "playbook is a lens for preparation and follow-up.",
            ),
        }
    localized_lenses = _localized_default_items(succession_lenses, language)
    localized_red_flags = _localized_default_items(red_flags, language)
    localized_research = _localized_default_research_items(
        external_research,
        language,
    )
    lines = [
        f"# {copy['title']} - {manifest['client']}",
        "",
        copy["intro"],
        "",
        f"{copy['generated']}: {generated_at}",
        "",
        f"## {copy['snapshot']}",
        "",
        f"- {copy['project']}: {manifest['project']}",
        f"- {copy['objective']}: {manifest['objective']}",
        f"- {copy['audience']}: {manifest['audience']}",
        f"- {copy['output_language']}: {manifest['output_language']}",
        "",
        f"## {copy['materials']}",
        "",
    ]
    if material_anchors:
        lines.extend(
            f"- `{item['id']}` {item['title']} [{item['material_type']}]: {item['summary']}"
            for item in material_anchors
        )
    else:
        lines.append(f"- {copy['no_materials']}")
    lines.extend(["", f"## {copy['lenses']}", ""])
    lines.extend(f"- {item}" for item in localized_lenses)
    lines.extend(["", f"## {copy['red_flags']}", ""])
    lines.extend(f"- {item}" for item in localized_red_flags)
    lines.extend(["", f"## {copy['industry']}", ""])
    if industry_context:
        lines.extend(f"- {item}" for item in industry_context)
    else:
        lines.append(f"- {copy['no_industry']}")
    lines.extend(["", f"## {copy['research']}", ""])
    if localized_research:
        for item in localized_research:
            label = item.get("title") or item.get("url") or copy["source"]
            takeaway = item.get("takeaway") or copy["no_takeaway"]
            url = item.get("url", "")
            source = f" ({url})" if url else ""
            lines.append(f"- {label}{source}: {takeaway}")
    else:
        lines.append(f"- {copy['no_research']}")
    lines.extend(
        [
            "",
            f"## {copy['posture']}",
            "",
            *(f"- {item}" for item in copy["posture_items"]),
            "",
        ]
    )
    return "\n".join(lines)


def prepare_clara_kickoff(
    case_dir: Path,
    *,
    industry_context: Sequence[str] = (),
    external_research: Sequence[Mapping[str, Any]] = (),
    now: datetime | None = None,
) -> ClaraKickoffPreparationResult:
    """Prepare Clara's first partner kickoff from case state and playbook.

    This is deterministic preparation. Codex may browse public or authorized
    sources before calling it, but the helper itself only records supplied
    source takeaways and controlled playbook lenses.
    """

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    timestamp = _now_iso(now)
    manifest = _read_json(_case_path(case_dir, "manifest"))
    materials = _read_json(_case_path(case_dir, "materials"))["materials"]
    mandate_path = _case_path(case_dir, "clara_mandate")
    mandate = _read_json(mandate_path)
    succession_lenses = _clean_string_list(_clara_playbook_value("succession_lenses"))
    red_flags = _clean_string_list(_clara_playbook_value("red_flags"))
    research_items = _clean_research_items(
        external_research or _clara_playbook_value("baseline_sources")
    )
    material_items = _material_anchors(materials)
    industry_items = _clean_string_list(industry_context)
    language = _ui_language(manifest)
    preparation = {
        "case_snapshot": {
            "client": manifest["client"],
            "project": manifest["project"],
            "objective": manifest["objective"],
            "audience": manifest["audience"],
            "output_language": manifest["output_language"],
        },
        "material_anchors": material_items,
        "succession_lenses": succession_lenses,
        "red_flags": red_flags,
        "industry_context": industry_items,
        "external_research": research_items,
        "preparation_note": (
            "Clara está preparada para el briefing del socio. El socio debe explicar "
            "el caso con naturalidad; Clara debe escuchar y formular únicamente las "
            "preguntas de aclaración imprescindibles."
            if language == "es"
            else (
                "Clara is prepared for a partner briefing. The partner should "
                "explain the case naturally; Clara should listen and ask only "
                "essential clarifications."
            )
        ),
    }
    mandate["status"] = "prepared"
    mandate["prepared_at"] = timestamp
    mandate["updated_at"] = timestamp
    mandate["preparation"] = preparation
    mandate["source_material_ids"] = [item["id"] for item in materials]
    _write_json(mandate_path, mandate)
    _touch_manifest(case_dir, timestamp)

    preparation_path = case_dir / CLARA_KICKOFF_PREPARATION_FILENAME
    preparation_path.write_text(
        _render_clara_kickoff_preparation_markdown(
            manifest=manifest,
            material_anchors=material_items,
            succession_lenses=succession_lenses,
            red_flags=red_flags,
            industry_context=industry_items,
            external_research=research_items,
            generated_at=timestamp,
        ),
        encoding="utf-8",
    )
    refresh_case_brief(case_dir, now=now)
    return ClaraKickoffPreparationResult(
        mandate_path=mandate_path,
        preparation_path=preparation_path,
        material_count=len(materials),
        baseline_source_count=len(research_items),
    )


def _merge_clean_strings(existing: Sequence[str], incoming: Any) -> list[str]:
    return _dedupe_preserve_order([*existing, *_clean_string_list(incoming)])


def update_clara_mandate_from_kickoff(
    case_dir: Path,
    kickoff_payload: Mapping[str, Any],
    *,
    material_id: str = "",
    session_path: str = "",
    now: datetime | None = None,
) -> dict[str, Any]:
    """Update Clara's mandate after a partner kickoff briefing."""

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))
    timestamp = _now_iso(now)
    mandate_path = _case_path(case_dir, "clara_mandate")
    mandate = _read_json(mandate_path)
    extracted = kickoff_payload.get("clara_mandate", kickoff_payload)
    if not isinstance(extracted, Mapping):
        extracted = {}
    body = mandate["mandate"]

    for field in (
        "engagement_objective",
        "client_decision",
        "clara_understanding",
        "partner_starting_orientation",
    ):
        value = str(extracted.get(field, "")).strip()
        if value:
            body[field] = value
    for field in (
        "sensitive_points",
        "what_clara_should_investigate",
        "what_clara_should_not_waste_time_on",
        "essential_clarifications",
        "next_steps",
    ):
        body[field] = _merge_clean_strings(
            body.get(field, []), extracted.get(field, [])
        )

    mandate["status"] = "kickoff_imported"
    mandate["kickoff_imported_at"] = timestamp
    mandate["updated_at"] = timestamp
    if material_id:
        mandate["source_material_ids"] = _dedupe_preserve_order(
            [*mandate.get("source_material_ids", []), material_id]
        )
    if session_path:
        mandate["voice_session_paths"] = _dedupe_preserve_order(
            [*mandate.get("voice_session_paths", []), session_path]
        )
    mandate["mandate"] = body
    _write_json(mandate_path, mandate)
    _touch_manifest(case_dir, timestamp)
    refresh_case_brief(case_dir, now=now)
    return mandate


def _html_escape(value: Any) -> str:
    return html.escape(str(value), quote=True)


def _ui_language(manifest: Mapping[str, Any]) -> str:
    language = str(manifest.get("output_language", "en")).strip().lower()
    return language if language in SUPPORTED_LANGUAGES else "en"


def _html_copy(language: str) -> dict[str, str]:
    if language == "it":
        return {
            "already_requested_next_steps": "Prossimi passi gia' richiesti",
            "approved_entries": "Voci gia' approvate per il pack cliente",
            "approved_metric": "voci approvate per il pack cliente",
            "brief_eyebrow": "Briefing partner",
            "brief_lead": (
                "Scheda interna per distinguere contenuti pronti, ipotesi "
                "candidate e prossime verifiche."
            ),
            "candidate_count": "voci candidate",
            "candidate_empty": "Nessun giudizio candidato registrato.",
            "candidate_judgement": "Giudizio candidato",
            "case_state": "Stato del caso",
            "clarification_empty": "Nessun chiarimento essenziale registrato.",
            "client_decision": "Decisione cliente",
            "control_note": "Nota di controllo",
            "control_note_body": (
                "Materiale di lavoro locale. Il giudizio professionale resta al "
                "partner senior."
            ),
            "deck_eyebrow": "Kickoff",
            "deck_note": "Materiale di lavoro interno.",
            "direction": "Direzione",
            "direction_item_include": "Quali voci candidate includere, escludere o correggere.",
            "direction_item_output": "Quale output il partner vuole come prossimo passo.",
            "engagement_objective": "Obiettivo dell'incarico",
            "essential_clarifications": "Chiarimenti essenziali",
            "evidence_gaps": "Punti da chiarire",
            "executive_storyline": "Traccia esecutiva",
            "facts_empty": "Nessun fatto acquisito registrato.",
            "facts_title": "Fatti acquisiti",
            "hypotheses_empty": "Nessuna ipotesi candidata registrata.",
            "hypotheses_heading": "Ipotesi candidate da confermare o correggere.",
            "hypotheses_eyebrow": "Ipotesi da verificare",
            "initial_readout": "Sintesi iniziale",
            "initial_readout_heading": "Fatti e orientamento emersi.",
            "investigation": "Indagine",
            "investigation_empty": "Nessun punto di indagine registrato.",
            "investigation_prompt": "Quali evidenze recuperare per prime.",
            "investigation_title": "Verifiche necessarie",
            "industry_context": "Contesto settoriale",
            "industry_context_empty": "Nessun contesto settoriale registrato.",
            "judgement": "Giudizio",
            "judgement_item_client": "Cosa puo' essere usato con il cliente.",
            "judgement_item_right": "Cosa e' gia' corretto.",
            "judgement_item_sensitive": "Cosa e' politicamente sensibile.",
            "known_facts": "Fatti acquisiti",
            "lenses_empty": "Nessuna lente di preparazione registrata.",
            "mandate": "Mandato",
            "materials_metric": "materiali indicizzati",
            "next_conversation": "Prossima conversazione",
            "next_conversation_heading": "Decisioni per il prossimo passaggio.",
            "next_step_empty": "Nessun prossimo passo registrato.",
            "next_steps": "Prossimi passi",
            "open_questions_metric": "domande aperte",
            "partner_orientation": "Orientamento del partner",
            "partner_orientation_empty": (
                "Orientamento del partner non ancora registrato."
            ),
            "partner_orientation_waiting": (
                "Orientamento del partner non ancora registrato; partire dalle "
                "voci candidate."
            ),
            "pending_metric": "voci di giudizio pending",
            "pending_note": (
                "Voci approvate per il pack cliente: {approved_count}. "
                "Il giudizio pending resta materiale di lavoro."
            ),
            "preparation_lenses": "Controlli di preparazione",
            "preparation_lenses_heading": "Punti da testare nella prima conversazione.",
            "priority_checks": "Verifiche prioritarie",
            "questions_empty": "Nessuna domanda aperta registrata.",
            "questions_heading": (
                "Domande da risolvere prima che qualcosa diventi materiale cliente."
            ),
            "red_flags": "Punti di attenzione",
            "red_flags_empty": "Nessun punto di attenzione registrato.",
            "research_empty": "Nessuna ricerca esterna registrata.",
            "research_notes": "Note di ricerca esterna",
            "ready_judgement": "Contenuto pronto per il pack cliente",
            "ready_judgement_empty": "Nessun contenuto approvato registrato.",
            "sensitive_empty": "Nessun punto sensibile registrato.",
            "sensitive_points": "Punti sensibili",
            "sensitive_default_gate": (
                "Il giudizio candidato non diventa materiale cliente finche' "
                "il partner non lo approva."
            ),
            "sensitive_default_paths": (
                "L'output cliente non deve esporre percorsi locali o meccaniche interne."
            ),
            "what_clara_needs": "Decisioni per il prossimo passaggio.",
        }
    if language == "es":
        return {
            "already_requested_next_steps": "Próximos pasos ya solicitados",
            "approved_entries": "Elementos aprobados para el paquete del cliente",
            "approved_metric": "elementos listos para el paquete del cliente",
            "brief_eyebrow": "Briefing del socio",
            "brief_lead": (
                "Documento interno para separar contenido listo, hipótesis "
                "candidatas y próximas comprobaciones."
            ),
            "candidate_count": "elementos candidatos",
            "candidate_empty": "Aún no se ha registrado ningún juicio candidato.",
            "candidate_judgement": "Juicio candidato",
            "case_state": "Estado del caso",
            "clarification_empty": "Aún no se han registrado aclaraciones esenciales.",
            "client_decision": "Decisión del cliente",
            "control_note": "Nota de control",
            "control_note_body": (
                "Material de trabajo local. El juicio profesional corresponde al "
                "socio sénior."
            ),
            "deck_eyebrow": "Kickoff",
            "deck_note": "Material de trabajo interno.",
            "direction": "Dirección",
            "direction_item_include": (
                "Qué elementos candidatos incluir, excluir o corregir."
            ),
            "direction_item_output": "Qué entregable quiere el socio a continuación.",
            "engagement_objective": "Objetivo del encargo",
            "essential_clarifications": "Aclaraciones esenciales",
            "evidence_gaps": "Lagunas de evidencia",
            "executive_storyline": "Línea argumental ejecutiva",
            "facts_empty": "Aún no se han registrado hechos.",
            "facts_title": "Hechos conocidos",
            "hypotheses_empty": "Aún no se han registrado hipótesis candidatas.",
            "hypotheses_heading": "Hipótesis candidatas que confirmar o corregir.",
            "hypotheses_eyebrow": "Hipótesis que comprobar",
            "initial_readout": "Síntesis inicial",
            "initial_readout_heading": "Hechos y orientación recogidos hasta ahora.",
            "investigation": "Investigación",
            "investigation_empty": "Aún no se han registrado puntos de investigación.",
            "investigation_prompt": "Qué evidencia obtener primero.",
            "investigation_title": "Comprobaciones necesarias",
            "industry_context": "Contexto sectorial",
            "industry_context_empty": "Aún no se ha registrado contexto sectorial.",
            "judgement": "Juicio",
            "judgement_item_client": "Qué puede utilizarse con el cliente.",
            "judgement_item_right": "Qué ya es correcto.",
            "judgement_item_sensitive": "Qué es políticamente sensible.",
            "known_facts": "Hechos conocidos",
            "lenses_empty": "No se han registrado perspectivas de preparación.",
            "mandate": "Mandato",
            "materials_metric": "materiales indexados",
            "next_conversation": "Próxima conversación",
            "next_conversation_heading": "Decisiones para el siguiente paso.",
            "next_step_empty": "Aún no se han registrado próximos pasos.",
            "next_steps": "Próximos pasos",
            "open_questions_metric": "preguntas abiertas",
            "partner_orientation": "Orientación del socio",
            "partner_orientation_empty": "Aún no se ha registrado la orientación del socio.",
            "partner_orientation_waiting": (
                "Aún no se ha registrado la orientación del socio; empezar por "
                "los elementos candidatos."
            ),
            "pending_metric": "elementos de juicio pendientes",
            "pending_note": (
                "Elementos aprobados para el paquete del cliente: {approved_count}. "
                "El juicio pendiente es solo material de trabajo."
            ),
            "preparation_lenses": "Comprobaciones de preparación",
            "preparation_lenses_heading": (
                "Puntos que comprobar en la primera conversación con el socio."
            ),
            "priority_checks": "Comprobaciones prioritarias",
            "questions_empty": "Aún no se han registrado preguntas abiertas.",
            "questions_heading": (
                "Preguntas que Clara debe resolver antes de que el contenido esté "
                "listo para el cliente."
            ),
            "red_flags": "Señales de alerta",
            "red_flags_empty": "No se han registrado señales de alerta.",
            "research_empty": "No se ha registrado investigación externa.",
            "research_notes": "Notas de investigación externa",
            "ready_judgement": "Contenido listo para el paquete del cliente",
            "ready_judgement_empty": "No se ha registrado contenido aprobado.",
            "sensitive_empty": "No se han registrado puntos sensibles.",
            "sensitive_points": "Puntos sensibles",
            "sensitive_default_gate": (
                "El juicio candidato no está listo para el cliente hasta que el "
                "socio lo marque como tal."
            ),
            "sensitive_default_paths": (
                "El entregable para el cliente no debe mostrar rutas de archivos "
                "locales ni mecanismos internos de trabajo."
            ),
            "what_clara_needs": "Decisiones para el siguiente paso.",
        }
    return {
        "already_requested_next_steps": "Already requested next steps",
        "approved_entries": "Approved client-pack entries",
        "approved_metric": "client-pack ready items",
        "brief_eyebrow": "Partner Brief",
        "brief_lead": (
            "Internal brief to separate ready content, candidate hypotheses, "
            "and the next checks."
        ),
        "candidate_count": "candidate items",
        "candidate_empty": "No candidate judgement captured yet.",
        "candidate_judgement": "Candidate Judgement",
        "case_state": "Case State",
        "clarification_empty": "No essential clarifications recorded yet.",
        "client_decision": "Client decision",
        "control_note": "Control Note",
        "control_note_body": (
            "Local working material. Professional judgement remains with the "
            "senior partner."
        ),
        "deck_eyebrow": "Kickoff",
        "deck_note": "Internal working material.",
        "direction": "Direction",
        "direction_item_include": "Which candidate items to include, exclude, or correct.",
        "direction_item_output": "What output the partner wants next.",
        "engagement_objective": "Engagement objective",
        "essential_clarifications": "Essential Clarifications",
        "evidence_gaps": "Evidence gaps",
        "executive_storyline": "Executive Storyline",
        "facts_empty": "No facts captured yet.",
        "facts_title": "Known facts",
        "hypotheses_empty": "No candidate hypotheses captured yet.",
        "hypotheses_heading": "Candidate hypotheses to confirm or correct.",
        "hypotheses_eyebrow": "Hypotheses to test",
        "initial_readout": "Initial synthesis",
        "initial_readout_heading": "Facts and orientation captured so far.",
        "investigation": "Investigation",
        "investigation_empty": "No investigation items recorded yet.",
        "investigation_prompt": "Which evidence to get first.",
        "investigation_title": "Required checks",
        "industry_context": "Industry Context",
        "industry_context_empty": "No industry context recorded yet.",
        "judgement": "Judgement",
        "judgement_item_client": "What can be used with the client.",
        "judgement_item_right": "What is already right.",
        "judgement_item_sensitive": "What is politically sensitive.",
        "known_facts": "Known facts",
        "lenses_empty": "No succession lenses recorded.",
        "mandate": "Mandate",
        "materials_metric": "materials indexed",
        "next_conversation": "Next conversation",
        "next_conversation_heading": "Decisions for the next step.",
        "next_step_empty": "No next steps recorded yet.",
        "next_steps": "Next Steps",
        "open_questions_metric": "open questions",
        "partner_orientation": "Partner orientation",
        "partner_orientation_empty": ("Partner orientation not recorded yet."),
        "partner_orientation_waiting": (
            "Partner orientation not recorded yet; start from the candidate items."
        ),
        "pending_metric": "pending judgement entries",
        "pending_note": (
            "Approved client-pack entries: {approved_count}. "
            "Pending judgement is working material only."
        ),
        "preparation_lenses": "Preparation Checks",
        "preparation_lenses_heading": "Points to test in the first partner conversation.",
        "priority_checks": "Priority Checks",
        "questions_empty": "No open questions registered yet.",
        "questions_heading": (
            "Questions Clara should resolve before anything becomes client-ready."
        ),
        "red_flags": "Red flags",
        "red_flags_empty": "No red flags recorded.",
        "research_empty": "No external research recorded.",
        "research_notes": "External Research Notes",
        "ready_judgement": "Client-Pack Ready Content",
        "ready_judgement_empty": "No approved content recorded yet.",
        "sensitive_empty": "No sensitive points recorded yet.",
        "sensitive_points": "Sensitive Points",
        "sensitive_default_gate": (
            "Candidate judgement is not client-ready until the partner marks it ready."
        ),
        "sensitive_default_paths": (
            "Client-facing output should not expose local file paths or internal "
            "workpaper mechanics."
        ),
        "what_clara_needs": "Decisions for the next step.",
    }


def _same_text_sequence(left: Sequence[str], right: Sequence[str]) -> bool:
    return _clean_string_list(left) == _clean_string_list(right)


ITALIAN_DEFAULT_TEXT = {
    "Ownership transfer and economic rights": "Trasferimento proprieta' e diritti economici",
    "Control, vetoes, and real decision rights": "Controllo, veto e poteri decisionali reali",
    "Operating leadership and successor credibility": (
        "Leadership operativa e credibilita' del successore"
    ),
    "Founder role after transition": "Ruolo del fondatore dopo la transizione",
    "Family alignment and politically acceptable narrative": (
        "Allineamento familiare e narrativa politicamente accettabile"
    ),
    "Governance mechanisms that make agreement operational": (
        "Meccanismi di governance che rendono operativo l'accordo"
    ),
    "Tax, legal, and financial constraints": "Vincoli fiscali, legali e finanziari",
    "Industry-specific capabilities needed by the next leader": (
        "Capacita' settoriali richieste al prossimo leader"
    ),
    "A successor is named but authority is not actually transferred.": (
        "Viene nominato un successore, ma l'autorita' non viene realmente trasferita."
    ),
    "Family agreement is stated but not translated into governance mechanisms.": (
        "L'accordo familiare e' dichiarato, ma non tradotto in meccanismi di governance."
    ),
    "A tax or legal structure solves wealth transfer but not leadership continuity.": (
        "Una struttura fiscale o legale risolve il trasferimento patrimoniale, "
        "ma non la continuita' della leadership."
    ),
    "The founder's future role is emotionally or operationally ambiguous.": (
        "Il ruolo futuro del fondatore resta ambiguo sul piano emotivo o operativo."
    ),
    "The case contains a private preference that is being disguised as analysis.": (
        "Il caso contiene una preferenza privata mascherata da analisi."
    ),
    "Customer, supplier, or employee trust is personal to the founder and not institutionalized.": (
        "La fiducia di clienti, fornitori o persone chiave e' personale al fondatore "
        "e non istituzionalizzata."
    ),
    "The next client decision is unclear, so a final recommendation would be premature.": (
        "La prossima decisione del cliente non e' chiara, quindi una raccomandazione "
        "finale sarebbe prematura."
    ),
    "Separate family, ownership, and business governance roles.": (
        "Separare ruoli di famiglia, proprieta' e governo dell'impresa."
    ),
    "CEO succession in family companies requires both capability and family-system fit.": (
        "La successione dell'AD nelle aziende familiari richiede capacita' e coerenza "
        "con il sistema familiare."
    ),
    "Family-business negotiation must protect both the economic decision and the family relationship.": (
        "La negoziazione familiare deve proteggere sia la decisione economica sia "
        "la relazione familiare."
    ),
    "Family-business governance needs explicit structures, roles, and policies.": (
        "La governance dell'azienda familiare richiede strutture, ruoli e policy espliciti."
    ),
}

SPANISH_DEFAULT_TEXT = {
    "Ownership transfer and economic rights": (
        "Transferencia de la propiedad y derechos económicos"
    ),
    "Control, vetoes, and real decision rights": (
        "Control, vetos y derechos de decisión reales"
    ),
    "Operating leadership and successor credibility": (
        "Liderazgo operativo y credibilidad del sucesor"
    ),
    "Founder role after transition": "Papel del fundador tras la transición",
    "Family alignment and politically acceptable narrative": (
        "Alineación familiar y narrativa políticamente aceptable"
    ),
    "Governance mechanisms that make agreement operational": (
        "Mecanismos de gobernanza que hacen operativo el acuerdo"
    ),
    "Tax, legal, and financial constraints": (
        "Restricciones fiscales, jurídicas y financieras"
    ),
    "Industry-specific capabilities needed by the next leader": (
        "Capacidades sectoriales que necesita el próximo líder"
    ),
    "A successor is named but authority is not actually transferred.": (
        "Se nombra a un sucesor, pero la autoridad no se transfiere realmente."
    ),
    "Family agreement is stated but not translated into governance mechanisms.": (
        "Se declara un acuerdo familiar, pero no se traduce en mecanismos de gobernanza."
    ),
    "A tax or legal structure solves wealth transfer but not leadership continuity.": (
        "Una estructura fiscal o jurídica resuelve la transferencia patrimonial, "
        "pero no la continuidad del liderazgo."
    ),
    "The founder's future role is emotionally or operationally ambiguous.": (
        "El futuro papel del fundador es ambiguo en términos emocionales u operativos."
    ),
    "The case contains a private preference that is being disguised as analysis.": (
        "El caso contiene una preferencia privada disfrazada de análisis."
    ),
    "Customer, supplier, or employee trust is personal to the founder and not institutionalized.": (
        "La confianza de clientes, proveedores o empleados depende personalmente "
        "del fundador y no está institucionalizada."
    ),
    "The next client decision is unclear, so a final recommendation would be premature.": (
        "La próxima decisión del cliente no está clara, por lo que una "
        "recomendación final sería prematura."
    ),
    "Separate family, ownership, and business governance roles.": (
        "Separar las funciones de familia, propiedad y gobierno empresarial."
    ),
    "CEO succession in family companies requires both capability and family-system fit.": (
        "La sucesión del consejero delegado en empresas familiares exige tanto "
        "capacidad como encaje en el sistema familiar."
    ),
    "Family-business negotiation must protect both the economic decision and the family relationship.": (
        "La negociación en la empresa familiar debe proteger tanto la decisión "
        "económica como la relación familiar."
    ),
    "Family-business governance needs explicit structures, roles, and policies.": (
        "La gobernanza de la empresa familiar necesita estructuras, funciones y "
        "políticas explícitas."
    ),
}


def _localized_default_text(value: str, language: str) -> str:
    if language == "it":
        return ITALIAN_DEFAULT_TEXT.get(value, value)
    if language == "es":
        return SPANISH_DEFAULT_TEXT.get(value, value)
    return value


def _localized_default_items(items: Any, language: str) -> list[str]:
    return [
        _localized_default_text(item, language) for item in _clean_string_list(items)
    ]


def _localized_default_research_items(
    items: Sequence[Mapping[str, Any]], language: str
) -> list[dict[str, str]]:
    localized: list[dict[str, str]] = []
    for item in _clean_research_items(items):
        localized.append(
            {
                "title": item["title"],
                "url": item["url"],
                "takeaway": _localized_default_text(item["takeaway"], language),
            }
        )
    return localized


def _default_partner_next_steps(language: str) -> list[str]:
    if language == "it":
        return [
            "Aprire il kickoff deck e reagire alle ipotesi iniziali.",
            "Decidere quali voci candidate includere, escludere, correggere o espandere.",
            "Registrare orientamento del partner e priorita' del primo output.",
        ]
    if language == "es":
        return [
            "Abrir el deck de kickoff y reaccionar a las hipótesis iniciales.",
            "Decidir qué elementos candidatos incluir, excluir, corregir o ampliar.",
            "Registrar la orientación del socio y la prioridad del primer entregable.",
        ]
    return [
        "Open the kickoff deck and react to the initial hypotheses.",
        "Decide which candidate items to include, exclude, correct, or expand.",
        "Record the partner's orientation and first output priority.",
    ]


def _html_list(items: Sequence[str], *, empty: str) -> str:
    cleaned = _clean_string_list(items)
    if not cleaned:
        return f'<p class="empty">{_html_escape(empty)}</p>'
    return (
        "<ul>" + "".join(f"<li>{_html_escape(item)}</li>" for item in cleaned) + "</ul>"
    )


def _html_research_list(
    items: Sequence[Mapping[str, Any]], *, empty: str = "No external research recorded."
) -> str:
    cleaned = _clean_research_items(items)
    if not cleaned:
        return f'<p class="empty">{_html_escape(empty)}</p>'
    rows: list[str] = []
    for item in cleaned:
        title = _html_escape(item.get("title") or "Source")
        takeaway = _html_escape(item.get("takeaway") or "No takeaway recorded.")
        url = str(item.get("url", "")).strip()
        source = f'<a href="{_html_escape(url)}">{_html_escape(url)}</a>' if url else ""
        rows.append(
            "<li>"
            f"<strong>{title}</strong>"
            f"<span>{takeaway}</span>"
            f"{source}"
            "</li>"
        )
    return '<ul class="research-list">' + "".join(rows) + "</ul>"


def _html_compact_list(items: Sequence[str], *, limit: int, empty: str) -> str:
    return _html_list(list(items)[:limit], empty=empty)


def _entry_texts_by_status(
    entries: Sequence[Mapping[str, Any]],
    *,
    statuses: set[str],
    kinds: set[str] | None = None,
) -> list[str]:
    selected: list[str] = []
    for entry in entries:
        if entry.get("status") not in statuses:
            continue
        if kinds is not None and entry.get("kind") not in kinds:
            continue
        selected.append(str(entry.get("text", "")).strip())
    return [item for item in selected if item]


def _open_question_texts(questions: Sequence[Mapping[str, Any]]) -> list[str]:
    return [
        f"{item['question']} ({item['why_it_matters']})"
        for item in questions
        if item.get("status") == "open"
    ]


def _strip_local_paths(value: str) -> str:
    return re.sub(r"/Users/[^\s,)`]+", "[percorso locale]", value)


def _partner_approved_content(entries: Sequence[Mapping[str, Any]]) -> list[str]:
    return _entry_texts_by_status(
        entries,
        statuses={"approved"},
        kinds={"fact", "advisor_judgement", "codex_inference", "decision_implication"},
    )


def _indexed_industry_context_items(
    materials: Sequence[Mapping[str, Any]],
) -> list[str]:
    keywords = (
        "industry",
        "market",
        "mercato",
        "settore",
        "research",
        "ricerca",
        "premium",
        "successione",
        "family business",
    )
    items: list[str] = []
    for material in materials:
        haystack = " ".join(
            [
                str(material.get("title", "")),
                str(material.get("path", "")),
                str(material.get("material_type", "")),
            ]
        ).lower()
        if not any(keyword in haystack for keyword in keywords):
            continue
        summary = str(material.get("summary", "")).strip()
        if not summary or summary in {
            "No readable preview text.",
            "PDF indexed; semantic review remains in Codex.",
        }:
            continue
        title = str(material.get("title", "Source")).strip() or "Source"
        clean_summary = _strip_local_paths(_mechanical_text_summary(summary, limit=300))
        items.append(f"{title}: {clean_summary}")
    return _dedupe_preserve_order(items)[:5]


def render_clara_kickoff_deck(
    case_dir: Path,
    *,
    output_path: Path | None = None,
) -> ClaraKickoffDeckResult:
    """Render Clara's first quiet partner kickoff deck from case state."""

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))
    manifest = _read_json(_case_path(case_dir, "manifest"))
    language = _ui_language(manifest)
    copy = _html_copy(language)
    materials = _read_json(_case_path(case_dir, "materials"))["materials"]
    entries = _read_json(_case_path(case_dir, "judgement"))["entries"]
    questions = _read_json(_case_path(case_dir, "open_questions"))["questions"]
    mandate = _read_json(_case_path(case_dir, "clara_mandate"))
    preparation = mandate.get("preparation", {})
    body = mandate.get("mandate", {})
    target_path = output_path or case_dir / CLARA_KICKOFF_DECK_FILENAME
    target_path.parent.mkdir(parents=True, exist_ok=True)

    facts = _entry_texts_by_status(
        entries, statuses={"approved", "pending"}, kinds={"fact"}
    )
    hypotheses = _entry_texts_by_status(
        entries,
        statuses={"approved", "pending"},
        kinds={"advisor_judgement", "codex_inference", "decision_implication"},
    )
    approved_count = sum(1 for entry in entries if entry["status"] == "approved")
    pending_count = sum(1 for entry in entries if entry["status"] == "pending")
    open_question_items = _open_question_texts(questions)
    material_count = len(materials)
    client_decision = (
        str(body.get("client_decision", "")).strip()
        or str(manifest.get("objective", "")).strip()
    )
    partner_orientation = str(body.get("partner_starting_orientation", "")).strip()
    if not partner_orientation:
        partner_orientation = copy["partner_orientation_empty"]
    investigation_items = _clean_string_list(
        body.get("what_clara_should_investigate", [])
    )
    next_step_items = _clean_string_list(body.get("next_steps", []))
    preparation_lenses = _localized_default_items(
        preparation.get("succession_lenses", []), language
    )
    preparation_red_flags = _localized_default_items(
        preparation.get("red_flags", []), language
    )

    html_text = f"""<!doctype html>
<html lang="{_html_escape(manifest['output_language'])}">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{_html_escape(copy['deck_eyebrow'])} - {_html_escape(manifest['client'])}</title>
  <style>
    :root {{
      --paper: #ffffff;
      --ink: #243027;
      --text: #465148;
      --muted: #6d776f;
      --line: #dfe6dd;
      --soft: #f6f8f5;
      --accent: #496a60;
      font-family: "Segoe UI", -apple-system, BlinkMacSystemFont, sans-serif;
    }}
    * {{ box-sizing: border-box; }}
    body {{ margin: 0; background: var(--paper); color: var(--text); line-height: 1.5; }}
    main {{ width: min(1120px, calc(100% - 40px)); margin: 0 auto; padding: 42px 0 64px; }}
    section {{ min-height: min(720px, calc(100vh - 44px)); padding: 40px 0; border-bottom: 1px solid var(--line); display: grid; align-content: center; gap: 24px; }}
    h1, h2, h3 {{ margin: 0; color: var(--ink); line-height: 1.12; letter-spacing: 0; }}
    h1 {{ max-width: 840px; font-size: 2.65rem; font-weight: 620; }}
    h2 {{ max-width: 880px; font-size: 2rem; font-weight: 620; }}
    h3 {{ font-size: 1rem; font-weight: 650; }}
    p {{ margin: 0; }}
    .eyebrow {{ color: var(--accent); font-size: 12px; font-weight: 700; letter-spacing: .08em; text-transform: uppercase; }}
    .lead {{ max-width: 820px; font-size: 1.08rem; }}
    .meta {{ display: flex; flex-wrap: wrap; gap: 8px; color: var(--muted); }}
    .meta span {{ border: 1px solid var(--line); border-radius: 999px; padding: 6px 10px; }}
    .grid {{ display: grid; grid-template-columns: repeat(3, minmax(0, 1fr)); gap: 14px; }}
    .grid.two {{ grid-template-columns: repeat(2, minmax(0, 1fr)); }}
    .panel {{ border: 1px solid var(--line); border-radius: 8px; padding: 18px; background: var(--soft); }}
    .metric strong {{ display: block; color: var(--ink); font-size: 2rem; }}
    ul {{ margin: 12px 0 0; padding-left: 20px; }}
    li {{ margin: 7px 0; }}
    .note {{ border-left: 4px solid var(--accent); background: var(--soft); padding: 14px 16px; }}
    @media (max-width: 820px) {{
      main {{ width: auto; margin: 0 20px; }}
      section {{ min-height: auto; }}
      h1 {{ font-size: 2rem; }}
      h2 {{ font-size: 1.5rem; }}
      .grid, .grid.two {{ grid-template-columns: minmax(0, 1fr); }}
    }}
  </style>
</head>
<body>
<main>
  <section>
    <div>
      <p class="eyebrow">{_html_escape(copy['deck_eyebrow'])}</p>
      <h1>{_html_escape(manifest['client'])}</h1>
    </div>
    <p class="lead">{_html_escape(client_decision)}</p>
    <div class="meta">
      <span>{_html_escape(material_count)} {_html_escape(copy['materials_metric'])}</span>
      <span>{_html_escape(pending_count)} {_html_escape(copy['candidate_count'])}</span>
      <span>{_html_escape(approved_count)} {_html_escape(copy['approved_metric'])}</span>
    </div>
    <p class="note">{_html_escape(copy['deck_note'])}</p>
  </section>

  <section>
    <div>
      <p class="eyebrow">{_html_escape(copy['initial_readout'])}</p>
      <h2>{_html_escape(copy['initial_readout_heading'])}</h2>
    </div>
    <div class="grid two">
      <div class="panel"><h3>{_html_escape(copy['known_facts'])}</h3>{_html_compact_list(facts, limit=4, empty=copy['facts_empty'])}</div>
      <div class="panel"><h3>{_html_escape(copy['partner_orientation'])}</h3><p>{_html_escape(partner_orientation)}</p></div>
    </div>
  </section>

  <section>
    <div>
      <p class="eyebrow">{_html_escape(copy['hypotheses_eyebrow'])}</p>
      <h2>{_html_escape(copy['hypotheses_heading'])}</h2>
    </div>
    {_html_compact_list(hypotheses, limit=7, empty=copy['hypotheses_empty'])}
  </section>

  <section>
    <div>
      <p class="eyebrow">{_html_escape(copy['evidence_gaps'])}</p>
      <h2>{_html_escape(copy['questions_heading'])}</h2>
    </div>
    {_html_compact_list(open_question_items, limit=7, empty=copy['questions_empty'])}
  </section>

  <section>
    <div>
      <p class="eyebrow">{_html_escape(copy['preparation_lenses'])}</p>
      <h2>{_html_escape(copy['preparation_lenses_heading'])}</h2>
    </div>
    <div class="grid two">
      <div class="panel"><h3>{_html_escape(copy['preparation_lenses'])}</h3>{_html_compact_list(preparation_lenses, limit=5, empty=copy['lenses_empty'])}</div>
      <div class="panel"><h3>{_html_escape(copy['red_flags'])}</h3>{_html_compact_list(preparation_red_flags, limit=5, empty=copy['red_flags_empty'])}</div>
    </div>
  </section>

  <section>
    <div>
      <p class="eyebrow">{_html_escape(copy['next_conversation'])}</p>
      <h2>{_html_escape(copy['next_conversation_heading'])}</h2>
    </div>
    <div class="grid two">
      <div class="panel"><h3>{_html_escape(copy['judgement'])}</h3><ul><li>{_html_escape(copy['judgement_item_right'])}</li><li>{_html_escape(copy['judgement_item_sensitive'])}</li><li>{_html_escape(copy['judgement_item_client'])}</li></ul></div>
      <div class="panel"><h3>{_html_escape(copy['investigation'])}</h3>{_html_compact_list(investigation_items, limit=5, empty=copy['investigation_prompt'])}</div>
      <div class="panel"><h3>{_html_escape(copy['direction'])}</h3><ul><li>{_html_escape(copy['direction_item_include'])}</li><li>{_html_escape(copy['direction_item_output'])}</li></ul></div>
      <div class="panel"><h3>{_html_escape(copy['already_requested_next_steps'])}</h3>{_html_compact_list(next_step_items, limit=5, empty=copy['next_step_empty'])}</div>
    </div>
  </section>
</main>
</body>
</html>
"""
    try:
        html_text = apply_fixed_16_9_deck_runtime(html_text)
        assert_fixed_16_9_deck_runtime(html_text, label=CLARA_KICKOFF_DECK_FILENAME)
    except ValueError as exc:
        raise CaseWorkspaceError(str(exc)) from exc
    _assert_human_visible_document_quality(html_text, label=CLARA_KICKOFF_DECK_FILENAME)
    target_path.write_text(html_text, encoding="utf-8")
    return ClaraKickoffDeckResult(
        html_path=target_path,
        hypothesis_count=len(hypotheses),
        open_question_count=len(open_question_items),
    )


def render_clara_partner_brief(
    case_dir: Path,
    *,
    output_path: Path | None = None,
) -> ClaraPartnerBriefResult:
    """Render Clara's first local HTML brief for the senior partner."""

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))
    manifest = _read_json(_case_path(case_dir, "manifest"))
    language = _ui_language(manifest)
    copy = _html_copy(language)
    materials = _read_json(_case_path(case_dir, "materials"))["materials"]
    entries = _read_json(_case_path(case_dir, "judgement"))["entries"]
    questions = _read_json(_case_path(case_dir, "open_questions"))["questions"]
    mandate = _read_json(_case_path(case_dir, "clara_mandate"))
    body = mandate["mandate"]
    preparation = mandate["preparation"]
    pending_count = sum(1 for entry in entries if entry["status"] == "pending")
    approved_count = sum(1 for entry in entries if entry["status"] == "approved")
    open_question_count = sum(1 for item in questions if item["status"] == "open")
    candidate_items = _entry_texts_by_status(
        entries,
        statuses={"pending"},
        kinds={"advisor_judgement", "codex_inference", "decision_implication"},
    )
    approved_partner_items = _partner_approved_content(entries)
    judgement_section_title = copy["candidate_judgement"]
    judgement_empty = copy["candidate_empty"]
    displayed_judgement_items = candidate_items
    if not candidate_items and approved_partner_items:
        judgement_section_title = copy["ready_judgement"]
        judgement_empty = copy["ready_judgement_empty"]
        displayed_judgement_items = approved_partner_items
    open_question_items = _open_question_texts(questions)
    client_decision_text = (
        str(body.get("client_decision", "")).strip()
        or str(manifest.get("objective", "")).strip()
    )
    partner_orientation_text = str(body.get("partner_starting_orientation", "")).strip()
    if not partner_orientation_text:
        partner_orientation_text = copy["partner_orientation_waiting"]
    investigation_items = (
        _clean_string_list(body.get("what_clara_should_investigate", []))
        or open_question_items[:5]
    )
    clarification_items = (
        _clean_string_list(body.get("essential_clarifications", []))
        or open_question_items[:5]
    )
    next_step_items = _clean_string_list(
        body.get("next_steps", [])
    ) or _default_partner_next_steps(language)
    sensitive_items = _clean_string_list(body.get("sensitive_points", [])) or [
        copy["sensitive_default_gate"],
        copy["sensitive_default_paths"],
    ]
    preparation_lenses = _localized_default_items(
        preparation.get("succession_lenses", []), language
    )
    research_items = _localized_default_research_items(
        preparation.get("external_research", []), language
    )
    industry_context_items = _clean_string_list(
        preparation.get("industry_context", [])
    ) or _indexed_industry_context_items(materials)
    target_path = output_path or case_dir / CLARA_PARTNER_BRIEF_FILENAME
    target_path.parent.mkdir(parents=True, exist_ok=True)
    if _same_text_sequence(investigation_items, clarification_items):
        investigation_sections = f"""
  <section>
    <h2>{_html_escape(copy['priority_checks'])}</h2>
    {_html_compact_list(investigation_items, limit=7, empty=copy['investigation_empty'])}
  </section>
"""
    else:
        investigation_sections = f"""
  <section>
    <h2>{_html_escape(copy['investigation_title'])}</h2>
    {_html_compact_list(investigation_items, limit=7, empty=copy['investigation_empty'])}
  </section>

  <section>
    <h2>{_html_escape(copy['essential_clarifications'])}</h2>
    {_html_compact_list(clarification_items, limit=7, empty=copy['clarification_empty'])}
  </section>
"""

    html_text = f"""<!doctype html>
<html lang="{_html_escape(manifest['output_language'])}">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{_html_escape(copy['brief_eyebrow'])} - {_html_escape(manifest['client'])}</title>
  <style>
    :root {{
      --paper: #ffffff;
      --ink: #243027;
      --text: #465148;
      --muted: #6d776f;
      --line: #dfe6dd;
      --soft: #f6f8f5;
      --accent: #496a60;
      font-family: "Segoe UI", -apple-system, BlinkMacSystemFont, sans-serif;
    }}
    * {{ box-sizing: border-box; }}
    body {{ margin: 0; color: var(--text); background: var(--paper); line-height: 1.55; }}
    main {{ width: min(1080px, calc(100% - 40px)); margin: 0 auto; padding: 48px 0 64px; }}
    header {{ padding-bottom: 34px; border-bottom: 1px solid var(--line); }}
    .eyebrow {{ margin: 0 0 12px; color: var(--accent); font-size: 12px; font-weight: 700; letter-spacing: .08em; text-transform: uppercase; }}
    h1, h2, h3 {{ margin: 0; color: var(--ink); letter-spacing: 0; line-height: 1.12; }}
    h1 {{ max-width: 820px; font-size: 2.7rem; font-weight: 620; }}
    h2 {{ font-size: 1.55rem; font-weight: 620; }}
    h3 {{ font-size: 1rem; font-weight: 650; }}
    p {{ margin: 0; }}
    .lead {{ max-width: 780px; margin-top: 18px; font-size: 1.05rem; }}
    section {{ padding: 34px 0; border-bottom: 1px solid var(--line); }}
    .grid {{ display: grid; grid-template-columns: repeat(3, minmax(0, 1fr)); gap: 14px; margin-top: 18px; }}
    .metric, .panel {{ border: 1px solid var(--line); border-radius: 8px; padding: 18px; background: #fff; }}
    .metric strong {{ display: block; color: var(--ink); font-size: 1.8rem; line-height: 1; }}
    .metric span, .empty {{ color: var(--muted); }}
    .panel {{ background: var(--soft); }}
    ul {{ margin: 14px 0 0; padding-left: 20px; }}
    li {{ margin: 7px 0; }}
    .research-list {{ display: grid; gap: 10px; padding-left: 0; list-style: none; }}
    .research-list li {{ display: grid; gap: 4px; padding-bottom: 10px; border-bottom: 1px solid var(--line); }}
    a {{ color: var(--accent); overflow-wrap: anywhere; }}
    .note {{ margin-top: 16px; padding: 14px 16px; border-left: 4px solid var(--accent); background: var(--soft); }}
    @media (max-width: 760px) {{
      main {{ width: auto; margin: 0 20px; }}
      h1 {{ font-size: 2rem; }}
      .grid {{ grid-template-columns: minmax(0, 1fr); }}
    }}
  </style>
</head>
<body>
<main>
  <header>
    <p class="eyebrow">{_html_escape(copy['brief_eyebrow'])}</p>
    <h1>{_html_escape(manifest['client'])}</h1>
    <p class="lead">{_html_escape(body.get('clara_understanding') or copy['brief_lead'])}</p>
  </header>

  <section>
    <h2>{_html_escape(copy['case_state'])}</h2>
    <div class="grid">
      <div class="metric"><strong>{len(materials)}</strong><span>{_html_escape(copy['materials_metric'])}</span></div>
      <div class="metric"><strong>{pending_count}</strong><span>{_html_escape(copy['pending_metric'])}</span></div>
      <div class="metric"><strong>{open_question_count}</strong><span>{_html_escape(copy['open_questions_metric'])}</span></div>
    </div>
    <p class="note">{_html_escape(copy['pending_note'].format(approved_count=approved_count))}</p>
  </section>

  <section>
    <h2>{_html_escape(copy['mandate'])}</h2>
    <div class="grid">
      <div class="panel"><h3>{_html_escape(copy['engagement_objective'])}</h3><p>{_html_escape(body.get('engagement_objective') or manifest['objective'])}</p></div>
      <div class="panel"><h3>{_html_escape(copy['client_decision'])}</h3><p>{_html_escape(client_decision_text)}</p></div>
      <div class="panel"><h3>{_html_escape(copy['partner_orientation'])}</h3><p>{_html_escape(partner_orientation_text)}</p></div>
    </div>
  </section>

  <section>
    <h2>{_html_escape(judgement_section_title)}</h2>
    {_html_compact_list(displayed_judgement_items, limit=8, empty=judgement_empty)}
  </section>

  <section>
    <h2>{_html_escape(copy['sensitive_points'])}</h2>
    {_html_compact_list(sensitive_items, limit=5, empty=copy['sensitive_empty'])}
  </section>

{investigation_sections}

  <section>
    <h2>{_html_escape(copy['next_steps'])}</h2>
    {_html_compact_list(next_step_items, limit=5, empty=copy['next_step_empty'])}
  </section>

  <section>
    <h2>{_html_escape(copy['preparation_lenses'])}</h2>
    {_html_list(preparation_lenses, empty=copy['lenses_empty'])}
  </section>

  <section>
    <h2>{_html_escape(copy['industry_context'])}</h2>
    {_html_list(industry_context_items, empty=copy['industry_context_empty'])}
  </section>

  <section>
    <h2>{_html_escape(copy['research_notes'])}</h2>
    {_html_research_list(research_items, empty=copy['research_empty'])}
  </section>

  <section>
    <h2>{_html_escape(copy['control_note'])}</h2>
    <p>{_html_escape(copy['control_note_body'])}</p>
  </section>
</main>
</body>
</html>
"""
    _assert_human_visible_document_quality(
        html_text, label=CLARA_PARTNER_BRIEF_FILENAME
    )
    target_path.write_text(html_text, encoding="utf-8")
    return ClaraPartnerBriefResult(
        html_path=target_path,
        open_clarification_count=len(clarification_items),
        next_step_count=len(next_step_items),
    )


def _case_fingerprint(manifest: Mapping[str, Any]) -> str:
    stable_fields = {
        key: manifest.get(key, "")
        for key in (
            "client",
            "project",
            "objective",
            "audience",
            "created_at",
        )
    }
    raw = json.dumps(stable_fields, sort_keys=True, ensure_ascii=True)
    return sha256(raw.encode("utf-8")).hexdigest()[:16]


def _record_origin_key(
    record: Mapping[str, Any],
    *,
    source_case_id: str,
    record_type: str,
) -> str:
    existing = str(record.get("origin_key", "")).strip()
    if existing:
        return existing
    return f"{source_case_id}:{record_type}:{record['id']}"


def _case_relative_file(case_dir: Path, path_value: str) -> Path | None:
    path = Path(path_value).expanduser()
    if not path.is_absolute():
        path = case_dir / path
    try:
        resolved = path.resolve()
    except OSError:
        return None
    case_root = case_dir.resolve()
    if not resolved.is_file() or not resolved.is_relative_to(case_root):
        return None
    relative = resolved.relative_to(case_root)
    if relative.name in CASE_FILES.values():
        return None
    return relative


def _exchange_id(source_case_id: str, exported_at: str, exporter: str) -> str:
    raw = f"{source_case_id}|{exported_at}|{exporter}"
    return sha256(raw.encode("utf-8")).hexdigest()[:16]


def _exchange_package_name(exchange_id: str) -> str:
    return f"case-update-{exchange_id}.zip"


def _workspace_archive_name(case_dir: Path, exported_at: str) -> str:
    return f"{_slugify(case_dir.name, fallback='case')}-workspace-{_compact_timestamp(exported_at)}.zip"


def _support_package_name(case_dir: Path, exported_at: str) -> str:
    return f"{_slugify(case_dir.name, fallback='case')}-support-{_compact_timestamp(exported_at)}.zip"


def _is_runtime_dependency_dir(name: str) -> bool:
    return name.startswith(".codex_") and name.endswith("_py")


def _share_exclusion_reason(path: Path) -> str:
    name = path.name
    if name in SHARE_EXCLUDED_FILE_NAMES:
        return "local metadata file"
    if name.startswith("._"):
        return "macOS resource fork"
    if name in SHARE_EXCLUDED_DIR_NAMES:
        return "local runtime or cache directory"
    if _is_runtime_dependency_dir(name):
        return "local plugin dependency directory"
    return ""


def _path_file_count_and_size(path: Path) -> tuple[int, int]:
    try:
        if path.is_symlink():
            return 1, path.lstat().st_size
        if path.is_file():
            return 1, path.stat().st_size
        if not path.is_dir():
            return 0, 0
    except OSError:
        return 0, 0
    file_count = 0
    byte_count = 0
    for item in path.rglob("*"):
        try:
            if item.is_file() or item.is_symlink():
                file_count += 1
                byte_count += item.lstat().st_size
        except OSError:
            continue
    return file_count, byte_count


def _iter_share_archive_files(
    case_dir: Path,
    package_path: Path,
) -> tuple[list[Path], int, int]:
    included_files: list[Path] = []
    excluded_file_count = 0
    excluded_bytes = 0
    package_resolved = package_path.resolve()
    pending = [case_dir]
    while pending:
        current = pending.pop()
        for child in sorted(current.iterdir(), key=lambda item: item.name.lower()):
            try:
                if child.resolve() == package_resolved:
                    count, size = _path_file_count_and_size(child)
                    excluded_file_count += count
                    excluded_bytes += size
                    continue
            except OSError:
                continue
            reason = _share_exclusion_reason(child)
            if reason:
                count, size = _path_file_count_and_size(child)
                excluded_file_count += count
                excluded_bytes += size
                LOGGER.debug("Skipping %s from case share archive: %s", child, reason)
                continue
            if child.is_symlink():
                count, size = _path_file_count_and_size(child)
                excluded_file_count += count
                excluded_bytes += size
                LOGGER.debug("Skipping symlink from case share archive: %s", child)
            elif child.is_dir():
                pending.append(child)
            elif child.is_file():
                included_files.append(child)
    return sorted(included_files), excluded_file_count, excluded_bytes


def export_case_workspace_archive(
    case_dir: Path,
    *,
    package_path: Path | None = None,
    now: datetime | None = None,
) -> CaseWorkspaceArchiveResult:
    """Export a clean full-workspace ZIP that excludes local runtimes and caches."""

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    timestamp = _now_iso(now)
    target_path = package_path
    if target_path is None:
        target_dir = case_dir.parent / "case_share_exports"
        target_dir.mkdir(parents=True, exist_ok=True)
        target_path = target_dir / _workspace_archive_name(case_dir, timestamp)
    else:
        target_path.parent.mkdir(parents=True, exist_ok=True)

    files, excluded_file_count, excluded_bytes = _iter_share_archive_files(
        case_dir,
        target_path,
    )
    root_name = case_dir.name.strip() or "case-workspace"
    with ZipFile(target_path, "w", compression=ZIP_DEFLATED) as archive:
        for path in files:
            relative_path = path.relative_to(case_dir)
            archive.write(path, f"{root_name}/{relative_path.as_posix()}")

    return CaseWorkspaceArchiveResult(
        package_path=target_path,
        included_file_count=len(files),
        excluded_file_count=excluded_file_count,
        excluded_bytes=excluded_bytes,
    )


def _support_markdown_value(value: Any, *, missing: str = "Not recorded") -> str:
    text = str(value or "").strip()
    return text or missing


def _support_request_block(
    request: str,
    *,
    empty: str = "No request text provided.",
) -> str:
    lines = request.strip().splitlines() or [empty]
    return "\n".join(f"> {line}" if line.strip() else ">" for line in lines)


def _render_support_request_markdown(
    *,
    manifest: Mapping[str, Any],
    request: str,
    requested_by: str,
    recipient: str,
    created_at: str,
    included_paths: Sequence[str],
    excluded_file_count: int,
    excluded_bytes: int,
) -> str:
    language = _ui_language(manifest)
    if language == "es":
        copy = {
            "title": "Solicitud de apoyo para {recipient}",
            "created": "Creada el",
            "requested_by": "Solicitada por",
            "recipient": "Destinatario",
            "client": "Cliente",
            "project": "Proyecto",
            "objective": "Objetivo",
            "audience": "Destinatario del trabajo",
            "language": "Idioma",
            "needed": "Qué se necesita",
            "included_state": "Estado de Clara incluido",
            "included": "incluido",
            "not_present": "no presente",
            "generated_only": "generado únicamente para este paquete",
            "exclusions": "Exclusiones locales",
            "excluded_count": "Número de archivos locales o de ejecución excluidos",
            "excluded_bytes": "Bytes locales o de ejecución excluidos",
            "exclusion_note": (
                "No se incluyen las carpetas ocultas de dependencias de OCR o de "
                "ejecución, como `.codex_*_py`, los entornos virtuales, las cachés, "
                "los metadatos de macOS ni las exportaciones de intercambio anteriores."
            ),
            "operating_note": "Nota operativa",
            "operating_body": (
                "La carpeta local de Clara de quien realiza la solicitud sigue "
                "siendo la copia de trabajo autoritativa. El destinatario puede "
                "utilizar Clara, Pro, edición manual o cambios del plugin para "
                "desbloquear la entrega. El material devuelto que resulte útil debe "
                "registrarse de nuevo en el caso cuando sea pertinente."
            ),
            "missing": "No registrado",
            "empty_request": "No se ha facilitado el texto de la solicitud.",
            "default_recipient": "Revisor de apoyo",
        }
    else:
        copy = {
            "title": "Support Request For {recipient}",
            "created": "Created at",
            "requested_by": "Requested by",
            "recipient": "Recipient",
            "client": "Client",
            "project": "Project",
            "objective": "Objective",
            "audience": "Audience",
            "language": "Language",
            "needed": "What Is Needed",
            "included_state": "Included Clara State",
            "included": "included",
            "not_present": "not present",
            "generated_only": "generated for this package only",
            "exclusions": "Local Exclusions",
            "excluded_count": "Excluded local/runtime file count",
            "excluded_bytes": "Excluded local/runtime bytes",
            "exclusion_note": (
                "Hidden OCR/runtime dependency folders such as `.codex_*_py`, "
                "virtual environments, caches, macOS metadata, and prior exchange "
                "exports are not included."
            ),
            "operating_note": "Operating Note",
            "operating_body": (
                "The requester's local Clara folder remains the authoritative "
                "working copy. The recipient may use Clara, Pro, manual editing, "
                "or plugin changes to unblock delivery. Useful returned material "
                "should be registered back into the case when it matters."
            ),
            "missing": "Not recorded",
            "empty_request": "No request text provided.",
            "default_recipient": "Support reviewer",
        }
    tracked_artifacts = (
        "case_manifest.json",
        "case_brief.md",
        "clara_mandate.json",
        "clara_kickoff_preparation.md",
        "clara_partner_brief.html",
        "material_registry.json",
        "judgement_log.json",
        "open_questions.json",
        "case_issues.json",
        "voice_sessions/",
        "outputs/",
    )
    included_set = set(included_paths)
    included_artifact_lines = []
    for artifact in tracked_artifacts:
        if artifact.endswith("/"):
            status = (
                copy["included"]
                if any(path.startswith(artifact) for path in included_set)
                else copy["not_present"]
            )
        else:
            status = (
                copy["included"] if artifact in included_set else copy["not_present"]
            )
        included_artifact_lines.append(f"- `{artifact}`: {status}")

    request_block = _support_request_block(request, empty=copy["empty_request"])
    requested_by_value = _support_markdown_value(
        requested_by,
        missing=copy["missing"],
    )
    recipient_value = _support_markdown_value(recipient, missing=copy["missing"])
    if language == "es" and recipient_value == "Support reviewer":
        recipient_value = copy["default_recipient"]
    return (
        f"# {copy['title'].format(recipient=recipient_value)}\n\n"
        f"- {copy['created']}: {created_at}\n"
        f"- {copy['requested_by']}: {requested_by_value}\n"
        f"- {copy['recipient']}: {recipient_value}\n"
        f"- {copy['client']}: {_support_markdown_value(manifest.get('client'), missing=copy['missing'])}\n"
        f"- {copy['project']}: {_support_markdown_value(manifest.get('project'), missing=copy['missing'])}\n"
        f"- {copy['objective']}: {_support_markdown_value(manifest.get('objective'), missing=copy['missing'])}\n"
        f"- {copy['audience']}: {_support_markdown_value(manifest.get('audience'), missing=copy['missing'])}\n"
        f"- {copy['language']}: {_support_markdown_value(manifest.get('output_language'), missing=copy['missing'])}\n\n"
        f"## {copy['needed']}\n\n"
        f"{request_block}\n\n"
        f"## {copy['included_state']}\n\n" + "\n".join(included_artifact_lines) + "\n"
        f"- `{SUPPORT_REQUEST_FILENAME}`: {copy['generated_only']}\n\n"
        f"## {copy['exclusions']}\n\n"
        f"- {copy['excluded_count']}: {excluded_file_count}\n"
        f"- {copy['excluded_bytes']}: {excluded_bytes}\n"
        f"- {copy['exclusion_note']}\n\n"
        f"## {copy['operating_note']}\n\n"
        f"{copy['operating_body']}\n"
    )


def prepare_support_package(
    case_dir: Path,
    *,
    request: str,
    requested_by: str = "",
    recipient: str = "Support reviewer",
    package_path: Path | None = None,
    now: datetime | None = None,
) -> CaseSupportPackageResult:
    """Create a clean support ZIP for delivery escalation."""

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    clean_request = request.strip()
    if not clean_request:
        raise CaseWorkspaceError("support request cannot be empty")

    timestamp = _now_iso(now)
    target_path = package_path
    if target_path is None:
        target_dir = case_dir.parent / "case_support_exports"
        target_dir.mkdir(parents=True, exist_ok=True)
        target_path = target_dir / _support_package_name(case_dir, timestamp)
    else:
        target_path.parent.mkdir(parents=True, exist_ok=True)

    files, excluded_file_count, excluded_bytes = _iter_share_archive_files(
        case_dir,
        target_path,
    )
    root_name = case_dir.name.strip() or "case-workspace"
    support_archive_path = f"{root_name}/{SUPPORT_REQUEST_FILENAME}"
    included_paths: list[str] = []
    with ZipFile(target_path, "w", compression=ZIP_DEFLATED) as archive:
        for path in files:
            relative_path = path.relative_to(case_dir)
            relative_name = relative_path.as_posix()
            if relative_name == SUPPORT_REQUEST_FILENAME:
                continue
            archive.write(path, f"{root_name}/{relative_name}")
            included_paths.append(relative_name)
        manifest = _read_json(_case_path(case_dir, "manifest"))
        archive.writestr(
            support_archive_path,
            _render_support_request_markdown(
                manifest=manifest,
                request=clean_request,
                requested_by=requested_by,
                recipient=recipient,
                created_at=timestamp,
                included_paths=included_paths,
                excluded_file_count=excluded_file_count,
                excluded_bytes=excluded_bytes,
            ),
        )

    return CaseSupportPackageResult(
        package_path=target_path,
        support_request_archive_path=support_archive_path,
        included_file_count=len(included_paths) + 1,
        excluded_file_count=excluded_file_count,
        excluded_bytes=excluded_bytes,
    )


def export_case_update(
    case_dir: Path,
    *,
    package_path: Path | None = None,
    exporter: str = "",
    now: datetime | None = None,
) -> CaseExchangeExportResult:
    """Export a deterministic ZIP package for local-first case exchange."""

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    timestamp = _now_iso(now)
    manifest = _read_json(_case_path(case_dir, "manifest"))
    materials = _read_json(_case_path(case_dir, "materials"))["materials"]
    entries = _read_json(_case_path(case_dir, "judgement"))["entries"]
    questions = _read_json(_case_path(case_dir, "open_questions"))["questions"]
    source_case_id = str(manifest.get("case_id") or _case_fingerprint(manifest))
    exchange_id = _exchange_id(source_case_id, timestamp, exporter)

    included_files: list[dict[str, str]] = []
    for material in materials:
        relative_path = _case_relative_file(case_dir, str(material.get("path", "")))
        if relative_path is None:
            continue
        archive_path = f"files/{relative_path.as_posix()}"
        included_files.append(
            {
                "material_id": str(material["id"]),
                "relative_path": relative_path.as_posix(),
                "archive_path": archive_path,
            }
        )

    update_payload = {
        "schema_version": EXCHANGE_SCHEMA_VERSION,
        "source": EXCHANGE_SOURCE,
        "exchange_id": exchange_id,
        "source_case_id": source_case_id,
        "exported_at": timestamp,
        "exporter": exporter,
        "case_manifest": manifest,
        "materials": [
            {
                **material,
                "origin_key": _record_origin_key(
                    material,
                    source_case_id=source_case_id,
                    record_type="material",
                ),
            }
            for material in materials
        ],
        "judgement_entries": [
            {
                **entry,
                "origin_key": _record_origin_key(
                    entry,
                    source_case_id=source_case_id,
                    record_type="judgement",
                ),
            }
            for entry in entries
        ],
        "open_questions": [
            {
                **question,
                "origin_key": _record_origin_key(
                    question,
                    source_case_id=source_case_id,
                    record_type="open_question",
                ),
            }
            for question in questions
        ],
        "included_files": included_files,
    }

    target_path = package_path
    if target_path is None:
        target_dir = case_dir / "exchange_exports"
        target_dir.mkdir(parents=True, exist_ok=True)
        target_path = target_dir / _exchange_package_name(exchange_id)
    else:
        target_path.parent.mkdir(parents=True, exist_ok=True)

    with ZipFile(target_path, "w", compression=ZIP_DEFLATED) as archive:
        archive.writestr(
            "case_update.json",
            json.dumps(update_payload, indent=2, sort_keys=True, ensure_ascii=True)
            + "\n",
        )
        for file_item in included_files:
            archive.write(
                case_dir / file_item["relative_path"],
                file_item["archive_path"],
            )

    return CaseExchangeExportResult(
        package_path=target_path,
        exchange_id=exchange_id,
        material_count=len(materials),
        judgement_count=len(entries),
        open_question_count=len(questions),
        included_file_count=len(included_files),
    )


def _load_case_update(package_path: Path) -> dict[str, Any]:
    with ZipFile(package_path) as archive:
        try:
            raw = archive.read("case_update.json")
        except KeyError as exc:
            raise CaseWorkspaceError(
                "case update package is missing case_update.json"
            ) from exc
    payload = json.loads(raw.decode("utf-8"))
    if not isinstance(payload, dict):
        raise CaseWorkspaceError("case update package must contain a JSON object")
    if payload.get("schema_version") != EXCHANGE_SCHEMA_VERSION:
        raise CaseWorkspaceError("unsupported case update schema_version")
    if payload.get("source") != EXCHANGE_SOURCE:
        raise CaseWorkspaceError("unsupported case update source")
    return payload


def _origin_index(records: Sequence[Mapping[str, Any]]) -> dict[str, Mapping[str, Any]]:
    return {
        str(record["origin_key"]): record
        for record in records
        if str(record.get("origin_key", "")).strip()
    }


def _material_compare_fields(record: Mapping[str, Any]) -> dict[str, Any]:
    return {
        key: record.get(key) for key in ("title", "material_type", "status", "summary")
    }


def _judgement_compare_fields(record: Mapping[str, Any]) -> dict[str, Any]:
    return {
        key: record.get(key)
        for key in (
            "kind",
            "text",
            "status",
            "rationale",
            "reviewer",
            "review_note",
        )
    }


def _question_compare_fields(record: Mapping[str, Any]) -> dict[str, Any]:
    return {key: record.get(key) for key in ("question", "why_it_matters", "status")}


def _append_conflict_question(
    questions: list[dict[str, Any]],
    *,
    origin_key: str,
    record_type: str,
    label: str,
    timestamp: str,
) -> bool:
    conflict_origin = f"conflict:{origin_key}"
    if any(item.get("origin_key") == conflict_origin for item in questions):
        return False
    question = {
        "id": _next_id("q", [item["id"] for item in questions]),
        "question": f"Review imported {record_type} conflict: {label}",
        "why_it_matters": (
            "Case-update import is append-only and did not overwrite the "
            "existing local record."
        ),
        "status": "open",
        "source_entry_ids": [],
        "created_at": timestamp,
        "updated_at": timestamp,
        "origin_key": conflict_origin,
    }
    questions.append(question)
    return True


def _safe_extract_archive_file(
    archive: ZipFile,
    *,
    archive_path: str,
    target_path: Path,
) -> None:
    archive_target = Path(archive_path)
    if archive_target.is_absolute() or ".." in archive_target.parts:
        raise CaseWorkspaceError(f"unsafe archive path: {archive_path}")
    target_path.parent.mkdir(parents=True, exist_ok=True)
    target_path.write_bytes(archive.read(archive_path))


def _append_exchange_log(
    case_dir: Path,
    *,
    exchange_id: str,
    source_case_id: str,
    package_path: Path,
    timestamp: str,
    counts: Mapping[str, int],
) -> None:
    log_path = case_dir / "exchange_log.json"
    if log_path.exists():
        payload = _read_json(log_path)
        if payload.get("schema_version") != SCHEMA_VERSION:
            raise CaseWorkspaceError("exchange_log.json: unsupported schema_version")
    else:
        payload = {"schema_version": SCHEMA_VERSION, "imports": []}
    payload["imports"].append(
        {
            "exchange_id": exchange_id,
            "source_case_id": source_case_id,
            "package_path": str(package_path.resolve()),
            "imported_at": timestamp,
            "counts": dict(counts),
        }
    )
    _write_json(log_path, payload)


def import_case_update(
    case_dir: Path,
    package_path: Path,
    *,
    now: datetime | None = None,
) -> CaseExchangeImportResult:
    """Import a case update by appending records and logging conflicts."""

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    timestamp = _now_iso(now)
    update_payload = _load_case_update(package_path)
    exchange_id = str(update_payload["exchange_id"])
    source_case_id = str(update_payload["source_case_id"])

    materials_path = _case_path(case_dir, "materials")
    judgement_path = _case_path(case_dir, "judgement")
    questions_path = _case_path(case_dir, "open_questions")
    material_payload = _read_json(materials_path)
    judgement_payload = _read_json(judgement_path)
    questions_payload = _read_json(questions_path)
    materials = material_payload["materials"]
    entries = judgement_payload["entries"]
    questions = questions_payload["questions"]

    material_origins = _origin_index(materials)
    judgement_origins = _origin_index(entries)
    question_origins = _origin_index(questions)
    included_by_material_id = {
        str(item["material_id"]): item
        for item in update_payload.get("included_files", [])
        if isinstance(item, Mapping)
    }

    imported_materials = 0
    imported_entries = 0
    imported_questions = 0
    skipped = 0
    conflicts = 0
    material_id_map: dict[str, str] = {}

    with ZipFile(package_path) as archive:
        for raw_material in update_payload.get("materials", []):
            if not isinstance(raw_material, Mapping):
                continue
            origin_key = _record_origin_key(
                raw_material,
                source_case_id=source_case_id,
                record_type="material",
            )
            existing = material_origins.get(origin_key)
            if existing is not None:
                material_id_map[str(raw_material["id"])] = str(existing["id"])
                if _material_compare_fields(existing) != _material_compare_fields(
                    raw_material
                ):
                    added_conflict = _append_conflict_question(
                        questions,
                        origin_key=origin_key,
                        record_type="material",
                        label=str(raw_material.get("title", raw_material["id"])),
                        timestamp=timestamp,
                    )
                    conflicts += int(added_conflict)
                else:
                    skipped += 1
                continue

            new_id = _next_id("mat", [item["id"] for item in materials])
            source_material_id = str(raw_material["id"])
            included_file = included_by_material_id.get(source_material_id)
            imported_path = str(raw_material.get("path", ""))
            availability = "external_reference"
            if included_file is not None:
                relative_path = Path(str(included_file["relative_path"]))
                if relative_path.is_absolute() or ".." in relative_path.parts:
                    raise CaseWorkspaceError(
                        f"unsafe included file path: {relative_path}"
                    )
                imported_file_path = (
                    case_dir / "exchange_imports" / exchange_id / relative_path
                )
                _safe_extract_archive_file(
                    archive,
                    archive_path=str(included_file["archive_path"]),
                    target_path=imported_file_path,
                )
                imported_path = str(imported_file_path.resolve())
                availability = "local_copy"
            material = {
                **dict(raw_material),
                "id": new_id,
                "path": imported_path,
                "imported_at": timestamp,
                "updated_at": timestamp,
                "origin_key": origin_key,
                "origin_case_id": source_case_id,
                "origin_material_id": source_material_id,
                "exchange_id": exchange_id,
                "availability": availability,
            }
            materials.append(material)
            material_origins[origin_key] = material
            material_id_map[source_material_id] = new_id
            imported_materials += 1

    for raw_entry in update_payload.get("judgement_entries", []):
        if not isinstance(raw_entry, Mapping):
            continue
        origin_key = _record_origin_key(
            raw_entry,
            source_case_id=source_case_id,
            record_type="judgement",
        )
        existing = judgement_origins.get(origin_key)
        if existing is not None:
            if _judgement_compare_fields(existing) != _judgement_compare_fields(
                raw_entry
            ):
                added_conflict = _append_conflict_question(
                    questions,
                    origin_key=origin_key,
                    record_type="judgement",
                    label=str(raw_entry.get("text", raw_entry["id"]))[:120],
                    timestamp=timestamp,
                )
                conflicts += int(added_conflict)
            else:
                skipped += 1
            continue

        source_ids = [
            material_id_map[source_id]
            for source_id in raw_entry.get("source_material_ids", [])
            if source_id in material_id_map
        ]
        entry = {
            **dict(raw_entry),
            "id": _next_id("jud", [item["id"] for item in entries]),
            "source_material_ids": source_ids,
            "imported_at": timestamp,
            "origin_key": origin_key,
            "origin_case_id": source_case_id,
            "origin_judgement_id": str(raw_entry["id"]),
            "exchange_id": exchange_id,
        }
        _validate_choice(str(entry["kind"]), JUDGEMENT_KINDS, "kind")
        _validate_choice(str(entry["status"]), JUDGEMENT_STATUSES, "status")
        entries.append(entry)
        judgement_origins[origin_key] = entry
        imported_entries += 1

    for raw_question in update_payload.get("open_questions", []):
        if not isinstance(raw_question, Mapping):
            continue
        origin_key = _record_origin_key(
            raw_question,
            source_case_id=source_case_id,
            record_type="open_question",
        )
        existing = question_origins.get(origin_key)
        if existing is not None:
            if _question_compare_fields(existing) != _question_compare_fields(
                raw_question
            ):
                added_conflict = _append_conflict_question(
                    questions,
                    origin_key=origin_key,
                    record_type="open question",
                    label=str(raw_question.get("question", raw_question["id"]))[:120],
                    timestamp=timestamp,
                )
                conflicts += int(added_conflict)
            else:
                skipped += 1
            continue
        question = {
            **dict(raw_question),
            "id": _next_id("q", [item["id"] for item in questions]),
            "imported_at": timestamp,
            "updated_at": timestamp,
            "origin_key": origin_key,
            "origin_case_id": source_case_id,
            "origin_question_id": str(raw_question["id"]),
            "exchange_id": exchange_id,
        }
        _validate_choice(str(question["status"]), OPEN_QUESTION_STATUSES, "status")
        questions.append(question)
        question_origins[origin_key] = question
        imported_questions += 1

    _write_json(materials_path, material_payload)
    _write_json(judgement_path, judgement_payload)
    _write_json(questions_path, questions_payload)
    _touch_manifest(case_dir, timestamp)
    counts = {
        "imported_materials": imported_materials,
        "imported_judgement_entries": imported_entries,
        "imported_open_questions": imported_questions,
        "skipped": skipped,
        "conflicts": conflicts,
    }
    _append_exchange_log(
        case_dir,
        exchange_id=exchange_id,
        source_case_id=source_case_id,
        package_path=package_path,
        timestamp=timestamp,
        counts=counts,
    )

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))
    refresh_case_brief(case_dir, now=now)

    return CaseExchangeImportResult(
        exchange_id=exchange_id,
        imported_material_count=imported_materials,
        imported_judgement_count=imported_entries,
        imported_open_question_count=imported_questions,
        skipped_count=skipped,
        conflict_count=conflicts,
    )


def _entry_source_suffix(
    entry: Mapping[str, Any],
    material_labels: Mapping[str, str],
    *,
    source_label: str = "Sources",
) -> str:
    source_ids = entry.get("source_material_ids", [])
    labels = [
        material_labels[source_id]
        for source_id in source_ids
        if source_id in material_labels
    ]
    if labels:
        return f" {source_label}: {', '.join(labels)}."
    return ""


def _brief_entry_line(
    entry: Mapping[str, Any],
    material_labels: Mapping[str, str],
    *,
    include_kind: bool = False,
    source_label: str = "Sources",
    kind_labels: Mapping[str, str] | None = None,
) -> str:
    prefix = ""
    if include_kind:
        kind = str(entry["kind"])
        kind_label = (kind_labels or {}).get(kind, kind.replace("_", " "))
        prefix += f"[{kind_label}] "
    return (
        f"- {prefix}{str(entry['text']).strip()}"
        f"{_entry_source_suffix(entry, material_labels, source_label=source_label)}"
    )


def _brief_lines_for_entries(
    entries: Sequence[Mapping[str, Any]],
    material_labels: Mapping[str, str],
    *,
    empty_text: str,
    include_kind: bool = False,
    source_label: str = "Sources",
    kind_labels: Mapping[str, str] | None = None,
) -> list[str]:
    if not entries:
        return [f"- {empty_text}"]
    return [
        _brief_entry_line(
            entry,
            material_labels,
            include_kind=include_kind,
            source_label=source_label,
            kind_labels=kind_labels,
        )
        for entry in entries
    ]


def _brief_material_lines(
    materials: Sequence[Mapping[str, Any]],
    *,
    no_materials: str = "No materials indexed.",
) -> list[str]:
    if not materials:
        return [f"- {no_materials}"]
    lines = []
    for item in materials:
        summary = str(item.get("summary", "")).strip()
        summary_suffix = f" - {summary}" if summary else ""
        lines.append(
            f"- {item['title']} [{item['material_type']}, {item['status']}]"
            f"{summary_suffix}"
        )
    return lines


def _brief_question_lines(
    questions: Sequence[Mapping[str, Any]],
    *,
    no_questions: str = "No open questions registered.",
    why_label: str = "Why it matters",
) -> list[str]:
    open_items = [item for item in questions if item.get("status") == "open"]
    if not open_items:
        return [f"- {no_questions}"]
    return [
        f"- {item['question']} {why_label}: {item['why_it_matters']}"
        for item in open_items
    ]


def _brief_issue_lines(
    issues: Sequence[Mapping[str, Any]],
    *,
    no_issues: str = "No active case issues registered.",
    evidence_label: str = "Evidence links",
    supporting_label: str = "supporting",
    opposing_label: str = "opposing",
    open_test_label: str = "open test",
) -> list[str]:
    active_issues = [item for item in issues if item.get("status") != "resolved"]
    if not active_issues:
        return [f"- {no_issues}"]
    lines: list[str] = []
    for issue in active_issues:
        synthesis = str(issue.get("current_synthesis", "")).strip()
        synthesis_suffix = f" - {synthesis}" if synthesis else ""
        lines.append(
            f"- {issue['title']} [{issue['status']}, "
            f"{issue.get('decision_area', '')}]{synthesis_suffix}"
        )
        evidence_for = len(issue.get("evidence_for", []))
        evidence_against = len(issue.get("evidence_against", []))
        open_tests = len(issue.get("open_tests", []))
        lines.append(
            f"  {evidence_label}: {evidence_for} {supporting_label}, "
            f"{evidence_against} {opposing_label}, {open_tests} {open_test_label}."
        )
    return lines


def _case_brief_copy(manifest: Mapping[str, Any]) -> dict[str, Any]:
    if _ui_language(manifest) == "es":
        return {
            "title": "Resumen del caso",
            "derived": (
                "Resumen de trabajo derivado. La fuente de referencia sigue siendo "
                "el conjunto de archivos JSON del caso."
            ),
            "pending_gate": (
                "Los elementos pendientes sirven únicamente para revisión y no "
                "pueden alimentar el paquete de decisión del cliente hasta que se "
                "marquen como listos para dicho paquete."
            ),
            "generated": "Generado",
            "updated": "Caso actualizado",
            "snapshot": "Resumen del caso",
            "project": "Proyecto",
            "objective": "Objetivo",
            "audience": "Destinatario",
            "status": "Estado",
            "output_language": "Idioma de salida",
            "materials_indexed": "Materiales indexados",
            "approved_count": "Elementos de juicio listos para el paquete de decisión",
            "pending_count": "Elementos de juicio pendientes",
            "rejected_count": "Elementos de juicio rechazados",
            "active_issues_count": "Cuestiones activas del caso",
            "current": "Comprensión actual (lista para el paquete de decisión)",
            "facts": "Hechos",
            "advisor": "Juicio consultivo",
            "inferences": "Inferencias de Codex",
            "decisions": "Implicaciones para la decisión",
            "no_facts": "Aún no hay hechos listos para el paquete de decisión.",
            "no_advisor": (
                "Aún no hay juicios consultivos listos para el paquete de decisión."
            ),
            "no_inferences": (
                "Aún no hay inferencias de Codex listas para el paquete de decisión."
            ),
            "no_decisions": (
                "Aún no hay implicaciones listas para el paquete de decisión."
            ),
            "candidate_review": (
                "Revisión de candidatos (no listos para el paquete de decisión)"
            ),
            "candidate_note": (
                "Estos elementos se han registrado para su revisión. No constituyen "
                "asesoramiento final."
            ),
            "no_pending": "No hay elementos de juicio pendientes.",
            "questions": "Preguntas abiertas",
            "no_questions": "No se han registrado preguntas abiertas.",
            "why": "Por qué es importante",
            "issues": "Cuestiones del caso",
            "no_issues": "No se han registrado cuestiones activas del caso.",
            "evidence": "Vínculos de evidencia",
            "supporting": "a favor",
            "opposing": "en contra",
            "open_test": "prueba abierta",
            "mandate": "Mandato de Clara",
            "current_understanding": "Comprensión actual",
            "partner_orientation": "Orientación del socio",
            "not_recorded": "Aún no registrado.",
            "clarifications": "Aclaraciones imprescindibles",
            "no_clarifications": "Aún no se han registrado aclaraciones de Clara.",
            "next_steps": "Próximos pasos de Clara",
            "no_next_steps": "Aún no se han registrado próximos pasos de Clara.",
            "materials": "Materiales",
            "no_materials": "No hay materiales indexados.",
            "control": "Notas de control",
            "control_brief": (
                "`case_brief.md` se deriva de `case_manifest.json`, "
                "`material_registry.json`, `judgement_log.json`, "
                "`open_questions.json`, `case_issues.json` y `clara_mandate.json`."
            ),
            "control_pack": (
                "`decision_pack.md` y `decision_pack.docx` utilizan únicamente "
                "juicios listos para el paquete de decisión."
            ),
            "control_rebuild": (
                "Vuelve a generar este resumen con `scripts/build_case_brief.py` si "
                "se han realizado cambios externos en los archivos JSON."
            ),
            "sources": "Fuentes",
            "kind_labels": {
                "fact": "hecho",
                "advisor_judgement": "juicio consultivo",
                "codex_inference": "inferencia de Codex",
                "decision_implication": "implicación para la decisión",
            },
        }
    return {
        "title": "Case Brief",
        "derived": "Derived working brief. The source of truth remains the JSON case files.",
        "pending_gate": (
            "Pending items are review material only and cannot feed the client "
            "decision pack until marked decision-pack ready."
        ),
        "generated": "Generated",
        "updated": "Case updated",
        "snapshot": "Case Snapshot",
        "project": "Project",
        "objective": "Objective",
        "audience": "Audience",
        "status": "Status",
        "output_language": "Output language",
        "materials_indexed": "Materials indexed",
        "approved_count": "Decision-pack-ready judgement entries",
        "pending_count": "Pending judgement entries",
        "rejected_count": "Rejected judgement entries",
        "active_issues_count": "Active case issues",
        "current": "Current Understanding (decision-pack ready)",
        "facts": "Facts",
        "advisor": "Consultant Judgement",
        "inferences": "Codex Inferences",
        "decisions": "Decision Implications",
        "no_facts": "No decision-pack-ready facts yet.",
        "no_advisor": "No decision-pack-ready consultant judgement yet.",
        "no_inferences": "No decision-pack-ready Codex inferences yet.",
        "no_decisions": "No decision-pack-ready decision implications yet.",
        "candidate_review": "Candidate Review (not decision-pack ready)",
        "candidate_note": "These items are captured for review. They are not final advice.",
        "no_pending": "No pending judgement entries.",
        "questions": "Open Questions",
        "no_questions": "No open questions registered.",
        "why": "Why it matters",
        "issues": "Case Issues",
        "no_issues": "No active case issues registered.",
        "evidence": "Evidence links",
        "supporting": "supporting",
        "opposing": "opposing",
        "open_test": "open test",
        "mandate": "Clara Mandate",
        "current_understanding": "Current understanding",
        "partner_orientation": "Partner orientation",
        "not_recorded": "Not recorded yet.",
        "clarifications": "Essential Clarifications",
        "no_clarifications": "No Clara clarifications recorded yet.",
        "next_steps": "Clara Next Steps",
        "no_next_steps": "No Clara next steps recorded yet.",
        "materials": "Materials",
        "no_materials": "No materials indexed.",
        "control": "Control Notes",
        "control_brief": (
            "`case_brief.md` is derived from `case_manifest.json`, "
            "`material_registry.json`, `judgement_log.json`, "
            "`open_questions.json`, `case_issues.json`, and `clara_mandate.json`."
        ),
        "control_pack": (
            "`decision_pack.md` and `decision_pack.docx` use decision-pack-ready "
            "judgement only."
        ),
        "control_rebuild": (
            "Rebuild this brief with `scripts/build_case_brief.py` if external "
            "edits were made to the JSON files."
        ),
        "sources": "Sources",
        "kind_labels": {},
    }


def _render_case_brief(
    *,
    manifest: Mapping[str, Any],
    materials: Sequence[Mapping[str, Any]],
    entries: Sequence[Mapping[str, Any]],
    open_questions: Sequence[Mapping[str, Any]],
    issues: Sequence[Mapping[str, Any]],
    clara_mandate: Mapping[str, Any],
    generated_at: str,
) -> str:
    copy = _case_brief_copy(manifest)
    grouped = _group_approved_entries(entries)
    material_labels = _material_label_by_id(materials)
    pending_entries = [entry for entry in entries if entry["status"] == "pending"]
    rejected_count = sum(1 for entry in entries if entry["status"] == "rejected")
    approved_count = sum(1 for entry in entries if entry["status"] == "approved")
    mandate_body = clara_mandate.get("mandate", {})
    mandate_status = str(clara_mandate.get("status", "not_started"))
    clarification_lines = [
        f"- {item}"
        for item in _clean_string_list(mandate_body.get("essential_clarifications", []))
    ] or [f"- {copy['no_clarifications']}"]
    next_step_lines = [
        f"- {item}" for item in _clean_string_list(mandate_body.get("next_steps", []))
    ] or [f"- {copy['no_next_steps']}"]

    lines = [
        f"# {copy['title']} - {manifest['client']}",
        "",
        copy["derived"],
        copy["pending_gate"],
        "",
        f"{copy['generated']}: {generated_at}",
        f"{copy['updated']}: {manifest.get('updated_at', '')}",
        "",
        f"## {copy['snapshot']}",
        "",
        f"- {copy['project']}: {manifest['project']}",
        f"- {copy['objective']}: {manifest['objective']}",
        f"- {copy['audience']}: {manifest['audience']}",
        f"- {copy['status']}: {manifest['status']}",
        f"- {copy['output_language']}: {manifest['output_language']}",
        f"- {copy['materials_indexed']}: {len(materials)}",
        f"- {copy['approved_count']}: {approved_count}",
        f"- {copy['pending_count']}: {len(pending_entries)}",
        f"- {copy['rejected_count']}: {rejected_count}",
        f"- {copy['active_issues_count']}: {sum(1 for issue in issues if issue.get('status') == 'active')}",
        "",
        f"## {copy['current']}",
        "",
        f"### {copy['facts']}",
        "",
        *_brief_lines_for_entries(
            grouped["fact"],
            material_labels,
            empty_text=copy["no_facts"],
            source_label=copy["sources"],
        ),
        "",
        f"### {copy['advisor']}",
        "",
        *_brief_lines_for_entries(
            grouped["advisor_judgement"],
            material_labels,
            empty_text=copy["no_advisor"],
            source_label=copy["sources"],
        ),
        "",
        f"### {copy['inferences']}",
        "",
        *_brief_lines_for_entries(
            grouped["codex_inference"],
            material_labels,
            empty_text=copy["no_inferences"],
            source_label=copy["sources"],
        ),
        "",
        f"### {copy['decisions']}",
        "",
        *_brief_lines_for_entries(
            grouped["decision_implication"],
            material_labels,
            empty_text=copy["no_decisions"],
            source_label=copy["sources"],
        ),
        "",
        f"## {copy['candidate_review']}",
        "",
        copy["candidate_note"],
        "",
        *_brief_lines_for_entries(
            pending_entries,
            material_labels,
            empty_text=copy["no_pending"],
            include_kind=True,
            source_label=copy["sources"],
            kind_labels=copy["kind_labels"],
        ),
        "",
        f"## {copy['questions']}",
        "",
        *_brief_question_lines(
            open_questions,
            no_questions=copy["no_questions"],
            why_label=copy["why"],
        ),
        "",
        f"## {copy['issues']}",
        "",
        *_brief_issue_lines(
            issues,
            no_issues=copy["no_issues"],
            evidence_label=copy["evidence"],
            supporting_label=copy["supporting"],
            opposing_label=copy["opposing"],
            open_test_label=copy["open_test"],
        ),
        "",
        f"## {copy['mandate']}",
        "",
        f"- {copy['status']}: {mandate_status}",
        f"- {copy['current_understanding']}: {str(mandate_body.get('clara_understanding', '')).strip() or copy['not_recorded']}",
        f"- {copy['partner_orientation']}: {str(mandate_body.get('partner_starting_orientation', '')).strip() or copy['not_recorded']}",
        "",
        f"### {copy['clarifications']}",
        "",
        *clarification_lines,
        "",
        f"### {copy['next_steps']}",
        "",
        *next_step_lines,
        "",
        f"## {copy['materials']}",
        "",
        *_brief_material_lines(materials, no_materials=copy["no_materials"]),
        "",
        f"## {copy['control']}",
        "",
        f"- {copy['control_brief']}",
        f"- {copy['control_pack']}",
        f"- {copy['control_rebuild']}",
        "",
    ]
    return "\n".join(lines)


def refresh_case_brief(
    case_dir: Path,
    *,
    now: datetime | None = None,
) -> CaseBriefResult:
    """Deterministically rebuild the readable brief from canonical case JSON.

    The brief is a convenience view, not a semantic source of truth. It uses
    fixed formatting so inclusion counts and pending gates remain auditable.
    """

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    generated_at = _now_iso(now)
    manifest = _read_json(_case_path(case_dir, "manifest"))
    materials = _read_json(_case_path(case_dir, "materials"))["materials"]
    entries = _read_json(_case_path(case_dir, "judgement"))["entries"]
    questions = _read_json(_case_path(case_dir, "open_questions"))["questions"]
    issues = _read_json(_case_path(case_dir, "issues"))["issues"]
    clara_mandate = _read_json(_case_path(case_dir, "clara_mandate"))
    brief_path = case_dir / CASE_BRIEF_FILENAME
    brief_text = _render_case_brief(
        manifest=manifest,
        materials=materials,
        entries=entries,
        open_questions=questions,
        issues=issues,
        clara_mandate=clara_mandate,
        generated_at=generated_at,
    )
    brief_path.write_text(brief_text, encoding="utf-8")

    return CaseBriefResult(
        brief_path=brief_path,
        approved_count=sum(1 for entry in entries if entry["status"] == "approved"),
        pending_count=sum(1 for entry in entries if entry["status"] == "pending"),
        rejected_count=sum(1 for entry in entries if entry["status"] == "rejected"),
        open_question_count=sum(
            1 for question in questions if question["status"] == "open"
        ),
        issue_count=sum(1 for issue in issues if issue["status"] == "active"),
    )


def _review_status_entries(
    entries: Sequence[Mapping[str, Any]],
    status: str,
) -> list[Mapping[str, Any]]:
    return [entry for entry in entries if entry["status"] == status]


def _review_source_text(
    entry: Mapping[str, Any],
    material_labels: Mapping[str, str],
) -> str:
    source_ids = entry.get("source_material_ids", [])
    labels = [
        material_labels.get(str(source_id), str(source_id))
        for source_id in source_ids
        if str(source_id).strip()
    ]
    if not labels:
        return "None recorded"
    return ", ".join(labels)


def _slugify_bundle_id(value: str, *, fallback: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "-", value.strip().lower()).strip("-")
    return slug or fallback


def _inclusion_bundles_path(case_dir: Path) -> Path:
    return case_dir / INCLUSION_BUNDLES_FILENAME


def _read_inclusion_bundles(case_dir: Path) -> list[dict[str, Any]]:
    path = _inclusion_bundles_path(case_dir)
    if not path.exists():
        return []
    payload = _read_json(path)
    if payload.get("schema_version") != SCHEMA_VERSION:
        raise CaseWorkspaceError(
            f"{INCLUSION_BUNDLES_FILENAME}: unsupported schema_version"
        )
    bundles = payload.get("bundles")
    if not isinstance(bundles, list):
        raise CaseWorkspaceError(
            f"{INCLUSION_BUNDLES_FILENAME}: bundles must be a list"
        )
    return [bundle for bundle in bundles if isinstance(bundle, dict)]


def apply_inclusion_bundles(
    case_dir: Path,
    bundles: Sequence[Mapping[str, Any]],
    *,
    now: datetime | None = None,
) -> InclusionBundleApplyResult:
    """Persist semantic inclusion-review bundles supplied by Codex/Clara.

    Bundle themes are semantic judgement, so this helper does not infer them.
    It only validates entry IDs, de-duplicates mechanical identifiers, and writes
    the review plan that later approval commands can apply audibly.
    """

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))
    judgement = _read_json(_case_path(case_dir, "judgement"))
    known_entries = {
        str(entry["id"]): entry
        for entry in judgement.get("entries", [])
        if entry.get("id")
    }
    cleaned_bundles: list[dict[str, Any]] = []
    used_bundle_ids: set[str] = set()
    used_entry_ids: dict[str, str] = {}

    for index, raw_bundle in enumerate(bundles, start=1):
        title = str(raw_bundle.get("title", "")).strip()
        if not title:
            raise CaseWorkspaceError(f"inclusion bundle {index}: title is required")
        bundle_id = _slugify_bundle_id(
            str(raw_bundle.get("id", "")).strip() or title,
            fallback=f"bundle-{index}",
        )
        if bundle_id in used_bundle_ids:
            raise CaseWorkspaceError(f"inclusion bundle id is duplicated: {bundle_id}")
        used_bundle_ids.add(bundle_id)
        entry_ids = _dedupe_preserve_order(
            str(value).strip() for value in raw_bundle.get("entry_ids", [])
        )
        if not entry_ids:
            raise CaseWorkspaceError(
                f"inclusion bundle {title!r}: at least one entry_id is required"
            )
        missing_ids = [
            entry_id for entry_id in entry_ids if entry_id not in known_entries
        ]
        if missing_ids:
            raise CaseWorkspaceError(
                f"inclusion bundle {title!r}: unknown entry_ids "
                + ", ".join(missing_ids)
            )
        duplicate_entries = [
            entry_id for entry_id in entry_ids if entry_id in used_entry_ids
        ]
        if duplicate_entries:
            duplicate_text = ", ".join(
                f"{entry_id} already in {used_entry_ids[entry_id]}"
                for entry_id in duplicate_entries
            )
            raise CaseWorkspaceError(
                f"inclusion bundle {title!r}: duplicate entry assignment: "
                + duplicate_text
            )
        for entry_id in entry_ids:
            used_entry_ids[entry_id] = bundle_id
        cleaned_bundles.append(
            {
                "id": bundle_id,
                "title": title,
                "description": str(raw_bundle.get("description", "")).strip(),
                "entry_ids": entry_ids,
            }
        )

    timestamp = _now_iso(now)
    payload = {
        "schema_version": SCHEMA_VERSION,
        "updated_at": timestamp,
        "bundles": cleaned_bundles,
    }
    bundles_path = _inclusion_bundles_path(case_dir)
    _write_json(bundles_path, payload)
    _touch_manifest(case_dir, timestamp)
    return InclusionBundleApplyResult(
        bundles_path=bundles_path,
        bundle_count=len(cleaned_bundles),
        bundled_entry_count=sum(len(bundle["entry_ids"]) for bundle in cleaned_bundles),
    )


def _pending_inclusion_bundle_groups(
    case_dir: Path,
    pending_entries: Sequence[Mapping[str, Any]],
) -> list[dict[str, Any]]:
    pending_by_id = {str(entry["id"]): entry for entry in pending_entries}
    groups: list[dict[str, Any]] = []
    for bundle in _read_inclusion_bundles(case_dir):
        entry_ids = [
            str(entry_id)
            for entry_id in bundle.get("entry_ids", [])
            if str(entry_id) in pending_by_id
        ]
        if not entry_ids:
            continue
        groups.append(
            {
                "id": str(bundle.get("id", "")).strip(),
                "title": str(bundle.get("title", "")).strip(),
                "description": str(bundle.get("description", "")).strip(),
                "entries": [pending_by_id[entry_id] for entry_id in entry_ids],
            }
        )
    return groups


def _compact_review_text(value: str, *, limit: int = 150) -> str:
    text = re.sub(r"\s+", " ", value).strip()
    if len(text) <= limit:
        return text
    return f"{text[:limit].rstrip()}..."


def _inclusion_bundle_lines(
    groups: Sequence[Mapping[str, Any]],
    pending_item_numbers: Mapping[str, int],
    *,
    copy: Mapping[str, str],
) -> list[str]:
    if not groups:
        return [copy["no_bundles"], ""]

    lines: list[str] = []
    for index, group in enumerate(groups, start=1):
        entries = [
            entry
            for entry in group.get("entries", [])
            if isinstance(entry, Mapping)
            and str(entry.get("id", "")) in pending_item_numbers
        ]
        item_numbers = [
            str(pending_item_numbers[str(entry["id"])]) for entry in entries
        ]
        lines.extend(
            [
                f"### {copy['bundle']} {index}: {group['title']}",
                "",
                f"- {copy['bundle_items']}: {', '.join(item_numbers)}",
            ]
        )
        description = str(group.get("description", "")).strip()
        if description:
            lines.append(f"- {copy['bundle_description']}: {description}")
        lines.append(
            f"- {copy['bundle_action'].format(index=index, bundle_id=group['id'])}"
        )
        for entry in entries:
            item_number = pending_item_numbers[str(entry["id"])]
            kind_label = str(entry["kind"]).replace("_", " ")
            lines.append(
                f"  - {copy['item']} {item_number} [{kind_label}]: "
                f"{_compact_review_text(str(entry['text']))}"
            )
        lines.append("")
    return lines


def _inclusion_entry_lines(
    entries: Sequence[Mapping[str, Any]],
    material_labels: Mapping[str, str],
    *,
    include_actions: bool,
    copy: Mapping[str, str],
) -> list[str]:
    if not entries:
        return [copy["no_entries"], ""]

    lines: list[str] = []
    for index, entry in enumerate(entries, start=1):
        kind_label = str(entry["kind"]).replace("_", " ")
        lines.extend(
            [
                f"### {index}. {kind_label}",
                "",
                str(entry["text"]).strip(),
                "",
                f"- {copy['sources']}: {_review_source_text(entry, material_labels)}",
            ]
        )
        rationale = str(entry.get("rationale", "")).strip()
        if rationale:
            lines.append(f"- {copy['rationale']}: {rationale}")
        reviewer = str(entry.get("reviewer", "") or "").strip()
        if reviewer:
            lines.append(f"- {copy['recorded_by']}: {reviewer}")
        review_note = str(entry.get("review_note", "") or "").strip()
        if review_note:
            lines.append(f"- {copy['review_note']}: {review_note}")
        if include_actions:
            lines.append(f"- {copy['entry_action'].format(index=index)}")
        lines.append("")
    return lines


def _inclusion_review_copy(language: str) -> dict[str, str]:
    if language == "it":
        return {
            "already_approved": "Voci gia' pronte per il pack cliente",
            "already_approved_count": "Voci gia' pronte per il pack cliente",
            "ask_evidence": "Chiedi evidenza",
            "bundle": "Pacchetto",
            "bundle_action": (
                "Risposta: includi pacchetto {index}, escludi pacchetto {index}, "
                "oppure mostrami evidenza sul pacchetto {index}. ID: `{bundle_id}`."
            ),
            "bundle_description": "Nota",
            "bundle_items": "Voci",
            "bundles": "Pacchetti pending",
            "bundles_count": "Pacchetti pending disponibili",
            "control_heading": "Note di controllo",
            "control_pending": (
                "Le voci pending e rejected non devono alimentare "
                "`decision_pack.md` o `decision_pack.docx`."
            ),
            "control_regenerated": (
                "Questo file deriva dai JSON canonici del caso e puo' essere rigenerato."
            ),
            "create_bundles": (
                "Per revisioni lunghe, Clara/Codex puo' creare "
                "`inclusion_bundles.json` con pacchetti tematici e rigenerare "
                "questa checklist."
            ),
            "correct_one": "Correggi una voce",
            "correct_one_example": "correggi voce 9: ...",
            "excluded": "Voci escluse",
            "excluded_count": "Voci gia' escluse dal pack cliente",
            "exclude_one": "Escludi una voce",
            "exclude_one_example": "escludi voce 7",
            "entry_action": (
                "Risposta: includi voce {index}, escludi voce {index}, "
                "correggi voce {index}, oppure mostrami evidenza sulla voce {index}."
            ),
            "generated": "Generato",
            "how_to_respond": "Come rispondere",
            "include_all": "Includi tutto",
            "include_all_example": "includi tutte le voci pending",
            "include_bundle": "Includi un pacchetto",
            "include_bundle_example": "includi pacchetto 1",
            "include_one": "Includi una voce",
            "include_one_example": "includi voce 4",
            "intro": (
                "Checklist di lavoro per decidere quali voci strutturate del caso "
                "sono pronte per il pack cliente."
            ),
            "intro_gate": (
                "Questo file non cambia lo stato delle voci. Il partner indica "
                "cosa includere, escludere, correggere o approfondire; la "
                "decisione viene registrata meccanicamente dopo conferma."
            ),
            "no_entries": "Nessuna voce in questo stato.",
            "no_bundles": "Nessun pacchetto tematico configurato.",
            "pending": "Voci pending",
            "pending_count": "Voci pending che richiedono decisione del partner",
            "item": "Voce",
            "project": "Progetto",
            "audience": "Destinatario",
            "rationale": "Motivo",
            "recorded_by": "Registrato da",
            "review_note": "Nota di revisione",
            "show_more_example": "mostrami evidenza sulla voce 3",
            "sources": "Fonti",
            "status_counts": "Conteggi",
            "title": "Revisione inclusione",
        }
    if language == "es":
        return {
            "already_approved": "Elementos ya listos para el paquete del cliente",
            "already_approved_count": (
                "Elementos ya marcados como listos para el paquete del cliente"
            ),
            "ask_evidence": "Pedir evidencia",
            "bundle": "Paquete",
            "bundle_action": (
                "Respuesta: incluye el paquete {index}, excluye el paquete {index} "
                "o muestra más sobre el paquete {index}. ID: `{bundle_id}`."
            ),
            "bundle_description": "Nota",
            "bundle_items": "Elementos",
            "bundles": "Paquetes pendientes",
            "bundles_count": "Paquetes pendientes disponibles",
            "control_heading": "Notas de control",
            "control_pending": (
                "Los elementos pendientes y rechazados no deben alimentar "
                "`decision_pack.md` ni `decision_pack.docx`."
            ),
            "control_regenerated": (
                "Este archivo deriva de los JSON canónicos del caso y puede "
                "regenerarse de forma segura."
            ),
            "create_bundles": (
                "Para revisiones largas, Clara/Codex puede crear "
                "`inclusion_bundles.json` con paquetes temáticos y regenerar "
                "esta lista."
            ),
            "correct_one": "Corregir un elemento",
            "correct_one_example": "corrige el elemento 9: ...",
            "excluded": "Elementos excluidos",
            "excluded_count": "Elementos ya excluidos del paquete del cliente",
            "exclude_one": "Excluir un elemento",
            "exclude_one_example": "excluye el elemento 7",
            "entry_action": (
                "Respuesta: incluye el elemento {index}, excluye el elemento "
                "{index}, corrige el elemento {index} o muestra más sobre el "
                "elemento {index}."
            ),
            "generated": "Generado",
            "how_to_respond": "Cómo responder",
            "include_all": "Incluir todo",
            "include_all_example": "incluye todos los elementos pendientes",
            "include_bundle": "Incluir un paquete",
            "include_bundle_example": "incluye el paquete 1",
            "include_one": "Incluir un elemento",
            "include_one_example": "incluye el elemento 4",
            "intro": (
                "Lista de trabajo para decidir qué elementos estructurados del "
                "caso están listos para el paquete del cliente."
            ),
            "intro_gate": (
                "Este documento no cambia el estado de los juicios. Indica qué "
                "elementos incluir, excluir, corregir o ampliar; la decisión se "
                "registra de forma mecánica tras la confirmación."
            ),
            "no_entries": "No hay elementos en este estado.",
            "no_bundles": "No hay paquetes temáticos configurados.",
            "pending": "Elementos pendientes",
            "pending_count": "Elementos pendientes que requieren decisión del socio",
            "item": "Elemento",
            "project": "Proyecto",
            "audience": "Destinatario",
            "rationale": "Motivo",
            "recorded_by": "Registrado por",
            "review_note": "Nota de revisión",
            "show_more_example": "muestra más sobre el elemento 3",
            "sources": "Fuentes",
            "status_counts": "Recuento por estado",
            "title": "Revisión de inclusión",
        }
    return {
        "already_approved": "Already Ready For Client Pack",
        "already_approved_count": "Already marked ready for client pack",
        "ask_evidence": "Ask for evidence",
        "bundle": "Bundle",
        "bundle_action": (
            "Response: include bundle {index}, exclude bundle {index}, "
            "or show more on bundle {index}. ID: `{bundle_id}`."
        ),
        "bundle_description": "Note",
        "bundle_items": "Items",
        "bundles": "Pending Approval Bundles",
        "bundles_count": "Pending approval bundles available",
        "control_heading": "Control Notes",
        "control_pending": (
            "Pending and rejected entries must not feed `decision_pack.md` "
            "or `decision_pack.docx`."
        ),
        "control_regenerated": (
            "This file is derived from canonical JSON case files and can be regenerated safely."
        ),
        "create_bundles": (
            "For long reviews, Clara/Codex can create `inclusion_bundles.json` "
            "with thematic bundles and regenerate this checklist."
        ),
        "correct_one": "Correct one item",
        "correct_one_example": "correct item 9: ...",
        "excluded": "Excluded Entries",
        "excluded_count": "Already excluded from client pack",
        "exclude_one": "Exclude one item",
        "exclude_one_example": "exclude item 7",
        "entry_action": (
            "Response: include item {index}, exclude item {index}, "
            "correct item {index}, or show more on item {index}."
        ),
        "generated": "Generated",
        "how_to_respond": "How To Respond",
        "include_all": "Include everything",
        "include_all_example": "include all pending items",
        "include_bundle": "Include one bundle",
        "include_bundle_example": "include bundle 1",
        "include_one": "Include one item",
        "include_one_example": "include item 4",
        "intro": (
            "Read-only checklist for deciding which structured case entries "
            "are ready for the client pack."
        ),
        "intro_gate": (
            "This artifact does not change judgement status. Indicate which "
            "items to include, exclude, correct, or expand; the decision is "
            "recorded mechanically after confirmation."
        ),
        "no_entries": "No entries in this status.",
        "no_bundles": "No thematic bundles configured.",
        "pending": "Pending Entries",
        "pending_count": "Pending entries requiring advisor decision",
        "item": "Item",
        "project": "Project",
        "audience": "Audience",
        "rationale": "Rationale",
        "recorded_by": "Recorded by",
        "review_note": "Review note",
        "show_more_example": "show more on item 3",
        "sources": "Sources",
        "status_counts": "Status Counts",
        "title": "Inclusion Review",
    }


def _render_inclusion_review(
    case_dir: Path,
    manifest: Mapping[str, Any],
    materials: Sequence[Mapping[str, Any]],
    entries: Sequence[Mapping[str, Any]],
    *,
    generated_at: str,
) -> str:
    """Render a mechanical status checklist without changing inclusion state."""

    material_labels = _material_label_by_id(materials)
    copy = _inclusion_review_copy(_ui_language(manifest))
    pending_entries = _review_status_entries(entries, "pending")
    approved_entries = _review_status_entries(entries, "approved")
    rejected_entries = _review_status_entries(entries, "rejected")
    pending_item_numbers = {
        str(entry["id"]): index for index, entry in enumerate(pending_entries, start=1)
    }
    pending_bundles = _pending_inclusion_bundle_groups(case_dir, pending_entries)
    lines = [
        f"# {copy['title']} - {manifest['client']}",
        "",
        copy["intro"],
        copy["intro_gate"],
        "",
        f"{copy['generated']}: {generated_at}",
        f"{copy['project']}: {manifest['project']}",
        f"{copy['audience']}: {manifest['audience']}",
        "",
        f"## {copy['status_counts']}",
        "",
        f"- {copy['pending_count']}: {len(pending_entries)}",
        f"- {copy['bundles_count']}: {len(pending_bundles)}",
        f"- {copy['already_approved_count']}: {len(approved_entries)}",
        f"- {copy['excluded_count']}: {len(rejected_entries)}",
        "",
        f"## {copy['how_to_respond']}",
        "",
        f"- {copy['include_all']}: `{copy['include_all_example']}`.",
        f"- {copy['include_bundle']}: `{copy['include_bundle_example']}`.",
        f"- {copy['include_one']}: `{copy['include_one_example']}`.",
        f"- {copy['exclude_one']}: `{copy['exclude_one_example']}`.",
        f"- {copy['correct_one']}: `{copy['correct_one_example']}`.",
        f"- {copy['ask_evidence']}: `{copy['show_more_example']}`.",
        "",
        f"## {copy['bundles']}",
        "",
        *_inclusion_bundle_lines(
            pending_bundles,
            pending_item_numbers,
            copy=copy,
        ),
        f"## {copy['pending']}",
        "",
        *_inclusion_entry_lines(
            pending_entries,
            material_labels,
            include_actions=True,
            copy=copy,
        ),
        f"## {copy['already_approved']}",
        "",
        *_inclusion_entry_lines(
            approved_entries,
            material_labels,
            include_actions=False,
            copy=copy,
        ),
        f"## {copy['excluded']}",
        "",
        *_inclusion_entry_lines(
            rejected_entries,
            material_labels,
            include_actions=False,
            copy=copy,
        ),
        f"## {copy['control_heading']}",
        "",
        f"- {copy['control_pending']}",
        f"- {copy['control_regenerated']}",
        f"- {copy['create_bundles']}",
        "",
    ]
    return "\n".join(lines)


def build_inclusion_review(
    case_dir: Path,
    *,
    output_path: Path | None = None,
    now: datetime | None = None,
) -> InclusionReviewResult:
    """Write a deterministic inclusion checklist without mutating case status.

    The inclusion gate is audit-sensitive mechanical work: this renderer only
    mirrors canonical JSON state and records actionable item numbers.
    """

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    manifest = _read_json(_case_path(case_dir, "manifest"))
    materials = _read_json(_case_path(case_dir, "materials"))["materials"]
    entries = _read_json(_case_path(case_dir, "judgement"))["entries"]
    review_path = output_path or case_dir / INCLUSION_REVIEW_FILENAME
    review_path.parent.mkdir(parents=True, exist_ok=True)
    review_text = _render_inclusion_review(
        case_dir,
        manifest=manifest,
        materials=materials,
        entries=entries,
        generated_at=_now_iso(now),
    )
    _assert_human_visible_document_quality(review_text, label=INCLUSION_REVIEW_FILENAME)
    review_path.write_text(review_text, encoding="utf-8")

    return InclusionReviewResult(
        review_path=review_path,
        approved_count=sum(1 for entry in entries if entry["status"] == "approved"),
        pending_count=sum(1 for entry in entries if entry["status"] == "pending"),
        rejected_count=sum(1 for entry in entries if entry["status"] == "rejected"),
    )


def _group_approved_entries(
    entries: Sequence[Mapping[str, Any]],
) -> dict[str, list[Mapping[str, Any]]]:
    grouped = {kind: [] for kind in sorted(JUDGEMENT_KINDS)}
    for entry in entries:
        if entry["status"] == "approved":
            grouped[entry["kind"]].append(entry)
    return grouped


def _material_label_by_id(materials: Sequence[Mapping[str, Any]]) -> dict[str, str]:
    return {item["id"]: str(item["title"]) for item in materials}


def _entry_line(
    entry: Mapping[str, Any],
    material_labels: Mapping[str, str],
    *,
    include_sources: bool = True,
    sources_label: str = "sources",
) -> str:
    text = str(entry["text"]).strip()
    source_ids = entry.get("source_material_ids", [])
    labels = [
        material_labels[source_id]
        for source_id in source_ids
        if source_id in material_labels
    ]
    if include_sources and labels:
        return f"{text} ({sources_label}: {', '.join(labels)})"
    return text


def _markdown_section(
    title: str,
    entries: Sequence[Mapping[str, Any]],
    material_labels: Mapping[str, str],
    *,
    no_items: str = "No decision-pack-ready items.",
    sources_label: str = "sources",
) -> list[str]:
    lines = [f"## {title}", ""]
    if not entries:
        lines.extend([no_items, ""])
        return lines
    lines.extend(
        f"- {_entry_line(entry, material_labels, sources_label=sources_label)}"
        for entry in entries
    )
    lines.append("")
    return lines


LANGUAGE_COPY = {
    "en": {
        "title": "Decision Pack",
        "project": "Project",
        "objective": "Objective",
        "audience": "Audience",
        "status": "Draft Status",
        "ready_note": "This draft uses only judgement marked ready for the client pack.",
        "not_ready": "The client pack is not ready yet. Review the candidate judgement entries before sharing this output.",
        "facts": "Facts",
        "storyline": "Executive Storyline",
        "advisor": "Consultant Readout",
        "inferences": "Analytical Implications",
        "decisions": "Decision Implications",
        "questions": "Open Questions To Resolve",
        "no_items": "No reviewed items yet.",
        "why": "Why it matters",
        "workpaper_note": "Detailed provenance and inclusion control are in decision_pack_workpaper.md.",
        "workpaper_title": "Decision Pack Workpaper",
        "materials": "Materials Indexed",
        "control": "Inclusion Control",
        "used": "Decision-pack-ready judgement entries used",
        "pending": "Pending judgement entries excluded",
        "rejected": "Rejected judgement entries excluded",
    },
    "it": {
        "title": "Decision Pack",
        "project": "Progetto",
        "objective": "Obiettivo",
        "audience": "Destinatario",
        "status": "Stato della bozza",
        "ready_note": "Questa bozza usa solo giudizi segnati come pronti per il pack cliente.",
        "not_ready": "Il pack cliente non e' ancora pronto. Rivedere le voci candidate prima di condividere questo output.",
        "facts": "Fatti acquisiti",
        "storyline": "Traccia esecutiva",
        "advisor": "Lettura consulenziale",
        "inferences": "Implicazioni analitiche",
        "decisions": "Implicazioni per la decisione",
        "questions": "Questioni aperte da risolvere",
        "no_items": "Nessuna voce rivista.",
        "why": "Perche conta",
        "workpaper_note": "Provenienza dettagliata e controlli di inclusione sono in decision_pack_workpaper.md.",
        "workpaper_title": "Workpaper del decision pack",
        "materials": "Materiali indicizzati",
        "control": "Controllo inclusione",
        "used": "Voci giudizio pronte usate",
        "pending": "Voci pending escluse",
        "rejected": "Voci escluse",
    },
    "fr": {
        "title": "Decision Pack",
        "project": "Projet",
        "objective": "Objectif",
        "audience": "Audience",
        "status": "Statut du brouillon",
        "ready_note": "Ce brouillon utilise uniquement les jugements marques comme prets pour le client.",
        "not_ready": "Le pack client n'est pas encore pret. Revoyez les elements candidats avant partage.",
        "facts": "Faits etablis",
        "storyline": "Fil conducteur executif",
        "advisor": "Lecture conseil",
        "inferences": "Implications analytiques",
        "decisions": "Implications pour la decision",
        "questions": "Questions ouvertes",
        "no_items": "Aucun element revu.",
        "why": "Pourquoi c'est important",
        "workpaper_note": "La provenance detaillee et le controle d'inclusion sont dans decision_pack_workpaper.md.",
        "workpaper_title": "Workpaper du decision pack",
        "materials": "Documents indexes",
        "control": "Controle d'inclusion",
        "used": "Jugements prets utilises",
        "pending": "Jugements en attente exclus",
        "rejected": "Jugements rejetes exclus",
    },
    "de": {
        "title": "Decision Pack",
        "project": "Projekt",
        "objective": "Ziel",
        "audience": "Zielgruppe",
        "status": "Entwurfsstatus",
        "ready_note": "Dieser Entwurf nutzt nur fuer den Client-Pack freigegebene Einschaetzungen.",
        "not_ready": "Der Client-Pack ist noch nicht bereit. Pruefen Sie die Kandidaten vor der Weitergabe.",
        "facts": "Gesicherte Fakten",
        "storyline": "Executive Storyline",
        "advisor": "Beraterliche Einordnung",
        "inferences": "Analytische Implikationen",
        "decisions": "Entscheidungsimplikationen",
        "questions": "Offene Fragen",
        "no_items": "Noch keine geprueften Punkte.",
        "why": "Warum es wichtig ist",
        "workpaper_note": "Detaillierte Herkunft und Inklusionskontrolle stehen in decision_pack_workpaper.md.",
        "workpaper_title": "Decision-Pack-Workpaper",
        "materials": "Indexierte Materialien",
        "control": "Inklusionskontrolle",
        "used": "Verwendete freigegebene Einschaetzungen",
        "pending": "Ausgeschlossene pending Einschaetzungen",
        "rejected": "Ausgeschlossene abgelehnte Einschaetzungen",
    },
    "es": {
        "title": "Paquete de decisión",
        "project": "Proyecto",
        "objective": "Objetivo",
        "audience": "Destinatario",
        "status": "Estado del borrador",
        "ready_note": (
            "Este borrador utiliza únicamente juicios marcados como listos para "
            "el paquete del cliente."
        ),
        "not_ready": (
            "El paquete del cliente aún no está listo. Revisa los elementos "
            "candidatos antes de compartir este entregable."
        ),
        "facts": "Hechos",
        "storyline": "Línea argumental ejecutiva",
        "advisor": "Lectura consultiva",
        "inferences": "Implicaciones analíticas",
        "decisions": "Implicaciones para la decisión",
        "questions": "Preguntas abiertas que resolver",
        "no_items": "Aún no hay elementos revisados.",
        "why": "Por qué es importante",
        "workpaper_note": (
            "La procedencia detallada y el control de inclusión están en "
            "decision_pack_workpaper.md."
        ),
        "workpaper_title": "Documento de trabajo del paquete de decisión",
        "materials": "Materiales indexados",
        "control": "Control de inclusión",
        "used": "Elementos de juicio listos utilizados",
        "pending": "Elementos de juicio pendientes excluidos",
        "rejected": "Elementos de juicio rechazados excluidos",
    },
}


def _copy(manifest: Mapping[str, Any], key: str) -> str:
    language = str(manifest.get("output_language") or "en")
    return LANGUAGE_COPY.get(language, LANGUAGE_COPY["en"])[key]


def _workpaper_output_copy(manifest: Mapping[str, Any]) -> dict[str, str]:
    if _ui_language(manifest) == "es":
        return {
            "explanation": (
                "Este documento de trabajo incluye las rutas de las fuentes y el "
                "control de inclusión. No está destinado al cliente."
            ),
            "no_materials": "No hay materiales indexados.",
            "no_questions": "No se han registrado preguntas abiertas.",
            "no_items": "No hay elementos listos para el paquete de decisión.",
            "sources": "fuentes",
            "title_header": "Título",
            "type_header": "Tipo",
            "status_header": "Estado",
        }
    return {
        "explanation": (
            "This workpaper includes source paths and inclusion control. It is not "
            "client-facing."
        ),
        "no_materials": "No materials indexed.",
        "no_questions": "No open questions registered.",
        "no_items": "No decision-pack-ready items.",
        "sources": "sources",
        "title_header": "Title",
        "type_header": "Type",
        "status_header": "Status",
    }


def _client_section(
    title: str,
    entries: Sequence[Mapping[str, Any]],
    material_labels: Mapping[str, str],
    manifest: Mapping[str, Any],
) -> list[str]:
    lines = [f"## {title}", ""]
    if not entries:
        lines.extend([_copy(manifest, "no_items"), ""])
        return lines
    lines.extend(
        f"- {_entry_line(entry, material_labels, include_sources=False)}"
        for entry in entries
    )
    lines.append("")
    return lines


def _entry_text_values(entries: Sequence[Mapping[str, Any]], limit: int) -> list[str]:
    values: list[str] = []
    for entry in entries[:limit]:
        text = str(entry.get("text", "")).strip()
        if text:
            values.append(text)
    return values


def _join_story_sentences(values: Sequence[str]) -> str:
    text = " ".join(value.rstrip(".") + "." for value in values if value.strip())
    return " ".join(text.split())


def _storyline_paragraphs(
    manifest: Mapping[str, Any],
    grouped: Mapping[str, Sequence[Mapping[str, Any]]],
    open_questions: Sequence[Mapping[str, Any]],
) -> list[str]:
    fact_text = _join_story_sentences(_entry_text_values(grouped["fact"], 2))
    advisor_text = _join_story_sentences(
        _entry_text_values(grouped["advisor_judgement"], 2)
    )
    decision_text = _join_story_sentences(
        [
            *_entry_text_values(grouped["decision_implication"], 1),
            *_entry_text_values(grouped["codex_inference"], 1),
        ]
    )
    open_items = [item for item in open_questions if item["status"] == "open"]
    first_open_question = str(open_items[0]["question"]).strip() if open_items else ""
    if not any((fact_text, advisor_text, decision_text)):
        return []

    language = _ui_language(manifest)
    if language == "it":
        paragraphs: list[str] = []
        if fact_text:
            paragraphs.append(f"Punto di partenza. {fact_text}")
        if advisor_text:
            paragraphs.append(f"Lettura del caso. {advisor_text}")
        if decision_text:
            paragraphs.append(f"Percorso consigliato. {decision_text}")
        if first_open_question:
            paragraphs.append(
                "Prima di finalizzare la raccomandazione, resta da chiudere "
                f"questo punto: {first_open_question}"
            )
        return paragraphs
    if language == "es":
        paragraphs = []
        if fact_text:
            paragraphs.append(f"Punto de partida. {fact_text}")
        if advisor_text:
            paragraphs.append(f"Lectura del caso. {advisor_text}")
        if decision_text:
            paragraphs.append(f"Ruta recomendada. {decision_text}")
        if first_open_question:
            paragraphs.append(
                "Antes de cerrar la recomendación, queda por resolver este punto: "
                f"{first_open_question}"
            )
        return paragraphs

    paragraphs = []
    if fact_text:
        paragraphs.append(f"Starting point. {fact_text}")
    if advisor_text:
        paragraphs.append(f"Case readout. {advisor_text}")
    if decision_text:
        paragraphs.append(f"Recommended path. {decision_text}")
    if first_open_question:
        paragraphs.append(
            "Before finalizing the recommendation, this point remains open: "
            f"{first_open_question}"
        )
    return paragraphs


def _render_markdown(
    manifest: Mapping[str, Any],
    materials: Sequence[Mapping[str, Any]],
    entries: Sequence[Mapping[str, Any]],
    open_questions: Sequence[Mapping[str, Any]],
) -> str:
    grouped = _group_approved_entries(entries)
    material_labels = _material_label_by_id(materials)
    pending_count = sum(1 for entry in entries if entry["status"] == "pending")
    rejected_count = sum(1 for entry in entries if entry["status"] == "rejected")
    approved_count = sum(1 for entry in entries if entry["status"] == "approved")

    lines = [
        f"# {manifest['client']} - {_copy(manifest, 'title')}",
        "",
        f"**{_copy(manifest, 'project')}:** {manifest['project']}",
        f"**{_copy(manifest, 'objective')}:** {manifest['objective']}",
        f"**{_copy(manifest, 'audience')}:** {manifest['audience']}",
        "",
        f"## {_copy(manifest, 'status')}",
        "",
        _copy(manifest, "ready_note"),
        "",
    ]
    if approved_count == 0:
        lines.extend([_copy(manifest, "not_ready"), ""])

    storyline = _storyline_paragraphs(manifest, grouped, open_questions)
    if storyline:
        lines.extend([f"## {_copy(manifest, 'storyline')}", ""])
        for paragraph in storyline:
            lines.extend([paragraph, ""])

    sections = [
        (_copy(manifest, "facts"), grouped["fact"]),
        (_copy(manifest, "advisor"), grouped["advisor_judgement"]),
        (_copy(manifest, "inferences"), grouped["codex_inference"]),
        (_copy(manifest, "decisions"), grouped["decision_implication"]),
    ]
    for title, section_entries in sections:
        lines.extend(_client_section(title, section_entries, material_labels, manifest))

    lines.extend([f"## {_copy(manifest, 'questions')}", ""])
    open_items = [item for item in open_questions if item["status"] == "open"]
    if open_items:
        lines.extend(
            f"- {item['question']} ({_copy(manifest, 'why')}: {item['why_it_matters']})"
            for item in open_items
        )
    else:
        lines.append(_copy(manifest, "no_items"))
    lines.extend(
        [
            "",
            f"_{_copy(manifest, 'workpaper_note')}_",
            "",
        ]
    )
    return "\n".join(lines)


def _render_workpaper_markdown(
    manifest: Mapping[str, Any],
    materials: Sequence[Mapping[str, Any]],
    entries: Sequence[Mapping[str, Any]],
    open_questions: Sequence[Mapping[str, Any]],
) -> str:
    workpaper_copy = _workpaper_output_copy(manifest)
    grouped = _group_approved_entries(entries)
    material_labels = _material_label_by_id(materials)
    pending_count = sum(1 for entry in entries if entry["status"] == "pending")
    rejected_count = sum(1 for entry in entries if entry["status"] == "rejected")
    approved_count = sum(1 for entry in entries if entry["status"] == "approved")

    lines = [
        f"# {_copy(manifest, 'workpaper_title')} - {manifest['client']}",
        "",
        f"{_copy(manifest, 'project')}: {manifest['project']}",
        f"{_copy(manifest, 'objective')}: {manifest['objective']}",
        f"{_copy(manifest, 'audience')}: {manifest['audience']}",
        "",
        workpaper_copy["explanation"],
        "",
        f"## {_copy(manifest, 'materials')}",
        "",
    ]
    if materials:
        lines.extend(
            f"- {item['title']} [{item['material_type']}, {item['status']}] - {item['path']}"
            for item in materials
        )
    else:
        lines.append(workpaper_copy["no_materials"])
    lines.append("")

    sections = [
        (_copy(manifest, "facts"), grouped["fact"]),
        (_copy(manifest, "advisor"), grouped["advisor_judgement"]),
        (_copy(manifest, "inferences"), grouped["codex_inference"]),
        (_copy(manifest, "decisions"), grouped["decision_implication"]),
    ]
    for title, section_entries in sections:
        lines.extend(
            _markdown_section(
                title,
                section_entries,
                material_labels,
                no_items=workpaper_copy["no_items"],
                sources_label=workpaper_copy["sources"],
            )
        )

    lines.extend([f"## {_copy(manifest, 'questions')}", ""])
    open_items = [item for item in open_questions if item["status"] == "open"]
    if open_items:
        lines.extend(
            f"- {item['question']} {_copy(manifest, 'why')}: {item['why_it_matters']}"
            for item in open_items
        )
    else:
        lines.append(workpaper_copy["no_questions"])
    lines.extend(
        [
            "",
            f"## {_copy(manifest, 'control')}",
            "",
            f"- {_copy(manifest, 'used')}: {approved_count}",
            f"- {_copy(manifest, 'pending')}: {pending_count}",
            f"- {_copy(manifest, 'rejected')}: {rejected_count}",
            "",
        ]
    )
    return "\n".join(lines)


def _add_docx_bullets(
    document: Any,
    items: Sequence[str],
    *,
    empty_text: str = "No decision-pack-ready items.",
) -> None:
    if not items:
        document.add_paragraph(empty_text)
        return
    for item in items:
        document.add_paragraph(item, style="List Bullet")


def _render_docx(
    path: Path,
    manifest: Mapping[str, Any],
    materials: Sequence[Mapping[str, Any]],
    entries: Sequence[Mapping[str, Any]],
    open_questions: Sequence[Mapping[str, Any]],
) -> None:
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Pt, RGBColor
    except ImportError as exc:
        raise CaseWorkspaceError(
            "python-docx is required to render decision_pack.docx"
        ) from exc

    grouped = _group_approved_entries(entries)
    workpaper_copy = _workpaper_output_copy(manifest)
    material_labels = _material_label_by_id(materials)
    approved_count = sum(1 for entry in entries if entry["status"] == "approved")
    document = Document()
    styles = document.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(10.5)
    styles["Title"].font.color.rgb = RGBColor(49, 95, 85)

    title = document.add_paragraph()
    title.style = styles["Title"]
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    title.add_run(f"{manifest['client']} - {_copy(manifest, 'title')}")

    for label, value in (
        (_copy(manifest, "project"), manifest["project"]),
        (_copy(manifest, "objective"), manifest["objective"]),
        (_copy(manifest, "audience"), manifest["audience"]),
    ):
        paragraph = document.add_paragraph()
        paragraph.add_run(f"{label}: ").bold = True
        paragraph.add_run(str(value))

    document.add_heading(_copy(manifest, "status"), level=1)
    document.add_paragraph(_copy(manifest, "ready_note"))
    if approved_count == 0:
        document.add_paragraph(_copy(manifest, "not_ready"))

    storyline = _storyline_paragraphs(manifest, grouped, open_questions)
    if storyline:
        document.add_heading(_copy(manifest, "storyline"), level=1)
        for paragraph_text in storyline:
            document.add_paragraph(paragraph_text)

    for title_text, section_entries in (
        (_copy(manifest, "facts"), grouped["fact"]),
        (_copy(manifest, "advisor"), grouped["advisor_judgement"]),
        (_copy(manifest, "inferences"), grouped["codex_inference"]),
        (_copy(manifest, "decisions"), grouped["decision_implication"]),
    ):
        document.add_heading(title_text, level=1)
        _add_docx_bullets(
            document,
            [
                _entry_line(entry, material_labels, include_sources=False)
                for entry in section_entries
            ],
            empty_text=workpaper_copy["no_items"],
        )

    document.add_heading(_copy(manifest, "questions"), level=1)
    open_items = [item for item in open_questions if item["status"] == "open"]
    _add_docx_bullets(
        document,
        [
            f"{item['question']} ({_copy(manifest, 'why')}: {item['why_it_matters']})"
            for item in open_items
        ],
        empty_text=workpaper_copy["no_questions"],
    )
    document.add_paragraph(_copy(manifest, "workpaper_note"))

    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def _render_workpaper_docx(
    path: Path,
    manifest: Mapping[str, Any],
    materials: Sequence[Mapping[str, Any]],
    entries: Sequence[Mapping[str, Any]],
    open_questions: Sequence[Mapping[str, Any]],
) -> None:
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Pt, RGBColor
    except ImportError as exc:
        raise CaseWorkspaceError(
            "python-docx is required to render decision_pack_workpaper.docx"
        ) from exc

    grouped = _group_approved_entries(entries)
    workpaper_copy = _workpaper_output_copy(manifest)
    material_labels = _material_label_by_id(materials)
    document = Document()
    styles = document.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(10.5)
    styles["Title"].font.color.rgb = RGBColor(49, 95, 85)

    title = document.add_paragraph()
    title.style = styles["Title"]
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    title.add_run(f"{_copy(manifest, 'workpaper_title')} - {manifest['client']}")

    for label, value in (
        (_copy(manifest, "project"), manifest["project"]),
        (_copy(manifest, "objective"), manifest["objective"]),
        (_copy(manifest, "audience"), manifest["audience"]),
    ):
        paragraph = document.add_paragraph()
        paragraph.add_run(f"{label}: ").bold = True
        paragraph.add_run(str(value))

    document.add_paragraph(workpaper_copy["explanation"])
    document.add_heading(_copy(manifest, "materials"), level=1)
    if materials:
        table = document.add_table(rows=1, cols=3)
        table.style = "Table Grid"
        header = table.rows[0].cells
        header[0].text = workpaper_copy["title_header"]
        header[1].text = workpaper_copy["type_header"]
        header[2].text = workpaper_copy["status_header"]
        for item in materials:
            row = table.add_row().cells
            row[0].text = str(item["title"])
            row[1].text = str(item["material_type"])
            row[2].text = str(item["status"])
    else:
        document.add_paragraph(workpaper_copy["no_materials"])

    for title_text, section_entries in (
        (_copy(manifest, "facts"), grouped["fact"]),
        (_copy(manifest, "advisor"), grouped["advisor_judgement"]),
        (_copy(manifest, "inferences"), grouped["codex_inference"]),
        (_copy(manifest, "decisions"), grouped["decision_implication"]),
    ):
        document.add_heading(title_text, level=1)
        _add_docx_bullets(
            document,
            [
                _entry_line(
                    entry,
                    material_labels,
                    sources_label=workpaper_copy["sources"],
                )
                for entry in section_entries
            ],
            empty_text=workpaper_copy["no_items"],
        )

    document.add_heading(_copy(manifest, "questions"), level=1)
    open_items = [item for item in open_questions if item["status"] == "open"]
    _add_docx_bullets(
        document,
        [
            f"{item['question']} {_copy(manifest, 'why')}: {item['why_it_matters']}"
            for item in open_items
        ],
        empty_text=workpaper_copy["no_questions"],
    )

    pending_count = sum(1 for entry in entries if entry["status"] == "pending")
    rejected_count = sum(1 for entry in entries if entry["status"] == "rejected")
    approved_count = sum(1 for entry in entries if entry["status"] == "approved")
    document.add_heading(_copy(manifest, "control"), level=1)
    _add_docx_bullets(
        document,
        [
            f"{_copy(manifest, 'used')}: {approved_count}",
            f"{_copy(manifest, 'pending')}: {pending_count}",
            f"{_copy(manifest, 'rejected')}: {rejected_count}",
        ],
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    document.save(path)


def build_decision_pack(
    case_dir: Path,
    *,
    output_dir: Path | None = None,
) -> DecisionPackResult:
    """Render Markdown and Word packs using client-pack-ready judgement only."""

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))
    refresh_case_brief(case_dir)

    manifest = _read_json(_case_path(case_dir, "manifest"))
    materials = _read_json(_case_path(case_dir, "materials"))["materials"]
    entries = _read_json(_case_path(case_dir, "judgement"))["entries"]
    questions = _read_json(_case_path(case_dir, "open_questions"))["questions"]
    target_dir = (
        output_dir
        or case_dir.expanduser().resolve().parent / "output" / "clara" / case_dir.name
    )
    target_dir.mkdir(parents=True, exist_ok=True)

    markdown_path = target_dir / "decision_pack.md"
    docx_path = target_dir / "decision_pack.docx"
    workpaper_markdown_path = target_dir / "decision_pack_workpaper.md"
    workpaper_docx_path = target_dir / "decision_pack_workpaper.docx"
    markdown_text = _render_markdown(manifest, materials, entries, questions)
    _assert_human_visible_document_quality(markdown_text, label="decision_pack.md")
    markdown_path.write_text(markdown_text, encoding="utf-8")
    workpaper_markdown_path.write_text(
        _render_workpaper_markdown(manifest, materials, entries, questions),
        encoding="utf-8",
    )
    _render_docx(docx_path, manifest, materials, entries, questions)
    _render_workpaper_docx(
        workpaper_docx_path,
        manifest,
        materials,
        entries,
        questions,
    )

    return DecisionPackResult(
        markdown_path=markdown_path,
        docx_path=docx_path,
        workpaper_markdown_path=workpaper_markdown_path,
        workpaper_docx_path=workpaper_docx_path,
        approved_count=sum(1 for entry in entries if entry["status"] == "approved"),
        pending_count=sum(1 for entry in entries if entry["status"] == "pending"),
        rejected_count=sum(1 for entry in entries if entry["status"] == "rejected"),
    )
