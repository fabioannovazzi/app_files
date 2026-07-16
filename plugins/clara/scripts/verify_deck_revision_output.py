"""Verify a corrected PPTX against a Clara deck-revision plan."""

from __future__ import annotations

import argparse
import json
import logging
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Mapping

from advisor_case_core import CaseWorkspaceError, validate_case_workspace
from deck_revision_text_match import target_text_matches

__all__ = [
    "DeckRevisionVerificationResult",
    "verify_deck_revision_output",
    "main",
]

LOGGER = logging.getLogger(__name__)


@dataclass(frozen=True)
class DeckRevisionVerificationResult:
    """Verification artifacts for a corrected deck."""

    session_dir: Path
    report_path: Path
    review_path: Path


def _now_iso(now: datetime | None = None) -> str:
    value = now or datetime.now(timezone.utc)
    return value.astimezone(timezone.utc).replace(microsecond=0).isoformat()


def _read_json(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise CaseWorkspaceError(f"expected JSON object in {path}")
    return payload


def _write_json(path: Path, payload: Mapping[str, Any]) -> None:
    path.write_text(
        json.dumps(payload, indent=2, ensure_ascii=True) + "\n",
        encoding="utf-8",
    )


def _relative_path(case_dir: Path, path: Path) -> str:
    try:
        return str(path.resolve().relative_to(case_dir.resolve()))
    except ValueError:
        return str(path.resolve())


def _resolve_voice_session_dir(case_dir: Path, voice_session: Path | None) -> Path:
    sessions_root = case_dir / "voice_sessions"
    if voice_session is None:
        if not sessions_root.is_dir():
            raise CaseWorkspaceError("case has no voice_sessions folder")
        sessions = sorted(path for path in sessions_root.iterdir() if path.is_dir())
        if not sessions:
            raise CaseWorkspaceError("case has no imported voice sessions")
        return sessions[-1].resolve()

    candidate = voice_session.expanduser()
    candidates = (
        [candidate]
        if candidate.is_absolute()
        else [case_dir / candidate, sessions_root / candidate]
    )
    for path in candidates:
        if path.is_dir():
            resolved = path.resolve()
            try:
                resolved.relative_to(sessions_root.resolve())
            except ValueError as error:
                raise CaseWorkspaceError(
                    f"voice session must live under {sessions_root}: {resolved}"
                ) from error
            return resolved
    raise CaseWorkspaceError(f"voice session does not exist: {voice_session}")


def _resolve_case_path(case_dir: Path, path: Path) -> Path:
    candidate = path.expanduser()
    if not candidate.is_absolute():
        candidate = case_dir / candidate
    if not candidate.is_file():
        raise CaseWorkspaceError(f"file does not exist: {candidate}")
    return candidate.resolve()


def _load_plan(
    session_dir: Path, plan_path: Path | None, case_dir: Path
) -> tuple[Path, dict[str, Any]]:
    candidate = plan_path or (session_dir / "deck_revision_changes.normalized.json")
    return _resolve_case_path(case_dir, candidate), _read_json(
        _resolve_case_path(case_dir, candidate)
    )


def _shape_text(shape: Any) -> str:
    text = str(getattr(shape, "text", "") or "")
    if text:
        return text
    table = getattr(shape, "table", None)
    if table is None:
        return ""
    values: list[str] = []
    for row in getattr(table, "rows", []):
        for cell in getattr(row, "cells", []):
            cell_text = str(getattr(cell, "text", "") or "").strip()
            if cell_text:
                values.append(cell_text)
    return "\n".join(values)


def _slide_text(slide: Any) -> str:
    return "\n".join(_shape_text(shape) for shape in slide.shapes if _shape_text(shape))


def _shape_by_index(slide: Any, shape_index: int) -> Any | None:
    shapes = list(slide.shapes)
    if shape_index < 1 or shape_index > len(shapes):
        return None
    return shapes[shape_index - 1]


def _target_text_matches(shape: Any, expected_text: str) -> bool:
    return target_text_matches(_shape_text(shape), expected_text)


def _verify_patch(slide: Any, patch: Mapping[str, Any]) -> dict[str, Any]:
    patch_id = patch.get("patch_id")
    operation = patch.get("operation")
    target = patch.get("target") if isinstance(patch.get("target"), dict) else {}
    value = patch.get("value") if isinstance(patch.get("value"), dict) else {}
    expected = ""
    actual = ""
    passed = False
    note = ""

    if operation == "set_title_text":
        expected = str(value.get("text", ""))
        title_shape = getattr(slide.shapes, "title", None)
        actual = _shape_text(title_shape) if title_shape is not None else ""
        passed = actual == expected
    elif operation == "set_shape_text":
        expected = str(value.get("text", ""))
        shape = _shape_by_index(slide, int(target.get("shape_index", 0)))
        actual = _shape_text(shape) if shape is not None else ""
        passed = actual == expected
    elif operation == "replace_text":
        shape = _shape_by_index(slide, int(target.get("shape_index", 0)))
        if shape is None:
            note = "target shape is missing"
            text = ""
        else:
            text = _shape_text(shape)
        old_text = str(value.get("old_text", ""))
        new_text = str(value.get("new_text", ""))
        expected = f"`{old_text}` replaced with `{new_text}`"
        actual = text
        passed = (
            shape is not None
            and bool(new_text)
            and new_text in text
            and old_text not in text
        )
    elif operation == "add_textbox":
        expected = str(value.get("text", ""))
        actual = _slide_text(slide)
        passed = bool(expected) and expected in actual
    elif operation == "delete_shape":
        expected_absent = str(value.get("expected_absent_text", ""))
        if expected_absent:
            expected = f"`{expected_absent}` absent"
            actual = _slide_text(slide)
            passed = expected_absent not in actual
        else:
            note = (
                "delete_shape has no expected_absent_text; verify deletion visually "
                "or provide a stronger patch assertion"
            )
    elif operation == "move_shape":
        shape = _shape_by_index(slide, int(target.get("shape_index", 0)))
        if shape is None:
            note = "target shape is missing"
        else:
            expected_parts: list[str] = []
            actual_parts: list[str] = []
            passed = True
            expected_text = str(target.get("expected_text", ""))
            if expected_text and not _target_text_matches(shape, expected_text):
                note = (
                    "target text mismatch after move; verify that the intended "
                    "shape was moved"
                )
                passed = False
            if "left" in value:
                expected_parts.append(f"left={value['left']}")
                actual_parts.append(f"left={int(getattr(shape, 'left', 0) or 0)}")
                passed = passed and int(getattr(shape, "left", 0) or 0) == value["left"]
            if "top" in value:
                expected_parts.append(f"top={value['top']}")
                actual_parts.append(f"top={int(getattr(shape, 'top', 0) or 0)}")
                passed = passed and int(getattr(shape, "top", 0) or 0) == value["top"]
            expected = ", ".join(expected_parts)
            actual = ", ".join(actual_parts)
    else:
        note = f"unsupported operation `{operation}`"

    return {
        "patch_id": patch_id,
        "operation": operation,
        "passed": passed,
        "expected": expected,
        "actual": actual,
        "note": note,
    }


def _verify_success_criterion(
    presentation: Any,
    slide: Any | None,
    criterion: Mapping[str, Any],
) -> dict[str, Any]:
    criterion_id = criterion.get("criterion_id")
    check_type = criterion.get("check_type")
    expected = ""
    actual = ""
    passed = False
    note = ""
    review_required = False

    if check_type == "title_equals":
        expected = str(criterion.get("expected_text", ""))
        title_shape = (
            getattr(slide.shapes, "title", None) if slide is not None else None
        )
        actual = _shape_text(title_shape) if title_shape is not None else ""
        passed = bool(expected) and actual == expected
    elif check_type == "text_present":
        expected = str(criterion.get("expected_text", ""))
        actual = _slide_text(slide) if slide is not None else ""
        passed = bool(expected) and expected in actual
    elif check_type == "text_absent":
        expected = str(
            criterion.get("absent_text") or criterion.get("expected_text") or ""
        )
        actual = _slide_text(slide) if slide is not None else ""
        passed = bool(expected) and expected not in actual
    elif check_type == "slide_count_equals":
        expected_count = criterion.get("expected_count")
        expected = str(expected_count)
        actual = str(len(presentation.slides))
        passed = (
            isinstance(expected_count, int)
            and len(presentation.slides) == expected_count
        )
    elif check_type == "shape_position":
        shape = (
            _shape_by_index(slide, int(criterion.get("shape_index", 0)))
            if slide is not None
            else None
        )
        if shape is None:
            note = "target shape is missing"
        else:
            expected_parts: list[str] = []
            actual_parts: list[str] = []
            passed = True
            if "expected_left" in criterion:
                expected_left = int(criterion["expected_left"])
                actual_left = int(getattr(shape, "left", 0) or 0)
                expected_parts.append(f"left={expected_left}")
                actual_parts.append(f"left={actual_left}")
                passed = passed and actual_left == expected_left
            if "expected_top" in criterion:
                expected_top = int(criterion["expected_top"])
                actual_top = int(getattr(shape, "top", 0) or 0)
                expected_parts.append(f"top={expected_top}")
                actual_parts.append(f"top={actual_top}")
                passed = passed and actual_top == expected_top
            expected = ", ".join(expected_parts)
            actual = ", ".join(actual_parts)
    elif check_type == "manual_review":
        note = criterion.get("note") or "manual review required"
        review_required = True
    elif check_type == "semantic_review":
        note = criterion.get("note") or "model-assisted semantic review required"
        review_required = True
    else:
        note = f"unsupported success criterion `{check_type}`"

    if slide is None and check_type not in {
        "slide_count_equals",
        "manual_review",
        "semantic_review",
    }:
        note = note or "slide is missing in corrected deck"

    return {
        "criterion_id": criterion_id,
        "check_type": check_type,
        "passed": passed,
        "expected": expected,
        "actual": actual,
        "note": note,
        "review_required": review_required,
    }


def _render_markdown(payload: Mapping[str, Any]) -> str:
    lines = [
        "# Deck Revision Verification",
        "",
        f"Status: `{payload['summary']['status']}`",
        "",
        f"- Corrected deck: `{payload['corrected_deck_path']}`",
        f"- Passed patches: {payload['summary']['passed_patches']}",
        f"- Failed patches: {payload['summary']['failed_patches']}",
        f"- Manual-review patches: {payload['summary']['manual_review_patches']}",
        f"- Passed success criteria: {payload['summary']['passed_success_criteria']}",
        f"- Failed success criteria: {payload['summary']['failed_success_criteria']}",
        f"- Manual-review success criteria: {payload['summary']['manual_review_success_criteria']}",
        "",
    ]
    for change in payload["changes"]:
        lines.extend(
            [
                f"## Slide {change['slide_number']} - {change['change_id']}",
                "",
                f"Status: `{change['status']}`",
                "",
            ]
        )
        for patch in change["patches"]:
            result = "passed" if patch["passed"] else "failed"
            if patch["note"] and not patch["passed"]:
                result = "manual_review"
            lines.append(f"- `{patch['patch_id']}` `{patch['operation']}`: {result}")
            if patch["note"]:
                lines.append(f"  - Note: {patch['note']}")
        if change["success_criteria"]:
            lines.append("Success criteria:")
            for criterion in change["success_criteria"]:
                result = "passed" if criterion["passed"] else "failed"
                if criterion.get("review_required"):
                    result = "manual_review"
                lines.append(
                    f"- `{criterion['criterion_id']}` `{criterion['check_type']}`: {result}"
                )
                if criterion["note"]:
                    lines.append(f"  - Note: {criterion['note']}")
        lines.append("")
    return "\n".join(lines)


def verify_deck_revision_output(
    case_dir: Path,
    corrected_deck_path: Path,
    *,
    voice_session: Path | None = None,
    plan_path: Path | None = None,
    report_path: Path | None = None,
    review_path: Path | None = None,
    now: datetime | None = None,
) -> DeckRevisionVerificationResult:
    """Verify supported patches in a corrected PPTX."""

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    try:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
    except ImportError as error:
        raise CaseWorkspaceError(
            "python-pptx is required to verify PPTX edits"
        ) from error

    case_dir = case_dir.resolve()
    session_dir = _resolve_voice_session_dir(case_dir, voice_session)
    corrected_deck = _resolve_case_path(case_dir, corrected_deck_path)
    resolved_plan_path, plan = _load_plan(session_dir, plan_path, case_dir)
    try:
        presentation = Presentation(str(corrected_deck))
    except PackageNotFoundError as error:
        raise CaseWorkspaceError(
            f"could not read corrected PPTX: {corrected_deck}"
        ) from error

    change_records: list[dict[str, Any]] = []
    passed_count = 0
    failed_count = 0
    manual_count = 0
    total_count = 0
    passed_criteria_count = 0
    failed_criteria_count = 0
    manual_criteria_count = 0
    total_criteria_count = 0
    for change in plan.get("changes", []):
        if not isinstance(change, dict):
            continue
        slide_number = int(change.get("slide_number", 0) or 0)
        patches = change.get("application_patches")
        if not isinstance(patches, list):
            patches = []
        success_criteria = change.get("success_criteria")
        if not isinstance(success_criteria, list):
            success_criteria = []
        slide = (
            presentation.slides[slide_number - 1]
            if 1 <= slide_number <= len(presentation.slides)
            else None
        )
        patch_records: list[dict[str, Any]] = []
        for patch in patches:
            total_count += 1
            if slide is None:
                record = {
                    "patch_id": (
                        patch.get("patch_id") if isinstance(patch, dict) else None
                    ),
                    "operation": (
                        patch.get("operation") if isinstance(patch, dict) else None
                    ),
                    "passed": False,
                    "expected": "",
                    "actual": "",
                    "note": "slide is missing in corrected deck",
                }
            elif isinstance(patch, dict):
                record = _verify_patch(slide, patch)
            else:
                record = {
                    "patch_id": None,
                    "operation": None,
                    "passed": False,
                    "expected": "",
                    "actual": "",
                    "note": "patch is not an object",
                }
            patch_records.append(record)
            if record["passed"]:
                passed_count += 1
            elif record["note"]:
                manual_count += 1
            else:
                failed_count += 1
        criteria_records: list[dict[str, Any]] = []
        for criterion in success_criteria:
            total_criteria_count += 1
            if isinstance(criterion, dict):
                record = _verify_success_criterion(presentation, slide, criterion)
            else:
                record = {
                    "criterion_id": None,
                    "check_type": None,
                    "passed": False,
                    "expected": "",
                    "actual": "",
                    "note": "success criterion is not an object",
                    "review_required": False,
                }
            criteria_records.append(record)
            if record["passed"]:
                passed_criteria_count += 1
            elif record.get("review_required"):
                manual_criteria_count += 1
            else:
                failed_criteria_count += 1
        has_failed_patch = any(
            not record["passed"] and not record["note"] for record in patch_records
        )
        has_manual_patch = any(
            not record["passed"] and bool(record["note"]) for record in patch_records
        )
        has_failed_criterion = any(
            not record["passed"] and not record.get("review_required")
            for record in criteria_records
        )
        has_manual_criterion = any(
            not record["passed"] and bool(record.get("review_required"))
            for record in criteria_records
        )
        status = (
            "verified"
            if (patch_records or criteria_records)
            and not has_failed_patch
            and not has_manual_patch
            and not has_failed_criterion
            and not has_manual_criterion
            else (
                "manual_review_required"
                if has_manual_patch or has_manual_criterion
                else (
                    "failed"
                    if has_failed_patch or has_failed_criterion
                    else "no_verification_checks"
                )
            )
        )
        change_records.append(
            {
                "change_id": change.get("change_id"),
                "slide_number": slide_number,
                "execution_strategy": change.get("execution_strategy"),
                "status": status,
                "patches": patch_records,
                "success_criteria": criteria_records,
            }
        )

    status = (
        "verified"
        if (total_count > 0 or total_criteria_count > 0)
        and failed_count == 0
        and manual_count == 0
        and failed_criteria_count == 0
        and manual_criteria_count == 0
        else (
            "no_verification_checks"
            if total_count == 0 and total_criteria_count == 0
            else (
                "manual_review_required"
                if failed_count == 0 and failed_criteria_count == 0
                else "failed"
            )
        )
    )
    payload: dict[str, Any] = {
        "schema_version": 1,
        "source": "clara_deck_revision_verification",
        "created_at": _now_iso(now),
        "voice_session": _relative_path(case_dir, session_dir),
        "plan_path": _relative_path(case_dir, resolved_plan_path),
        "corrected_deck_path": _relative_path(case_dir, corrected_deck),
        "summary": {
            "status": status,
            "passed_patches": passed_count,
            "failed_patches": failed_count,
            "manual_review_patches": manual_count,
            "total_patches": total_count,
            "passed_success_criteria": passed_criteria_count,
            "failed_success_criteria": failed_criteria_count,
            "manual_review_success_criteria": manual_criteria_count,
            "total_success_criteria": total_criteria_count,
        },
        "changes": change_records,
    }

    resolved_report_path = report_path or (
        session_dir / "deck_revision_verification.json"
    )
    resolved_review_path = review_path or (
        session_dir / "deck_revision_verification.md"
    )
    if not resolved_report_path.is_absolute():
        resolved_report_path = case_dir / resolved_report_path
    if not resolved_review_path.is_absolute():
        resolved_review_path = case_dir / resolved_review_path
    _write_json(resolved_report_path, payload)
    resolved_review_path.write_text(_render_markdown(payload), encoding="utf-8")
    return DeckRevisionVerificationResult(
        session_dir=session_dir,
        report_path=resolved_report_path,
        review_path=resolved_review_path,
    )


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Verify a corrected PPTX against a Clara deck-revision plan.",
    )
    parser.add_argument("case_dir", type=Path)
    parser.add_argument("corrected_deck", type=Path)
    parser.add_argument(
        "--voice-session",
        type=Path,
        default=None,
        help="Voice session folder name/path. Defaults to latest voice session.",
    )
    parser.add_argument(
        "--plan",
        type=Path,
        default=None,
        help="Normalized change plan. Defaults to deck_revision_changes.normalized.json.",
    )
    args = parser.parse_args(argv)

    logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")
    result = verify_deck_revision_output(
        args.case_dir,
        args.corrected_deck,
        voice_session=args.voice_session,
        plan_path=args.plan,
    )
    LOGGER.info("wrote deck revision verification to %s", result.report_path)
    LOGGER.info("wrote deck revision verification review to %s", result.review_path)
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
