"""Apply supported Clara deck-revision patches to a copied PPTX."""

from __future__ import annotations

import argparse
import hashlib
import json
import logging
import shutil
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Mapping

from advisor_case_core import CaseWorkspaceError, validate_case_workspace
from deck_revision_execution_contract import PATCH_EXECUTION_STRATEGY
from deck_revision_text_match import target_text_matches
from verify_deck_revision_output import verify_deck_revision_output

__all__ = [
    "DeckRevisionApplyResult",
    "apply_deck_revision_plan",
    "main",
]

LOGGER = logging.getLogger(__name__)
DEFAULT_TEXTBOX = {
    "left": 914400,
    "top": 1371600,
    "width": 7315200,
    "height": 914400,
}


@dataclass(frozen=True)
class DeckRevisionApplyResult:
    """Artifacts created when supported deck-revision patches are applied."""

    session_dir: Path
    corrected_deck_path: Path
    apply_report_path: Path
    apply_review_path: Path
    verification_report_path: Path
    verification_review_path: Path
    output_review_path: Path
    output_review_markdown_path: Path


def _now_iso(now: datetime | None = None) -> str:
    value = now or datetime.now(timezone.utc)
    return value.astimezone(timezone.utc).replace(microsecond=0).isoformat()


def _read_json(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise CaseWorkspaceError(f"expected JSON object in {path}")
    return payload


def _write_json(path: Path, payload: Mapping[str, Any]) -> None:
    path.write_text(
        json.dumps(payload, indent=2, ensure_ascii=True) + "\n",
        encoding="utf-8",
    )


def _deck_titles(presentation: Any) -> list[dict[str, Any]]:
    titles: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is None:
            title_text = ""
        else:
            title_text = _shape_text(title_shape)
        if not title_text:
            for shape in slide.shapes:
                text = _shape_text(shape).strip()
                if text:
                    title_text = text
                    break
        titles.append({"slide_number": index, "title": title_text})
    return titles


def _render_output_review_markdown(payload: Mapping[str, Any]) -> str:
    lines = [
        "# Deck Revision Final Output Review",
        "",
        f"Status: `{payload['summary']['status']}`",
        "",
        f"- Corrected deck: `{payload['corrected_deck_path']}`",
        f"- Mechanical verification: `{payload['verification_status']}`",
        "",
        "## Required Clara/Codex Loop",
        "",
        "Before this corrected PPTX is final, inspect the rendered deck and complete this semantic review:",
        "",
    ]
    for check in payload["checks"]:
        lines.append(f"- `{check['check_id']}`: {check['question']}")
    lines.extend(["", "## Slide Titles", ""])
    for title in payload["slide_titles"]:
        lines.append(
            f"- Slide {title['slide_number']}: {title['title'] or '[missing title]'}"
        )
    lines.extend(
        [
            "",
            "## Loop Rule",
            "",
            "If any title, subtitle, or visible copy reads like an internal instruction, workpaper note, prompt, implementation rule, or construction rationale, revise the deck and rerun render/verification before final approval.",
            "",
        ]
    )
    return "\n".join(lines)


def _write_output_review(
    *,
    case_dir: Path,
    session_dir: Path,
    corrected_deck_path: Path,
    plan_path: Path,
    verification_status: str,
    presentation: Any,
    now: datetime | None = None,
) -> tuple[Path, Path]:
    payload: dict[str, Any] = {
        "schema_version": 1,
        "source": "clara_deck_revision_output_review",
        "created_at": _now_iso(now),
        "voice_session": _relative_path(case_dir, session_dir),
        "plan_path": _relative_path(case_dir, plan_path),
        "corrected_deck_path": _relative_path(case_dir, corrected_deck_path),
        "verification_status": verification_status,
        "summary": {
            "status": "requires_clara_codex_review",
            "complete": False,
            "reason": "semantic and audience-facing deck review has not been completed",
        },
        "checks": [
            {
                "check_id": "audience_copy",
                "question": "Do all visible titles and subtitles read as audience-facing slide copy rather than internal instructions?",
                "status": "pending",
            },
            {
                "check_id": "process_language",
                "question": "Has all workpaper, prompt, construction, approval, and implementation language been kept out of visible slides?",
                "status": "pending",
            },
            {
                "check_id": "requested_structure",
                "question": "Does the final deck structure match the approved interpretation, including inserted/deleted/reordered slides?",
                "status": "pending",
            },
            {
                "check_id": "semantic_evidence_fit",
                "question": "Do the edits preserve the intended meaning of the transcript evidence and avoid unsupported claims?",
                "status": "pending",
            },
            {
                "check_id": "visual_render",
                "question": "Has the rendered deck been inspected for clipping, overlap, stale artifacts, wrong footers, and broken layout?",
                "status": "pending",
            },
        ],
        "slide_titles": _deck_titles(presentation),
    }
    review_path = session_dir / "deck_revision_output_review.json"
    markdown_path = session_dir / "deck_revision_output_review.md"
    _write_json(review_path, payload)
    markdown_path.write_text(_render_output_review_markdown(payload), encoding="utf-8")
    return review_path, markdown_path


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _relative_path(case_dir: Path, path: Path) -> str:
    try:
        return str(path.resolve().relative_to(case_dir.resolve()))
    except ValueError:
        return str(path.resolve())


def _resolve_voice_session_dir(case_dir: Path, voice_session: Path | None) -> Path:
    sessions_root = case_dir / "voice_sessions"
    if voice_session is None:
        if not sessions_root.is_dir():
            raise CaseWorkspaceError("case has no voice_sessions folder")
        sessions = sorted(path for path in sessions_root.iterdir() if path.is_dir())
        if not sessions:
            raise CaseWorkspaceError("case has no imported voice sessions")
        return sessions[-1].resolve()

    candidate = voice_session.expanduser()
    candidates = (
        [candidate]
        if candidate.is_absolute()
        else [case_dir / candidate, sessions_root / candidate]
    )
    for path in candidates:
        if path.is_dir():
            resolved = path.resolve()
            try:
                resolved.relative_to(sessions_root.resolve())
            except ValueError as error:
                raise CaseWorkspaceError(
                    f"voice session must live under {sessions_root}: {resolved}"
                ) from error
            return resolved
    raise CaseWorkspaceError(f"voice session does not exist: {voice_session}")


def _resolve_case_path(case_dir: Path, raw_path: str | Path, *, label: str) -> Path:
    candidate = Path(raw_path).expanduser()
    if not candidate.is_absolute():
        candidate = case_dir / candidate
    if not candidate.is_file():
        raise CaseWorkspaceError(f"{label} does not exist: {candidate}")
    return candidate.resolve()


def _load_plan(
    session_dir: Path, plan_path: Path | None, case_dir: Path
) -> tuple[Path, dict[str, Any]]:
    candidate = plan_path or (session_dir / "deck_revision_changes.normalized.json")
    if not candidate.is_absolute():
        candidate = case_dir / candidate
    if not candidate.is_file():
        raise CaseWorkspaceError(
            f"normalized deck revision plan is missing: {candidate}; run finalize_deck_revision_plan.py first"
        )
    return candidate.resolve(), _read_json(candidate)


def _load_approval(
    case_dir: Path,
    session_dir: Path,
    plan_path: Path,
    approval_path: Path | None,
) -> tuple[Path, dict[str, Any], Path]:
    candidate = approval_path or (session_dir / "deck_revision_approval.json")
    if not candidate.is_absolute():
        candidate = case_dir / candidate
    if not candidate.is_file():
        raise CaseWorkspaceError(
            f"deck revision approval is missing: {candidate}; run approve_deck_revision_plan.py after reviewing deck_revision_changes.md"
        )
    approval = _read_json(candidate)
    if not bool(approval.get("approved")):
        raise CaseWorkspaceError("deck revision approval is not marked approved")
    expected_hash = _sha256(plan_path)
    if approval.get("plan_sha256") != expected_hash:
        raise CaseWorkspaceError(
            "deck revision approval does not match the current normalized plan; re-review and re-approve the plan"
        )
    if not bool(approval.get("understanding_reviewed")):
        raise CaseWorkspaceError(
            "deck revision approval does not certify the understanding checkpoint; "
            "show deck_revision_understanding.md to the user and re-approve with "
            "--understanding-reviewed"
        )
    raw_understanding_path = approval.get("understanding_path")
    if (
        not isinstance(raw_understanding_path, str)
        or not raw_understanding_path.strip()
    ):
        raise CaseWorkspaceError(
            "deck revision approval is missing understanding_path; re-review "
            "deck_revision_understanding.md and re-approve the plan"
        )
    understanding_path = Path(raw_understanding_path).expanduser()
    if not understanding_path.is_absolute():
        understanding_path = case_dir / understanding_path
    if not understanding_path.is_file():
        raise CaseWorkspaceError(
            f"deck revision understanding checkpoint is missing: {understanding_path}; "
            "re-run finalization and approval before applying"
        )
    expected_understanding_hash = _sha256(understanding_path)
    if approval.get("understanding_sha256") != expected_understanding_hash:
        raise CaseWorkspaceError(
            "deck revision understanding checkpoint changed after approval; "
            "show the updated deck_revision_understanding.md and re-approve the plan"
        )
    return candidate.resolve(), approval, understanding_path.resolve()


def _load_workbench(session_dir: Path) -> dict[str, Any]:
    workbench_path = session_dir / "deck_revision_workbench.json"
    if not workbench_path.is_file():
        raise CaseWorkspaceError(
            f"deck revision workbench is missing: {workbench_path}; run build_deck_revision_workbench.py first"
        )
    return _read_json(workbench_path)


def _shape_by_index(slide: Any, shape_index: int) -> Any | None:
    shapes = list(slide.shapes)
    if shape_index < 1 or shape_index > len(shapes):
        return None
    return shapes[shape_index - 1]


def _shape_text(shape: Any) -> str:
    return str(getattr(shape, "text", "") or "")


def _target_text_matches(shape: Any, expected_text: str) -> bool:
    return target_text_matches(_shape_text(shape), expected_text)


def _target_mismatch_result(
    patch: Mapping[str, Any],
    actual_text: str,
    expected_text: str,
) -> dict[str, Any]:
    return {
        "patch_id": patch.get("patch_id"),
        "operation": patch.get("operation"),
        "status": "blocked",
        "message": (
            "target text mismatch; expected "
            f"{expected_text!r}, found {actual_text!r}"
        ),
    }


def _replace_text_in_shape(shape: Any, old_text: str, new_text: str) -> bool:
    if not hasattr(shape, "text_frame"):
        return False
    current = _shape_text(shape)
    if old_text not in current:
        return False
    shape.text = current.replace(old_text, new_text)
    return True


def _delete_shape(shape: Any) -> None:
    element = shape._element  # noqa: SLF001 - python-pptx exposes no public delete API.
    parent = element.getparent()
    parent.remove(element)


def _apply_patch(slide: Any, patch: Mapping[str, Any]) -> dict[str, Any]:
    patch_id = patch.get("patch_id")
    operation = patch.get("operation")
    target = patch.get("target") if isinstance(patch.get("target"), dict) else {}
    value = patch.get("value") if isinstance(patch.get("value"), dict) else {}
    try:
        if operation == "set_title_text":
            title_shape = getattr(slide.shapes, "title", None)
            if title_shape is None:
                return {
                    "patch_id": patch_id,
                    "operation": operation,
                    "status": "blocked",
                    "message": "slide has no title shape",
                }
            expected_text = str(target.get("expected_text", ""))
            if expected_text and not _target_text_matches(title_shape, expected_text):
                return _target_mismatch_result(
                    patch,
                    _shape_text(title_shape),
                    expected_text,
                )
            title_shape.text = str(value["text"])
        elif operation == "set_shape_text":
            shape = _shape_by_index(slide, int(target["shape_index"]))
            if shape is None:
                return {
                    "patch_id": patch_id,
                    "operation": operation,
                    "status": "blocked",
                    "message": "target shape is missing",
                }
            expected_text = str(target.get("expected_text", ""))
            if expected_text and not _target_text_matches(shape, expected_text):
                return _target_mismatch_result(
                    patch,
                    _shape_text(shape),
                    expected_text,
                )
            shape.text = str(value["text"])
        elif operation == "replace_text":
            old_text = str(value["old_text"])
            new_text = str(value["new_text"])
            target_shape_index = target.get("shape_index")
            replaced = False
            if isinstance(target_shape_index, int):
                shape = _shape_by_index(slide, target_shape_index)
                if shape is not None:
                    expected_text = str(target.get("expected_text", ""))
                    if expected_text and not _target_text_matches(shape, expected_text):
                        return _target_mismatch_result(
                            patch,
                            _shape_text(shape),
                            expected_text,
                        )
                    replaced = _replace_text_in_shape(shape, old_text, new_text)
            else:
                return {
                    "patch_id": patch_id,
                    "operation": operation,
                    "status": "blocked",
                    "message": "replace_text requires target.shape_index",
                }
            if not replaced:
                return {
                    "patch_id": patch_id,
                    "operation": operation,
                    "status": "blocked",
                    "message": "old_text was not found in the target slide",
                }
        elif operation == "add_textbox":
            left = int(value.get("left", DEFAULT_TEXTBOX["left"]))
            top = int(value.get("top", DEFAULT_TEXTBOX["top"]))
            width = int(value.get("width", DEFAULT_TEXTBOX["width"]))
            height = int(value.get("height", DEFAULT_TEXTBOX["height"]))
            textbox = slide.shapes.add_textbox(left, top, width, height)
            textbox.text = str(value["text"])
        elif operation == "delete_shape":
            shape = _shape_by_index(slide, int(target["shape_index"]))
            if shape is None:
                return {
                    "patch_id": patch_id,
                    "operation": operation,
                    "status": "blocked",
                    "message": "target shape is missing",
                }
            expected_text = str(target.get("expected_text", ""))
            if expected_text and not _target_text_matches(shape, expected_text):
                return _target_mismatch_result(
                    patch,
                    _shape_text(shape),
                    expected_text,
                )
            _delete_shape(shape)
        elif operation == "move_shape":
            shape = _shape_by_index(slide, int(target["shape_index"]))
            if shape is None:
                return {
                    "patch_id": patch_id,
                    "operation": operation,
                    "status": "blocked",
                    "message": "target shape is missing",
                }
            expected_text = str(target.get("expected_text", ""))
            if expected_text and not _target_text_matches(shape, expected_text):
                return _target_mismatch_result(
                    patch,
                    _shape_text(shape),
                    expected_text,
                )
            if "left" in value:
                shape.left = int(value["left"])
            if "top" in value:
                shape.top = int(value["top"])
        else:
            return {
                "patch_id": patch_id,
                "operation": operation,
                "status": "blocked",
                "message": f"unsupported operation `{operation}`",
            }
    except (KeyError, TypeError, ValueError) as error:
        return {
            "patch_id": patch_id,
            "operation": operation,
            "status": "blocked",
            "message": f"invalid patch payload: {error}",
        }
    return {
        "patch_id": patch_id,
        "operation": operation,
        "status": "applied",
        "message": "applied",
    }


def _render_markdown(payload: Mapping[str, Any]) -> str:
    lines = [
        "# Deck Revision Apply Report",
        "",
        f"Status: `{payload['summary']['status']}`",
        "",
        f"- Source deck: `{payload['source_deck_path']}`",
        f"- Corrected deck: `{payload['corrected_deck_path']}`",
        f"- Approval: `{payload.get('approval_path') or 'not used'}`",
        f"- Approved by: {payload.get('approved_by') or 'not recorded'}",
        f"- Reviewed understanding: `{payload.get('understanding_path') or 'not used'}`",
        f"- Applied patches: {payload['summary']['applied_patches']}",
        f"- Blocked patches: {payload['summary']['blocked_patches']}",
        f"- Verification status: `{payload['summary']['verification_status']}`",
        f"- Verification: `{payload['verification_report_path']}`",
        f"- Final output review: `{payload['output_review_path']}`",
        "",
    ]
    for change in payload["changes"]:
        lines.extend(
            [
                f"## Slide {change['slide_number']} - {change['change_id']}",
                "",
                f"Status: `{change['status']}`",
                f"Strategy: `{change.get('execution_strategy') or 'not recorded'}`",
                "",
            ]
        )
        if change["patches"]:
            for patch in change["patches"]:
                lines.append(
                    f"- `{patch['patch_id']}` `{patch['operation']}`: {patch['status']} - {patch['message']}"
                )
        else:
            lines.append("- No executable patches; not applied automatically.")
        lines.append("")
    return "\n".join(lines)


def apply_deck_revision_plan(
    case_dir: Path,
    *,
    voice_session: Path | None = None,
    plan_path: Path | None = None,
    approval_path: Path | None = None,
    output_path: Path | None = None,
    allow_unapproved: bool = False,
    now: datetime | None = None,
) -> DeckRevisionApplyResult:
    """Apply supported patches from a reviewed deck-revision plan."""

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))
    try:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
    except ImportError as error:
        raise CaseWorkspaceError(
            "python-pptx is required to apply PPTX edits"
        ) from error

    case_dir = case_dir.resolve()
    session_dir = _resolve_voice_session_dir(case_dir, voice_session)
    workbench = _load_workbench(session_dir)
    resolved_plan_path, plan = _load_plan(session_dir, plan_path, case_dir)
    resolved_approval_path: Path | None = None
    reviewed_understanding_path: Path | None = None
    approval: dict[str, Any] | None = None
    if not allow_unapproved:
        (
            resolved_approval_path,
            approval,
            reviewed_understanding_path,
        ) = _load_approval(
            case_dir,
            session_dir,
            resolved_plan_path,
            approval_path,
        )
    source_deck_path = _resolve_case_path(
        case_dir,
        workbench.get("source_paths", {}).get("deck_path", ""),
        label="source PPTX deck",
    )
    corrected_deck_path = output_path or (session_dir / "deck_revision_corrected.pptx")
    if not corrected_deck_path.is_absolute():
        corrected_deck_path = case_dir / corrected_deck_path
    corrected_deck_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source_deck_path, corrected_deck_path)
    try:
        presentation = Presentation(str(corrected_deck_path))
    except PackageNotFoundError as error:
        raise CaseWorkspaceError(
            f"could not read PPTX deck: {corrected_deck_path}"
        ) from error

    change_records: list[dict[str, Any]] = []
    applied_count = 0
    blocked_count = 0
    model_or_human_count = 0
    for change in plan.get("changes", []):
        if not isinstance(change, dict):
            continue
        slide_number = int(change.get("slide_number", 0) or 0)
        execution_strategy = str(change.get("execution_strategy", "")).strip()
        patches = change.get("application_patches")
        if not isinstance(patches, list):
            patches = []
        if not patches and execution_strategy != PATCH_EXECUTION_STRATEGY:
            model_or_human_count += 1
        if not 1 <= slide_number <= len(presentation.slides):
            patch_records = [
                {
                    "patch_id": (
                        patch.get("patch_id") if isinstance(patch, dict) else None
                    ),
                    "operation": (
                        patch.get("operation") if isinstance(patch, dict) else None
                    ),
                    "status": "blocked",
                    "message": "slide is missing in source deck",
                }
                for patch in patches
            ]
        else:
            slide = presentation.slides[slide_number - 1]
            patch_records = [
                _apply_patch(slide, patch)
                for patch in patches
                if isinstance(patch, dict)
            ]
        applied_count += sum(
            1 for patch in patch_records if patch["status"] == "applied"
        )
        blocked_count += sum(
            1 for patch in patch_records if patch["status"] != "applied"
        )
        status = (
            "applied"
            if patch_records
            and all(patch["status"] == "applied" for patch in patch_records)
            else (
                "blocked"
                if patch_records
                else (
                    f"requires_{execution_strategy}"
                    if execution_strategy
                    and execution_strategy != PATCH_EXECUTION_STRATEGY
                    else "no_executable_patches"
                )
            )
        )
        change_records.append(
            {
                "change_id": change.get("change_id"),
                "slide_number": slide_number,
                "execution_strategy": execution_strategy,
                "status": status,
                "patches": patch_records,
            }
        )

    presentation.save(corrected_deck_path)
    verification = verify_deck_revision_output(
        case_dir,
        corrected_deck_path,
        voice_session=session_dir,
        plan_path=resolved_plan_path,
        report_path=session_dir / "deck_revision_verification.json",
        review_path=session_dir / "deck_revision_verification.md",
        now=now,
    )
    verification_payload = _read_json(verification.report_path)
    verification_status = (
        verification_payload.get("summary", {}).get("status")
        if isinstance(verification_payload.get("summary"), dict)
        else "unknown"
    )
    output_review_path, output_review_markdown_path = _write_output_review(
        case_dir=case_dir,
        session_dir=session_dir,
        corrected_deck_path=corrected_deck_path,
        plan_path=resolved_plan_path,
        verification_status=str(verification_status),
        presentation=presentation,
        now=now,
    )
    status = (
        "applied_verified_pending_output_review"
        if applied_count > 0
        and blocked_count == 0
        and verification_status == "verified"
        else (
            "applied_but_verification_failed"
            if applied_count > 0 and blocked_count == 0
            else (
                "model_or_human_execution_required"
                if model_or_human_count > 0
                and applied_count == 0
                and blocked_count == 0
                else (
                    "no_executable_patches"
                    if applied_count == 0 and blocked_count == 0
                    else "partial_or_blocked"
                )
            )
        )
    )
    payload: dict[str, Any] = {
        "schema_version": 1,
        "source": "clara_deck_revision_apply_report",
        "created_at": _now_iso(now),
        "voice_session": _relative_path(case_dir, session_dir),
        "plan_path": _relative_path(case_dir, resolved_plan_path),
        "approval_path": (
            _relative_path(case_dir, resolved_approval_path)
            if resolved_approval_path is not None
            else None
        ),
        "approved_by": approval.get("approved_by") if approval is not None else None,
        "understanding_path": (
            _relative_path(case_dir, reviewed_understanding_path)
            if reviewed_understanding_path is not None
            else None
        ),
        "source_deck_path": _relative_path(case_dir, source_deck_path),
        "corrected_deck_path": _relative_path(case_dir, corrected_deck_path),
        "verification_report_path": _relative_path(case_dir, verification.report_path),
        "verification_review_path": _relative_path(case_dir, verification.review_path),
        "output_review_path": _relative_path(case_dir, output_review_path),
        "output_review_markdown_path": _relative_path(
            case_dir, output_review_markdown_path
        ),
        "summary": {
            "status": status,
            "applied_patches": applied_count,
            "blocked_patches": blocked_count,
            "model_or_human_changes": model_or_human_count,
            "changes": len(change_records),
            "verification_status": verification_status,
            "final_output_review_status": "requires_clara_codex_review",
        },
        "changes": change_records,
    }
    apply_report_path = session_dir / "deck_revision_apply_report.json"
    apply_review_path = session_dir / "deck_revision_apply_report.md"
    _write_json(apply_report_path, payload)
    apply_review_path.write_text(_render_markdown(payload), encoding="utf-8")
    return DeckRevisionApplyResult(
        session_dir=session_dir,
        corrected_deck_path=corrected_deck_path,
        apply_report_path=apply_report_path,
        apply_review_path=apply_review_path,
        verification_report_path=verification.report_path,
        verification_review_path=verification.review_path,
        output_review_path=output_review_path,
        output_review_markdown_path=output_review_markdown_path,
    )


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Apply supported Clara deck-revision patches to a copied PPTX.",
    )
    parser.add_argument("case_dir", type=Path)
    parser.add_argument(
        "--voice-session",
        type=Path,
        default=None,
        help="Voice session folder name/path. Defaults to latest voice session.",
    )
    parser.add_argument(
        "--plan",
        type=Path,
        default=None,
        help="Normalized change plan. Defaults to deck_revision_changes.normalized.json.",
    )
    parser.add_argument(
        "--approval",
        type=Path,
        default=None,
        help="Approval JSON path. Defaults to deck_revision_approval.json.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=None,
        help="Corrected PPTX path. Defaults to deck_revision_corrected.pptx in the voice session.",
    )
    parser.add_argument(
        "--allow-unapproved",
        action="store_true",
        help="Apply for an internal dry run even if approved_for_pptx_revision is false.",
    )
    args = parser.parse_args(argv)

    logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")
    result = apply_deck_revision_plan(
        args.case_dir,
        voice_session=args.voice_session,
        plan_path=args.plan,
        approval_path=args.approval,
        output_path=args.output,
        allow_unapproved=args.allow_unapproved,
    )
    LOGGER.info("wrote corrected deck to %s", result.corrected_deck_path)
    LOGGER.info("wrote apply report to %s", result.apply_report_path)
    LOGGER.info("wrote verification report to %s", result.verification_report_path)
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
