"""Prepare local deck-revision intake artifacts from a Clara voice session."""

from __future__ import annotations

import argparse
import json
import logging
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Mapping

from advisor_case_core import (
    CaseWorkspaceError,
    prepare_editable_pptx_merge_input,
    validate_case_workspace,
)
from match_feedback_frames_to_deck_slides import (
    SlideFrameMatchError,
    match_feedback_timeline_to_deck,
)

__all__ = [
    "VoiceDeckRevisionIntakeResult",
    "prepare_voice_deck_revision_intake",
    "main",
]

LOGGER = logging.getLogger(__name__)

VIDEO_SUFFIXES = {".webm", ".mp4", ".mov", ".mkv"}
STYLE_SPEC_FILENAMES = {
    "ag": "ag-style-spec.md",
    "bain": "bain-style-spec.md",
}
COMPANY_PROFILE_FILENAMES = (
    "company_profile.json",
    "clara_company_profile.json",
)
ADVISORY_OUTPUT_SHAPER_PATH = (
    Path(".agents") / "skills" / "advisory-output-shaper" / "SKILL.md"
)


@dataclass(frozen=True)
class VoiceDeckRevisionIntakeResult:
    """Artifacts created for local deck-correction review."""

    session_dir: Path
    intake_path: Path
    gate_path: Path
    deck_snapshot_path: Path | None
    merge_input_report_path: Path | None
    style_spec_snapshot_path: Path | None


def _now_iso(now: datetime | None = None) -> str:
    value = now or datetime.now(timezone.utc)
    return value.astimezone(timezone.utc).replace(microsecond=0).isoformat()


def _read_json(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise CaseWorkspaceError(f"expected JSON object in {path}")
    return payload


def _write_json(path: Path, payload: Mapping[str, Any]) -> None:
    path.write_text(
        json.dumps(payload, indent=2, ensure_ascii=True) + "\n",
        encoding="utf-8",
    )


def _relative_path(case_dir: Path, path: Path) -> str:
    try:
        return str(path.resolve().relative_to(case_dir.resolve()))
    except ValueError:
        return str(path.resolve())


def _repo_root() -> Path:
    return Path(__file__).resolve().parents[3]


def _style_spec_root() -> Path:
    return _repo_root() / "docs" / "specs" / "pptx_templates"


def _normalize_style_key(value: str) -> str:
    compact = "".join(character for character in value.lower() if character.isalnum())
    if compact in {"ag", "ang"}:
        return "ag"
    return compact


def _style_key_from_path(path: Path) -> str:
    stem = path.stem.replace("-style-spec", "")
    return _normalize_style_key(stem)


def _first_string(payload: Mapping[str, Any], keys: tuple[str, ...]) -> str:
    for key in keys:
        value = payload.get(key)
        if isinstance(value, str) and value.strip():
            return value.strip()
    return ""


def _resolve_existing_file(
    raw_path: Path,
    *,
    case_dir: Path,
    base_dir: Path | None = None,
    description: str,
) -> Path:
    expanded = raw_path.expanduser()
    candidates: list[Path] = []
    if expanded.is_absolute():
        candidates.append(expanded)
    else:
        if base_dir is not None:
            candidates.append(base_dir / expanded)
        candidates.extend((case_dir / expanded, _repo_root() / expanded))
    for candidate in candidates:
        if candidate.is_file():
            return candidate.resolve()
    raise CaseWorkspaceError(f"{description} does not exist: {raw_path}")


def _company_profile_candidates(case_dir: Path) -> list[Path]:
    candidates: list[Path] = []
    for base_dir in (case_dir, case_dir.parent):
        for filename in COMPANY_PROFILE_FILENAMES:
            candidates.append(base_dir / filename)
        candidates.append(base_dir / ".clara" / "company_profile.json")
    return candidates


def _load_company_profile(
    case_dir: Path,
    company_profile_path: Path | None,
) -> tuple[Path | None, dict[str, Any]]:
    if company_profile_path is not None:
        resolved = _resolve_existing_file(
            company_profile_path,
            case_dir=case_dir,
            description="company profile",
        )
        return resolved, _read_json(resolved)

    for candidate in _company_profile_candidates(case_dir):
        if candidate.is_file():
            return candidate.resolve(), _read_json(candidate)
    return None, {}


def _company_profile_payload(
    *,
    case_dir: Path,
    profile_path: Path | None,
    profile: Mapping[str, Any],
) -> dict[str, Any]:
    if profile_path is None:
        return {
            "status": "missing",
            "path": None,
            "search_rule": "case folder, then parent company folder",
        }
    return {
        "status": "available",
        "path": _relative_path(case_dir, profile_path),
        "company": _first_string(
            profile,
            ("company", "company_name", "firm", "firm_name", "advisor_firm"),
        )
        or None,
        "default_deck_style": _first_string(
            profile,
            (
                "default_deck_style",
                "deck_style",
                "deck_style_key",
                "pptx_style",
                "pptx_style_key",
            ),
        )
        or None,
        "advisory_method": _first_string(
            profile,
            ("advisory_method", "advisory_skill", "default_advisory_skill"),
        )
        or None,
    }


def _style_spec_for_key(style_key: str) -> tuple[str, Path]:
    normalized = _normalize_style_key(style_key)
    filename = STYLE_SPEC_FILENAMES.get(normalized)
    if filename is None:
        allowed = ", ".join(sorted(STYLE_SPEC_FILENAMES))
        raise CaseWorkspaceError(
            f"unknown deck style '{style_key}'; expected one of: {allowed}"
        )
    path = _style_spec_root() / filename
    if not path.is_file():
        raise CaseWorkspaceError(f"deck style spec is missing: {path}")
    return normalized, path


def _style_title(spec_text: str) -> str | None:
    for line in spec_text.splitlines():
        stripped = line.strip()
        if stripped.startswith("#"):
            return stripped.lstrip("#").strip() or None
    return None


def _resolve_deck_style(
    *,
    case_dir: Path,
    session_dir: Path,
    manifest: Mapping[str, Any],
    company_profile_path: Path | None,
    company_profile: Mapping[str, Any],
    deck_style: str,
    style_spec_path: Path | None,
) -> tuple[dict[str, Any], Path | None]:
    if deck_style and style_spec_path is not None:
        raise CaseWorkspaceError("pass either --deck-style or --style-spec, not both")

    source = ""
    style_key = ""
    spec_path: Path | None = None
    if style_spec_path is not None:
        spec_path = _resolve_existing_file(
            style_spec_path,
            case_dir=case_dir,
            description="deck style spec",
        )
        style_key = _style_key_from_path(spec_path)
        source = "explicit_style_spec"
    elif deck_style:
        style_key, spec_path = _style_spec_for_key(deck_style)
        source = "explicit_deck_style"
    else:
        profile_style_spec = _first_string(
            company_profile,
            ("deck_style_spec_path", "pptx_style_spec_path", "style_spec_path"),
        )
        profile_style_key = _first_string(
            company_profile,
            (
                "default_deck_style",
                "deck_style",
                "deck_style_key",
                "pptx_style",
                "pptx_style_key",
            ),
        )
        manifest_style_spec = _first_string(
            manifest,
            ("deck_style_spec_path", "pptx_style_spec_path", "style_spec_path"),
        )
        manifest_style_key = _first_string(
            manifest,
            (
                "deck_style",
                "deck_style_key",
                "pptx_style",
                "pptx_style_key",
            ),
        )
        if profile_style_spec:
            spec_path = _resolve_existing_file(
                Path(profile_style_spec),
                case_dir=case_dir,
                base_dir=(
                    company_profile_path.parent
                    if company_profile_path is not None
                    else None
                ),
                description="company deck style spec",
            )
            style_key = profile_style_key or _style_key_from_path(spec_path)
            source = "company_profile_style_spec"
        elif profile_style_key:
            style_key, spec_path = _style_spec_for_key(profile_style_key)
            source = "company_profile_deck_style"
        elif manifest_style_spec:
            spec_path = _resolve_existing_file(
                Path(manifest_style_spec),
                case_dir=case_dir,
                description="case deck style spec",
            )
            style_key = manifest_style_key or _style_key_from_path(spec_path)
            source = "case_manifest_style_spec"
        elif manifest_style_key:
            style_key, spec_path = _style_spec_for_key(manifest_style_key)
            source = "case_manifest_deck_style"

    if spec_path is None:
        return (
            {
                "status": "missing",
                "source": "not_found",
                "style_key": None,
                "spec_path": None,
                "snapshot_path": None,
                "required_before_editing": True,
                "resolution_order": [
                    "explicit --style-spec",
                    "explicit --deck-style",
                    "company profile in case folder or parent folder",
                    "case_manifest.json deck_style fields",
                ],
            },
            None,
        )

    spec_text = spec_path.read_text(encoding="utf-8")
    snapshot_path = session_dir / "deck_style_spec.md"
    snapshot_path.write_text(spec_text, encoding="utf-8")
    return (
        {
            "status": "available",
            "source": source,
            "style_key": _normalize_style_key(style_key),
            "style_name": _style_title(spec_text),
            "spec_path": _relative_path(case_dir, spec_path),
            "snapshot_path": _relative_path(case_dir, snapshot_path),
            "required_before_editing": True,
        },
        snapshot_path,
    )


def _advisory_method_payload(case_dir: Path) -> dict[str, Any]:
    skill_path = _repo_root() / ADVISORY_OUTPUT_SHAPER_PATH
    return {
        "status": "required",
        "skill": "advisory-output-shaper",
        "skill_path": (
            _relative_path(case_dir, skill_path) if skill_path.is_file() else None
        ),
        "method_scope": (
            "Use for advisory/family/governance wording, tension handling, "
            "meeting usefulness, and anti-BS shaping before deck edits."
        ),
    }


def _resolve_voice_session_dir(case_dir: Path, voice_session: Path | None) -> Path:
    sessions_root = case_dir / "voice_sessions"
    if voice_session is None:
        if not sessions_root.is_dir():
            raise CaseWorkspaceError("case has no voice_sessions folder")
        sessions = sorted(path for path in sessions_root.iterdir() if path.is_dir())
        if not sessions:
            raise CaseWorkspaceError("case has no imported voice sessions")
        return sessions[-1].resolve()

    candidate = voice_session.expanduser()
    candidates = (
        [candidate]
        if candidate.is_absolute()
        else [case_dir / candidate, sessions_root / candidate]
    )
    for path in candidates:
        if path.is_dir():
            resolved = path.resolve()
            try:
                resolved.relative_to(sessions_root.resolve())
            except ValueError as error:
                raise CaseWorkspaceError(
                    f"voice session must live under {sessions_root}: {resolved}"
                ) from error
            return resolved
    raise CaseWorkspaceError(f"voice session does not exist: {voice_session}")


def _find_material(
    case_dir: Path,
    material_id: str,
) -> dict[str, Any]:
    registry = _read_json(case_dir / "material_registry.json")
    for material in registry.get("materials", []):
        if isinstance(material, dict) and material.get("id") == material_id:
            return material
    raise CaseWorkspaceError(f"material not found: {material_id}")


def _resolve_deck_path(
    *,
    case_dir: Path,
    deck_path: Path | None,
    deck_material_id: str,
) -> tuple[Path | None, str]:
    if deck_path is not None and deck_material_id:
        raise CaseWorkspaceError("pass either --deck or --deck-material-id, not both")
    if deck_material_id:
        material = _find_material(case_dir, deck_material_id)
        material_path = Path(str(material.get("path", ""))).expanduser()
        if not material_path.is_file():
            raise CaseWorkspaceError(
                f"deck material path does not exist: {material_path}"
            )
        return material_path.resolve(), deck_material_id
    if deck_path is None:
        return None, ""
    resolved = deck_path.expanduser().resolve()
    if not resolved.is_file():
        raise CaseWorkspaceError(f"deck file does not exist: {resolved}")
    return resolved, ""


def _clean_text(value: str) -> str:
    return "\n".join(line.strip() for line in value.splitlines() if line.strip())


def _shape_text(shape: Any) -> str:
    text = _clean_text(str(getattr(shape, "text", "") or ""))
    if text:
        return text
    if not bool(getattr(shape, "has_table", False)):
        return ""
    table = shape.table
    cell_texts: list[str] = []
    for row in getattr(table, "rows", []):
        for cell in getattr(row, "cells", []):
            cell_text = _clean_text(str(getattr(cell, "text", "") or ""))
            if cell_text:
                cell_texts.append(cell_text)
    return "\n".join(cell_texts)


def _extract_pptx_snapshot(
    *,
    deck_path: Path,
    deck_material_id: str,
    case_dir: Path,
    output_path: Path,
    now: datetime | None,
) -> None:
    """Extract mechanical PPTX text/object context for later model judgement."""

    if deck_path.suffix.lower() != ".pptx":
        raise CaseWorkspaceError(
            f"deck revision intake currently requires a .pptx deck: {deck_path}"
        )
    try:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
    except ImportError as error:
        raise CaseWorkspaceError(
            "python-pptx is required to extract PPTX deck snapshots"
        ) from error
    try:
        presentation = Presentation(str(deck_path))
    except PackageNotFoundError as error:
        raise CaseWorkspaceError(f"could not read PPTX deck: {deck_path}") from error

    slides: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        shape_records: list[dict[str, Any]] = []
        text_values: list[str] = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            text = _shape_text(shape)
            if not text:
                continue
            text_values.append(text)
            shape_records.append(
                {
                    "shape_index": shape_index,
                    "text": text,
                    "left": int(getattr(shape, "left", 0) or 0),
                    "top": int(getattr(shape, "top", 0) or 0),
                    "width": int(getattr(shape, "width", 0) or 0),
                    "height": int(getattr(shape, "height", 0) or 0),
                }
            )
        title_shape = getattr(slide.shapes, "title", None)
        title = _shape_text(title_shape) if title_shape is not None else ""
        if not title and text_values:
            title = text_values[0]
        slides.append(
            {
                "slide_number": slide_number,
                "title": title,
                "texts": text_values,
                "shapes": shape_records,
            }
        )

    payload = {
        "schema_version": 1,
        "source": "clara_voice_deck_snapshot",
        "created_at": _now_iso(now),
        "deck_path": _relative_path(case_dir, deck_path),
        "deck_material_id": deck_material_id or None,
        "slide_count": len(slides),
        "slides": slides,
    }
    _write_json(output_path, payload)


def _detect_video_path(case_dir: Path, session_dir: Path) -> Path | None:
    timeline_path = session_dir / "video_timeline.json"
    if timeline_path.is_file():
        timeline = _read_json(timeline_path)
        raw_path = str(timeline.get("video_path", "")).strip()
        if raw_path:
            candidate = Path(raw_path).expanduser()
            if not candidate.is_absolute():
                candidate = case_dir / candidate
            if candidate.is_file():
                return candidate.resolve()
    for path in sorted(session_dir.iterdir()):
        name = path.name.lower()
        if (
            path.is_file()
            and path.suffix.lower() in VIDEO_SUFFIXES
            and ("screen" in name or "video" in name)
        ):
            return path.resolve()
    return None


def _path_if_file(path: Path) -> Path | None:
    return path.resolve() if path.is_file() else None


def _path_from_metadata(case_dir: Path, value: Any) -> Path | None:
    clean = str(value or "").strip()
    if not clean:
        return None
    path = Path(clean).expanduser()
    if not path.is_absolute():
        path = case_dir / path
    return path.resolve() if path.is_file() else None


def _speaker_attribution_status(
    *,
    case_dir: Path,
    session_dir: Path,
    transcript_material_id: str,
    attributed_transcript_path: Path | None,
) -> dict[str, Any]:
    if attributed_transcript_path is not None:
        resolved = attributed_transcript_path.expanduser().resolve()
        if not resolved.is_file():
            raise CaseWorkspaceError(
                f"attributed transcript does not exist: {resolved}"
            )
        return {
            "status": "available",
            "attributed_transcript_path": _relative_path(case_dir, resolved),
            "source": "explicit_path",
        }

    material = (
        _find_material(case_dir, transcript_material_id)
        if transcript_material_id
        else None
    )
    attribution_task_path: Path | None = None
    if material is not None:
        metadata = material.get("source_metadata")
        material_path = Path(str(material.get("path", ""))).expanduser()
        if isinstance(metadata, dict):
            attribution_task_path = _path_from_metadata(
                case_dir,
                metadata.get("speaker_attribution_task"),
            )
        if (
            isinstance(metadata, dict)
            and str(metadata.get("speaker_attribution", "")).strip()
            and material_path.is_file()
        ):
            return {
                "status": "available",
                "attributed_transcript_path": _relative_path(case_dir, material_path),
                "source": "transcript_material_metadata",
                "note": str(metadata["speaker_attribution"]),
            }

    default_attributed = _path_if_file(session_dir / "attributed_transcript.md")
    if default_attributed is not None:
        return {
            "status": "available",
            "attributed_transcript_path": _relative_path(case_dir, default_attributed),
            "source": "voice_session_file",
        }
    default_task = _path_if_file(session_dir / "speaker_attribution_task.md")
    if default_task is not None:
        attribution_task_path = default_task
    return {
        "status": "pending",
        "attributed_transcript_path": None,
        "attribution_task_path": (
            _relative_path(case_dir, attribution_task_path)
            if attribution_task_path is not None
            else None
        ),
        "source": "not_found",
        "required_before_gate": True,
    }


def _merge_input_status(
    *,
    deck_path: Path,
    session_dir: Path,
    case_dir: Path,
    now: datetime | None,
) -> tuple[dict[str, Any], Path | None]:
    report_path = session_dir / "deck_merge_input_report.json"
    try:
        result = prepare_editable_pptx_merge_input(
            deck_path,
            report_path=report_path,
            now=now,
        )
    except CaseWorkspaceError as error:
        return (
            {
                "status": "blocked",
                "error": str(error),
                "report_path": (
                    _relative_path(case_dir, report_path)
                    if report_path.exists()
                    else None
                ),
            },
            report_path if report_path.exists() else None,
        )
    return (
        {
            "status": result.status,
            "merge_base_path": _relative_path(case_dir, result.merge_base_path),
            "report_path": _relative_path(case_dir, result.report_path),
            "editable_merge_guard": (
                "normalized_required"
                if result.source_legacy_media
                else "normalization_not_required"
            ),
            "source_legacy_media_count": len(result.source_legacy_media),
            "merge_base_legacy_media_count": len(result.merge_base_legacy_media),
        },
        result.report_path,
    )


def _next_actions(
    *,
    speaker_attribution: Mapping[str, Any],
    deck_attached: bool,
    deck_style: Mapping[str, Any],
    feedback_timeline_path: Path | None,
) -> list[str]:
    actions: list[str] = []
    if speaker_attribution.get("status") != "available":
        actions.append(
            "Create a local speaker-attributed transcript from transcript text and source metadata."
        )
    actions.append(
        "Run a model-led gate: decide whether this is a deck-correction call for an existing deck."
    )
    if not deck_attached:
        actions.append(
            "If the gate passes, request or attach the PPTX deck and rerun this intake with --deck or --deck-material-id."
        )
    if deck_style.get("status") != "available":
        actions.append(
            "Resolve the company/deck style authority before producing edit instructions or changing the PPTX."
        )
    else:
        actions.append(
            "Apply the resolved deck style spec to every added, rebuilt, or materially reformatted slide."
        )
    if feedback_timeline_path is None:
        actions.append(
            "Use the raw video manually for visual context because no feedback timeline is available."
        )
    else:
        actions.append(
            "Use feedback timeline rows only as visual evidence when row-level evidence status allows it."
        )
    actions.append(
        "After the gate passes and the deck is attached, produce slide-specific edit instructions before editing the PPTX."
    )
    actions.append(
        "Run build_deck_revision_workbench.py to create the local workbench, output schema, and Codex prompt for those slide-specific instructions."
    )
    actions.append(
        "Use advisory-output-shaper rules for useful, room-safe advisory wording before writing slide changes."
    )
    return actions


def _gate_markdown(payload: Mapping[str, Any]) -> str:
    evidence = payload["evidence"]
    deck = payload["deck"]
    slide_matching = evidence.get("slide_matching")
    slide_matching_status = (
        slide_matching.get("status")
        if isinstance(slide_matching, Mapping)
        else "not available"
    )
    speaker = payload["speaker_attribution"]
    company = payload["company_profile"]
    deck_style = payload["deck_style"]
    advisory_method = payload["advisory_method"]
    lines = [
        "# Deck Revision Gate",
        "",
        "This is a local Codex/Clara review artifact. The hosted server captured files only; it did not decide whether this is a deck-correction call.",
        "",
        "## Required Order",
        "",
        "1. Finish speaker attribution if it is still pending.",
        "2. Decide whether the session is actually about correcting an existing deck.",
        "3. If the gate passes and no deck is attached, request or attach the PPTX.",
        "4. Resolve the company/deck style authority before creating or editing slides.",
        "5. Use the attributed transcript, feedback timeline, screen video, deck snapshot, and style spec to create slide-specific edit instructions.",
        "",
        "## Current Inputs",
        "",
        f"- Speaker attribution: `{speaker['status']}`",
        f"- Raw transcript: `{evidence.get('raw_transcript_path') or 'not found'}`",
        f"- Attributed transcript: `{speaker.get('attributed_transcript_path') or 'not found'}`",
        f"- Speaker attribution task: `{speaker.get('attribution_task_path') or 'not found'}`",
        f"- Feedback timeline: `{evidence.get('feedback_timeline_path') or 'not found'}`",
        f"- Slide matching: `{slide_matching_status or 'not available'}`",
        f"- Screen video: `{evidence.get('screen_video_path') or 'not found'}`",
        f"- Deck status: `{deck['status']}`",
        f"- Deck path: `{deck.get('path') or 'not attached'}`",
        f"- Deck snapshot: `{deck.get('snapshot_path') or 'not available'}`",
        f"- Company profile: `{company['status']}`",
        f"- Company profile path: `{company.get('path') or 'not found'}`",
        f"- Deck style status: `{deck_style['status']}`",
        f"- Deck style key: `{deck_style.get('style_key') or 'not resolved'}`",
        f"- Deck style name: `{deck_style.get('style_name') or 'not resolved'}`",
        f"- Deck style spec: `{deck_style.get('snapshot_path') or deck_style.get('spec_path') or 'not found'}`",
        f"- Advisory method: `{advisory_method['skill']}`",
        "",
        "## Gate Question",
        "",
        "Is this imported voice session a call whose purpose is correcting an existing deck?",
        "",
        "Return a structured answer with: `is_deck_correction_call`, `confidence`, `rationale`, `needs_deck_file`, `deck_file_required_reason`, and `next_action`.",
        "",
        "Rules:",
        "",
        "- Do not decide the gate with keyword matching. This is semantic judgement.",
        "- Do not infer requested deck changes with deterministic rules. Semantic understanding belongs to Codex/Clara reasoning.",
        "- If the call is not a deck-correction call, ignore the video for deck editing.",
        "- If the deck is not visible or no deck file is attached, say what is missing instead of inventing slide edits.",
        "- If no deck style is resolved, stop before editing and ask for the company profile, deck style, or style spec.",
        "- Use the resolved style spec as visual authority for added, rebuilt, or materially reformatted slides.",
        "- Use advisory-output-shaper as the advisory method: convert feedback into useful, room-safe, evidence-aware slide changes.",
        "- Distinguish explicit requested changes from Codex/Clara inference.",
        "- Do not apply PPTX edits in this gate step.",
        "",
        "## Next Actions",
        "",
    ]
    lines.extend(f"- {action}" for action in payload["next_actions"])
    lines.append("")
    return "\n".join(lines)


def prepare_voice_deck_revision_intake(
    case_dir: Path,
    *,
    voice_session: Path | None = None,
    transcript_material_id: str = "",
    deck_path: Path | None = None,
    deck_material_id: str = "",
    attributed_transcript_path: Path | None = None,
    company_profile_path: Path | None = None,
    deck_style: str = "",
    style_spec_path: Path | None = None,
    now: datetime | None = None,
) -> VoiceDeckRevisionIntakeResult:
    """Write local artifacts connecting a voice session to optional deck context.

    Deterministic code only collects evidence, validates paths, extracts PPTX
    text/object context, and enforces output shape. The deck-correction gate is
    semantic, so Codex/Clara must decide it from the produced review pack.
    """

    errors = validate_case_workspace(case_dir)
    if errors:
        raise CaseWorkspaceError("; ".join(errors))

    session_dir = _resolve_voice_session_dir(case_dir, voice_session)
    raw_transcript_path = _path_if_file(session_dir / "raw_transcript.md")
    feedback_timeline_path = _path_if_file(session_dir / "feedback_timeline.json")
    video_timeline_path = _path_if_file(session_dir / "video_timeline.json")
    screen_video_path = _detect_video_path(case_dir, session_dir)
    manifest = _read_json(case_dir / "case_manifest.json")
    company_profile_path_resolved, company_profile = _load_company_profile(
        case_dir,
        company_profile_path,
    )
    company_profile_status = _company_profile_payload(
        case_dir=case_dir,
        profile_path=company_profile_path_resolved,
        profile=company_profile,
    )
    deck_style_status, style_spec_snapshot_path = _resolve_deck_style(
        case_dir=case_dir,
        session_dir=session_dir,
        manifest=manifest,
        company_profile_path=company_profile_path_resolved,
        company_profile=company_profile,
        deck_style=deck_style,
        style_spec_path=style_spec_path,
    )
    advisory_method = _advisory_method_payload(case_dir)
    speaker_attribution = _speaker_attribution_status(
        case_dir=case_dir,
        session_dir=session_dir,
        transcript_material_id=transcript_material_id,
        attributed_transcript_path=attributed_transcript_path,
    )

    resolved_deck_path, resolved_deck_material_id = _resolve_deck_path(
        case_dir=case_dir,
        deck_path=deck_path,
        deck_material_id=deck_material_id,
    )
    deck_snapshot_path: Path | None = None
    merge_input_report_path: Path | None = None
    deck_payload: dict[str, Any]
    if resolved_deck_path is None:
        deck_payload = {
            "status": "missing",
            "path": None,
            "material_id": None,
            "snapshot_path": None,
            "editable_merge_input": {"status": "not_prepared"},
        }
    else:
        deck_snapshot_path = session_dir / "deck_snapshot.json"
        _extract_pptx_snapshot(
            deck_path=resolved_deck_path,
            deck_material_id=resolved_deck_material_id,
            case_dir=case_dir,
            output_path=deck_snapshot_path,
            now=now,
        )
        merge_status, merge_input_report_path = _merge_input_status(
            deck_path=resolved_deck_path,
            session_dir=session_dir,
            case_dir=case_dir,
            now=now,
        )
        deck_payload = {
            "status": "attached",
            "path": _relative_path(case_dir, resolved_deck_path),
            "material_id": resolved_deck_material_id or None,
            "snapshot_path": _relative_path(case_dir, deck_snapshot_path),
            "editable_merge_input": merge_status,
        }

    slide_matching_status: dict[str, Any] | None = None
    if (
        resolved_deck_path is not None
        and deck_snapshot_path is not None
        and feedback_timeline_path is not None
    ):
        try:
            matched_timeline = match_feedback_timeline_to_deck(
                feedback_timeline_path=feedback_timeline_path,
                deck_path=resolved_deck_path,
                deck_snapshot_path=deck_snapshot_path,
                base_dir=case_dir,
                now=now,
            )
            raw_status = matched_timeline.get("slide_matching")
            if isinstance(raw_status, dict):
                slide_matching_status = {
                    "status": raw_status.get("status"),
                    "method": raw_status.get("method"),
                    "slide_count": raw_status.get("slide_count"),
                    "summary": raw_status.get("summary"),
                    "reason": raw_status.get("reason"),
                }
        except (OSError, SlideFrameMatchError) as error:
            LOGGER.warning("Could not match feedback frames to deck slides: %s", error)
            slide_matching_status = {
                "status": "error",
                "reason": str(error),
            }

    intake_path = session_dir / "deck_revision_intake.json"
    gate_path = session_dir / "deck_revision_gate.md"
    payload: dict[str, Any] = {
        "schema_version": 1,
        "source": "clara_voice_deck_revision_intake",
        "created_at": _now_iso(now),
        "voice_session": _relative_path(case_dir, session_dir),
        "transcript_material_id": transcript_material_id or None,
        "company_profile": company_profile_status,
        "deck_style": deck_style_status,
        "advisory_method": advisory_method,
        "speaker_attribution": speaker_attribution,
        "deck_correction_gate": {
            "status": "pending_model_review",
            "question": "Is this imported voice session a call whose purpose is correcting an existing deck?",
            "model_review_required": True,
            "deterministic_gate_used": False,
            "required_output_schema": {
                "is_deck_correction_call": "boolean",
                "confidence": "high | medium | low",
                "rationale": "string",
                "needs_deck_file": "boolean",
                "deck_file_required_reason": "string",
                "next_action": "ignore_video | request_deck | produce_edit_instructions",
            },
        },
        "evidence": {
            "raw_transcript_path": (
                _relative_path(case_dir, raw_transcript_path)
                if raw_transcript_path is not None
                else None
            ),
            "feedback_timeline_path": (
                _relative_path(case_dir, feedback_timeline_path)
                if feedback_timeline_path is not None
                else None
            ),
            "video_timeline_path": (
                _relative_path(case_dir, video_timeline_path)
                if video_timeline_path is not None
                else None
            ),
            "screen_video_path": (
                _relative_path(case_dir, screen_video_path)
                if screen_video_path is not None
                else None
            ),
            "slide_matching": slide_matching_status,
        },
        "deck": deck_payload,
        "next_actions": _next_actions(
            speaker_attribution=speaker_attribution,
            deck_attached=resolved_deck_path is not None,
            deck_style=deck_style_status,
            feedback_timeline_path=feedback_timeline_path,
        ),
    }
    _write_json(intake_path, payload)
    gate_path.write_text(_gate_markdown(payload), encoding="utf-8")

    return VoiceDeckRevisionIntakeResult(
        session_dir=session_dir,
        intake_path=intake_path,
        gate_path=gate_path,
        deck_snapshot_path=deck_snapshot_path,
        merge_input_report_path=merge_input_report_path,
        style_spec_snapshot_path=style_spec_snapshot_path,
    )


def main() -> int:
    """Run deck-revision intake preparation."""

    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("case_dir", type=Path)
    parser.add_argument("--voice-session", type=Path)
    parser.add_argument("--transcript-material-id", default="")
    parser.add_argument("--deck", type=Path)
    parser.add_argument("--deck-material-id", default="")
    parser.add_argument("--attributed-transcript", type=Path)
    parser.add_argument("--company-profile", type=Path)
    parser.add_argument("--deck-style", default="")
    parser.add_argument("--style-spec", type=Path)
    args = parser.parse_args()
    logging.basicConfig(level=logging.INFO, format="%(message)s")

    result = prepare_voice_deck_revision_intake(
        args.case_dir,
        voice_session=args.voice_session,
        transcript_material_id=args.transcript_material_id,
        deck_path=args.deck,
        deck_material_id=args.deck_material_id,
        attributed_transcript_path=args.attributed_transcript,
        company_profile_path=args.company_profile,
        deck_style=args.deck_style,
        style_spec_path=args.style_spec,
    )
    LOGGER.info("Deck revision intake: %s", result.intake_path)
    LOGGER.info("Deck revision gate: %s", result.gate_path)
    if result.deck_snapshot_path is not None:
        LOGGER.info("Deck snapshot: %s", result.deck_snapshot_path)
    if result.style_spec_snapshot_path is not None:
        LOGGER.info("Deck style spec: %s", result.style_spec_snapshot_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
