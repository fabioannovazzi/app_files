#!/usr/bin/env python3
"""Prepare and mechanically audit Clara advisory-deliverable validation runs.

The helper is deterministic because supported-format extraction, hashing,
declared-schema validation, cross-field consistency, and packaging are
mechanically verifiable. It does not select material claims or judge support,
reasoning, recommendations, corrections, or professional readiness.
"""

from __future__ import annotations

import argparse
import hashlib
import html
import json
import logging
import re
import shutil
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

__all__ = [
    "AdvisoryValidationError",
    "prepare_validation",
    "package_validation",
    "read_supported_deliverable",
    "validate_advisory_contract",
    "validate_review_record",
]

LOGGER = logging.getLogger(__name__)
CLARA_ROOT = Path(__file__).resolve().parents[3]
REPOSITORY_ROOT = Path(__file__).resolve().parents[5]
SUPPORTED_PRIMARY_SUFFIXES = {
    ".docx",
    ".htm",
    ".html",
    ".markdown",
    ".md",
    ".pdf",
    ".pptx",
    ".txt",
}
SUPPORTED_SOURCE_SUFFIXES = SUPPORTED_PRIMARY_SUFFIXES | {".csv", ".parquet", ".xlsx"}
LANGUAGES = {"it", "en", "fr", "de", "es"}
REVIEW_DIMENSIONS = (
    "contract_conformance",
    "factual_source_support",
    "calculations_data_provenance",
    "reasoning_assumptions",
    "contradictions_missing_evidence",
    "recommendation_evidence_decision_fit",
    "professional_judgement_boundaries",
    "correction_needs",
    "residual_uncertainty",
    "delivery_readiness",
)
REQUIRED_CONTRACT_FIELDS = {
    "schema_version",
    "decision",
    "purpose",
    "audience",
    "deliverable_type",
    "output_language",
    "scope_included",
    "scope_excluded",
    "available_inputs",
    "evidence_requirements",
    "analysis_plan",
    "assumptions",
    "unresolved_questions",
    "success_criteria",
    "selected_clara_workflow",
    "validation_profile",
    "validation_scope",
    "correction_policy",
    "professional_judgement_policy",
}
DIMENSION_STATUSES = {
    "conforms",
    "partially_conforms",
    "does_not_conform",
    "contradicted",
    "uncertain",
    "judgment_required",
    "not_applicable",
}
CORRECTION_STATUSES = {
    "not_needed",
    "proposed",
    "completed",
    "blocked",
    "professional_review_required",
}
FORMAT_CHECK_STATUSES = {
    "passed",
    "issues_found",
    "blocked",
    "not_run",
    "not_applicable",
}
READINESS_STATUSES = {
    "ready",
    "ready_with_residual_uncertainty",
    "not_ready",
    "blocked",
}
URL_RE = re.compile(r"https?://[^\s)\]>\"']+", re.IGNORECASE)
CITATION_RE = re.compile(r"\[(?:\^?[A-Za-z0-9_-]+|\d+(?:,\s*\d+)*)\]")
FOOTNOTE_RE = re.compile(r"^\[\^?([A-Za-z0-9_-]+)\]:\s*(.+)$", re.MULTILINE)
NUMBER_RE = re.compile(
    r"(?<![\w.])(?:(?:EUR|USD|GBP|[$€£])\s*)?-?\d[\d.,]*(?:\s*%|\s*(?:bps|EUR|USD|GBP))?(?!\w)",
    re.IGNORECASE,
)
FORMULA_RE = re.compile(
    r"(?:[$€£]?\s*-?\d[\d.,]*\s*)(?:[+*/=]|\s-\s)(?:\s*[$€£]?\s*-?\d[\d.,%]*)"
)


class AdvisoryValidationError(ValueError):
    """Raised when a mechanically verifiable validation contract is invalid."""


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag == "a":
            for key, value in attrs:
                if key.casefold() == "href" and value:
                    self.parts.append(f" {value} ")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"p", "div", "br", "li", "section", "article", "h1", "h2", "h3"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        self.parts.append(html.unescape(data))

    def text(self) -> str:
        return re.sub(r"\n{3,}", "\n\n", "".join(self.parts)).strip()


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _write_json(path: Path, payload: Any) -> None:
    path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def _load_json(path: Path) -> Any:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        raise AdvisoryValidationError(f"invalid JSON in {path}: {exc}") from exc


def _is_non_empty_string(value: Any) -> bool:
    return isinstance(value, str) and bool(value.strip())


def _is_string_list(value: Any, *, allow_empty: bool) -> bool:
    return (
        isinstance(value, list)
        and (allow_empty or bool(value))
        and all(_is_non_empty_string(item) for item in value)
    )


def _within(path: Path, parent: Path) -> bool:
    try:
        path.resolve().relative_to(parent.resolve())
    except ValueError:
        return False
    return True


def _ensure_output_location(output_dir: Path) -> None:
    if _within(output_dir, REPOSITORY_ROOT):
        raise AdvisoryValidationError(
            "validation outputs must be outside the Git workspace"
        )


def validate_advisory_contract(payload: Any) -> list[str]:
    """Return mechanical shape errors for advisory_contract.json."""

    if not isinstance(payload, dict):
        return ["advisory contract must be an object"]
    errors: list[str] = []
    missing = sorted(REQUIRED_CONTRACT_FIELDS - payload.keys())
    if missing:
        errors.append(f"missing contract fields: {', '.join(missing)}")
        return errors
    if payload["schema_version"] != "1.0":
        errors.append('schema_version must be "1.0"')
    for field in ("decision", "purpose", "deliverable_type", "selected_clara_workflow"):
        if not _is_non_empty_string(payload[field]):
            errors.append(f"{field} must be a non-empty string")
    if payload["output_language"] not in LANGUAGES:
        errors.append("output_language must be it, en, fr, de, or es")
    list_rules = {
        "audience": False,
        "scope_included": False,
        "scope_excluded": True,
        "available_inputs": False,
        "evidence_requirements": False,
        "analysis_plan": False,
        "assumptions": True,
        "unresolved_questions": True,
        "success_criteria": False,
    }
    for field, allow_empty in list_rules.items():
        if not _is_string_list(payload[field], allow_empty=allow_empty):
            qualifier = "possibly empty" if allow_empty else "non-empty"
            errors.append(f"{field} must be a {qualifier} array of non-empty strings")

    profile = payload["validation_profile"]
    if not isinstance(profile, dict):
        errors.append("validation_profile must be an object")
    else:
        dimensions = profile.get("review_dimensions")
        if not isinstance(dimensions, list) or tuple(dimensions) != REVIEW_DIMENSIONS:
            errors.append(
                "validation_profile.review_dimensions must use the fixed ordered set"
            )
        format_checks = profile.get("format_checks")
        if not isinstance(format_checks, list):
            errors.append("validation_profile.format_checks must be an array")
        else:
            workflows: list[str] = []
            for index, item in enumerate(format_checks):
                subject = f"validation_profile.format_checks[{index}]"
                if not isinstance(item, dict):
                    errors.append(f"{subject} must be an object")
                    continue
                if not all(
                    _is_non_empty_string(item.get(field))
                    for field in ("workflow", "reason")
                ):
                    errors.append(f"{subject} requires workflow and reason")
                if item.get("requirement") not in {
                    "required",
                    "if_applicable",
                    "not_required",
                }:
                    errors.append(f"{subject}.requirement is invalid")
                if not _is_string_list(item.get("artifact_refs"), allow_empty=True):
                    errors.append(f"{subject}.artifact_refs must be a string array")
                if _is_non_empty_string(item.get("workflow")):
                    workflows.append(item["workflow"])
            if len(workflows) != len(set(workflows)):
                errors.append(
                    "validation_profile.format_checks workflows must be unique"
                )

    scope = payload["validation_scope"]
    if not isinstance(scope, dict):
        errors.append("validation_scope must be an object")
    else:
        if scope.get("coverage") not in {
            "all_material_content",
            "selected_material_content",
            "limited",
        }:
            errors.append("validation_scope.coverage is invalid")
        for field, allow_empty in (
            ("included_sections", False),
            ("excluded_sections", True),
            ("limitations", True),
        ):
            if not _is_string_list(scope.get(field), allow_empty=allow_empty):
                errors.append(f"validation_scope.{field} must be a string array")

    correction = payload["correction_policy"]
    if not isinstance(correction, dict):
        errors.append("correction_policy must be an object")
    else:
        if correction.get("mode") != "separate_artifact":
            errors.append("correction_policy.mode must be separate_artifact")
        if correction.get("preserve_original") is not True:
            errors.append("correction_policy.preserve_original must be true")
        for field in ("allowed", "approval_required_before_delivery"):
            if not isinstance(correction.get(field), bool):
                errors.append(f"correction_policy.{field} must be boolean")

    judgement = payload["professional_judgement_policy"]
    if not isinstance(judgement, dict):
        errors.append("professional_judgement_policy must be an object")
    else:
        for field in ("owner", "model_role"):
            if not _is_non_empty_string(judgement.get(field)):
                errors.append(f"professional_judgement_policy.{field} is required")
        if not isinstance(judgement.get("approval_required_before_delivery"), bool):
            errors.append(
                "professional_judgement_policy.approval_required_before_delivery must be boolean"
            )
    return errors


def _read_pdf(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError as exc:
        raise AdvisoryValidationError(
            "PyMuPDF is required for PDF extraction; run Clara's dependency check"
        ) from exc
    with fitz.open(path) as document:
        pages = [page.get_text("text") for page in document]
    text = "\n\n".join(pages).strip()
    if not text:
        raise AdvisoryValidationError(
            "PDF has no readable text layer; run Clara's input-aware OCR preflight"
        )
    return text, {"parser": "pymupdf_text", "page_count": len(pages)}


def _read_docx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from docx import Document  # type: ignore[import-not-found]
    except ImportError as exc:
        raise AdvisoryValidationError(
            "python-docx is required for DOCX extraction; run Clara's dependency check"
        ) from exc
    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    table_cell_count = 0
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            table_cell_count += len(values)
            parts.append(" | ".join(values))
    return "\n\n".join(parts).strip(), {
        "parser": "python_docx",
        "paragraph_count": len(document.paragraphs),
        "table_count": len(document.tables),
        "table_cell_count": table_cell_count,
    }


def _read_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:
        raise AdvisoryValidationError(
            "python-pptx is required for PPTX extraction; run Clara's dependency check"
        ) from exc
    presentation = Presentation(path)
    parts: list[str] = []
    text_shape_count = 0
    table_cell_count = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"# Slide {slide_number}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                parts.append(shape.text.strip())
                text_shape_count += 1
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    table_cell_count += len(values)
                    parts.append(" | ".join(values))
    return "\n\n".join(parts).strip(), {
        "parser": "python_pptx_visible_text",
        "slide_count": len(presentation.slides),
        "text_shape_count": text_shape_count,
        "table_cell_count": table_cell_count,
        "limitations": [
            "Visible text and table cells were extracted; visual meaning, charts, images, and speaker notes require separate review."
        ],
    }


def read_supported_deliverable(path: Path) -> tuple[str, dict[str, Any]]:
    """Extract text and mechanical parser metadata from one supported artifact."""

    suffix = path.suffix.casefold()
    if suffix not in SUPPORTED_PRIMARY_SUFFIXES:
        supported = ", ".join(sorted(SUPPORTED_PRIMARY_SUFFIXES))
        raise AdvisoryValidationError(
            f"unsupported primary deliverable format {suffix or '<none>'}; supported: {supported}"
        )
    if suffix in {".md", ".markdown", ".txt"}:
        return path.read_text(encoding="utf-8", errors="replace").strip(), {
            "parser": "plain_text"
        }
    if suffix in {".html", ".htm"}:
        extractor = _HTMLTextExtractor()
        extractor.feed(path.read_text(encoding="utf-8", errors="replace"))
        return extractor.text(), {"parser": "html_text"}
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    return _read_pptx(path)


def _ordered_unique(
    items: list[str], *, limit: int, strip_terminal_punctuation: bool = True
) -> list[str]:
    seen: set[str] = set()
    output: list[str] = []
    for item in items:
        cleaned = item.strip()
        if strip_terminal_punctuation:
            cleaned = cleaned.rstrip(".,;:!?)\\]}>'\"")
        if not cleaned or cleaned.casefold() in seen:
            continue
        seen.add(cleaned.casefold())
        output.append(cleaned)
        if len(output) >= limit:
            break
    return output


def _citation_inventory(text: str) -> dict[str, Any]:
    return {
        "schema_version": "1.0",
        "purpose": "Mechanical navigation inventory; not semantic support assessment.",
        "urls": _ordered_unique(URL_RE.findall(text), limit=500),
        "citation_markers": _ordered_unique(
            CITATION_RE.findall(text),
            limit=1000,
            strip_terminal_punctuation=False,
        ),
        "footnotes": [
            {"id": match.group(1), "text": match.group(2).strip()}
            for match in FOOTNOTE_RE.finditer(text)
        ][:500],
    }


def _calculation_inventory(text: str) -> dict[str, Any]:
    return {
        "schema_version": "1.0",
        "purpose": "Mechanical numeric-token inventory; not calculation verification or semantic materiality selection.",
        "numeric_tokens": _ordered_unique(NUMBER_RE.findall(text), limit=2000),
        "formula_like_fragments": _ordered_unique(FORMULA_RE.findall(text), limit=500),
    }


def _copy_contract(source: Path, output_dir: Path) -> Path:
    target = output_dir / "advisory_contract.json"
    if source.resolve() == target.resolve():
        return target
    if target.exists() and target.read_bytes() != source.read_bytes():
        raise AdvisoryValidationError(
            f"refusing to overwrite a different advisory contract: {target}"
        )
    shutil.copyfile(source, target)
    return target


def _source_inventory(paths: list[Path]) -> dict[str, Any]:
    items: list[dict[str, Any]] = []
    seen: set[Path] = set()
    for source in paths:
        resolved = source.resolve()
        if resolved in seen:
            continue
        seen.add(resolved)
        if not source.is_file():
            raise AdvisoryValidationError(f"source file does not exist: {source}")
        items.append(
            {
                "path": str(resolved),
                "name": source.name,
                "suffix": source.suffix.casefold(),
                "supported_source_type": source.suffix.casefold()
                in SUPPORTED_SOURCE_SUFFIXES,
                "byte_count": source.stat().st_size,
                "sha256": _sha256(source),
            }
        )
    return {
        "schema_version": "1.0",
        "purpose": "Selected-source identity inventory; contents remain subject to model-led review.",
        "sources": items,
    }


def prepare_validation(
    deliverable: Path,
    advisory_contract: Path,
    output_dir: Path,
    *,
    source_files: list[Path] | None = None,
) -> dict[str, Path]:
    """Write deterministic preparation artifacts for a validator run."""

    _ensure_output_location(output_dir)
    if not deliverable.is_file():
        raise AdvisoryValidationError(f"deliverable does not exist: {deliverable}")
    if not advisory_contract.is_file():
        raise AdvisoryValidationError(
            "advisory_contract.json is required; Clara must create it from explicit or user-confirmed context"
        )
    contract_payload = _load_json(advisory_contract)
    contract_errors = validate_advisory_contract(contract_payload)
    if contract_errors:
        raise AdvisoryValidationError("; ".join(contract_errors))
    text, parser_metadata = read_supported_deliverable(deliverable)
    if not text:
        raise AdvisoryValidationError(
            "deliverable extraction produced no readable text"
        )

    output_dir.mkdir(parents=True, exist_ok=True)
    contract_copy = _copy_contract(advisory_contract, output_dir)
    source_inventory = _source_inventory(source_files or [])
    deliverable_hash = _sha256(deliverable)
    inventory = {
        "schema_version": "1.0",
        "source_path": str(deliverable.resolve()),
        "source_name": deliverable.name,
        "source_suffix": deliverable.suffix.casefold(),
        "source_sha256": deliverable_hash,
        "source_byte_count": deliverable.stat().st_size,
        "advisory_contract_path": str(contract_copy.resolve()),
        "advisory_contract_sha256": _sha256(contract_copy),
        "extraction": {
            **parser_metadata,
            "character_count": len(text),
            "word_count": len(re.findall(r"\S+", text)),
        },
        "boundary": {
            "semantic_selection": "model_led",
            "semantic_assessment": "model_led",
            "hidden_model_api_calls": False,
            "original_preserved": True,
        },
    }
    extracted_path = output_dir / "extracted_deliverable.md"
    extracted_path.write_text(text.rstrip() + "\n", encoding="utf-8")
    paths = {
        "advisory_contract": contract_copy,
        "deliverable_inventory": output_dir / "deliverable_inventory.json",
        "extracted_deliverable": extracted_path,
        "citation_inventory": output_dir / "citation_inventory.json",
        "calculation_inventory": output_dir / "calculation_inventory.json",
        "source_inventory": output_dir / "source_inventory.json",
    }
    _write_json(paths["deliverable_inventory"], inventory)
    _write_json(paths["citation_inventory"], _citation_inventory(text))
    _write_json(paths["calculation_inventory"], _calculation_inventory(text))
    _write_json(paths["source_inventory"], source_inventory)
    return paths


def _validate_dimension_review(value: Any, *, subject: str) -> list[str]:
    if not isinstance(value, dict):
        return [f"{subject} must be an object"]
    errors: list[str] = []
    required = {
        "status",
        "analysis",
        "evidence_refs",
        "issues",
        "correction_status",
        "professional_review_required",
    }
    missing = sorted(required - value.keys())
    if missing:
        errors.append(f"{subject} missing fields: {', '.join(missing)}")
        return errors
    if value["status"] not in DIMENSION_STATUSES:
        errors.append(f"{subject}.status is invalid")
    if not _is_non_empty_string(value["analysis"]):
        errors.append(f"{subject}.analysis is required")
    for field in ("evidence_refs", "issues"):
        if not _is_string_list(value[field], allow_empty=True):
            errors.append(f"{subject}.{field} must be a string array")
    if value["correction_status"] not in CORRECTION_STATUSES:
        errors.append(f"{subject}.correction_status is invalid")
    if not isinstance(value["professional_review_required"], bool):
        errors.append(f"{subject}.professional_review_required must be boolean")
    return errors


def validate_review_record(payload: Any, contract: dict[str, Any]) -> list[str]:
    """Return mechanical shape and internal-consistency errors for a review."""

    if not isinstance(payload, dict):
        return ["review record must be an object"]
    errors: list[str] = []
    required = {
        "schema_version",
        "language",
        "advisory_contract_sha256",
        "deliverable_sha256",
        "coverage_review",
        "dimension_reviews",
        "findings",
        "format_specific_checks",
        "correction",
        "overall_assessment",
        "delivery_readiness",
    }
    missing = sorted(required - payload.keys())
    if missing:
        return [f"review record missing fields: {', '.join(missing)}"]
    if payload["schema_version"] != "1.0":
        errors.append('review schema_version must be "1.0"')
    if payload["language"] != contract["output_language"]:
        errors.append("review language must match the advisory contract")
    for field in ("advisory_contract_sha256", "deliverable_sha256"):
        value = payload[field]
        if not isinstance(value, str) or re.fullmatch(r"[0-9a-f]{64}", value) is None:
            errors.append(f"{field} must be a lowercase SHA-256")

    coverage = payload["coverage_review"]
    if not isinstance(coverage, dict):
        errors.append("coverage_review must be an object")
    else:
        if coverage.get("selection_method") != "model_led_materiality_review":
            errors.append(
                "coverage_review.selection_method must be model_led_materiality_review"
            )
        if coverage.get("scope") != contract["validation_scope"]["coverage"]:
            errors.append(
                "coverage_review.scope must match advisory_contract.validation_scope"
            )
        if not _is_string_list(coverage.get("reviewed_sections"), allow_empty=False):
            errors.append("coverage_review.reviewed_sections must be non-empty")
        for field in ("omitted_sections", "limitations"):
            if not _is_string_list(coverage.get(field), allow_empty=True):
                errors.append(f"coverage_review.{field} must be a string array")
        if not _is_non_empty_string(coverage.get("analysis")):
            errors.append("coverage_review.analysis is required")

    dimensions = payload["dimension_reviews"]
    if not isinstance(dimensions, dict) or set(dimensions) != set(REVIEW_DIMENSIONS):
        errors.append(
            "dimension_reviews must contain exactly the ten required dimensions"
        )
    else:
        for dimension in REVIEW_DIMENSIONS:
            errors.extend(
                _validate_dimension_review(
                    dimensions[dimension], subject=f"dimension_reviews.{dimension}"
                )
            )

    findings = payload["findings"]
    finding_ids: list[str] = []
    if not isinstance(findings, list):
        errors.append("findings must be an array")
    else:
        for index, finding in enumerate(findings):
            subject = f"findings[{index}]"
            if not isinstance(finding, dict):
                errors.append(f"{subject} must be an object")
                continue
            if not _is_non_empty_string(finding.get("id")):
                errors.append(f"{subject}.id is required")
            else:
                finding_ids.append(finding["id"])
            if finding.get("dimension") not in REVIEW_DIMENSIONS:
                errors.append(f"{subject}.dimension is invalid")
            if not _is_non_empty_string(finding.get("finding")):
                errors.append(f"{subject}.finding is required")
            if finding.get("status") not in DIMENSION_STATUSES:
                errors.append(f"{subject}.status is invalid")
            if not _is_string_list(finding.get("evidence_refs"), allow_empty=True):
                errors.append(f"{subject}.evidence_refs must be a string array")
            if not isinstance(finding.get("correction_action"), str):
                errors.append(f"{subject}.correction_action must be a string")
            if finding.get("correction_status") not in CORRECTION_STATUSES:
                errors.append(f"{subject}.correction_status is invalid")
            if not isinstance(finding.get("professional_review_required"), bool):
                errors.append(f"{subject}.professional_review_required must be boolean")
        if len(finding_ids) != len(set(finding_ids)):
            errors.append("finding ids must be unique")

    checks = payload["format_specific_checks"]
    check_by_workflow: dict[str, dict[str, Any]] = {}
    if not isinstance(checks, list):
        errors.append("format_specific_checks must be an array")
    else:
        for index, check in enumerate(checks):
            subject = f"format_specific_checks[{index}]"
            if not isinstance(check, dict):
                errors.append(f"{subject} must be an object")
                continue
            workflow = check.get("workflow")
            if not _is_non_empty_string(workflow):
                errors.append(f"{subject}.workflow is required")
            elif workflow in check_by_workflow:
                errors.append(f"duplicate format-specific check: {workflow}")
            else:
                check_by_workflow[workflow] = check
            if check.get("status") not in FORMAT_CHECK_STATUSES:
                errors.append(f"{subject}.status is invalid")
            if not _is_string_list(check.get("artifact_refs"), allow_empty=True):
                errors.append(f"{subject}.artifact_refs must be a string array")
            if not _is_non_empty_string(check.get("analysis")):
                errors.append(f"{subject}.analysis is required")
            if check.get("status") == "passed" and not check.get("artifact_refs"):
                errors.append(f"{subject} passed without an artifact reference")

    for declared in contract["validation_profile"]["format_checks"]:
        if declared["requirement"] != "required":
            continue
        workflow = declared["workflow"]
        actual = check_by_workflow.get(workflow)
        if actual is None:
            errors.append(f"required format check is missing: {workflow}")
        elif actual.get("status") != "passed":
            errors.append(f"required format check did not pass: {workflow}")

    correction = payload["correction"]
    if not isinstance(correction, dict):
        errors.append("correction must be an object")
    else:
        if correction.get("status") not in {
            "not_required",
            "required",
            "completed",
            "blocked",
        }:
            errors.append("correction.status is invalid")
        if not _is_non_empty_string(correction.get("summary")):
            errors.append("correction.summary is required")
        if not isinstance(correction.get("corrected_artifact"), str):
            errors.append("correction.corrected_artifact must be a string")
        if not _is_string_list(correction.get("unresolved_changes"), allow_empty=True):
            errors.append("correction.unresolved_changes must be a string array")
        if (
            correction.get("status") == "completed"
            and not contract["correction_policy"]["allowed"]
        ):
            errors.append(
                "correction completed even though the contract disallows correction"
            )

    overall = payload["overall_assessment"]
    readiness = payload["delivery_readiness"]
    if not isinstance(overall, dict):
        errors.append("overall_assessment must be an object")
    if not isinstance(readiness, dict):
        errors.append("delivery_readiness must be an object")
    if isinstance(overall, dict) and isinstance(readiness, dict):
        outcome = overall.get("outcome")
        status = readiness.get("status")
        if outcome not in READINESS_STATUSES:
            errors.append("overall_assessment.outcome is invalid")
        if status not in READINESS_STATUSES:
            errors.append("delivery_readiness.status is invalid")
        if outcome != status:
            errors.append("overall assessment and delivery readiness must match")
        if not _is_non_empty_string(overall.get("analysis")):
            errors.append("overall_assessment.analysis is required")
        for field in ("residual_uncertainties", "professional_review_items"):
            if not _is_string_list(overall.get(field), allow_empty=True):
                errors.append(f"overall_assessment.{field} must be a string array")
        if not _is_string_list(readiness.get("conditions"), allow_empty=True):
            errors.append("delivery_readiness.conditions must be a string array")

        attention_statuses = {
            "does_not_conform",
            "contradicted",
            "uncertain",
            "judgment_required",
        }
        unresolved_attention = isinstance(dimensions, dict) and any(
            isinstance(dimensions.get(dimension), dict)
            and (
                dimensions[dimension].get("status") in attention_statuses
                or dimensions[dimension].get("correction_status")
                in {"proposed", "blocked", "professional_review_required"}
            )
            for dimension in REVIEW_DIMENSIONS
        )
        if status == "ready" and unresolved_attention:
            errors.append(
                "ready status cannot coexist with unresolved review attention"
            )
        if status == "ready" and (
            overall.get("residual_uncertainties")
            or overall.get("professional_review_items")
        ):
            errors.append(
                "ready status cannot carry residual uncertainty or professional-review items"
            )
        if (
            isinstance(correction, dict)
            and correction.get("status")
            in {
                "required",
                "blocked",
            }
            and status in {"ready", "ready_with_residual_uncertainty"}
        ):
            errors.append("unresolved correction cannot be delivery-ready")
    return errors


def _render_package(
    inventory: dict[str, Any],
    review: dict[str, Any],
    audit: dict[str, Any],
) -> str:
    dimensions = review.get("dimension_reviews", {})
    lines = [
        "# Advisory deliverable validation",
        "",
        f"Original: {inventory.get('source_name', 'unknown')}",
        f"Original SHA-256: {inventory.get('source_sha256', 'unknown')}",
        f"Contract SHA-256: {inventory.get('advisory_contract_sha256', 'unknown')}",
        f"Record complete: {'yes' if audit['record_complete'] else 'no'}",
        f"Delivery readiness: {audit['effective_delivery_readiness']}",
        "",
        "## Review dimensions",
        "",
    ]
    for dimension in REVIEW_DIMENSIONS:
        record = dimensions.get(dimension, {}) if isinstance(dimensions, dict) else {}
        lines.append(
            f"- **{dimension.replace('_', ' ').title()}** — {record.get('status', 'missing')}: {record.get('analysis', 'No analysis recorded.')}"
        )
    lines.extend(["", "## Findings", ""])
    findings = review.get("findings", [])
    if findings:
        for finding in findings:
            lines.append(
                f"- **{finding.get('id', 'finding')}** ({finding.get('dimension', 'unknown')}): {finding.get('finding', '')}"
            )
    else:
        lines.append("- No individual findings recorded.")
    lines.extend(["", "## Format-specific checks", ""])
    checks = review.get("format_specific_checks", [])
    if checks:
        for check in checks:
            lines.append(
                f"- **{check.get('workflow', 'unknown')}** — {check.get('status', 'missing')}: {check.get('analysis', '')}"
            )
    else:
        lines.append("- No format-specific checks recorded.")
    lines.extend(
        [
            "",
            "## Correction",
            "",
            f"Status: {review.get('correction', {}).get('status', 'missing')}",
            "",
            review.get("correction", {}).get(
                "summary", "No correction summary recorded."
            ),
            "",
            "## Residual uncertainty and professional review",
            "",
        ]
    )
    overall = review.get("overall_assessment", {})
    for item in overall.get("residual_uncertainties", []) or []:
        lines.append(f"- Residual uncertainty: {item}")
    for item in overall.get("professional_review_items", []) or []:
        lines.append(f"- Professional review: {item}")
    if not overall.get("residual_uncertainties") and not overall.get(
        "professional_review_items"
    ):
        lines.append("- None recorded.")
    if audit["errors"]:
        lines.extend(["", "## Mechanical audit errors", ""])
        lines.extend(f"- {error}" for error in audit["errors"])
    return "\n".join(lines).rstrip() + "\n"


def package_validation(
    inventory_path: Path,
    review_draft_path: Path,
    advisory_contract_path: Path,
    output_dir: Path,
    *,
    corrected_deliverable: Path | None = None,
) -> tuple[dict[str, Path], dict[str, Any]]:
    """Audit a model-authored review and write the validation package."""

    _ensure_output_location(output_dir)
    inventory = _load_json(inventory_path)
    review = _load_json(review_draft_path)
    contract = _load_json(advisory_contract_path)
    errors = validate_advisory_contract(contract)
    if not isinstance(inventory, dict):
        errors.append("deliverable inventory must be an object")
        inventory = {}
    if isinstance(review, dict) and isinstance(contract, dict) and not errors:
        errors.extend(validate_review_record(review, contract))
    elif not isinstance(review, dict):
        errors.append("review record must be an object")
        review = {}

    original_path = Path(str(inventory.get("source_path", "")))
    expected_original_hash = inventory.get("source_sha256")
    original_unchanged = (
        original_path.is_file() and _sha256(original_path) == expected_original_hash
    )
    if not original_unchanged:
        errors.append("original deliverable is missing or changed since preparation")
    contract_hash = _sha256(advisory_contract_path)
    if inventory.get("advisory_contract_sha256") != contract_hash:
        errors.append("advisory contract no longer matches the preparation inventory")
    if review.get("advisory_contract_sha256") != contract_hash:
        errors.append("review is not bound to the current advisory contract")
    if review.get("deliverable_sha256") != expected_original_hash:
        errors.append("review is not bound to the prepared deliverable")

    corrected_metadata: dict[str, Any] | None = None
    correction_status = review.get("correction", {}).get("status")
    if correction_status == "completed":
        if corrected_deliverable is None or not corrected_deliverable.is_file():
            errors.append("completed correction requires --corrected-deliverable")
        elif corrected_deliverable.resolve() == original_path.resolve():
            errors.append("corrected deliverable must use a separate path")
        else:
            corrected_hash = _sha256(corrected_deliverable)
            if corrected_hash == expected_original_hash:
                errors.append("completed correction has the same bytes as the original")
            corrected_metadata = {
                "path": str(corrected_deliverable.resolve()),
                "name": corrected_deliverable.name,
                "sha256": corrected_hash,
                "byte_count": corrected_deliverable.stat().st_size,
            }
    elif corrected_deliverable is not None:
        errors.append(
            "a corrected deliverable was supplied but correction.status is not completed"
        )

    declared_readiness = review.get("delivery_readiness", {}).get("status", "blocked")
    audit = {
        "schema_version": "1.0",
        "record_complete": not errors,
        "errors": errors,
        "checks": {
            "advisory_contract_shape": not validate_advisory_contract(contract),
            "contract_hash_bound": review.get("advisory_contract_sha256")
            == contract_hash,
            "deliverable_hash_bound": review.get("deliverable_sha256")
            == expected_original_hash,
            "original_unchanged": original_unchanged,
            "separate_corrected_artifact": correction_status != "completed"
            or corrected_metadata is not None,
            "hidden_model_api_calls": False,
        },
        "declared_delivery_readiness": declared_readiness,
        "effective_delivery_readiness": declared_readiness if not errors else "blocked",
        "corrected_artifact": corrected_metadata,
        "interpretation": "Record completeness proves mechanical consistency only; semantic support, reasoning, recommendation quality, and professional judgement remain model-led and professionally reviewed.",
    }
    output_dir.mkdir(parents=True, exist_ok=True)
    paths = {
        "review": output_dir / "advisory_validation_review.json",
        "audit": output_dir / "validation_audit.json",
        "package": output_dir / "advisory_validation_package.md",
    }
    _write_json(paths["review"], review)
    _write_json(paths["audit"], audit)
    paths["package"].write_text(
        _render_package(inventory, review, audit), encoding="utf-8"
    )
    return paths, audit


def _prepare_command(args: argparse.Namespace) -> int:
    paths = prepare_validation(
        args.deliverable,
        args.advisory_contract,
        args.output_dir,
        source_files=args.source_file,
    )
    LOGGER.info("Prepared advisory validation: %s", paths["deliverable_inventory"])
    return 0


def _package_command(args: argparse.Namespace) -> int:
    paths, audit = package_validation(
        args.deliverable_inventory,
        args.review_draft,
        args.advisory_contract,
        args.output_dir,
        corrected_deliverable=args.corrected_deliverable,
    )
    LOGGER.info("Wrote advisory validation package: %s", paths["package"])
    return 0 if audit["record_complete"] else 1


def _parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    subparsers = parser.add_subparsers(dest="command", required=True)
    prepare = subparsers.add_parser(
        "prepare", help="Extract and inventory a deliverable."
    )
    prepare.add_argument("deliverable", type=Path)
    prepare.add_argument("--advisory-contract", type=Path, required=True)
    prepare.add_argument("--output-dir", type=Path, required=True)
    prepare.add_argument("--source-file", type=Path, action="append", default=[])
    prepare.set_defaults(handler=_prepare_command)

    package = subparsers.add_parser(
        "package", help="Audit and package a model-led review."
    )
    package.add_argument("deliverable_inventory", type=Path)
    package.add_argument("review_draft", type=Path)
    package.add_argument("--advisory-contract", type=Path, required=True)
    package.add_argument("--output-dir", type=Path, required=True)
    package.add_argument("--corrected-deliverable", type=Path)
    package.set_defaults(handler=_package_command)
    return parser


def main() -> int:
    logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")
    parser = _parser()
    args = parser.parse_args()
    try:
        return int(args.handler(args))
    except AdvisoryValidationError as exc:
        parser.error(str(exc))
    return 2


if __name__ == "__main__":
    raise SystemExit(main())
