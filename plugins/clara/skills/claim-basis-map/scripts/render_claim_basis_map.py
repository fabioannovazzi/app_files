#!/usr/bin/env python3
"""Render a deterministic Markdown claim-basis map from deck claims JSON."""

from __future__ import annotations

import argparse
import json
import logging
from collections import defaultdict
from pathlib import Path
from typing import Any

__all__ = [
    "compare_current_deck_snapshot",
    "extract_current_deck_snapshot",
    "main",
    "render_claim_basis_map",
]

LOGGER = logging.getLogger(__name__)

BASIS_ORDER = (
    "source-backed",
    "calculated",
    "claim-linked",
    "reasoned",
    "assumption",
    "ungrounded",
)
CLAIM_SHAPE_NAME_PREFIX = "clara-claim:"
TEXT_DRIFT_REQUIRES_REFRESH = {
    "edited",
    "missing-or-edited",
    "duplicate-current-key",
    "reference-broken",
    "untracked-current-text",
}


def _clean_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, str):
        return " ".join(value.split())
    return " ".join(str(value).split())


def _non_empty_list(value: Any) -> list[Any]:
    if not isinstance(value, list):
        return []
    return [item for item in value if item not in (None, "", [], {})]


def _non_empty_dict(value: Any) -> dict[str, Any]:
    if not isinstance(value, dict):
        return {}
    return {
        str(key): item for key, item in value.items() if item not in (None, "", [], {})
    }


def _basis_type(claim: dict[str, Any]) -> str:
    """Classify from explicit metadata only; no semantic judgement is made here."""
    if _non_empty_list(claim.get("source_refs")):
        return "source-backed"
    if _non_empty_dict(claim.get("calculation_ref")):
        return "calculated"
    if _non_empty_list(claim.get("claim_refs")):
        return "claim-linked"
    if _non_empty_list(claim.get("reasoning_inputs")):
        return "reasoned"
    if _clean_text(claim.get("assumption_basis")):
        return "assumption"
    return "ungrounded"


def _source_ref_text(ref: Any) -> str:
    if not isinstance(ref, dict):
        return _clean_text(ref)
    parts = [
        _clean_text(ref.get("title")),
        _clean_text(ref.get("locator")),
        _clean_text(ref.get("url")),
        _clean_text(ref.get("path")),
    ]
    text = ", ".join(part for part in parts if part)
    quote = _clean_text(ref.get("quote"))
    if quote:
        text = f"{text} | Excerpt: {quote}" if text else f"Excerpt: {quote}"
    return text


def _reasoning_input_text(item: Any) -> str:
    if not isinstance(item, dict):
        return _clean_text(item)
    parts = [
        _clean_text(item.get("label")),
        _clean_text(item.get("title")),
        _clean_text(item.get("locator")),
        _clean_text(item.get("url")),
        _clean_text(item.get("path")),
        _clean_text(item.get("summary")),
    ]
    return ", ".join(part for part in parts if part)


def _calculation_inputs_text(calculation: dict[str, Any]) -> str:
    inputs = calculation.get("inputs")
    if isinstance(inputs, list):
        return "; ".join(_clean_text(item) for item in inputs if _clean_text(item))
    return _clean_text(inputs)


def _claim_ref_text(ref: Any) -> str:
    if not isinstance(ref, dict):
        return _clean_text(ref)
    claim_key = _clean_text(ref.get("claim_key"))
    slide_number = ref.get("slide_number")
    claim_text = _clean_text(ref.get("claim"))
    if claim_text and isinstance(slide_number, int):
        return f'Slide {slide_number} - "{claim_text}"'
    if claim_key:
        return f"claim_key {claim_key}"
    return _clean_text(ref)


def _shape_claim_key(shape_name: Any) -> str:
    name = _clean_text(shape_name)
    if name.startswith(CLAIM_SHAPE_NAME_PREFIX):
        return name[len(CLAIM_SHAPE_NAME_PREFIX) :].strip()
    return ""


def _iter_text_shapes(shapes: Any) -> Any:
    for shape in shapes:
        if hasattr(shape, "shapes"):
            yield from _iter_text_shapes(shape.shapes)
            continue
        if hasattr(shape, "text") and _clean_text(shape.text):
            yield shape


def extract_current_deck_snapshot(pptx_path: Path) -> dict[str, Any]:
    """Extract visible text plus invisible claim keys from a current PPTX.

    This is deterministic because it only reads the PPTX text layer and shape
    names. It does not decide whether text is semantically a claim.
    """

    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Install python-pptx to check a current PPTX.") from exc

    presentation = Presentation(str(pptx_path))
    slides: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts: list[dict[str, Any]] = []
        for shape in _iter_text_shapes(slide.shapes):
            item = {"text": _clean_text(shape.text)}
            claim_key = _shape_claim_key(getattr(shape, "name", ""))
            if claim_key:
                item["claim_key"] = claim_key
            texts.append(item)
        slides.append({"slide_number": slide_number, "texts": texts})
    return {"deck": str(pptx_path), "slides": slides}


def _current_items(current_snapshot: dict[str, Any]) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for slide in current_snapshot.get("slides", []):
        if not isinstance(slide, dict):
            continue
        slide_number = slide.get("slide_number")
        if not isinstance(slide_number, int):
            continue
        for text_item in _non_empty_list(slide.get("texts")):
            if isinstance(text_item, dict):
                text = _clean_text(text_item.get("text"))
                claim_key = _clean_text(text_item.get("claim_key"))
            else:
                text = _clean_text(text_item)
                claim_key = ""
            if text:
                item = {"slide_number": slide_number, "text": text}
                if claim_key:
                    item["claim_key"] = claim_key
                items.append(item)
        for claim in _non_empty_list(slide.get("claims")):
            if not isinstance(claim, dict):
                continue
            text = _clean_text(claim.get("claim"))
            claim_key = _clean_text(claim.get("claim_key"))
            if text:
                item = {"slide_number": slide_number, "text": text}
                if claim_key:
                    item["claim_key"] = claim_key
                items.append(item)
    return items


def _claim_identity(claim: dict[str, Any]) -> tuple[int, str]:
    return int(claim["_slide_number"]), _clean_text(claim["claim"])


def _text_contains_claim(current_text: str, claim_text: str) -> bool:
    return current_text == claim_text or claim_text in current_text


def _find_exact_text_matches(
    current_items: list[dict[str, Any]], claim_text: str
) -> list[dict[str, Any]]:
    return [
        item
        for item in current_items
        if _text_contains_claim(_clean_text(item.get("text")), claim_text)
    ]


def _looks_like_untracked_claim_text(text: str) -> bool:
    if not text:
        return False
    if text.isdigit():
        return False
    if "informazioni confidenziali preparate" in text.lower():
        return False
    words = text.split()
    if len(words) < 6:
        return False
    if text.upper() == text and len(words) <= 8:
        return False
    return any(marker in text for marker in (".", ":", ";", "?"))


def _drift_status_for_claim(
    claim: dict[str, Any],
    current_items: list[dict[str, Any]],
    current_by_key: dict[str, list[dict[str, Any]]],
) -> dict[str, Any]:
    slide_number, claim_text = _claim_identity(claim)
    claim_key = _clean_text(claim.get("claim_key"))
    if claim_key:
        keyed_items = current_by_key.get(claim_key, [])
        if len(keyed_items) > 1:
            return {
                "status": "duplicate-current-key",
                "slide_number": slide_number,
                "claim": claim_text,
                "detail": f"current PPTX contains {len(keyed_items)} shapes named {CLAIM_SHAPE_NAME_PREFIX}{claim_key}",
            }
        if keyed_items:
            item = keyed_items[0]
            current_text = _clean_text(item.get("text"))
            current_slide = int(item["slide_number"])
            if _text_contains_claim(current_text, claim_text):
                status = "unchanged" if current_slide == slide_number else "moved"
                return {
                    "status": status,
                    "slide_number": slide_number,
                    "claim": claim_text,
                    "current_slide_number": current_slide,
                }
            return {
                "status": "edited",
                "slide_number": slide_number,
                "claim": claim_text,
                "current_slide_number": current_slide,
                "current_text": current_text,
            }

    matches = _find_exact_text_matches(current_items, claim_text)
    if matches:
        same_slide = [
            item for item in matches if int(item["slide_number"]) == slide_number
        ]
        if same_slide:
            return {
                "status": "unchanged",
                "slide_number": slide_number,
                "claim": claim_text,
            }
        first = matches[0]
        return {
            "status": "moved",
            "slide_number": slide_number,
            "claim": claim_text,
            "current_slide_number": int(first["slide_number"]),
        }
    return {
        "status": "missing-or-edited",
        "slide_number": slide_number,
        "claim": claim_text,
        "detail": "claim text not found exactly in the current PPTX text layer",
    }


def compare_current_deck_snapshot(
    payload: dict[str, Any], current_snapshot: dict[str, Any]
) -> dict[str, Any]:
    """Compare generated claim text to a current deck snapshot.

    The check is deterministic and intentionally shallow: exact normalized text
    and invisible claim keys identify drift. Semantic support decisions remain
    outside this script.
    """

    errors = _validate_payload(payload)
    if errors:
        raise ValueError("\n".join(errors))
    slides = sorted(payload["slides"], key=lambda slide: slide["slide_number"])
    by_key, by_slide_claim, index_errors = _claim_indexes(slides)
    if index_errors:
        raise ValueError("\n".join(index_errors))

    items = _current_items(current_snapshot)
    current_by_key: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for item in items:
        claim_key = _clean_text(item.get("claim_key"))
        if claim_key:
            current_by_key[claim_key].append(item)

    statuses: dict[tuple[int, str], dict[str, Any]] = {}
    for slide in slides:
        for claim in slide["claims"]:
            status = _drift_status_for_claim(claim, items, current_by_key)
            statuses[_claim_identity(claim)] = status

    issues = [
        status
        for status in statuses.values()
        if status["status"] in TEXT_DRIFT_REQUIRES_REFRESH
    ]
    moved = [status for status in statuses.values() if status["status"] == "moved"]

    for slide in slides:
        current_slide_number = int(slide["slide_number"])
        for claim in slide["claims"]:
            for ref in _non_empty_list(claim.get("claim_refs")):
                target, issue = _resolve_claim_ref(
                    ref, current_slide_number, by_key, by_slide_claim
                )
                if issue or target is None:
                    continue
                target_status = statuses.get(_claim_identity(target))
                if target_status is None:
                    continue
                if target_status["status"] not in ("unchanged",):
                    issues.append(
                        {
                            "status": "reference-broken",
                            "slide_number": current_slide_number,
                            "claim": _clean_text(claim["claim"]),
                            "detail": (
                                "depends on claim now marked "
                                f"{target_status['status']}: "
                                f"Slide {target_status['slide_number']} - "
                                f"\"{target_status['claim']}\""
                            ),
                        }
                    )

    generated_claim_texts = {
        _clean_text(claim["claim"]) for slide in slides for claim in slide["claims"]
    }
    slide_titles = {_clean_text(slide.get("slide_title")) for slide in slides}
    untracked: list[dict[str, Any]] = []
    for item in items:
        if _clean_text(item.get("claim_key")):
            continue
        text = _clean_text(item.get("text"))
        if text in slide_titles:
            continue
        if any(
            _text_contains_claim(text, claim_text)
            for claim_text in generated_claim_texts
        ):
            continue
        if _looks_like_untracked_claim_text(text):
            untracked.append(
                {
                    "status": "untracked-current-text",
                    "slide_number": int(item["slide_number"]),
                    "claim": text,
                    "detail": "current text is not present in the generation-time claim snapshot",
                }
            )
    issues.extend(untracked)

    return {
        "checked_deck": _clean_text(current_snapshot.get("deck")),
        "claims_checked": len(statuses),
        "issues": issues,
        "moved": moved,
        "untracked_current_text": untracked,
    }


def _claim_indexes(
    slides: list[dict[str, Any]],
) -> tuple[dict[str, dict[str, Any]], dict[tuple[int, str], dict[str, Any]], list[str]]:
    by_key: dict[str, dict[str, Any]] = {}
    by_slide_claim: dict[tuple[int, str], dict[str, Any]] = {}
    errors: list[str] = []
    for slide in slides:
        slide_number = slide["slide_number"]
        for claim in slide["claims"]:
            claim["_slide_number"] = slide_number
            claim_key = _clean_text(claim.get("claim_key"))
            claim_text = _clean_text(claim["claim"])
            if claim_key:
                if claim_key in by_key:
                    errors.append(f"Duplicate claim_key: {claim_key}.")
                by_key[claim_key] = claim
            slide_claim_key = (slide_number, claim_text)
            if slide_claim_key in by_slide_claim:
                errors.append(
                    f'Duplicate claim text on slide {slide_number}: "{claim_text}".'
                )
            by_slide_claim[slide_claim_key] = claim
    return by_key, by_slide_claim, errors


def _resolve_claim_ref(
    ref: Any,
    current_slide_number: int,
    by_key: dict[str, dict[str, Any]],
    by_slide_claim: dict[tuple[int, str], dict[str, Any]],
) -> tuple[dict[str, Any] | None, str | None]:
    if not isinstance(ref, dict):
        return None, f"claim reference is not an object: {_claim_ref_text(ref)}"
    claim_key = _clean_text(ref.get("claim_key"))
    target: dict[str, Any] | None = None
    if claim_key:
        target = by_key.get(claim_key)
        if target is None:
            return None, f"unresolved claim_key reference: {claim_key}"
    else:
        slide_number = ref.get("slide_number")
        claim_text = _clean_text(ref.get("claim"))
        if not isinstance(slide_number, int) or not claim_text:
            return None, "claim reference requires claim_key or slide_number + claim"
        target = by_slide_claim.get((slide_number, claim_text))
        if target is None:
            return (
                None,
                f'unresolved claim reference: Slide {slide_number} - "{claim_text}"',
            )
    target_slide_number = target["_slide_number"]
    if target_slide_number >= current_slide_number:
        return (
            target,
            "claim reference must point to an earlier slide: "
            f"target slide {target_slide_number}, current slide {current_slide_number}",
        )
    return target, None


def _claim_grounding_issue(
    claim: dict[str, Any],
    by_key: dict[str, dict[str, Any]],
    by_slide_claim: dict[tuple[int, str], dict[str, Any]],
    seen: set[int] | None = None,
) -> str:
    seen = seen or set()
    claim_identity = id(claim)
    if claim_identity in seen:
        return "cyclic claim reference"
    seen.add(claim_identity)
    if _basis_type(claim) == "ungrounded":
        return "no captured source, calculation, reasoning input, or assumption"
    current_slide_number = int(claim["_slide_number"])
    for ref in _non_empty_list(claim.get("claim_refs")):
        target, issue = _resolve_claim_ref(
            ref, current_slide_number, by_key, by_slide_claim
        )
        if issue:
            return issue
        if target is None:
            return f"unresolved claim reference: {_claim_ref_text(ref)}"
        upstream_issue = _claim_grounding_issue(
            target, by_key, by_slide_claim, seen=set(seen)
        )
        if upstream_issue:
            return (
                f'depends on Slide {target["_slide_number"]} claim '
                f'"{_clean_text(target["claim"])}": {upstream_issue}'
            )
    return ""


def _validate_payload(payload: Any) -> list[str]:
    errors: list[str] = []
    if not isinstance(payload, dict):
        return ["Top-level JSON value must be an object."]
    slides = payload.get("slides")
    if not isinstance(slides, list):
        return ["Top-level 'slides' must be a list."]
    seen_slide_numbers: set[int] = set()
    for slide_index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            errors.append(f"slides[{slide_index}] must be an object.")
            continue
        slide_number = slide.get("slide_number")
        if not isinstance(slide_number, int):
            errors.append(f"slides[{slide_index}].slide_number must be an integer.")
        elif slide_number in seen_slide_numbers:
            errors.append(f"Duplicate slide_number: {slide_number}.")
        else:
            seen_slide_numbers.add(slide_number)
        claims = slide.get("claims")
        if not isinstance(claims, list):
            errors.append(f"slides[{slide_index}].claims must be a list.")
            continue
        for claim_index, claim in enumerate(claims, start=1):
            if not isinstance(claim, dict):
                errors.append(
                    f"slides[{slide_index}].claims[{claim_index}] must be an object."
                )
                continue
            if not _clean_text(claim.get("claim")):
                errors.append(
                    f"slides[{slide_index}].claims[{claim_index}].claim is required."
                )
    if not errors:
        _, _, index_errors = _claim_indexes(slides)
        errors.extend(index_errors)
    return errors


def _claim_block(claim: dict[str, Any]) -> list[str]:
    claim_text = _clean_text(claim["claim"])
    basis = _basis_type(claim)
    lines = [f'- "{claim_text}"']
    if basis == "source-backed":
        for ref in _non_empty_list(claim.get("source_refs")):
            text = _source_ref_text(ref)
            lines.append(
                f"  Source: {text or 'source reference captured without details'}"
            )
    elif basis == "calculated":
        calculation = _non_empty_dict(claim.get("calculation_ref"))
        inputs = _calculation_inputs_text(calculation)
        method = _clean_text(calculation.get("method"))
        output = _clean_text(calculation.get("output"))
        if inputs:
            lines.append(f"  Inputs: {inputs}")
        if method:
            lines.append(f"  Method: {method}")
        if output:
            lines.append(f"  Output: {output}")
    elif basis == "claim-linked":
        for ref in _non_empty_list(claim.get("claim_refs")):
            lines.append(f"  Based on claim: {_claim_ref_text(ref)}")
        reasoning = _clean_text(claim.get("reasoning"))
        if reasoning:
            lines.append(f"  Reasoning: {reasoning}")
    elif basis == "reasoned":
        inputs = [
            text
            for text in (
                _reasoning_input_text(item)
                for item in _non_empty_list(claim.get("reasoning_inputs"))
            )
            if text
        ]
        if inputs:
            lines.append(f"  Inputs: {'; '.join(inputs)}")
        reasoning = _clean_text(claim.get("reasoning"))
        if reasoning:
            lines.append(f"  Reasoning: {reasoning}")
    elif basis == "assumption":
        lines.append(f"  Assumption: {_clean_text(claim.get('assumption_basis'))}")
    else:
        lines.append(
            "  Basis: no captured source, calculation, reasoning input, or assumption"
        )
    return lines


def _render_drift_check(report: dict[str, Any]) -> list[str]:
    lines = ["## Current Deck Check", ""]
    checked_deck = _clean_text(report.get("checked_deck"))
    if checked_deck:
        lines.extend([f"Checked deck: {checked_deck}", ""])
    issues = report.get("issues", [])
    if issues:
        lines.extend(["### Claims Requiring Refresh", ""])
        for issue in issues:
            lines.append(
                f'- Slide {issue["slide_number"]}: "{_clean_text(issue["claim"])}"'
            )
            lines.append(f"  Status: {issue['status']}")
            current_slide = issue.get("current_slide_number")
            if isinstance(current_slide, int):
                lines.append(f"  Current slide: {current_slide}")
            current_text = _clean_text(issue.get("current_text"))
            if current_text:
                lines.append(f'  Current text: "{current_text}"')
            detail = _clean_text(issue.get("detail"))
            if detail:
                lines.append(f"  Detail: {detail}")
            lines.append("")
    else:
        lines.extend(
            ["No edited, missing, untracked, or broken-reference claims detected.", ""]
        )

    moved = report.get("moved", [])
    if moved:
        lines.extend(["### Moved Claims", ""])
        for item in moved:
            lines.append(
                f'- Slide {item["slide_number"]}: "{_clean_text(item["claim"])}"'
            )
            lines.append(f"  Current slide: {item['current_slide_number']}")
            lines.append("")
    return lines


def render_claim_basis_map(
    payload: dict[str, Any], current_snapshot: dict[str, Any] | None = None
) -> str:
    errors = _validate_payload(payload)
    if errors:
        raise ValueError("\n".join(errors))

    deck_name = _clean_text(payload.get("deck"))
    slides = sorted(payload["slides"], key=lambda slide: slide["slide_number"])
    by_key, by_slide_claim, index_errors = _claim_indexes(slides)
    if index_errors:
        raise ValueError("\n".join(index_errors))
    ungrounded: list[tuple[int, str, str]] = []
    for slide in slides:
        for claim in slide["claims"]:
            issue = _claim_grounding_issue(claim, by_key, by_slide_claim)
            if issue:
                ungrounded.append(
                    (slide["slide_number"], _clean_text(claim["claim"]), issue)
                )

    lines = ["# Claim Basis Map", ""]
    if deck_name:
        lines.extend([f"Deck: {deck_name}", ""])

    lines.extend(["## Ungrounded Claims", ""])
    if ungrounded:
        for slide_number, claim_text, issue in ungrounded:
            lines.append(f'- Slide {slide_number}: "{claim_text}"')
            lines.append(f"  Basis: {issue}")
            lines.append("")
    else:
        lines.extend(["No ungrounded or unresolved claim dependencies captured.", ""])

    if current_snapshot is not None:
        lines.append("")
        lines.extend(
            _render_drift_check(
                compare_current_deck_snapshot(payload, current_snapshot)
            )
        )

    for slide in slides:
        slide_number = slide["slide_number"]
        slide_title = _clean_text(slide.get("slide_title"))
        heading = f"## Slide {slide_number}"
        if slide_title:
            heading = f"{heading} - {slide_title}"
        lines.extend([heading, ""])

        claims_by_basis: dict[str, list[dict[str, Any]]] = {
            basis: [] for basis in BASIS_ORDER
        }
        for claim in slide["claims"]:
            claims_by_basis[_basis_type(claim)].append(claim)

        emitted = False
        for basis in BASIS_ORDER:
            claims = claims_by_basis[basis]
            if not claims:
                continue
            emitted = True
            lines.extend([f"### {basis.title()}", ""])
            for claim in claims:
                lines.extend(_claim_block(claim))
                lines.append("")
        if not emitted:
            lines.extend(["No claims captured for this slide.", ""])

    return "\n".join(lines).rstrip() + "\n"


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Render deck.claims.md from structured deck.claims.json."
    )
    parser.add_argument("claims_json", type=Path)
    parser.add_argument("--output", "-o", type=Path)
    parser.add_argument(
        "--current-pptx",
        type=Path,
        help="Optional current PPTX to compare against the generation-time claims JSON.",
    )
    parser.add_argument(
        "--current-claims-json",
        type=Path,
        help="Optional current claim/text snapshot JSON to compare against.",
    )
    parser.add_argument(
        "--snapshot-output",
        type=Path,
        help="Optional path for the extracted current PPTX text snapshot.",
    )
    return parser.parse_args()


def main() -> int:
    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
    args = _parse_args()
    if args.current_pptx and args.current_claims_json:
        raise ValueError("Use only one of --current-pptx or --current-claims-json.")
    payload = json.loads(args.claims_json.read_text(encoding="utf-8"))
    current_snapshot = None
    if args.current_pptx:
        current_snapshot = extract_current_deck_snapshot(args.current_pptx)
        if args.snapshot_output:
            args.snapshot_output.write_text(
                json.dumps(current_snapshot, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            LOGGER.info("Wrote %s", args.snapshot_output)
    elif args.current_claims_json:
        current_snapshot = json.loads(
            args.current_claims_json.read_text(encoding="utf-8")
        )
    markdown = render_claim_basis_map(payload, current_snapshot=current_snapshot)
    output = args.output or args.claims_json.with_suffix(".md")
    output.write_text(markdown, encoding="utf-8")
    LOGGER.info("Wrote %s", output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
