from __future__ import annotations

import json
import shutil
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from src.slides.pptx_template_manifest import (
    DECK_PPTX_TEMPLATE_FILENAME,
    build_pptx_template_manifest,
    deck_pptx_template_manifest_path,
    ensure_deck_pptx_template_manifest,
    load_deck_pptx_template_manifest,
)
from src.slides.semantic_pptx import (
    SLIDES_PPTX_SPEC_FILENAME,
    render_slides_pptx_from_template,
)


def test_build_pptx_template_manifest_extracts_supported_roles() -> None:
    template_path = Path("src/slides/pptx_templates/uniform.pptx")

    manifest = build_pptx_template_manifest(template_path)

    title_only = manifest.layout_for_role("title_only")
    title_body = manifest.layout_for_role("title_body")
    text_visual = manifest.layout_for_role("text_visual")

    assert title_only is not None
    assert title_only.layout_name == "Title Only"
    assert title_body is not None
    assert title_body.layout_name == "Title and Content"
    assert text_visual is not None
    assert text_visual.layout_name == "Content with Caption"
    assert text_visual.text_placeholder_idx is not None
    assert text_visual.visual_placeholder_idx is not None


def test_ensure_deck_pptx_template_manifest_persists_manifest(tmp_path: Path) -> None:
    deck_path = tmp_path / "deckTemplate"
    deck_path.mkdir(parents=True, exist_ok=True)
    shutil.copy(
        Path("src/slides/pptx_templates/uniform.pptx"),
        deck_path / DECK_PPTX_TEMPLATE_FILENAME,
    )

    manifest = ensure_deck_pptx_template_manifest(deck_path)

    assert manifest is not None
    assert deck_pptx_template_manifest_path(deck_path).exists()


def test_load_deck_pptx_template_manifest_preserves_zero_title_placeholder_idx(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckTemplateLoad"
    deck_path.mkdir(parents=True, exist_ok=True)
    shutil.copy(
        Path("src/slides/pptx_templates/uniform.pptx"),
        deck_path / DECK_PPTX_TEMPLATE_FILENAME,
    )

    manifest = ensure_deck_pptx_template_manifest(deck_path)
    loaded_manifest = load_deck_pptx_template_manifest(deck_path)

    assert manifest is not None
    assert loaded_manifest is not None
    title_body = loaded_manifest.layout_for_role("title_body")
    text_visual = loaded_manifest.layout_for_role("text_visual")
    assert title_body is not None
    assert text_visual is not None
    assert title_body.title_placeholder_idx == 0
    assert text_visual.title_placeholder_idx == 0


def test_build_pptx_template_manifest_extracts_text_visual_prototype(
    tmp_path: Path,
) -> None:
    template_path = tmp_path / "prototype_template.pptx"
    _write_text_visual_prototype_template(template_path)

    manifest = build_pptx_template_manifest(template_path)

    prototype = manifest.prototype_for_role("text_visual")

    assert prototype is not None
    assert prototype.slide_index >= 0
    assert prototype.title_shape_index is not None
    assert prototype.body_label_shape_index is not None
    assert prototype.body_shape_index is not None
    assert prototype.visual_width > 0
    assert prototype.visual_height > 0


def test_render_slides_pptx_from_custom_template_uses_template_layouts(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckCustomTemplate"
    deck_path.mkdir(parents=True, exist_ok=True)
    shutil.copy(
        Path("src/slides/pptx_templates/uniform.pptx"),
        deck_path / DECK_PPTX_TEMPLATE_FILENAME,
    )
    image_path = deck_path / "pptx_assets" / "visual.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (640, 360), "white").save(image_path)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "bullets_visual",
                "layoutVariant": "bullets_visual_right",
                "densityHint": "light",
                "title": "Strategy options",
                "body": "",
                "bullets": ["Option A", "Option B"],
                "visualPath": "pptx_assets/visual.png",
                "visualType": "figure",
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide = presentation.slides[0]
    text_shapes = [
        shape for shape in slide.shapes if hasattr(shape, "text") and "Option A" in shape.text
    ]
    title_shapes = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "text") and "Strategy options" in shape.text
    ]
    picture_shapes = [shape for shape in slide.shapes if int(shape.shape_type) == 13]

    assert len(title_shapes) == 1
    assert len(text_shapes) == 1
    assert len(picture_shapes) == 1
    title_shape = title_shapes[0]
    text_shape = text_shapes[0]
    picture_shape = picture_shapes[0]
    assert int(title_shape.left) == 457200
    assert int(text_shape.left) == 457200
    assert int(picture_shape.left) >= 3575050
    assert int(text_shape.left + text_shape.width) <= int(picture_shape.left)


def test_render_slides_pptx_from_custom_template_uses_cached_manifest_for_title(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckCustomTemplateCachedManifest"
    deck_path.mkdir(parents=True, exist_ok=True)
    shutil.copy(
        Path("src/slides/pptx_templates/uniform.pptx"),
        deck_path / DECK_PPTX_TEMPLATE_FILENAME,
    )
    ensure_deck_pptx_template_manifest(deck_path)
    image_path = deck_path / "pptx_assets" / "visual.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (640, 360), "white").save(image_path)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "bullets_visual",
                "layoutVariant": "bullets_visual_right",
                "densityHint": "light",
                "title": "Strategy options",
                "body": "",
                "bullets": ["Option A", "Option B"],
                "visualPath": "pptx_assets/visual.png",
                "visualType": "figure",
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide = next(iter(presentation.slides))
    texts = [
        getattr(shape, "text", "").strip()
        for shape in slide.shapes
        if hasattr(shape, "text") and getattr(shape, "text", "").strip()
    ]

    assert "Strategy options" in texts


def test_render_slides_pptx_from_custom_template_discards_template_sample_slides(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckCustomTemplateWithSamples"
    deck_path.mkdir(parents=True, exist_ok=True)
    template_path = deck_path / DECK_PPTX_TEMPLATE_FILENAME
    template = Presentation(str(Path("src/slides/pptx_templates/uniform.pptx")))
    sample_slide = template.slides.add_slide(template.slide_layouts[1])
    sample_slide.shapes.title.text = "TEMPLATE SAMPLE SLIDE"
    content_placeholder = sample_slide.placeholders[1]
    content_placeholder.text = "This text must not survive in exported decks."
    template.save(str(template_path))

    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "text_only",
                "layoutVariant": "text_full_width",
                "densityHint": "light",
                "title": "Generated slide",
                "body": "Only this generated slide should remain.",
                "bullets": [],
                "visualPath": "",
                "visualType": "",
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)

    slides = list(presentation.slides)
    assert len(slides) == 1
    texts = [
        getattr(shape, "text", "").strip()
        for shape in slides[0].shapes
        if hasattr(shape, "text") and getattr(shape, "text", "").strip()
    ]
    assert "Generated slide" in texts
    assert "TEMPLATE SAMPLE SLIDE" not in texts


def test_render_slides_pptx_from_custom_template_uses_text_visual_prototype(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckPrototypeTemplate"
    deck_path.mkdir(parents=True, exist_ok=True)
    template_path = deck_path / DECK_PPTX_TEMPLATE_FILENAME
    _write_text_visual_prototype_template(template_path)
    image_path = deck_path / "pptx_assets" / "visual.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (800, 500), "white").save(image_path)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "bullets_visual",
                "layoutVariant": "bullets_visual_right",
                "densityHint": "light",
                "title": "Prototype exhibit page",
                "body": "",
                "bullets": ["Point A", "Point B"],
                "visualPath": "pptx_assets/visual.png",
                "visualType": "figure",
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide = presentation.slides[0]
    texts = [
        getattr(shape, "text", "").strip()
        for shape in slide.shapes
        if hasattr(shape, "text") and getattr(shape, "text", "").strip()
    ]
    picture_shapes = [shape for shape in slide.shapes if int(shape.shape_type) == 13]

    assert len(picture_shapes) == 1
    assert "Prototype exhibit page" in texts
    assert "Key takeaways" in texts
    assert any("Point A" in text for text in texts)
    assert "CLIENT WORDMARK" not in texts
    assert "CONFIDENTIAL | CLIENT NAME" not in texts
    assert "BACKUP EXHIBIT / LARGE TABLE / IMAGE" not in texts
    assert "OPTIONAL MINI-EXHIBIT" not in texts


def _write_text_visual_prototype_template(template_path: Path) -> None:
    presentation = Presentation(str(Path("src/slides/pptx_templates/uniform.pptx")))
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_test_textbox(
        slide,
        text="Appendix exhibit / backup page",
        left=0.60,
        top=0.20,
        width=8.00,
        height=0.35,
    )
    _add_test_textbox(
        slide,
        text="Use appendix pages for definitions, backup cuts, sensitivity checks, or raw tables.",
        left=0.60,
        top=0.58,
        width=9.40,
        height=0.24,
    )
    _add_test_textbox(
        slide,
        text="CLIENT WORDMARK",
        left=10.80,
        top=0.24,
        width=1.70,
        height=0.22,
    )
    _add_test_textbox(
        slide,
        text="CONFIDENTIAL | CLIENT NAME",
        left=0.60,
        top=6.86,
        width=3.00,
        height=0.18,
    )
    _add_test_textbox(
        slide,
        text="19",
        left=12.20,
        top=6.82,
        width=0.30,
        height=0.18,
    )
    _add_test_textbox(
        slide,
        text="BACKUP EXHIBIT / LARGE TABLE / IMAGE",
        left=0.92,
        top=3.55,
        width=6.45,
        height=0.28,
    )
    _add_test_textbox(
        slide,
        text="Supporting notes",
        left=9.08,
        top=1.58,
        width=2.80,
        height=0.22,
    )
    _add_test_textbox(
        slide,
        text="• Use this panel for definitions, caveats, or assumptions.",
        left=9.08,
        top=1.98,
        width=2.75,
        height=1.30,
    )
    _add_test_textbox(
        slide,
        text="OPTIONAL MINI-EXHIBIT",
        left=9.28,
        top=4.32,
        width=2.30,
        height=0.28,
    )
    presentation.save(str(template_path))


def _add_test_textbox(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
) -> None:
    textbox = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    textbox.text_frame.text = text
