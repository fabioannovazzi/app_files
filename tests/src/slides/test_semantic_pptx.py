from __future__ import annotations

import json
import zipfile
from pathlib import Path

import pytest
from PIL import Image, ImageDraw

import src.slides.semantic_pptx as semantic_pptx
from src.slides.models import Deck, Slide
from src.slides.semantic_pptx import (
    SLIDES_PPTX_SPEC_FILENAME,
    SlidesPptxSlide,
    build_slides_pptx_spec,
    build_slides_pptx_spec_from_report_payload,
    render_slides_pptx_from_template,
    write_slides_pptx_spec,
)


def test_build_slides_pptx_spec_uses_analysis_blocks_and_writes_visual_crop(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckSemantic",
        prompt_style="editorial",
        slides=[
            Slide(
                id="slide0.html",
                title_html="Ignored",
                body_html="<img src='assets/source.png' alt='source' />",
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    asset_path = deck_path / "assets" / "source.png"
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (1200, 700), "white")
    ImageDraw.Draw(image).rectangle((640, 120, 1120, 500), fill=(20, 20, 20))
    image.save(asset_path)
    analysis_payload = {
        "deckId": deck.deck_id,
        "lang": "ita",
        "generatedAt": "2026-03-17T10:00:00+00:00",
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "assetPath": "assets/source.png",
                "titleText": "Market entry strategy",
                "bulletTexts": ["First priority", "Second priority"],
                "figureRegions": [{"x": 640, "y": 120, "w": 480, "h": 380}],
                "blocks": [
                    {
                        "blockId": "block-0",
                        "type": "figure",
                        "text": "",
                        "items": [],
                        "bbox": {"x": 640, "y": 120, "w": 480, "h": 380},
                    },
                    {
                        "blockId": "block-1",
                        "type": "title",
                        "text": "Market\nentry\nstrategy",
                        "items": [],
                        "bbox": {"x": 60, "y": 40, "w": 420, "h": 180},
                    },
                    {
                        "blockId": "block-2",
                        "type": "list",
                        "text": "First priority",
                        "items": ["First priority"],
                        "bbox": {"x": 60, "y": 260, "w": 420, "h": 100},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    assert spec.template_key == "uniform"
    assert spec.prompt_style == "uniform"
    assert len(spec.slides) == 1
    slide = spec.slides[0]
    assert slide.kind == "bullets_visual"
    assert slide.layout_variant == "bullets_visual_right"
    assert slide.density_hint == "light"
    assert slide.title == "Market entry strategy"
    assert slide.bullets == ["First priority", "Second priority"]
    assert slide.visual_path.startswith("pptx_assets/")
    assert (deck_path / slide.visual_path).exists()


def test_build_slides_pptx_spec_preserves_section_header_slides(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckSections",
        prompt_style="uniform",
        slides=[
            Slide(
                id="slide0.html",
                title_html="",
                body_html=(
                    '<section class="section-header">'
                    '<ol class="section-header__sections">'
                    '<li class="section-header__section is-current">'
                    '<span class="section-header__section-label">Data Pipeline</span>'
                    '<ul class="section-header__subsections">'
                    '<li class="section-header__subsection is-current">Batch Processing</li>'
                    '<li class="section-header__subsection">Stream Processing</li>'
                    "</ul>"
                    "</li>"
                    '<li class="section-header__section">'
                    '<span class="section-header__section-label">Quality Controls</span>'
                    "</li>"
                    "</ol>"
                    "</section>"
                ),
                kind="sectionHeader",
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=None)

    assert len(spec.slides) == 1
    slide = spec.slides[0]
    assert slide.kind == "section_header"
    assert slide.layout_variant == "section_header_agenda"


def test_build_slides_pptx_spec_resolves_pdf_asset_url_for_image_slide(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckPdf",
        prompt_style="uniform",
        slides=[
            Slide(
                id="slide0.html",
                title_html="",
                body_html=(
                    '<img src="/slides/deck/deckPdf/assets/pages/page-1.png" '
                    'alt="Slide image 1" />'
                ),
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    asset_path = deck_path / "assets" / "pages" / "page-1.png"
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (1280, 720), "white").save(asset_path)

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=None)
    write_slides_pptx_spec(deck_path, spec)
    buffer = render_slides_pptx_from_template(deck_path)

    archive = zipfile.ZipFile(buffer)
    media_entries = [
        name for name in archive.namelist() if name.startswith("ppt/media/")
    ]
    assert spec.slides[0].kind == "visual_only"
    assert spec.slides[0].layout_variant == "visual_full_width"
    assert spec.slides[0].visual_path == "assets/pages/page-1.png"
    assert media_entries


def test_build_slides_pptx_spec_from_report_payload_selects_generic_layouts(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckLaunchSpec"
    deck_path.mkdir(parents=True, exist_ok=True)
    image_path = deck_path / "pptx_assets" / "launch-grid.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (640, 360), "white").save(image_path)
    payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0",
                "title": "Lipstick Launch Signal Test",
                "body": [
                    "Testing whether launch attributes are signal, not noise.",
                    "This is a constrained read on audited Ulta launches.",
                ],
                "footerText": "Ulta launch cohort | March 2026",
            },
            {
                "slideId": "slide1",
                "title": "Included launches",
                "bullets": [
                    "Hydrating launches over-index versus the older base.",
                    "Long-wear remains common.",
                ],
                "visualPath": "pptx_assets/launch-grid.png",
                "visualType": "figure",
            },
            {
                "slideId": "slide2",
                "title": "What survived the audit",
                "body": "The launch cohort still points to a comfort-led signal.",
                "comparisonColumns": [
                    {
                        "title": "What the data supports",
                        "items": ["Hydrating over-indexes", "Long-wear over-indexes"],
                    },
                    {
                        "title": "What the data does not support",
                        "items": [
                            "No dramatic new era",
                            "Refillable does not survive audit",
                        ],
                    },
                ],
                "calloutTitle": "Bottom line",
                "calloutBody": "The surviving signal is modest but defensible.",
            },
        ],
    }

    spec = build_slides_pptx_spec_from_report_payload(payload, deck_path=deck_path)

    assert spec.template_key == "uniform"
    assert spec.prompt_style == "uniform"
    assert [slide.layout_variant for slide in spec.slides] == [
        "cover_with_footer",
        "bullets_visual_right",
        "comparison_columns",
    ]
    assert spec.slides[0].footer_text == "Ulta launch cohort | March 2026"
    assert spec.slides[1].visual_path == "pptx_assets/launch-grid.png"
    assert (
        spec.slides[2].callout_body == "The surviving signal is modest but defensible."
    )


def test_title_layout_for_slide_uses_repair_scale() -> None:
    slide = SlidesPptxSlide(
        slide_id="slide0.html",
        kind="text_only",
        layout_variant="text_full_width",
        density_hint="medium",
        title="A deliberately long title that would otherwise remain too large on the slide",
        repair_hints={"title_scale": 0.8},
    )

    layout = semantic_pptx._title_layout_for_slide(26.0, slide)

    assert layout["font_size"] < 26.0
    assert layout["height"] > 0.86


def test_implication_banner_box_uses_repair_scale() -> None:
    slide = SlidesPptxSlide(
        slide_id="slide0.html",
        kind="text_visual",
        layout_variant="text_visual_right",
        title="Title",
        implication="IMPLICATION: Keep early pilots in the test cohort.",
        repair_hints={"banner_scale": 1.2},
    )

    banner_box = semantic_pptx._implication_banner_box(slide)

    assert banner_box is not None
    assert banner_box["height"] > 0.84


def test_build_slides_pptx_spec_from_report_payload_renders_launch_report(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckLaunchRender"
    deck_path.mkdir(parents=True, exist_ok=True)
    image_path = deck_path / "pptx_assets" / "launches.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (640, 360), "white").save(image_path)
    payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0",
                "title": "Launch thesis",
                "body": "Launches are the cleanest place to test whether attribute signal survives audit.",
                "footerText": "Ulta | new arrivals",
            },
            {
                "slideId": "slide1",
                "title": "Included launches",
                "bullets": [
                    "Hydrating launches over-index.",
                    "Long-wear remains common.",
                ],
                "visualPath": "pptx_assets/launches.png",
                "visualType": "figure",
            },
        ],
    }

    spec = build_slides_pptx_spec_from_report_payload(payload, deck_path=deck_path)
    write_slides_pptx_spec(deck_path, spec)

    buffer = render_slides_pptx_from_template(deck_path)

    presentation = Presentation(buffer)
    slide_texts = [
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and str(shape.text or "").strip()
    ]

    assert len(presentation.slides) == 2
    assert any("Launch thesis" in text for text in slide_texts)
    assert any("Ulta | new arrivals" in text for text in slide_texts)
    assert any("Included launches" in text for text in slide_texts)


def test_build_slides_pptx_spec_from_report_payload_rejects_invalid_ast(
    tmp_path: Path,
) -> None:
    payload = {
        "slides": [
            {
                "title": "Broken launch slide",
                "layoutVariant": "freeform_masterpiece",
            }
        ]
    }

    with pytest.raises(ValueError, match="unsupported layoutVariant"):
        build_slides_pptx_spec_from_report_payload(payload, deck_path=tmp_path)


def test_build_slides_pptx_spec_extracts_div_only_cover_metadata(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckCoverMetadata",
        prompt_style="uniform",
        slides=[
            Slide(
                id="slide0.html",
                title_html="<h1>Iniziativa Sintetica - Analisi</h1>",
                body_html=(
                    '<div class="deck-title__meta">'
                    '<div class="deck-title__subtitle">Example Advisory</div>'
                    '<div class="deck-title__date">Gennaio 2026</div>'
                    "</div>"
                ),
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=None)

    assert len(spec.slides) == 1
    slide = spec.slides[0]
    assert slide.kind == "text_only"
    assert slide.layout_variant == "text_full_width"
    assert slide.title == "Iniziativa Sintetica - Analisi"
    assert slide.body == "Example Advisory\n\nGennaio 2026"


def test_build_slides_pptx_spec_sorts_analysis_bullets_by_bbox_and_strips_markers(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckBulletOrder",
        prompt_style="uniform",
        slides=[Slide(id="slide0.html", title_html="", body_html="")],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)
    analysis_payload = {
        "deckId": deck.deck_id,
        "lang": "ita",
        "generatedAt": "2026-03-19T10:00:00+00:00",
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "titleText": "Modalità di elaborazione per una pipeline dati",
                "bulletTexts": [
                    "Batch: Elaborazioni pianificate",
                    "• Streaming: Elaborazioni continue",
                    "Interactive: Query esplorative",
                    "Archive: Conservazione a lungo termine",
                ],
                "blocks": [
                    {
                        "blockId": "block-0",
                        "type": "list",
                        "text": "Batch: Elaborazioni pianificate",
                        "bbox": {"x": 90.0, "y": 480.0, "w": 200.0, "h": 40.0},
                    },
                    {
                        "blockId": "block-1",
                        "type": "list",
                        "text": "• Streaming: Elaborazioni continue",
                        "bbox": {"x": 90.0, "y": 230.0, "w": 200.0, "h": 40.0},
                    },
                    {
                        "blockId": "block-2",
                        "type": "list",
                        "text": "Interactive: Query esplorative",
                        "bbox": {"x": 90.0, "y": 330.0, "w": 200.0, "h": 40.0},
                    },
                    {
                        "blockId": "block-3",
                        "type": "list",
                        "text": "Archive: Conservazione a lungo termine",
                        "bbox": {"x": 90.0, "y": 580.0, "w": 200.0, "h": 40.0},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.bullets == [
        "Streaming: Elaborazioni continue",
        "Interactive: Query esplorative",
        "Batch: Elaborazioni pianificate",
        "Archive: Conservazione a lungo termine",
    ]


def test_build_slides_pptx_spec_groups_exhibit_blocks_into_one_visual_crop(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckGroupedExhibit",
        prompt_style="uniform",
        slides=[
            Slide(
                id="slide0.html",
                title_html="",
                body_html="<img src='assets/source.png' alt='source' />",
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    asset_path = deck_path / "assets" / "source.png"
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (1200, 700), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((640, 100, 1050, 520), outline="black", width=4)
    draw.text((700, 160), "99.9%", fill="black")
    draw.text((690, 450), "SLA", fill="black")
    image.save(asset_path)

    analysis_payload = {
        "deckId": deck.deck_id,
        "lang": "ita",
        "generatedAt": "2026-03-19T10:00:00+00:00",
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "assetPath": "assets/source.png",
                "titleText": "Il servizio espone un livello operativo verificabile",
                "blocks": [
                    {
                        "blockId": "bullet-0",
                        "type": "bullet_item",
                        "text": "Disponibilità target 99,9%",
                        "listLevel": 0,
                        "bbox": {"x": 60.0, "y": 180.0, "w": 420.0, "h": 48.0},
                    },
                    {
                        "blockId": "figure-0",
                        "type": "figure",
                        "groupId": "exhibit-1",
                        "renderMode": "group_as_image",
                        "bbox": {"x": 640.0, "y": 100.0, "w": 410.0, "h": 420.0},
                    },
                    {
                        "blockId": "metric-0",
                        "type": "metric",
                        "groupId": "exhibit-1",
                        "renderMode": "group_as_image",
                        "text": "99.9%",
                        "bbox": {"x": 690.0, "y": 140.0, "w": 180.0, "h": 120.0},
                    },
                    {
                        "blockId": "label-0",
                        "type": "exhibit_label",
                        "groupId": "exhibit-1",
                        "renderMode": "group_as_image",
                        "text": "SLA",
                        "bbox": {"x": 700.0, "y": 430.0, "w": 120.0, "h": 40.0},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.layout_variant == "bullets_visual_right"
    assert slide.visual_path.startswith("pptx_assets/")
    assert (deck_path / slide.visual_path).exists()
    assert slide.bullets == ["Disponibilità target 99,9%"]


def test_build_slides_pptx_spec_preserves_bullet_levels_from_analysis(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckNestedBullets",
        prompt_style="uniform",
        slides=[Slide(id="slide0.html", title_html="", body_html="")],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)
    analysis_payload = {
        "deckId": deck.deck_id,
        "lang": "ita",
        "generatedAt": "2026-03-19T10:00:00+00:00",
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "titleText": "Modalità operative e relativi controlli",
                "blocks": [
                    {
                        "blockId": "bullet-0",
                        "type": "bullet_item",
                        "text": "Pipeline Standard",
                        "listLevel": 0,
                        "bbox": {"x": 80.0, "y": 180.0, "w": 260.0, "h": 40.0},
                    },
                    {
                        "blockId": "bullet-1",
                        "type": "bullet_item",
                        "text": "Batch: Caricamenti pianificati, validazione, tracciabilità.",
                        "listLevel": 1,
                        "bbox": {"x": 110.0, "y": 230.0, "w": 520.0, "h": 40.0},
                    },
                    {
                        "blockId": "bullet-2",
                        "type": "bullet_item",
                        "text": "Streaming: Eventi continui, checkpoint, bassa latenza.",
                        "listLevel": 1,
                        "bbox": {"x": 110.0, "y": 280.0, "w": 520.0, "h": 40.0},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.bullets == [
        "Pipeline Standard",
        "\tBatch: Caricamenti pianificati, validazione, tracciabilità.",
        "\tStreaming: Eventi continui, checkpoint, bassa latenza.",
    ]


def test_render_slides_pptx_from_template_creates_native_presentation(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckRender"
    deck_path.mkdir(parents=True, exist_ok=True)
    image_path = deck_path / "pptx_assets" / "visual.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (640, 360), "white").save(image_path)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "bullets_visual",
                "layoutVariant": "bullets_visual_right",
                "title": "Strategy options",
                "body": "",
                "bullets": ["Option A", "Option B"],
                "visualPath": "pptx_assets/visual.png",
                "visualType": "figure",
            },
            {
                "slideId": "slide1.html",
                "kind": "text_only",
                "layoutVariant": "text_full_width",
                "title": "Conclusion",
                "body": "Use the template-backed export first.",
                "bullets": [],
                "visualPath": "",
                "visualType": "",
            },
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)

    archive = zipfile.ZipFile(buffer)
    slide_entries = [
        name
        for name in archive.namelist()
        if name.startswith("ppt/slides/slide") and name.endswith(".xml")
    ]
    media_entries = [
        name for name in archive.namelist() if name.startswith("ppt/media/")
    ]
    assert len(slide_entries) == 2
    assert media_entries


def test_render_slides_pptx_section_header_uses_section_content(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckSectionRender"
    deck_path.mkdir(parents=True, exist_ok=True)
    (deck_path / "slide0.html").write_text(
        "<html><body><div class='slide-container'>"
        "<h1 class='slide-title' data-role='title'></h1>"
        "<div class='slide-body'>"
        "<section class='section-header'>"
        "<ol class='section-header__sections'>"
        "<li class='section-header__section is-current'>"
        "<span class='section-header__section-label'>Data Pipeline</span>"
        "<ul class='section-header__subsections'>"
        "<li class='section-header__subsection is-current'>Batch Processing</li>"
        "<li class='section-header__subsection'>Stream Processing</li>"
        "</ul>"
        "</li>"
        "<li class='section-header__section'>"
        "<span class='section-header__section-label'>Quality Controls</span>"
        "</li>"
        "<li class='section-header__section'>"
        "<span class='section-header__section-label'>Monitoring</span>"
        "</li>"
        "</ol>"
        "</section>"
        "</div></div></body></html>",
        encoding="utf-8",
    )
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "section_header",
                "layoutVariant": "section_header_agenda",
                "title": "Data Pipeline",
                "body": "",
                "bullets": [],
                "visualPath": "",
                "visualType": "",
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)

    presentation = Presentation(buffer)
    slide = presentation.slides[0]
    texts = [
        shape.text
        for shape in slide.shapes
        if hasattr(shape, "text") and str(shape.text or "").strip()
    ]

    assert "Data Pipeline" in texts
    assert "Batch Processing" in texts
    assert "Stream Processing" in texts
    assert "Quality Controls" in texts
    assert "Monitoring" in texts
    assert all("Slide 1" not in text for text in texts)


def test_render_slides_pptx_right_variant_uses_disjoint_text_and_visual_columns(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckLayoutCheck"
    deck_path.mkdir(parents=True, exist_ok=True)
    image_path = deck_path / "pptx_assets" / "visual.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (640, 360), "white").save(image_path)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "bullets_visual",
                "layoutVariant": "bullets_visual_right",
                "title": "Strategy options",
                "body": "",
                "bullets": ["Option A", "Option B", "Option C"],
                "visualPath": "pptx_assets/visual.png",
                "visualType": "figure",
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide = presentation.slides[0]
    text_shapes = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "text") and "Option A" in shape.text
    ]
    picture_shapes = [shape for shape in slide.shapes if int(shape.shape_type) == 13]

    assert len(text_shapes) == 1
    assert len(picture_shapes) == 1
    text_shape = text_shapes[0]
    picture_shape = picture_shapes[0]
    assert int(text_shape.left + text_shape.width) <= int(picture_shape.left)


def test_render_slides_pptx_bottom_variant_uses_disjoint_vertical_regions(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckBottomLayoutCheck"
    deck_path.mkdir(parents=True, exist_ok=True)
    image_path = deck_path / "pptx_assets" / "visual.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (640, 360), "white").save(image_path)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "bullets_visual",
                "layoutVariant": "bullets_visual_bottom",
                "title": "Strategy options",
                "body": "",
                "bullets": [
                    "Option A with supporting rationale",
                    "Option B with supporting rationale",
                    "Option C with supporting rationale",
                    "Option D with supporting rationale",
                ],
                "visualPath": "pptx_assets/visual.png",
                "visualType": "figure",
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide = presentation.slides[0]
    text_shapes = [
        shape
        for shape in slide.shapes
        if hasattr(shape, "text") and "Option A with supporting rationale" in shape.text
    ]
    picture_shapes = [shape for shape in slide.shapes if int(shape.shape_type) == 13]

    assert len(text_shapes) == 1
    assert len(picture_shapes) == 1
    text_shape = text_shapes[0]
    picture_shape = picture_shapes[0]
    assert int(text_shape.top + text_shape.height) <= int(picture_shape.top)


def test_render_slides_pptx_dense_bullets_use_smaller_font_size(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckDenseText"
    deck_path.mkdir(parents=True, exist_ok=True)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "bullets_only",
                "layoutVariant": "bullets_full_width",
                "densityHint": "light",
                "title": "Light slide",
                "body": "",
                "bullets": ["Short bullet", "Another short bullet"],
                "visualPath": "",
                "visualType": "",
            },
            {
                "slideId": "slide1.html",
                "kind": "bullets_only",
                "layoutVariant": "bullets_full_width",
                "densityHint": "dense",
                "title": "Dense slide",
                "body": "",
                "bullets": [
                    "A much longer bullet point with more content to force denser fitting.",
                    "Another much longer bullet point with more content to force denser fitting.",
                    "A third long bullet point to ensure the density-specific font rule is used.",
                    "A fourth long bullet point to keep the text frame under pressure.",
                ],
                "visualPath": "",
                "visualType": "",
            },
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    first_slide_text = [
        shape
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text") and "Short bullet" in shape.text
    ][0]
    second_slide_text = [
        shape
        for shape in presentation.slides[1].shapes
        if hasattr(shape, "text") and "A much longer bullet point" in shape.text
    ][0]

    first_size = first_slide_text.text_frame.paragraphs[0].runs[0].font.size.pt
    second_size = second_slide_text.text_frame.paragraphs[0].runs[0].font.size.pt

    assert second_size < first_size


def test_render_slides_pptx_strips_duplicate_leading_bullet_markers(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckBulletMarkers"
    deck_path.mkdir(parents=True, exist_ok=True)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "bullets_only",
                "layoutVariant": "bullets_full_width",
                "title": "Bullet hygiene",
                "body": "",
                "bullets": ["• Streaming: Elaborazioni continue"],
                "visualPath": "",
                "visualType": "",
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    text_shape = next(
        shape
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text")
        and "Streaming: Elaborazioni continue" in str(shape.text or "")
    )
    paragraph = text_shape.text_frame.paragraphs[0]
    paragraph_xml = paragraph._p.xml

    assert "• • Streaming" not in str(text_shape.text or "")
    assert paragraph.text == "Streaming: Elaborazioni continue"
    assert '<a:buChar char="•"/>' in paragraph_xml
    assert 'indent="-' in paragraph_xml
    assert 'marL="' in paragraph_xml


def test_render_slides_pptx_uses_native_hanging_indent_for_wrapped_bullets(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckNativeBullets"
    deck_path.mkdir(parents=True, exist_ok=True)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "bullets_only",
                "layoutVariant": "bullets_full_width",
                "title": "Definitions",
                "body": "",
                "bullets": [
                    "SLA (Service Level Agreement): The service target is 99.9% availability with a four-hour response window."
                ],
                "visualPath": "",
                "visualType": "",
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    text_shape = next(
        shape
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text")
        and "SLA (Service Level Agreement)" in str(shape.text or "")
    )
    paragraph = text_shape.text_frame.paragraphs[0]
    paragraph_xml = paragraph._p.xml

    assert paragraph.text.startswith("SLA (Service Level Agreement)")
    assert '<a:buChar char="•"/>' in paragraph_xml
    assert 'indent="-' in paragraph_xml
    assert 'marL="' in paragraph_xml


def test_write_slides_pptx_spec_persists_expected_json(tmp_path: Path) -> None:
    deck = Deck(
        deck_id="deckHtmlFallback",
        prompt_style="uniform",
        slides=[
            Slide(
                id="slide0.html",
                title_html="<h1>Title</h1>",
                body_html="<ul><li>Alpha</li><li>Beta</li></ul>",
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=None)
    output_path = write_slides_pptx_spec(deck_path, spec)

    payload = json.loads(output_path.read_text(encoding="utf-8"))
    assert output_path.name == SLIDES_PPTX_SPEC_FILENAME
    assert payload["templateKey"] == "uniform"
    assert payload["slides"][0]["kind"] == "bullets_only"
    assert payload["slides"][0]["layout_variant"] == "bullets_full_width"
    assert payload["slides"][0]["bullets"] == ["Alpha", "Beta"]


def test_build_slides_pptx_spec_selects_bottom_variant_for_bottom_figure(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckBottomFigure",
        prompt_style="uniform",
        slides=[Slide(id="slide0.html", title_html="Bottom figure", body_html="")],
    )
    deck_path = tmp_path / deck.deck_id
    asset_path = deck_path / "assets" / "source.png"
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (1200, 700), "white").save(asset_path)
    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "assetPath": "assets/source.png",
                "titleText": "Bottom figure",
                "bulletTexts": [
                    "First bullet",
                    "Second bullet",
                    "Third bullet",
                ],
                "blocks": [
                    {
                        "blockId": "b0",
                        "type": "figure",
                        "text": "",
                        "items": [],
                        "bbox": {"x": 120, "y": 360, "w": 980, "h": 300},
                    }
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    assert spec.slides[0].kind == "bullets_visual"
    assert spec.slides[0].layout_variant == "bullets_visual_bottom"
    assert spec.slides[0].density_hint == "light"


def test_build_slides_pptx_spec_prefers_right_variant_for_tall_right_side_figure(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckTallRightFigure",
        prompt_style="uniform",
        slides=[Slide(id="slide0.html", title_html="Tall right figure", body_html="")],
    )
    deck_path = tmp_path / deck.deck_id
    asset_path = deck_path / "assets" / "source.png"
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (1200, 700), "white").save(asset_path)
    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "assetPath": "assets/source.png",
                "titleText": "Tall right figure",
                "bulletTexts": [
                    "First bullet",
                    "Second bullet",
                    "Third bullet",
                ],
                "blocks": [
                    {
                        "blockId": "b0",
                        "type": "figure",
                        "text": "",
                        "items": [],
                        "bbox": {"x": 720, "y": 120, "w": 360, "h": 500},
                    }
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    assert spec.slides[0].kind == "bullets_visual"
    assert spec.slides[0].layout_variant == "bullets_visual_right"


def test_build_slides_pptx_spec_selects_table_focus_for_table_visuals(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckTableFocus",
        prompt_style="uniform",
        slides=[Slide(id="slide0.html", title_html="Table", body_html="")],
    )
    deck_path = tmp_path / deck.deck_id
    asset_path = deck_path / "assets" / "source.png"
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (1200, 700), "white").save(asset_path)
    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "assetPath": "assets/source.png",
                "titleText": "Comparison table",
                "blocks": [
                    {
                        "blockId": "b0",
                        "type": "table",
                        "text": "a b c",
                        "items": [],
                        "bbox": {"x": 80, "y": 200, "w": 1000, "h": 360},
                    }
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    assert spec.slides[0].kind == "visual_only"
    assert spec.slides[0].layout_variant == "table_focus"
    assert spec.slides[0].table_model is None


def test_render_slides_pptx_table_focus_uses_native_table_when_table_model_available(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckNativeTable"
    deck_path.mkdir(parents=True, exist_ok=True)
    image_path = deck_path / "pptx_assets" / "table.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (640, 360), "white").save(image_path)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "visual_only",
                "layoutVariant": "table_focus",
                "title": "Comparison table",
                "body": "Illustrative values",
                "bullets": [],
                "visualPath": "pptx_assets/table.png",
                "visualType": "table",
                "tableModel": {
                    "source": "deterministic_simple",
                    "confidence": 0.88,
                    "rowCount": 3,
                    "columnCount": 2,
                    "headerRows": 1,
                    "columnWidths": [0.6, 0.4],
                    "rows": [
                        {
                            "cells": [
                                {
                                    "text": "Country",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": True,
                                    "align": "center",
                                },
                                {
                                    "text": "Revenue",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": True,
                                    "align": "center",
                                },
                            ]
                        },
                        {
                            "cells": [
                                {
                                    "text": "Segment A",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "25",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "right",
                                },
                            ]
                        },
                        {
                            "cells": [
                                {
                                    "text": "Segment B",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "31",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "right",
                                },
                            ]
                        },
                    ],
                },
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide = presentation.slides[0]
    table_shapes = [shape for shape in slide.shapes if int(shape.shape_type) == 19]
    picture_shapes = [shape for shape in slide.shapes if int(shape.shape_type) == 13]

    assert len(table_shapes) == 1
    assert len(picture_shapes) == 0
    table = table_shapes[0].table
    assert table.cell(0, 0).text == "Country"
    assert table.cell(2, 1).text == "31"


def test_render_slides_pptx_table_focus_falls_back_to_image_for_suspicious_table_models(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckSuspiciousNativeTable"
    deck_path.mkdir(parents=True, exist_ok=True)
    image_path = deck_path / "pptx_assets" / "table.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (640, 360), "white").save(image_path)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "visual_only",
                "layoutVariant": "table_focus",
                "title": "Comparison table",
                "body": "Illustrative values",
                "bullets": [],
                "visualPath": "pptx_assets/table.png",
                "visualType": "table",
                "tableModel": {
                    "source": "deterministic_simple",
                    "confidence": 0.94,
                    "rowCount": 3,
                    "columnCount": 6,
                    "headerRows": 1,
                    "columnWidths": [0.08, 0.12, 0.28, 0.09, 0.16, 0.27],
                    "rows": [
                        {
                            "cells": [
                                {
                                    "text": "",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": True,
                                    "align": "center",
                                },
                                {
                                    "text": "",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": True,
                                    "align": "center",
                                },
                                {
                                    "text": "metrica operativa frammentata tra più colonne.",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": True,
                                    "align": "center",
                                },
                                {
                                    "text": "",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": True,
                                    "align": "center",
                                },
                                {
                                    "text": "",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": True,
                                    "align": "center",
                                },
                                {
                                    "text": "",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": True,
                                    "align": "center",
                                },
                            ]
                        },
                        {
                            "cells": [
                                {
                                    "text": "Segmento",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "(soglia incompleta)",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "Classificazione (livello,",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "Canale",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "Accesso Dati",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "Requisiti Operativi Chiave",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                            ]
                        },
                        {
                            "cells": [
                                {
                                    "text": "Segmento A",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "(fino a 500 unità)",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "Classe standard",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "Sì",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "Nessun requisito aggiuntivo",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "Configurazione predefinita,",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                            ]
                        },
                    ],
                },
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide = presentation.slides[0]
    table_shapes = [shape for shape in slide.shapes if int(shape.shape_type) == 19]
    picture_shapes = [shape for shape in slide.shapes if int(shape.shape_type) == 13]

    assert len(table_shapes) == 0
    assert len(picture_shapes) == 1


def test_build_slides_pptx_spec_separates_cover_footer_from_body(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckCoverFooter",
        prompt_style="uniform",
        slides=[Slide(id="slide0.html", title_html="", body_html="")],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)
    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "titleText": "Synthetic case: Service operations review",
                "blocks": [
                    {
                        "blockId": "title-0",
                        "type": "title",
                        "text": "Synthetic case: Service operations review",
                        "bbox": {"x": 80, "y": 80, "w": 900, "h": 160},
                    },
                    {
                        "blockId": "body-0",
                        "type": "body_text",
                        "text": "Evaluation of a synthetic document-processing workflow.",
                        "bbox": {"x": 90, "y": 340, "w": 920, "h": 90},
                    },
                    {
                        "blockId": "footer-0",
                        "type": "body_text",
                        "text": "Prepared for the Example Review Committee January 2026",
                        "bbox": {"x": 92, "y": 1100, "w": 420, "h": 50},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.kind == "cover_with_footer"
    assert slide.layout_variant == "cover_with_footer"
    assert slide.body == "Evaluation of a synthetic document-processing workflow."
    assert slide.footer_text == "Prepared for the Example Review Committee January 2026"


def test_build_slides_pptx_spec_preserves_comparison_columns_callout_and_implication(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckComparisonColumns",
        prompt_style="uniform",
        slides=[Slide(id="slide0.html", title_html="", body_html="")],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)
    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "titleText": "Service demand differs between standard and priority tiers.",
                "blocks": [
                    {
                        "blockId": "body-0",
                        "type": "body_text",
                        "text": "Governing thought: priority users optimize for response time.",
                        "bbox": {"x": 80, "y": 180, "w": 1000, "h": 70},
                    },
                    {
                        "blockId": "left-title",
                        "type": "group_label",
                        "text": "Standard Tier",
                        "groupKind": "comparison",
                        "bbox": {"x": 100, "y": 320, "w": 320, "h": 40},
                    },
                    {
                        "blockId": "left-b0",
                        "type": "bullet_item",
                        "text": "Designed for scheduled batch processing",
                        "groupKind": "comparison",
                        "bbox": {"x": 100, "y": 390, "w": 420, "h": 45},
                    },
                    {
                        "blockId": "left-b1",
                        "type": "body_text",
                        "text": "- Driven by predictable request volumes",
                        "bbox": {"x": 100, "y": 460, "w": 430, "h": 45},
                    },
                    {
                        "blockId": "right-title",
                        "type": "bullet_item",
                        "text": "Priority Tier",
                        "groupKind": "comparison",
                        "bbox": {"x": 700, "y": 320, "w": 320, "h": 40},
                    },
                    {
                        "blockId": "right-b0",
                        "type": "bullet_item",
                        "text": "Priority requests require defined response targets",
                        "groupKind": "comparison",
                        "bbox": {"x": 700, "y": 390, "w": 420, "h": 45},
                    },
                    {
                        "blockId": "right-b1",
                        "type": "bullet_item",
                        "text": "Deep reliance on monitoring and alerting",
                        "groupKind": "comparison",
                        "bbox": {"x": 700, "y": 460, "w": 430, "h": 45},
                    },
                    {
                        "blockId": "callout-title",
                        "type": "group_label",
                        "text": "Central Synthesis / The Gap",
                        "groupKind": "comparison",
                        "bbox": {"x": 420, "y": 760, "w": 320, "h": 40},
                    },
                    {
                        "blockId": "callout-body",
                        "type": "group_label",
                        "text": "Operators need modular services with stable APIs across critical workflows.",
                        "groupKind": "callout",
                        "bbox": {"x": 160, "y": 840, "w": 820, "h": 80},
                    },
                    {
                        "blockId": "implication-0",
                        "type": "body_text",
                        "text": "IMPLICATION: Pilot with priority workflows first.",
                        "bbox": {"x": 120, "y": 1040, "w": 900, "h": 60},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.kind == "comparison_columns"
    assert slide.layout_variant == "comparison_columns"
    assert slide.implication == "IMPLICATION: Pilot with priority workflows first."
    assert slide.callout_title == "Central Synthesis / The Gap"
    assert (
        slide.callout_body
        == "Operators need modular services with stable APIs across critical workflows."
    )
    assert slide.comparison_columns == [
        {
            "title": "Standard Tier",
            "bullets": [
                "Designed for scheduled batch processing",
                "Driven by predictable request volumes",
            ],
        },
        {
            "title": "Priority Tier",
            "bullets": [
                "Priority requests require defined response targets",
                "Deep reliance on monitoring and alerting",
            ],
        },
    ]


def test_build_slides_pptx_spec_supports_grouped_comparison_panels_and_callout_items(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckComparisonPanel",
        prompt_style="uniform",
        slides=[Slide(id="slide0.html", title_html="", body_html="")],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)
    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "titleText": "Service demand differs between standard and priority tiers.",
                "blocks": [
                    {
                        "blockId": "body-0",
                        "type": "body_text",
                        "text": "Governing thought: priority users optimize for response time.",
                        "bbox": {"x": 80, "y": 180, "w": 1000, "h": 70},
                    },
                    {
                        "blockId": "left-title",
                        "type": "group_label",
                        "text": "Standard Tier",
                        "groupKind": "comparison_panel",
                        "renderMode": "group_as_image",
                        "bbox": {"x": 100, "y": 320, "w": 320, "h": 40},
                    },
                    {
                        "blockId": "left-b0",
                        "type": "bullet_item",
                        "text": "Designed for scheduled batch processing",
                        "groupKind": "comparison_panel",
                        "renderMode": "group_as_image",
                        "bbox": {"x": 100, "y": 390, "w": 420, "h": 45},
                    },
                    {
                        "blockId": "right-title",
                        "type": "group_label",
                        "text": "Priority Tier",
                        "groupKind": "comparison_panel",
                        "renderMode": "group_as_image",
                        "bbox": {"x": 700, "y": 320, "w": 320, "h": 40},
                    },
                    {
                        "blockId": "right-b0",
                        "type": "bullet_item",
                        "text": "Priority requests require defined response targets",
                        "groupKind": "comparison_panel",
                        "renderMode": "group_as_image",
                        "bbox": {"x": 700, "y": 390, "w": 420, "h": 45},
                    },
                    {
                        "blockId": "callout-body",
                        "type": "callout_banner",
                        "text": (
                            "Central Synthesis / The Gap\n"
                            "Operators need modular services with stable APIs across critical workflows."
                        ),
                        "items": [
                            "Central Synthesis / The Gap",
                            "Operators need modular services with stable APIs across critical workflows.",
                        ],
                        "bbox": {"x": 160, "y": 840, "w": 820, "h": 80},
                    },
                    {
                        "blockId": "implication-0",
                        "type": "body_text",
                        "text": "IMPLICATION: Pilot with priority workflows first.",
                        "bbox": {"x": 120, "y": 1040, "w": 900, "h": 60},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.kind == "comparison_columns"
    assert slide.callout_title == "Central Synthesis / The Gap"
    assert (
        slide.callout_body
        == "Operators need modular services with stable APIs across critical workflows."
    )
    assert slide.comparison_columns == [
        {
            "title": "Standard Tier",
            "bullets": ["Designed for scheduled batch processing"],
        },
        {
            "title": "Priority Tier",
            "bullets": ["Priority requests require defined response targets"],
        },
    ]


def test_build_slides_pptx_spec_keeps_callout_body_when_comparison_extraction_fails(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckComparisonFallback",
        prompt_style="uniform",
        slides=[Slide(id="slide0.html", title_html="", body_html="")],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)
    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "titleText": "Standard and priority service needs differ materially.",
                "blocks": [
                    {
                        "blockId": "body-0",
                        "type": "body_text",
                        "text": "Governing thought: priority users optimize for response time.",
                        "bbox": {"x": 80, "y": 180, "w": 1000, "h": 70},
                    },
                    {
                        "blockId": "left-title",
                        "type": "group_label",
                        "text": "Standard Tier",
                        "groupKind": "comparison_panel",
                        "renderMode": "group_as_image",
                        "bbox": {"x": 100, "y": 320, "w": 320, "h": 40},
                    },
                    {
                        "blockId": "left-b0",
                        "type": "bullet_item",
                        "text": "Designed for scheduled batch processing",
                        "groupKind": "comparison_panel",
                        "renderMode": "group_as_image",
                        "bbox": {"x": 100, "y": 390, "w": 420, "h": 45},
                    },
                    {
                        "blockId": "callout-body",
                        "type": "callout_banner",
                        "text": "Operators need modular services with stable APIs across critical workflows.",
                        "bbox": {"x": 160, "y": 840, "w": 820, "h": 80},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.kind != "comparison_columns"
    assert slide.callout_body == ""
    assert (
        "Operators need modular services with stable APIs across critical workflows."
        in slide.body
    )


def test_build_slides_pptx_spec_keeps_grouped_exhibit_text_inside_visual_crop(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckGroupedVisualText",
        prompt_style="uniform",
        slides=[
            Slide(
                id="slide0.html",
                title_html="",
                body_html="<img src='assets/source.png' alt='source' />",
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    asset_path = deck_path / "assets" / "source.png"
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (1000, 700), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((520, 120, 880, 470), outline="black", width=4)
    draw.text((610, 220), "Card title", fill="black")
    draw.text((610, 320), "Bullet text", fill="black")
    image.save(asset_path)

    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "assetPath": "assets/source.png",
                "titleText": "Strategic priorities",
                "blocks": [
                    {
                        "blockId": "title-0",
                        "type": "title",
                        "text": "Strategic priorities",
                        "bbox": {"x": 80, "y": 40, "w": 500, "h": 80},
                    },
                    {
                        "blockId": "body-0",
                        "type": "body_text",
                        "text": "We focus on the most actionable segment first.",
                        "bbox": {"x": 80, "y": 150, "w": 700, "h": 70},
                    },
                    {
                        "blockId": "figure-0",
                        "type": "figure",
                        "groupId": "exhibit-1",
                        "groupKind": "exhibit",
                        "visualText": "Card title\nBullet text",
                        "visualItems": ["Card title", "Bullet text"],
                        "bbox": {"x": 520, "y": 120, "w": 360, "h": 350},
                    },
                    {
                        "blockId": "label-0",
                        "type": "exhibit_label",
                        "text": "Card title",
                        "groupId": "exhibit-1",
                        "groupKind": "exhibit",
                        "renderMode": "group_as_image",
                        "bbox": {"x": 610, "y": 220, "w": 120, "h": 40},
                    },
                    {
                        "blockId": "bullet-0",
                        "type": "bullet_item",
                        "text": "Bullet text",
                        "groupId": "exhibit-1",
                        "groupKind": "exhibit",
                        "parentId": "label-0",
                        "bbox": {"x": 610, "y": 320, "w": 140, "h": 40},
                    },
                    {
                        "blockId": "implication-0",
                        "type": "body_text",
                        "text": "IMPLICATION: Preserve the exhibit as one visual unit.",
                        "bbox": {"x": 90, "y": 610, "w": 800, "h": 55},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.kind == "text_visual"
    assert slide.layout_variant == "text_visual_right"
    assert slide.bullets == []
    assert slide.implication == "IMPLICATION: Preserve the exhibit as one visual unit."
    crop_path = deck_path / slide.visual_path
    assert crop_path.exists()
    with Image.open(crop_path) as cropped:
        assert cropped.width >= 320
        assert cropped.height >= 300


def test_build_slides_pptx_spec_extracts_unlabeled_bottom_banner_and_trims_visual_crop(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckBottomBannerTrim",
        prompt_style="uniform",
        slides=[
            Slide(
                id="slide0.html",
                title_html="",
                body_html="<img src='assets/source.png' alt='source' />",
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    asset_path = deck_path / "assets" / "source.png"
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (1000, 700), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((130, 210, 870, 505), outline=(70, 110, 150), width=4)
    draw.rectangle((92, 560, 908, 640), fill=(236, 240, 246))
    draw.text(
        (120, 585),
        "Early usage validates the service before broader rollout.",
        fill="black",
    )
    image.save(asset_path)

    implication_text = "Early usage validates the service before broader rollout."
    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "assetPath": "assets/source.png",
                "titleText": "Near-term value is captured through staged service adoption.",
                "blocks": [
                    {
                        "blockId": "title-0",
                        "type": "title",
                        "text": "Near-term value is captured through staged service adoption.",
                        "bbox": {"x": 40, "y": 28, "w": 840, "h": 72},
                    },
                    {
                        "blockId": "body-0",
                        "type": "body_text",
                        "text": "The synthetic service supports teams with different operating needs.",
                        "bbox": {"x": 60, "y": 130, "w": 840, "h": 56},
                    },
                    {
                        "blockId": "figure-0",
                        "type": "figure",
                        "bbox": {"x": 80, "y": 180, "w": 840, "h": 470},
                    },
                    {
                        "blockId": "banner-0",
                        "type": "body_text",
                        "text": implication_text,
                        "bbox": {"x": 92, "y": 560, "w": 816, "h": 80},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.implication == implication_text
    assert slide.visual_path
    crop_path = deck_path / slide.visual_path
    assert crop_path.exists()
    with Image.open(crop_path) as cropped:
        assert cropped.height < 400


def test_build_slides_pptx_spec_rebuilds_simple_cards_row_as_native_visual(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckNativeCards",
        prompt_style="uniform",
        slides=[
            Slide(
                id="slide0.html",
                title_html="",
                body_html="<img src='assets/source.png' alt='source' />",
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    asset_path = deck_path / "assets" / "source.png"
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (1200, 720), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((90, 250, 350, 560), outline=(98, 118, 145), width=4)
    draw.rectangle((470, 250, 730, 560), outline=(68, 112, 88), width=4)
    draw.rectangle((850, 250, 1110, 560), outline=(68, 112, 88), width=4)
    draw.polygon(
        [
            (370, 395),
            (420, 395),
            (420, 380),
            (450, 410),
            (420, 440),
            (420, 425),
            (370, 425),
        ],
        fill=(145, 153, 166),
    )
    draw.polygon(
        [
            (750, 395),
            (800, 395),
            (800, 380),
            (830, 410),
            (800, 440),
            (800, 425),
            (750, 425),
        ],
        fill=(145, 153, 166),
    )
    image.save(asset_path)

    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "assetPath": "assets/source.png",
                "titleText": "Near-term value is captured through staged service adoption.",
                "blocks": [
                    {
                        "blockId": "title-0",
                        "type": "title",
                        "text": "Near-term value is captured through staged service adoption.",
                        "bbox": {"x": 40, "y": 28, "w": 1040, "h": 72},
                    },
                    {
                        "blockId": "body-0",
                        "type": "body_text",
                        "text": "The synthetic service supports teams with different operating needs.",
                        "bbox": {"x": 60, "y": 130, "w": 1020, "h": 56},
                    },
                    {
                        "blockId": "figure-0",
                        "type": "figure",
                        "groupKind": "flow_diagram",
                        "bbox": {"x": 70, "y": 220, "w": 1060, "h": 360},
                        "visualLines": [
                            {
                                "text": "Guided Setup",
                                "bbox": {"x": 140, "y": 300, "w": 150, "h": 28},
                            },
                            {
                                "text": "(Initial)",
                                "bbox": {"x": 165, "y": 334, "w": 105, "h": 24},
                            },
                            {
                                "text": "Initial validation occurs",
                                "bbox": {"x": 118, "y": 390, "w": 200, "h": 22},
                            },
                            {
                                "text": "through a controlled pilot.",
                                "bbox": {"x": 122, "y": 420, "w": 190, "h": 22},
                            },
                            {
                                "text": "Core Workflow",
                                "bbox": {"x": 532, "y": 300, "w": 136, "h": 28},
                            },
                            {
                                "text": "(Recurring)",
                                "bbox": {"x": 512, "y": 334, "w": 176, "h": 24},
                            },
                            {
                                "text": "Standard processing handles",
                                "bbox": {"x": 495, "y": 390, "w": 206, "h": 22},
                            },
                            {
                                "text": "recurring service requests.",
                                "bbox": {"x": 500, "y": 420, "w": 198, "h": 22},
                            },
                            {
                                "text": "Analytics Layer",
                                "bbox": {"x": 905, "y": 300, "w": 160, "h": 28},
                            },
                            {
                                "text": "(Later)",
                                "bbox": {"x": 920, "y": 334, "w": 136, "h": 24},
                            },
                            {
                                "text": "Monitoring summarizes usage",
                                "bbox": {"x": 885, "y": 390, "w": 210, "h": 22},
                            },
                            {
                                "text": "and overall service quality.",
                                "bbox": {"x": 898, "y": 420, "w": 188, "h": 22},
                            },
                        ],
                    },
                    {
                        "blockId": "implication-0",
                        "type": "implication_banner",
                        "text": "IMPLICATION: Early usage validates the service before broader rollout.",
                        "bbox": {"x": 70, "y": 602, "w": 1040, "h": 66},
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.native_visual is not None
    assert slide.native_visual["kind"] == "cards_row"
    assert len(slide.native_visual["cards"]) == 3
    assert slide.native_visual["connectors"] is True
    assert slide.visual_path == ""
    assert slide.kind == "text_visual"


def test_build_slides_pptx_spec_keeps_complex_matrix_visual_as_image(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckComplexMatrix",
        prompt_style="uniform",
        slides=[
            Slide(
                id="slide0.html",
                title_html="",
                body_html="<img src='assets/source.png' alt='source' />",
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    asset_path = deck_path / "assets" / "source.png"
    asset_path.parent.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (1200, 720), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((620, 250, 760, 360), outline=(120, 120, 120), width=3)
    draw.rectangle((850, 250, 990, 360), outline=(120, 120, 120), width=3)
    draw.rectangle((620, 430, 760, 540), outline=(120, 120, 120), width=3)
    draw.rectangle((850, 430, 990, 540), outline=(120, 120, 120), width=3)
    image.save(asset_path)

    analysis_payload = {
        "deckId": deck.deck_id,
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "assetPath": "assets/source.png",
                "titleText": "The service differentiates through a modular architecture.",
                "blocks": [
                    {
                        "blockId": "title-0",
                        "type": "title",
                        "text": "The service differentiates through a modular architecture.",
                        "bbox": {"x": 40, "y": 28, "w": 1040, "h": 72},
                    },
                    {
                        "blockId": "body-0",
                        "type": "body_text",
                        "text": "Governing thought: reliability and observability matter.",
                        "bbox": {"x": 60, "y": 130, "w": 1020, "h": 56},
                    },
                    {
                        "blockId": "figure-0",
                        "type": "figure",
                        "bbox": {"x": 560, "y": 210, "w": 520, "h": 360},
                        "visualLines": [
                            {
                                "text": "Option A",
                                "bbox": {"x": 670, "y": 260, "w": 60, "h": 24},
                            },
                            {
                                "text": "Option B",
                                "bbox": {"x": 895, "y": 260, "w": 86, "h": 24},
                            },
                            {
                                "text": "Option C",
                                "bbox": {"x": 675, "y": 445, "w": 42, "h": 24},
                            },
                            {
                                "text": "Option D",
                                "bbox": {"x": 900, "y": 445, "w": 62, "h": 24},
                            },
                        ],
                    },
                ],
            }
        ],
    }

    spec = build_slides_pptx_spec(deck, deck_path, slide_analysis=analysis_payload)

    slide = spec.slides[0]
    assert slide.native_visual is None
    assert slide.visual_path


def test_render_slides_pptx_from_template_renders_native_cards_visual_text(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckRenderNativeCards"
    deck_path.mkdir(parents=True, exist_ok=True)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "text_visual",
                "layoutVariant": "text_visual_bottom",
                "title": "Near-term value is captured through staged service adoption.",
                "body": "The synthetic service supports teams with different operating needs.",
                "visualPath": "",
                "visualType": "figure",
                "nativeVisual": {
                    "kind": "cards_row",
                    "connectors": True,
                    "cards": [
                        {
                            "left": 0.02,
                            "top": 0.1,
                            "width": 0.27,
                            "height": 0.76,
                            "title": "Guided Setup (Initial)",
                            "body": "Initial validation occurs through a controlled pilot.",
                            "items": [],
                            "accentRgb": [98, 118, 145],
                            "fullBorder": True,
                        },
                        {
                            "left": 0.365,
                            "top": 0.1,
                            "width": 0.27,
                            "height": 0.76,
                            "title": "Core Workflow (Recurring)",
                            "body": "Standard processing handles recurring service requests.",
                            "items": [],
                            "accentRgb": [68, 112, 88],
                            "fullBorder": False,
                        },
                        {
                            "left": 0.71,
                            "top": 0.1,
                            "width": 0.27,
                            "height": 0.76,
                            "title": "Analytics Layer (Later)",
                            "body": "Monitoring summarizes usage and overall service quality.",
                            "items": [],
                            "accentRgb": [68, 112, 88],
                            "fullBorder": False,
                        },
                    ],
                },
                "implication": "IMPLICATION: Early usage validates the service before broader rollout.",
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide = presentation.slides[0]
    texts = [
        shape.text
        for shape in slide.shapes
        if hasattr(shape, "text") and str(shape.text or "").strip()
    ]

    assert any("Guided Setup (Initial)" in text for text in texts)
    assert any("Core Workflow (Recurring)" in text for text in texts)
    assert any("Analytics Layer (Later)" in text for text in texts)
    assert any(
        "IMPLICATION: Early usage validates the service before broader rollout." in text
        for text in texts
    )


def test_render_slides_pptx_from_template_renders_launch_product_tiles(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckLaunchTiles"
    deck_path.mkdir(parents=True, exist_ok=True)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "text_visual",
                "layoutVariant": "text_visual_bottom",
                "title": "The included launches already show the split inside the cohort",
                "body": "Three launches make the proposition visible without asking the reader to infer the story from a chart alone.",
                "visualPath": "",
                "visualType": "figure",
                "nativeVisual": {
                    "kind": "launch_product_tiles",
                    "tiles": [
                        {
                            "brand": "SYNTHETIC BRAND A",
                            "product": "Product Alpha",
                            "body": "Counterexample inside the cohort.",
                            "tags": ["matte", "full", "long-wear"],
                            "badge": "counterexample",
                            "accentRgb": [83, 54, 48],
                        },
                        {
                            "brand": "SYNTHETIC BRAND B",
                            "product": "Product Beta",
                            "body": "Comfort and care are explicit.",
                            "tags": ["hydrating", "sheer", "long-wear"],
                            "accentRgb": [128, 84, 64],
                        },
                        {
                            "brand": "SYNTHETIC BRAND C",
                            "product": "Product Gamma",
                            "body": "Hydration sits at the center of the story.",
                            "tags": ["hydrating", "matte", "long-wear"],
                            "accentRgb": [103, 121, 88],
                        },
                    ],
                },
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide = presentation.slides[0]
    texts = [
        shape.text
        for shape in slide.shapes
        if hasattr(shape, "text") and str(shape.text or "").strip()
    ]

    assert any("Product Alpha" in text for text in texts)
    assert any("Product Beta" in text for text in texts)
    assert any("Product Gamma" in text for text in texts)
    assert any("counterexample" in text.lower() for text in texts)


def test_render_slides_pptx_table_focus_renders_table_title_and_implication_banner(
    tmp_path: Path,
) -> None:
    from pptx import Presentation

    deck_path = tmp_path / "deckTableFocusBanner"
    deck_path.mkdir(parents=True, exist_ok=True)
    image_path = deck_path / "pptx_assets" / "table.png"
    image_path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", (640, 360), "white").save(image_path)
    spec_payload = {
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "slideId": "slide0.html",
                "kind": "text_visual",
                "layoutVariant": "table_focus",
                "title": "Comparison table",
                "body": "Illustrative values",
                "tableTitle": "Service-Level Matrix",
                "implication": "IMPLICATION: Use the pilot cohort for the controlled rollout.",
                "bullets": [],
                "visualPath": "pptx_assets/table.png",
                "visualType": "table",
                "tableModel": {
                    "source": "deterministic_simple",
                    "confidence": 0.88,
                    "rowCount": 3,
                    "columnCount": 2,
                    "headerRows": 1,
                    "columnWidths": [0.6, 0.4],
                    "rows": [
                        {
                            "cells": [
                                {
                                    "text": "Country",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": True,
                                    "align": "center",
                                },
                                {
                                    "text": "Revenue",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": True,
                                    "align": "center",
                                },
                            ]
                        },
                        {
                            "cells": [
                                {
                                    "text": "Segment A",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "25",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "right",
                                },
                            ]
                        },
                        {
                            "cells": [
                                {
                                    "text": "Segment B",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "left",
                                },
                                {
                                    "text": "31",
                                    "rowSpan": 1,
                                    "colSpan": 1,
                                    "isHeader": False,
                                    "align": "right",
                                },
                            ]
                        },
                    ],
                },
            }
        ],
    }
    (deck_path / SLIDES_PPTX_SPEC_FILENAME).write_text(
        json.dumps(spec_payload, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide = presentation.slides[0]
    texts = [
        shape.text
        for shape in slide.shapes
        if hasattr(shape, "text") and str(shape.text or "").strip()
    ]

    assert any("Illustrative values" in text for text in texts)
    assert any("Service-Level Matrix" in text for text in texts)
    assert any(
        "IMPLICATION: Use the pilot cohort for the controlled rollout." in text
        for text in texts
    )
