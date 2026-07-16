from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from src.slides.launch_brief import (
    build_report_payload_from_launch_brief,
    validate_launch_brief_payload,
)
from src.slides.semantic_pptx import (
    build_slides_pptx_spec_from_report_payload,
    render_slides_pptx_from_template,
    write_slides_pptx_spec,
)


def _sample_launch_brief() -> dict[str, object]:
    return {
        "version": "launch_brief/1",
        "deckName": "Ulta Lipstick Launch Brief Experiment",
        "templateKey": "uniform",
        "promptStyle": "uniform",
        "slides": [
            {
                "role": "cover",
                "title": "Testing whether launch attributes are signal, not noise",
                "body": [
                    "Ulta lipstick launches are the cleanest stress test for early signal.",
                    "The point is to see what survives after basic audit.",
                ],
                "footerText": "Ulta new arrivals | April 2026",
            },
            {
                "role": "launch_tiles",
                "title": "Included launches",
                "body": "The launch set is still comfort-led rather than novelty-led.",
                "implication": "Implication: care language matters more than novelty language.",
                "products": [
                    {
                        "brand": "Brand A",
                        "product": "Hydrating Lip Stick",
                        "body": "Hydrating stick form with high brand familiarity.",
                        "tags": ["Hydrating", "Stick"],
                        "badge": "Core signal",
                    },
                    {
                        "brand": "Brand B",
                        "product": "Daily Lip Tint",
                        "body": "Tint-forward form framed as easy daily wear.",
                        "tags": ["Tint", "Daily wear"],
                    },
                ],
            },
            {
                "role": "comparison",
                "title": "What survived vs what failed",
                "body": "Audit keeps one modest comfort-led signal and removes the sustainability false positive.",
                "left": {
                    "heading": "Survived",
                    "items": [
                        "Hydrating over-indexes",
                        "Long-wear remains a credible support feature",
                    ],
                },
                "right": {
                    "heading": "Failed",
                    "items": [
                        "Refillable does not survive audit",
                        "Sustainability boilerplate inflated the first read",
                    ],
                },
                "calloutTitle": "Bottom line",
                "calloutBody": "The signal is real, but smaller than the first pass implied.",
            },
        ],
    }


def test_validate_launch_brief_payload_accepts_supported_roles() -> None:
    payload = _sample_launch_brief()

    validated = validate_launch_brief_payload(payload)

    assert validated["version"] == "launch_brief/1"
    assert len(validated["slides"]) == 3


def test_validate_launch_brief_payload_rejects_unsupported_role() -> None:
    payload = _sample_launch_brief()
    payload["slides"] = [
        {
            "role": "bespoke_masterpiece",
            "title": "Unsupported",
            "body": "Still unsupported.",
        }
    ]

    with pytest.raises(ValueError, match="unsupported role"):
        validate_launch_brief_payload(payload)


def test_validate_launch_brief_payload_rejects_launch_tiles_with_missing_product_name() -> (
    None
):
    payload = _sample_launch_brief()
    payload["slides"] = [
        {
            "role": "launch_tiles",
            "title": "Included launches",
            "body": "Missing products should fail.",
            "products": [
                {"brand": "Brand only"},
                {"brand": "Another brand", "product": "Real product"},
            ],
        }
    ]

    with pytest.raises(ValueError, match="missing a product name"):
        validate_launch_brief_payload(payload)


def test_validate_launch_brief_payload_rejects_launch_tiles_with_too_few_products() -> (
    None
):
    payload = _sample_launch_brief()
    payload["slides"] = [
        {
            "role": "launch_tiles",
            "title": "Included launches",
            "body": "One product should fail.",
            "products": [{"brand": "Only brand", "product": "Only product"}],
        }
    ]

    with pytest.raises(ValueError, match="at least two products"):
        validate_launch_brief_payload(payload)


def test_build_report_payload_from_launch_brief_compiles_supported_slide_shapes() -> (
    None
):
    payload = _sample_launch_brief()

    report_payload = build_report_payload_from_launch_brief(payload)

    assert report_payload["deckName"] == "Ulta Lipstick Launch Brief Experiment"
    assert report_payload["templateKey"] == "uniform"
    assert len(report_payload["slides"]) == 3
    assert report_payload["slides"][0]["footerText"] == "Ulta new arrivals | April 2026"
    assert report_payload["slides"][1]["layoutVariant"] == "text_visual_bottom"
    assert report_payload["slides"][1]["nativeVisual"]["kind"] == "launch_product_tiles"
    assert report_payload["slides"][2]["comparisonColumns"][0]["title"] == "Survived"


def test_build_report_payload_from_launch_brief_renders_with_shared_semantic_engine(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckLaunchBrief"
    deck_path.mkdir(parents=True, exist_ok=True)
    report_payload = build_report_payload_from_launch_brief(_sample_launch_brief())

    spec = build_slides_pptx_spec_from_report_payload(
        report_payload, deck_path=deck_path
    )
    write_slides_pptx_spec(deck_path, spec)

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide_texts = [
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and str(shape.text or "").strip()
    ]

    assert len(presentation.slides) == 3
    assert any(
        "Testing whether launch attributes are signal, not noise" in text
        for text in slide_texts
    )
    assert any("Included launches" in text for text in slide_texts)
    assert any("What survived vs what failed" in text for text in slide_texts)
