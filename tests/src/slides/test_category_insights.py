from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from src.slides.category_insights import (
    build_launch_brief_from_category_insights,
    validate_category_insights_payload,
)
from src.slides.launch_brief import build_report_payload_from_launch_brief
from src.slides.semantic_pptx import (
    build_slides_pptx_spec_from_report_payload,
    render_slides_pptx_from_template,
    write_slides_pptx_spec,
)


def _sample_category_insights() -> dict[str, object]:
    return {
        "version": "category_insights/1",
        "deckName": "Ulta Lipstick Category Insights",
        "retailer": "Ulta",
        "category": "Lipstick",
        "observationWindow": "April 2026",
        "thesis": "Hydration survives audit as the clearest launch signal.",
        "summary": [
            "New lipstick launches skew comfort-led rather than novelty-led.",
            "The initial sustainability read weakens once retailer boilerplate is excluded.",
        ],
        "evidenceIntro": "Representative launches showing the current shape of the category.",
        "evidenceExamples": [
            {
                "brand": "Brand A",
                "product": "Hydrating Lip Stick",
                "body": "Hydrating stick form with strong brand familiarity.",
                "tags": ["Hydrating", "Stick"],
                "badge": "Core signal",
            },
            {
                "brand": "Brand B",
                "product": "Daily Lip Tint",
                "body": "Tint-led form framed as easy daily wear.",
                "tags": ["Tint", "Daily wear"],
            },
        ],
        "survivingSignals": [
            "Hydrating language over-indexes versus the older base.",
            "Long-wear remains a credible support feature.",
        ],
        "droppedSignals": [
            "Refillable does not survive audit.",
            "Retailer sustainability boilerplate inflated the first read.",
        ],
        "bottomLine": "The signal is real, but smaller than the first pass suggested.",
    }


def test_validate_category_insights_payload_accepts_supported_shape() -> None:
    validated = validate_category_insights_payload(_sample_category_insights())

    assert validated["version"] == "category_insights/1"


def test_validate_category_insights_payload_rejects_missing_evidence() -> None:
    payload = _sample_category_insights()
    payload["evidenceExamples"] = [
        {"brand": "Example Brand", "product": "Hydrating Lip Stick"}
    ]

    with pytest.raises(ValueError, match="at least two 'evidenceExamples'"):
        validate_category_insights_payload(payload)


def test_build_launch_brief_from_category_insights_compiles_three_slide_brief() -> None:
    brief = build_launch_brief_from_category_insights(_sample_category_insights())

    assert brief["version"] == "launch_brief/1"
    assert len(brief["slides"]) == 3
    assert brief["slides"][0]["role"] == "cover"
    assert brief["slides"][1]["role"] == "launch_tiles"
    assert brief["slides"][2]["role"] == "comparison"
    assert brief["slides"][2]["right"]["heading"] == "Failed"


def test_build_launch_brief_from_category_insights_renders_with_shared_semantic_engine(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckCategoryInsights"
    deck_path.mkdir(parents=True, exist_ok=True)
    brief = build_launch_brief_from_category_insights(_sample_category_insights())
    report_payload = build_report_payload_from_launch_brief(brief)
    spec = build_slides_pptx_spec_from_report_payload(
        report_payload, deck_path=deck_path
    )
    write_slides_pptx_spec(deck_path, spec)

    buffer = render_slides_pptx_from_template(deck_path)
    presentation = Presentation(buffer)
    slide_texts = [
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and str(shape.text or "").strip()
    ]

    assert len(presentation.slides) == 3
    assert any(
        "Hydration survives audit as the clearest launch signal." in text
        for text in slide_texts
    )
    assert any("Illustrative launches" in text for text in slide_texts)
