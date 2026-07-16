from __future__ import annotations

import io
import json
import os
import sys
import types
import zipfile
from datetime import UTC, datetime, timedelta
from pathlib import Path
from typing import Callable

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from PIL import Image, ImageDraw

from modules.auth.dependencies import (
    require_authenticated_user,
    require_authenticated_user_for_site,
)
from modules.auth.session import AuthenticatedUser
from tests_tagging_stub import ensure_tagging_stub


def _ensure_jinja_stub() -> None:
    """Provide a minimal jinja2 stub so module imports remain self-contained."""

    if "jinja2" in sys.modules:
        return
    jinja2_stub = types.ModuleType("jinja2")

    def pass_context(func: Callable[..., object]) -> Callable[..., object]:
        return func

    class FileSystemLoader:
        def __init__(self, directory: object) -> None:
            self.directory = directory

    class _Template:
        def render(self, context: dict[str, object]) -> str:
            return ""

    class Environment:
        def __init__(self, **_: object) -> None:
            self.globals: dict[str, object] = {}

        def get_template(self, _name: str) -> _Template:
            return _Template()

    jinja2_stub.pass_context = pass_context
    jinja2_stub.FileSystemLoader = FileSystemLoader
    jinja2_stub.Environment = Environment
    sys.modules["jinja2"] = jinja2_stub


_ensure_jinja_stub()
ensure_tagging_stub()

from modules.slides import api as slides_api
from src.slides.models import Deck, Section, Slide, Subsection
from src.slides.layout_service import build_deck_layout_payload
from src.slides.pptx_template_manifest import (
    DECK_PPTX_TEMPLATE_FILENAME,
    DECK_PPTX_TEMPLATE_MANIFEST_FILENAME,
    build_pptx_template_manifest,
    write_deck_pptx_template_manifest,
)
from src.slides.storage import DeckStorage
from src.slides.tagging import apply_enrichment_patch

INDEX_TEMPLATE = "<html><body><a href='slides/slide0.html'>slide</a></body></html>"
SLIDE_TEMPLATE = "<div class='slide-container'><h1 data-role='title'>{title}</h1><div class='slide-body'><p>Body</p></div></div>"


@pytest.fixture
def deck_storage(tmp_path: Path):
    original = slides_api.get_storage()
    storage = DeckStorage(tmp_path)
    slides_api.set_storage(storage)
    yield storage
    slides_api.set_storage(original)


@pytest.fixture
def client(deck_storage: DeckStorage) -> TestClient:
    app = FastAPI()
    app.include_router(slides_api.router)
    app.dependency_overrides[require_authenticated_user] = lambda: None
    return TestClient(app)


@pytest.fixture
def stub_uploaded_deck_processing(monkeypatch: pytest.MonkeyPatch) -> None:
    def _fake_build_layout(
        deck: Deck,
        deck_path: Path,
        *,
        lang: str = "eng",
        progress_callback: Callable[[int, int], None] | None = None,
    ) -> dict[str, object]:
        total = len(deck.slides)
        if progress_callback is not None:
            progress_callback(total, total)
        return {
            "deckId": deck.deck_id,
            "lang": lang,
            "generatedAt": "2026-03-18T00:00:00+00:00",
            "slides": [
                {
                    "slideId": slide.id,
                    "slideNumber": index + 1,
                    "pageNumber": index + 1,
                    "assetPath": f"assets/{index + 1}.png",
                    "blocks": [
                        {
                            "blockId": "block-0",
                            "type": "title",
                            "text": "",
                            "items": [],
                            "bbox": {"x": 10.0, "y": 20.0, "w": 300.0, "h": 80.0},
                            "confidence": 0.5,
                        },
                        {
                            "blockId": "block-1",
                            "type": "table",
                            "text": "",
                            "items": [],
                            "bbox": {"x": 40.0, "y": 140.0, "w": 360.0, "h": 220.0},
                        },
                    ],
                    "titleText": str(slide.title_html or ""),
                    "bulletTexts": [],
                    "figureRegions": [],
                }
                for index, slide in enumerate(deck.slides)
            ],
        }

    def _fake_build_ocr(
        deck: Deck,
        deck_path: Path,
        *,
        lang: str = "eng",
        include_bboxes: bool = True,
        layout_payload: dict[str, object] | None = None,
        pdf_path: Path | None = None,
        progress_callback: Callable[[int, int], None] | None = None,
    ) -> dict[str, object]:
        raw_slides = (
            layout_payload.get("slides")
            if isinstance(layout_payload, dict)
            and isinstance(layout_payload.get("slides"), list)
            else []
        )
        total = len(raw_slides)
        if progress_callback is not None:
            progress_callback(total, total)
        slides = []
        for index, slide in enumerate(raw_slides):
            if not isinstance(slide, dict):
                continue
            slide_id = str(slide.get("slideId") or "")
            text = f"OCR {slide_id}"
            slides.append(
                {
                    "slide_id": slide_id,
                    "slide_number": index + 1,
                    "page_number": index + 1,
                    "ocr_text": text,
                    "lines": [{"text": text}],
                    "blocks": [
                        {
                            "block_id": "block-0",
                            "type": "list",
                            "text": text,
                            "items": [text],
                            "bbox": {"x": 10.0, "y": 20.0, "w": 300.0, "h": 80.0},
                            "confidence": 0.9,
                        }
                    ],
                    "title_text": text,
                    "bullet_texts": [text],
                    "figure_regions": [],
                }
            )
        return {
            "deck_id": deck.deck_id,
            "lang": lang,
            "ocr_strategy": slides_api.OCR_STRATEGY_LAYOUT_GUIDED,
            "generated_at": "2026-03-18T00:00:00+00:00",
            "slides": slides,
        }

    monkeypatch.setattr(slides_api, "build_deck_layout_payload", _fake_build_layout)
    monkeypatch.setattr(slides_api, "build_deck_ocr_payload", _fake_build_ocr)

    class _ImmediateExecutor:
        def submit(self, fn, *args, **kwargs):
            fn(*args, **kwargs)
            return None

    monkeypatch.setattr(slides_api, "_OCR_EXECUTOR", _ImmediateExecutor())


@pytest.fixture
def site_client() -> TestClient:
    app = FastAPI()
    app.include_router(slides_api.site_router)
    app.dependency_overrides[require_authenticated_user_for_site] = lambda: None
    return TestClient(app)


def _seed_deck(
    storage: DeckStorage,
    deck_id: str,
    *,
    prompt_style: str = "uniform",
    owner_email: str | None = None,
    shared_with: list[str] | None = None,
) -> None:
    storage.save_deck(
        Deck(
            deck_id=deck_id,
            prompt_style=prompt_style,
            owner_email=owner_email,
            shared_with=shared_with or [],
            slides=[
                Slide(id="slide0.html", title_html="First", body_html="<p>Alpha</p>"),
                Slide(id="slide1.html", title_html="Second", body_html="<p>Beta</p>"),
            ],
        )
    )


def _seed_image_deck(
    storage: DeckStorage,
    deck_id: str,
    *,
    owner_email: str | None = None,
    shared_with: list[str] | None = None,
) -> None:
    storage.save_deck(
        Deck(
            deck_id=deck_id,
            owner_email=owner_email,
            shared_with=shared_with or [],
            slides=[
                Slide(
                    id="slide0.html",
                    title_html="Image slide",
                    body_html="<img src='chart.png' alt='chart' />",
                )
            ],
        )
    )


def _seed_custom_deck(
    storage: DeckStorage,
    deck_id: str,
    titles: list[str],
    *,
    include_sections: bool = False,
    prompt_style: str = "uniform",
    owner_email: str | None = None,
    shared_with: list[str] | None = None,
) -> None:
    slides = [
        Slide(id=f"slide{index}.html", title_html=title, body_html=f"<p>{title}</p>")
        for index, title in enumerate(titles)
    ]
    sections: list[Section] = []
    if include_sections and slides:
        sections.append(
            Section(
                id=f"{deck_id}-section",
                title=f"{deck_id} Section",
                start_slide=slides[0].id,
                subsections=[
                    Subsection(
                        id=f"{deck_id}-sub",
                        title=f"{deck_id} Subsection",
                        start_slide=slides[-1].id,
                    )
                ],
            )
        )
    storage.save_deck(
        Deck(
            deck_id=deck_id,
            prompt_style=prompt_style,
            owner_email=owner_email,
            shared_with=shared_with or [],
            slides=slides,
            sections=sections,
        )
    )


def test_get_deck_returns_slides(client: TestClient, deck_storage: DeckStorage) -> None:
    _seed_deck(deck_storage, "deckA")

    response = client.get("/slides/deck/deckA")

    assert response.status_code == 200
    data = response.json()
    assert data["deckId"] == "deckA"
    assert data["promptStyle"] == "uniform"
    assert len(data["slides"]) == 2
    assert data["sections"] == []


def test_list_decks_filters_to_owner_and_shared(
    deck_storage: DeckStorage,
) -> None:
    user = AuthenticatedUser(email="viewer@example.com")
    _seed_deck(deck_storage, "owned", owner_email="viewer@example.com")
    _seed_deck(
        deck_storage,
        "shared",
        owner_email="owner@example.com",
        shared_with=["viewer@example.com"],
    )
    _seed_deck(deck_storage, "hidden", owner_email="other@example.com")
    app = FastAPI()
    app.include_router(slides_api.router)
    app.dependency_overrides[require_authenticated_user] = lambda: user
    client = TestClient(app)

    response = client.get("/slides/decks")

    assert response.status_code == 200
    decks = {deck["deckId"] for deck in response.json()["decks"]}
    assert decks == {"owned", "shared"}


def test_list_decks_returns_decks_without_chart_remap_recovery_step(
    client: TestClient,
    deck_storage: DeckStorage,
) -> None:
    _seed_deck(deck_storage, "deckA")

    response = client.get("/slides/decks")

    assert response.status_code == 200
    payload = response.json()
    assert payload["decks"]
    assert payload["decks"][0]["deckId"] == "deckA"


def test_get_deck_missing_prompt_style_defaults_to_uniform(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    deck_dir = deck_storage.root / "legacyDeck"
    deck_dir.mkdir(parents=True, exist_ok=True)
    (deck_dir / "index.json").write_text(
        json.dumps({"slides": ["slide0.html"]}), encoding="utf-8"
    )
    (deck_dir / "slide0.html").write_text(
        SLIDE_TEMPLATE.format(title="Legacy"),
        encoding="utf-8",
    )

    response = client.get("/slides/deck/legacyDeck")

    assert response.status_code == 200
    assert response.json()["promptStyle"] == "uniform"


def test_save_deck_updates_files(client: TestClient, deck_storage: DeckStorage) -> None:
    _seed_deck(deck_storage, "deckB")

    payload = {
        "slides": [
            {
                "id": "slide0.html",
                "titleHtml": "Updated",
                "bodyHtml": "<p>Alpha</p>",
                "kind": "normal",
            },
            {
                "id": "slide2.html",
                "titleHtml": "",
                "bodyHtml": "",
                "kind": "sectionHeader",
                "sectionId": "A",
            },
        ],
        "sections": [
            {
                "id": "A",
                "title": "Section A",
                "startSlide": "slide0.html",
                "subsections": [],
            }
        ],
    }

    response = client.post("/slides/deck/deckB/save", json=payload)

    assert response.status_code == 200
    assert response.json()["promptStyle"] == "uniform"
    deck_dir = deck_storage.root / "deckB"
    assert (deck_dir / "slide2.html").exists()
    assert "Updated" in (deck_dir / "slide0.html").read_text(encoding="utf-8")
    index_payload = json.loads((deck_dir / "index.json").read_text(encoding="utf-8"))
    assert index_payload["promptStyle"] == "uniform"
    assert any(
        isinstance(entry, dict) and entry.get("kind") == "sectionHeader"
        for entry in index_payload["slides"]
    )
    assert index_payload["sections"][0]["id"] == "A"


def test_save_deck_rebuilds_slide_analysis_from_pruned_layout_and_ocr(
    client: TestClient,
    deck_storage: DeckStorage,
) -> None:
    _seed_deck(deck_storage, "deckArtifacts")
    deck_storage.save_layout_payload(
        "deckArtifacts",
        {
            "deckId": "deckArtifacts",
            "lang": "eng",
            "generatedAt": "2026-03-18T00:00:00+00:00",
            "slides": [
                {
                    "slideId": "slide0.html",
                    "slideNumber": 1,
                    "pageNumber": 1,
                    "assetPath": "assets/a.png",
                    "blocks": [
                        {
                            "blockId": "block-0",
                            "type": "title",
                            "text": "",
                            "items": [],
                            "bbox": {"x": 0.0, "y": 0.0, "w": 10.0, "h": 10.0},
                            "confidence": 0.5,
                        }
                    ],
                    "titleText": "",
                    "bulletTexts": [],
                    "figureRegions": [],
                },
                {
                    "slideId": "slide1.html",
                    "slideNumber": 2,
                    "pageNumber": 2,
                    "assetPath": "assets/b.png",
                    "blocks": [
                        {
                            "blockId": "block-0",
                            "type": "title",
                            "text": "",
                            "items": [],
                            "bbox": {"x": 0.0, "y": 0.0, "w": 10.0, "h": 10.0},
                            "confidence": 0.5,
                        }
                    ],
                    "titleText": "",
                    "bulletTexts": [],
                    "figureRegions": [],
                },
            ],
        },
    )
    deck_storage.save_ocr_payload(
        "deckArtifacts",
        {
            "deck_id": "deckArtifacts",
            "lang": "eng",
            "ocr_strategy": slides_api.OCR_STRATEGY_LAYOUT_GUIDED,
            "generated_at": "2026-03-18T00:00:00+00:00",
            "slides": [
                {
                    "slide_id": "slide0.html",
                    "slide_number": 1,
                    "page_number": 1,
                    "ocr_text": "First OCR",
                    "lines": [{"text": "First OCR"}],
                    "blocks": [
                        {
                            "block_id": "block-0",
                            "type": "list",
                            "text": "First OCR",
                            "items": ["First OCR"],
                            "bbox": {"x": 0.0, "y": 0.0, "w": 10.0, "h": 10.0},
                            "confidence": 0.9,
                        }
                    ],
                    "title_text": "First OCR",
                    "bullet_texts": ["First OCR"],
                    "figure_regions": [],
                },
                {
                    "slide_id": "slide1.html",
                    "slide_number": 2,
                    "page_number": 2,
                    "ocr_text": "Second OCR",
                    "lines": [{"text": "Second OCR"}],
                    "blocks": [
                        {
                            "block_id": "block-0",
                            "type": "list",
                            "text": "Second OCR",
                            "items": ["Second OCR"],
                            "bbox": {"x": 0.0, "y": 0.0, "w": 10.0, "h": 10.0},
                            "confidence": 0.9,
                        }
                    ],
                    "title_text": "Second OCR",
                    "bullet_texts": ["Second OCR"],
                    "figure_regions": [],
                },
            ],
        },
    )
    deck_storage.save_slide_analysis_payload(
        "deckArtifacts",
        {
            "deckId": "deckArtifacts",
            "lang": "eng",
            "generatedAt": "2026-03-18T00:00:00+00:00",
            "slides": [],
        },
    )

    response = client.post(
        "/slides/deck/deckArtifacts/save",
        json={
            "slides": [
                {
                    "id": "slide1.html",
                    "titleHtml": "Second",
                    "bodyHtml": "<p>Beta</p>",
                    "kind": "normal",
                }
            ],
            "sections": [],
        },
    )

    assert response.status_code == 200
    layout_payload = deck_storage.load_layout_payload("deckArtifacts")
    ocr_payload = deck_storage.load_ocr_payload("deckArtifacts")
    analysis_payload = deck_storage.load_slide_analysis_payload("deckArtifacts")
    assert layout_payload is not None
    assert ocr_payload is not None
    assert analysis_payload is not None
    assert [slide["slideId"] for slide in layout_payload["slides"]] == ["slide1.html"]
    assert [slide["slide_id"] for slide in ocr_payload["slides"]] == ["slide1.html"]
    assert analysis_payload["slides"][0]["slideId"] == "slide1.html"
    assert analysis_payload["slides"][0]["blocks"][0]["type"] == "bullet_item"
    assert analysis_payload["slides"][0]["blocks"][0]["text"] == "Second OCR"


def test_build_slide_analysis_payload_preserves_visual_ocr_fields() -> None:
    layout_payload = {
        "deckId": "deckVisualFields",
        "lang": "eng",
        "generatedAt": "2026-03-18T00:00:00+00:00",
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "assetPath": "assets/slide0.png",
                "blocks": [
                    {
                        "blockId": "figure-0",
                        "type": "figure",
                        "text": "",
                        "items": [],
                        "bbox": {"x": 20.0, "y": 40.0, "w": 300.0, "h": 180.0},
                    }
                ],
                "titleText": "",
                "bulletTexts": [],
                "figureRegions": [{"x": 20.0, "y": 40.0, "w": 300.0, "h": 180.0}],
            }
        ],
    }
    ocr_payload = {
        "deck_id": "deckVisualFields",
        "lang": "eng",
        "generated_at": "2026-03-18T00:00:00+00:00",
        "slides": [
            {
                "slide_id": "slide0.html",
                "slide_number": 1,
                "page_number": 1,
                "ocr_text": "",
                "lines": [],
                "blocks": [
                    {
                        "block_id": "figure-0",
                        "type": "figure",
                        "text": "",
                        "items": [],
                        "visual_text": "Phase 1: Italy",
                        "visual_items": ["Phase 1: Italy", "Pilot fleets"],
                        "visual_lines": [
                            {
                                "text": "Phase 1: Italy",
                                "bbox": {"x": 42.0, "y": 62.0, "w": 110.0, "h": 18.0},
                            }
                        ],
                        "bbox": {"x": 20.0, "y": 40.0, "w": 300.0, "h": 180.0},
                    }
                ],
                "title_text": "",
                "bullet_texts": [],
                "figure_regions": [{"x": 20.0, "y": 40.0, "w": 300.0, "h": 180.0}],
            }
        ],
    }

    merged = slides_api._build_slide_analysis_payload(
        layout_payload,
        ocr_payload,
        deck_id="deckVisualFields",
        lang="eng",
    )

    assert merged is not None
    blocks = merged["slides"][0]["blocks"]
    assert blocks[0]["visualText"] == "Phase 1: Italy"
    assert blocks[0]["visualItems"] == ["Phase 1: Italy", "Pilot fleets"]
    assert blocks[0]["visualLines"][0]["text"] == "Phase 1: Italy"


def test_deck_ocr_sends_completion_email_notification_after_deck_processing(
    deck_storage: DeckStorage, monkeypatch: pytest.MonkeyPatch
) -> None:
    user = AuthenticatedUser(email="viewer@example.com")
    app = FastAPI()
    app.include_router(slides_api.router)
    app.dependency_overrides[require_authenticated_user] = lambda: user
    client = TestClient(app)
    _seed_image_deck(deck_storage, "deckOcrNotify", owner_email=user.email)

    fake_payload = {
        "deck_id": "deckOcrNotify",
        "lang": "ita",
        "generated_at": "2026-02-17T00:00:00+00:00",
        "slides": [
            {
                "slide_id": "slide0.html",
                "slide_number": 1,
                "page_number": 1,
                "ocr_text": "hello",
                "lines": [],
            }
        ],
    }
    monkeypatch.setattr(
        slides_api,
        "build_deck_ocr_payload",
        lambda *args, **kwargs: fake_payload,
    )
    monkeypatch.setattr(
        slides_api,
        "build_deck_layout_payload",
        lambda *args, **kwargs: {
            "deck_id": "deckOcrNotify",
            "lang": "ita",
            "generated_at": "2026-02-17T00:00:00+00:00",
            "slides": [
                {
                    "slide_id": "slide0.html",
                    "slide_number": 1,
                    "page_number": 1,
                    "asset_path": "assets/chart.png",
                    "blocks": [],
                    "title_text": "",
                    "bullet_texts": [],
                    "figure_regions": [],
                }
            ],
        },
    )

    notifications: list[tuple[float, str, dict[str, str]]] = []

    def _capture_notify(
        elapsed_sec: float, step: str, session_data: dict[str, str], notifier=None
    ) -> None:
        notifications.append((elapsed_sec, step, session_data))

    monkeypatch.setattr(slides_api, "notify_finished", _capture_notify, raising=False)

    response = client.post("/slides/deck/deckOcrNotify/ocr", json={"lang": "ita"})

    assert response.status_code == 200
    assert len(notifications) == 1
    elapsed_sec, step, session_data = notifications[0]
    assert elapsed_sec >= 0
    assert step == "deck_processing"
    assert session_data["notify_email"] == user.email
    assert session_data["notify_lang"] == "it"


def test_slides_ocr_returns_raw_payload_and_plain_text(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    raw_ocr = [
        [
            [
                [0.0, 0.0],
                [16.0, 0.0],
                [16.0, 8.0],
                [0.0, 8.0],
            ],
            ["Hello world", 0.96],
        ]
    ]

    monkeypatch.setattr(
        slides_api,
        "extract_raw_ocr_from_data_url",
        lambda data_url, lang: raw_ocr,
    )

    response = client.post(
        "/slides/ocr",
        json={
            "imageDataUrl": "data:image/png;base64,ZmFrZQ==",
            "lang": "eng",
        },
    )

    assert response.status_code == 200
    assert response.json() == {
        "ocrText": "Hello world",
        "rawOcr": raw_ocr,
    }


def test_slides_layout_returns_raw_layout_and_blocks(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    raw_layout = [
        {
            "type": "title",
            "bbox": [0.0, 0.0, 120.0, 36.0],
            "res": [["Main title", 0.97]],
        },
        {
            "type": "figure",
            "bbox": [140.0, 80.0, 400.0, 260.0],
            "res": [],
        },
    ]

    monkeypatch.setattr(
        slides_api,
        "extract_raw_layout_from_data_url",
        lambda data_url, lang: raw_layout,
    )
    monkeypatch.setattr(
        slides_api,
        "extract_layout_summary_from_raw_layout",
        lambda raw_layout_arg, slide_id=None, slide_number=None: {
            "blocks": [
                {
                    "block_id": "block-0",
                    "type": "title",
                    "text": "Main title",
                    "items": ["Main title"],
                    "bbox": {"x": 0.0, "y": 0.0, "w": 120.0, "h": 36.0},
                    "confidence": 0.97,
                },
                {
                    "block_id": "block-1",
                    "type": "figure",
                    "text": "",
                    "items": [],
                    "bbox": {"x": 140.0, "y": 80.0, "w": 260.0, "h": 180.0},
                    "confidence": 0.9,
                },
            ],
            "title_text": "Main title",
            "bullet_texts": ["First point"],
            "figure_regions": [{"x": 140.0, "y": 80.0, "w": 260.0, "h": 180.0}],
        },
    )

    response = client.post(
        "/slides/layout",
        json={
            "imageDataUrl": "data:image/png;base64,ZmFrZQ==",
            "lang": "eng",
            "slideId": "slide-1",
            "slideNumber": 1,
        },
    )

    assert response.status_code == 200
    body = response.json()
    assert body["rawLayout"] == raw_layout
    assert body["titleText"] == "Main title"
    assert body["bulletTexts"] == ["First point"]
    assert body["figureRegions"] == [{"x": 140.0, "y": 80.0, "w": 260.0, "h": 180.0}]
    assert body["blocks"][0]["blockId"] == "block-0"
    assert body["blocks"][0]["type"] == "title"
    assert body["blocks"][0]["text"] == "Main title"
    assert body["blocks"][0]["items"] == ["Main title"]
    assert body["blocks"][0]["bbox"] == {"x": 0.0, "y": 0.0, "w": 120.0, "h": 36.0}
    assert body["blocks"][0]["confidence"] == 0.97
    assert body["blocks"][1]["blockId"] == "block-1"
    assert body["blocks"][1]["type"] == "figure"
    assert body["blocks"][1]["text"] == ""
    assert body["blocks"][1]["items"] == []
    assert body["blocks"][1]["bbox"] == {"x": 140.0, "y": 80.0, "w": 260.0, "h": 180.0}
    assert body["blocks"][1]["confidence"] == 0.9


def test_deck_layout_builds_persists_and_marks_deck(
    client: TestClient, deck_storage: DeckStorage, monkeypatch: pytest.MonkeyPatch
) -> None:
    _seed_image_deck(deck_storage, "deckLayout")
    fake_payload = {
        "deck_id": "deckLayout",
        "lang": "eng",
        "generated_at": "2026-03-13T00:00:00+00:00",
        "slides": [
            {
                "slide_id": "slide0.html",
                "slide_number": 1,
                "page_number": 1,
                "asset_path": "assets/chart.png",
                "blocks": [
                    {
                        "block_id": "block-0",
                        "type": "title",
                        "text": "Main title",
                        "items": ["Main title"],
                        "bbox": {"x": 0.0, "y": 0.0, "w": 120.0, "h": 36.0},
                        "confidence": 0.97,
                    }
                ],
                "title_text": "Main title",
                "bullet_texts": ["First point"],
                "figure_regions": [{"x": 140.0, "y": 80.0, "w": 260.0, "h": 180.0}],
            }
        ],
    }
    monkeypatch.setattr(
        slides_api,
        "build_deck_layout_payload",
        lambda deck, deck_path, lang="eng": fake_payload,
    )

    response = client.post("/slides/deck/deckLayout/layout", json={"lang": "eng"})

    assert response.status_code == 200
    body = response.json()
    assert body["cached"] is False
    assert body["payload"]["deckId"] == "deckLayout"
    assert body["payload"]["slides"][0]["slideId"] == "slide0.html"
    assert deck_storage.load_layout_payload("deckLayout") == body["payload"]
    assert deck_storage.load_slide_analysis_payload("deckLayout") == body["payload"]

    deck_response = client.get("/slides/deck/deckLayout")
    assert deck_response.status_code == 200
    assert deck_response.json()["hasLayout"] is True


def test_deck_layout_reuses_cached_payload(
    client: TestClient, deck_storage: DeckStorage, monkeypatch: pytest.MonkeyPatch
) -> None:
    _seed_image_deck(deck_storage, "deckLayoutCached")
    cached_payload = {
        "deckId": "deckLayoutCached",
        "lang": "eng",
        "generatedAt": "2026-03-13T00:00:00+00:00",
        "slides": [
            {
                "slideId": "slide0.html",
                "slideNumber": 1,
                "pageNumber": 1,
                "assetPath": "assets/chart.png",
                "blocks": [],
                "titleText": "",
                "bulletTexts": [],
                "figureRegions": [],
            }
        ],
    }
    deck_storage.save_layout_payload("deckLayoutCached", cached_payload)

    def _unexpected(*args, **kwargs):
        raise AssertionError(
            "build_deck_layout_payload should not run for cached payloads"
        )

    monkeypatch.setattr(slides_api, "build_deck_layout_payload", _unexpected)

    response = client.post("/slides/deck/deckLayoutCached/layout", json={"lang": "eng"})

    assert response.status_code == 200
    assert response.json() == {
        "cached": True,
        "payload": cached_payload,
    }
    assert (
        deck_storage.load_slide_analysis_payload("deckLayoutCached") == cached_payload
    )


def test_deck_layout_merges_cached_ocr_text_into_cached_layout(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_image_deck(deck_storage, "deckLayoutMerged")
    deck_storage.save_layout_payload(
        "deckLayoutMerged",
        {
            "deckId": "deckLayoutMerged",
            "lang": "eng",
            "generatedAt": "2026-03-13T00:00:00+00:00",
            "slides": [
                {
                    "slideId": "slide0.html",
                    "slideNumber": 1,
                    "pageNumber": 1,
                    "assetPath": "assets/chart.png",
                    "blocks": [
                        {
                            "blockId": "block-0",
                            "type": "title",
                            "text": "",
                            "items": [],
                            "bbox": {"x": 10.0, "y": 20.0, "w": 300.0, "h": 80.0},
                            "confidence": 0.5,
                        },
                        {
                            "blockId": "block-1",
                            "type": "table",
                            "text": "",
                            "items": [],
                            "bbox": {"x": 40.0, "y": 140.0, "w": 360.0, "h": 220.0},
                        },
                    ],
                    "titleText": "",
                    "bulletTexts": [],
                    "figureRegions": [],
                }
            ],
        },
    )
    deck_storage.save_ocr_payload(
        "deckLayoutMerged",
        {
            "deck_id": "deckLayoutMerged",
            "lang": "eng",
            "generated_at": "2026-03-13T00:00:00+00:00",
            "slides": [
                {
                    "slide_id": "slide0.html",
                    "slide_number": 1,
                    "page_number": 1,
                    "ocr_text": "Main title",
                    "lines": [],
                    "blocks": [
                        {
                            "block_id": "block-0",
                            "type": "list",
                            "text": "Main title",
                            "items": ["Main title"],
                            "bbox": {"x": 10.0, "y": 20.0, "w": 300.0, "h": 80.0},
                            "confidence": 0.9,
                        },
                        {
                            "block_id": "block-1",
                            "type": "table",
                            "text": "Country Revenue",
                            "items": [],
                            "bbox": {"x": 40.0, "y": 140.0, "w": 360.0, "h": 220.0},
                            "table_model": {
                                "source": "deterministic_simple",
                                "confidence": 0.84,
                                "row_count": 2,
                                "column_count": 2,
                                "header_rows": 1,
                                "column_widths": [0.6, 0.4],
                                "has_merged_cells": False,
                                "rows": [
                                    {
                                        "cells": [
                                            {
                                                "text": "Country",
                                                "row_span": 1,
                                                "col_span": 1,
                                                "is_header": True,
                                                "align": "center",
                                            },
                                            {
                                                "text": "Revenue",
                                                "row_span": 1,
                                                "col_span": 1,
                                                "is_header": True,
                                                "align": "center",
                                            },
                                        ]
                                    },
                                    {
                                        "cells": [
                                            {
                                                "text": "Italy",
                                                "row_span": 1,
                                                "col_span": 1,
                                                "is_header": False,
                                                "align": "left",
                                            },
                                            {
                                                "text": "25",
                                                "row_span": 1,
                                                "col_span": 1,
                                                "is_header": False,
                                                "align": "right",
                                            },
                                        ]
                                    },
                                ],
                            },
                        },
                    ],
                    "title_text": "Main title",
                    "bullet_texts": ["Main title"],
                    "figure_regions": [],
                }
            ],
        },
    )

    response = client.post("/slides/deck/deckLayoutMerged/layout", json={"lang": "eng"})

    assert response.status_code == 200
    body = response.json()
    assert body["cached"] is True
    assert body["payload"]["slides"][0]["titleText"] == "Main title"
    assert body["payload"]["slides"][0]["blocks"][0]["type"] == "bullet_item"
    assert body["payload"]["slides"][0]["blocks"][0]["text"] == "Main title"
    assert body["payload"]["slides"][0]["blocks"][0]["items"] == ["Main title"]
    assert body["payload"]["slides"][0]["blocks"][0]["confidence"] == 0.9
    assert body["payload"]["slides"][0]["blocks"][1]["tableModel"] is not None
    assert body["payload"]["slides"][0]["blocks"][1]["tableModel"]["column_count"] == 2
    assert body["payload"]["slides"][0]["bulletTexts"] == ["Main title"]
    cached_layout = deck_storage.load_layout_payload("deckLayoutMerged")
    assert cached_layout is not None
    assert cached_layout["deckId"] == "deckLayoutMerged"
    assert cached_layout["slides"][0]["slideId"] == "slide0.html"
    assert cached_layout["slides"][0]["assetPath"] == "assets/chart.png"
    assert cached_layout["slides"][0]["blocks"][0]["type"] == "title"
    assert cached_layout["slides"][0]["blocks"][1]["type"] == "table"
    assert (
        deck_storage.load_slide_analysis_payload("deckLayoutMerged") == body["payload"]
    )


def test_ocr_payload_uses_layout_guided_strategy_accepts_versioned_variants() -> None:
    assert slides_api._ocr_payload_uses_layout_guided_strategy(
        {"ocr_strategy": slides_api.OCR_STRATEGY_LAYOUT_GUIDED}
    )
    assert slides_api._ocr_payload_uses_layout_guided_strategy(
        {"ocr_strategy": "layout_guided_text_region_assignment_v5"}
    )
    assert slides_api._ocr_payload_uses_layout_guided_strategy(
        {"ocrStrategy": "layout_guided_text_region_assignment_v6"}
    )
    assert not slides_api._ocr_payload_uses_layout_guided_strategy(
        {"ocr_strategy": "full_slide_ocr_v1"}
    )


def test_build_deck_layout_payload_normalizes_relative_deck_path(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    deck_root = Path("slide_decks") / "relativeLayoutDeck"
    absolute_deck_path = (Path.cwd() / deck_root).resolve()
    assets_dir = absolute_deck_path / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    image_path = assets_dir / "chart.png"
    image_path.write_bytes(b"fake-png")

    deck = Deck(
        deck_id="relativeLayoutDeck",
        slides=[
            Slide(
                id="slide0.html",
                title_html="Image slide",
                body_html="<img src='chart.png' alt='chart' />",
            )
        ],
    )

    monkeypatch.setattr(
        "src.slides.layout_service.extract_raw_layout_from_image_path",
        lambda _image_path, _lang: [{"type": "figure"}],
    )
    monkeypatch.setattr(
        "src.slides.layout_service.extract_layout_summary_from_raw_layout",
        lambda _raw_layout, slide_id=None, slide_number=None: {
            "blocks": [],
            "title_text": "",
            "bullet_texts": [],
            "figure_regions": [],
        },
    )

    try:
        payload = build_deck_layout_payload(
            deck, deck_root, lang="eng", apply_semantic_correction=False
        )
    finally:
        if absolute_deck_path.exists():
            for child in sorted(absolute_deck_path.rglob("*"), reverse=True):
                if child.is_file():
                    child.unlink()
                elif child.is_dir():
                    child.rmdir()
            absolute_deck_path.rmdir()

    assert payload["slides"][0]["asset_path"] == "assets/chart.png"


def test_build_deck_layout_payload_prunes_redundant_unknown_blocks(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    deck_root = Path("slide_decks") / "layoutPruneDeck"
    absolute_deck_path = (Path.cwd() / deck_root).resolve()
    assets_dir = absolute_deck_path / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    image_path = assets_dir / "chart.png"
    image_path.write_bytes(b"fake-png")

    deck = Deck(
        deck_id="layoutPruneDeck",
        slides=[
            Slide(
                id="slide0.html",
                title_html="Image slide",
                body_html="<img src='chart.png' alt='chart' />",
            )
        ],
    )

    monkeypatch.setattr(
        "src.slides.layout_service.extract_raw_layout_from_image_path",
        lambda _image_path, _lang: [{"type": "title"}],
    )
    monkeypatch.setattr(
        "src.slides.layout_service.extract_layout_summary_from_raw_layout",
        lambda _raw_layout, slide_id=None, slide_number=None: {
            "blocks": [
                {
                    "block_id": "title-0",
                    "type": "title",
                    "text": "Main title",
                    "bbox": {"x": 10.0, "y": 10.0, "w": 300.0, "h": 60.0},
                    "confidence": 0.9,
                },
                {
                    "block_id": "unknown-duplicate",
                    "type": "unknown",
                    "text": "",
                    "bbox": {"x": 10.0, "y": 10.0, "w": 300.0, "h": 60.0},
                    "confidence": 0.8,
                },
                {
                    "block_id": "unknown-tiny",
                    "type": "unknown",
                    "text": "",
                    "bbox": {"x": 400.0, "y": 400.0, "w": 8.0, "h": 8.0},
                    "confidence": 0.6,
                },
            ],
            "title_text": "Main title",
            "bullet_texts": [],
            "figure_regions": [],
        },
    )

    try:
        payload = build_deck_layout_payload(
            deck, deck_root, lang="eng", apply_semantic_correction=False
        )
    finally:
        if absolute_deck_path.exists():
            for child in sorted(absolute_deck_path.rglob("*"), reverse=True):
                if child.is_file():
                    child.unlink()
                elif child.is_dir():
                    child.rmdir()
            absolute_deck_path.rmdir()

    blocks = payload["slides"][0]["blocks"]
    assert len(blocks) == 1
    assert blocks[0]["type"] == "title"
    assert blocks[0]["block_id"] == "title-0"


def test_start_deck_layout_returns_completed_when_cached(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_image_deck(deck_storage, "deckLayoutStatusCached")
    cached_payload = {
        "deckId": "deckLayoutStatusCached",
        "lang": "eng",
        "generatedAt": "2026-03-13T00:00:00+00:00",
        "slides": [],
    }
    deck_storage.save_layout_payload("deckLayoutStatusCached", cached_payload)

    response = client.post(
        "/slides/deck/deckLayoutStatusCached/layout/start", json={"lang": "eng"}
    )

    assert response.status_code == 202
    body = response.json()
    assert body["status"] == "completed"
    assert body["builtPages"] == 0
    assert body["totalPages"] == 0
    assert body["message"] == "Layout payload already available."
    assert body["step"] == "complete"
    assert body["lastCompletedStep"] == "layout"
    assert body["updatedAt"]


def test_start_deck_layout_enqueues_background_work(
    client: TestClient, deck_storage: DeckStorage, monkeypatch: pytest.MonkeyPatch
) -> None:
    _seed_image_deck(deck_storage, "deckLayoutAsync")
    submitted: list[tuple[object, tuple[object, ...], dict[str, object]]] = []

    class _Executor:
        def submit(self, fn, *args, **kwargs):
            submitted.append((fn, args, kwargs))
            return None

    monkeypatch.setattr(slides_api, "_LAYOUT_EXECUTOR", _Executor())

    response = client.post(
        "/slides/deck/deckLayoutAsync/layout/start", json={"lang": "eng"}
    )

    assert response.status_code == 202
    body = response.json()
    assert body["status"] == "running"
    assert body["totalPages"] == 1
    assert body["step"] == "layout"
    assert body["startedAt"]
    assert body["updatedAt"]
    assert len(submitted) == 1
    assert submitted[0][1] == ("deckLayoutAsync",)
    assert submitted[0][2] == {"lang": "eng"}


def test_deck_ocr_status_keeps_layout_running_state_when_ocr_payload_exists(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_image_deck(deck_storage, "deckOcrLayoutRunning")
    deck_storage.save_ocr_payload(
        "deckOcrLayoutRunning",
        {
            "deckId": "deckOcrLayoutRunning",
            "lang": "eng",
            "generatedAt": "2026-03-14T00:00:00+00:00",
            "slides": [
                {
                    "slideId": "slide0.html",
                    "slideNumber": 1,
                    "pageNumber": 1,
                    "ocrText": "hello",
                    "lines": [],
                }
            ],
        },
    )
    slides_api._set_ocr_progress(
        "deckOcrLayoutRunning",
        status="running",
        built_pages=0,
        total_pages=1,
        message="Analyzing layout. Slides done: 0 of 1.",
        lang="ita",
        step="layout",
        last_completed_step="ocr",
    )

    response = client.get("/slides/deck/deckOcrLayoutRunning/ocr/status")

    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "running"
    assert body["message"] == "Analyzing layout. Slides done: 0 of 1."
    assert body["lang"] == "ita"
    assert body["step"] == "layout"
    assert body["lastCompletedStep"] == "ocr"
    assert body["updatedAt"]


def test_deck_ocr_sends_failure_email_notification(
    deck_storage: DeckStorage, monkeypatch: pytest.MonkeyPatch
) -> None:
    user = AuthenticatedUser(email="viewer@example.com")
    app = FastAPI()
    app.include_router(slides_api.router)
    app.dependency_overrides[require_authenticated_user] = lambda: user
    client = TestClient(app)
    _seed_image_deck(deck_storage, "deckOcrNotifyFail", owner_email=user.email)

    def _raise_error(*args, **kwargs):
        raise ValueError("ocr failed")

    monkeypatch.setattr(slides_api, "build_deck_ocr_payload", _raise_error)

    notifications: list[tuple[str, dict[str, str]]] = []

    def _capture_notify(step: str, session_data: dict[str, str], notifier=None) -> None:
        notifications.append((step, session_data))

    monkeypatch.setattr(slides_api, "notify_failed", _capture_notify)

    response = client.post("/slides/deck/deckOcrNotifyFail/ocr", json={"lang": "eng"})

    assert response.status_code == 400
    assert len(notifications) == 1
    step, session_data = notifications[0]
    assert step == "deck_processing"
    assert session_data["notify_email"] == user.email
    assert session_data["notify_lang"] == "en"


def test_deck_ocr_unexpected_error_sets_failed_status(
    deck_storage: DeckStorage, monkeypatch: pytest.MonkeyPatch
) -> None:
    user = AuthenticatedUser(email="viewer@example.com")
    app = FastAPI()
    app.include_router(slides_api.router)
    app.dependency_overrides[require_authenticated_user] = lambda: user
    client = TestClient(app)
    _seed_image_deck(deck_storage, "deckOcrUnexpectedFail", owner_email=user.email)

    def _raise_unexpected(*args, **kwargs):
        raise RuntimeError("unexpected ocr failure")

    monkeypatch.setattr(slides_api, "build_deck_ocr_payload", _raise_unexpected)

    notifications: list[tuple[str, dict[str, str]]] = []

    def _capture_notify(step: str, session_data: dict[str, str], notifier=None) -> None:
        notifications.append((step, session_data))

    monkeypatch.setattr(slides_api, "notify_failed", _capture_notify)

    response = client.post(
        "/slides/deck/deckOcrUnexpectedFail/ocr", json={"lang": "eng"}
    )

    assert response.status_code == 500
    assert "unexpected ocr failure" in response.json()["detail"]
    status_response = client.get("/slides/deck/deckOcrUnexpectedFail/ocr/status")
    assert status_response.status_code == 200
    assert status_response.json()["status"] == "failed"
    assert "unexpected ocr failure" in str(status_response.json()["error"])
    assert len(notifications) == 1
    step, session_data = notifications[0]
    assert step == "deck_processing"
    assert session_data["notify_email"] == user.email
    assert session_data["notify_lang"] == "en"


def test_save_deck_prompt_style_round_trips_through_api_responses(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    payload = {
        "promptStyle": "editorial",
        "slides": [
            {
                "id": "slide0.html",
                "titleHtml": "Style",
                "bodyHtml": "<p>Editorial</p>",
                "kind": "normal",
            }
        ],
        "sections": [],
    }

    response = client.post("/slides/deck/deckStyle/save", json=payload)

    assert response.status_code == 200
    assert response.json()["promptStyle"] == "editorial"
    index_payload = json.loads(
        (deck_storage.root / "deckStyle" / "index.json").read_text(encoding="utf-8")
    )
    assert index_payload["promptStyle"] == "editorial"
    get_response = client.get("/slides/deck/deckStyle")
    assert get_response.status_code == 200
    assert get_response.json()["promptStyle"] == "editorial"
    list_response = client.get("/slides/decks")
    assert list_response.status_code == 200
    deck_summary = next(
        deck for deck in list_response.json()["decks"] if deck["deckId"] == "deckStyle"
    )
    assert deck_summary["promptStyle"] == "editorial"


def test_save_deck_prompt_style_change_ignored(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_deck(deck_storage, "deckImmutable", prompt_style="uniform")
    payload = {
        "promptStyle": "editorial",
        "slides": [
            {
                "id": "slide0.html",
                "titleHtml": "First",
                "bodyHtml": "<p>Alpha</p>",
                "kind": "normal",
            },
            {
                "id": "slide1.html",
                "titleHtml": "Second",
                "bodyHtml": "<p>Beta</p>",
                "kind": "normal",
            },
        ],
        "sections": [],
    }

    response = client.post("/slides/deck/deckImmutable/save", json=payload)

    assert response.status_code == 200
    assert response.json()["promptStyle"] == "uniform"
    deck_dir = deck_storage.root / "deckImmutable"
    index_payload = json.loads((deck_dir / "index.json").read_text(encoding="utf-8"))
    assert index_payload["promptStyle"] == "uniform"


def test_save_deck_preserves_full_html(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_deck(deck_storage, "deckFullHtml")
    full_html = (
        "<!DOCTYPE html><html><body>"
        "<div class='slide-container' data-custom-flag='keep-me'>"
        "<div class='slide-body'><img src='image.png' alt='Image' /></div>"
        "</div></body></html>"
    )
    payload = {
        "slides": [
            {
                "id": "slide0.html",
                "titleHtml": "Updated",
                "bodyHtml": "<img src='image.png' alt='Image' />",
                "fullHtml": full_html,
                "kind": "normal",
            }
        ],
        "sections": [],
    }

    response = client.post("/slides/deck/deckFullHtml/save", json=payload)

    assert response.status_code == 200
    reloaded = client.get("/slides/deck/deckFullHtml")
    assert reloaded.status_code == 200
    data = reloaded.json()
    assert 'data-custom-flag="keep-me"' in data["slides"][0]["fullHtml"]


def test_save_deck_persists_converted_section_header(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_deck(deck_storage, "deckConvert")

    payload = {
        "slides": [
            {
                "id": "slide0.html",
                "titleHtml": "",
                "bodyHtml": "",
                "kind": "sectionHeader",
                "sectionId": "Intro",
            },
            {
                "id": "slide1.html",
                "titleHtml": "Second",
                "bodyHtml": "<p>Beta</p>",
                "kind": "normal",
            },
        ],
        "sections": [
            {
                "id": "Intro",
                "title": "Intro",
                "startSlide": "slide1.html",
                "subsections": [],
            }
        ],
    }

    response = client.post("/slides/deck/deckConvert/save", json=payload)

    assert response.status_code == 200
    reloaded = client.get("/slides/deck/deckConvert")
    assert reloaded.status_code == 200
    data = reloaded.json()
    assert data["slides"][0]["kind"] == "sectionHeader"
    assert data["slides"][0]["sectionId"] == "Intro"
    assert data["sections"][0]["startSlide"] == "slide1.html"


def test_import_slide_creates_new_file(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_deck(deck_storage, "deckSrc")
    _seed_deck(deck_storage, "deckDest")

    response = client.post(
        "/slides/deck/deckDest/import",
        json={
            "sourceDeckId": "deckSrc",
            "sourceSlideId": "slide0.html",
            "afterSlideId": "slide0.html",
            "currentOrder": ["slide0.html", "slide1.html"],
        },
    )

    assert response.status_code == 200
    data = response.json()
    assert data["slide"]["id"].startswith("slide")
    assert len(data["order"]) == 3
    new_slide_path = deck_storage.root / "deckDest" / data["slide"]["id"]
    assert new_slide_path.exists()


def test_import_slide_different_prompt_styles_returns_400(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_deck(deck_storage, "deckSrcEditorial", prompt_style="editorial")
    _seed_deck(deck_storage, "deckDestUniform", prompt_style="uniform")

    response = client.post(
        "/slides/deck/deckDestUniform/import",
        json={
            "sourceDeckId": "deckSrcEditorial",
            "sourceSlideId": "slide0.html",
            "afterSlideId": "slide0.html",
            "currentOrder": ["slide0.html", "slide1.html"],
        },
    )

    assert response.status_code == 400
    detail = response.json()["detail"]
    assert "uniform" in detail
    assert "editorial" in detail


def test_import_slide_copies_layout_and_ocr_artifacts(
    client: TestClient,
    deck_storage: DeckStorage,
) -> None:
    _seed_deck(deck_storage, "deckSrc")
    _seed_deck(deck_storage, "deckDest")
    deck_storage.save_layout_payload(
        "deckSrc",
        {
            "deckId": "deckSrc",
            "lang": "eng",
            "generatedAt": "2026-03-18T00:00:00+00:00",
            "slides": [
                {
                    "slideId": "slide0.html",
                    "slideNumber": 1,
                    "pageNumber": 1,
                    "assetPath": "assets/source.png",
                    "blocks": [
                        {
                            "blockId": "block-0",
                            "type": "title",
                            "text": "",
                            "items": [],
                            "bbox": {"x": 10.0, "y": 20.0, "w": 50.0, "h": 20.0},
                            "confidence": 0.5,
                        }
                    ],
                    "titleText": "",
                    "bulletTexts": [],
                    "figureRegions": [],
                }
            ],
        },
    )
    deck_storage.save_ocr_payload(
        "deckSrc",
        {
            "deck_id": "deckSrc",
            "lang": "eng",
            "ocr_strategy": slides_api.OCR_STRATEGY_LAYOUT_GUIDED,
            "generated_at": "2026-03-18T00:00:00+00:00",
            "slides": [
                {
                    "slide_id": "slide0.html",
                    "slide_number": 1,
                    "page_number": 1,
                    "ocr_text": "Imported OCR",
                    "lines": [{"text": "Imported OCR"}],
                    "blocks": [
                        {
                            "block_id": "block-0",
                            "type": "list",
                            "text": "Imported OCR",
                            "items": ["Imported OCR"],
                            "bbox": {"x": 10.0, "y": 20.0, "w": 50.0, "h": 20.0},
                            "confidence": 0.9,
                            "audit_status": "ok",
                        }
                    ],
                    "title_text": "Imported OCR",
                    "bullet_texts": ["Imported OCR"],
                    "figure_regions": [],
                }
            ],
        },
    )
    deck_storage.save_layout_payload(
        "deckDest",
        {
            "deckId": "deckDest",
            "lang": "eng",
            "generatedAt": "2026-03-18T00:00:00+00:00",
            "slides": [],
        },
    )
    deck_storage.save_ocr_payload(
        "deckDest",
        {
            "deck_id": "deckDest",
            "lang": "eng",
            "ocr_strategy": slides_api.OCR_STRATEGY_LAYOUT_GUIDED,
            "generated_at": "2026-03-18T00:00:00+00:00",
            "slides": [],
        },
    )

    response = client.post(
        "/slides/deck/deckDest/import",
        json={
            "sourceDeckId": "deckSrc",
            "sourceSlideId": "slide0.html",
            "afterSlideId": "slide0.html",
            "currentOrder": ["slide0.html", "slide1.html"],
        },
    )

    assert response.status_code == 200
    data = response.json()
    imported_slide_id = data["slide"]["id"]
    target_layout = deck_storage.load_layout_payload("deckDest")
    target_ocr = deck_storage.load_ocr_payload("deckDest")
    target_analysis = deck_storage.load_slide_analysis_payload("deckDest")
    assert target_layout is not None
    assert target_ocr is not None
    assert target_analysis is not None
    imported_layout = next(
        slide
        for slide in target_layout["slides"]
        if slide["slideId"] == imported_slide_id
    )
    imported_ocr = next(
        slide
        for slide in target_ocr["slides"]
        if slide["slide_id"] == imported_slide_id
    )
    imported_analysis = next(
        slide
        for slide in target_analysis["slides"]
        if slide["slideId"] == imported_slide_id
    )
    assert imported_layout["assetPath"] == "assets/source.png"
    assert imported_ocr["blocks"][0]["audit_status"] == "ok"
    assert imported_analysis["blocks"][0]["type"] == "bullet_item"
    assert imported_analysis["blocks"][0]["text"] == "Imported OCR"


def test_concatenate_decks_creates_combined_deck(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_custom_deck(
        deck_storage, "deckOne", ["One-A", "Two-A"], include_sections=True
    )
    _seed_custom_deck(deck_storage, "deckTwo", ["One-B"], include_sections=True)

    response = client.post(
        "/slides/deck/concatenate",
        json={"newDeckId": "combo", "sourceDeckIds": ["deckOne", "deckTwo"]},
    )

    assert response.status_code == 201
    data = response.json()
    assert data["deckId"] == "combo"
    assert data["promptStyle"] == "uniform"
    assert [slide["titleHtml"] for slide in data["slides"]] == [
        "One-A",
        "Two-A",
        "One-B",
    ]
    assert len(data["sections"]) == 2
    assert data["sections"][0]["startSlide"] == data["slides"][0]["id"]
    assert data["sections"][1]["startSlide"] == data["slides"][2]["id"]


def test_concatenate_decks_merges_full_layout_and_ocr_artifacts(
    client: TestClient,
    deck_storage: DeckStorage,
) -> None:
    _seed_custom_deck(deck_storage, "deckOneArtifacts", ["One-A"])
    _seed_custom_deck(deck_storage, "deckTwoArtifacts", ["Two-B"])
    deck_storage.save_layout_payload(
        "deckOneArtifacts",
        {
            "deckId": "deckOneArtifacts",
            "lang": "eng",
            "generatedAt": "2026-03-18T00:00:00+00:00",
            "slides": [
                {
                    "slideId": "slide0.html",
                    "slideNumber": 1,
                    "pageNumber": 1,
                    "assetPath": "assets/one.png",
                    "blocks": [
                        {
                            "blockId": "block-0",
                            "type": "title",
                            "text": "",
                            "items": [],
                            "bbox": {"x": 10.0, "y": 20.0, "w": 60.0, "h": 20.0},
                            "confidence": 0.5,
                        }
                    ],
                    "titleText": "",
                    "bulletTexts": [],
                    "figureRegions": [],
                }
            ],
        },
    )
    deck_storage.save_layout_payload(
        "deckTwoArtifacts",
        {
            "deckId": "deckTwoArtifacts",
            "lang": "eng",
            "generatedAt": "2026-03-18T00:00:00+00:00",
            "slides": [
                {
                    "slideId": "slide0.html",
                    "slideNumber": 1,
                    "pageNumber": 1,
                    "assetPath": "assets/two.png",
                    "blocks": [
                        {
                            "blockId": "block-0",
                            "type": "title",
                            "text": "",
                            "items": [],
                            "bbox": {"x": 10.0, "y": 20.0, "w": 60.0, "h": 20.0},
                            "confidence": 0.5,
                        }
                    ],
                    "titleText": "",
                    "bulletTexts": [],
                    "figureRegions": [],
                }
            ],
        },
    )
    deck_storage.save_ocr_payload(
        "deckOneArtifacts",
        {
            "deck_id": "deckOneArtifacts",
            "lang": "eng",
            "ocr_strategy": slides_api.OCR_STRATEGY_LAYOUT_GUIDED,
            "generated_at": "2026-03-18T00:00:00+00:00",
            "slides": [
                {
                    "slide_id": "slide0.html",
                    "slide_number": 1,
                    "page_number": 1,
                    "ocr_text": "One OCR",
                    "lines": [{"text": "One OCR"}],
                    "blocks": [
                        {
                            "block_id": "block-0",
                            "type": "list",
                            "text": "One OCR",
                            "items": ["One OCR"],
                            "bbox": {"x": 10.0, "y": 20.0, "w": 60.0, "h": 20.0},
                            "confidence": 0.9,
                            "audit_status": "ok",
                        }
                    ],
                    "title_text": "One OCR",
                    "bullet_texts": ["One OCR"],
                    "figure_regions": [],
                }
            ],
        },
    )
    deck_storage.save_ocr_payload(
        "deckTwoArtifacts",
        {
            "deck_id": "deckTwoArtifacts",
            "lang": "eng",
            "ocr_strategy": slides_api.OCR_STRATEGY_LAYOUT_GUIDED,
            "generated_at": "2026-03-18T00:00:00+00:00",
            "slides": [
                {
                    "slide_id": "slide0.html",
                    "slide_number": 1,
                    "page_number": 1,
                    "ocr_text": "Two OCR",
                    "lines": [{"text": "Two OCR"}],
                    "blocks": [
                        {
                            "block_id": "block-0",
                            "type": "list",
                            "text": "Two OCR",
                            "items": ["Two OCR"],
                            "bbox": {"x": 10.0, "y": 20.0, "w": 60.0, "h": 20.0},
                            "confidence": 0.9,
                            "visual_status": "corrected",
                        }
                    ],
                    "title_text": "Two OCR",
                    "bullet_texts": ["Two OCR"],
                    "figure_regions": [],
                }
            ],
        },
    )

    response = client.post(
        "/slides/deck/concatenate",
        json={
            "newDeckId": "comboArtifacts",
            "sourceDeckIds": ["deckOneArtifacts", "deckTwoArtifacts"],
        },
    )

    assert response.status_code == 201
    merged_layout = deck_storage.load_layout_payload("comboArtifacts")
    merged_ocr = deck_storage.load_ocr_payload("comboArtifacts")
    merged_analysis = deck_storage.load_slide_analysis_payload("comboArtifacts")
    assert merged_layout is not None
    assert merged_ocr is not None
    assert merged_analysis is not None
    assert [slide["assetPath"] for slide in merged_layout["slides"]] == [
        "assets/deckOneArtifacts/one.png",
        "assets/deckTwoArtifacts/two.png",
    ]
    assert [slide["ocr_text"] for slide in merged_ocr["slides"]] == [
        "One OCR",
        "Two OCR",
    ]
    assert merged_ocr["slides"][1]["blocks"][0]["visual_status"] == "corrected"
    assert merged_analysis["slides"][0]["blocks"][0]["type"] == "bullet_item"
    assert merged_analysis["slides"][1]["blocks"][0]["text"] == "Two OCR"


def test_concatenate_decks_keeps_assets_separate_when_filenames_collide(
    client: TestClient,
    deck_storage: DeckStorage,
) -> None:
    deck_storage.save_deck(
        Deck(
            deck_id="deckAssetsOne",
            prompt_style="uniform",
            slides=[
                Slide(
                    id="slide0.html",
                    title_html="One",
                    body_html=(
                        "<img src='/slides/deck/deckAssetsOne/assets/shared.png' alt='one' />"
                    ),
                )
            ],
        )
    )
    deck_storage.save_deck(
        Deck(
            deck_id="deckAssetsTwo",
            prompt_style="uniform",
            slides=[
                Slide(
                    id="slide0.html",
                    title_html="Two",
                    body_html=(
                        "<img src='/slides/deck/deckAssetsTwo/assets/shared.png' alt='two' />"
                    ),
                )
            ],
        )
    )
    first_asset = deck_storage.root / "deckAssetsOne" / "assets" / "shared.png"
    first_asset.parent.mkdir(parents=True, exist_ok=True)
    first_asset.write_bytes(b"one")
    second_asset = deck_storage.root / "deckAssetsTwo" / "assets" / "shared.png"
    second_asset.parent.mkdir(parents=True, exist_ok=True)
    second_asset.write_bytes(b"two")

    response = client.post(
        "/slides/deck/concatenate",
        json={
            "newDeckId": "comboAssets",
            "sourceDeckIds": ["deckAssetsOne", "deckAssetsTwo"],
        },
    )

    assert response.status_code == 201
    combined = deck_storage.load_deck("comboAssets")
    assert (
        "/slides/deck/comboAssets/assets/deckAssetsOne/shared.png"
        in combined.slides[0].body_html
    )
    assert (
        "/slides/deck/comboAssets/assets/deckAssetsTwo/shared.png"
        in combined.slides[1].body_html
    )
    assert (
        deck_storage.root / "comboAssets" / "assets" / "deckAssetsOne" / "shared.png"
    ).read_bytes() == b"one"
    assert (
        deck_storage.root / "comboAssets" / "assets" / "deckAssetsTwo" / "shared.png"
    ).read_bytes() == b"two"


def test_concatenate_decks_preserves_shared_custom_template(
    client: TestClient,
    deck_storage: DeckStorage,
) -> None:
    _seed_custom_deck(deck_storage, "deckTemplateOne", ["One"])
    _seed_custom_deck(deck_storage, "deckTemplateTwo", ["Two"])
    template_bytes = Path("src/review_brief/pptx_templates/uniform.pptx").read_bytes()
    for deck_id in ("deckTemplateOne", "deckTemplateTwo"):
        deck_path = deck_storage.root / deck_id
        template_path = deck_path / DECK_PPTX_TEMPLATE_FILENAME
        template_path.write_bytes(template_bytes)
        write_deck_pptx_template_manifest(
            deck_path,
            build_pptx_template_manifest(template_path),
        )

    response = client.post(
        "/slides/deck/concatenate",
        json={
            "newDeckId": "comboTemplate",
            "sourceDeckIds": ["deckTemplateOne", "deckTemplateTwo"],
        },
    )

    assert response.status_code == 201
    combined_path = deck_storage.root / "comboTemplate"
    assert (combined_path / DECK_PPTX_TEMPLATE_FILENAME).exists()
    assert (combined_path / DECK_PPTX_TEMPLATE_MANIFEST_FILENAME).exists()


def test_concatenate_decks_conflict_when_target_exists(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_custom_deck(deck_storage, "combo", ["Existing"])
    _seed_custom_deck(deck_storage, "sourceDeck", ["Source"])

    response = client.post(
        "/slides/deck/concatenate",
        json={"newDeckId": "combo", "sourceDeckIds": ["sourceDeck"]},
    )

    assert response.status_code == 409


def test_concatenate_decks_returns_404_for_missing_source(client: TestClient) -> None:
    response = client.post(
        "/slides/deck/concatenate",
        json={"newDeckId": "newDeck", "sourceDeckIds": ["missingDeck"]},
    )

    assert response.status_code == 404


def test_concatenate_decks_different_prompt_styles_returns_400(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    _seed_custom_deck(deck_storage, "deckUniform", ["Uniform"], prompt_style="uniform")
    _seed_custom_deck(
        deck_storage, "deckEditorial", ["Editorial"], prompt_style="editorial"
    )

    response = client.post(
        "/slides/deck/concatenate",
        json={
            "newDeckId": "comboMismatch",
            "sourceDeckIds": ["deckUniform", "deckEditorial"],
        },
    )

    assert response.status_code == 400
    detail = response.json()["detail"]
    assert "uniform" in detail
    assert "editorial" in detail


def test_concatenate_decks_skips_duplicate_slide_content(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    duplicate_body = "<p>Shared</p>"
    storage = deck_storage
    storage.save_deck(
        Deck(
            deck_id="deck1",
            slides=[
                Slide(id="slide0.html", title_html="Unique A", body_html="<p>A</p>"),
                Slide(
                    id="slide1.html",
                    title_html="Shared",
                    body_html=duplicate_body,
                    full_html=duplicate_body,
                ),
            ],
        )
    )
    storage.save_deck(
        Deck(
            deck_id="deck2",
            slides=[
                Slide(
                    id="slide0.html",
                    title_html="Shared",
                    body_html=duplicate_body,
                    full_html=duplicate_body,
                ),
                Slide(id="slide1.html", title_html="Unique B", body_html="<p>B</p>"),
            ],
        )
    )

    response = client.post(
        "/slides/deck/concatenate",
        json={"newDeckId": "comboDup", "sourceDeckIds": ["deck1", "deck2"]},
    )

    assert response.status_code == 201
    data = response.json()
    titles = [slide["titleHtml"] for slide in data["slides"]]
    assert titles == ["Unique A", "Shared", "Unique B"]


def test_rewrite_asset_sources_for_print_inlines_section_header_styles(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckPrintCss"
    deck_path.mkdir(parents=True, exist_ok=True)
    css_text = (
        ":root{--notebooklm-title-size-px:42px;} .section-header{font-size:42px;}"
    )
    (deck_path / "section_header.css").write_text(css_text, encoding="utf-8")
    html = (
        "<!DOCTYPE html><html><head></head><body>"
        "<section class='section-header'>"
        "<link rel='stylesheet' href='./section_header.css' />"
        "<ol><li>One</li></ol></section></body></html>"
    )

    rewritten = slides_api._rewrite_asset_sources_for_print(
        html, deck_id="deckPrintCss", deck_path=deck_path, prompt_style="uniform"
    )

    assert css_text in rewritten
    assert 'data-inline-stylesheet="section_header.css"' in rewritten
    assert "<link" not in rewritten


def test_rewrite_asset_sources_for_print_inlines_section_header_styles_from_assets_dir(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckPrintCssAssets"
    assets_path = deck_path / "assets"
    assets_path.mkdir(parents=True, exist_ok=True)
    css_text = ".section-header{font-size:40px;}"
    (assets_path / "section_header.css").write_text(css_text, encoding="utf-8")
    html = (
        "<!DOCTYPE html><html><head></head><body>"
        "<section class='section-header'>"
        "<link rel='stylesheet' href='./section_header.css' />"
        "<ol><li>One</li></ol></section></body></html>"
    )

    rewritten = slides_api._rewrite_asset_sources_for_print(
        html, deck_id="deckPrintCssAssets", deck_path=deck_path, prompt_style="uniform"
    )

    assert css_text in rewritten
    assert 'data-inline-stylesheet="section_header.css"' in rewritten


def test_rewrite_asset_sources_for_print_applies_fixed_geometry_to_image_only_slides(
    tmp_path: Path,
) -> None:
    deck_id = "deckImageOnlyPrint"
    deck_path = tmp_path / deck_id
    assets_path = deck_path / "assets"
    assets_path.mkdir(parents=True, exist_ok=True)
    image_name = "chart.png"
    (assets_path / image_name).write_bytes(b"png")
    html = (
        "<!DOCTYPE html><html><head></head><body>"
        "<div class='slide-container'>"
        "<h1 class='slide-title' data-role='title'></h1>"
        "<div class='slide-body'><div>"
        f"<img src='/slides/deck/{deck_id}/assets/{image_name}' alt='Slide image 1' />"
        "</div></div>"
        "</div></body></html>"
    )

    rewritten = slides_api._rewrite_asset_sources_for_print(
        html,
        deck_id=deck_id,
        deck_path=deck_path,
        prompt_style="uniform",
    )

    assert 'data-inline-stylesheet="slides_export_base.css"' in rewritten
    assert 'data-inline-stylesheet="image_slide_export.css"' in rewritten
    assert "slide-image-only-export" in rewritten
    assert 'data-inline-stylesheet="intro_slide.css"' not in rewritten


def test_rewrite_asset_sources_for_print_applies_fixed_geometry_to_html_slides(
    tmp_path: Path,
) -> None:
    deck_path = tmp_path / "deckHtmlPrint"
    deck_path.mkdir(parents=True, exist_ok=True)
    html = (
        "<!DOCTYPE html><html><head></head><body>"
        "<div class='slide-container'>"
        "<h1 class='slide-title' data-role='title'>Title</h1>"
        "<div class='slide-body'><p>Body text</p></div>"
        "</div></body></html>"
    )

    rewritten = slides_api._rewrite_asset_sources_for_print(
        html,
        deck_id="deckHtmlPrint",
        deck_path=deck_path,
        prompt_style="uniform",
    )

    assert 'data-inline-stylesheet="slides_export_base.css"' in rewritten
    assert 'data-inline-stylesheet="intro_slide.css"' in rewritten
    assert 'data-inline-stylesheet="image_slide_export.css"' not in rewritten


def test_print_deck_job_flow(
    client: TestClient, deck_storage: DeckStorage, monkeypatch: pytest.MonkeyPatch
) -> None:
    _seed_deck(deck_storage, "deckPrint")

    class ImmediateExecutor:
        def submit(self, fn, *args, **kwargs):
            fn(*args, **kwargs)

    dummy_executor = ImmediateExecutor()
    monkeypatch.setattr(slides_api, "HAS_PLAYWRIGHT", True)
    monkeypatch.setattr(
        slides_api, "_render_deck_pdf", lambda deck, path: io.BytesIO(b"%PDF-1.4 test")
    )
    monkeypatch.setattr(slides_api, "_PRINT_EXECUTOR", dummy_executor)

    response = client.post("/slides/deck/deckPrint/print")
    assert response.status_code == 202
    data = response.json()
    job_id = data["jobId"]

    status_response = client.get(f"/slides/deck/print/{job_id}")
    assert status_response.status_code == 200
    assert status_response.json()["status"] == "succeeded"

    download = client.get(f"/slides/deck/print/{job_id}/download")
    assert download.status_code == 200
    assert download.headers["content-type"] == "application/pdf"


def test_print_deck_returns_503_when_playwright_missing(
    client: TestClient, deck_storage: DeckStorage, monkeypatch: pytest.MonkeyPatch
) -> None:
    _seed_deck(deck_storage, "deckPrintMissing")
    monkeypatch.setattr(slides_api, "HAS_PLAYWRIGHT", False)

    response = client.post("/slides/deck/deckPrintMissing/print")

    assert response.status_code == 503


def test_render_deck_pdf_uses_raster_path_when_vector_is_disabled(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    deck = Deck(deck_id="deckRasterOnly", prompt_style="uniform", slides=[])
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)

    def _fail_vector(*_args: object, **_kwargs: object) -> io.BytesIO:
        raise AssertionError("Vector PDF renderer must not run when disabled.")

    monkeypatch.setattr(slides_api, "_SLIDES_EXPORT_ENABLE_VECTOR_PDF", False)
    monkeypatch.setattr(slides_api, "_render_deck_pdf_with_playwright", _fail_vector)
    monkeypatch.setattr(
        slides_api,
        "_render_deck_slide_images",
        lambda *_args, **_kwargs: [b"slide-image"],
    )
    monkeypatch.setattr(
        slides_api,
        "_render_deck_pdf_with_reportlab",
        lambda *_args, **_kwargs: io.BytesIO(b"raster-pdf"),
    )

    rendered = slides_api._render_deck_pdf(deck, deck_path)

    assert rendered.getvalue() == b"raster-pdf"


def test_resolve_deck_export_page_size_reads_source_pdf_when_html_metadata_missing(
    tmp_path: Path,
) -> None:
    deck = Deck(
        deck_id="deckSourceSize",
        slides=[Slide(id="slide0.html", title_html="", body_html="<p>No image</p>")],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)

    source_doc = slides_api.fitz.open()
    source_doc.new_page(width=1376.0, height=768.0)
    source_doc.save(deck_path / "source.pdf")
    source_doc.close()

    resolved = slides_api._resolve_deck_export_page_size(deck, deck_path=deck_path)

    assert resolved == pytest.approx((1376.0, 768.0))


def test_resolve_deck_export_page_size_prefers_slide_metadata_over_source_pdf(
    tmp_path: Path,
) -> None:
    body_html = (
        "<div><img src='/slides/deck/deckMeta/assets/slide.png' "
        "data-pdf-crop-w-pt='1376' data-pdf-crop-h-pt='768' /></div>"
    )
    deck = Deck(
        deck_id="deckMeta",
        slides=[
            Slide(
                id="slide0.html",
                title_html="",
                body_html=body_html,
                full_html=(
                    "<!DOCTYPE html><html><body><div class='slide-container'>"
                    "<div class='slide-body'>"
                    f"{body_html}"
                    "</div></div></body></html>"
                ),
            )
        ],
    )
    deck_path = tmp_path / deck.deck_id
    deck_path.mkdir(parents=True, exist_ok=True)

    source_doc = slides_api.fitz.open()
    source_doc.new_page(width=1280.0, height=720.0)
    source_doc.save(deck_path / "source.pdf")
    source_doc.close()

    resolved = slides_api._resolve_deck_export_page_size(deck, deck_path=deck_path)

    assert resolved == pytest.approx((1376.0, 768.0))


def test_ensure_pptx_available_reports_runtime_python(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(
        slides_api,
        "_get_pptx_import_error",
        lambda: "ModuleNotFoundError: No module named 'pptx'",
    )

    with pytest.raises(slides_api.HTTPException) as excinfo:
        slides_api._ensure_pptx_available()

    assert excinfo.value.status_code == 503
    assert "python-pptx" in str(excinfo.value.detail)
    assert sys.executable in str(excinfo.value.detail)


def test_enqueue_export_pptx_uses_runtime_dependency_check(
    client: TestClient,
    deck_storage: DeckStorage,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    _seed_deck(deck_storage, "deckPptxDep")
    monkeypatch.setattr(slides_api, "HAS_PLAYWRIGHT", True)

    def _raise_missing() -> None:
        raise slides_api.HTTPException(
            status_code=503,
            detail="runtime pptx missing",
        )

    monkeypatch.setattr(slides_api, "_ensure_pptx_available", _raise_missing)

    response = client.post("/slides/deck/deckPptxDep/export-pptx")

    assert response.status_code == 503
    assert response.json()["detail"] == "runtime pptx missing"


def test_enqueue_export_pptx_does_not_require_playwright_for_semantic_export(
    client: TestClient,
    deck_storage: DeckStorage,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    _seed_deck(deck_storage, "deckPptxSemantic")
    monkeypatch.setattr(slides_api, "HAS_PLAYWRIGHT", False)
    monkeypatch.setattr(slides_api, "_ensure_pptx_available", lambda: None)

    monkeypatch.setattr(
        slides_api,
        "_start_pptx_job",
        lambda deck_id, *, source="rendered": "pptx-semantic-job",
    )

    response = client.post("/slides/deck/deckPptxSemantic/export-pptx?source=template")

    assert response.status_code == 202
    assert response.json()["jobId"] == "pptx-semantic-job"


@pytest.mark.parametrize(
    ("query_suffix", "expected_source"),
    [
        ("", "rendered"),
        ("?source=rendered", "rendered"),
        ("?source=template", "template"),
    ],
)
def test_enqueue_export_pptx_passes_source_to_job_start(
    client: TestClient,
    deck_storage: DeckStorage,
    monkeypatch: pytest.MonkeyPatch,
    query_suffix: str,
    expected_source: str,
) -> None:
    _seed_deck(deck_storage, "deckPptxSource")
    monkeypatch.setattr(slides_api, "HAS_PLAYWRIGHT", True)
    monkeypatch.setattr(slides_api, "_ensure_pptx_available", lambda: None)
    captured: dict[str, str] = {}

    def _fake_start(deck_id: str, *, source: str = "rendered") -> str:
        captured["deck_id"] = deck_id
        captured["source"] = source
        return "pptx-job-1"

    monkeypatch.setattr(slides_api, "_start_pptx_job", _fake_start)

    response = client.post(f"/slides/deck/deckPptxSource/export-pptx{query_suffix}")

    assert response.status_code == 202
    payload = response.json()
    assert payload["jobId"] == "pptx-job-1"
    assert payload["source"] == expected_source
    assert captured == {"deck_id": "deckPptxSource", "source": expected_source}


def test_render_deck_pptx_from_slide_images_creates_full_slide_images() -> None:
    from pptx import Presentation

    image_buffer = io.BytesIO()
    Image.new("RGB", (1280, 720), "navy").save(image_buffer, format="PNG")

    pptx_buffer = slides_api._render_deck_pptx_from_slide_images(
        [image_buffer.getvalue()],
        page_size_pt=(1280.0, 720.0),
    )

    presentation = Presentation(pptx_buffer)
    slide = presentation.slides[0]
    picture_shapes = [shape for shape in slide.shapes if int(shape.shape_type) == 13]

    assert len(presentation.slides) == 1
    assert len(picture_shapes) == 1
    assert int(picture_shapes[0].left) == 0
    assert int(picture_shapes[0].top) == 0
    assert int(picture_shapes[0].width) == int(presentation.slide_width)
    assert int(picture_shapes[0].height) == int(presentation.slide_height)


def test_run_pptx_job_uses_rendered_export_pipeline_for_rendered_source(
    deck_storage: DeckStorage,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    _seed_deck(deck_storage, "deckPptxRendered")
    job = slides_api.PptxJob(
        job_id="pptx-rendered-job",
        deck_id="deckPptxRendered",
        source="rendered",
        status="pending",
        created_at=datetime.now(UTC),
    )
    monkeypatch.setattr(slides_api, "get_pptx_job_record", lambda _: job)
    status_updates: list[tuple[str, str | None]] = []
    monkeypatch.setattr(
        slides_api,
        "update_pptx_job_status",
        lambda job_id, status, detail=None: status_updates.append((status, detail)),
    )
    output_paths: list[Path] = []
    monkeypatch.setattr(
        slides_api,
        "set_pptx_output_path",
        lambda _job_id, output_path: output_paths.append(output_path),
    )
    pptx_dir = tmp_path / "pptx_jobs"
    pptx_dir.mkdir(parents=True, exist_ok=True)
    monkeypatch.setattr(slides_api, "_PPTX_JOB_DIR", pptx_dir)
    captured: dict[str, object] = {}
    monkeypatch.setattr(
        slides_api,
        "_render_deck_pptx",
        lambda deck, deck_path: captured.update(
            {"deck_id": deck.deck_id, "deck_path": deck_path}
        )
        or io.BytesIO(b"pptx-rendered"),
    )

    def _unexpected_semantic_export(*_args, **_kwargs):
        raise AssertionError("Semantic PPTX export path should not be used.")

    monkeypatch.setattr(
        slides_api, "build_slides_pptx_spec", _unexpected_semantic_export
    )
    monkeypatch.setattr(
        slides_api,
        "render_slides_pptx_from_template",
        _unexpected_semantic_export,
    )

    slides_api._run_pptx_job("pptx-rendered-job")

    assert [status for status, _ in status_updates] == ["running", "succeeded"]
    assert output_paths == [pptx_dir / "pptx-rendered-job.pptx"]
    assert output_paths[0].read_bytes() == b"pptx-rendered"
    assert captured == {
        "deck_id": "deckPptxRendered",
        "deck_path": deck_storage.root / "deckPptxRendered",
    }


def test_run_pptx_job_uses_semantic_export_pipeline(
    deck_storage: DeckStorage,
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    _seed_deck(deck_storage, "deckPptxTemplate")
    job = slides_api.PptxJob(
        job_id="pptx-template-job",
        deck_id="deckPptxTemplate",
        source="template",
        status="pending",
        created_at=datetime.now(UTC),
    )
    monkeypatch.setattr(slides_api, "get_pptx_job_record", lambda _: job)
    status_updates: list[tuple[str, str | None]] = []
    monkeypatch.setattr(
        slides_api,
        "update_pptx_job_status",
        lambda job_id, status, detail=None: status_updates.append((status, detail)),
    )
    output_paths: list[Path] = []
    monkeypatch.setattr(
        slides_api,
        "set_pptx_output_path",
        lambda _job_id, output_path: output_paths.append(output_path),
    )
    pptx_dir = tmp_path / "pptx_jobs"
    pptx_dir.mkdir(parents=True, exist_ok=True)
    monkeypatch.setattr(slides_api, "_PPTX_JOB_DIR", pptx_dir)
    captured: dict[str, object] = {}
    monkeypatch.setattr(
        slides_api,
        "_load_cached_slide_analysis_payload",
        lambda deck_id, _storage: {
            "deckId": deck_id,
            "slides": [],
        },
    )
    monkeypatch.setattr(
        slides_api,
        "build_slides_pptx_spec",
        lambda deck, deck_path, *, slide_analysis: captured.update(
            {
                "deck_id": deck.deck_id,
                "deck_path": deck_path,
                "analysis": slide_analysis,
            }
        )
        or object(),
    )
    monkeypatch.setattr(
        slides_api,
        "write_slides_pptx_spec",
        lambda deck_path, spec: deck_path / "slides_pptx_spec.json",
    )
    monkeypatch.setattr(
        slides_api,
        "render_slides_pptx_from_template",
        lambda _deck_path: io.BytesIO(b"pptx-semantic"),
    )

    def _unexpected_render(*_args, **_kwargs):
        raise AssertionError("Screenshot PPTX export path should not be used.")

    monkeypatch.setattr(slides_api, "_render_deck_slide_images", _unexpected_render)

    slides_api._run_pptx_job("pptx-template-job")

    assert [status for status, _ in status_updates] == ["running", "succeeded"]
    assert output_paths == [pptx_dir / "pptx-template-job.pptx"]
    assert output_paths[0].read_bytes() == b"pptx-semantic"
    assert captured["deck_id"] == "deckPptxTemplate"
    assert captured["deck_path"] == deck_storage.root / "deckPptxTemplate"
    assert captured["analysis"] == {"deckId": "deckPptxTemplate", "slides": []}


def test_merge_adjacent_title_blocks_combines_split_title() -> None:
    blocks = [
        {
            "text": "Blush sales have nearly tripled since 2022,",
            "x": 28.0,
            "y": 42.0,
            "w": 1280.0,
            "h": 72.0,
            "is_title": True,
        },
        {
            "text": "driven by volume growth and premiumized baseline.",
            "x": 32.0,
            "y": 122.0,
            "w": 1240.0,
            "h": 68.0,
            "is_title": True,
        },
        {
            "text": "• Monthly sales rose from $16.4M to $47.4M.",
            "x": 44.0,
            "y": 248.0,
            "w": 760.0,
            "h": 62.0,
            "is_title": False,
        },
    ]

    merged = slides_api._merge_adjacent_title_blocks(
        blocks=blocks,
        image_size=(1920, 1080),
    )

    title_blocks = [item for item in merged if bool(item.get("is_title"))]
    assert len(title_blocks) == 1
    assert title_blocks[0]["text"] == (
        "Blush sales have nearly tripled since 2022,\n"
        "driven by volume growth and premiumized baseline."
    )


def test_run_pptx_job_applies_post_render_compare_when_enabled(
    deck_storage: DeckStorage,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    deck = Deck(
        deck_id="deckPptxRepair",
        slides=[Slide(id="slide0.html", title_html="Title", body_html="<p>Body</p>")],
    )
    deck_storage.save_deck(deck)
    job = slides_api.PptxJob(
        job_id="pptx-repair-job",
        deck_id=deck.deck_id,
        source="template",
        status="queued",
        created_at=datetime.now(UTC),
        output_path=None,
    )
    monkeypatch.setattr(slides_api, "get_pptx_job_record", lambda _: job)
    status_updates: list[tuple[str, str | None]] = []
    monkeypatch.setattr(
        slides_api,
        "update_pptx_job_status",
        lambda _job_id, status, detail=None: status_updates.append((status, detail)),
    )
    output_paths: list[Path] = []
    monkeypatch.setattr(
        slides_api,
        "set_pptx_output_path",
        lambda _job_id, output_path: output_paths.append(output_path),
    )
    pptx_dir = deck_storage.root / "pptx_jobs"
    pptx_dir.mkdir(parents=True, exist_ok=True)
    monkeypatch.setattr(slides_api, "_PPTX_JOB_DIR", pptx_dir)
    monkeypatch.setattr(
        slides_api,
        "_load_cached_slide_analysis_payload",
        lambda deck_id, _storage: {"deckId": deck_id, "slides": []},
    )
    fake_spec = object()
    monkeypatch.setattr(
        slides_api,
        "build_slides_pptx_spec",
        lambda *_args, **_kwargs: fake_spec,
    )
    monkeypatch.setattr(
        slides_api,
        "write_slides_pptx_spec",
        lambda deck_path, spec: deck_path / "slides_pptx_spec.json",
    )
    monkeypatch.setattr(
        slides_api,
        "render_slides_pptx_from_template",
        lambda _deck_path: io.BytesIO(b"initial-pptx"),
    )
    compare_calls: list[tuple[Path, object, str]] = []
    monkeypatch.setattr(
        slides_api,
        "apply_post_render_compare_loop",
        lambda *, deck, deck_path, spec, pptx_path, job_id: compare_calls.append(
            (pptx_path, spec, job_id)
        )
        or (spec, {"jobId": job_id}),
    )
    monkeypatch.setattr(slides_api, "_SLIDES_PPTX_POST_RENDER_COMPARE_ENABLED", True)

    slides_api._run_pptx_job("pptx-repair-job")

    assert [status for status, _ in status_updates] == ["running", "succeeded"]
    assert output_paths == [pptx_dir / "pptx-repair-job.pptx"]
    assert output_paths[0].read_bytes() == b"initial-pptx"
    assert compare_calls == [
        (pptx_dir / "pptx-repair-job.pptx.tmp", fake_spec, "pptx-repair-job")
    ]


def test_launch_playwright_chromium_falls_back_to_channel(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    missing_executable_error = (
        "Executable doesn't exist at /tmp/chromium\n"
        "Please run the following command to download new browsers:\n"
        "playwright install"
    )

    class FakeChromium:
        def __init__(self) -> None:
            self.calls: list[dict[str, object]] = []

        def launch(self, **kwargs) -> str:
            self.calls.append(dict(kwargs))
            if not kwargs:
                raise RuntimeError(missing_executable_error)
            if kwargs.get("channel") == "chrome":
                return "browser-from-channel"
            raise RuntimeError("channel unavailable")

    monkeypatch.setattr(slides_api, "PlaywrightError", RuntimeError)
    monkeypatch.setattr(
        slides_api, "_discover_playwright_browser_executables", lambda: []
    )

    fake_playwright = types.SimpleNamespace(chromium=FakeChromium())

    launched = slides_api._launch_playwright_chromium(fake_playwright)

    assert launched == "browser-from-channel"
    assert fake_playwright.chromium.calls == [{}, {"channel": "chrome"}]


def test_launch_playwright_chromium_falls_back_to_cached_executable(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    missing_executable_error = (
        "Executable doesn't exist at /tmp/chromium\n"
        "Please run the following command to download new browsers:\n"
        "playwright install"
    )

    class FallbackChromium:
        def __init__(self) -> None:
            self.calls: list[dict[str, object]] = []

        def launch(self, **kwargs) -> str:
            self.calls.append(dict(kwargs))
            if not kwargs:
                raise RuntimeError(missing_executable_error)
            if kwargs.get("channel") in {"chrome", "msedge"}:
                raise RuntimeError("channel unavailable")
            if kwargs.get("executable_path") == "/tmp/fake-chrome":
                return "browser-from-executable"
            raise RuntimeError("unexpected launch arguments")

    monkeypatch.setattr(slides_api, "PlaywrightError", RuntimeError)
    monkeypatch.setattr(
        slides_api,
        "_discover_playwright_browser_executables",
        lambda: [Path("/tmp/fake-chrome")],
    )

    fake_playwright = types.SimpleNamespace(chromium=FallbackChromium())

    launched = slides_api._launch_playwright_chromium(fake_playwright)

    assert launched == "browser-from-executable"
    assert fake_playwright.chromium.calls == [
        {},
        {"channel": "chrome"},
        {"channel": "msedge"},
        {"executable_path": "/tmp/fake-chrome"},
    ]


def test_launch_playwright_chromium_uses_slides_executable_override(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    missing_executable_error = (
        "Executable doesn't exist at /tmp/chromium\n"
        "Please run the following command to download new browsers:\n"
        "playwright install"
    )
    executable = tmp_path / "slides-chrome"
    executable.write_text("", encoding="utf-8")

    class OverrideChromium:
        def __init__(self) -> None:
            self.calls: list[dict[str, object]] = []

        def launch(self, **kwargs) -> str:
            self.calls.append(dict(kwargs))
            if not kwargs:
                raise RuntimeError(missing_executable_error)
            if kwargs.get("executable_path") == str(executable):
                return "browser-from-override"
            raise RuntimeError("unexpected launch arguments")

    monkeypatch.setattr(slides_api, "PlaywrightError", RuntimeError)
    monkeypatch.setattr(
        slides_api, "_discover_playwright_browser_executables", lambda: []
    )
    monkeypatch.setenv("SLIDES_PLAYWRIGHT_EXECUTABLE_PATH", str(executable))

    fake_playwright = types.SimpleNamespace(chromium=OverrideChromium())

    launched = slides_api._launch_playwright_chromium(fake_playwright)

    assert launched == "browser-from-override"
    assert fake_playwright.chromium.calls == [
        {},
        {"executable_path": str(executable)},
    ]


def test_launch_playwright_chromium_raises_when_no_fallback_available(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    missing_executable_error = (
        "Executable doesn't exist at /tmp/chromium\n"
        "Please run the following command to download new browsers:\n"
        "playwright install"
    )

    class AlwaysFailChromium:
        def launch(self, **_kwargs) -> str:
            raise RuntimeError(missing_executable_error)

    monkeypatch.setattr(slides_api, "PlaywrightError", RuntimeError)
    monkeypatch.setattr(
        slides_api, "_discover_playwright_browser_executables", lambda: []
    )

    fake_playwright = types.SimpleNamespace(chromium=AlwaysFailChromium())

    with pytest.raises(RuntimeError, match="browser executable is unavailable"):
        slides_api._launch_playwright_chromium(fake_playwright)


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_pdf_rejects_non_notebooklm_dimensions(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    document = slides_api.fitz.open()
    document.new_page(width=1278, height=720)
    payload = document.tobytes()
    document.close()

    response = client.post(
        "/slides/deck/upload-pdf",
        data={"deckId": "badPdfDims", "promptStyle": "uniform"},
        files={"file": ("bad.pdf", payload, "application/pdf")},
    )

    assert response.status_code == 400
    detail = response.json()["detail"]
    assert "Uploaded PDF must use NotebookLM slide pages" in detail
    assert "Page 1 is 1278.0x720.0." in detail
    assert not (deck_storage.root / "badPdfDims").exists()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_pdf_accepts_notebooklm_1376x768_dimensions(
    client: TestClient,
    deck_storage: DeckStorage,
    stub_uploaded_deck_processing: None,
) -> None:
    document = slides_api.fitz.open()
    document.new_page(width=1376, height=768)
    document.new_page(width=1376, height=768)
    payload = document.tobytes()
    document.close()

    response = client.post(
        "/slides/deck/upload-pdf",
        data={"deckId": "goodPdfDims1376", "promptStyle": "uniform"},
        files={"file": ("good.pdf", payload, "application/pdf")},
    )

    assert response.status_code == 201
    body = response.json()
    assert body["deckId"] == "goodPdfDims1376"
    assert body["hasLayout"] is True
    assert len(body["slides"]) == 2
    deck_path = deck_storage.root / "goodPdfDims1376"
    assert deck_path.exists()
    assert (deck_path / "layout.json").exists()
    assert (deck_path / "ocr.json").exists()
    assert (deck_path / "slide_analysis.json").exists()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_pdf_enqueues_server_side_processing(
    client: TestClient,
    deck_storage: DeckStorage,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    submitted: list[tuple[object, tuple[object, ...], dict[str, object]]] = []

    class _Executor:
        def submit(self, fn, *args, **kwargs):
            submitted.append((fn, args, kwargs))
            return None

    monkeypatch.setattr(slides_api, "_OCR_EXECUTOR", _Executor())

    document = slides_api.fitz.open()
    document.new_page(width=1376, height=768)
    payload = document.tobytes()
    document.close()

    response = client.post(
        "/slides/deck/upload-pdf",
        data={"deckId": "queuedPdfDeck", "promptStyle": "uniform"},
        files={"file": ("queued.pdf", payload, "application/pdf")},
    )

    assert response.status_code == 201
    body = response.json()
    assert body["deckId"] == "queuedPdfDeck"
    assert body["hasLayout"] is False
    deck_path = deck_storage.root / "queuedPdfDeck"
    assert deck_path.exists()
    assert not (deck_path / "layout.json").exists()
    assert not (deck_path / "ocr.json").exists()
    assert not (deck_path / "slide_analysis.json").exists()
    assert len(submitted) == 1
    assert submitted[0][1] == ("queuedPdfDeck",)
    assert submitted[0][2] == {"lang": "eng"}

    status_response = client.get("/slides/deck/queuedPdfDeck/ocr/status")

    assert status_response.status_code == 200
    status_payload = status_response.json()
    assert status_payload["status"] == "running"
    assert status_payload["builtPages"] == 0
    assert status_payload["totalPages"] == 1
    assert status_payload["step"] == "layout"
    assert status_payload["lang"] == "eng"
    assert status_payload["updatedAt"]


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_pdf_can_skip_server_side_processing(
    client: TestClient,
    deck_storage: DeckStorage,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    submitted: list[tuple[object, tuple[object, ...], dict[str, object]]] = []

    class _Executor:
        def submit(self, fn, *args, **kwargs):
            submitted.append((fn, args, kwargs))
            return None

    monkeypatch.setattr(slides_api, "_OCR_EXECUTOR", _Executor())

    document = slides_api.fitz.open()
    document.new_page(width=1376, height=768)
    payload = document.tobytes()
    document.close()

    response = client.post(
        "/slides/deck/upload-pdf",
        data={
            "deckId": "plainUploadedPdf",
            "promptStyle": "uniform",
            "runOcr": "false",
        },
        files={"file": ("plain.pdf", payload, "application/pdf")},
    )

    assert response.status_code == 201
    body = response.json()
    assert body["deckId"] == "plainUploadedPdf"
    assert body["hasLayout"] is False
    deck_path = deck_storage.root / "plainUploadedPdf"
    assert deck_path.exists()
    assert not (deck_path / "layout.json").exists()
    assert not (deck_path / "ocr.json").exists()
    assert not (deck_path / "slide_analysis.json").exists()
    assert submitted == []

    status_response = client.get("/slides/deck/plainUploadedPdf/ocr/status")

    assert status_response.status_code == 200
    status_payload = status_response.json()
    assert status_payload["status"] == "idle"
    assert status_payload["message"] == "OCR not started."
    assert status_payload["updatedAt"]


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_pdf_rejects_1280x720_dimensions(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    document = slides_api.fitz.open()
    document.new_page(width=1280, height=720)
    payload = document.tobytes()
    document.close()

    response = client.post(
        "/slides/deck/upload-pdf",
        data={"deckId": "badPdfDims1280", "promptStyle": "uniform"},
        files={"file": ("good1280.pdf", payload, "application/pdf")},
    )

    assert response.status_code == 400
    detail = response.json()["detail"]
    assert "Uploaded PDF must use NotebookLM slide pages" in detail
    assert "Page 1 is 1280.0x720.0." in detail
    assert not (deck_storage.root / "badPdfDims1280").exists()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_pdf_rejects_mixed_notebooklm_dimensions(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    document = slides_api.fitz.open()
    document.new_page(width=1376, height=768)
    document.new_page(width=1280, height=720)
    payload = document.tobytes()
    document.close()

    response = client.post(
        "/slides/deck/upload-pdf",
        data={"deckId": "mixedPdfDims", "promptStyle": "uniform"},
        files={"file": ("mixed.pdf", payload, "application/pdf")},
    )

    assert response.status_code == 400
    detail = response.json()["detail"]
    assert "Uploaded PDF must use NotebookLM slide pages" in detail
    assert "Page 2 is 1280.0x720.0." in detail
    assert not (deck_storage.root / "mixedPdfDims").exists()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for PPTX template uploads",
)
def test_upload_pptx_template_lists_and_sets_default(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    template_bytes = Path("src/review_brief/pptx_templates/uniform.pptx").read_bytes()

    upload_response = client.post(
        "/slides/pptx-templates/upload",
        data={"setDefault": "true"},
        files={
            "file": (
                "Corporate Template.pptx",
                template_bytes,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    assert upload_response.status_code == 201
    uploaded = upload_response.json()
    assert uploaded["name"] == "Corporate Template"
    assert uploaded["isDefault"] is True

    list_response = client.get("/slides/pptx-templates")

    assert list_response.status_code == 200
    payload = list_response.json()
    assert payload["defaultTemplateId"] == uploaded["templateId"]
    assert payload["templates"][0]["templateId"] == uploaded["templateId"]
    assert payload["templates"][0]["isDefault"] is True


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for PPTX template uploads",
)
def test_upload_pdf_applies_selected_saved_pptx_template(
    client: TestClient,
    deck_storage: DeckStorage,
    stub_uploaded_deck_processing: None,
) -> None:
    template_bytes = Path("src/review_brief/pptx_templates/uniform.pptx").read_bytes()
    template_response = client.post(
        "/slides/pptx-templates/upload",
        data={"setDefault": "false"},
        files={
            "file": (
                "Selected Template.pptx",
                template_bytes,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    template_id = template_response.json()["templateId"]
    document = slides_api.fitz.open()
    document.new_page(width=1376, height=768)
    payload = document.tobytes()
    document.close()

    response = client.post(
        "/slides/deck/upload-pdf",
        data={"deckId": "templatedDeck", "pptxTemplateId": template_id},
        files={"file": ("deck.pdf", payload, "application/pdf")},
    )

    assert response.status_code == 201
    deck_path = deck_storage.root / "templatedDeck"
    assert (deck_path / "pptx_template.pptx").exists()
    assert (deck_path / "pptx_template_manifest.json").exists()
    assert (deck_path / "layout.json").exists()
    assert (deck_path / "ocr.json").exists()
    assert (deck_path / "slide_analysis.json").exists()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for PPTX template uploads",
)
def test_upload_pdf_uses_default_saved_pptx_template_when_no_choice_provided(
    client: TestClient,
    deck_storage: DeckStorage,
    stub_uploaded_deck_processing: None,
) -> None:
    template_bytes = Path("src/review_brief/pptx_templates/uniform.pptx").read_bytes()
    client.post(
        "/slides/pptx-templates/upload",
        data={"setDefault": "true"},
        files={
            "file": (
                "Default Template.pptx",
                template_bytes,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    document = slides_api.fitz.open()
    document.new_page(width=1376, height=768)
    payload = document.tobytes()
    document.close()

    response = client.post(
        "/slides/deck/upload-pdf",
        data={"deckId": "defaultTemplatedDeck"},
        files={"file": ("deck.pdf", payload, "application/pdf")},
    )

    assert response.status_code == 201
    deck_path = deck_storage.root / "defaultTemplatedDeck"
    assert (deck_path / "pptx_template.pptx").exists()
    assert (deck_path / "pptx_template_manifest.json").exists()
    assert (deck_path / "layout.json").exists()
    assert (deck_path / "ocr.json").exists()
    assert (deck_path / "slide_analysis.json").exists()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_deck_accepts_zip_archive(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    archive = io.BytesIO()
    with zipfile.ZipFile(archive, "w") as zf:
        zf.writestr(
            "index.html", "<html><body><a href='slide0.html'>slide</a></body></html>"
        )
        zf.writestr(
            "slide0.html",
            "<div class='slide-container'><h1 data-role='title'>Zip</h1><div class='slide-body'><p>Body</p></div></div>",
        )
    archive.seek(0)

    response = client.post(
        "/slides/deck/upload",
        data={"deckId": "uploadedZip"},
        files={"files": ("deck.zip", archive, "application/zip")},
    )

    assert response.status_code == 410
    assert response.json()["detail"] == (
        "HTML deck uploads are no longer supported. Use the PDF/image upload instead."
    )
    assert not (deck_storage.root / "uploadedZip").exists()


def test_slides_page_includes_preview_configuration(
    site_client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    slides_api.get_preview_bundle_config.cache_clear()
    try:
        monkeypatch.setenv(
            "SLIDES_PREVIEW_STYLE_BUNDLES",
            "https://cdn.example.com/base.css https://cdn.example.com/theme.css",
        )
        monkeypatch.setenv(
            "SLIDES_PREVIEW_SCRIPT_BUNDLES",
            "https://cdn.example.com/app.js",
        )
        monkeypatch.setenv(
            "SLIDES_PREVIEW_ORIGIN_ALLOWLIST",
            "https://cdn.example.com, https://cdn2.example.com",
        )
        response = site_client.get("/slides/page")
        assert response.status_code == 200
        html = response.text
        marker = '<script id="slidesEditorBootstrap" type="application/json">'
        assert marker in html
        payload_raw = html.split(marker, 1)[1].split("</script>", 1)[0]
        payload = json.loads(payload_raw.strip())
        assert payload["preview_styles"] == [
            "https://cdn.example.com/base.css",
            "https://cdn.example.com/theme.css",
        ]
        assert payload["preview_scripts"] == ["https://cdn.example.com/app.js"]
        assert payload["preview_allowlist"] == [
            "https://cdn.example.com",
            "https://cdn2.example.com",
        ]
    finally:
        slides_api.get_preview_bundle_config.cache_clear()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_deck_with_file_list(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    index_html = "<html><body><a href='slide0.html'>slide</a></body></html>"
    slide_html = "<div class='slide-container'><h1 data-role='title'>List</h1><div class='slide-body'><p>Body</p></div></div>"

    response = client.post(
        "/slides/deck/upload",
        data={"deckId": "uploadedFiles"},
        files=[
            ("files", ("index.html", index_html.encode("utf-8"), "text/html")),
            ("files", ("slide0.html", slide_html.encode("utf-8"), "text/html")),
        ],
    )

    assert response.status_code == 410
    assert response.json()["detail"] == (
        "HTML deck uploads are no longer supported. Use the PDF/image upload instead."
    )
    assert not (deck_storage.root / "uploadedFiles").exists()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_deck_accepts_nested_directory_zip(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    archive = io.BytesIO()
    with zipfile.ZipFile(archive, "w") as zf:
        zf.writestr("nested/index.html", INDEX_TEMPLATE)
        zf.writestr(
            "nested/slides/slide0.html",
            SLIDE_TEMPLATE.format(title="Nested Zip"),
        )
    archive.seek(0)

    response = client.post(
        "/slides/deck/upload",
        data={"deckId": "nestedZip"},
        files={"files": ("deck.zip", archive, "application/zip")},
    )

    assert response.status_code == 410
    assert response.json()["detail"] == (
        "HTML deck uploads are no longer supported. Use the PDF/image upload instead."
    )
    assert not (deck_storage.root / "nestedZip").exists()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_deck_promotes_folder_upload(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    response = client.post(
        "/slides/deck/upload",
        data={"deckId": "folderUpload"},
        files=[
            (
                "files",
                ("Folder/index.html", INDEX_TEMPLATE.encode("utf-8"), "text/html"),
            ),
            (
                "files",
                (
                    "Folder/slides/slide0.html",
                    SLIDE_TEMPLATE.format(title="Folder").encode("utf-8"),
                    "text/html",
                ),
            ),
        ],
    )

    assert response.status_code == 410
    assert response.json()["detail"] == (
        "HTML deck uploads are no longer supported. Use the PDF/image upload instead."
    )
    assert not (deck_storage.root / "folderUpload").exists()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_deck_reports_unexpected_error(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    response = client.post(
        "/slides/deck/upload",
        data={"deckId": "boomDeck"},
        files=[
            ("files", ("index.html", INDEX_TEMPLATE.encode("utf-8"), "text/html")),
            (
                "files",
                ("slides/slide0.html", SLIDE_TEMPLATE.encode("utf-8"), "text/html"),
            ),
        ],
    )

    assert response.status_code == 410
    assert response.json()["detail"] == (
        "HTML deck uploads are no longer supported. Use the PDF/image upload instead."
    )
    assert not (deck_storage.root / "boomDeck").exists()


@pytest.mark.skipif(
    not getattr(slides_api, "HAS_MULTIPART", False),
    reason="python-multipart is required for deck uploads",
)
def test_upload_deck_missing_index_returns_error(
    client: TestClient, deck_storage: DeckStorage
) -> None:
    response = client.post(
        "/slides/deck/upload",
        data={"deckId": "invalidDeck"},
        files=[
            (
                "files",
                (
                    "slide0.html",
                    "<div class='slide-container'><div class='slide-body'></div></div>".encode(
                        "utf-8"
                    ),
                    "text/html",
                ),
            )
        ],
    )

    assert response.status_code == 410
    assert response.json()["detail"] == (
        "HTML deck uploads are no longer supported. Use the PDF/image upload instead."
    )
    assert not (deck_storage.root / "invalidDeck").exists()


def test_tagging_stub_runs_two_stage_flow(deck_storage: DeckStorage) -> None:
    from src.slides import tagging

    deck_id = "deckTag"
    slide_body = """
    <h2>Next steps</h2>
    <ul>
      <li>Launch pricing floor</li>
    </ul>
    <div class="metric">100</div>
    """
    full_html = (
        "<div class='slide-container'>" "<h1>Title</h1>" f"{slide_body}" "</div>"
    )
    deck_storage.save_deck(
        Deck(
            deck_id=deck_id,
            slides=[
                Slide(
                    id="slide0.html",
                    title_html="Title",
                    body_html=slide_body,
                    full_html=full_html,
                )
            ],
        )
    )

    slide = deck_storage.load_deck(deck_id).slides[0]
    stamped = tagging.stamp_slide(slide)
    patch = {
        "slide_topic": "pricing",
        "slide_kind": "insight",
        "metrics": [
            {
                "index": 0,
                "label": "Revenue",
                "unit": "M USD",
                "year": "2024",
                "canonical_slide": stamped.slide_id,
            }
        ],
        "recommendations": [
            {
                "index": 0,
                "priority": "1",
                "owner": "Ops",
                "relates_to": stamped.slide_id,
                "canonical_slide": stamped.slide_id,
            }
        ],
    }
    enriched = tagging.apply_enrichment_patch(stamped, patch)
    summary = tagging.summarize_tagged_slides([enriched])

    assert summary.metric_duplicates["slide0.html"] == ["slide0.html"]
    assert 'data-slide-topic="pricing"' in enriched.html
    assert 'data-metric-canonical="slide0.html"' in enriched.html
    assert 'data-relates-to="slide0.html"' in enriched.html


def test_extract_slide_thumbnail_uses_slide_image() -> None:
    slide = Slide(
        id="slide0.html",
        title_html="Title",
        body_html="<img src='image.png' alt='Image' />",
        full_html=(
            "<!DOCTYPE html><html><body>"
            "<div class='slide-container'><div class='slide-body'>"
            "<img src='image.png' alt='Image' />"
            "</div></div></body></html>"
        ),
    )

    thumbnail = slides_api._extract_slide_thumbnail(slide)

    assert thumbnail == '<img src="image.png" alt="Image" />'


def test_cover_notebooklm_logo_covers_to_bottom_right_corner() -> None:
    width, height = 1280, 720
    image = Image.new("RGB", (width, height), (0, 0, 0))
    image.putpixel((width - 201, height - 32), (255, 0, 0))

    slides_api._cover_notebooklm_logo(image)

    assert image.getpixel((width - 1, height - 1)) == (0, 0, 0)
    assert image.getpixel((width - 200, height - 30)) == (0, 0, 0)
    assert image.getpixel((width - 201, height - 31)) == (0, 0, 0)
    assert image.getpixel((width - 201, height - 32)) == (255, 0, 0)
    assert image.getpixel((0, 0)) == (0, 0, 0)


def test_render_image_deck_applies_notebooklm_logo_cover(
    deck_storage: DeckStorage,
) -> None:
    deck_id = "deckImage"
    deck_path = deck_storage.root / deck_id
    deck_path.mkdir(parents=True, exist_ok=True)

    width, height = 1280, 720
    source = Image.new("RGB", (width, height), (0, 0, 0))
    source.putpixel((width - 201, height - 32), (255, 0, 0))
    buffer = io.BytesIO()
    source.save(buffer, format="PNG")

    slides_api._render_image_deck(
        deck_id,
        deck_path,
        buffer.getvalue(),
        deck_storage,
        "slide.png",
        prompt_style=slides_api.resolve_prompt_style_key(None),
        owner_email=None,
        shared_with=[],
    )

    assets = list((deck_path / "assets").iterdir())
    assert len(assets) == 1
    with Image.open(assets[0]) as rendered:
        assert rendered.getpixel((width - 1, height - 1)) == (0, 0, 0)
        assert rendered.getpixel((width - 200, height - 30)) == (0, 0, 0)
        assert rendered.getpixel((width - 201, height - 31)) == (0, 0, 0)
        assert rendered.getpixel((width - 201, height - 32)) == (255, 0, 0)
        assert rendered.getpixel((0, 0)) == (0, 0, 0)


def test_notebook_chart_remap_routes_removed(client: TestClient) -> None:
    response = client.post("/slides/notebook-remap/jobs")
    assert response.status_code == 404

    response = client.get("/slides/notebook-remap/jobs/some-job")
    assert response.status_code == 404

    response = client.get("/slides/notebook-remap/jobs/some-job/unmatched/download")
    assert response.status_code == 404


def test_notebook_chart_remap_unmatched_download_route_removed(
    client: TestClient,
    deck_storage: DeckStorage,
) -> None:
    deck_id = "deckA"
    _seed_deck(deck_storage, deck_id)
    response = client.get(f"/slides/deck/{deck_id}/chart-remap/unmatched/download")
    assert response.status_code == 404


def test_get_or_create_ocr_payload_force_rebuild_ignores_cached_payload(
    deck_storage: DeckStorage,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    deck_id = "deckForceRebuild"
    deck = Deck(
        deck_id=deck_id,
        prompt_style="uniform",
        slides=[Slide(id="slide0.html", title_html="A", body_html="<p>A</p>")],
    )
    deck_storage.save_deck(deck)
    deck_storage.save_ocr_payload(
        deck_id,
        {
            "deck_id": deck_id,
            "lang": "eng",
            "generated_at": "2026-02-20T00:00:00+00:00",
            "slides": [{"slide_id": "cached-slide"}],
        },
    )

    observed: dict[str, object] = {}
    fresh_payload = {
        "deck_id": deck_id,
        "lang": "eng",
        "generated_at": "2026-02-24T00:00:00+00:00",
        "slides": [{"slide_id": "fresh-slide"}],
    }

    def _fake_ensure(
        deck_arg: Deck,
        deck_path_arg: Path,
        *,
        lang: str = "eng",
        include_bboxes: bool = True,
        cached_payload: dict[str, object] | None = None,
        layout_payload: dict[str, object] | None = None,
        progress_callback: Callable[[int, int], None] | None = None,
        event_callback: Callable[[str, dict[str, object] | None], None] | None = None,
    ) -> dict[str, object]:
        observed["deck_id"] = deck_arg.deck_id
        observed["deck_path"] = str(deck_path_arg)
        observed["lang"] = lang
        observed["include_bboxes"] = include_bboxes
        observed["cached_payload"] = cached_payload
        observed["layout_payload"] = layout_payload
        observed["progress_callback"] = progress_callback
        observed["event_callback"] = event_callback
        return fresh_payload

    monkeypatch.setattr(slides_api, "ensure_deck_ocr_payload", _fake_ensure)

    result = slides_api._get_or_create_ocr_payload(
        deck,
        deck_storage,
        force_rebuild=True,
    )

    assert result == fresh_payload
    assert observed["deck_id"] == deck_id
    assert observed["cached_payload"] is None
    assert observed["layout_payload"] is None
    assert deck_storage.load_ocr_payload(deck_id) == fresh_payload


def test_get_or_create_ocr_payload_uses_cached_when_completion_crashes(
    deck_storage: DeckStorage,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    deck_id = "deckCachedFallback"
    slide_id = "slide0.html"
    deck = Deck(
        deck_id=deck_id,
        prompt_style="uniform",
        slides=[Slide(id=slide_id, title_html="A", body_html="<p>A</p>")],
    )
    deck_storage.save_deck(deck)
    cached_payload = {
        "deck_id": deck_id,
        "lang": "eng",
        "ocr_strategy": "layout_guided_text_region_assignment_v6",
        "generated_at": "2026-03-01T00:00:00+00:00",
        "slides": [
            {
                "slide_id": slide_id,
                "slide_number": 1,
                "page_number": 1,
                "ocr_text": "cached line",
                "lines": [{"text": "cached line"}],
                "blocks": [],
            }
        ],
    }
    deck_storage.save_ocr_payload(deck_id, cached_payload)

    def _raise_completion(
        deck_arg: Deck,
        deck_path_arg: Path,
        *,
        lang: str = "eng",
        include_bboxes: bool = True,
        cached_payload: dict[str, object] | None = None,
        layout_payload: dict[str, object] | None = None,
    ) -> dict[str, object]:
        assert deck_arg.deck_id == deck_id
        assert deck_path_arg == deck_storage.root / deck_id
        assert lang == "eng"
        assert include_bboxes is True
        assert cached_payload is not None
        assert layout_payload is None
        raise RuntimeError("5 != 6 for key class_ids!")

    monkeypatch.setattr(slides_api, "ensure_deck_ocr_payload", _raise_completion)

    result = slides_api._get_or_create_ocr_payload(deck, deck_storage)

    assert result is not None
    assert result["deck_id"] == deck_id
    slides = result.get("slides")
    assert isinstance(slides, list)
    assert slides and slides[0].get("slide_id") == slide_id
    assert deck_storage.load_ocr_payload(deck_id) == result
