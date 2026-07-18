from __future__ import annotations

import base64
import importlib.util
import json
import os
import subprocess
import sys
import zlib
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from urllib.parse import parse_qs, urlsplit
from zipfile import ZIP_DEFLATED, ZipFile

import pytest

ROOT = Path(__file__).resolve().parents[2]
PLUGIN_ROOT = ROOT / "plugins" / "clara"
SCRIPTS_DIR = PLUGIN_ROOT / "scripts"


@pytest.mark.parametrize(
    ("skill_name", "display_name", "old_skill_name"),
    [
        ("beautify-deck", "Beautify Deck", "clara-beautify-deck"),
        ("claim-basis-map", "Claim Basis Map", "clara-claim-basis-map"),
        ("interview", "Interview", "clara-interview"),
        ("transcribe", "Transcribe", "clara-transcribe"),
        ("deck-correction", "Deck Correction", "clara-deck-correction"),
    ],
)
def test_clara_skill_identity_uses_namespace_without_redundant_prefix(
    skill_name: str, display_name: str, old_skill_name: str
) -> None:
    skill_root = PLUGIN_ROOT / "skills" / skill_name
    skill_text = (skill_root / "SKILL.md").read_text(encoding="utf-8")
    agent_metadata = (skill_root / "agents" / "openai.yaml").read_text(encoding="utf-8")

    assert f"\nname: {skill_name}\n" in f"\n{skill_text}"
    assert f'display_name: "{display_name}"' in agent_metadata
    assert f"Use ${skill_name}" in agent_metadata
    assert not (PLUGIN_ROOT / "skills" / old_skill_name).exists()


def load_core() -> Any:
    script_path = SCRIPTS_DIR / "advisor_case_core.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_core", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_hosted_interview_manager() -> Any:
    script_path = SCRIPTS_DIR / "manage_hosted_interview.py"
    spec = importlib.util.spec_from_file_location(
        "clara_manage_hosted_interview", script_path
    )
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def load_html_deck_runtime() -> Any:
    script_path = SCRIPTS_DIR / "html_deck_runtime.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_html_deck_runtime", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_hosted_voice_importer() -> Any:
    script_path = SCRIPTS_DIR / "import_hosted_voice_bundle.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_hosted_voice_importer", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_voice_feedback_timeline_builder() -> Any:
    script_path = SCRIPTS_DIR / "build_voice_feedback_timeline.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_voice_feedback_timeline", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_feedback_slide_matcher() -> Any:
    script_path = SCRIPTS_DIR / "match_feedback_frames_to_deck_slides.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_feedback_slide_matcher", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_hosted_transcript_finalizer() -> Any:
    script_path = SCRIPTS_DIR / "finalize_hosted_transcript.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_hosted_transcript_finalizer", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_voice_deck_revision_preparer() -> Any:
    script_path = SCRIPTS_DIR / "prepare_voice_deck_revision.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_voice_deck_revision_preparer", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_workbench_builder() -> Any:
    script_path = SCRIPTS_DIR / "build_deck_revision_workbench.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_workbench", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_plan_finalizer() -> Any:
    script_path = SCRIPTS_DIR / "finalize_deck_revision_plan.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_plan_finalizer", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_material_analyzer() -> Any:
    script_path = SCRIPTS_DIR / "analyze_deck_revision_materials.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_material_analyzer", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_quote_matrix_builder() -> Any:
    script_path = SCRIPTS_DIR / "build_deck_revision_quote_candidate_matrix.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_quote_matrix", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_interpretation_packet_builder() -> Any:
    script_path = SCRIPTS_DIR / "build_deck_revision_interpretation_packets.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_interpretation_packets", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_execution_planner() -> Any:
    script_path = SCRIPTS_DIR / "build_deck_revision_execution_plan.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_execution_planner", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_execution_packet_builder() -> Any:
    script_path = SCRIPTS_DIR / "build_deck_revision_execution_packets.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_execution_packets", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_approver() -> Any:
    script_path = SCRIPTS_DIR / "approve_deck_revision_plan.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_approver", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_applier() -> Any:
    script_path = SCRIPTS_DIR / "apply_deck_revision_plan.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_applier", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_output_review_completer() -> Any:
    script_path = SCRIPTS_DIR / "complete_deck_revision_output_review.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_output_review_completer",
            script_path,
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_fixture_runner() -> Any:
    script_path = SCRIPTS_DIR / "run_deck_revision_fixture.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_fixture_runner", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_revision_verifier() -> Any:
    script_path = SCRIPTS_DIR / "verify_deck_revision_output.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_revision_verifier", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_latest_hosted_voice_importer() -> Any:
    script_path = SCRIPTS_DIR / "import_latest_hosted_voice_bundle.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_latest_hosted_voice_importer", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_hosted_audio_uploader() -> Any:
    script_path = SCRIPTS_DIR / "upload_hosted_audio.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_hosted_audio_uploader", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_hosted_voice_launcher() -> Any:
    script_path = SCRIPTS_DIR / "launch_hosted_voice.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_hosted_voice_launcher", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_deck_feedback_starter() -> Any:
    script_path = SCRIPTS_DIR / "start_deck_feedback.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_deck_feedback_starter", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def load_dependency_checker() -> Any:
    script_path = SCRIPTS_DIR / "check_dependencies.py"
    sys.path.insert(0, str(SCRIPTS_DIR))
    try:
        spec = importlib.util.spec_from_file_location(
            "advisor_case_workspace_dependency_checker", script_path
        )
        assert spec and spec.loader
        module = importlib.util.module_from_spec(spec)
        sys.modules[spec.name] = module
        spec.loader.exec_module(module)
        return module
    finally:
        sys.path.remove(str(SCRIPTS_DIR))


def fixed_now() -> datetime:
    return datetime(2026, 1, 2, 10, 30, tzinfo=timezone.utc)


def init_case(tmp_path: Path) -> tuple[Any, Path]:
    core = load_core()
    case_dir = tmp_path / "case"
    core.initialize_case(
        case_dir,
        client="ClientCo",
        project="Succession advisory",
        objective="Prepare a decision pack",
        audience="Owner",
        output_language="it",
        now=fixed_now(),
    )
    return core, case_dir


def write_minimal_pptx_package(
    path: Path,
    *,
    media_parts: dict[str, bytes] | None = None,
    custom_properties: bytes | None = None,
) -> None:
    content_types = (
        b'<?xml version="1.0" encoding="UTF-8"?>'
        b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        b'<Default Extension="xml" ContentType="application/xml"/>'
        b'<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        b'<Override PartName="/ppt/presentation.xml" '
        b'ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
    )
    if custom_properties is not None:
        content_types += (
            b'<Override PartName="/docProps/custom.xml" '
            b'ContentType="application/vnd.openxmlformats-officedocument.custom-properties+xml"/>'
        )
    content_types += b"</Types>"
    with ZipFile(path, "w", compression=ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", b"<Relationships/>")
        archive.writestr("ppt/presentation.xml", b"<p:presentation/>")
        for name, data in (media_parts or {}).items():
            archive.writestr(name, data)
        if custom_properties is not None:
            archive.writestr("docProps/custom.xml", custom_properties)


def test_fixed_html_deck_runtime_adds_geometry_invariants() -> None:
    runtime = load_html_deck_runtime()
    html = (
        "<!doctype html><html><head><style>main{width:100%;}</style></head>"
        "<body><main><section><svg viewBox='0 0 16 9'></svg></section></main></body></html>"
    )

    fixed_html = runtime.apply_fixed_16_9_deck_runtime(html)

    assert 'data-clara-fixed-16-9-deck="true"' in fixed_html
    assert 'class="clara-fixed-16-9-deck"' in fixed_html
    assert "aspect-ratio: 16 / 9" in fixed_html
    assert "--clara-deck-width" in fixed_html
    assert "widthFromHeight" in fixed_html
    assert "preserveAspectRatio" in fixed_html
    assert "setCaptureHandleConfig" in fixed_html
    assert "clara_html_deck" in fixed_html
    assert "slide_id" in fixed_html
    assert "slide_title" in fixed_html
    runtime.assert_fixed_16_9_deck_runtime(fixed_html, label="test.html")
    assert runtime.apply_fixed_16_9_deck_runtime(fixed_html) == fixed_html

    with pytest.raises(ValueError, match="deck marker"):
        runtime.assert_fixed_16_9_deck_runtime(html, label="broken.html")


def write_downloaded_voice_bundle(
    path: Path,
    *,
    captured_at: str,
    transcript: str,
    source: str = "case_notes_hosted_voice",
    mtime: int = 1,
) -> None:
    path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": source,
                "captured_at": captured_at,
                "model": "gpt-realtime-whisper",
                "user_transcript": transcript,
                "assistant_transcript": "",
                "extraction_json": {
                    "cleaned_notes_markdown": "",
                    "entries": [],
                    "open_questions": [],
                },
                "extraction_text": "",
            }
        ),
        encoding="utf-8",
    )
    os.utime(path, (mtime, mtime))


def downloaded_voice_bundle_payload(
    *,
    captured_at: str,
    transcript: str,
    source: str = "case_notes_hosted_voice",
    audio_file_name: str = "",
    video_file_name: str = "",
) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "schema_version": 1,
        "source": source,
        "captured_at": captured_at,
        "model": "gpt-realtime-whisper",
        "user_transcript": transcript,
        "assistant_transcript": "",
        "extraction_json": {
            "cleaned_notes_markdown": "",
            "entries": [],
            "open_questions": [],
        },
        "extraction_text": "",
    }
    if audio_file_name:
        payload["audio_file_name"] = audio_file_name
        payload["audio_content_type"] = "audio/wav"
    if video_file_name:
        payload["video_file_name"] = video_file_name
        payload["video_content_type"] = "video/webm"
        payload["video_chunks"] = 1
        payload["screen_capture_metadata"] = {
            "started_at": "2026-01-02T10:30:00.000Z",
            "mime_type": "video/webm",
            "width": 1440,
            "height": 900,
            "display_surface": "browser",
        }
    return payload


def write_downloaded_voice_zip_bundle(
    path: Path,
    *,
    captured_at: str,
    transcript: str,
    audio_file_name: str = "meeting.wav",
    audio_bytes: bytes = b"audio bytes",
    video_file_name: str = "",
    video_bytes: bytes = b"video bytes",
    source: str = "case_notes_hosted_voice",
    mtime: int = 1,
) -> None:
    payload = downloaded_voice_bundle_payload(
        captured_at=captured_at,
        transcript=transcript,
        source=source,
        audio_file_name=audio_file_name,
        video_file_name=video_file_name,
    )
    json_name = path.with_suffix(".json").name
    with ZipFile(path, "w") as archive:
        archive.writestr(json_name, json.dumps(payload))
        archive.writestr(audio_file_name, audio_bytes)
        if video_file_name:
            archive.writestr(video_file_name, video_bytes)
    os.utime(path, (mtime, mtime))


def test_manifest_and_skill_are_generic() -> None:
    manifest = json.loads(
        (PLUGIN_ROOT / ".codex-plugin" / "plugin.json").read_text(encoding="utf-8")
    )
    skill = (PLUGIN_ROOT / "skills" / "clara" / "SKILL.md").read_text(encoding="utf-8")
    normalized_skill = " ".join(skill.split())
    plugin_text = "\n".join(
        path.read_text(encoding="utf-8")
        for path in PLUGIN_ROOT.rglob("*")
        if path.is_file() and path.suffix in {".json", ".md", ".html", ".py", ".txt"}
    )

    assert manifest["name"] == "clara"
    assert manifest["skills"] == "./skills/"
    assert "advisor_judgement" in skill
    assert "Human-Visible Document Quality Gate" in skill
    assert "run repeated model-led" in skill
    assert "looking for bullshit" in skill
    assert "idiotic style figures" in skill
    assert "Two-Loop Advisory Delivery Model" in skill
    assert "Advisory Intelligence Loop" in skill
    assert "Presentation Excellence Loop" in skill
    assert "best current advisory position" in normalized_skill
    assert "responsible decision posture now" in normalized_skill
    assert "The first deck is not a summary" in normalized_skill
    assert "advisory_evidence_map.md" in skill
    assert "living Loop 1 control artifact" in skill
    assert "claim-by-claim, not source-by-source" in skill
    assert "supports, weakens, contradicts, or creates" in skill
    assert "what this evidence proves" in skill
    assert "what this evidence does not prove" in skill
    assert "directness, reliability, corroboration, bias or limitation" in skill
    assert "Update `advisory_evidence_map.md` whenever new material is indexed" in skill
    assert "Evidence-navigation test" in skill
    assert "Default advisor-time assumption: the advisor has no time." in skill
    assert "advisory_workpaper.md" in skill
    assert "judgement_checkpoint.md" in skill
    assert "presentation_storyline.md" in skill
    assert "presentation_review.md" in skill
    assert "Use `/goal` only for major phase gates" in skill
    assert "A beautiful" in skill
    assert "suppresses critical unknowns" in skill
    assert "Evidence-gap test" in skill
    assert "Do not turn unresolved evidence needs into generic" in skill
    default_prompts = manifest["interface"]["defaultPrompt"]
    default_prompt = "\n".join(default_prompts)
    assert all(len(prompt) <= 128 for prompt in default_prompts)
    assert "update the evidence map" in default_prompt
    assert "decisions that cannot wait" in default_prompt
    assert "render the right business chart" in default_prompt
    forbidden_terms = ("Al" + "fredo", "Car" + "lo", "Gal" + "loni")
    for term in forbidden_terms:
        assert term not in plugin_text


def test_conversation_capabilities_are_separate_and_discoverable() -> None:
    manifest = json.loads(
        (PLUGIN_ROOT / ".codex-plugin" / "plugin.json").read_text(encoding="utf-8")
    )
    fixture = json.loads(
        (PLUGIN_ROOT / "evals" / "trigger_fixtures.json").read_text(encoding="utf-8")
    )
    main_skill = (PLUGIN_ROOT / "skills" / "clara" / "SKILL.md").read_text(
        encoding="utf-8"
    )

    assert manifest["version"] == "0.1.92"
    assert manifest["interface"]["shortDescription"] == ("AI companion for consultants")
    assert len(manifest["interface"]["defaultPrompt"]) == 3
    assert "hosted-interviews" in manifest["keywords"]
    assert "voice-transcription" in manifest["keywords"]
    assert "deck-correction" in manifest["keywords"]
    assert "brand-fit" in manifest["keywords"]
    assert "reporting-engine" in manifest["keywords"]
    assert "dataset-semantic-layer" in manifest["keywords"]
    assert "Conversation Workflow Router" in main_skill
    assert "Hosted-interview bundles and Hosted Voice bundles" in main_skill
    fixture_ids = {item["id"] for item in fixture["should_trigger"]}
    assert {
        "hosted-participant-interview",
        "voice-note-transcription-import",
        "voice-led-pptx-correction",
        "one-command-deck-feedback-capture",
        "brand-fit-current-presence-and-owned-catalogue",
        "brand-fit-stored-snapshot-boundary",
        "reporting-engine-direct-chart-analysis",
        "reporting-engine-chart-contract",
        "reporting-engine-semantic-layer",
    }.issubset(fixture_ids)
    expected_routes = {
        item["id"]: item.get("expected_skill")
        for item in fixture["should_trigger"]
        if item["id"] in fixture_ids
    }
    assert expected_routes["hosted-participant-interview"] == "clara:interview"
    assert expected_routes["voice-note-transcription-import"] == "clara:transcribe"
    assert expected_routes["voice-led-pptx-correction"] == "clara:deck-correction"
    assert expected_routes["one-command-deck-feedback-capture"] == (
        "clara:deck-correction"
    )
    assert expected_routes["brand-fit-current-presence-and-owned-catalogue"] == (
        "clara:brand-fit"
    )
    assert expected_routes["brand-fit-stored-snapshot-boundary"] == ("clara:brand-fit")
    assert expected_routes["reporting-engine-direct-chart-analysis"] == (
        "clara:reporting-engine"
    )
    assert expected_routes["reporting-engine-chart-contract"] == (
        "clara:reporting-engine"
    )
    assert expected_routes["reporting-engine-semantic-layer"] == (
        "clara:reporting-engine"
    )
    main_description = main_skill.split("---", 2)[1]
    assert "hosted interview" not in main_description.lower()
    assert "transcrib" not in main_description.lower()
    assert "deck-correction" not in main_description.lower()


def test_hosted_interview_manager_prepares_exact_campaign(monkeypatch) -> None:
    manager = load_hosted_interview_manager()
    observed: dict[str, Any] = {}

    def fake_request_json(_opener, **kwargs):
        observed.update(kwargs)
        return {"public_url": "https://mparanza.com/case-notes/interview/token"}

    monkeypatch.setattr(manager, "_request_json", fake_request_json)

    result = manager.prepare_campaign_interview(
        object(),
        interview_campaign_id="research-program-v1",
        case_id="participant-001",
        participant_name="Participant",
        language="fr",
    )

    assert result["public_url"].endswith("/token")
    assert observed["method"] == "POST"
    assert observed["url"].endswith("/campaigns/research-program-v1/interviews")
    assert observed["payload"] == {
        "case_id": "participant-001",
        "participant_name": "Participant",
        "language": "fr",
        "interviewee_role": "",
        "expires_in_hours": 168,
    }


def test_hosted_interview_manager_status_uses_public_endpoint(monkeypatch) -> None:
    manager = load_hosted_interview_manager()
    observed: dict[str, Any] = {}

    def fake_request_json(_opener, **kwargs):
        observed.update(kwargs)
        return {"status": "completed"}

    monkeypatch.setattr(manager, "_request_json", fake_request_json)

    result = manager.get_interview_status(
        object(),
        token_or_url="https://mparanza.com/case-notes/interview/token-123/output",
    )

    assert result == {"status": "completed"}
    assert observed["url"].endswith("/case-notes/api/interviews/token-123/status")


def test_hosted_interview_manager_magic_link_returns_to_authorized_voice_route(
    monkeypatch,
) -> None:
    manager = load_hosted_interview_manager()
    observed: dict[str, Any] = {}

    def fake_request_json(_opener, **kwargs):
        observed.update(kwargs)
        return {}

    monkeypatch.setattr(manager, "_request_json", fake_request_json)

    manager.request_magic_link(object(), email="advisor@example.com")

    assert observed["payload"] == {
        "email": "advisor@example.com",
        "redirect_path": "/case-notes/voice/launch",
    }


def test_human_visible_document_quality_gate_rejects_scaffolding() -> None:
    core = load_core()

    violations = core.audit_human_visible_document_text("""
        Pagina 3
        jud-0007
        etichetta da confermare
        Il prossimo passo e' produrre un documento piu decidibile.
        Il documento deve aiutare a pensare.
        Questo passaggio puo' creare valore.
        """)

    assert "visible judgement/source id" in violations
    assert "page-number scaffolding" in violations
    assert "placeholder label" in violations
    assert "pseudo-editorial filler" in violations
    assert "generic value language" in violations


def test_initialize_case_creates_schema_files(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)

    errors = core.validate_case_workspace(case_dir)
    manifest = json.loads((case_dir / "case_manifest.json").read_text())
    materials = json.loads((case_dir / "material_registry.json").read_text())
    judgement = json.loads((case_dir / "judgement_log.json").read_text())
    questions = json.loads((case_dir / "open_questions.json").read_text())
    issues = json.loads((case_dir / "case_issues.json").read_text())
    clara_mandate = json.loads((case_dir / "clara_mandate.json").read_text())
    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert errors == []
    assert manifest["client"] == "ClientCo"
    assert manifest["output_language"] == "it"
    assert materials["materials"] == []
    assert judgement["entries"] == []
    assert questions["questions"] == []
    assert issues["issues"] == []
    assert clara_mandate["persona"] == "Clara"
    assert clara_mandate["status"] == "not_started"
    assert "Derived working brief" in brief
    assert "The source of truth remains the JSON case files." in brief
    assert "Pending judgement entries: 0" in brief
    assert "Active case issues: 0" in brief
    assert "## Clara Mandate" in brief


def test_validate_workspace_cli_accepts_clean_case(tmp_path: Path) -> None:
    _, case_dir = init_case(tmp_path)

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "validate_workspace.py"),
            str(case_dir),
        ],
        check=True,
        capture_output=True,
        text=True,
    )

    assert "validation_errors=[]" in result.stderr


def test_validate_workspace_cli_reports_errors(tmp_path: Path) -> None:
    _, case_dir = init_case(tmp_path)
    manifest_path = case_dir / "case_manifest.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    manifest["status"] = "not-a-status"
    manifest_path.write_text(json.dumps(manifest), encoding="utf-8")

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "validate_workspace.py"),
            str(case_dir),
        ],
        check=False,
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert "validation_error: case_manifest.json: invalid status" in result.stderr


def test_validate_workspace_cli_reports_stale_linked_audio_pointer(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    pointer_path = case_dir / "source_materials" / "interviews" / "audio.md"
    pointer_path.parent.mkdir(parents=True)
    pointer_path.write_text(
        "# Audio pointer\n\n"
        "- Stato trascrizione: non ancora trascritto in Clara\n\n"
        "Quando sara' prodotta la trascrizione, conservarla in `notes/`.\n",
        encoding="utf-8",
    )
    transcript_path = (
        case_dir
        / "voice_sessions"
        / "20260102103000Z"
        / "raw_transcript_rule_attributed.md"
    )
    transcript_path.parent.mkdir(parents=True)
    transcript_path.write_text("# Transcript\n\nMorgan: text.\n", encoding="utf-8")
    registry_path = case_dir / "material_registry.json"
    registry = json.loads(registry_path.read_text(encoding="utf-8"))
    registry["materials"].extend(
        [
            {
                "added_at": fixed_now().isoformat(),
                "id": "mat-0001",
                "last_reviewed": None,
                "material_type": "source",
                "path": str(pointer_path.resolve()),
                "status": "indexed",
                "summary": "Pointer to raw audio.",
                "title": "Audio pointer",
                "updated_at": fixed_now().isoformat(),
            },
            {
                "added_at": fixed_now().isoformat(),
                "id": "mat-0002",
                "last_reviewed": None,
                "material_type": "transcript",
                "path": str(transcript_path.resolve()),
                "source_metadata": {
                    "raw_audio_pointer_material_id": "mat-0001",
                },
                "status": "indexed",
                "summary": "Attributed transcript.",
                "title": "Interview transcript",
                "updated_at": fixed_now().isoformat(),
            },
        ]
    )
    registry_path.write_text(json.dumps(registry), encoding="utf-8")

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "validate_workspace.py"),
            str(case_dir),
        ],
        check=False,
        capture_output=True,
        text=True,
    )

    assert result.returncode == 1
    assert (
        "validation_error: material_registry.json: audio pointer mat-0001 "
        "linked to transcript mat-0002 is not marked transcribed" in result.stderr
    )
    assert (
        "validation_error: material_registry.json: audio pointer mat-0001 "
        "linked to transcript mat-0002 still says not transcribed" in result.stderr
    )


def test_delete_materials_removes_registry_and_source_references(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    kept_source = tmp_path / "kept-source.md"
    kept_source.write_text("# Kept\n\nValid source.", encoding="utf-8")
    session_dir = case_dir / "voice_sessions" / "20260102103000Z"
    session_dir.mkdir(parents=True)
    wrong_transcript = session_dir / "raw_transcript.md"
    wrong_transcript.write_text("wrong transcript", encoding="utf-8")
    kept_material = core.register_material(
        case_dir,
        kept_source,
        title="Kept source",
        now=fixed_now(),
    )
    wrong_material = core.register_material(
        case_dir,
        wrong_transcript,
        material_type="transcript",
        title="Wrong transcript",
        summary="Wrong transcript import.",
        now=fixed_now(),
    )
    core.prepare_clara_kickoff(case_dir, now=fixed_now())
    core.update_clara_mandate_from_kickoff(
        case_dir,
        {"clara_mandate": {"clara_understanding": "Kickoff imported."}},
        material_id=wrong_material["id"],
        session_path="voice_sessions/20260102103000Z",
        now=fixed_now(),
    )
    entries = core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "fact",
                "text": "This claim has another valid source.",
                "status": "pending",
                "source_material_ids": [wrong_material["id"], kept_material["id"]],
                "rationale": "Mixed evidence.",
            },
            {
                "kind": "codex_inference",
                "text": "This claim depends only on the wrong transcript.",
                "status": "pending",
                "source_material_ids": [wrong_material["id"]],
                "rationale": "Wrong transcript.",
            },
        ],
        now=fixed_now(),
    )

    result = core.delete_materials(case_dir, [wrong_material["id"]], now=fixed_now())
    registry = json.loads((case_dir / "material_registry.json").read_text())
    judgement = json.loads((case_dir / "judgement_log.json").read_text())
    mandate = json.loads((case_dir / "clara_mandate.json").read_text())
    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert result.removed_material_ids == (wrong_material["id"],)
    assert [item["id"] for item in registry["materials"]] == [kept_material["id"]]
    assert judgement["entries"][0]["source_material_ids"] == [kept_material["id"]]
    assert judgement["entries"][1]["source_material_ids"] == []
    assert set(result.updated_judgement_ids) == {
        entries[0]["id"],
        entries[1]["id"],
    }
    assert result.unanchored_judgement_ids == (entries[1]["id"],)
    assert wrong_material["id"] not in mandate["source_material_ids"]
    assert all(
        item["id"] != wrong_material["id"]
        for item in mandate["preparation"]["material_anchors"]
    )
    assert result.removed_mandate_voice_session_paths == (
        "voice_sessions/20260102103000Z",
    )
    assert session_dir in result.orphan_candidate_paths
    assert wrong_transcript.exists()
    assert session_dir.exists()
    assert "Wrong transcript" not in brief
    assert "Kept source" in brief


def test_delete_material_cli_reports_summary_and_refreshes_brief(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    wrong_source = tmp_path / "wrong-source.md"
    wrong_source.write_text("# Wrong\n\nBad material.", encoding="utf-8")
    wrong_material = core.register_material(
        case_dir,
        wrong_source,
        title="Wrong source",
        summary="Bad material.",
        now=fixed_now(),
    )

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "delete_material.py"),
            str(case_dir),
            wrong_material["id"],
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    summary = json.loads(result.stderr)
    registry = json.loads((case_dir / "material_registry.json").read_text())
    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert summary["removed_material_ids"] == [wrong_material["id"]]
    assert summary["updated_judgement_ids"] == []
    assert registry["materials"] == []
    assert "Wrong source" not in brief


def test_add_open_questions_cli_imports_questions_json(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    source_entries = core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "fact",
                "text": "The senior owner still controls quality decisions.",
                "status": "pending",
                "source_material_ids": [],
                "rationale": "Interview note.",
            }
        ],
    )
    source_entry_id = source_entries[0]["id"]
    questions_json = tmp_path / "questions.json"
    questions_json.write_text(
        json.dumps(
            {
                "questions": [
                    {
                        "question": "Who owns quality after the transition?",
                        "why_it_matters": "Quality ownership is a succession gate.",
                        "source_entry_ids": [source_entry_id],
                    },
                    {
                        "question": "Which decisions require board approval?",
                        "why_it_matters": "Decision rights must be explicit.",
                    },
                ]
            }
        ),
        encoding="utf-8",
    )

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "add_open_questions.py"),
            str(case_dir),
            "--questions-json",
            str(questions_json),
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    questions = json.loads((case_dir / "open_questions.json").read_text())
    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert "Added 2 open question(s)." in result.stderr
    assert [item["id"] for item in questions["questions"]] == ["q-0001", "q-0002"]
    assert questions["questions"][0]["source_entry_ids"] == [source_entry_id]
    assert questions["questions"][1]["source_entry_ids"] == []
    assert "Who owns quality after the transition?" in brief
    assert "Open Questions" in brief


def test_upsert_case_issues_cli_tracks_cross_interview_issue(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    source_entries = core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "codex_inference",
                "text": "Quality ownership is the main transition gate.",
                "status": "pending",
                "source_material_ids": [],
                "rationale": "Interview note.",
            }
        ],
    )
    source_entry_id = source_entries[0]["id"]
    source_question = core.add_open_question(
        case_dir,
        question="Who owns quality after the founder steps back?",
        why_it_matters="It tests whether the mandate has operational substance.",
    )
    issues_json = tmp_path / "issues.json"
    issues_json.write_text(
        json.dumps(
            {
                "issues": [
                    {
                        "id": "production_quality_transition",
                        "title": "Production and quality transition",
                        "decision_area": "Operating transition",
                        "current_synthesis": (
                            "The transition is not credible until quality ownership is explicit."
                        ),
                        "evidence_for": [source_entry_id],
                        "evidence_against": [],
                        "open_tests": [source_question["id"]],
                    }
                ]
            }
        ),
        encoding="utf-8",
    )

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "upsert_case_issues.py"),
            str(case_dir),
            "--issues-json",
            str(issues_json),
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    issues = json.loads((case_dir / "case_issues.json").read_text())
    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert "Upserted 1 case issue(s)." in result.stderr
    assert issues["issues"][0]["id"] == "production_quality_transition"
    assert issues["issues"][0]["evidence_for"] == [source_entry_id]
    assert issues["issues"][0]["open_tests"] == [source_question["id"]]
    assert "Production and quality transition" in brief
    assert "Evidence links: 1 supporting, 0 opposing, 1 open test." in brief
    assert "jud-0001" not in brief
    assert "q-0001" not in brief


def test_integrate_transcript_review_cli_applies_evidence_chain(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    session_dir = case_dir / "voice_sessions" / "20260102T103000Z"
    session_dir.mkdir(parents=True)
    raw_transcript = session_dir / "raw_transcript.md"
    attributed_transcript = session_dir / "raw_transcript_rule_attributed.md"
    backup_transcript = session_dir / "raw_transcript_unattributed.md"
    review_note = session_dir / "clara_review.md"
    raw_transcript.write_text("raw transcript", encoding="utf-8")
    attributed_transcript.write_text("attributed transcript", encoding="utf-8")
    backup_transcript.write_text("unattributed transcript", encoding="utf-8")
    review_note.write_text(
        "\n".join(
            [
                "# Clara Audio Review",
                "",
                "## Neutral Summary",
                "",
                "Pending local Clara review.",
                "",
                "## Clara Opinion",
                "",
                "Pending local Clara review.",
                "",
                "## Cleaned Transcript",
                "",
                "Transcript body.",
            ]
        ),
        encoding="utf-8",
    )
    material = core.register_material(
        case_dir,
        raw_transcript,
        material_type="transcript",
        title="Interview transcript",
        now=fixed_now(),
    )
    question = core.add_open_question(
        case_dir,
        question="Who owns quality after the transition?",
        why_it_matters="Quality ownership is the main gate.",
        now=fixed_now(),
    )
    existing_judgement = core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "fact",
                "text": "Existing evidence says quality is a transition gate.",
                "status": "pending",
                "source_material_ids": [material["id"]],
                "rationale": "Existing interview note.",
            }
        ],
        now=fixed_now(),
    )[0]
    core.upsert_case_issues(
        case_dir,
        [
            {
                "id": "quality_transition",
                "title": "Quality transition",
                "decision_area": "Operating transition",
                "current_synthesis": "Quality ownership is unresolved.",
                "evidence_for": [existing_judgement["id"]],
                "evidence_against": [],
                "open_tests": [question["id"]],
            }
        ],
        now=fixed_now(),
    )
    plan_json = tmp_path / "integration_plan.json"
    plan_json.write_text(
        json.dumps(
            {
                "material_registry_updates": [
                    {
                        "material_id": material["id"],
                        "path": str(attributed_transcript),
                        "source_metadata": {
                            "speaker_attribution": "text-only local pass",
                            "unattributed_transcript_backup": str(
                                backup_transcript.relative_to(case_dir)
                            ),
                        },
                    }
                ],
                "review_notes": [
                    {
                        "path": str(review_note.relative_to(case_dir)),
                        "sections": {
                            "Neutral Summary": (
                                "Text-only review finds one clear transition claim."
                            ),
                            "Clara Opinion": (
                                "Do not mark this interview decision-pack-ready yet."
                            ),
                        },
                    }
                ],
                "judgements": [
                    {
                        "key": "quality_claim",
                        "kind": "codex_inference",
                        "text": (
                            "The interview makes quality ownership a testable transition gate."
                        ),
                        "status": "pending",
                        "source_material_ids": [material["id"]],
                        "rationale": "Transcript review.",
                    }
                ],
                "open_question_links": [
                    {
                        "question_id": question["id"],
                        "source_entry_refs": ["quality_claim"],
                    }
                ],
                "case_issue_updates": [
                    {
                        "issue_id": "quality_transition",
                        "current_synthesis": (
                            "Quality ownership is still unresolved but now has "
                            "a transcript-backed test."
                        ),
                        "evidence_for_refs": ["quality_claim"],
                    }
                ],
            }
        ),
        encoding="utf-8",
    )

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "integrate_transcript_review.py"),
            str(case_dir),
            "--plan-json",
            str(plan_json),
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    registry = json.loads((case_dir / "material_registry.json").read_text())
    judgement = json.loads((case_dir / "judgement_log.json").read_text())
    questions = json.loads((case_dir / "open_questions.json").read_text())
    issues = json.loads((case_dir / "case_issues.json").read_text())
    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")
    review_text = review_note.read_text(encoding="utf-8")
    summary = json.loads(result.stderr)
    new_entry_id = summary["judgement_key_map"]["quality_claim"]

    assert summary["validation_errors"] == []
    assert summary["material_updates"] == [material["id"]]
    assert registry["materials"][0]["path"] == str(attributed_transcript.resolve())
    assert (
        registry["materials"][0]["source_metadata"]["speaker_attribution"]
        == "text-only local pass"
    )
    assert raw_transcript.exists()
    assert backup_transcript.exists()
    assert "Text-only review finds one clear transition claim." in review_text
    assert "Do not mark this interview decision-pack-ready yet." in review_text
    assert judgement["entries"][-1]["id"] == new_entry_id
    assert judgement["entries"][-1]["status"] == "pending"
    assert judgement["entries"][-1]["source_material_ids"] == [material["id"]]
    assert new_entry_id in questions["questions"][0]["source_entry_ids"]
    assert new_entry_id in issues["issues"][0]["evidence_for"]
    assert "The interview makes quality ownership a testable transition gate." in brief
    assert "Decision-pack-ready judgement entries: 0" in brief


def test_index_materials_keeps_original_file_path(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    source = tmp_path / "source.md"
    source.write_text("# Source\n\nThe owner still controls quality.", encoding="utf-8")

    indexed = core.index_materials(case_dir, [source], now=fixed_now())
    registry = json.loads((case_dir / "material_registry.json").read_text())

    assert len(indexed) == 1
    assert registry["materials"][0]["path"] == str(source.resolve())
    assert registry["materials"][0]["material_type"] == "source"
    assert "owner still controls quality" in registry["materials"][0]["summary"]


def test_index_materials_merges_source_metadata(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    source = tmp_path / "source.md"
    source.write_text(
        "# Source\n\nThe corrected deck is authoritative.", encoding="utf-8"
    )

    indexed = core.index_materials(
        case_dir,
        [source],
        source_metadata={
            "provenance_note": "Corrected by the advisor.",
            "planned_presentation_date": "2026-07-06",
        },
        now=fixed_now(),
    )
    core.index_materials(
        case_dir,
        [source],
        source_metadata={"advisor_perspective": "Authoritative version."},
        now=fixed_now(),
    )
    registry = json.loads((case_dir / "material_registry.json").read_text())

    assert len(indexed) == 1
    assert len(registry["materials"]) == 1
    assert registry["materials"][0]["source_metadata"] == {
        "advisor_perspective": "Authoritative version.",
        "planned_presentation_date": "2026-07-06",
        "provenance_note": "Corrected by the advisor.",
    }


def test_index_materials_cli_accepts_provenance_and_source_metadata(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    source = tmp_path / "source.md"
    source.write_text(
        "# Source\n\nThe corrected deck is authoritative.", encoding="utf-8"
    )

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "index_materials.py"),
            str(case_dir),
            str(source),
            "--source-metadata",
            '{"planned_presentation_date": "2026-07-06"}',
            "--provenance-note",
            "Corrected by the advisor.",
        ],
        check=False,
        capture_output=True,
        text=True,
    )
    registry = json.loads((case_dir / "material_registry.json").read_text())

    assert result.returncode == 0, result.stderr
    assert registry["materials"][0]["source_metadata"] == {
        "planned_presentation_date": "2026-07-06",
        "provenance_note": "Corrected by the advisor.",
    }


def test_index_materials_supports_pptx_preview(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    core, case_dir = init_case(tmp_path)
    deck = tmp_path / "succession_plan.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Succession Plan"
    slide.placeholders[1].text = "Quality governance"
    presentation.save(deck)

    indexed = core.index_materials(case_dir, [deck], now=fixed_now())
    registry = json.loads((case_dir / "material_registry.json").read_text())

    assert len(indexed) == 1
    assert registry["materials"][0]["path"] == str(deck.resolve())
    assert registry["materials"][0]["material_type"] == "source"
    assert "Succession Plan" in registry["materials"][0]["summary"]
    assert "Quality governance" in registry["materials"][0]["summary"]


def test_inspect_pptx_legacy_media_finds_wmf_and_emf(tmp_path: Path) -> None:
    core = load_core()
    deck = tmp_path / "legacy.pptx"
    write_minimal_pptx_package(
        deck,
        media_parts={
            "ppt/media/image1.emf": b"emf",
            "ppt/media/image2.wmf": b"wmf",
            "ppt/media/image3.png": b"png",
        },
    )

    legacy_media = core.inspect_pptx_legacy_media(deck)

    assert legacy_media == ("ppt/media/image1.emf", "ppt/media/image2.wmf")


def test_normalize_legacy_pptx_roundtrips_and_preserves_custom_props(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    core = load_core()
    source = tmp_path / "legacy_deck.pptx"
    output = tmp_path / "legacy_deck_normalized.pptx"
    custom_xml = (
        b'<?xml version="1.0" encoding="UTF-8"?>'
        b'<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/custom-properties">'
        b'<property name="ClaraTranscriptPath"/></Properties>'
    )
    write_minimal_pptx_package(
        source,
        media_parts={"ppt/media/image1.emf": b"emf"},
        custom_properties=custom_xml,
    )

    def fake_roundtrip(
        source_path: Path,
        *,
        output_dir: Path,
        soffice_binary: Path,
    ) -> Path:
        assert source_path == source.resolve()
        assert soffice_binary == tmp_path / "soffice"
        converted = output_dir / source.name
        write_minimal_pptx_package(
            converted,
            media_parts={"ppt/media/image1.png": b"png"},
        )
        return converted

    monkeypatch.setattr(core, "_run_soffice_pptx_roundtrip", fake_roundtrip)
    monkeypatch.setattr(
        core, "resolve_soffice_binary", lambda configured=None: configured
    )

    result = core.normalize_legacy_pptx_for_editable_merge(
        source,
        output_path=output,
        soffice_binary=tmp_path / "soffice",
        now=fixed_now(),
    )
    report = json.loads(result.report_path.read_text(encoding="utf-8"))

    assert result.normalized is True
    assert result.legacy_media_before == ("ppt/media/image1.emf",)
    assert result.legacy_media_after == ()
    assert report["custom_properties_preserved"] is True
    assert report["legacy_media_before_count"] == 1
    assert report["legacy_media_after_count"] == 0
    with ZipFile(output) as archive:
        assert archive.read("docProps/custom.xml") == custom_xml
        content_types = archive.read("[Content_Types].xml").decode("utf-8")
    assert "/docProps/custom.xml" in content_types


def test_normalize_legacy_pptx_copies_clean_deck_without_soffice(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    core = load_core()
    source = tmp_path / "clean_deck.pptx"
    output = tmp_path / "clean_deck_normalized.pptx"
    write_minimal_pptx_package(
        source,
        media_parts={"ppt/media/image1.png": b"png"},
    )

    def fail_roundtrip(*_args: object, **_kwargs: object) -> Path:
        raise AssertionError("clean decks should not call LibreOffice")

    monkeypatch.setattr(core, "_run_soffice_pptx_roundtrip", fail_roundtrip)

    result = core.normalize_legacy_pptx_for_editable_merge(
        source,
        output_path=output,
        now=fixed_now(),
    )
    report = json.loads(result.report_path.read_text(encoding="utf-8"))

    assert result.normalized is False
    assert result.legacy_media_before == ()
    assert result.legacy_media_after == ()
    assert output.exists()
    assert report["soffice_binary"] is None


def test_prepare_editable_pptx_merge_input_uses_normalized_sibling(
    tmp_path: Path,
) -> None:
    core = load_core()
    source = tmp_path / "legacy_deck.pptx"
    normalized = tmp_path / "legacy_deck_normalized_for_merge.pptx"
    write_minimal_pptx_package(
        source,
        media_parts={"ppt/media/image1.emf": b"emf"},
    )
    write_minimal_pptx_package(
        normalized,
        media_parts={"ppt/media/image1.png": b"png"},
    )

    result = core.prepare_editable_pptx_merge_input(source, now=fixed_now())
    report = json.loads(result.report_path.read_text(encoding="utf-8"))

    assert result.merge_base_path == normalized.resolve()
    assert result.status == "normalized_merge_base_ready"
    assert result.source_legacy_media == ("ppt/media/image1.emf",)
    assert result.merge_base_legacy_media == ()
    assert report["editable_merge_guard"] == "normalized_required"
    assert report["normalized_path"] == str(normalized.resolve())


def test_prepare_editable_pptx_merge_input_requires_normalized_or_skip(
    tmp_path: Path,
) -> None:
    core = load_core()
    source = tmp_path / "legacy_deck.pptx"
    write_minimal_pptx_package(
        source,
        media_parts={"ppt/media/image1.emf": b"emf"},
    )

    with pytest.raises(core.CaseWorkspaceError, match="requires a normalized"):
        core.prepare_editable_pptx_merge_input(source, now=fixed_now())


def test_prepare_editable_pptx_merge_input_records_skip_reason(
    tmp_path: Path,
) -> None:
    core = load_core()
    source = tmp_path / "legacy_deck.pptx"
    write_minimal_pptx_package(
        source,
        media_parts={"ppt/media/image1.emf": b"emf"},
    )

    result = core.prepare_editable_pptx_merge_input(
        source,
        skip_normalization_reason="Only creating an image fallback for affected slides.",
        now=fixed_now(),
    )
    report = json.loads(result.report_path.read_text(encoding="utf-8"))

    assert result.merge_base_path == source.resolve()
    assert result.status == "normalization_skipped_with_reason"
    assert result.skip_normalization_reason == (
        "Only creating an image fallback for affected slides."
    )
    assert report["skip_normalization_reason"] == (
        "Only creating an image fallback for affected slides."
    )


def test_prepare_editable_pptx_merge_input_allows_clean_source(
    tmp_path: Path,
) -> None:
    core = load_core()
    source = tmp_path / "clean_deck.pptx"
    write_minimal_pptx_package(
        source,
        media_parts={"ppt/media/image1.png": b"png"},
    )

    result = core.prepare_editable_pptx_merge_input(source, now=fixed_now())
    report = json.loads(result.report_path.read_text(encoding="utf-8"))

    assert result.merge_base_path == source.resolve()
    assert result.status == "source_clean_for_merge"
    assert result.source_legacy_media == ()
    assert report["editable_merge_guard"] == "normalization_not_required"


def test_ingest_note_text_registers_pasted_note(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)

    material = core.ingest_note_text(
        case_dir,
        title="Advisor call",
        text="The leadership proposal needs a conflict process.",
        now=fixed_now(),
    )
    note_path = Path(material["path"])

    assert material["material_type"] == "note"
    assert note_path.exists()
    assert note_path.parent == case_dir / "notes"
    assert "conflict process" in note_path.read_text(encoding="utf-8")


def test_ingest_note_file_copies_external_note_into_case(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    source = tmp_path / "downloaded_note.md"
    source.write_text("A downloaded note should be durable.", encoding="utf-8")

    material = core.ingest_note_file(
        case_dir,
        title="Advisor file note",
        notes_file=source,
        now=fixed_now(),
    )
    note_path = Path(material["path"])

    assert material["material_type"] == "note"
    assert note_path.exists()
    assert note_path.parent == case_dir / "notes"
    assert note_path != source.resolve()
    assert note_path.read_text(encoding="utf-8") == source.read_text(encoding="utf-8")
    assert "downloaded note should be durable" in material["summary"]


def test_copy_case_file_routes_presentation_draft_without_registration(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    downloaded_deck = tmp_path / "20260706_Client_Incontro G3_V2.pptx"
    downloaded_deck.write_bytes(b"draft deck bytes")

    result = core.copy_case_file(case_dir, downloaded_deck, now=fixed_now())
    registry = json.loads((case_dir / "material_registry.json").read_text())

    assert result.kind == "presentation"
    assert result.copied is True
    assert result.registered_material is None
    assert (
        result.destination_path
        == (
            case_dir
            / "outputs"
            / "presentations"
            / "current"
            / "20260706_Client_Incontro G3_V2.pptx"
        ).resolve()
    )
    assert result.destination_path.read_bytes() == b"draft deck bytes"
    assert registry["materials"] == []
    assert result.legacy_pptx_normalization is None


def test_copy_case_file_auto_normalizes_legacy_pptx_presentation(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    core, case_dir = init_case(tmp_path)
    downloaded_deck = tmp_path / "20260706_Client_deck.pptx"
    custom_xml = (
        b'<?xml version="1.0" encoding="UTF-8"?>'
        b'<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/custom-properties">'
        b'<property name="ClaraTranscriptPath"/></Properties>'
    )
    write_minimal_pptx_package(
        downloaded_deck,
        media_parts={"ppt/media/image1.emf": b"emf"},
        custom_properties=custom_xml,
    )

    def fake_roundtrip(
        source_path: Path,
        *,
        output_dir: Path,
        soffice_binary: Path,
    ) -> Path:
        assert source_path.name == "20260706_Client_deck.pptx"
        converted = output_dir / source_path.name
        write_minimal_pptx_package(
            converted,
            media_parts={"ppt/media/image1.png": b"png"},
        )
        return converted

    monkeypatch.setattr(core, "_run_soffice_pptx_roundtrip", fake_roundtrip)
    monkeypatch.setattr(
        core, "resolve_soffice_binary", lambda configured=None: tmp_path / "soffice"
    )

    result = core.copy_case_file(case_dir, downloaded_deck, now=fixed_now())
    normalization = result.legacy_pptx_normalization

    assert result.kind == "presentation"
    assert normalization is not None
    assert normalization.source_path == result.destination_path
    assert (
        normalization.output_path
        == (
            case_dir
            / "outputs"
            / "presentations"
            / "current"
            / "20260706_Client_deck_normalized_for_merge.pptx"
        ).resolve()
    )
    assert normalization.output_path.exists()
    assert normalization.legacy_media_before == ("ppt/media/image1.emf",)
    assert normalization.legacy_media_after == ()
    report = json.loads(normalization.report_path.read_text(encoding="utf-8"))
    assert report["merge_guidance"] == (
        "Use this normalized PPTX as the base for editable slide merging."
    )


def test_copy_case_file_can_skip_legacy_pptx_normalization(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    core, case_dir = init_case(tmp_path)
    downloaded_deck = tmp_path / "20260706_Client_deck.pptx"
    write_minimal_pptx_package(
        downloaded_deck,
        media_parts={"ppt/media/image1.emf": b"emf"},
    )

    def fail_roundtrip(*_args: object, **_kwargs: object) -> Path:
        raise AssertionError("normalization should be skipped")

    monkeypatch.setattr(core, "_run_soffice_pptx_roundtrip", fail_roundtrip)

    result = core.copy_case_file(
        case_dir,
        downloaded_deck,
        normalize_legacy_pptx=False,
        now=fixed_now(),
    )

    assert result.kind == "presentation"
    assert result.legacy_pptx_normalization is None


def test_copy_case_file_registers_source_document(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    downloaded_pdf = tmp_path / "governance-background.pdf"
    downloaded_pdf.write_bytes(b"%PDF source reference")

    result = core.copy_case_file(
        case_dir,
        downloaded_pdf,
        register=True,
        title="Governance background",
        now=fixed_now(),
    )
    registry = json.loads((case_dir / "material_registry.json").read_text())
    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert result.kind == "source"
    assert (
        result.destination_path
        == (
            case_dir / "source_materials" / "project_docs" / "governance-background.pdf"
        ).resolve()
    )
    assert result.registered_material is not None
    assert result.registered_material["id"] == "mat-0001"
    assert registry["materials"][0]["path"] == str(result.destination_path)
    assert registry["materials"][0]["material_type"] == "source"
    assert registry["materials"][0]["title"] == "Governance background"
    assert (
        "PDF indexed; semantic review remains in Codex."
        in registry["materials"][0]["summary"]
    )
    assert "Governance background" in brief


def test_copy_case_file_uses_unique_name_for_conflicting_destination(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    first_download = tmp_path / "board-note.md"
    second_download = tmp_path / "other" / "board-note.md"
    first_download.write_text("First note.", encoding="utf-8")
    second_download.parent.mkdir()
    second_download.write_text("Second note.", encoding="utf-8")

    first_result = core.copy_case_file(case_dir, first_download, kind="note")
    second_result = core.copy_case_file(case_dir, second_download, kind="note")

    assert first_result.destination_path == (case_dir / "notes" / "board-note.md")
    assert second_result.destination_path == (case_dir / "notes" / "board-note-2.md")
    assert first_result.destination_path.read_text(encoding="utf-8") == "First note."
    assert second_result.destination_path.read_text(encoding="utf-8") == "Second note."


def test_clara_kickoff_preparation_and_partner_brief(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    source = tmp_path / "company.md"
    source.write_text(
        "The founder still owns key customer relationships.",
        encoding="utf-8",
    )
    core.index_materials(case_dir, [source], now=fixed_now())

    prepared = core.prepare_clara_kickoff(
        case_dir,
        industry_context=[
            "Customer trust is relationship-led, so succession must test commercial continuity."
        ],
        external_research=[
            {
                "title": "Public source",
                "url": "https://example.com/source",
                "takeaway": "Separate ownership from operating authority.",
            }
        ],
        now=fixed_now(),
    )
    mandate = core.update_clara_mandate_from_kickoff(
        case_dir,
        {
            "clara_mandate": {
                "engagement_objective": "Prepare the first working readout.",
                "client_decision": "Choose how authority transfers.",
                "clara_understanding": (
                    "The case is about commercial continuity after founder step-back."
                ),
                "partner_starting_orientation": "Keep the transition pragmatic.",
                "sensitive_points": [
                    "Do not make the founder's preference sound arbitrary."
                ],
                "what_clara_should_investigate": ["Customer concentration risk."],
                "what_clara_should_not_waste_time_on": ["Generic succession theory."],
                "essential_clarifications": ["Who owns key customer relationships?"],
                "next_steps": ["Prepare a partner questions list."],
            }
        },
        material_id="mat-9999",
        session_path="voice_sessions/20260102T103000Z",
        now=fixed_now(),
    )
    brief = core.render_clara_partner_brief(case_dir)
    deck = core.render_clara_kickoff_deck(case_dir)

    preparation_text = prepared.preparation_path.read_text(encoding="utf-8")
    html = brief.html_path.read_text(encoding="utf-8")
    deck_html = deck.html_path.read_text(encoding="utf-8")
    case_brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert prepared.material_count == 1
    assert prepared.baseline_source_count == 1
    assert mandate["status"] == "kickoff_imported"
    assert mandate["mandate"]["client_decision"] == "Choose how authority transfers."
    assert "Clara Kickoff Preparation" in preparation_text
    assert "Customer trust is relationship-led" in preparation_text
    assert "Briefing partner" in html
    assert "commercial continuity after founder step-back" in html
    assert "Who owns key customer relationships?" in html
    assert "Kickoff" in deck_html
    assert "Customer concentration risk." in deck_html
    assert "Decisioni per il prossimo passaggio." in deck_html
    assert 'data-clara-fixed-16-9-deck="true"' in deck_html
    assert "main.clara-fixed-16-9-deck" in deck_html
    assert "aspect-ratio: 16 / 9" in deck_html
    assert "widthFromHeight" in deck_html
    assert "preserveAspectRatio" in deck_html
    assert core.audit_human_visible_document_text(html) == []
    assert core.audit_human_visible_document_text(deck_html) == []
    assert deck.hypothesis_count == 0
    assert deck.open_question_count == 0
    assert brief.open_clarification_count == 1
    assert brief.next_step_count == 1
    assert "## Clara Mandate" in case_brief
    assert "Prepare a partner questions list." in case_brief


def test_partner_brief_surfaces_candidates_before_voice_kickoff(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    material = core.ingest_note_text(
        case_dir,
        title="Advisor note",
        text="The operating transition needs explicit quality ownership.",
        now=fixed_now(),
    )
    core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "advisor_judgement",
                "text": "Quality ownership is the central transition risk.",
                "source_material_ids": [material["id"]],
            }
        ],
        now=fixed_now(),
    )
    core.add_open_question(
        case_dir,
        question="Who owns quality decisions after the founder steps back?",
        why_it_matters="It defines whether authority has really transferred.",
        now=fixed_now(),
    )

    result = core.render_clara_partner_brief(case_dir)
    html = result.html_path.read_text(encoding="utf-8")

    assert "Quality ownership is the central transition risk." in html
    assert "Who owns quality decisions after the founder steps back?" in html
    assert "Decidere quali voci candidate includere" in html
    assert "Verifiche prioritarie" in html
    assert "Chiarimenti essenziali" not in html
    assert "Not explicit yet." not in html
    assert "No next steps recorded yet." not in html
    assert core.audit_human_visible_document_text(html) == []


def test_partner_brief_surfaces_approved_content_after_inclusion(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    material = core.ingest_note_text(
        case_dir,
        title="Advisor note",
        text="The transition needs explicit quality ownership.",
        now=fixed_now(),
    )
    core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "fact",
                "text": "The founder still owns quality authority.",
                "status": "approved",
                "source_material_ids": [material["id"]],
            },
            {
                "kind": "advisor_judgement",
                "text": "Approved quality ownership is the central transition risk.",
                "status": "approved",
                "source_material_ids": [material["id"]],
            },
        ],
        now=fixed_now(),
    )

    result = core.render_clara_partner_brief(case_dir)
    html = result.html_path.read_text(encoding="utf-8")

    assert "Contenuto pronto per il pack cliente" in html
    assert "Approved quality ownership is the central transition risk." in html
    assert "Nessun giudizio candidato registrato." not in html


def test_partner_brief_uses_indexed_research_when_industry_context_empty(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    research_dir = tmp_path / "research"
    research_dir.mkdir()
    research_file = research_dir / "parma_market.md"
    research_file.write_text(
        "Prosciutto di Parma market pressure makes quality and export resilience central.",
        encoding="utf-8",
    )
    core.index_materials(case_dir, [research_file], now=fixed_now())

    result = core.render_clara_partner_brief(case_dir)
    html = result.html_path.read_text(encoding="utf-8")

    assert "Contesto settoriale" in html
    assert "Prosciutto di Parma market pressure" in html
    assert "Nessun contesto settoriale registrato." not in html


def test_unapproved_judgement_is_excluded_from_decision_pack(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    source = tmp_path / "memo.md"
    source.write_text("Approved source", encoding="utf-8")
    material = core.index_materials(case_dir, [source], now=fixed_now())[0]

    core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "fact",
                "text": "Approved fact for the decision maker.",
                "status": "approved",
                "source_material_ids": [material["id"]],
                "reviewer": "Reviewer",
            },
            {
                "kind": "advisor_judgement",
                "text": "Pending sensitive judgement must stay out.",
                "status": "pending",
                "source_material_ids": [material["id"]],
            },
            {
                "kind": "decision_implication",
                "text": "Require explicit conflict escalation rules.",
                "status": "approved",
                "source_material_ids": [material["id"]],
                "reviewer": "Reviewer",
            },
        ],
        now=fixed_now(),
    )
    core.add_open_question(
        case_dir,
        question="Who owns quality after the transition?",
        why_it_matters="It changes the minimum governance conditions.",
        now=fixed_now(),
    )

    result = core.build_decision_pack(case_dir)
    markdown = result.markdown_path.read_text(encoding="utf-8")
    workpaper = result.workpaper_markdown_path.read_text(encoding="utf-8")

    assert "Approved fact for the decision maker." in markdown
    assert "Require explicit conflict escalation rules." in markdown
    assert "Who owns quality after the transition?" in markdown
    assert "Traccia esecutiva" in markdown
    assert "Percorso consigliato." in markdown
    assert "Pending sensitive judgement must stay out." not in markdown
    assert "Fatti acquisiti" in markdown
    assert "Lettura consulenziale" in markdown
    assert "Materials Indexed" not in markdown
    assert str(source.resolve()) not in markdown
    assert "Pending sensitive judgement must stay out." not in workpaper
    assert "Voci pending escluse: 1" in workpaper
    assert str(source.resolve()) in workpaper
    assert result.docx_path.exists()
    assert result.workpaper_docx_path.exists()


def test_pending_only_decision_pack_has_no_executive_storyline(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    material = core.ingest_note_text(
        case_dir,
        title="Advisor note",
        text="The proposed transition is still only a hypothesis.",
        now=fixed_now(),
    )
    core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "advisor_judgement",
                "text": "Pending judgement must not become narrative.",
                "source_material_ids": [material["id"]],
            }
        ],
        now=fixed_now(),
    )
    core.add_open_question(
        case_dir,
        question="Which evidence would prove the hypothesis?",
        why_it_matters="It determines whether the client pack can be written.",
        now=fixed_now(),
    )

    result = core.build_decision_pack(case_dir)
    markdown = result.markdown_path.read_text(encoding="utf-8")

    assert "Il pack cliente non e' ancora pronto." in markdown
    assert "Traccia esecutiva" not in markdown
    assert "Pending judgement must not become narrative." not in markdown


def test_case_brief_updates_from_case_state_and_review_status(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    material = core.ingest_note_text(
        case_dir,
        title="Advisor note",
        text="The transition should depend on formal governance.",
        now=fixed_now(),
    )

    entries = core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "fact",
                "text": "The current owner keeps final approval rights.",
                "status": "approved",
                "source_material_ids": [material["id"]],
                "reviewer": "Advisor",
            },
            {
                "kind": "advisor_judgement",
                "text": "The transition is risky without a formal veto rule.",
                "status": "pending",
                "source_material_ids": [material["id"]],
            },
        ],
        now=fixed_now(),
    )
    core.add_open_question(
        case_dir,
        question="Who can veto exceptions?",
        why_it_matters="It identifies the real control owner.",
        now=fixed_now(),
    )

    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert "Current Understanding (decision-pack ready)" in brief
    assert "The current owner keeps final approval rights." in brief
    assert "Candidate Review (not decision-pack ready)" in brief
    assert "The transition is risky without a formal veto rule." in brief
    assert "Pending judgement entries: 1" in brief
    assert "Who can veto exceptions?" in brief

    core.set_judgement_status(
        case_dir,
        entries[1]["id"],
        status="approved",
        reviewer="Advisor",
        now=fixed_now(),
    )
    updated_brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert "Decision-pack-ready judgement entries: 2" in updated_brief
    assert "Pending judgement entries: 0" in updated_brief
    assert "No pending judgement entries." in updated_brief


def test_set_judgement_statuses_updates_multiple_reviewed_entries(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)

    entries = core.add_judgement_entries(
        case_dir,
        [
            {"kind": "fact", "text": "The owner wants a decision pack."},
            {
                "kind": "advisor_judgement",
                "text": "The governance proposal needs conditions.",
            },
            {
                "kind": "decision_implication",
                "text": "The open conflict process still needs review.",
            },
        ],
        now=fixed_now(),
    )

    updated = core.set_judgement_statuses(
        case_dir,
        [entries[1]["id"], entries[0]["id"]],
        status="approved",
        reviewer="Advisor",
        review_note="Reviewed together.",
        now=fixed_now(),
    )
    judgement = json.loads((case_dir / "judgement_log.json").read_text())
    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert updated[0]["id"] == entries[1]["id"]
    assert updated[1]["id"] == entries[0]["id"]
    assert judgement["entries"][0]["status"] == "approved"
    assert judgement["entries"][0]["reviewer"] == "Advisor"
    assert judgement["entries"][0]["review_note"] == "Reviewed together."
    assert judgement["entries"][1]["status"] == "approved"
    assert judgement["entries"][2]["status"] == "pending"
    assert "Decision-pack-ready judgement entries: 2" in brief
    assert "Pending judgement entries: 1" in brief


def test_set_judgement_statuses_rejects_missing_id_without_partial_write(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)

    entries = core.add_judgement_entries(
        case_dir,
        [
            {"kind": "fact", "text": "A pending fact."},
            {"kind": "advisor_judgement", "text": "A pending judgement."},
        ],
        now=fixed_now(),
    )

    with pytest.raises(core.CaseWorkspaceError):
        core.set_judgement_statuses(
            case_dir,
            [entries[0]["id"], "jud-9999"],
            status="approved",
            reviewer="Advisor",
            now=fixed_now(),
        )
    judgement = json.loads((case_dir / "judgement_log.json").read_text())

    assert judgement["entries"][0]["status"] == "pending"
    assert judgement["entries"][0]["reviewed_at"] is None
    assert judgement["entries"][1]["status"] == "pending"


def test_bulk_approval_script_lists_summary_and_approves_all_pending(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    material = core.ingest_note_text(
        case_dir,
        title="Advisor note",
        text="The transition plan needs explicit governance tests.",
        now=fixed_now(),
    )
    core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "fact",
                "text": "The project is a 30-day decision-pack process.",
                "source_material_ids": [material["id"]],
            },
            {
                "kind": "advisor_judgement",
                "text": "The proposal should not proceed without a written mandate.",
                "source_material_ids": [material["id"]],
            },
        ],
        now=fixed_now(),
    )

    list_result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "approve_judgements.py"),
            str(case_dir),
        ],
        cwd=ROOT,
        capture_output=True,
        check=False,
        text=True,
    )
    list_output = list_result.stdout + list_result.stderr

    assert list_result.returncode == 0
    assert "Candidate decision-pack entries (2):" in list_output
    assert "1. [fact]" in list_output
    assert "1. jud-0001 [fact]" not in list_output
    assert "Sources: Advisor note." in list_output
    assert "Tell Clara: include all pending items." in list_output
    assert "Tell Clara: include item <number>" in list_output
    assert "python scripts/approve_judgements.py" not in list_output

    approve_result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "approve_judgements.py"),
            str(case_dir),
            "--all-pending",
            "--recorded-by",
            "Advisor",
            "--review-note",
            "Reviewed together.",
        ],
        cwd=ROOT,
        capture_output=True,
        check=False,
        text=True,
    )
    approve_output = approve_result.stdout + approve_result.stderr
    judgement = json.loads((case_dir / "judgement_log.json").read_text())

    assert approve_result.returncode == 0
    assert "Judgement entries marked: 2 status=approved." in approve_output
    assert judgement["entries"][0]["status"] == "approved"
    assert judgement["entries"][0]["reviewer"] == "Advisor"
    assert judgement["entries"][0]["review_note"] == "Reviewed together."
    assert judgement["entries"][1]["status"] == "approved"


def test_bulk_approval_script_updates_one_numbered_summary_item(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    core.add_judgement_entries(
        case_dir,
        [
            {"kind": "fact", "text": "A first pending entry."},
            {"kind": "advisor_judgement", "text": "A second pending entry."},
        ],
        now=fixed_now(),
    )

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "approve_judgements.py"),
            str(case_dir),
            "--item",
            "2",
            "--exclude",
            "--recorded-by",
            "Advisor",
            "--review-note",
            "Not accepted.",
        ],
        cwd=ROOT,
        capture_output=True,
        check=False,
        text=True,
    )
    output = result.stdout + result.stderr
    judgement = json.loads((case_dir / "judgement_log.json").read_text())

    assert result.returncode == 0
    assert "Judgement entries marked: 1 status=rejected." in output
    assert judgement["entries"][0]["status"] == "pending"
    assert judgement["entries"][1]["status"] == "rejected"
    assert judgement["entries"][1]["reviewer"] == "Advisor"
    assert judgement["entries"][1]["review_note"] == "Not accepted."


def test_inclusion_review_renders_status_checklist_without_mutating_entries(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    material = core.ingest_note_text(
        case_dir,
        title="Advisor note",
        text="The transition needs an explicit mandate.",
        now=fixed_now(),
    )
    entries = core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "fact",
                "text": "The owner requested a decision pack.",
                "source_material_ids": [material["id"]],
                "rationale": "Captured in the intake note.",
            },
            {
                "kind": "advisor_judgement",
                "text": "The mandate is ready for the client pack.",
                "status": "approved",
                "source_material_ids": [material["id"]],
                "reviewer": "Advisor",
            },
            {
                "kind": "decision_implication",
                "text": "The rejected path should stay out of the pack.",
                "status": "rejected",
                "source_material_ids": [material["id"]],
                "reviewer": "Advisor",
                "review_note": "Too speculative.",
            },
        ],
        now=fixed_now(),
    )

    result = core.build_inclusion_review(case_dir, now=fixed_now())
    review = result.review_path.read_text(encoding="utf-8")
    judgement = json.loads((case_dir / "judgement_log.json").read_text())

    assert result.review_path == case_dir / "inclusion_review.md"
    assert result.pending_count == 1
    assert result.approved_count == 1
    assert result.rejected_count == 1
    assert "# Revisione inclusione - ClientCo" in review
    assert "Voci pending che richiedono decisione del partner: 1" in review
    assert "### 1. fact" in review
    assert "jud-0001" not in review
    assert "## Come rispondere" in review
    assert "Risposta: includi voce 1, escludi voce 1" in review
    assert "includi tutte le voci pending" in review
    assert "python scripts/approve_judgements.py" not in review
    assert "Fonti: Advisor note" in review
    assert "Motivo: Captured in the intake note." in review
    assert "The mandate is ready for the client pack." in review
    assert "Too speculative." in review
    assert core.audit_human_visible_document_text(review) == []
    assert [entry["status"] for entry in judgement["entries"]] == [
        entries[0]["status"],
        entries[1]["status"],
        entries[2]["status"],
    ]


def test_inclusion_review_renders_semantic_approval_bundles(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    entries = core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "fact",
                "text": "Deleghe operative da formalizzare prima del passaggio.",
            },
            {
                "kind": "decision_implication",
                "text": "Le deleghe devono avere soglie e revoche esplicite.",
            },
            {
                "kind": "advisor_judgement",
                "text": "Lo scenario Jordan AD richiede un mandato separato.",
            },
        ],
        now=fixed_now(),
    )

    bundle_result = core.apply_inclusion_bundles(
        case_dir,
        [
            {
                "id": "deleghe",
                "title": "Deleghe",
                "description": "Deleghe e limiti operativi.",
                "entry_ids": [entries[0]["id"], entries[1]["id"]],
            },
            {
                "id": "scenario-jordan-ad",
                "title": "Scenario Jordan AD",
                "entry_ids": [entries[2]["id"]],
            },
        ],
        now=fixed_now(),
    )
    result = core.build_inclusion_review(case_dir, now=fixed_now())
    review = result.review_path.read_text(encoding="utf-8")
    judgement = json.loads((case_dir / "judgement_log.json").read_text())

    assert bundle_result.bundle_count == 2
    assert bundle_result.bundled_entry_count == 3
    assert "Pacchetti pending disponibili: 2" in review
    assert "### Pacchetto 1: Deleghe" in review
    assert "Voci: 1, 2" in review
    assert "ID: `deleghe`" in review
    assert "### Pacchetto 2: Scenario Jordan AD" in review
    assert "Risposta: includi pacchetto 2, escludi pacchetto 2" in review
    assert "### 1. fact" in review
    assert [entry["status"] for entry in judgement["entries"]] == [
        "pending",
        "pending",
        "pending",
    ]


def test_apply_inclusion_bundles_script_persists_semantic_plan(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    entries = core.add_judgement_entries(
        case_dir,
        [
            {"kind": "fact", "text": "A delegation fact."},
            {"kind": "advisor_judgement", "text": "A scenario judgement."},
        ],
        now=fixed_now(),
    )
    bundles_json = tmp_path / "bundles.json"
    bundles_json.write_text(
        json.dumps(
            {
                "bundles": [
                    {
                        "title": "Deleghe",
                        "entry_ids": [entries[0]["id"]],
                    },
                    {
                        "id": "scenario-jordan-ad",
                        "title": "Scenario Jordan AD",
                        "entry_ids": [entries[1]["id"]],
                    },
                ]
            }
        ),
        encoding="utf-8",
    )

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "apply_inclusion_bundles.py"),
            str(case_dir),
            "--bundles-json",
            str(bundles_json),
        ],
        cwd=ROOT,
        capture_output=True,
        check=False,
        text=True,
    )
    output = result.stdout + result.stderr
    bundle_payload = json.loads((case_dir / "inclusion_bundles.json").read_text())

    assert result.returncode == 0
    assert "Bundle count: 2" in output
    assert "Bundled entries: 2" in output
    assert bundle_payload["bundles"][0]["id"] == "deleghe"
    assert bundle_payload["bundles"][1]["id"] == "scenario-jordan-ad"


def test_bulk_approval_script_updates_numbered_bundle(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    entries = core.add_judgement_entries(
        case_dir,
        [
            {"kind": "fact", "text": "A delegation fact."},
            {"kind": "advisor_judgement", "text": "A delegation judgement."},
            {"kind": "decision_implication", "text": "A scenario implication."},
        ],
        now=fixed_now(),
    )
    core.apply_inclusion_bundles(
        case_dir,
        [
            {
                "id": "deleghe",
                "title": "Deleghe",
                "entry_ids": [entries[0]["id"], entries[1]["id"]],
            }
        ],
        now=fixed_now(),
    )

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "approve_judgements.py"),
            str(case_dir),
            "--bundle",
            "1",
            "--include",
            "--recorded-by",
            "Advisor",
            "--review-note",
            "Approved bundle.",
        ],
        cwd=ROOT,
        capture_output=True,
        check=False,
        text=True,
    )
    output = result.stdout + result.stderr
    judgement = json.loads((case_dir / "judgement_log.json").read_text())

    assert result.returncode == 0
    assert "Candidate inclusion bundles (1):" in output
    assert "Bundle 1. Deleghe (2 items): items 1, 2" in output
    assert "Judgement entries marked: 2 status=approved." in output
    assert [entry["status"] for entry in judgement["entries"]] == [
        "approved",
        "approved",
        "pending",
    ]
    assert judgement["entries"][0]["review_note"] == "Approved bundle."


def test_invalid_judgement_status_raises(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)

    with pytest.raises(core.CaseWorkspaceError):
        core.add_judgement_entries(
            case_dir,
            [{"kind": "fact", "text": "A fact.", "status": "draft"}],
            now=fixed_now(),
        )


def test_hosted_voice_launcher_carries_case_context(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    launcher = load_hosted_voice_launcher()
    core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "advisor_judgement",
                "text": "Morgan ExampleCo is worried the succession will stay informal.",
                "rationale": "The family has not agreed a clear veto rule.",
            }
        ],
        now=fixed_now(),
    )

    context = launcher.build_case_context(case_dir)
    launch_url = launcher.build_launch_url(
        "https://mparanza.com/case-notes/voice/launch",
        case_context=context,
    )
    encoded = parse_qs(urlsplit(launch_url).query)["case_context_z"][0]
    decoded_context = zlib.decompress(base64.urlsafe_b64decode(encoded)).decode("utf-8")

    assert "ClientCo" in decoded_context
    assert "Morgan ExampleCo is worried" in decoded_context
    assert "case_context_z=" in launch_url
    assert "mode=" not in launch_url


def test_hosted_voice_launcher_builds_transcription_launch_url(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    launcher = load_hosted_voice_launcher()

    launch_url, _context_limit = launcher.build_limited_launch_url(
        case_dir,
        base_url="https://mparanza.com/case-notes/voice/launch",
    )
    query = parse_qs(urlsplit(launch_url).query)
    decoded_context = zlib.decompress(
        base64.urlsafe_b64decode(query["case_context_z"][0])
    ).decode("utf-8")

    assert "ClientCo" in decoded_context
    assert "mode" not in query
    assert not (case_dir / "clara_kickoff_preparation.md").exists()


def test_hosted_voice_launcher_rejects_unsupported_voice_purpose(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    launcher = load_hosted_voice_launcher()

    with pytest.raises(launcher.CaseWorkspaceError, match="Unsupported voice purpose"):
        launcher.build_limited_launch_url(
            case_dir,
            base_url="https://mparanza.com/case-notes/voice/launch",
            purpose="unsupported_capture",
        )


def test_hosted_voice_launcher_rejects_kickoff_context(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    launcher = load_hosted_voice_launcher()

    with pytest.raises(launcher.CaseWorkspaceError, match="Unsupported voice purpose"):
        launcher.build_case_context(
            case_dir,
            purpose="kickoff",
            max_chars=2500,
        )


def test_hosted_voice_launcher_keeps_context_within_url_budget(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    launcher = load_hosted_voice_launcher()
    core.add_judgement_entries(
        case_dir,
        [
            {
                "kind": "codex_inference",
                "text": (
                    "Scenario "
                    f"{index}: validate decision rights, operating controls, "
                    f"risk register item {index}, and implementation trigger {index}."
                ),
            }
            for index in range(1, 35)
        ],
        now=fixed_now(),
    )

    launch_url, context_limit = launcher.build_limited_launch_url(
        case_dir,
        max_context_chars=2500,
        max_url_chars=1000,
    )

    assert len(launch_url) <= 1000
    assert context_limit <= 2500
    assert "case_context_z=" in launch_url


def test_hosted_voice_launcher_builds_clean_chrome_args(tmp_path: Path) -> None:
    launcher = load_hosted_voice_launcher()
    profile_dir = tmp_path / "chrome-profile"

    args = launcher.build_chrome_launch_args(
        "http://127.0.0.1:8000/case-notes/voice/launch?case_context_z=abc",
        profile_dir=profile_dir,
        remote_debugging_port=9224,
    )

    assert f"--user-data-dir={profile_dir}" in args
    assert "--no-first-run" in args
    assert "--remote-debugging-port=9224" in args
    assert "--use-fake-ui-for-media-stream" in args
    assert "--unsafely-treat-insecure-origin-as-secure=http://127.0.0.1:8000" in args
    assert args[-1].endswith("case_context_z=abc")


def test_hosted_voice_launcher_can_leave_chrome_microphone_prompt_manual(
    tmp_path: Path,
) -> None:
    launcher = load_hosted_voice_launcher()

    args = launcher.build_chrome_launch_args(
        "https://mparanza.com/case-notes/voice/launch",
        profile_dir=tmp_path / "chrome-profile",
        auto_accept_microphone=False,
    )

    assert "--use-fake-ui-for-media-stream" not in args


def test_start_deck_feedback_imports_new_bundle_and_records_html_target(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    starter = load_deck_feedback_starter()
    downloads_dir = tmp_path / "Downloads"
    downloads_dir.mkdir()
    target_deck = tmp_path / "client-deck.html"
    target_deck.write_text("<html><body>Deck</body></html>", encoding="utf-8")
    bundle_path = downloads_dir / "case-notes-voice-20260102.zip"
    wrote_bundle = False

    def complete_capture(_seconds: float) -> None:
        nonlocal wrote_bundle
        if wrote_bundle:
            return
        write_downloaded_voice_zip_bundle(
            bundle_path,
            captured_at="2026-01-02T10:30:00+00:00",
            transcript="Please strengthen the headline on this slide.",
            audio_file_name="deck-feedback.webm",
            video_file_name="deck-screen.webm",
            mtime=100,
        )
        wrote_bundle = True

    result = starter.start_deck_feedback(
        case_dir,
        target_deck_path=target_deck,
        downloads_dir=downloads_dir,
        open_browser=False,
        poll_seconds=0.01,
        timeout_seconds=1,
        sleep=complete_capture,
    )
    handoff = json.loads(result.handoff_path.read_text(encoding="utf-8"))

    assert result.target_kind == "html"
    assert result.selected_bundle_path == bundle_path
    assert result.import_result.raw_transcript_path.exists()
    assert result.deck_revision_intake_path is None
    assert handoff["source"] == "clara_deck_feedback_capture"
    assert handoff["status"] == "imported"
    assert handoff["target_deck"] == {
        "path": str(target_deck),
        "kind": "html",
    }
    assert handoff["screen_video_path"].endswith("deck-screen.webm")
    assert handoff["data_posture"]["hosted_capture_used"] is True


def test_dependency_checker_maps_core_and_ocr_package_imports() -> None:
    checker = load_dependency_checker()

    assert checker.requirement_name("Pillow>=11.1") == "pillow"
    assert checker.requirement_name("PyMuPDF>=1.24") == "pymupdf"
    assert checker.requirement_name("python-docx>=1.1") == "python-docx"
    assert checker.requirement_name("paddlepaddle>=3.0") == "paddlepaddle"
    assert checker.requirement_name("opencv-python>=4.8") == "opencv-python"
    assert checker.import_name("Pillow") == "PIL"
    assert checker.import_name("PyMuPDF") == "fitz"
    assert checker.import_name("imageio-ffmpeg") == "imageio_ffmpeg"
    assert checker.import_name("python-docx") == "docx"
    assert checker.import_name("paddlepaddle") == "paddle"
    assert checker.import_name("opencv-python") == "cv2"


def test_dependency_checker_selects_optional_ocr_requirements() -> None:
    checker = load_dependency_checker()

    selected = checker.selected_requirement_files(None, include_optional=True)
    selected_names = {path.name for path in selected}
    ocr_requirements = (PLUGIN_ROOT / "requirements-ocr.txt").read_text(
        encoding="utf-8"
    )

    assert "requirements.txt" in selected_names
    assert "requirements-ocr.txt" in selected_names
    assert "paddleocr" in ocr_requirements
    assert "paddlepaddle" in ocr_requirements
    assert "opencv-python" in ocr_requirements


def test_dependency_checker_checks_multiple_requirement_files(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    checker = load_dependency_checker()
    core_requirements = tmp_path / "requirements.txt"
    ocr_requirements = tmp_path / "requirements-ocr.txt"
    core_requirements.write_text("Pillow>=11.1\n", encoding="utf-8")
    ocr_requirements.write_text("paddlepaddle>=3.0\n", encoding="utf-8")
    checked_imports: list[str] = []

    def fake_find_spec(import_target: str) -> object:
        checked_imports.append(import_target)
        return object()

    monkeypatch.setattr(checker.importlib.util, "find_spec", fake_find_spec)

    missing = checker.check_dependencies([core_requirements, ocr_requirements])

    assert missing == []
    assert checked_imports == ["PIL", "paddle"]


def test_case_workspace_archive_excludes_local_runtime_dirs(tmp_path: Path) -> None:
    core, case_dir = init_case(tmp_path)
    runtime_dir = case_dir / ".codex_audit_reconciliation_py"
    runtime_dir.mkdir()
    runtime_file = runtime_dir / "library.py"
    runtime_file.write_text("local dependency cache", encoding="utf-8")
    (case_dir / ".DS_Store").write_text("local mac metadata", encoding="utf-8")
    note_dir = case_dir / "notes"
    note_dir.mkdir()
    note_path = note_dir / "advisor-note.md"
    note_path.write_text("Case-owned note.", encoding="utf-8")

    result = core.export_case_workspace_archive(case_dir, now=fixed_now())
    with ZipFile(result.package_path) as archive:
        archive_names = set(archive.namelist())

    assert result.package_path.parent == tmp_path / "case_share_exports"
    assert "case/case_manifest.json" in archive_names
    assert "case/case_brief.md" in archive_names
    assert "case/notes/advisor-note.md" in archive_names
    assert all(".codex_audit_reconciliation_py" not in name for name in archive_names)
    assert all(".DS_Store" not in name for name in archive_names)
    assert result.excluded_file_count == 2
    assert result.excluded_bytes >= runtime_file.stat().st_size


def test_support_package_includes_request_and_excludes_runtime_dirs(
    tmp_path: Path,
) -> None:
    core, case_dir = init_case(tmp_path)
    runtime_dir = case_dir / ".codex_audit_reconciliation_py"
    runtime_dir.mkdir()
    runtime_file = runtime_dir / "ocr.py"
    runtime_file.write_text("local OCR dependency", encoding="utf-8")
    (case_dir / ".DS_Store").write_text("local mac metadata", encoding="utf-8")
    note_dir = case_dir / "notes"
    note_dir.mkdir()
    note_path = note_dir / "advisor-note.md"
    note_path.write_text("The draft slides are too generic.", encoding="utf-8")

    result = core.prepare_support_package(
        case_dir,
        request="The slides are not good enough; the support reviewer should improve them.",
        requested_by="Advisor",
        now=fixed_now(),
    )
    with ZipFile(result.package_path) as archive:
        archive_names = set(archive.namelist())
        support_note = archive.read("case/support_request.md").decode("utf-8")

    assert result.package_path.parent == tmp_path / "case_support_exports"
    assert result.support_request_archive_path == "case/support_request.md"
    assert "case/case_manifest.json" in archive_names
    assert "case/case_brief.md" in archive_names
    assert "case/notes/advisor-note.md" in archive_names
    assert "case/support_request.md" in archive_names
    assert "The slides are not good enough" in support_note
    assert "Recipient: Support reviewer" in support_note
    assert "Requested by: Advisor" in support_note
    assert (
        "The requester's local Clara folder remains the authoritative" in support_note
    )
    assert all(".codex_audit_reconciliation_py" not in name for name in archive_names)
    assert all(".DS_Store" not in name for name in archive_names)
    assert result.excluded_file_count == 2
    assert result.excluded_bytes >= runtime_file.stat().st_size


def test_support_package_script_writes_default_package_path(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)

    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS_DIR / "prepare_support_package.py"),
            str(case_dir),
            "--request",
            "Clara needs a support reviewer to improve the output deck.",
            "--requested-by",
            "Advisor",
        ],
        cwd=ROOT,
        capture_output=True,
        check=False,
        text=True,
    )
    output = result.stdout + result.stderr

    assert result.returncode == 0
    assert "Prepared Clara support package:" in output
    assert "Support note: case/support_request.md" in output
    assert (tmp_path / "case_support_exports").exists()


def test_import_hosted_voice_bundle_adds_local_pending_judgement(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    importer = load_hosted_voice_importer()
    bundle_path = tmp_path / "case-notes-voice.json"
    bundle_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "case_notes_hosted_voice",
                "captured_at": "2026-01-02T10:30:00+00:00",
                "model": "gpt-realtime-2",
                "user_transcript": "The informal governance risk is material.",
                "assistant_transcript": "Which condition would reduce it?",
                "extraction_json": {
                    "cleaned_notes_markdown": "# Notes\n\nGovernance risk.",
                    "entries": [
                        {
                            "kind": "advisor_judgement",
                            "text": "Informal governance is a material risk.",
                            "rationale": "It affects implementation credibility.",
                        },
                        {"kind": "unsupported", "text": "Ignore this."},
                    ],
                    "open_questions": [
                        {
                            "question": "Who can veto governance exceptions?",
                            "why_it_matters": "It defines the real control owner.",
                        }
                    ],
                },
            }
        ),
        encoding="utf-8",
    )

    result = importer.import_hosted_voice_bundle(case_dir, bundle_path)
    registry = json.loads((case_dir / "material_registry.json").read_text())
    judgement = json.loads((case_dir / "judgement_log.json").read_text())
    questions = json.loads((case_dir / "open_questions.json").read_text())

    assert result.judgement_count == 1
    assert result.open_question_count == 1
    assert result.discussion_review_pack_path.exists()
    assert result.clara_review_path.exists()
    assert registry["materials"][0]["material_type"] == "transcript"
    assert judgement["entries"][0]["status"] == "pending"
    assert judgement["entries"][0]["source_material_ids"] == [result.material_id]
    assert judgement["entries"][0]["text"] == "Informal governance is a material risk."
    assert (
        questions["questions"][0]["question"] == "Who can veto governance exceptions?"
    )
    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")
    assert "Informal governance is a material risk." in brief
    assert "Pending judgement entries: 1" in brief
    assert "Who can veto governance exceptions?" in brief
    review_pack = result.discussion_review_pack_path.read_text(encoding="utf-8")
    clara_review = result.clara_review_path.read_text(encoding="utf-8")
    assert "Local Codex Discussion Review Pack" in review_pack
    assert "The server handled hosted audio processing only." in review_pack
    assert "Informal governance is a material risk." in review_pack
    assert f"Clara review: `{result.clara_review_path.relative_to(case_dir)}`" in (
        review_pack
    )
    assert "Clara Audio Review" in clara_review
    assert "Generated locally by the Clara plugin" in clara_review
    assert "## Required Transcript Processing" in clara_review
    assert "assign speaker attribution" in clara_review
    assert "check transcript quality" in clara_review
    assert "obviously wrong transcription words" in clara_review
    assert "Pending local Clara review." in clara_review
    assert "The informal governance risk is material." in clara_review
    assert (
        f"Source proposed entries to material id `{result.material_id}`" in review_pack
    )
    assert "First assign speakers" in review_pack
    assert "correct only obvious transcription errors" in review_pack
    assert "Do not mark anything decision-pack ready." in review_pack


def test_find_latest_hosted_voice_bundle_skips_imported_and_invalid(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    latest = load_latest_hosted_voice_importer()
    downloads_dir = tmp_path / "Downloads"
    downloads_dir.mkdir()
    older = downloads_dir / "case-notes-voice-older.json"
    selected = downloads_dir / "case-notes-audio-selected.zip"
    already_imported = downloads_dir / "case-notes-voice-newer.json"
    invalid = downloads_dir / "case-notes-voice-invalid.json"

    write_downloaded_voice_bundle(
        older,
        captured_at="2026-01-02T10:30:00+00:00",
        transcript="Older transcript.",
        mtime=100,
    )
    write_downloaded_voice_zip_bundle(
        selected,
        captured_at="2026-01-03T10:30:00+00:00",
        transcript="Selected transcript.",
        mtime=200,
    )
    write_downloaded_voice_bundle(
        already_imported,
        captured_at="2026-01-04T10:30:00+00:00",
        transcript="Already imported transcript.",
        mtime=300,
    )
    write_downloaded_voice_bundle(
        invalid,
        captured_at="2026-01-05T10:30:00+00:00",
        transcript="Wrong source transcript.",
        source="other_source",
        mtime=400,
    )
    (case_dir / "voice_sessions" / "20260104103000Z").mkdir(parents=True)

    found = latest.find_latest_hosted_voice_bundle(
        case_dir,
        downloads_dir=downloads_dir,
    )
    found_with_imported = latest.find_latest_hosted_voice_bundle(
        case_dir,
        downloads_dir=downloads_dir,
        include_imported=True,
    )

    assert found.path == selected
    assert found.already_imported is False
    assert found_with_imported.path == already_imported
    assert found_with_imported.already_imported is True


def test_import_latest_hosted_voice_bundle_imports_downloaded_transcript(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    latest = load_latest_hosted_voice_importer()
    downloads_dir = tmp_path / "Downloads"
    downloads_dir.mkdir()
    bundle_path = downloads_dir / "case-notes-voice-20260102.zip"
    write_downloaded_voice_zip_bundle(
        bundle_path,
        captured_at="2026-01-02T10:30:00+00:00",
        transcript="Latest downloaded transcript for Clara.",
        audio_file_name="latest-capture.wav",
        mtime=100,
    )

    result = latest.import_latest_hosted_voice_bundle(
        case_dir,
        downloads_dir=downloads_dir,
    )
    raw_transcript = result.raw_transcript_path.read_text(encoding="utf-8")
    clara_review = result.clara_review_path.read_text(encoding="utf-8")

    assert result.raw_transcript_path.exists()
    assert result.audio_path == result.session_dir / "latest-capture.wav"
    assert result.audio_path.read_bytes() == b"audio bytes"
    assert "Latest downloaded transcript for Clara." in raw_transcript
    assert "Generated locally by the Clara plugin" in clara_review
    assert "Pending local Clara review." in clara_review


def test_import_uploaded_audio_bundle_adds_pending_knowledge(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    importer = load_hosted_voice_importer()
    bundle_path = tmp_path / "case-notes-audio.json"
    companion_audio_path = tmp_path / "family-meeting.wav"
    companion_audio_path.write_bytes(b"audio bytes")
    bundle_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "case_notes_hosted_voice",
                "capture_source": "uploaded_audio",
                "captured_at": "2026-01-02T10:30:00+00:00",
                "source_metadata": {
                    "source_type": "interview",
                    "title": "ExampleCo - CFO interview",
                    "interview_date": "2026-06-20",
                    "participants": "CFO",
                    "role": "Finance",
                    "interviewer": "Reviewer",
                    "notes": "First management interview.",
                },
                "model": "gpt-5.5",
                "transcription_model": "gpt-4o-transcribe",
                "audio_file_name": "family-meeting.wav",
                "audio_content_type": "audio/wav",
                "user_transcript": "The meeting shows unclear decision rights.",
                "transcript_text_prompted": (
                    "The meeting shows unclear decision rights."
                ),
                "assistant_transcript": "",
                "extraction_json": {
                    "cleaned_notes_markdown": "# Meeting Notes\n\nDecision rights unclear.",
                    "entries": [
                        {
                            "kind": "codex_inference",
                            "text": "Decision rights need validation after the meeting.",
                            "rationale": "The uploaded meeting audio was ambiguous.",
                        }
                    ],
                    "open_questions": [
                        {
                            "question": "Who has final authority after the meeting?",
                            "why_it_matters": "It determines whether the governance path is executable.",
                        }
                    ],
                },
                "extraction_text": "{}",
            }
        ),
        encoding="utf-8",
    )

    result = importer.import_hosted_voice_bundle(case_dir, bundle_path)
    raw_transcript = result.raw_transcript_path.read_text(encoding="utf-8")
    clara_review = result.clara_review_path.read_text(encoding="utf-8")
    attribution_report = json.loads(
        result.speaker_attribution_report_path.read_text(encoding="utf-8")
    )
    attribution_task = result.speaker_attribution_task_path.read_text(encoding="utf-8")
    registry = json.loads((case_dir / "material_registry.json").read_text())
    judgement = json.loads((case_dir / "judgement_log.json").read_text())
    questions = json.loads((case_dir / "open_questions.json").read_text())

    assert result.judgement_count == 1
    assert result.audio_path == result.session_dir / "family-meeting.wav"
    assert result.audio_path.read_bytes() == b"audio bytes"
    assert "Capture source: uploaded_audio" in raw_transcript
    assert "Audio file: family-meeting.wav" in raw_transcript
    assert "Source title: ExampleCo - CFO interview" in raw_transcript
    assert "Interviewee / participants: CFO" in raw_transcript
    assert "Role / function" not in raw_transcript
    assert clara_review.startswith("# Clara Audio Review")
    assert "## Required Transcript Processing" in clara_review
    assert "assign speaker attribution" in clara_review
    assert "check transcript quality" in clara_review
    assert "## Well-Supported Points" in clara_review
    assert "Generated locally by the Clara plugin" in clara_review
    assert "Pending local Clara review." in clara_review
    assert "The meeting shows unclear decision rights." in clara_review
    assert registry["materials"][0]["title"] == "ExampleCo - CFO interview"
    assert registry["materials"][0]["source_metadata"]["interviewer"] == "Reviewer"
    assert "role" not in registry["materials"][0]["source_metadata"]
    assert "confidentiality" not in registry["materials"][0]["source_metadata"]
    assert result.attributed_transcript_path is None
    assert result.speaker_attribution_task_path == (
        result.session_dir / "speaker_attribution_task.md"
    )
    assert "Write `attributed_transcript.md`" in attribution_task
    assert "Candidate speaker names: CFO, Reviewer" in attribution_task
    assert "Speaker 1" in attribution_task
    assert "Speaker 2" in attribution_task
    assert attribution_report["status"] == "needs_model_attribution"
    assert attribution_report["method"] == "requires_clara_codex_attribution"
    assert attribution_report["speaker_labels"] == []
    assert attribution_report["candidate_speaker_names"] == ["CFO", "Reviewer"]
    assert attribution_report["speaker_labels_are_provisional"] is False
    assert attribution_report["requires_clara_codex_attribution"] is True
    assert attribution_report["requires_review"] is True
    assert registry["materials"][0]["path"] == str(result.raw_transcript_path.resolve())
    assert "speaker_attribution" not in registry["materials"][0]["source_metadata"]
    assert "attributed_transcript" not in registry["materials"][0]["source_metadata"]
    assert registry["materials"][0]["source_metadata"][
        "speaker_attribution_report"
    ].endswith("speaker_attribution_report.json")
    assert registry["materials"][0]["source_metadata"][
        "speaker_attribution_task"
    ].endswith("speaker_attribution_task.md")
    assert judgement["entries"][0]["status"] == "pending"
    assert judgement["entries"][0]["kind"] == "codex_inference"
    assert (
        questions["questions"][0]["question"]
        == "Who has final authority after the meeting?"
    )


def test_import_hosted_voice_bundle_auto_attributes_single_metadata_speaker(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    importer = load_hosted_voice_importer()
    bundle_path = tmp_path / "case-notes-voice.json"
    bundle_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "case_notes_hosted_voice",
                "capture_source": "live_capture",
                "captured_at": "2026-01-02T10:30:00+00:00",
                "source_metadata": {
                    "title": "Deck correction walkthrough",
                    "participants": "Facilitator",
                    "notes": "Facilitator is the only person speaking.",
                },
                "model": "gpt-realtime-whisper",
                "transcription_model": "gpt-4o-transcribe",
                "user_transcript": "This slide needs a stronger headline.",
                "assistant_transcript": "",
                "extraction_json": {"cleaned_notes_markdown": ""},
            }
        ),
        encoding="utf-8",
    )

    result = importer.import_hosted_voice_bundle(case_dir, bundle_path)
    attribution = result.attributed_transcript_path.read_text(encoding="utf-8")
    report = json.loads(
        result.speaker_attribution_report_path.read_text(encoding="utf-8")
    )
    registry = json.loads((case_dir / "material_registry.json").read_text())
    transcript = registry["materials"][0]
    backup_path = result.session_dir / "raw_transcript_unattributed.md"
    review_pack = result.discussion_review_pack_path.read_text(encoding="utf-8")

    assert result.attributed_transcript_path == (
        result.session_dir / "attributed_transcript.md"
    )
    assert "Facilitator: This slide needs a stronger headline." in attribution
    assert backup_path.exists()
    assert report["status"] == "attributed"
    assert report["method"] == "metadata_single_speaker"
    assert report["speaker_labels"] == ["Facilitator"]
    assert report["requires_review"] is False
    assert result.speaker_attribution_task_path is None
    assert transcript["path"] == str(result.attributed_transcript_path.resolve())
    assert (
        transcript["source_metadata"]["speaker_attribution"]
        == "single-speaker local text attribution using metadata_single_speaker; "
        "no audio or voice diarization model used"
    )
    assert transcript["source_metadata"]["attributed_transcript"].endswith(
        "attributed_transcript.md"
    )
    assert transcript["source_metadata"]["speaker_attribution_report"].endswith(
        "speaker_attribution_report.json"
    )
    assert (
        "Attributed transcript: `voice_sessions/20260102103000Z/attributed_transcript.md`"
        in review_pack
    )
    assert (
        "Speaker attribution report: `voice_sessions/20260102103000Z/speaker_attribution_report.json`"
        in review_pack
    )


def test_import_hosted_voice_bundle_requires_clara_attribution_without_metadata(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    importer = load_hosted_voice_importer()
    bundle_path = tmp_path / "case-notes-voice.json"
    bundle_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "case_notes_hosted_voice",
                "capture_source": "live_capture",
                "captured_at": "2026-01-02T10:30:00+00:00",
                "source_metadata": {"title": "Unlabelled deck correction"},
                "model": "gpt-realtime-whisper",
                "transcription_model": "gpt-4o-transcribe",
                "user_transcript": "This page needs a clearer title.",
                "assistant_transcript": "",
                "extraction_json": {"cleaned_notes_markdown": ""},
            }
        ),
        encoding="utf-8",
    )

    result = importer.import_hosted_voice_bundle(case_dir, bundle_path)
    report = json.loads(
        result.speaker_attribution_report_path.read_text(encoding="utf-8")
    )
    attribution_task = result.speaker_attribution_task_path.read_text(encoding="utf-8")
    registry = json.loads((case_dir / "material_registry.json").read_text())
    transcript = registry["materials"][0]

    assert result.attributed_transcript_path is None
    assert result.speaker_attribution_task_path == (
        result.session_dir / "speaker_attribution_task.md"
    )
    assert "This page needs a clearer title." in attribution_task
    assert "If real names are unknown" in attribution_task
    assert report["status"] == "needs_model_attribution"
    assert report["method"] == "requires_clara_codex_attribution"
    assert report["speaker_labels"] == []
    assert report["candidate_speaker_names"] == []
    assert report["speaker_labels_are_provisional"] is False
    assert report["requires_clara_codex_attribution"] is True
    assert report["requires_review"] is True
    assert transcript["path"] == str(result.raw_transcript_path.resolve())
    assert "attributed_transcript" not in transcript["source_metadata"]
    assert transcript["source_metadata"]["speaker_attribution_task"].endswith(
        "speaker_attribution_task.md"
    )


def test_finalize_hosted_transcript_records_local_attribution_and_audio_pointer(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    importer = load_hosted_voice_importer()
    finalizer = load_hosted_transcript_finalizer()
    bundle_path = tmp_path / "case-notes-audio.json"
    bundle_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "case_notes_hosted_voice",
                "capture_source": "uploaded_audio",
                "captured_at": "2026-01-02T10:30:00+00:00",
                "source_metadata": {
                    "source_type": "interview",
                    "title": "Family interview",
                    "interview_date": "2026-01-02",
                    "participants": "Interviewee; Advisor",
                    "interviewer": "Advisor",
                },
                "audio_file_name": "family.m4a",
                "user_transcript": "Question? Answer.",
                "extraction_json": {"cleaned_notes_markdown": ""},
            }
        ),
        encoding="utf-8",
    )
    imported = importer.import_hosted_voice_bundle(case_dir, bundle_path)
    attributed_path = imported.session_dir / "raw_transcript_rule_attributed.md"
    attributed_path.write_text(
        "# Hosted Voice Transcript\n\n"
        "## Speaker-Attributed Transcript\n\n"
        "[001] Advisor: Question?\n\n"
        "[002] Interviewee: Answer.\n",
        encoding="utf-8",
    )
    audio_pointer_path = case_dir / "source_materials" / "interviews" / "audio.md"
    audio_pointer_path.parent.mkdir(parents=True)
    audio_pointer_path.write_text(
        "# Audio pointer\n\n"
        "- Stato trascrizione: non ancora trascritto in Clara\n\n"
        "Quando sara' prodotta la trascrizione, conservarla in `notes/`.\n",
        encoding="utf-8",
    )
    raw_text = imported.raw_transcript_path.read_text(encoding="utf-8")

    result = finalizer.finalize_hosted_transcript(
        case_dir,
        imported.material_id,
        attributed_path,
        audio_pointer_path=audio_pointer_path,
        audio_pointer_title="Audio interview pointer",
        now=fixed_now(),
    )

    registry = json.loads((case_dir / "material_registry.json").read_text())
    transcript = next(
        material
        for material in registry["materials"]
        if material["id"] == imported.material_id
    )
    pointer = next(
        material
        for material in registry["materials"]
        if material["id"] == result.audio_pointer_material_id
    )
    backup_path = imported.session_dir / "raw_transcript_unattributed.md"
    brief = (case_dir / "case_brief.md").read_text(encoding="utf-8")

    assert result.unattributed_transcript_backup_path == backup_path
    assert backup_path.read_text(encoding="utf-8") == raw_text
    assert transcript["path"] == str(attributed_path.resolve())
    assert transcript["material_type"] == "transcript"
    assert (
        transcript["source_metadata"]["speaker_attribution"]
        == finalizer.DEFAULT_SPEAKER_ATTRIBUTION_NOTE
    )
    assert (
        transcript["source_metadata"]["unattributed_transcript_backup"]
        == "voice_sessions/20260102103000Z/raw_transcript_unattributed.md"
    )
    assert (
        transcript["source_metadata"]["raw_audio_pointer_material_id"] == pointer["id"]
    )
    assert pointer["path"] == str(audio_pointer_path.resolve())
    assert pointer["source_metadata"]["transcription_status"] == "transcribed"
    assert (
        pointer["source_metadata"]["linked_transcript_material_id"]
        == imported.material_id
    )
    assert (
        pointer["source_metadata"]["linked_transcript_path"]
        == "voice_sessions/20260102103000Z/raw_transcript_rule_attributed.md"
    )
    assert pointer["source_metadata"]["transcribed_at"] == fixed_now().isoformat()
    assert "Transcribed in Clara" in pointer["summary"]
    pointer_text = audio_pointer_path.read_text(encoding="utf-8")
    assert "- Stato trascrizione: trascritto in Clara" in pointer_text
    assert f"- Materiale trascrizione: `{imported.material_id}`" in pointer_text
    assert (
        "- Trascrizione collegata: "
        "`voice_sessions/20260102103000Z/raw_transcript_rule_attributed.md`"
        in pointer_text
    )
    assert "non ancora trascritto" not in pointer_text
    assert "Quando sara' prodotta la trascrizione" not in pointer_text
    assert "Family interview [transcript, indexed]" in brief
    assert "Audio interview pointer [source, indexed]" in brief

    result_again = finalizer.finalize_hosted_transcript(
        case_dir,
        imported.material_id,
        attributed_path,
        audio_pointer_path=audio_pointer_path,
        audio_pointer_title="Audio interview pointer",
        now=fixed_now(),
    )
    registry_again = json.loads((case_dir / "material_registry.json").read_text())

    assert result_again.audio_pointer_material_id == pointer["id"]
    assert [material["path"] for material in registry_again["materials"]].count(
        str(audio_pointer_path.resolve())
    ) == 1
    assert (
        audio_pointer_path.read_text(encoding="utf-8").count("## Trascrizione Clara")
        == 1
    )


def test_import_hosted_voice_bundle_uses_explicit_companion_audio_path(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    importer = load_hosted_voice_importer()
    bundle_path = tmp_path / "case-notes-audio.json"
    source_audio_path = tmp_path / "recordings" / "family-meeting.wav"
    source_audio_path.parent.mkdir()
    source_audio_path.write_bytes(b"source audio bytes")
    bundle_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "case_notes_hosted_voice",
                "capture_source": "uploaded_audio",
                "captured_at": "2026-01-02T10:30:00+00:00",
                "audio_file_name": "family-meeting.wav",
                "user_transcript": "Uploaded meeting transcript.",
                "extraction_json": {"cleaned_notes_markdown": ""},
            }
        ),
        encoding="utf-8",
    )

    result = importer.import_hosted_voice_bundle(
        case_dir,
        bundle_path,
        companion_audio_path=source_audio_path,
    )

    assert result.audio_path == result.session_dir / "family-meeting.wav"
    assert result.audio_path.read_bytes() == b"source audio bytes"


def test_upload_hosted_audio_saves_and_imports_bundle(
    tmp_path: Path,
    monkeypatch,
) -> None:
    _, case_dir = init_case(tmp_path)
    uploader = load_hosted_audio_uploader()
    audio_path = tmp_path / "morgan.m4a"
    audio_path.write_bytes(b"audio bytes")
    run_dir = tmp_path / "hosted-run"
    captured: dict[str, Any] = {}

    monkeypatch.setattr(uploader, "_new_opener", lambda: object())

    def fake_authenticate_with_magic_link(
        _opener,
        magic_link,
        *,
        launch_url,
        timeout_seconds,
    ):
        captured["auth"] = {
            "magic_link": magic_link,
            "launch_url": launch_url,
            "timeout_seconds": timeout_seconds,
        }
        return f"launch:{magic_link}"

    monkeypatch.setattr(
        uploader,
        "authenticate_with_magic_link",
        fake_authenticate_with_magic_link,
    )

    def fake_upload_audio_file(*_args, **kwargs):
        captured["upload"] = kwargs
        return {"status": "queued", "job_id": "job-123"}

    def fake_poll_upload_job(*_args, **kwargs):
        captured["poll"] = kwargs
        return {
            "status": "done",
            "bundle": {
                "schema_version": 1,
                "source": "case_notes_hosted_voice",
                "capture_source": "uploaded_audio",
                "captured_at": "2026-01-02T10:30:00+00:00",
                "audio_file_name": "morgan.m4a",
                "user_transcript": "Transcript from hosted audio.",
                "extraction_json": {"cleaned_notes_markdown": ""},
            },
        }

    class FakeImportResult:
        material_id = "mat-test"
        session_dir = case_dir / "voice_sessions" / "20260102103000Z"
        audio_path = session_dir / "morgan.m4a"

    def fake_import_hosted_voice_bundle(*args, **kwargs):
        captured["import_args"] = args
        captured["import_kwargs"] = kwargs
        return FakeImportResult()

    monkeypatch.setattr(uploader, "upload_audio_file", fake_upload_audio_file)
    monkeypatch.setattr(uploader, "poll_upload_job", fake_poll_upload_job)
    monkeypatch.setattr(
        uploader,
        "import_hosted_voice_bundle",
        fake_import_hosted_voice_bundle,
    )

    result = uploader.upload_hosted_audio(
        case_dir=case_dir,
        audio_path=audio_path,
        magic_link="https://mparanza.com/auth/magic/consume?token=test",
        output_dir=run_dir,
        source_metadata=uploader.build_source_metadata(
            source_type="interview",
            title="Morgan interview",
            participants="Morgan",
        ),
        poll_seconds=30,
        poll_interval_seconds=0.01,
    )

    bundle = json.loads(result.bundle_path.read_text(encoding="utf-8"))
    upload_kwargs = captured["upload"]

    assert bundle["user_transcript"] == "Transcript from hosted audio."
    assert result.run_dir == run_dir
    assert result.import_result.material_id == "mat-test"
    assert upload_kwargs["launch_token"].startswith("launch:")
    assert upload_kwargs["audio_path"] == audio_path
    assert upload_kwargs["source_metadata"]["title"] == "Morgan interview"
    auth_query = parse_qs(urlsplit(captured["auth"]["launch_url"]).query)
    decoded_context = zlib.decompress(
        base64.urlsafe_b64decode(auth_query["case_context_z"][0])
    ).decode("utf-8")
    assert "mode" not in auth_query
    assert "ClientCo" in decoded_context
    assert captured["poll"]["job_id"] == "job-123"
    assert captured["import_args"] == (case_dir, result.bundle_path)
    assert captured["import_kwargs"]["companion_audio_path"] == audio_path


def test_upload_hosted_audio_refreshes_context_before_launch_url(
    tmp_path: Path,
    monkeypatch,
) -> None:
    _, case_dir = init_case(tmp_path)
    uploader = load_hosted_audio_uploader()
    (case_dir / "case_brief.md").write_text("STALE brief", encoding="utf-8")
    (case_dir / "clara_kickoff_preparation.md").write_text(
        "STALE kickoff prep",
        encoding="utf-8",
    )
    calls: list[str] = []

    def fake_refresh_case_brief(refreshed_case_dir: Path) -> None:
        assert refreshed_case_dir == case_dir
        calls.append("brief")
        (refreshed_case_dir / "case_brief.md").write_text(
            "FRESH brief for ClientCo",
            encoding="utf-8",
        )

    monkeypatch.setattr(uploader, "refresh_case_brief", fake_refresh_case_brief)

    debrief_url = uploader._build_context_launch_url(
        case_dir,
        base_url="https://mparanza.com",
    )
    debrief_query = parse_qs(urlsplit(debrief_url).query)
    debrief_context = zlib.decompress(
        base64.urlsafe_b64decode(debrief_query["case_context_z"][0])
    ).decode("utf-8")

    assert calls == ["brief"]
    assert "FRESH brief for ClientCo" in debrief_context
    assert "STALE brief" not in debrief_context
    assert "mode" not in debrief_query


def test_upload_hosted_audio_can_reuse_session_cookie(
    tmp_path: Path,
    monkeypatch,
) -> None:
    _, case_dir = init_case(tmp_path)
    uploader = load_hosted_audio_uploader()
    audio_path = tmp_path / "morgan.m4a"
    audio_path.write_bytes(b"audio bytes")
    captured: dict[str, Any] = {}

    monkeypatch.setattr(uploader, "_new_opener", lambda: object())

    def fake_authenticate_with_session_cookie(*_args, **kwargs):
        captured["auth"] = kwargs
        return "launch-from-cookie"

    def fake_upload_audio_file(*_args, **kwargs):
        captured["upload"] = kwargs
        return {"status": "queued", "job_id": "job-cookie"}

    def fake_poll_upload_job(*_args, **kwargs):
        captured["poll"] = kwargs
        return {
            "status": "done",
            "bundle": {
                "schema_version": 1,
                "source": "case_notes_hosted_voice",
                "capture_source": "uploaded_audio",
                "captured_at": "2026-01-02T10:30:00+00:00",
                "audio_file_name": "morgan.m4a",
                "user_transcript": "Transcript from cookie-authenticated upload.",
                "extraction_json": {"cleaned_notes_markdown": ""},
            },
        }

    monkeypatch.setattr(
        uploader,
        "authenticate_with_session_cookie",
        fake_authenticate_with_session_cookie,
    )
    monkeypatch.setattr(uploader, "upload_audio_file", fake_upload_audio_file)
    monkeypatch.setattr(uploader, "poll_upload_job", fake_poll_upload_job)

    result = uploader.upload_hosted_audio(
        case_dir=case_dir,
        audio_path=audio_path,
        cookie_header="auth_session=test-cookie",
        source_metadata={"title": "Morgan interview"},
        import_bundle=False,
    )

    bundle = json.loads(result.bundle_path.read_text(encoding="utf-8"))

    assert bundle["user_transcript"] == "Transcript from cookie-authenticated upload."
    assert result.import_result is None
    assert captured["auth"]["cookie_header"] == "auth_session=test-cookie"
    auth_query = parse_qs(urlsplit(captured["auth"]["launch_url"]).query)
    decoded_context = zlib.decompress(
        base64.urlsafe_b64decode(auth_query["case_context_z"][0])
    ).decode("utf-8")
    assert "mode" not in auth_query
    assert "ClientCo" in decoded_context
    assert captured["upload"]["launch_token"] == "launch-from-cookie"
    assert captured["poll"]["job_id"] == "job-cookie"


def test_hosted_audio_cookie_header_strips_optional_cookie_prefix() -> None:
    uploader = load_hosted_audio_uploader()
    opener = uploader._new_opener()

    uploader._set_cookie_header(opener, "Cookie: auth_session=test-cookie")

    assert opener.addheaders[-1] == ("Cookie", "auth_session=test-cookie")

    uploader._set_cookie_header(opener, "auth_session=next-cookie")

    assert opener.addheaders[-1] == ("Cookie", "auth_session=next-cookie")
    assert [
        header for header in opener.addheaders if header[0].lower() == "cookie"
    ] == [("Cookie", "auth_session=next-cookie")]


def test_hosted_audio_magic_link_can_relaunch_with_case_context() -> None:
    uploader = load_hosted_audio_uploader()
    seen_urls: list[str] = []

    class FakeResponse:
        def __init__(self, final_url: str) -> None:
            self.final_url = final_url

        def __enter__(self):
            return self

        def __exit__(self, *_args) -> None:
            return None

        def geturl(self) -> str:
            return self.final_url

        def read(self) -> bytes:
            return b""

    class FakeOpener:
        def open(self, request, timeout):
            seen_urls.append(request.full_url)
            if len(seen_urls) == 1:
                return FakeResponse("https://mparanza.com/case-notes/voice")
            return FakeResponse(
                "https://mparanza.com/case-notes/voice?session=context-token"
            )

    token = uploader.authenticate_with_magic_link(
        FakeOpener(),
        "https://mparanza.com/auth/magic/consume?token=test",
        launch_url=(
            "https://mparanza.com/case-notes/voice/launch?" "case_context_z=abc"
        ),
    )

    assert token == "context-token"
    assert seen_urls == [
        "https://mparanza.com/auth/magic/consume?token=test",
        ("https://mparanza.com/case-notes/voice/launch?" "case_context_z=abc"),
    ]


def test_upload_hosted_audio_retries_large_recording_as_chunks_after_413(
    tmp_path: Path,
    monkeypatch,
) -> None:
    uploader = load_hosted_audio_uploader()
    audio_path = tmp_path / "morgan.m4a"
    audio_path.write_bytes(b"abcdefghijk")
    calls: list[tuple[str, bytes]] = []
    chunk_bodies: list[bytes] = []

    monkeypatch.setattr(uploader, "CHUNKED_AUDIO_UPLOAD_CHUNK_BYTES", 5)

    def fake_read_json_response(_opener, request, *, timeout_seconds):
        body = request.data.read() if request.data is not None else b""
        calls.append((request.full_url, body))
        if request.full_url == "https://example.test/case-notes/api/voice/upload":
            return uploader.JsonResponse(413, {"detail": "request body too large"})
        if request.full_url.endswith("/case-notes/api/voice/upload/chunks/start"):
            assert b'name="total_bytes"\r\n\r\n11' in body
            assert b'name="total_chunks"\r\n\r\n3' in body
            return uploader.JsonResponse(
                201,
                {
                    "status": "ready",
                    "upload_id": "upload-123",
                    "chunk_size": 5,
                },
            )
        if request.full_url.endswith("/case-notes/api/voice/upload/chunks/upload-123"):
            chunk_bodies.append(body)
            return uploader.JsonResponse(200, {"status": "received"})
        if request.full_url.endswith(
            "/case-notes/api/voice/upload/chunks/upload-123/finish"
        ):
            return uploader.JsonResponse(
                202,
                {"status": "queued", "job_id": "job-chunk"},
            )
        raise AssertionError(f"unexpected URL: {request.full_url}")

    monkeypatch.setattr(uploader, "_read_json_response", fake_read_json_response)

    payload = uploader.upload_audio_file(
        object(),
        base_url="https://example.test",
        launch_token="launch-token",
        audio_path=audio_path,
        source_metadata={"title": "Morgan interview"},
        audio_content_type="audio/mp4",
    )

    assert payload["job_id"] == "job-chunk"
    assert [url for url, _body in calls] == [
        "https://example.test/case-notes/api/voice/upload",
        "https://example.test/case-notes/api/voice/upload/chunks/start",
        "https://example.test/case-notes/api/voice/upload/chunks/upload-123",
        "https://example.test/case-notes/api/voice/upload/chunks/upload-123",
        "https://example.test/case-notes/api/voice/upload/chunks/upload-123",
        "https://example.test/case-notes/api/voice/upload/chunks/upload-123/finish",
    ]
    assert len(chunk_bodies) == 3
    assert b"abcde" in chunk_bodies[0]
    assert b"fghij" in chunk_bodies[1]
    assert b"k" in chunk_bodies[2]


def test_upload_hosted_audio_retries_large_recording_as_chunks_after_transport_error(
    tmp_path: Path,
    monkeypatch,
) -> None:
    uploader = load_hosted_audio_uploader()
    audio_path = tmp_path / "morgan.m4a"
    audio_path.write_bytes(b"abcdefghijk")
    captured: dict[str, Any] = {}

    monkeypatch.setattr(uploader, "CHUNKED_AUDIO_UPLOAD_CHUNK_BYTES", 5)

    def fake_submit_single_audio_upload(*_args, **_kwargs):
        raise uploader.CaseWorkspaceError("hosted audio request failed: timed out")

    def fake_submit_chunked_audio_upload(*_args, **kwargs):
        captured.update(kwargs)
        return uploader.JsonResponse(202, {"status": "queued", "job_id": "job-retry"})

    monkeypatch.setattr(
        uploader,
        "_submit_single_audio_upload",
        fake_submit_single_audio_upload,
    )
    monkeypatch.setattr(
        uploader,
        "_submit_chunked_audio_upload",
        fake_submit_chunked_audio_upload,
    )

    payload = uploader.upload_audio_file(
        object(),
        base_url="https://example.test",
        launch_token="launch-token",
        audio_path=audio_path,
        source_metadata={"title": "Morgan interview"},
        audio_content_type="audio/mp4",
    )

    assert payload["job_id"] == "job-retry"
    assert captured["audio_size_bytes"] == 11
    assert captured["audio_path"] == audio_path


def test_upload_hosted_audio_poll_fails_fast_on_non_200(
    tmp_path: Path,
    monkeypatch,
) -> None:
    uploader = load_hosted_audio_uploader()
    calls: list[str] = []

    def fake_read_json_response(_opener, request, *, timeout_seconds):
        calls.append(request.full_url)
        return uploader.JsonResponse(401, {"detail": "Not authenticated"})

    monkeypatch.setattr(uploader, "_read_json_response", fake_read_json_response)

    with pytest.raises(
        uploader.CaseWorkspaceError,
        match=r"hosted audio job poll failed \(401\): Not authenticated",
    ):
        uploader.poll_upload_job(
            object(),
            base_url="https://example.test",
            job_id="missing-job",
            run_dir=tmp_path,
            poll_seconds=60,
            poll_interval_seconds=0.01,
        )

    latest_payload = json.loads(
        (tmp_path / "latest_job_payload.json").read_text(encoding="utf-8")
    )
    assert calls == ["https://example.test/case-notes/api/voice/upload/missing-job"]
    assert latest_payload == {"detail": "Not authenticated"}


def test_import_hosted_voice_zip_bundle_copies_packaged_audio(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    importer = load_hosted_voice_importer()
    bundle_path = tmp_path / "case-notes-voice-20260102.zip"
    payload = downloaded_voice_bundle_payload(
        captured_at="2026-01-02T10:30:00+00:00",
        transcript="The live capture has a transcript and audio recording.",
        audio_file_name="case-notes-voice-20260102.webm",
    )
    with ZipFile(bundle_path, "w") as archive:
        archive.writestr("case-notes-voice-20260102.json", json.dumps(payload))
        archive.writestr("case-notes-voice-20260102.webm", b"webm audio bytes")

    result = importer.import_hosted_voice_bundle(case_dir, bundle_path)
    raw_transcript = result.raw_transcript_path.read_text(encoding="utf-8")
    review_pack = result.discussion_review_pack_path.read_text(encoding="utf-8")

    assert result.audio_path == result.session_dir / "case-notes-voice-20260102.webm"
    assert result.audio_path.read_bytes() == b"webm audio bytes"
    assert "The live capture has a transcript and audio recording." in raw_transcript
    assert (
        "Audio file: `voice_sessions/20260102103000Z/case-notes-voice-20260102.webm`"
        in (review_pack)
    )


def test_import_hosted_voice_zip_bundle_copies_screen_video_provenance(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    importer = load_hosted_voice_importer()
    bundle_path = tmp_path / "case-notes-voice-20260102.zip"
    payload = downloaded_voice_bundle_payload(
        captured_at="2026-01-02T10:30:00+00:00",
        transcript="[00:10:33] Reviewer: On this slide move the bridge.",
        audio_file_name="case-notes-voice-20260102.webm",
        video_file_name="case-notes-screen-20260102.webm",
    )
    with ZipFile(bundle_path, "w") as archive:
        archive.writestr("case-notes-voice-20260102.json", json.dumps(payload))
        archive.writestr("case-notes-voice-20260102.webm", b"webm audio bytes")
        archive.writestr("case-notes-screen-20260102.webm", b"webm video bytes")

    result = importer.import_hosted_voice_bundle(case_dir, bundle_path)
    raw_transcript = result.raw_transcript_path.read_text(encoding="utf-8")
    review_pack = result.discussion_review_pack_path.read_text(encoding="utf-8")
    registry = json.loads((case_dir / "material_registry.json").read_text())
    timeline = json.loads(result.video_timeline_path.read_text())

    assert result.video_path == result.session_dir / "case-notes-screen-20260102.webm"
    assert result.video_path.read_bytes() == b"webm video bytes"
    assert result.video_timeline_path == (result.session_dir / "video_timeline.json")
    assert "Screen video: case-notes-screen-20260102.webm" in raw_transcript
    assert "Screen capture display surface: browser" in raw_transcript
    assert (
        "Screen video: `voice_sessions/20260102103000Z/case-notes-screen-20260102.webm`"
        in review_pack
    )
    assert (
        "Video timeline: `voice_sessions/20260102103000Z/video_timeline.json`"
        in review_pack
    )
    assert "use it as provenance for visual references" in review_pack
    assert registry["materials"][0]["source_metadata"]["screen_video"].endswith(
        "case-notes-screen-20260102.webm"
    )
    assert registry["materials"][0]["source_metadata"][
        "screen_video_timeline"
    ].endswith("video_timeline.json")
    assert "mode" not in timeline
    assert timeline["entries"] == []
    assert timeline["video_path"].endswith("case-notes-screen-20260102.webm")


def test_import_hosted_voice_zip_bundle_builds_feedback_timeline_from_timing(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    importer = load_hosted_voice_importer()
    bundle_path = tmp_path / "case-notes-voice-20260102.zip"
    payload = downloaded_voice_bundle_payload(
        captured_at="2026-01-02T10:30:00+00:00",
        transcript=(
            "On this slide the headline is too weak and does not say the answer. "
            "Please make the margin expansion point the lead message. Then delete "
            "the second content box because it distracts from pricing discipline."
        ),
        audio_file_name="case-notes-voice-20260102.webm",
        video_file_name="case-notes-screen-20260102.webm",
    )
    payload["timed_transcript_segments"] = [
        {
            "segment_id": "R1",
            "start_ms": 10000,
            "end_ms": 18000,
            "text": "on this slide the headline is weak",
            "active_slide_id": "margin-story",
            "active_slide_title": "Margin expansion is the lead message",
            "active_slide_index": 3,
            "active_slide_number": 4,
            "active_deck_title": "Commercial review",
            "active_slide_relative_ms": 8000,
        },
        {
            "segment_id": "R2",
            "start_ms": 18000,
            "end_ms": 26000,
            "text": "make margin expansion the lead message",
            "active_slide_id": "margin-story",
            "active_slide_title": "Margin expansion is the lead message",
            "active_slide_index": 3,
            "active_slide_number": 4,
            "active_deck_title": "Commercial review",
            "active_slide_relative_ms": 8000,
        },
        {
            "segment_id": "R3",
            "start_ms": 26000,
            "end_ms": 34000,
            "text": "delete the second content box",
            "active_slide_id": "pricing-discipline",
            "active_slide_title": "Pricing discipline supports the plan",
            "active_slide_index": 4,
            "active_slide_number": 5,
            "active_deck_title": "Commercial review",
            "active_slide_relative_ms": 25000,
        },
    ]
    payload["active_slide_capture"] = {
        "strategy": "capture_handle",
        "status": "captured",
        "supported": True,
        "identified_event_count": 2,
    }
    payload["active_slide_timeline"] = [
        {
            "event_type": "active_slide",
            "relative_ms": 8000,
            "slide_id": "margin-story",
            "slide_title": "Margin expansion is the lead message",
        },
        {
            "event_type": "active_slide",
            "relative_ms": 25000,
            "slide_id": "pricing-discipline",
            "slide_title": "Pricing discipline supports the plan",
        },
    ]
    with ZipFile(bundle_path, "w") as archive:
        archive.writestr("case-notes-voice-20260102.json", json.dumps(payload))
        archive.writestr("case-notes-voice-20260102.webm", b"webm audio bytes")
        archive.writestr("case-notes-screen-20260102.webm", b"webm video bytes")

    result = importer.import_hosted_voice_bundle(case_dir, bundle_path)
    review_pack = result.discussion_review_pack_path.read_text(encoding="utf-8")
    registry = json.loads((case_dir / "material_registry.json").read_text())
    timeline = json.loads(result.feedback_timeline_path.read_text(encoding="utf-8"))
    video_timeline = json.loads(result.video_timeline_path.read_text(encoding="utf-8"))
    first_entry = timeline["entries"][0]

    assert result.feedback_timeline_path == (
        result.session_dir / "feedback_timeline.json"
    )
    assert (
        "Feedback timeline: `voice_sessions/20260102103000Z/feedback_timeline.json`"
        in review_pack
    )
    assert "`use_as_visual_evidence` is true" in review_pack
    assert "Treat `timing_only`, `weak_alignment`" in review_pack
    assert registry["materials"][0]["source_metadata"]["feedback_timeline"].endswith(
        "feedback_timeline.json"
    )
    assert timeline["source"] == "case_notes_voice_feedback_timeline"
    assert timeline["alignment"]["method"] == "monotonic_fuzzy_text_alignment"
    assert (
        timeline["alignment"]["timing_source"] == "realtime_timed_transcript_segments"
    )
    assert first_entry["realtime_segment_ids"] == ["R1", "R2", "R3"]
    assert first_entry["start_ms"] == 10000
    assert first_entry["end_ms"] == 34000
    assert first_entry["active_slide_context_status"] == "changed"
    assert first_entry["active_slide_at_start"]["slide_id"] == "margin-story"
    assert [slide["slide_id"] for slide in first_entry["active_slides"]] == [
        "margin-story",
        "pricing-discipline",
    ]
    assert timeline["alignment"]["active_slide_segment_count"] == 3
    assert timeline["evidence_summary"]["active_slide_changed_entries"] == 1
    assert video_timeline["active_slide_capture"]["status"] == "captured"
    assert video_timeline["active_slide_timeline"][1]["slide_id"] == (
        "pricing-discipline"
    )
    assert first_entry["frames"]
    assert first_entry["frames"][0]["frame_time_ms"] > 0
    assert first_entry["visual_evidence_status"] == "timing_only"
    assert first_entry["use_as_visual_evidence"] is False
    assert "no extracted frame path" in first_entry["evidence_note"]
    assert timeline["evidence_summary"]["visual_evidence_entries"] == 0
    assert timeline["evidence_summary"]["timing_only_entries"] == 1
    assert result.deck_revision_intake_path == (
        result.session_dir / "deck_revision_intake.json"
    )
    assert result.deck_revision_gate_path == (
        result.session_dir / "deck_revision_gate.md"
    )
    assert (
        "Deck revision gate: `voice_sessions/20260102103000Z/deck_revision_gate.md`"
        in review_pack
    )
    assert registry["materials"][0]["source_metadata"]["deck_revision_intake"].endswith(
        "deck_revision_intake.json"
    )
    intake = json.loads(result.deck_revision_intake_path.read_text(encoding="utf-8"))
    assert intake["deck_correction_gate"]["status"] == "pending_model_review"
    assert intake["deck_correction_gate"]["deterministic_gate_used"] is False
    assert intake["deck"]["status"] == "missing"


def test_prepare_voice_deck_revision_intake_marks_gate_and_missing_deck(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    preparer = load_voice_deck_revision_preparer()
    session_dir = case_dir / "voice_sessions" / "20260102103000Z"
    session_dir.mkdir(parents=True)
    (session_dir / "raw_transcript.md").write_text(
        "Facilitator: This slide needs a stronger headline.\n",
        encoding="utf-8",
    )
    (session_dir / "speaker_attribution_task.md").write_text(
        "# Speaker Attribution Task\n",
        encoding="utf-8",
    )
    (session_dir / "case-notes-screen-20260102.webm").write_bytes(b"video")
    (session_dir / "feedback_timeline.json").write_text(
        json.dumps({"entries": []}),
        encoding="utf-8",
    )

    result = preparer.prepare_voice_deck_revision_intake(
        case_dir,
        voice_session=Path("20260102103000Z"),
        now=fixed_now(),
    )
    intake = json.loads(result.intake_path.read_text(encoding="utf-8"))
    gate = result.gate_path.read_text(encoding="utf-8")

    assert result.deck_snapshot_path is None
    assert intake["source"] == "clara_voice_deck_revision_intake"
    assert intake["speaker_attribution"]["status"] == "pending"
    assert intake["speaker_attribution"]["attribution_task_path"].endswith(
        "speaker_attribution_task.md"
    )
    assert intake["deck_correction_gate"]["status"] == "pending_model_review"
    assert intake["deck_correction_gate"]["deterministic_gate_used"] is False
    assert intake["deck"]["status"] == "missing"
    assert intake["company_profile"]["status"] == "missing"
    assert intake["deck_style"]["status"] == "missing"
    assert intake["deck_style"]["required_before_editing"] is True
    assert "Resolve the company/deck style authority" in " ".join(
        intake["next_actions"]
    )
    assert "advisory-output-shaper" in gate
    assert intake["evidence"]["feedback_timeline_path"].endswith(
        "feedback_timeline.json"
    )
    assert "request or attach the PPTX" in " ".join(intake["next_actions"])
    assert "Do not decide the gate with keyword matching" in gate


def test_prepare_voice_deck_revision_intake_inherits_parent_company_style_profile(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    (tmp_path / "company_profile.json").write_text(
        json.dumps(
            {
                "schema_version": 1,
                "company": "A&G",
                "default_deck_style": "A&G",
                "advisory_method": "advisory-output-shaper",
            }
        ),
        encoding="utf-8",
    )
    preparer = load_voice_deck_revision_preparer()
    session_dir = case_dir / "voice_sessions" / "20260102103000Z"
    session_dir.mkdir(parents=True)
    (session_dir / "raw_transcript.md").write_text(
        "Facilitator: Add one slide and keep it in our A&G format.\n",
        encoding="utf-8",
    )

    result = preparer.prepare_voice_deck_revision_intake(
        case_dir,
        voice_session=Path("20260102103000Z"),
        now=fixed_now(),
    )
    intake = json.loads(result.intake_path.read_text(encoding="utf-8"))
    gate = result.gate_path.read_text(encoding="utf-8")

    assert result.style_spec_snapshot_path == session_dir / "deck_style_spec.md"
    assert intake["company_profile"]["status"] == "available"
    assert intake["company_profile"]["company"] == "A&G"
    assert intake["deck_style"]["status"] == "available"
    assert intake["deck_style"]["source"] == "company_profile_deck_style"
    assert intake["deck_style"]["style_key"] == "ag"
    assert intake["deck_style"]["style_name"] == "A&G PPTX Style Spec"
    assert intake["deck_style"]["snapshot_path"].endswith("deck_style_spec.md")
    snapshot_text = result.style_spec_snapshot_path.read_text(encoding="utf-8")
    assert "A&G PPTX Style Spec" in snapshot_text
    assert "#002060" in snapshot_text
    assert "A&G PPTX Style Spec" in gate
    assert "deck_style_spec.md" in gate
    assert "advisory-output-shaper" in gate
    assert "Apply the resolved deck style spec" in " ".join(intake["next_actions"])


def test_prepare_voice_deck_revision_intake_accepts_explicit_deck_style(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    preparer = load_voice_deck_revision_preparer()
    session_dir = case_dir / "voice_sessions" / "20260102103000Z"
    session_dir.mkdir(parents=True)
    (session_dir / "raw_transcript.md").write_text(
        "Facilitator: Add one slide in Bain format.\n",
        encoding="utf-8",
    )

    result = preparer.prepare_voice_deck_revision_intake(
        case_dir,
        voice_session=Path("20260102103000Z"),
        deck_style="bain",
        now=fixed_now(),
    )
    intake = json.loads(result.intake_path.read_text(encoding="utf-8"))
    snapshot_text = result.style_spec_snapshot_path.read_text(encoding="utf-8")

    assert intake["deck_style"]["status"] == "available"
    assert intake["deck_style"]["source"] == "explicit_deck_style"
    assert intake["deck_style"]["style_key"] == "bain"
    assert intake["deck_style"]["style_name"] == "Bain Style Spec"
    assert "bain-style-spec.md" in intake["deck_style"]["spec_path"]
    assert "Bain Style Spec" in snapshot_text


def test_prepare_voice_deck_revision_intake_attaches_pptx_snapshot(
    tmp_path: Path,
) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    _, case_dir = init_case(tmp_path)
    preparer = load_voice_deck_revision_preparer()
    session_dir = case_dir / "voice_sessions" / "20260102103000Z"
    session_dir.mkdir(parents=True)
    (session_dir / "raw_transcript.md").write_text(
        "Partner: Move the margin point up on this page.\n",
        encoding="utf-8",
    )
    (session_dir / "attributed_transcript.md").write_text(
        "Partner: Move the margin point up on this page.\n",
        encoding="utf-8",
    )
    deck = tmp_path / "margin_story.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Margin expansion opportunity"
    slide.placeholders[1].text = "Pricing discipline drives EBITDA improvement"
    presentation.save(deck)

    result = preparer.prepare_voice_deck_revision_intake(
        case_dir,
        voice_session=session_dir,
        deck_path=deck,
        now=fixed_now(),
    )
    intake = json.loads(result.intake_path.read_text(encoding="utf-8"))
    snapshot = json.loads(result.deck_snapshot_path.read_text(encoding="utf-8"))
    merge_report = json.loads(result.merge_input_report_path.read_text())

    assert intake["speaker_attribution"]["status"] == "available"
    assert intake["deck"]["status"] == "attached"
    assert intake["deck"]["snapshot_path"].endswith("deck_snapshot.json")
    assert intake["deck"]["editable_merge_input"]["status"] == "source_clean_for_merge"
    assert snapshot["slide_count"] == 1
    assert snapshot["slides"][0]["title"] == "Margin expansion opportunity"
    assert (
        "Pricing discipline drives EBITDA improvement" in snapshot["slides"][0]["texts"]
    )
    assert merge_report["status"] == "source_clean_for_merge"


def test_prepare_voice_deck_revision_shape_text_skips_non_table_graphics() -> None:
    preparer = load_voice_deck_revision_preparer()

    class GraphicFrame:
        has_table = False
        text = ""

        @property
        def table(self) -> Any:
            raise ValueError("shape does not contain a table")

    assert preparer._shape_text(GraphicFrame()) == ""


def test_prepare_voice_deck_revision_intake_runs_slide_matching(
    tmp_path: Path,
    monkeypatch,
) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    _, case_dir = init_case(tmp_path)
    preparer = load_voice_deck_revision_preparer()
    session_dir = case_dir / "voice_sessions" / "20260102103000Z"
    session_dir.mkdir(parents=True)
    (session_dir / "raw_transcript.md").write_text(
        "Partner: Make the title stronger on this page.\n",
        encoding="utf-8",
    )
    (session_dir / "feedback_timeline.json").write_text(
        json.dumps({"schema_version": 1, "entries": []}),
        encoding="utf-8",
    )
    deck = tmp_path / "title_story.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Initial title"
    slide.placeholders[1].text = "Supporting point"
    presentation.save(deck)
    captured: dict[str, Path] = {}

    def fake_match_feedback_timeline_to_deck(**kwargs):
        captured["feedback_timeline_path"] = kwargs["feedback_timeline_path"]
        captured["deck_path"] = kwargs["deck_path"]
        captured["deck_snapshot_path"] = kwargs["deck_snapshot_path"]
        timeline = json.loads(kwargs["feedback_timeline_path"].read_text())
        timeline["slide_matching"] = {
            "status": "complete",
            "method": "test",
            "slide_count": 1,
            "summary": {"matched_entries": 0},
        }
        kwargs["feedback_timeline_path"].write_text(
            json.dumps(timeline),
            encoding="utf-8",
        )
        return timeline

    monkeypatch.setattr(
        preparer,
        "match_feedback_timeline_to_deck",
        fake_match_feedback_timeline_to_deck,
    )

    result = preparer.prepare_voice_deck_revision_intake(
        case_dir,
        voice_session=session_dir,
        deck_path=deck,
        now=fixed_now(),
    )
    intake = json.loads(result.intake_path.read_text(encoding="utf-8"))
    gate = result.gate_path.read_text(encoding="utf-8")

    assert captured["feedback_timeline_path"] == session_dir / "feedback_timeline.json"
    assert captured["deck_path"] == deck
    assert captured["deck_snapshot_path"] == result.deck_snapshot_path
    assert intake["evidence"]["slide_matching"]["status"] == "complete"
    assert "Slide matching: `complete`" in gate


def build_ready_deck_revision_intake(tmp_path: Path) -> tuple[Path, Path]:
    pytest.importorskip("pptx")
    from pptx import Presentation

    _, case_dir = init_case(tmp_path)
    preparer = load_voice_deck_revision_preparer()
    session_dir = case_dir / "voice_sessions" / "20260102103000Z"
    session_dir.mkdir(parents=True)
    transcript = (
        "Partner: Move the margin point up on this page and make the title stronger.\n"
    )
    (session_dir / "raw_transcript.md").write_text(transcript, encoding="utf-8")
    (session_dir / "attributed_transcript.md").write_text(
        transcript,
        encoding="utf-8",
    )
    (session_dir / "case-notes-screen-20260102.webm").write_bytes(b"video")
    (session_dir / "feedback_timeline.json").write_text(
        json.dumps(
            {
                "entries": [
                    {
                        "text": "Move the margin point up on this page.",
                        "start_ms": 26000,
                        "end_ms": 34000,
                        "visual_evidence_status": "usable",
                        "use_as_visual_evidence": True,
                    }
                ]
            }
        ),
        encoding="utf-8",
    )
    deck = tmp_path / "margin_story.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Margin expansion opportunity"
    slide.placeholders[1].text = "Pricing discipline drives EBITDA improvement"
    presentation.save(deck)
    preparer.prepare_voice_deck_revision_intake(
        case_dir,
        voice_session=session_dir,
        deck_path=deck,
        deck_style="ag",
        now=fixed_now(),
    )
    return case_dir, session_dir


def executable_deck_revision_change(**overrides: Any) -> dict[str, Any]:
    change: dict[str, Any] = {
        "change_id": "chg-001",
        "slide_number": 1,
        "change_scope": "text",
        "change_type": "rewrite_headline",
        "requested_change": (
            "Rewrite the title so the margin point is the lead message."
        ),
        "interpretation": (
            "The slide should lead with the conclusion that margin expansion is "
            "the main story."
        ),
        "rationale": (
            "The partner asked to move the margin point up and make the title stronger."
        ),
        "explicitness": "explicit",
        "confidence": "high",
        "execution_strategy": "deterministic_patch",
        "transcript_evidence": [
            {
                "speaker": "Partner",
                "timestamp_ms": 26000,
                "quote": (
                    "Move the margin point up on this page and make the title stronger."
                ),
            }
        ],
        "visual_evidence": [
            {
                "evidence_type": "deck_snapshot",
                "path": "voice_sessions/20260102103000Z/deck_snapshot.json",
                "note": "Slide 1 contains the margin expansion title.",
            }
        ],
        "success_criteria": [
            {
                "criterion_id": "chg-001-c001",
                "check_type": "title_equals",
                "description": "The slide title is rewritten to lead with margin.",
                "expected_text": "Margin expansion should lead the page story",
            }
        ],
        "application_patches": [
            {
                "patch_id": "chg-001-p001",
                "operation": "set_title_text",
                "target": {"expected_text": "Margin expansion opportunity"},
                "value": {"text": "Margin expansion should lead the page story"},
            }
        ],
    }
    change.update(overrides)
    return change


def semantic_deck_revision_change(**overrides: Any) -> dict[str, Any]:
    change = executable_deck_revision_change(
        change_id="chg-002",
        change_scope="storyline",
        change_type="sharpen_storyline",
        requested_change="Make the slide tell a stronger margin story.",
        interpretation=(
            "The page needs a stronger conclusion and may need wording/layout "
            "judgement beyond an exact text replacement."
        ),
        execution_strategy="model_assisted_edit",
        material_requirements=["Codex must draft the stronger wording before edit."],
        success_criteria=[
            {
                "criterion_id": "chg-002-c001",
                "check_type": "semantic_review",
                "description": (
                    "The final slide makes margin expansion the dominant conclusion."
                ),
                "note": "Requires Codex/consultant semantic review.",
            }
        ],
        application_patches=[],
    )
    change.update(overrides)
    return change


def test_build_deck_revision_workbench_creates_prompt_schema_and_review(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    builder = load_deck_revision_workbench_builder()

    result = builder.build_deck_revision_workbench(
        case_dir,
        voice_session=Path("20260102103000Z"),
        now=fixed_now(),
    )
    workbench = json.loads(result.workbench_path.read_text(encoding="utf-8"))
    schema = json.loads(result.schema_path.read_text(encoding="utf-8"))
    prompt = result.prompt_path.read_text(encoding="utf-8")
    review = result.review_path.read_text(encoding="utf-8")

    assert result.session_dir == session_dir.resolve()
    assert workbench["source"] == "clara_deck_revision_workbench"
    assert workbench["status"]["pptx_edits_made"] is False
    assert workbench["deck"]["slide_count"] == 1
    assert workbench["deck_style"]["style_key"] == "ag"
    assert workbench["source_paths"]["changes_output_path"].endswith(
        "deck_revision_changes.json"
    )
    assert workbench["source_paths"]["understanding_path"].endswith(
        "deck_revision_understanding.md"
    )
    assert workbench["source_paths"]["execution_plan_path"].endswith(
        "deck_revision_execution_plan.json"
    )
    assert workbench["source_paths"]["interpretation_packets_path"].endswith(
        "deck_revision_interpretation_packets.json"
    )
    assert workbench["source_paths"]["execution_packets_path"].endswith(
        "deck_revision_execution_packets.json"
    )
    assert schema["properties"]["changes"]["type"] == "array"
    assert "execution_strategy" in schema["properties"]["changes"]["items"]["required"]
    assert "success_criteria" in schema["properties"]["changes"]["items"]["required"]
    change_schema = schema["properties"]["changes"]["items"]["properties"]
    assert change_schema["packet_scope"]["enum"] == ["deck", "slide", "slide_cluster"]
    assert change_schema["affected_slide_numbers"]["type"] == "array"
    assert "Do not edit the PPTX" in prompt
    assert "build_deck_revision_interpretation_packets.py" in prompt
    assert "Assign one execution strategy" in prompt
    assert "consultant checkpoint" in review


def test_finalize_deck_revision_plan_renders_consultant_change_list(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    builder = load_deck_revision_workbench_builder()
    finalizer = load_deck_revision_plan_finalizer()
    builder.build_deck_revision_workbench(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    changes_path = session_dir / "deck_revision_changes.json"
    changes_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "codex_deck_revision_plan",
                "approved_for_pptx_revision": False,
                "changes": [
                    executable_deck_revision_change(
                        style_notes=["Keep the resolved A&G deck style."]
                    )
                ],
                "open_questions": [],
            }
        ),
        encoding="utf-8",
    )

    result = finalizer.finalize_deck_revision_plan(
        case_dir,
        changes_path,
        voice_session=Path("20260102103000Z"),
        now=fixed_now(),
    )
    normalized = json.loads(result.normalized_plan_path.read_text(encoding="utf-8"))
    review = result.review_path.read_text(encoding="utf-8")
    handoff = result.handoff_path.read_text(encoding="utf-8")

    assert normalized["source"] == "clara_deck_revision_changes"
    assert normalized["approved_for_pptx_revision"] is False
    assert normalized["model_requested_pptx_revision_approval"] is False
    assert normalized["changes"][0]["slide_title"] == "Margin expansion opportunity"
    assert normalized["changes"][0]["affected_slide_numbers"] == [1]
    assert normalized["changes"][0]["packet_scope"] is None
    assert normalized["changes"][0]["change_scope"] == "text"
    assert normalized["changes"][0]["execution_strategy"] == "deterministic_patch"
    assert normalized["changes"][0]["success_criteria"][0]["check_type"] == (
        "title_equals"
    )
    assert normalized["changes"][0]["application_patches"][0]["operation"] == (
        "set_title_text"
    )
    understanding = result.understanding_path.read_text(encoding="utf-8")
    assert "Status: consultant review required before PPTX revision." in review
    assert "## Slide 1 - Margin expansion opportunity" in review
    assert "Clara interpretation:" in review
    assert "Success criteria:" in review
    assert "Partner at 00:26" in review
    assert "`chg-001-p001`: `set_title_text`" in review
    assert "Deck Revision Understanding" in understanding
    assert "Should change:" in understanding
    assert "waiting for separate plan approval" in handoff
    assert "deck_revision_execution_plan.json" in handoff
    assert "approve_deck_revision_plan.py" in handoff


def test_deck_revision_interpretation_packets_route_matched_feedback(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    (session_dir / "feedback_timeline.json").write_text(
        json.dumps(
            {
                "schema_version": 1,
                "entries": [
                    {
                        "feedback_unit_id": "F001",
                        "clean_text": "Move the margin point up on this page.",
                        "start_ms": 26000,
                        "end_ms": 34000,
                        "slide_match": {
                            "status": "matched",
                            "best_slide_number": 1,
                            "confidence": "high",
                            "score": 0.94,
                        },
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    workbench_builder = load_deck_revision_workbench_builder()
    packet_builder = load_deck_revision_interpretation_packet_builder()
    workbench_builder.build_deck_revision_workbench(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )

    result = packet_builder.build_deck_revision_interpretation_packets(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    index = json.loads(result.packet_index_path.read_text(encoding="utf-8"))
    review = result.review_path.read_text(encoding="utf-8")
    packet_path = case_dir / index["packets"][0]["packet_path"]
    packet = json.loads(packet_path.read_text(encoding="utf-8"))

    assert index["source"] == "clara_deck_revision_interpretation_packets"
    assert index["summary"]["slide_packets"] == 1
    assert index["summary"]["deck_packets"] == 0
    assert index["packets"][0]["packet_scope"] == "slide"
    assert packet["slides"][0]["slide_number"] == 1
    assert packet["feedback_units"][0]["text"] == (
        "Move the margin point up on this page."
    )
    assert "not as one huge semantic prompt" in review


def write_executable_deck_revision_plan(
    case_dir: Path,
    session_dir: Path,
    *,
    approved: bool = True,
) -> Path:
    builder = load_deck_revision_workbench_builder()
    finalizer = load_deck_revision_plan_finalizer()
    builder.build_deck_revision_workbench(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    changes_path = session_dir / "deck_revision_changes.json"
    changes_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "codex_deck_revision_plan",
                "approved_for_pptx_revision": approved,
                "changes": [executable_deck_revision_change()],
                "open_questions": [],
            }
        ),
        encoding="utf-8",
    )
    finalizer.finalize_deck_revision_plan(
        case_dir,
        changes_path,
        voice_session=session_dir,
        now=fixed_now(),
    )
    return session_dir / "deck_revision_changes.normalized.json"


def write_quote_backed_deck_revision_plan(case_dir: Path, session_dir: Path) -> Path:
    builder = load_deck_revision_workbench_builder()
    finalizer = load_deck_revision_plan_finalizer()
    builder.build_deck_revision_workbench(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    changes_path = session_dir / "deck_revision_changes.json"
    changes_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "codex_deck_revision_plan",
                "approved_for_pptx_revision": False,
                "changes": [
                    semantic_deck_revision_change(
                        change_id="chg-quote-001",
                        slide_number=1,
                        change_scope="content",
                        change_type="source_better_quotes",
                        requested_change=(
                            "Pull better quotes from project interview transcripts "
                            "about production quality, roles, deleghe, and conflict."
                        ),
                        interpretation=(
                            "The evidence slide should use stronger interview "
                            "quotes for the four evidence points rather than "
                            "reusing shallow existing quotes."
                        ),
                        material_requirements=[
                            "Build deck_revision_quote_candidate_matrix.json before selecting quotes."
                        ],
                    )
                ],
                "open_questions": [],
            }
        ),
        encoding="utf-8",
    )
    finalizer.finalize_deck_revision_plan(
        case_dir,
        changes_path,
        voice_session=session_dir,
        now=fixed_now(),
    )
    return session_dir / "deck_revision_changes.normalized.json"


def test_deck_revision_material_analyzer_marks_executable_patches_ready(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    plan_path = write_executable_deck_revision_plan(case_dir, session_dir)
    analyzer = load_deck_revision_material_analyzer()
    approver = load_deck_revision_approver()

    result = analyzer.analyze_deck_revision_materials(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    needs = json.loads(result.needs_path.read_text(encoding="utf-8"))
    review = result.review_path.read_text(encoding="utf-8")

    assert needs["source"] == "clara_deck_revision_material_needs"
    assert needs["summary"]["status"] == "ready_for_approval"
    assert needs["approval"]["status"] == "missing"
    assert needs["changes"][0]["ready_for_auto_apply"] is True
    assert "`chg-001-p001` `set_title_text`: ready" in review

    approver.approve_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        plan_path=plan_path,
        reviewer="Reviewer",
        understanding_reviewed=True,
        now=fixed_now(),
    )
    approved_result = analyzer.analyze_deck_revision_materials(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    approved_needs = json.loads(approved_result.needs_path.read_text(encoding="utf-8"))

    assert approved_needs["summary"]["status"] == "ready_for_auto_apply"
    assert approved_needs["approval"]["status"] == "approved"
    assert approved_needs["approval"]["understanding_reviewed"] is True
    assert approved_needs["approval"]["understanding_hash_matches"] is True


def test_deck_revision_quote_candidate_matrix_extracts_transcript_candidates(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    source_dir = case_dir / "source_materials" / "interviews"
    source_dir.mkdir(parents=True, exist_ok=True)
    (source_dir / "interview-alex.md").write_text(
        "\n".join(
            [
                "Alex: La produzione qualità è il nodo più delicato.",
                "Alex: Le deleghe e i ruoli non sono chiari e questo crea conflitto.",
                "Alex: Un passaggio serio deve rendere oggettive le decisioni.",
            ]
        ),
        encoding="utf-8",
    )
    write_quote_backed_deck_revision_plan(case_dir, session_dir)
    quote_builder = load_deck_revision_quote_matrix_builder()

    result = quote_builder.build_deck_revision_quote_candidate_matrix(
        case_dir,
        voice_session=session_dir,
        max_candidates_per_change=5,
        now=fixed_now(),
    )
    matrix = json.loads(result.matrix_path.read_text(encoding="utf-8"))
    review = result.review_path.read_text(encoding="utf-8")

    assert matrix["source"] == "clara_deck_revision_quote_candidate_matrix"
    assert matrix["summary"]["status"] == "ready_for_codex_review"
    assert matrix["summary"]["required_changes"] == 1
    assert matrix["changes"][0]["change_id"] == "chg-quote-001"
    assert matrix["changes"][0]["candidate_count"] >= 1
    assert "interview-alex.md" in matrix["changes"][0]["candidates"][0]["source_path"]
    assert "This is an evidence-preparation artifact" in review


def test_deck_revision_quote_candidate_matrix_ignores_existing_quote_page_reference() -> (
    None
):
    quote_builder = load_deck_revision_quote_matrix_builder()

    assert (
        quote_builder.change_needs_quote_candidate_matrix(
            {
                "change_id": "chg-quote-page-palette",
                "change_type": "align_quote_card_palette",
                "requested_change": (
                    "Keep the quote content, but align the four quote-card "
                    "colors with the evidence points."
                ),
                "interpretation": (
                    "The interview quote slide already has sourced quotes; "
                    "this change only fixes visual hierarchy."
                ),
            }
        )
        is False
    )


def test_deck_revision_material_analyzer_blocks_quote_changes_until_matrix_exists(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    source_dir = case_dir / "source_materials" / "interviews"
    source_dir.mkdir(parents=True, exist_ok=True)
    (source_dir / "interview-jordan.md").write_text(
        "Jordan: Produzione qualità, deleghe, ruoli e conflitto vanno trattati con evidenze migliori.",
        encoding="utf-8",
    )
    write_quote_backed_deck_revision_plan(case_dir, session_dir)
    analyzer = load_deck_revision_material_analyzer()

    missing_result = analyzer.analyze_deck_revision_materials(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    missing_needs = json.loads(missing_result.needs_path.read_text(encoding="utf-8"))

    assert missing_needs["quote_candidate_matrix"]["status"] == "required_missing"
    assert missing_needs["quote_candidate_matrix"]["required_change_ids"] == [
        "chg-quote-001"
    ]
    assert any(
        "deck_revision_quote_candidate_matrix" in item
        for item in missing_needs["changes"][0]["missing"]
    )

    quote_builder = load_deck_revision_quote_matrix_builder()
    quote_builder.build_deck_revision_quote_candidate_matrix(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    available_result = analyzer.analyze_deck_revision_materials(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    available_needs = json.loads(
        available_result.needs_path.read_text(encoding="utf-8")
    )

    assert available_needs["quote_candidate_matrix"]["status"] == "available"
    assert not any(
        "deck_revision_quote_candidate_matrix" in item
        for item in available_needs["changes"][0]["missing"]
    )


def test_deck_revision_execution_plan_routes_patch_and_semantic_changes(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    builder = load_deck_revision_workbench_builder()
    finalizer = load_deck_revision_plan_finalizer()
    planner = load_deck_revision_execution_planner()
    builder.build_deck_revision_workbench(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    changes_path = session_dir / "deck_revision_changes.json"
    changes_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "codex_deck_revision_plan",
                "changes": [
                    executable_deck_revision_change(),
                    semantic_deck_revision_change(),
                ],
            }
        ),
        encoding="utf-8",
    )
    finalizer.finalize_deck_revision_plan(
        case_dir,
        changes_path,
        voice_session=session_dir,
        now=fixed_now(),
    )

    result = planner.build_deck_revision_execution_plan(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    execution_plan = json.loads(result.execution_plan_path.read_text(encoding="utf-8"))
    review = result.review_path.read_text(encoding="utf-8")

    assert execution_plan["summary"]["status"] == "mixed_execution_required"
    assert execution_plan["summary"]["ready_for_patcher"] == 1
    assert execution_plan["summary"]["requires_model_or_human_work"] == 1
    assert execution_plan["changes"][0]["status"] == "ready_for_patcher"
    assert execution_plan["changes"][1]["status"] == "requires_model_assisted_edit"
    assert "Use Codex/presentation editing" in review


def test_deck_revision_execution_packets_group_slide_and_deck_changes(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    builder = load_deck_revision_workbench_builder()
    finalizer = load_deck_revision_plan_finalizer()
    planner = load_deck_revision_execution_planner()
    packet_builder = load_deck_revision_execution_packet_builder()
    builder.build_deck_revision_workbench(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    changes_path = session_dir / "deck_revision_changes.json"
    changes_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "codex_deck_revision_plan",
                "changes": [
                    executable_deck_revision_change(),
                    semantic_deck_revision_change(
                        change_id="chg-global-font",
                        change_scope="visual",
                        change_type="increase_global_font_size",
                        requested_change="Make the font bigger in all slides.",
                        interpretation=(
                            "Increase readable body-copy font sizes across the "
                            "deck while preserving layout and hierarchy."
                        ),
                        packet_scope="deck",
                        affected_slide_numbers=[1],
                        execution_group_id="global-font-size",
                        dependency_change_ids=["chg-001"],
                        execution_strategy="model_assisted_edit",
                    ),
                ],
            }
        ),
        encoding="utf-8",
    )
    finalizer.finalize_deck_revision_plan(
        case_dir,
        changes_path,
        voice_session=session_dir,
        now=fixed_now(),
    )
    planner.build_deck_revision_execution_plan(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )

    result = packet_builder.build_deck_revision_execution_packets(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    index = json.loads(result.packet_index_path.read_text(encoding="utf-8"))
    review = result.review_path.read_text(encoding="utf-8")
    deck_packet_summary = next(
        packet for packet in index["packets"] if packet["packet_scope"] == "deck"
    )
    deck_packet = json.loads(
        (case_dir / deck_packet_summary["packet_path"]).read_text(encoding="utf-8")
    )
    slide_packet_summary = next(
        packet for packet in index["packets"] if packet["packet_scope"] == "slide"
    )
    slide_packet = json.loads(
        (case_dir / slide_packet_summary["packet_path"]).read_text(encoding="utf-8")
    )

    assert index["source"] == "clara_deck_revision_execution_packets"
    assert index["summary"]["deck_packets"] == 1
    assert index["summary"]["slide_packets"] == 1
    assert deck_packet["change_ids"] == ["chg-global-font"]
    assert deck_packet["dependency_change_ids"] == ["chg-001"]
    assert slide_packet["change_ids"] == ["chg-001"]
    assert slide_packet["related_change_summaries"][0]["change_id"] == (
        "chg-global-font"
    )
    assert "one deck-editing prompt" in review


def test_apply_deck_revision_plan_writes_corrected_pptx_and_verifies(
    tmp_path: Path,
) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    plan_path = write_executable_deck_revision_plan(case_dir, session_dir)
    approver = load_deck_revision_approver()
    approver.approve_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        plan_path=plan_path,
        reviewer="Reviewer",
        understanding_reviewed=True,
        now=fixed_now(),
    )
    applier = load_deck_revision_applier()

    result = applier.apply_deck_revision_plan(
        case_dir,
        voice_session=Path("20260102103000Z"),
        now=fixed_now(),
    )
    corrected = Presentation(str(result.corrected_deck_path))
    apply_report = json.loads(result.apply_report_path.read_text(encoding="utf-8"))
    verification = json.loads(
        result.verification_report_path.read_text(encoding="utf-8")
    )

    assert corrected.slides[0].shapes.title.text == (
        "Margin expansion should lead the page story"
    )
    assert apply_report["summary"]["status"] == (
        "applied_verified_pending_output_review"
    )
    assert apply_report["summary"]["applied_patches"] == 1
    assert apply_report["summary"]["final_output_review_status"] == (
        "requires_clara_codex_review"
    )
    assert apply_report["approved_by"] == "Reviewer"
    assert apply_report["understanding_path"].endswith("deck_revision_understanding.md")
    assert verification["summary"]["status"] == "verified"
    assert verification["summary"]["passed_patches"] == 1
    assert verification["summary"]["passed_success_criteria"] == 1
    output_review = json.loads(result.output_review_path.read_text(encoding="utf-8"))
    output_review_markdown = result.output_review_markdown_path.read_text(
        encoding="utf-8"
    )

    assert output_review["summary"]["status"] == "requires_clara_codex_review"
    assert output_review["verification_status"] == "verified"
    assert output_review["slide_titles"] == [
        {"slide_number": 1, "title": "Margin expansion should lead the page story"}
    ]
    assert [check["check_id"] for check in output_review["checks"]] == [
        "audience_copy",
        "process_language",
        "requested_structure",
        "semantic_evidence_fit",
        "visual_render",
    ]
    assert "audience-facing slide copy" in output_review_markdown


def test_apply_deck_revision_plan_accepts_whitespace_normalized_target_text(
    tmp_path: Path,
) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation

    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    plan_path = write_executable_deck_revision_plan(case_dir, session_dir)
    normalized = json.loads(plan_path.read_text(encoding="utf-8"))
    normalized["changes"][0]["application_patches"][0]["target"][
        "expected_text"
    ] = "Margin\n\n expansion\t opportunity"
    plan_path.write_text(json.dumps(normalized), encoding="utf-8")
    approver = load_deck_revision_approver()
    approver.approve_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        plan_path=plan_path,
        reviewer="Reviewer",
        understanding_reviewed=True,
        now=fixed_now(),
    )
    applier = load_deck_revision_applier()

    result = applier.apply_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    corrected = Presentation(str(result.corrected_deck_path))
    apply_report = json.loads(result.apply_report_path.read_text(encoding="utf-8"))
    verification = json.loads(
        result.verification_report_path.read_text(encoding="utf-8")
    )

    assert corrected.slides[0].shapes.title.text == (
        "Margin expansion should lead the page story"
    )
    assert apply_report["summary"]["status"] == (
        "applied_verified_pending_output_review"
    )
    assert apply_report["summary"]["applied_patches"] == 1
    assert apply_report["summary"]["blocked_patches"] == 0
    assert verification["summary"]["status"] == "verified"
    assert verification["summary"]["passed_patches"] == 1


def test_approve_deck_revision_plan_requires_understanding_review(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    plan_path = write_executable_deck_revision_plan(case_dir, session_dir)
    approver = load_deck_revision_approver()

    with pytest.raises(
        approver.CaseWorkspaceError,
        match="understanding checkpoint has not been marked reviewed",
    ):
        approver.approve_deck_revision_plan(
            case_dir,
            voice_session=session_dir,
            plan_path=plan_path,
            reviewer="Reviewer",
            now=fixed_now(),
        )


def test_deck_revision_fixture_runner_exercises_full_patch_harness(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    fixture_dir = tmp_path / "deck_revision_fixture"
    fixture_dir.mkdir()
    changes_path = fixture_dir / "deck_revision_changes.json"
    changes_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "codex_deck_revision_plan",
                "changes": [executable_deck_revision_change()],
            }
        ),
        encoding="utf-8",
    )
    (fixture_dir / "fixture.json").write_text(
        json.dumps(
            {
                "case_dir": str(case_dir),
                "voice_session": str(session_dir),
                "changes_json": str(changes_path),
                "approve": True,
                "reviewer": "Reviewer",
                "apply": True,
                "expected": {
                    "change_count": 1,
                    "execution_status": "ready_for_deterministic_apply",
                    "material_status": "ready_for_approval",
                    "apply_status": "applied_verified_pending_output_review",
                    "final_output_review_status": "requires_clara_codex_review",
                    "verification_status": "verified",
                    "change_strategies": ["deterministic_patch"],
                },
            }
        ),
        encoding="utf-8",
    )
    runner = load_deck_revision_fixture_runner()

    result = runner.run_deck_revision_fixture(fixture_dir, now=fixed_now())
    report = json.loads(result.report_path.read_text(encoding="utf-8"))
    review = result.review_path.read_text(encoding="utf-8")

    assert report["summary"]["status"] == "passed"
    assert all(record["passed"] for record in report["expectations"])
    assert report["artifacts"]["understanding"].endswith(
        "deck_revision_understanding.md"
    )
    assert report["artifacts"]["output_review"].endswith(
        "deck_revision_output_review.json"
    )
    assert "Deck Revision Fixture Report" in review


def test_apply_deck_revision_plan_requires_approval_by_default(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    write_executable_deck_revision_plan(case_dir, session_dir, approved=True)
    applier = load_deck_revision_applier()

    with pytest.raises(
        applier.CaseWorkspaceError,
        match="approval is missing",
    ):
        applier.apply_deck_revision_plan(
            case_dir,
            voice_session=session_dir,
            now=fixed_now(),
        )


def test_complete_deck_revision_output_review_requires_all_confirmations(
    tmp_path: Path,
) -> None:
    pytest.importorskip("pptx")
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    plan_path = write_executable_deck_revision_plan(case_dir, session_dir)
    approver = load_deck_revision_approver()
    applier = load_deck_revision_applier()
    completer = load_deck_revision_output_review_completer()
    approver.approve_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        plan_path=plan_path,
        reviewer="Reviewer",
        understanding_reviewed=True,
        now=fixed_now(),
    )
    applier.apply_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )

    with pytest.raises(
        completer.CaseWorkspaceError,
        match="final output review is incomplete",
    ):
        completer.complete_deck_revision_output_review(
            case_dir,
            voice_session=session_dir,
            reviewer="Reviewer",
            audience_copy_reviewed=True,
            process_language_reviewed=True,
            requested_structure_reviewed=True,
            semantic_evidence_fit_reviewed=True,
            now=fixed_now(),
        )


def test_complete_deck_revision_output_review_records_completion(
    tmp_path: Path,
) -> None:
    pytest.importorskip("pptx")
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    plan_path = write_executable_deck_revision_plan(case_dir, session_dir)
    approver = load_deck_revision_approver()
    applier = load_deck_revision_applier()
    completer = load_deck_revision_output_review_completer()
    approver.approve_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        plan_path=plan_path,
        reviewer="Reviewer",
        understanding_reviewed=True,
        now=fixed_now(),
    )
    applier.apply_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )

    result = completer.complete_deck_revision_output_review(
        case_dir,
        voice_session=session_dir,
        reviewer="Reviewer",
        note="Rendered deck inspected.",
        audience_copy_reviewed=True,
        process_language_reviewed=True,
        requested_structure_reviewed=True,
        semantic_evidence_fit_reviewed=True,
        visual_render_reviewed=True,
        now=fixed_now(),
    )
    completion = json.loads(result.completion_path.read_text(encoding="utf-8"))
    completion_markdown = result.completion_markdown_path.read_text(encoding="utf-8")

    assert completion["summary"]["status"] == "complete"
    assert completion["summary"]["final_delivery_allowed"] is True
    assert completion["completed_by"] == "Reviewer"
    assert completion["review_note"] == "Rendered deck inspected."
    assert all(completion["confirmations"].values())
    assert completion["output_review_sha256"]
    assert completion["corrected_deck_sha256"]
    assert "Final Output Review Completion" in completion_markdown


def test_apply_deck_revision_plan_rejects_stale_approval_hash(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    plan_path = write_executable_deck_revision_plan(case_dir, session_dir)
    approver = load_deck_revision_approver()
    applier = load_deck_revision_applier()
    approver.approve_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        plan_path=plan_path,
        reviewer="Reviewer",
        understanding_reviewed=True,
        now=fixed_now(),
    )
    normalized = json.loads(plan_path.read_text(encoding="utf-8"))
    normalized["changes"][0]["requested_change"] = "Changed after approval."
    plan_path.write_text(json.dumps(normalized), encoding="utf-8")

    with pytest.raises(
        applier.CaseWorkspaceError,
        match="approval does not match",
    ):
        applier.apply_deck_revision_plan(
            case_dir,
            voice_session=session_dir,
            now=fixed_now(),
        )


def test_apply_deck_revision_plan_rejects_stale_understanding_hash(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    plan_path = write_executable_deck_revision_plan(case_dir, session_dir)
    approver = load_deck_revision_approver()
    applier = load_deck_revision_applier()
    approver.approve_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        plan_path=plan_path,
        reviewer="Reviewer",
        understanding_reviewed=True,
        now=fixed_now(),
    )
    understanding_path = session_dir / "deck_revision_understanding.md"
    understanding_path.write_text(
        understanding_path.read_text(encoding="utf-8") + "\nChanged after approval.\n",
        encoding="utf-8",
    )

    with pytest.raises(
        applier.CaseWorkspaceError,
        match="understanding checkpoint changed after approval",
    ):
        applier.apply_deck_revision_plan(
            case_dir,
            voice_session=session_dir,
            now=fixed_now(),
        )


def test_apply_deck_revision_plan_blocks_wrong_target_text(tmp_path: Path) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    plan_path = write_executable_deck_revision_plan(case_dir, session_dir)
    normalized = json.loads(plan_path.read_text(encoding="utf-8"))
    normalized["changes"][0]["application_patches"][0]["target"][
        "expected_text"
    ] = "Wrong title"
    plan_path.write_text(json.dumps(normalized), encoding="utf-8")
    approver = load_deck_revision_approver()
    applier = load_deck_revision_applier()
    approver.approve_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        plan_path=plan_path,
        reviewer="Reviewer",
        understanding_reviewed=True,
        now=fixed_now(),
    )

    result = applier.apply_deck_revision_plan(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    apply_report = json.loads(result.apply_report_path.read_text(encoding="utf-8"))

    assert apply_report["summary"]["status"] == "partial_or_blocked"
    assert apply_report["summary"]["blocked_patches"] == 1
    assert "target text mismatch" in apply_report["changes"][0]["patches"][0]["message"]


def test_finalize_deck_revision_plan_rejects_untargeted_replace_text(
    tmp_path: Path,
) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    builder = load_deck_revision_workbench_builder()
    finalizer = load_deck_revision_plan_finalizer()
    builder.build_deck_revision_workbench(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    changes_path = session_dir / "deck_revision_changes.json"
    changes_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "codex_deck_revision_plan",
                "changes": [
                    executable_deck_revision_change(
                        change_type="rewrite_body",
                        requested_change="Replace generic EBITDA text.",
                        interpretation=(
                            "The body text should use sharper margin wording."
                        ),
                        rationale="The partner asked for sharper wording.",
                        application_patches=[
                            {
                                "operation": "replace_text",
                                "target": {
                                    "expected_text": (
                                        "Pricing discipline drives EBITDA improvement"
                                    )
                                },
                                "value": {
                                    "old_text": "EBITDA improvement",
                                    "new_text": "margin improvement",
                                },
                            }
                        ],
                    )
                ],
            }
        ),
        encoding="utf-8",
    )

    with pytest.raises(
        finalizer.CaseWorkspaceError,
        match="replace_text` requires target.shape_index",
    ):
        finalizer.finalize_deck_revision_plan(
            case_dir,
            changes_path,
            voice_session=session_dir,
            now=fixed_now(),
        )


def test_verify_deck_revision_output_detects_failed_patch(tmp_path: Path) -> None:
    pytest.importorskip("pptx")

    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    plan_path = write_executable_deck_revision_plan(case_dir, session_dir)
    verifier = load_deck_revision_verifier()
    unchanged_deck = tmp_path / "unchanged.pptx"
    workbench = json.loads(
        (session_dir / "deck_revision_workbench.json").read_text(encoding="utf-8")
    )
    source_deck = case_dir / workbench["source_paths"]["deck_path"]
    unchanged_deck.write_bytes(source_deck.read_bytes())

    result = verifier.verify_deck_revision_output(
        case_dir,
        unchanged_deck,
        voice_session=session_dir,
        plan_path=plan_path,
        now=fixed_now(),
    )
    verification = json.loads(result.report_path.read_text(encoding="utf-8"))

    assert verification["summary"]["status"] == "failed"
    assert verification["summary"]["failed_patches"] == 1
    assert verification["summary"]["failed_success_criteria"] == 1


def test_finalize_deck_revision_plan_rejects_unknown_slide(tmp_path: Path) -> None:
    case_dir, session_dir = build_ready_deck_revision_intake(tmp_path)
    builder = load_deck_revision_workbench_builder()
    finalizer = load_deck_revision_plan_finalizer()
    builder.build_deck_revision_workbench(
        case_dir,
        voice_session=session_dir,
        now=fixed_now(),
    )
    changes_path = session_dir / "deck_revision_changes.json"
    changes_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "codex_deck_revision_plan",
                "changes": [
                    executable_deck_revision_change(
                        slide_number=99,
                        requested_change="Rewrite the headline.",
                        interpretation="The headline should become stronger.",
                        rationale="The partner asked for a stronger headline.",
                        transcript_evidence=[{"quote": "Make the title stronger."}],
                        visual_evidence=[
                            {
                                "evidence_type": "deck_snapshot",
                                "note": "The deck snapshot does not have slide 99.",
                            }
                        ],
                    )
                ],
            }
        ),
        encoding="utf-8",
    )

    with pytest.raises(
        finalizer.CaseWorkspaceError,
        match="slide_number 99 is not in deck snapshot",
    ):
        finalizer.finalize_deck_revision_plan(
            case_dir,
            changes_path,
            voice_session=session_dir,
            now=fixed_now(),
        )


def test_voice_feedback_timeline_requires_video_file(tmp_path: Path) -> None:
    builder = load_voice_feedback_timeline_builder()

    with pytest.raises(
        builder.FeedbackTimelineError,
        match="screen video file is missing",
    ):
        builder.build_feedback_timeline_payload(
            clean_transcript="The headline is too weak. Move the margin point up.",
            timed_transcript_segments=[
                {"text": "the headline is weak", "start_ms": 1000, "end_ms": 2000}
            ],
            video_path=tmp_path / "missing.webm",
            output_path=tmp_path / "feedback_timeline.json",
        )


def test_voice_feedback_timeline_requires_timestamps(tmp_path: Path) -> None:
    builder = load_voice_feedback_timeline_builder()
    video_path = tmp_path / "screen.webm"
    video_path.write_bytes(b"fake video")

    with pytest.raises(
        builder.FeedbackTimelineError,
        match="timed transcript segments do not include timestamps",
    ):
        builder.build_feedback_timeline_payload(
            clean_transcript="The headline is too weak. Move the margin point up.",
            timed_transcript_segments=[{"text": "the headline is weak"}],
            video_path=video_path,
            output_path=tmp_path / "feedback_timeline.json",
        )


def test_voice_feedback_timeline_uses_imageio_ffmpeg_fallback(
    tmp_path: Path,
    monkeypatch,
) -> None:
    builder = load_voice_feedback_timeline_builder()
    video_path = tmp_path / "screen.webm"
    video_path.write_bytes(b"fake video")
    captured_ffmpeg_paths: list[str] = []
    bundled_ffmpeg_path = "/tmp/bundled-ffmpeg"

    class FakeImageioFfmpeg:
        @staticmethod
        def get_ffmpeg_exe() -> str:
            return bundled_ffmpeg_path

    def fake_extract_frame(**kwargs):
        captured_ffmpeg_paths.append(kwargs["ffmpeg_path"])
        frame_path = kwargs["frame_path"]
        frame_path.parent.mkdir(parents=True, exist_ok=True)
        frame_path.write_bytes(b"png")
        return True, ""

    monkeypatch.setattr(builder.shutil, "which", lambda _name: None)
    monkeypatch.setitem(sys.modules, "imageio_ffmpeg", FakeImageioFfmpeg)
    monkeypatch.setattr(builder, "_extract_frame", fake_extract_frame)

    timeline = builder.build_feedback_timeline_payload(
        clean_transcript=(
            "Move the headline up and make the margin story the lead message."
        ),
        timed_transcript_segments=[
            {
                "segment_id": "R1",
                "start_ms": 1000,
                "end_ms": 9000,
                "text": "move the headline up and make the margin story the lead message",
            },
        ],
        video_path=video_path,
        output_path=tmp_path / "feedback_timeline.json",
    )
    entry = timeline["entries"][0]

    assert captured_ffmpeg_paths
    assert set(captured_ffmpeg_paths) == {bundled_ffmpeg_path}
    assert entry["frame_extraction_status"] == "complete"
    assert {frame["status"] for frame in entry["frames"]} == {"extracted"}
    assert entry["use_as_visual_evidence"] is True


def test_voice_feedback_timeline_marks_low_confidence_as_weak_hint(
    tmp_path: Path,
) -> None:
    builder = load_voice_feedback_timeline_builder()
    video_path = tmp_path / "screen.webm"
    video_path.write_bytes(b"fake video")

    timeline = builder.build_feedback_timeline_payload(
        clean_transcript=(
            "Completely unrelated text about governance ownership and succession "
            "timing. Another unrelated sentence about cash flow and board approvals."
        ),
        timed_transcript_segments=[
            {
                "segment_id": "R1",
                "start_ms": 0,
                "end_ms": 8000,
                "text": "on this slide the headline is weak",
            },
            {
                "segment_id": "R2",
                "start_ms": 8000,
                "end_ms": 16000,
                "text": "make margin expansion the lead message",
            },
            {
                "segment_id": "R3",
                "start_ms": 16000,
                "end_ms": 24000,
                "text": "delete the second content box",
            },
            {
                "segment_id": "R4",
                "start_ms": 24000,
                "end_ms": 32000,
                "text": "next page looks fine",
            },
        ],
        video_path=video_path,
        output_path=tmp_path / "feedback_timeline.json",
        extract_frames=False,
    )
    entry = timeline["entries"][0]

    assert entry["alignment_confidence_label"] == "low"
    assert entry["visual_evidence_status"] == "weak_alignment"
    assert entry["use_as_visual_evidence"] is False
    assert "below the usable threshold" in entry["evidence_note"]
    assert entry["skipped_realtime_segment_ids_before_match"] == ["R1"]
    assert entry["realtime_segment_ids"] == ["R1", "R2", "R3", "R4"]
    assert timeline["evidence_summary"]["weak_alignment_entries"] == 1
    assert timeline["evidence_summary"]["skipped_realtime_segment_ids"] == ["R1"]


def _write_test_slide_image(path: Path, *, title: str, body: str, fill: str) -> None:
    from PIL import Image, ImageDraw

    image = Image.new("RGB", (960, 540), fill)
    draw = ImageDraw.Draw(image)
    draw.rectangle((36, 36, 924, 504), outline="#1F2937", width=6)
    draw.text((80, 80), title, fill="#111827")
    draw.text((80, 170), body, fill="#111827")
    image.save(path)


def test_feedback_slide_matcher_matches_bordered_frame_to_rendered_slide(
    tmp_path: Path,
) -> None:
    matcher = load_feedback_slide_matcher()
    case_dir = tmp_path / "case"
    session_dir = case_dir / "voice_sessions" / "20260102103000Z"
    frames_dir = session_dir / "frames"
    render_dir = session_dir / "slide_renders"
    frames_dir.mkdir(parents=True)
    render_dir.mkdir(parents=True)
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")
    slide_one = render_dir / "slide-001.png"
    slide_two = render_dir / "slide-002.png"
    _write_test_slide_image(
        slide_one,
        title="Governance options",
        body="Board roles and decision rights",
        fill="#CBD5E1",
    )
    _write_test_slide_image(
        slide_two,
        title="Margin expansion opportunity",
        body="Pricing discipline and SG&A efficiency",
        fill="#F8FAFC",
    )

    from PIL import Image

    frame = Image.new("RGB", (1280, 720), "#202124")
    pasted_slide = Image.open(slide_two).resize((880, 495))
    frame.paste(pasted_slide, (210, 118))
    frame_path = frames_dir / "F001_000010000.png"
    frame.save(frame_path)
    deck_snapshot_path = session_dir / "deck_snapshot.json"
    deck_snapshot_path.write_text(
        json.dumps(
            {
                "slides": [
                    {"slide_number": 1, "title": "Governance options"},
                    {"slide_number": 2, "title": "Margin expansion opportunity"},
                ]
            }
        ),
        encoding="utf-8",
    )
    timeline_path = session_dir / "feedback_timeline.json"
    timeline_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "entries": [
                    {
                        "feedback_unit_id": "F001",
                        "frames": [
                            {
                                "status": "extracted",
                                "path": "voice_sessions/20260102103000Z/frames/F001_000010000.png",
                                "frame_time_ms": 10000,
                            }
                        ],
                    }
                ],
            }
        ),
        encoding="utf-8",
    )

    result = matcher.match_feedback_timeline_to_deck(
        feedback_timeline_path=timeline_path,
        deck_path=deck_path,
        deck_snapshot_path=deck_snapshot_path,
        base_dir=case_dir,
        slide_render_dir=render_dir,
        now=fixed_now(),
    )
    entry = result["entries"][0]
    frame_match = entry["frames"][0]["slide_match"]

    assert result["slide_matching"]["status"] == "complete"
    assert entry["slide_match"]["best_slide_number"] == 2
    assert frame_match["best_slide_number"] == 2
    assert frame_match["confidence"] in {"high", "medium"}
    assert frame_match["best_crop_label"] != "full_frame"
    assert frame_match["candidates"][0]["slide_number"] == 2


def test_feedback_slide_matcher_invalidates_incomplete_render_cache(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    matcher = load_feedback_slide_matcher()
    render_dir = tmp_path / "slide_renders"
    render_dir.mkdir()
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")
    _write_test_slide_image(
        render_dir / "slide-001.png",
        title="Stale only slide",
        body="This cache is incomplete",
        fill="#FFFFFF",
    )

    def fake_find_soffice(explicit_path: str | None) -> str:
        return "/bin/true"

    def fake_run(
        command: list[str],
        check: bool,
        capture_output: bool,
        text: bool,
        timeout: int,
    ) -> Any:
        outdir = Path(command[command.index("--outdir") + 1])
        outdir.mkdir(parents=True, exist_ok=True)
        (outdir / "deck.pdf").write_bytes(b"%PDF-1.4\n")
        return type("Completed", (), {"returncode": 0, "stdout": "", "stderr": ""})()

    def fake_render_pdf_pages(
        pdf_path: Path,
        output_dir: Path,
        *,
        rendered_slide_numbers: list[int] | None = None,
    ) -> list[Any]:
        slide_one = output_dir / "slide-001.png"
        slide_two = output_dir / "slide-002.png"
        _write_test_slide_image(
            slide_one,
            title="Fresh slide one",
            body="The rerender replaced the stale cache",
            fill="#F8FAFC",
        )
        _write_test_slide_image(
            slide_two,
            title="Fresh slide two",
            body="The expected slide exists",
            fill="#E2E8F0",
        )
        return [
            matcher._SlideRender(slide_number=1, path=slide_one),
            matcher._SlideRender(slide_number=2, path=slide_two),
        ]

    monkeypatch.setattr(matcher, "_find_soffice", fake_find_soffice)
    monkeypatch.setattr(matcher.subprocess, "run", fake_run)
    monkeypatch.setattr(matcher, "_render_pdf_pages", fake_render_pdf_pages)

    renders = matcher._render_deck_slides(
        deck_path,
        render_dir,
        expected_slide_numbers={1, 2},
    )

    assert [render.slide_number for render in renders] == [1, 2]
    assert (render_dir / "slide-002.png").is_file()


def test_feedback_slide_matcher_marks_unrelated_frame_as_no_match(
    tmp_path: Path,
) -> None:
    matcher = load_feedback_slide_matcher()
    case_dir = tmp_path / "case"
    session_dir = case_dir / "voice_sessions" / "20260102103000Z"
    frames_dir = session_dir / "frames"
    render_dir = session_dir / "slide_renders"
    frames_dir.mkdir(parents=True)
    render_dir.mkdir(parents=True)
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")
    slide_one = render_dir / "slide-001.png"
    _write_test_slide_image(
        slide_one,
        title="Governance options",
        body="Board roles and decision rights",
        fill="#FFFFFF",
    )

    from PIL import Image, ImageDraw

    frame = Image.new("RGB", (1280, 720), "#111111")
    draw = ImageDraw.Draw(frame)
    draw.ellipse((380, 140, 900, 640), fill="#78350F")
    frame_path = frames_dir / "F001_000010000.png"
    frame.save(frame_path)
    deck_snapshot_path = session_dir / "deck_snapshot.json"
    deck_snapshot_path.write_text(
        json.dumps({"slides": [{"slide_number": 1, "title": "Governance options"}]}),
        encoding="utf-8",
    )
    timeline_path = session_dir / "feedback_timeline.json"
    timeline_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "entries": [
                    {
                        "feedback_unit_id": "F001",
                        "frames": [
                            {
                                "status": "extracted",
                                "path": "voice_sessions/20260102103000Z/frames/F001_000010000.png",
                                "frame_time_ms": 10000,
                            }
                        ],
                    }
                ],
            }
        ),
        encoding="utf-8",
    )

    result = matcher.match_feedback_timeline_to_deck(
        feedback_timeline_path=timeline_path,
        deck_path=deck_path,
        deck_snapshot_path=deck_snapshot_path,
        base_dir=case_dir,
        slide_render_dir=render_dir,
        now=fixed_now(),
    )

    assert result["entries"][0]["slide_match"]["status"] == "no_match"
    assert result["entries"][0]["frames"][0]["slide_match"]["best_slide_number"] is None


def test_import_hosted_voice_bundle_does_not_update_clara_mandate(
    tmp_path: Path,
) -> None:
    _, case_dir = init_case(tmp_path)
    importer = load_hosted_voice_importer()
    bundle_path = tmp_path / "case-notes-kickoff.json"
    bundle_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "source": "case_notes_hosted_voice",
                "captured_at": "2026-01-02T10:30:00+00:00",
                "model": "gpt-realtime-2",
                "user_transcript": "The first output is a partner readout.",
                "assistant_transcript": "Understood.",
                "extraction_json": {
                    "cleaned_notes_markdown": "# Kickoff\n\nPartner readout.",
                    "clara_mandate": {
                        "engagement_objective": "Create a working partner readout.",
                        "client_decision": "Decide the first transfer scenario.",
                        "clara_understanding": "The work starts from a pragmatic scenario.",
                        "partner_starting_orientation": "Avoid over-formalizing too early.",
                        "sensitive_points": ["Wording around family preference."],
                        "what_clara_should_investigate": ["Industry continuity risks."],
                        "essential_clarifications": [
                            "Which scenario is politically viable?"
                        ],
                        "next_steps": ["Build the first partner brief."],
                    },
                    "entries": [
                        {
                            "kind": "advisor_judgement",
                            "text": "The first output should remain internal.",
                            "rationale": "The client-facing narrative is not ready.",
                        }
                    ],
                    "open_questions": [],
                },
            }
        ),
        encoding="utf-8",
    )

    result = importer.import_hosted_voice_bundle(case_dir, bundle_path)
    mandate = json.loads((case_dir / "clara_mandate.json").read_text())
    judgement = json.loads((case_dir / "judgement_log.json").read_text())

    assert result.clara_mandate_updated is False
    assert mandate["status"] == "not_started"
    assert judgement["entries"][0]["status"] == "pending"


def test_case_update_import_appends_local_copy_and_judgement(
    tmp_path: Path,
) -> None:
    core = load_core()
    source_case = tmp_path / "source-case"
    target_case = tmp_path / "target-case"
    core.initialize_case(
        source_case,
        client="SourceCo",
        project="Advisory",
        objective="Share working case update",
        audience="Advisor",
        output_language="it",
        now=fixed_now(),
    )
    core.initialize_case(
        target_case,
        client="TargetCo",
        project="Advisory",
        objective="Receive working case update",
        audience="Advisor",
        output_language="it",
        now=fixed_now(),
    )
    material = core.ingest_note_text(
        source_case,
        title="Advisor note",
        text="The operating model needs an explicit veto rule.",
        now=fixed_now(),
    )
    core.add_judgement_entries(
        source_case,
        [
            {
                "kind": "decision_implication",
                "text": "The decision pack should state who can veto exceptions.",
                "status": "approved",
                "source_material_ids": [material["id"]],
                "reviewer": "Advisor",
            }
        ],
        now=fixed_now(),
    )
    core.add_open_question(
        source_case,
        question="Who can veto exceptions?",
        why_it_matters="It fixes the governance control owner.",
        now=fixed_now(),
    )

    exported = core.export_case_update(
        source_case,
        exporter="Advisor",
        now=fixed_now(),
    )
    imported = core.import_case_update(
        target_case,
        exported.package_path,
        now=fixed_now(),
    )
    target_registry = json.loads((target_case / "material_registry.json").read_text())
    target_judgement = json.loads((target_case / "judgement_log.json").read_text())
    target_questions = json.loads((target_case / "open_questions.json").read_text())

    assert exported.included_file_count == 1
    assert imported.imported_material_count == 1
    assert imported.imported_judgement_count == 1
    assert imported.imported_open_question_count == 1
    assert imported.conflict_count == 0
    assert Path(target_registry["materials"][0]["path"]).exists()
    assert target_registry["materials"][0]["availability"] == "local_copy"
    assert target_judgement["entries"][0]["status"] == "approved"
    assert target_judgement["entries"][0]["source_material_ids"] == [
        target_registry["materials"][0]["id"]
    ]
    assert target_questions["questions"][0]["question"] == "Who can veto exceptions?"


def test_case_update_import_is_idempotent(tmp_path: Path) -> None:
    core, source_case = init_case(tmp_path / "source")
    _, target_case = init_case(tmp_path / "target")
    material = core.ingest_note_text(
        source_case,
        title="Advisor note",
        text="The transition has a reputational risk.",
        now=fixed_now(),
    )
    core.add_judgement_entries(
        source_case,
        [
            {
                "kind": "advisor_judgement",
                "text": "The transition has a reputational risk.",
                "status": "pending",
                "source_material_ids": [material["id"]],
            }
        ],
        now=fixed_now(),
    )
    exported = core.export_case_update(source_case, now=fixed_now())

    first_import = core.import_case_update(target_case, exported.package_path)
    second_import = core.import_case_update(target_case, exported.package_path)
    target_judgement = json.loads((target_case / "judgement_log.json").read_text())

    assert first_import.imported_judgement_count == 1
    assert second_import.imported_material_count == 0
    assert second_import.imported_judgement_count == 0
    assert second_import.skipped_count == 2
    assert len(target_judgement["entries"]) == 1


def test_case_update_import_logs_conflict_without_overwriting(
    tmp_path: Path,
) -> None:
    core, source_case = init_case(tmp_path / "source")
    _, target_case = init_case(tmp_path / "target")
    core.add_judgement_entries(
        source_case,
        [
            {
                "kind": "advisor_judgement",
                "text": "Original judgement.",
                "status": "pending",
            }
        ],
        now=fixed_now(),
    )
    first_package = core.export_case_update(
        source_case,
        now=datetime(2026, 1, 2, 10, 31, tzinfo=timezone.utc),
    )
    core.import_case_update(target_case, first_package.package_path, now=fixed_now())

    source_judgement_path = source_case / "judgement_log.json"
    source_judgement = json.loads(source_judgement_path.read_text())
    source_judgement["entries"][0]["text"] = "Changed judgement."
    source_judgement_path.write_text(json.dumps(source_judgement), encoding="utf-8")
    second_package = core.export_case_update(
        source_case,
        now=datetime(2026, 1, 2, 10, 32, tzinfo=timezone.utc),
    )

    result = core.import_case_update(
        target_case,
        second_package.package_path,
        now=fixed_now(),
    )
    target_judgement = json.loads((target_case / "judgement_log.json").read_text())
    target_questions = json.loads((target_case / "open_questions.json").read_text())

    assert result.imported_judgement_count == 0
    assert result.conflict_count == 1
    assert target_judgement["entries"][0]["text"] == "Original judgement."
    assert (
        "Review imported judgement conflict"
        in target_questions["questions"][0]["question"]
    )
