from __future__ import annotations

import hashlib
import importlib.util
import json
from pathlib import Path
from typing import Any

import fitz
import jsonschema
import pytest
from docx import Document
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
CLARA_ROOT = ROOT / "plugins" / "clara"
SKILL_ROOT = CLARA_ROOT / "skills" / "advisory-deliverable-validator"
SCRIPT_PATH = SKILL_ROOT / "scripts" / "advisory_validation.py"
CONTRACT_SCHEMA_PATH = CLARA_ROOT / "contracts" / "advisory_contract.v1.schema.json"
REVIEW_SCHEMA_PATH = (
    SKILL_ROOT / "references" / "advisory_validation_review.schema.json"
)


def _validator_module() -> Any:
    spec = importlib.util.spec_from_file_location(
        "clara_advisory_deliverable_validator", SCRIPT_PATH
    )
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _contract(*, format_checks: list[dict[str, Any]] | None = None) -> dict[str, Any]:
    return {
        "schema_version": "1.0",
        "contract_status": "ready_for_handoff",
        "decision": "Whether to proceed with the proposed channel pilot.",
        "purpose": "Give the steering group a source-backed recommendation.",
        "audience": "Steering group and engagement partner",
        "deliverable_type": "advisory memo",
        "output_language": "en",
        "scope_included": ["Commercial evidence", "Implementation conditions"],
        "scope_excluded": ["Legal and tax advice"],
        "available_inputs": [
            {
                "id": "deliverable",
                "description": "Completed advisory memo and selected evidence",
                "status": "available",
                "source_ref": "advisory_memo.md",
            }
        ],
        "evidence_requirements": [
            {
                "id": "claim-support",
                "requirement": "Material factual claims identify support",
                "rationale": "Delivery readiness depends on traceable material claims.",
                "status": "available",
                "input_ids": ["deliverable"],
            }
        ],
        "analysis_plan": [
            {
                "id": "validate-deliverable",
                "objective": "Review support and challenge the recommendation",
                "method": "Apply the advisory deliverable validation workflow",
                "input_ids": ["deliverable"],
                "output": "Structured validation package",
            }
        ],
        "assumptions": [
            {
                "statement": "The supplied extract is complete",
                "status": "provisional",
                "materiality": "material",
            }
        ],
        "unresolved_questions": [],
        "success_criteria": ["The decision and conditions are explicit"],
        "selected_clara_workflow": "clara:clara",
        "validation_profile": {
            "review_dimensions": [
                "contract_conformance",
                "factual_source_support",
                "calculations_data_provenance",
                "reasoning_assumptions",
                "contradictions_missing_evidence",
                "recommendation_evidence_decision_fit",
                "professional_judgement_boundaries",
                "correction_needs",
                "residual_uncertainty",
                "delivery_readiness",
            ],
            "format_checks": format_checks or [],
        },
        "validation_scope": {
            "coverage": "all_material_content",
            "included_sections": ["Entire deliverable"],
            "excluded_sections": [],
            "limitations": [],
        },
        "correction_policy": {
            "mode": "separate_artifact",
            "preserve_original": True,
            "allowed": True,
            "approval_required_before_delivery": True,
        },
        "professional_judgement_policy": {
            "owner": "Engagement partner",
            "model_role": "Identify issues and draft corrections for review",
            "approval_required_before_delivery": True,
        },
        "source_facts": [
            {
                "category": "constraint",
                "text": "Legal and tax advice remain outside scope.",
                "source_anchor": "Legal and tax advice",
                "input_id": "deliverable",
            }
        ],
        "explicit_questions": [],
        "generation_handoff": {
            "workflow": "clara:clara",
            "objective": "Validate the completed advisory memo.",
            "input_ids": ["deliverable"],
            "instructions": ["Apply the declared validation profile."],
            "expected_outputs": ["Advisory validation package"],
            "preserve_specialist_authority": True,
        },
        "model_review": {
            "method": "model_led_assignment_contract_review",
            "dimensions": {
                "material_facts_dates_numbers_entities_constraints_questions": "conforms",
                "decision_purpose_audience_deliverable": "conforms",
                "scope_inputs_evidence": "conforms",
                "analysis_and_success_criteria": "conforms",
                "workflow_and_generation_handoff": "conforms",
                "validation_and_professional_judgement": "conforms",
            },
            "overall_status": "conforms",
        },
    }


def _dimension_review() -> dict[str, Any]:
    return {
        "status": "conforms",
        "analysis": "The reviewed content conforms on this dimension.",
        "evidence_refs": ["extracted_deliverable.md"],
        "issues": [],
        "correction_status": "not_needed",
        "professional_review_required": False,
    }


def _review(
    inventory: dict[str, Any],
    *,
    format_checks: list[dict[str, Any]] | None = None,
    correction_status: str = "not_required",
    corrected_artifact: Path | None = None,
) -> dict[str, Any]:
    dimensions = {
        dimension: _dimension_review()
        for dimension in _validator_module().REVIEW_DIMENSIONS
    }
    if correction_status == "completed":
        assert corrected_artifact is not None
        dimensions["correction_needs"]["correction_status"] = "completed"
    corrected_hash = (
        hashlib.sha256(corrected_artifact.read_bytes()).hexdigest()
        if corrected_artifact is not None
        else ""
    )
    coverage_inventory = json.loads(
        Path(inventory["coverage_inventory_path"]).read_text(encoding="utf-8")
    )
    considered_unit_ids = [unit["id"] for unit in coverage_inventory["units"]]
    provenance_mode = inventory["lineage"]["provenance_mode"]
    unit_assessments = [
        {
            "unit_id": unit_id,
            "status": (
                "reviewed_material_claims"
                if index == 0
                else "reviewed_no_material_claims"
            ),
            "material_claim_ids": [],
            "untracked_claim_ids": (["external-claim-0001"] if index == 0 else []),
            "analysis": (
                "The material recommendation was selected for claim review."
                if index == 0
                else "No additional material claim was identified in this unit."
            ),
        }
        for index, unit_id in enumerate(considered_unit_ids)
    ]
    return {
        "schema_version": "1.3",
        "language": "en",
        "advisory_contract_sha256": inventory["advisory_contract_sha256"],
        "deliverable_sha256": inventory["source_sha256"],
        "coverage_inventory_sha256": inventory["coverage_inventory_sha256"],
        "lineage_inventory_sha256": inventory["lineage"]["lineage_inventory_sha256"],
        "coverage_review": {
            "selection_method": "model_led_materiality_review",
            "scope": "all_material_content",
            "reviewed_sections": ["Entire deliverable"],
            "omitted_sections": [],
            "considered_unit_ids": considered_unit_ids,
            "omitted_unit_ids": [],
            "unit_assessments": unit_assessments,
            "limitations": [],
            "analysis": "All material content was reviewed.",
        },
        "lineage_review": {
            "provenance_mode": provenance_mode,
            "selection_method": "model_led_claim_chain_review",
            "reviewed_claim_ids": [],
            "chain_assessments": [],
            "untracked_material_claims": [
                {
                    "id": "external-claim-0001",
                    "statement": "Proceed with a bounded pilot.",
                    "deliverable_locations": ["Recommendation"],
                    "evidence_ids": ["source-0001"],
                    "dependency_claim_ids": [],
                    "support_status": "adequate",
                    "reasoning_status": "sound",
                    "contradiction_resolution": "",
                    "analysis": "The material recommendation and its stated basis were reviewed.",
                    "recheck": {
                        "required": False,
                        "kind": "none",
                        "status": "not_required",
                        "evidence_ids": [],
                        "analysis": "No targeted recheck was required for this fixture.",
                    },
                    "resolution": {
                        "status": "no_change",
                        "explanation": "No material weakness was identified in the fixture.",
                    },
                }
            ],
            "limitations": [],
            "analysis": "The review matched the final claim to the supplied support because no generation-time lineage was supplied.",
        },
        "dimension_reviews": dimensions,
        "findings": [],
        "format_specific_checks": format_checks or [],
        "correction": {
            "status": correction_status,
            "summary": (
                "A separate correction was completed."
                if correction_status == "completed"
                else "No correction was required."
            ),
            "corrected_artifact": (
                str(corrected_artifact.resolve())
                if corrected_artifact is not None
                else ""
            ),
            "corrected_artifact_sha256": corrected_hash,
            "corrected_inventory_sha256": (
                "0" * 64 if correction_status == "completed" else ""
            ),
            "corrected_review_sha256": (
                "0" * 64 if correction_status == "completed" else ""
            ),
            "unresolved_changes": [],
        },
        "approvals": {
            "professional_judgement": {
                "status": "approved",
                "approved_by": "Engagement partner",
                "evidence_refs": ["Explicit test-fixture approval"],
            },
            "correction": {
                "status": (
                    "approved" if correction_status == "completed" else "not_required"
                ),
                "approved_by": (
                    "Engagement partner" if correction_status == "completed" else ""
                ),
                "evidence_refs": (
                    ["Explicit corrected-artifact approval"]
                    if correction_status == "completed"
                    else []
                ),
            },
        },
        "overall_assessment": {
            "outcome": "ready",
            "analysis": "No material unresolved issue was identified.",
            "residual_uncertainties": [],
            "professional_review_items": [],
        },
        "delivery_readiness": {"status": "ready", "conditions": []},
    }


def _write_deliverable(path: Path) -> None:
    suffix = path.suffix
    text = "# Recommendation\n\nProceed with a bounded pilot. Source: [1]. Value: EUR 120,000."
    if suffix in {".md", ".markdown", ".txt"}:
        path.write_text(text, encoding="utf-8")
    elif suffix == ".html":
        path.write_text(
            '<html><body><h1>Recommendation</h1><p>Proceed with a bounded pilot.</p><a href="https://example.com/evidence">Evidence</a></body></html>',
            encoding="utf-8",
        )
    elif suffix == ".docx":
        document = Document()
        document.add_heading("Recommendation", level=1)
        document.add_paragraph("Proceed with a bounded pilot. Value: EUR 120,000.")
        document.save(path)
    elif suffix == ".pptx":
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Recommendation"
        slide.placeholders[1].text = "Proceed with a bounded pilot. Value: EUR 120,000."
        presentation.save(path)
    elif suffix == ".pdf":
        document = fitz.open()
        page = document.new_page()
        page.insert_text(
            (72, 72), "Recommendation: proceed with a bounded pilot. Value EUR 120,000."
        )
        document.save(path)
        document.close()
    else:
        raise AssertionError(f"unexpected test suffix: {suffix}")


def _prepare(tmp_path: Path, *, format_checks: list[dict[str, Any]] | None = None):
    validator = _validator_module()
    deliverable = tmp_path / "advisory_memo.md"
    contract_path = tmp_path / "advisory_contract.json"
    output_dir = tmp_path / "validation"
    _write_deliverable(deliverable)
    support = tmp_path / "selected_support.txt"
    support.write_text(
        "Reviewed support for the bounded-pilot recommendation.",
        encoding="utf-8",
    )
    contract_path.write_text(
        json.dumps(_contract(format_checks=format_checks)), encoding="utf-8"
    )
    paths = validator.prepare_validation(
        deliverable,
        contract_path,
        output_dir,
        source_files=[support],
    )
    inventory = json.loads(paths["deliverable_inventory"].read_text(encoding="utf-8"))
    return validator, deliverable, contract_path, output_dir, paths, inventory


def _prepare_corrected_review(
    validator: Any,
    corrected: Path,
    contract_path: Path,
    tmp_path: Path,
) -> tuple[Path, Path]:
    corrected_output = tmp_path / "corrected_validation"
    paths = validator.prepare_validation(
        corrected,
        contract_path,
        corrected_output,
        source_files=[tmp_path / "selected_support.txt"],
    )
    inventory = json.loads(paths["deliverable_inventory"].read_text(encoding="utf-8"))
    corrected_review = corrected_output / "corrected_review.json"
    corrected_review.write_text(json.dumps(_review(inventory)), encoding="utf-8")
    return paths["deliverable_inventory"], corrected_review


def _lineage_receipt(evidence_id: str, observation: str) -> dict[str, Any]:
    return {
        "id": evidence_id,
        "evidence_type": "management_assertion",
        "recorded_at": "2026-08-18T08:00:00+00:00",
        "recorded_by": "clara:clara",
        "capture_status": "assertion_only",
        "source": {
            "material_ids": [],
            "url": "",
            "locator": "management interview",
            "artifact_refs": [],
        },
        "observation": observation,
        "scope": "What management stated in the interview.",
        "limitations": ["The assertion is not independent verification."],
        "verification": {
            "status": "not_checked",
            "checked_at": "",
            "method": "",
            "notes": [],
        },
        "rechecks_evidence_id": "",
        "supersedes_evidence_id": "",
    }


def _lineage_claim(
    claim_id: str,
    statement: str,
    *,
    evidence_ids: list[str] | None = None,
    dependency_ids: list[str] | None = None,
    claim_type: str = "assertion",
) -> dict[str, Any]:
    evidence_ids = evidence_ids or []
    dependency_ids = dependency_ids or []
    return {
        "id": claim_id,
        "statement": statement,
        "claim_type": claim_type,
        "recorded_at": "2026-08-18T08:00:00+00:00",
        "recorded_by": "clara:clara",
        "provenance": {
            "workflow": "clara:clara",
            "step": "analysis",
            "artifact": "analysis.md",
            "locator": claim_id,
        },
        "evidence_links": [
            {
                "evidence_id": evidence_id,
                "relationship": "supports",
                "analysis": "The receipt records the stated premise.",
                "proves": statement,
                "does_not_prove": "Any broader or downstream conclusion.",
            }
            for evidence_id in evidence_ids
        ],
        "dependency": {
            "mode": "all_of" if dependency_ids else "none",
            "claim_ids": dependency_ids,
            "derivation_type": "reasoning" if dependency_ids else "direct",
            "explanation": (
                "Every named premise is required for this conclusion."
                if dependency_ids
                else "Directly stated premise."
            ),
            "calculation_evidence_id": "",
        },
        "decision_use": "direct" if dependency_ids else "supporting",
        "uncertainty": [],
        "professional_judgement_required": False,
        "appearances": [],
        "state": "active",
        "supersedes_claim_id": "",
    }


def _chain_assessment(claim: dict[str, Any]) -> dict[str, Any]:
    return {
        "claim_id": claim["id"],
        "statement": claim["statement"],
        "deliverable_locations": ["Recommendation"],
        "evidence_ids": [
            link["evidence_id"] for link in claim.get("evidence_links", [])
        ],
        "dependency_claim_ids": claim["dependency"]["claim_ids"],
        "support_status": "adequate",
        "reasoning_status": "sound",
        "contradiction_resolution": "",
        "analysis": "The stated basis and dependency were challenged in context.",
        "recheck": {
            "required": False,
            "kind": "none",
            "status": "not_required",
            "evidence_ids": [],
            "analysis": "No targeted recheck was required.",
        },
        "resolution": {
            "status": "no_change",
            "explanation": "No material weakness was detected.",
        },
    }


def test_advisory_contract_schema_covers_the_stable_cross_workflow_boundary() -> None:
    payload = _contract()
    schema = json.loads(CONTRACT_SCHEMA_PATH.read_text(encoding="utf-8"))

    jsonschema.Draft202012Validator(schema).validate(payload)

    validator = _validator_module()
    assert validator.validate_advisory_contract(payload) == []
    assert set(validator.REQUIRED_CONTRACT_FIELDS) == {
        "schema_version",
        "decision",
        "purpose",
        "audience",
        "deliverable_type",
        "output_language",
        "scope_included",
        "scope_excluded",
        "available_inputs",
        "evidence_requirements",
        "analysis_plan",
        "assumptions",
        "unresolved_questions",
        "success_criteria",
        "selected_clara_workflow",
        "validation_profile",
        "validation_scope",
        "correction_policy",
        "professional_judgement_policy",
    }


@pytest.mark.parametrize(
    "suffix", [".md", ".markdown", ".txt", ".html", ".pdf", ".docx", ".pptx"]
)
def test_prepare_extracts_every_supported_primary_format(
    tmp_path: Path, suffix: str
) -> None:
    validator = _validator_module()
    deliverable = tmp_path / f"deliverable{suffix}"
    _write_deliverable(deliverable)

    text, metadata = validator.read_supported_deliverable(deliverable)

    assert "Recommendation" in text
    assert metadata["parser"] in {
        "plain_text",
        "html_text",
        "pymupdf_text",
        "python_docx",
        "python_pptx_visible_text",
    }


def test_html_extraction_excludes_non_visible_script_style_and_template_content(
    tmp_path: Path,
) -> None:
    validator = _validator_module()
    deliverable = tmp_path / "deliverable.html"
    deliverable.write_text(
        "<html><style>.private{display:none}</style><body>Visible"
        "<script>hiddenToken=42</script><template>Draft secret</template>"
        "</body></html>",
        encoding="utf-8",
    )

    text, _ = validator.read_supported_deliverable(deliverable)

    assert text == "Visible"


@pytest.mark.parametrize("suffix", [".docx", ".pptx", ".pdf"])
def test_corrupt_supported_document_reports_a_readable_validation_error(
    tmp_path: Path,
    suffix: str,
) -> None:
    validator = _validator_module()
    deliverable = tmp_path / f"corrupt{suffix}"
    deliverable.write_bytes(b"not a valid document package")

    with pytest.raises(
        validator.AdvisoryValidationError,
        match="unreadable or damaged",
    ):
        validator.read_supported_deliverable(deliverable)


def test_prepare_rejects_spreadsheet_as_primary_but_allows_it_as_evidence(
    tmp_path: Path,
) -> None:
    validator = _validator_module()
    spreadsheet = tmp_path / "analysis.xlsx"
    spreadsheet.write_bytes(b"synthetic workbook fixture")
    contract_path = tmp_path / "advisory_contract.json"
    contract_path.write_text(json.dumps(_contract()), encoding="utf-8")

    with pytest.raises(
        validator.AdvisoryValidationError,
        match="unsupported primary deliverable format",
    ):
        validator.prepare_validation(
            spreadsheet, contract_path, tmp_path / "unsupported-primary"
        )

    deliverable = tmp_path / "memo.md"
    _write_deliverable(deliverable)
    paths = validator.prepare_validation(
        deliverable,
        contract_path,
        tmp_path / "with-evidence",
        source_files=[spreadsheet],
    )
    source_inventory = json.loads(paths["source_inventory"].read_text(encoding="utf-8"))
    assert source_inventory["sources"][0]["supported_source_type"] is True


def test_prepare_inventories_links_values_and_preserves_original(
    tmp_path: Path,
) -> None:
    validator, deliverable, _, _, paths, inventory = _prepare(tmp_path)
    original_bytes = deliverable.read_bytes()

    citations = json.loads(paths["citation_inventory"].read_text(encoding="utf-8"))
    calculations = json.loads(
        paths["calculation_inventory"].read_text(encoding="utf-8")
    )

    assert inventory["source_sha256"] == validator._sha256(deliverable)
    assert citations["citation_markers"] == ["[1]"]
    assert "EUR 120,000" in calculations["numeric_tokens"]
    assert deliverable.read_bytes() == original_bytes
    assert inventory["boundary"]["semantic_selection"] == "model_led"
    assert inventory["boundary"]["hidden_model_api_calls"] is False


def test_prepare_rejects_an_output_path_that_aliases_the_primary_deliverable(
    tmp_path: Path,
) -> None:
    validator = _validator_module()
    output_dir = tmp_path / "validation"
    output_dir.mkdir()
    deliverable = output_dir / "extracted_deliverable.md"
    original_bytes = b"# Original with trailing spaces   \n"
    deliverable.write_bytes(original_bytes)
    contract_path = tmp_path / "advisory_contract.json"
    contract_path.write_text(json.dumps(_contract()), encoding="utf-8")

    with pytest.raises(
        validator.AdvisoryValidationError,
        match="refusing to overwrite primary deliverable",
    ):
        validator.prepare_validation(deliverable, contract_path, output_dir)

    assert deliverable.read_bytes() == original_bytes


def test_prepare_rejects_an_output_path_that_aliases_the_input_contract(
    tmp_path: Path,
) -> None:
    validator = _validator_module()
    output_dir = tmp_path / "validation"
    output_dir.mkdir()
    deliverable = tmp_path / "advisory_memo.md"
    _write_deliverable(deliverable)
    contract_path = output_dir / "extracted_deliverable.md"
    original_bytes = json.dumps(_contract()).encode("utf-8")
    contract_path.write_bytes(original_bytes)

    with pytest.raises(
        validator.AdvisoryValidationError,
        match="refusing to overwrite advisory contract",
    ):
        validator.prepare_validation(deliverable, contract_path, output_dir)

    assert contract_path.read_bytes() == original_bytes


def test_package_accepts_a_complete_model_led_review(tmp_path: Path) -> None:
    validator, _, contract_path, output_dir, paths, inventory = _prepare(tmp_path)
    review = _review(inventory)
    review_path = output_dir / "advisory_validation_review_draft.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    schema = json.loads(REVIEW_SCHEMA_PATH.read_text(encoding="utf-8"))
    jsonschema.Draft202012Validator(schema).validate(review)
    package_paths, audit = validator.package_validation(
        paths["deliverable_inventory"],
        review_path,
        contract_path,
        output_dir,
    )

    assert audit["record_complete"] is True
    assert audit["effective_delivery_readiness"] == "ready"
    assert package_paths["review"].is_file()
    assert package_paths["audit"].is_file()
    assert "Delivery readiness: ready" in package_paths["package"].read_text(
        encoding="utf-8"
    )


def test_ready_review_requires_explicit_professional_judgement_approval(
    tmp_path: Path,
) -> None:
    validator, _, contract_path, _, _, inventory = _prepare(tmp_path)
    contract = json.loads(contract_path.read_text(encoding="utf-8"))
    review = _review(inventory)
    review["approvals"]["professional_judgement"] = {
        "status": "pending",
        "approved_by": "",
        "evidence_refs": [],
    }

    errors = validator.validate_review_record(
        review,
        contract,
        provenance_mode="matched_support",
        matched_support_source_ids={"source-0001"},
    )

    assert "professional-judgement approval is required before delivery" in errors


def test_ready_review_rejects_partially_conforming_dimension(
    tmp_path: Path,
) -> None:
    validator, _, contract_path, _, _, inventory = _prepare(tmp_path)
    contract = json.loads(contract_path.read_text(encoding="utf-8"))
    review = _review(inventory)
    review["dimension_reviews"]["factual_source_support"][
        "status"
    ] = "partially_conforms"

    errors = validator.validate_review_record(review, contract)

    assert "ready status cannot coexist with unresolved review attention" in errors


def test_delivery_ready_review_rejects_unresolved_professional_review_flag(
    tmp_path: Path,
) -> None:
    validator, _, contract_path, _, _, inventory = _prepare(tmp_path)
    contract = json.loads(contract_path.read_text(encoding="utf-8"))
    review = _review(inventory)
    review["dimension_reviews"]["professional_judgement_boundaries"][
        "professional_review_required"
    ] = True
    review["overall_assessment"]["outcome"] = "ready_with_residual_uncertainty"
    review["overall_assessment"]["residual_uncertainties"] = [
        "Professional review remains pending."
    ]
    review["delivery_readiness"]["status"] = "ready_with_residual_uncertainty"

    errors = validator.validate_review_record(review, contract)

    assert (
        "delivery-ready status cannot coexist with a dimension requiring professional review"
        in errors
    )


def test_delivery_ready_review_rejects_a_finding_requiring_professional_review(
    tmp_path: Path,
) -> None:
    validator, _, contract_path, _, _, inventory = _prepare(tmp_path)
    contract = json.loads(contract_path.read_text(encoding="utf-8"))
    review = _review(inventory)
    review["findings"] = [
        {
            "id": "finding-1",
            "dimension": "professional_judgement_boundaries",
            "finding": "The final risk acceptance remains professionally owned.",
            "status": "judgment_required",
            "evidence_refs": ["advisory_contract.json"],
            "correction_action": "Obtain the engagement partner's decision.",
            "correction_status": "professional_review_required",
            "professional_review_required": True,
        }
    ]
    review["overall_assessment"]["outcome"] = "ready_with_residual_uncertainty"
    review["overall_assessment"]["residual_uncertainties"] = [
        "Professional review remains pending."
    ]
    review["delivery_readiness"]["status"] = "ready_with_residual_uncertainty"

    errors = validator.validate_review_record(review, contract)

    assert (
        "delivery-ready status cannot coexist with a finding requiring professional review"
        in errors
    )


def test_partially_conforming_review_can_be_ready_with_residual_uncertainty(
    tmp_path: Path,
) -> None:
    validator, _, contract_path, _, _, inventory = _prepare(tmp_path)
    contract = json.loads(contract_path.read_text(encoding="utf-8"))
    review = _review(inventory)
    review["dimension_reviews"]["factual_source_support"][
        "status"
    ] = "partially_conforms"
    review["overall_assessment"]["outcome"] = "ready_with_residual_uncertainty"
    review["overall_assessment"]["residual_uncertainties"] = [
        "One immaterial source limitation remains explicit."
    ]
    review["delivery_readiness"]["status"] = "ready_with_residual_uncertainty"

    errors = validator.validate_review_record(
        review,
        contract,
        provenance_mode="matched_support",
        matched_support_source_ids={"source-0001"},
    )

    assert errors == []


def test_completed_correction_requires_explicit_correction_approval(
    tmp_path: Path,
) -> None:
    validator, deliverable, contract_path, _, _, inventory = _prepare(tmp_path)
    contract = json.loads(contract_path.read_text(encoding="utf-8"))
    corrected = tmp_path / "advisory_memo_corrected.md"
    corrected.write_text(
        deliverable.read_text(encoding="utf-8") + "\nCondition added.\n",
        encoding="utf-8",
    )
    review = _review(
        inventory,
        correction_status="completed",
        corrected_artifact=corrected,
    )
    review["approvals"]["correction"] = {
        "status": "pending",
        "approved_by": "",
        "evidence_refs": [],
    }

    errors = validator.validate_review_record(review, contract)

    assert "correction approval is required before delivery" in errors


def test_package_blocks_a_missing_required_format_check(tmp_path: Path) -> None:
    required_check = {
        "workflow": "clara:reporting-engine",
        "requirement": "required",
        "reason": "The recommendation relies on spreadsheet calculations.",
        "artifact_refs": [],
    }
    validator, _, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path, format_checks=[required_check]
    )
    review = _review(inventory)
    review["overall_assessment"]["outcome"] = "blocked"
    review["delivery_readiness"]["status"] = "blocked"
    review_path = output_dir / "advisory_validation_review_draft.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    _, audit = validator.package_validation(
        paths["deliverable_inventory"], review_path, contract_path, output_dir
    )

    assert audit["record_complete"] is False
    assert audit["effective_delivery_readiness"] == "blocked"
    assert "required format check is missing: clara:reporting-engine" in audit["errors"]


def test_package_blocks_a_passed_format_check_with_a_missing_artifact(
    tmp_path: Path,
) -> None:
    artifact_name = "reporting_engine_validation.json"
    required_check = {
        "workflow": "clara:reporting-engine",
        "requirement": "required",
        "reason": "The recommendation relies on spreadsheet calculations.",
        "artifact_refs": [artifact_name],
    }
    validator, _, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path, format_checks=[required_check]
    )
    review = _review(
        inventory,
        format_checks=[
            {
                "workflow": "clara:reporting-engine",
                "status": "passed",
                "artifact_refs": [artifact_name],
                "analysis": "The Reporting Engine check passed.",
            }
        ],
    )
    review_path = output_dir / "advisory_validation_review_draft.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    _, audit = validator.package_validation(
        paths["deliverable_inventory"], review_path, contract_path, output_dir
    )

    expected_artifact = tmp_path / artifact_name
    assert audit["record_complete"] is False
    assert audit["checks"]["format_check_artifacts_exist"] is False
    assert (
        "passed format check artifact does not exist: "
        f"clara:reporting-engine: {expected_artifact}"
    ) in audit["errors"]


def test_package_hashes_a_present_required_format_check_artifact(
    tmp_path: Path,
) -> None:
    artifact_name = "reporting_engine_validation.json"
    artifact = tmp_path / artifact_name
    artifact.write_text('{"status":"passed"}\n', encoding="utf-8")
    required_check = {
        "workflow": "clara:reporting-engine",
        "requirement": "required",
        "reason": "The recommendation relies on spreadsheet calculations.",
        "artifact_refs": [artifact_name],
    }
    validator, _, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path, format_checks=[required_check]
    )
    review = _review(
        inventory,
        format_checks=[
            {
                "workflow": "clara:reporting-engine",
                "status": "passed",
                "artifact_refs": [artifact_name],
                "analysis": "The Reporting Engine check passed.",
            }
        ],
    )
    review_path = output_dir / "advisory_validation_review_draft.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    _, audit = validator.package_validation(
        paths["deliverable_inventory"], review_path, contract_path, output_dir
    )

    assert audit["record_complete"] is False
    assert audit["checks"]["format_check_artifacts_exist"] is False
    assert any(
        "lacks authoritative result artifacts" in error for error in audit["errors"]
    )


def test_package_accepts_authoritative_reporting_engine_result(
    tmp_path: Path,
) -> None:
    artifact_name = "render_manifest.json"
    artifact = tmp_path / artifact_name
    artifact.write_text(
        json.dumps(
            {
                "schema_version": "0.2",
                "owner": "clara.reporting-engine",
                "runner": {"returncode": 0, "status": "ok"},
                "render_proof": {"status": "rendered"},
            }
        ),
        encoding="utf-8",
    )
    required_check = {
        "workflow": "clara:reporting-engine",
        "requirement": "required",
        "reason": "The recommendation relies on spreadsheet calculations.",
        "artifact_refs": [artifact_name],
    }
    validator, _, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path, format_checks=[required_check]
    )
    review = _review(
        inventory,
        format_checks=[
            {
                "workflow": "clara:reporting-engine",
                "status": "passed",
                "artifact_refs": [artifact_name],
                "analysis": "The authoritative Reporting Engine result passed.",
            }
        ],
    )
    review_path = output_dir / "review.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    _, audit = validator.package_validation(
        paths["deliverable_inventory"], review_path, contract_path, output_dir
    )

    assert audit["record_complete"] is True
    assert audit["format_check_artifacts"][0]["authoritative_result"] == {
        "kind": "reporting_engine_render",
        "passed": True,
    }


def test_package_rejects_html_qa_reports_for_different_deliverable_bytes(
    tmp_path: Path,
) -> None:
    static_name = "html-build-report.json"
    browser_name = "browser-qa.json"
    required_check = {
        "workflow": "clara:html-deck",
        "requirement": "required",
        "reason": "The completed deliverable is an HTML deck.",
        "artifact_refs": [static_name, browser_name],
    }
    validator, _, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path, format_checks=[required_check]
    )
    wrong_sha256 = "0" * 64
    (tmp_path / static_name).write_text(
        json.dumps(
            {
                "schema_version": "clara.html_deck_build.v1",
                "result": "pass",
                "input": {"sha256": inventory["source_sha256"]},
                "deck": {},
                "checks": [],
                "summary": {},
            }
        ),
        encoding="utf-8",
    )
    (tmp_path / browser_name).write_text(
        json.dumps(
            {
                "schema_version": "clara.html_deck_browser_qa.v1",
                "result": "pass",
                "input": {"sha256": wrong_sha256},
                "browser": {},
                "viewports": [],
            }
        ),
        encoding="utf-8",
    )
    review = _review(
        inventory,
        format_checks=[
            {
                "workflow": "clara:html-deck",
                "status": "passed",
                "artifact_refs": [static_name, browser_name],
                "analysis": "Static and browser checks were supplied.",
            }
        ],
    )
    review_path = output_dir / "review.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    _, audit = validator.package_validation(
        paths["deliverable_inventory"], review_path, contract_path, output_dir
    )

    assert audit["record_complete"] is False
    assert (
        "passed HTML Deck result is not bound to the prepared deliverable: "
        "html_browser_qa"
    ) in audit["errors"]


def test_package_rejects_an_output_path_that_aliases_a_format_check_artifact(
    tmp_path: Path,
) -> None:
    artifact_ref = "validation/validation_audit.json"
    required_check = {
        "workflow": "clara:reporting-engine",
        "requirement": "required",
        "reason": "The recommendation relies on spreadsheet calculations.",
        "artifact_refs": [artifact_ref],
    }
    validator, _, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path, format_checks=[required_check]
    )
    format_artifact = output_dir / "validation_audit.json"
    original_bytes = b'{"authoritative":"format check"}\n'
    format_artifact.write_bytes(original_bytes)
    review = _review(
        inventory,
        format_checks=[
            {
                "workflow": "clara:reporting-engine",
                "status": "passed",
                "artifact_refs": [artifact_ref],
                "analysis": "The Reporting Engine check passed.",
            }
        ],
    )
    review_path = output_dir / "advisory_validation_review_draft.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    with pytest.raises(
        validator.AdvisoryValidationError,
        match="refusing to overwrite format-check artifact",
    ):
        validator.package_validation(
            paths["deliverable_inventory"], review_path, contract_path, output_dir
        )

    assert format_artifact.read_bytes() == original_bytes


def test_completed_correction_requires_a_separate_changed_artifact(
    tmp_path: Path,
) -> None:
    validator, deliverable, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path
    )
    corrected = tmp_path / "advisory_memo_corrected.md"
    corrected.write_text(
        deliverable.read_text(encoding="utf-8") + "\nCondition: confirm the owner.\n",
        encoding="utf-8",
    )
    review = _review(
        inventory,
        correction_status="completed",
        corrected_artifact=corrected,
    )
    corrected_inventory, corrected_review = _prepare_corrected_review(
        validator,
        corrected,
        contract_path,
        tmp_path,
    )
    review["correction"]["corrected_inventory_sha256"] = hashlib.sha256(
        corrected_inventory.read_bytes()
    ).hexdigest()
    review["correction"]["corrected_review_sha256"] = hashlib.sha256(
        corrected_review.read_bytes()
    ).hexdigest()
    review_path = output_dir / "advisory_validation_review_draft.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    _, audit = validator.package_validation(
        paths["deliverable_inventory"],
        review_path,
        contract_path,
        output_dir,
        corrected_deliverable=corrected,
        corrected_deliverable_inventory=corrected_inventory,
        corrected_review=corrected_review,
    )

    assert audit["record_complete"] is True
    assert audit["checks"]["original_unchanged"] is True
    assert audit["checks"]["separate_corrected_artifact"] is True
    assert audit["checks"]["corrected_artifact_re_reviewed"] is True
    assert audit["corrected_artifact"]["path"] == str(corrected.resolve())


def test_completed_correction_rejects_a_second_review_with_stale_coverage_hash(
    tmp_path: Path,
) -> None:
    validator, deliverable, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path
    )
    corrected = tmp_path / "advisory_memo_corrected.md"
    corrected.write_text(
        deliverable.read_text(encoding="utf-8") + "\nCondition: confirm the owner.\n",
        encoding="utf-8",
    )
    corrected_inventory, corrected_review = _prepare_corrected_review(
        validator,
        corrected,
        contract_path,
        tmp_path,
    )
    stale_review = json.loads(corrected_review.read_text(encoding="utf-8"))
    stale_review["coverage_inventory_sha256"] = "f" * 64
    corrected_review.write_text(json.dumps(stale_review), encoding="utf-8")
    review = _review(
        inventory,
        correction_status="completed",
        corrected_artifact=corrected,
    )
    review["correction"]["corrected_inventory_sha256"] = hashlib.sha256(
        corrected_inventory.read_bytes()
    ).hexdigest()
    review["correction"]["corrected_review_sha256"] = hashlib.sha256(
        corrected_review.read_bytes()
    ).hexdigest()
    review_path = output_dir / "advisory_validation_review_draft.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    _, audit = validator.package_validation(
        paths["deliverable_inventory"],
        review_path,
        contract_path,
        output_dir,
        corrected_deliverable=corrected,
        corrected_deliverable_inventory=corrected_inventory,
        corrected_review=corrected_review,
    )

    assert audit["record_complete"] is False
    assert "corrected review is not bound to its coverage inventory" in audit["errors"]


def test_completed_correction_is_bound_to_the_declared_path(tmp_path: Path) -> None:
    validator, deliverable, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path
    )
    corrected = tmp_path / "advisory_memo_corrected.md"
    corrected.write_text(
        deliverable.read_text(encoding="utf-8") + "\nCondition added.\n",
        encoding="utf-8",
    )
    other = tmp_path / "different_corrected_artifact.md"
    other.write_text("Different content.\n", encoding="utf-8")
    review = _review(
        inventory,
        correction_status="completed",
        corrected_artifact=corrected,
    )
    review["correction"]["corrected_artifact"] = str(other.resolve())
    review_path = output_dir / "advisory_validation_review_draft.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    _, audit = validator.package_validation(
        paths["deliverable_inventory"],
        review_path,
        contract_path,
        output_dir,
        corrected_deliverable=corrected,
    )

    assert audit["record_complete"] is False
    assert audit["checks"]["corrected_artifact_path_bound"] is False
    assert (
        "corrected deliverable does not match the path bound in the review"
        in audit["errors"]
    )


def test_completed_correction_is_bound_to_the_declared_hash(tmp_path: Path) -> None:
    validator, deliverable, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path
    )
    corrected = tmp_path / "advisory_memo_corrected.md"
    corrected.write_text(
        deliverable.read_text(encoding="utf-8") + "\nCondition added.\n",
        encoding="utf-8",
    )
    review = _review(
        inventory,
        correction_status="completed",
        corrected_artifact=corrected,
    )
    review["correction"]["corrected_artifact_sha256"] = "0" * 64
    review_path = output_dir / "advisory_validation_review_draft.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    _, audit = validator.package_validation(
        paths["deliverable_inventory"],
        review_path,
        contract_path,
        output_dir,
        corrected_deliverable=corrected,
    )

    assert audit["record_complete"] is False
    assert audit["checks"]["corrected_artifact_hash_bound"] is False
    assert (
        "corrected deliverable does not match the SHA-256 bound in the review"
        in audit["errors"]
    )


def test_package_rejects_an_output_path_that_aliases_the_original(
    tmp_path: Path,
) -> None:
    validator = _validator_module()
    output_dir = tmp_path / "validation"
    output_dir.mkdir()
    deliverable = output_dir / "advisory_validation_package.md"
    original_bytes = b"# Original advisory memo\n"
    deliverable.write_bytes(original_bytes)
    contract_path = tmp_path / "advisory_contract.json"
    contract_path.write_text(json.dumps(_contract()), encoding="utf-8")
    paths = validator.prepare_validation(deliverable, contract_path, output_dir)
    inventory = json.loads(paths["deliverable_inventory"].read_text(encoding="utf-8"))
    review = _review(inventory)
    review_path = output_dir / "advisory_validation_review_draft.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    with pytest.raises(
        validator.AdvisoryValidationError,
        match="refusing to overwrite original deliverable",
    ):
        validator.package_validation(
            paths["deliverable_inventory"], review_path, contract_path, output_dir
        )

    assert deliverable.read_bytes() == original_bytes


def test_semantic_evaluation_fixtures_cover_required_cases() -> None:
    payload = json.loads(
        (SKILL_ROOT / "evals" / "semantic_validation_cases.json").read_text(
            encoding="utf-8"
        )
    )
    case_ids = {case["id"] for case in payload["cases"]}

    assert {
        "supported-recommendation",
        "partially-supported-market-claim",
        "unsupported-claim",
        "contradicted-recommendation-premise",
        "reasoning-gap",
        "calculation-provenance-gap",
        "contract-failure",
        "judgment-dependent-choice",
        "external-document-without-contract",
        "unsupported-primary-format",
        "web-listings-do-not-prove-total-stock",
        "interview-quote-versus-underlying-truth",
        "derived-claim-requires-a-and-b",
        "calculation-recheck-after-input-change",
        "corrected-claim-supersedes-original",
        "two-hundred-page-coverage",
        "zero-claim-ready-review",
        "fake-completed-recheck",
        "corrected-artifact-without-second-review",
        "generic-passed-format-json",
    } <= case_ids
    assert "deterministic keyword or scorecard logic" in payload["purpose"]


def test_skill_and_public_page_keep_format_checks_and_model_data_explicit() -> None:
    skill = (SKILL_ROOT / "SKILL.md").read_text(encoding="utf-8")
    public_page = (
        ROOT
        / "static"
        / "shared"
        / "clara-advisory-deliverable-validator"
        / "index.html"
    ).read_text(encoding="utf-8")
    public_copy = (ROOT / "static" / "shared" / "product-function-pages.js").read_text(
        encoding="utf-8"
    )
    workflow_copy = public_copy.split('"clara-advisory-deliverable-validator":', 1)[
        1
    ].split('"clara-documents":', 1)[0]

    assert 'data-function-page="clara-advisory-deliverable-validator"' in public_page
    assert workflow_copy.count('modelDataStatus: "relevant"') == 5
    assert "Validate an advisory deliverable" in workflow_copy
    assert "No automatic anonymization is applied" in workflow_copy
    assert "mechanical closure remains partial" in workflow_copy
    assert "does not automatically open URLs" in workflow_copy
    assert "advisory_evidence_register.json" in workflow_copy
    assert "advisory_claim_register.json" in workflow_copy
    assert "walks every declared dependency" in workflow_copy
    assert "matched support" in workflow_copy
    assert "extraction omits script, style, and template blocks" in workflow_copy
    for workflow in (
        "clara:claim-basis-map",
        "clara:html-deck",
        "clara:reporting-engine",
        "clara:deck-correction",
    ):
        assert workflow in skill
    assert "keyword classifier" in skill
    assert "semantic scorecard" in skill
    assert "scripts make no model API calls" in skill


def test_generation_time_lineage_walks_all_dependencies_before_ready(
    tmp_path: Path,
) -> None:
    validator = _validator_module()
    case_dir = tmp_path / "case"
    lineage = validator._lineage_module()
    lineage.initialize_lineage(case_dir)
    evidence_a = _lineage_receipt("ev-a", "Management stated premise A.")
    evidence_b = _lineage_receipt("ev-b", "Management stated premise B.")
    claim_a = _lineage_claim("cl-a", "Premise A", evidence_ids=["ev-a"])
    claim_b = _lineage_claim("cl-b", "Premise B", evidence_ids=["ev-b"])
    claim_x = _lineage_claim(
        "cl-x",
        "Proceed with the pilot.",
        dependency_ids=["cl-a", "cl-b"],
        claim_type="recommendation",
    )
    lineage.record_evidence(case_dir, [evidence_a, evidence_b])
    lineage.record_claims(case_dir, [claim_a, claim_b, claim_x])
    deliverable = tmp_path / "memo.md"
    deliverable.write_text(
        "# Recommendation\n\nProceed with the pilot.", encoding="utf-8"
    )
    lineage.add_claim_appearances(
        case_dir,
        [
            {
                "claim_id": "cl-x",
                "appearance": {
                    "artifact": str(deliverable.resolve()),
                    "path_reference": "absolute",
                    "artifact_sha256": hashlib.sha256(
                        deliverable.read_bytes()
                    ).hexdigest(),
                    "artifact_byte_count": deliverable.stat().st_size,
                    "locator": "Recommendation",
                    "recorded_at": "2026-08-18T08:30:00+00:00",
                },
            }
        ],
    )
    contract_path = tmp_path / "advisory_contract.json"
    contract_path.write_text(json.dumps(_contract()), encoding="utf-8")
    output_dir = tmp_path / "validation"
    paths = validator.prepare_validation(
        deliverable,
        contract_path,
        output_dir,
        evidence_register=case_dir / "advisory_evidence_register.json",
        claim_register=case_dir / "advisory_claim_register.json",
    )
    inventory = json.loads(paths["deliverable_inventory"].read_text())
    review = _review(inventory)
    review["lineage_review"] = {
        "provenance_mode": "generation_time",
        "selection_method": "model_led_claim_chain_review",
        "reviewed_claim_ids": ["cl-x"],
        "chain_assessments": [
            _chain_assessment(claim_a),
            _chain_assessment(claim_b),
            _chain_assessment(claim_x),
        ],
        "untracked_material_claims": [],
        "limitations": [],
        "analysis": "The recommendation was walked back through both required premises.",
    }
    review["coverage_review"]["unit_assessments"][0]["material_claim_ids"] = ["cl-x"]
    review["coverage_review"]["unit_assessments"][0]["untracked_claim_ids"] = []
    review_path = tmp_path / "review.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    package_paths, audit = validator.package_validation(
        paths["deliverable_inventory"],
        review_path,
        paths["advisory_contract"],
        output_dir,
    )

    assert audit["record_complete"] is True
    assert audit["effective_delivery_readiness"] == "ready"
    assert package_paths["recheck_tasks"].is_file()


def test_generation_time_lineage_rejects_an_omitted_dependency(tmp_path: Path) -> None:
    validator = _validator_module()
    claim_a = _lineage_claim("cl-a", "Premise A")
    claim_x = _lineage_claim(
        "cl-x", "Conclusion X", dependency_ids=["cl-a"], claim_type="conclusion"
    )
    review = {
        "provenance_mode": "generation_time",
        "selection_method": "model_led_claim_chain_review",
        "reviewed_claim_ids": ["cl-x"],
        "chain_assessments": [_chain_assessment(claim_x)],
        "untracked_material_claims": [],
        "limitations": [],
        "analysis": "Only the conclusion was reviewed.",
    }
    claim_register = {"schema_version": "1.0", "claims": [claim_a, claim_x]}

    errors = validator._validate_lineage_review(
        review,
        provenance_mode="generation_time",
        claim_register=claim_register,
        evidence_register={"schema_version": "1.0", "evidence": []},
    )

    assert "lineage_review omitted dependency chain claims: cl-a" in errors


def test_generation_time_lineage_rejects_an_omitted_evidence_receipt() -> None:
    validator = _validator_module()
    evidence_a = _lineage_receipt("ev-a", "Management stated premise A.")
    evidence_b = _lineage_receipt("ev-b", "Management stated premise B.")
    evidence_b["rechecks_evidence_id"] = "ev-a"
    claim = _lineage_claim(
        "cl-a",
        "Premise A",
        evidence_ids=["ev-b"],
    )
    assessment = _chain_assessment(claim)
    assessment["evidence_ids"] = ["ev-b"]
    review = {
        "provenance_mode": "generation_time",
        "selection_method": "model_led_claim_chain_review",
        "reviewed_claim_ids": ["cl-a"],
        "chain_assessments": [assessment],
        "untracked_material_claims": [],
        "limitations": [],
        "analysis": "The claim was reviewed without its prior evidence receipt.",
    }

    errors = validator._validate_lineage_review(
        review,
        provenance_mode="generation_time",
        claim_register={"schema_version": "1.0", "claims": [claim]},
        evidence_register={
            "schema_version": "1.0",
            "evidence": [evidence_a, evidence_b],
        },
    )

    assert (
        "lineage_review.chain_assessments[0].evidence_ids must match lineage" in errors
    )


def test_ready_review_cannot_pass_without_any_material_claim_review(
    tmp_path: Path,
) -> None:
    validator, _, contract_path, _, paths, inventory = _prepare(tmp_path)
    review = _review(inventory)
    review["lineage_review"]["untracked_material_claims"] = []
    for assessment in review["coverage_review"]["unit_assessments"]:
        assessment["status"] = "reviewed_no_material_claims"
        assessment["material_claim_ids"] = []
        assessment["untracked_claim_ids"] = []
        assessment["analysis"] = "The model declared no material claim in this unit."
    coverage = json.loads(paths["coverage_inventory"].read_text(encoding="utf-8"))

    errors = validator.validate_review_record(
        review,
        json.loads(contract_path.read_text(encoding="utf-8")),
        provenance_mode="matched_support",
        coverage_inventory=coverage,
        matched_support_source_ids={"source-0001"},
    )

    assert (
        "delivery-ready status requires at least one model-reviewed material claim"
        in errors
    )


def test_matched_support_cannot_disguise_a_reconstructed_claim_as_lineage(
    tmp_path: Path,
) -> None:
    validator, _, contract_path, _, paths, inventory = _prepare(tmp_path)
    review = _review(inventory)
    reconstructed = review["lineage_review"]["untracked_material_claims"].pop()
    reconstructed["claim_id"] = reconstructed.pop("id")
    review["lineage_review"]["chain_assessments"] = [reconstructed]
    review["coverage_review"]["unit_assessments"][0]["untracked_claim_ids"] = []
    coverage = json.loads(paths["coverage_inventory"].read_text(encoding="utf-8"))

    errors = validator.validate_review_record(
        review,
        json.loads(contract_path.read_text(encoding="utf-8")),
        provenance_mode="matched_support",
        coverage_inventory=coverage,
        matched_support_source_ids={"source-0001"},
    )

    assert (
        "matched-support review must put reconstructed claims in untracked_material_claims"
        in errors
    )


def test_generation_time_review_rejects_a_false_deliverable_locator() -> None:
    validator = _validator_module()
    receipt = _lineage_receipt("ev-a", "Management stated premise A.")
    claim = _lineage_claim("cl-a", "Premise A", evidence_ids=["ev-a"])
    claim["appearances"] = [
        {
            "artifact": "/tmp/memo.md",
            "path_reference": "absolute",
            "artifact_sha256": "a" * 64,
            "artifact_byte_count": 10,
            "locator": "Section 4",
            "recorded_at": "2026-08-18T08:10:00+00:00",
        }
    ]
    assessment = _chain_assessment(claim)
    assessment["deliverable_locations"] = ["Section 99"]
    review = {
        "provenance_mode": "generation_time",
        "selection_method": "model_led_claim_chain_review",
        "reviewed_claim_ids": ["cl-a"],
        "chain_assessments": [assessment],
        "untracked_material_claims": [],
        "limitations": [],
        "analysis": "The selected claim was reviewed.",
    }

    errors = validator._validate_lineage_review(
        review,
        provenance_mode="generation_time",
        claim_register={"schema_version": "1.0", "claims": [claim]},
        evidence_register={"schema_version": "1.0", "evidence": [receipt]},
        deliverable_sha256="a" * 64,
    )

    assert any("must match hash-bound claim appearances" in error for error in errors)


def test_adequate_support_requires_resolution_of_contradicting_lineage() -> None:
    validator = _validator_module()
    receipt = _lineage_receipt("ev-a", "The source contradicts premise A.")
    claim = _lineage_claim("cl-a", "Premise A", evidence_ids=["ev-a"])
    claim["evidence_links"][0]["relationship"] = "contradicts"
    assessment = _chain_assessment(claim)
    assessment["contradiction_resolution"] = ""
    review = {
        "provenance_mode": "generation_time",
        "selection_method": "model_led_claim_chain_review",
        "reviewed_claim_ids": [],
        "chain_assessments": [assessment],
        "untracked_material_claims": [],
        "limitations": [],
        "analysis": "The contradiction was reviewed.",
    }

    errors = validator._validate_lineage_review(
        review,
        provenance_mode="generation_time",
        claim_register={"schema_version": "1.0", "claims": [claim]},
        evidence_register={"schema_version": "1.0", "evidence": [receipt]},
    )

    assert any(
        "adequate support requires contradiction_resolution" in error
        for error in errors
    )


def test_completed_recheck_requires_a_real_successor_receipt() -> None:
    validator = _validator_module()
    prior = _lineage_receipt("ev-old", "The prior web observation.")
    fake = _lineage_receipt("ev-new", "The claimed recheck result.")
    fake["evidence_type"] = "web_capture"
    fake["rechecks_evidence_id"] = "ev-old"
    claim = _lineage_claim("cl-a", "Premise A", evidence_ids=["ev-new"])
    assessment = _chain_assessment(claim)
    assessment["evidence_ids"] = ["ev-old", "ev-new"]
    assessment["recheck"] = {
        "required": True,
        "kind": "web",
        "status": "completed",
        "evidence_ids": ["ev-new"],
        "analysis": "The model declared the targeted web recheck complete.",
    }
    review = {
        "provenance_mode": "generation_time",
        "selection_method": "model_led_claim_chain_review",
        "reviewed_claim_ids": [],
        "chain_assessments": [assessment],
        "untracked_material_claims": [],
        "limitations": [],
        "analysis": "The recheck was reviewed.",
    }

    errors = validator._validate_lineage_review(
        review,
        provenance_mode="generation_time",
        claim_register={"schema_version": "1.0", "claims": [claim]},
        evidence_register={"schema_version": "1.0", "evidence": [prior, fake]},
    )

    assert any("lacks completed recheck verification" in error for error in errors)


def test_pending_model_selected_recheck_blocks_ready_and_is_packaged(
    tmp_path: Path,
) -> None:
    validator, _deliverable, contract_path, output_dir, paths, inventory = _prepare(
        tmp_path
    )
    review = _review(inventory)
    item = review["lineage_review"]["untracked_material_claims"][0]
    item["support_status"] = "uncertain"
    item["recheck"] = {
        "required": True,
        "kind": "web",
        "status": "pending",
        "evidence_ids": [],
        "analysis": "Reopen the captured public page because the observation may be stale.",
    }
    item["resolution"] = {
        "status": "pending",
        "explanation": "The claim remains unresolved until the targeted recheck.",
    }
    review_path = tmp_path / "review.json"
    review_path.write_text(json.dumps(review), encoding="utf-8")

    package_paths, audit = validator.package_validation(
        paths["deliverable_inventory"],
        review_path,
        contract_path,
        output_dir,
    )

    tasks = json.loads(package_paths["recheck_tasks"].read_text())
    assert audit["record_complete"] is False
    assert audit["effective_delivery_readiness"] == "blocked"
    assert tasks["tasks"][0]["kind"] == "web"
    assert tasks["tasks"][0]["status"] == "pending"


def test_two_hundred_page_scale_uses_bounded_coverage_units(tmp_path: Path) -> None:
    validator = _validator_module()
    deliverable = tmp_path / "long_report.md"
    deliverable.write_text(
        "\n\n".join(
            f"# Page {page:03d}\n\n" + (f"Reviewed content for page {page}. " * 45)
            for page in range(1, 201)
        ),
        encoding="utf-8",
    )
    contract_path = tmp_path / "advisory_contract.json"
    contract_path.write_text(json.dumps(_contract()), encoding="utf-8")
    output_dir = tmp_path / "validation"

    paths = validator.prepare_validation(deliverable, contract_path, output_dir)
    inventory = json.loads(paths["deliverable_inventory"].read_text())
    coverage = json.loads(paths["coverage_inventory"].read_text())
    review = _review(inventory)

    assert len(coverage["units"]) >= 20
    assert review["coverage_review"]["considered_unit_ids"] == [
        unit["id"] for unit in coverage["units"]
    ]
    review["coverage_review"]["considered_unit_ids"].pop()
    errors = validator.validate_review_record(
        review,
        _contract(),
        provenance_mode="matched_support",
        coverage_inventory=coverage,
    )
    assert any("does not account for units" in error for error in errors)
