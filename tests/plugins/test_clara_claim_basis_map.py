from __future__ import annotations

import importlib.util
from pathlib import Path
from types import ModuleType

import pytest

ROOT = Path(__file__).resolve().parents[2]
SCRIPT_PATH = (
    ROOT
    / "plugins"
    / "clara"
    / "skills"
    / "claim-basis-map"
    / "scripts"
    / "render_claim_basis_map.py"
)


def _load_renderer() -> ModuleType:
    spec = importlib.util.spec_from_file_location(
        "clara_claim_basis_map_renderer", SCRIPT_PATH
    )
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def test_render_claim_basis_map_surfaces_claim_link_to_prior_slide() -> None:
    renderer = _load_renderer()
    payload = {
        "deck": "deck.pptx",
        "slides": [
            {
                "slide_number": 4,
                "slide_title": "Market direction",
                "claims": [
                    {
                        "claim": "Premiumization remains the main growth vector.",
                        "reasoning_inputs": [
                            {"label": "Category report", "locator": "p. 14"}
                        ],
                        "reasoning": "Premium formats outgrow the category.",
                    }
                ],
            },
            {
                "slide_number": 6,
                "slide_title": "Retail implications",
                "claims": [
                    {
                        "claim": "Retailers will prioritize premium discovery space.",
                        "claim_refs": [
                            {
                                "slide_number": 4,
                                "claim": "Premiumization remains the main growth vector.",
                            }
                        ],
                        "reasoning": "Discovery space follows the main growth vector.",
                    }
                ],
            },
        ],
    }

    markdown = renderer.render_claim_basis_map(payload)

    assert "No ungrounded or unresolved claim dependencies captured." in markdown
    assert "### Claim-Linked" in markdown
    assert (
        'Based on claim: Slide 4 - "Premiumization remains the main growth vector."'
        in markdown
    )


def test_render_claim_basis_map_surfaces_unresolved_upstream_claim() -> None:
    renderer = _load_renderer()
    payload = {
        "slides": [
            {
                "slide_number": 3,
                "claims": [
                    {
                        "claim": "Retailers will prioritize premium discovery space.",
                        "claim_refs": [
                            {
                                "slide_number": 1,
                                "claim": "Premiumization remains the main growth vector.",
                            }
                        ],
                    }
                ],
            }
        ]
    }

    markdown = renderer.render_claim_basis_map(payload)

    assert (
        "Basis: unresolved claim reference: Slide 1 - "
        '"Premiumization remains the main growth vector."' in markdown
    )


def test_render_claim_basis_map_rejects_future_claim_reference() -> None:
    renderer = _load_renderer()
    payload = {
        "slides": [
            {
                "slide_number": 1,
                "claims": [
                    {
                        "claim": "Retailers will prioritize premium discovery space.",
                        "claim_refs": [
                            {
                                "slide_number": 2,
                                "claim": "Premiumization remains the main growth vector.",
                            }
                        ],
                    }
                ],
            },
            {
                "slide_number": 2,
                "claims": [{"claim": "Premiumization remains the main growth vector."}],
            },
        ]
    }

    markdown = renderer.render_claim_basis_map(payload)

    assert (
        "Basis: claim reference must point to an earlier slide: target slide 2, "
        "current slide 1" in markdown
    )


def test_render_claim_basis_map_rejects_duplicate_claim_key() -> None:
    renderer = _load_renderer()
    duplicate_key = "claim-growth-vector"
    payload = {
        "slides": [
            {
                "slide_number": 1,
                "claims": [
                    {"claim": "First claim.", "claim_key": duplicate_key},
                    {"claim": "Second claim.", "claim_key": duplicate_key},
                ],
            }
        ]
    }

    with pytest.raises(ValueError, match="Duplicate claim_key"):
        renderer.render_claim_basis_map(payload)


def test_render_claim_basis_map_surfaces_edited_keyed_claim() -> None:
    renderer = _load_renderer()
    payload = {
        "deck": "deck.pptx",
        "slides": [
            {
                "slide_number": 1,
                "claims": [
                    {
                        "claim": "The market grew 12% in 2025.",
                        "claim_key": "s01-c01",
                        "source_refs": [{"title": "Market report"}],
                    }
                ],
            }
        ],
    }
    current_snapshot = {
        "deck": "deck-edited.pptx",
        "slides": [
            {
                "slide_number": 1,
                "texts": [
                    {
                        "claim_key": "s01-c01",
                        "text": "The market grew 18% in 2025.",
                    }
                ],
            }
        ],
    }

    markdown = renderer.render_claim_basis_map(
        payload, current_snapshot=current_snapshot
    )

    assert "## Current Deck Check" in markdown
    assert "Status: edited" in markdown
    assert 'Current text: "The market grew 18% in 2025."' in markdown


def test_render_claim_basis_map_surfaces_moved_claim() -> None:
    renderer = _load_renderer()
    payload = {
        "slides": [
            {
                "slide_number": 1,
                "claims": [
                    {
                        "claim": "Premium SKUs contributed 62% of growth.",
                        "source_refs": [{"title": "Market report"}],
                    }
                ],
            },
            {"slide_number": 2, "claims": []},
        ],
    }
    current_snapshot = {
        "slides": [
            {"slide_number": 1, "texts": []},
            {
                "slide_number": 2,
                "texts": ["Premium SKUs contributed 62% of growth."],
            },
        ]
    }

    markdown = renderer.render_claim_basis_map(
        payload, current_snapshot=current_snapshot
    )

    assert "### Moved Claims" in markdown
    assert "Current slide: 2" in markdown


def test_render_claim_basis_map_surfaces_untracked_current_text() -> None:
    renderer = _load_renderer()
    payload = {
        "slides": [
            {
                "slide_number": 1,
                "slide_title": "Market direction",
                "claims": [
                    {
                        "claim": "The market grew 12% in 2025.",
                        "source_refs": [{"title": "Market report"}],
                    }
                ],
            }
        ],
    }
    current_snapshot = {
        "slides": [
            {
                "slide_number": 1,
                "texts": [
                    "Market direction",
                    "The market grew 12% in 2025.",
                    "A newly added claim has enough words to be reported.",
                ],
            }
        ]
    }

    markdown = renderer.render_claim_basis_map(
        payload, current_snapshot=current_snapshot
    )

    assert "Status: untracked-current-text" in markdown
    assert "A newly added claim has enough words to be reported." in markdown


def test_render_claim_basis_map_surfaces_broken_reference_after_upstream_edit() -> None:
    renderer = _load_renderer()
    payload = {
        "slides": [
            {
                "slide_number": 1,
                "claims": [
                    {
                        "claim": "Premiumization remains the main growth vector.",
                        "claim_key": "s01-c01",
                        "source_refs": [{"title": "Market report"}],
                    }
                ],
            },
            {
                "slide_number": 2,
                "claims": [
                    {
                        "claim": "Retailers will prioritize premium discovery space.",
                        "claim_refs": [
                            {
                                "slide_number": 1,
                                "claim": "Premiumization remains the main growth vector.",
                            }
                        ],
                    }
                ],
            },
        ],
    }
    current_snapshot = {
        "slides": [
            {
                "slide_number": 1,
                "texts": [
                    {
                        "claim_key": "s01-c01",
                        "text": "Value brands are now the main growth vector.",
                    }
                ],
            },
            {
                "slide_number": 2,
                "texts": ["Retailers will prioritize premium discovery space."],
            },
        ]
    }

    markdown = renderer.render_claim_basis_map(
        payload, current_snapshot=current_snapshot
    )

    assert "Status: reference-broken" in markdown
    assert "depends on claim now marked edited" in markdown


def test_extract_current_deck_snapshot_reads_hidden_claim_shape_name(
    tmp_path: Path,
) -> None:
    renderer = _load_renderer()
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Inches

    deck = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    shape.name = "clara-claim:s01-c01"
    shape.text = "The market grew 12% in 2025."
    presentation.save(deck)

    snapshot = renderer.extract_current_deck_snapshot(deck)

    assert snapshot["slides"][0]["texts"] == [
        {"text": "The market grew 12% in 2025.", "claim_key": "s01-c01"}
    ]
