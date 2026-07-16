from __future__ import annotations

import io
import json
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Mapping, Sequence
from urllib.parse import urlparse

from PIL import Image

from src.review_brief.slides_deck import ReviewBriefDeckSlide, ReviewBriefDeckSpec
from src.slides.notebooklm_style import load_notebooklm_style, resolve_prompt_style_key

__all__ = [
    "REVIEW_BRIEF_PPTX_SPEC_FILENAME",
    "ReviewBriefPptxSlide",
    "ReviewBriefPptxSpec",
    "build_review_brief_pptx_spec",
    "load_review_brief_pptx_spec",
    "render_review_brief_pptx_from_template",
    "review_brief_pptx_spec_path",
    "write_review_brief_pptx_spec",
]

REVIEW_BRIEF_PPTX_SPEC_FILENAME = "review_brief_pptx_spec.json"
_SLIDE_WIDTH_IN = 13.333333
_SLIDE_HEIGHT_IN = 7.5
_TITLE_TOP_IN = 0.62
_TITLE_LEFT_IN = 0.78
_TITLE_WIDTH_IN = 11.7
_TITLE_HEIGHT_IN = 0.7
_EYEBROW_TOP_IN = 0.34
_DIVIDER_TOP_IN = 1.22
_CONTENT_TOP_IN = 1.52
_CONTENT_BOTTOM_IN = 6.78
_UNIFORM_META_TOP_IN = 1.34
_UNIFORM_SUMMARY_BODY_TOP_IN = 1.82
_UNIFORM_SUMMARY_BULLETS_TOP_IN = 2.52
_UNIFORM_CHART_COPY_LEFT_IN = 0.94
_UNIFORM_CHART_COPY_WIDTH_IN = 3.15
_UNIFORM_CHART_LEFT_IN = 4.5
_UNIFORM_CHART_TOP_IN = 1.9
_UNIFORM_CHART_WIDTH_IN = 7.98
_UNIFORM_CHART_HEIGHT_IN = 4.78


@dataclass(frozen=True, slots=True)
class ReviewBriefPptxSlide:
    """Native slide content placed into the review-brief PPTX template."""

    kind: str
    title: str
    eyebrow: str = ""
    body: str = ""
    bullets: list[str] = field(default_factory=list)
    scope_items: list[str] = field(default_factory=list)
    chart_path: str = ""
    chart_caption: str = ""
    chart_id: str = ""


@dataclass(frozen=True, slots=True)
class ReviewBriefPptxSpec:
    """Serialized template spec stored alongside a generated review deck."""

    template_key: str
    prompt_style: str
    slides: list[ReviewBriefPptxSlide]


def review_brief_pptx_spec_path(deck_path: Path) -> Path:
    """Return the companion template-spec path for ``deck_path``."""

    return deck_path / REVIEW_BRIEF_PPTX_SPEC_FILENAME


def write_review_brief_pptx_spec(deck_path: Path, spec: ReviewBriefPptxSpec) -> Path:
    """Persist ``spec`` next to the generated deck assets."""

    payload = {
        "templateKey": spec.template_key,
        "promptStyle": spec.prompt_style,
        "slides": [asdict(slide) for slide in spec.slides],
    }
    output_path = review_brief_pptx_spec_path(deck_path)
    output_path.write_text(
        json.dumps(payload, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    return output_path


def load_review_brief_pptx_spec(deck_path: Path) -> ReviewBriefPptxSpec:
    """Load the persisted template spec for a generated review deck."""

    payload = json.loads(
        review_brief_pptx_spec_path(deck_path).read_text(encoding="utf-8")
    )
    if not isinstance(payload, dict):
        raise ValueError("Review brief PPTX spec must be a JSON object.")
    raw_slides = payload.get("slides")
    if not isinstance(raw_slides, list):
        raise ValueError("Review brief PPTX spec is missing slides.")
    slides = [
        ReviewBriefPptxSlide(
            kind=str(item.get("kind") or "").strip(),
            title=str(item.get("title") or "").strip(),
            eyebrow=str(item.get("eyebrow") or "").strip(),
            body=str(item.get("body") or "").strip(),
            bullets=_read_text_list(item.get("bullets")),
            scope_items=_read_text_list(item.get("scope_items")),
            chart_path=str(item.get("chart_path") or "").strip(),
            chart_caption=str(item.get("chart_caption") or "").strip(),
            chart_id=str(item.get("chart_id") or "").strip(),
        )
        for item in raw_slides
        if isinstance(item, Mapping)
    ]
    prompt_style = resolve_prompt_style_key(
        str(payload.get("promptStyle") or "").strip() or None
    )
    template_key = resolve_prompt_style_key(
        str(payload.get("templateKey") or "").strip() or prompt_style
    )
    return ReviewBriefPptxSpec(
        template_key=template_key,
        prompt_style=prompt_style,
        slides=slides,
    )


def build_review_brief_pptx_spec(
    brief_payload: Mapping[str, object],
    deck_spec: ReviewBriefDeckSpec,
    *,
    chart_image_urls: Mapping[str, str],
) -> ReviewBriefPptxSpec:
    """Build the native PPTX template spec from the review-brief JSON payload."""

    prompt_style = resolve_prompt_style_key(deck_spec.prompt_style)
    title_scope = _summary_scope_items(brief_payload)
    use_uniform_layout = prompt_style == "uniform"
    slides: list[ReviewBriefPptxSlide] = []
    chart_counter = 0
    for slide in deck_spec.slides:
        if slide.kind == "summary":
            slides.append(
                ReviewBriefPptxSlide(
                    kind="summary",
                    eyebrow="" if use_uniform_layout else "Executive summary",
                    title=slide.title,
                    body=slide.body,
                    bullets=list(slide.bullets),
                    scope_items=title_scope,
                )
            )
            continue
        chart_counter += 1
        chart_path = _relative_asset_path(
            chart_image_urls.get(str(slide.chart_id or "").strip())
        )
        slides.append(
            ReviewBriefPptxSlide(
                kind="chart",
                eyebrow="" if use_uniform_layout else f"Insight {chart_counter:02d}",
                title=slide.title,
                body=slide.body,
                bullets=list(slide.bullets),
                chart_path=chart_path,
                chart_caption=str(slide.chart_alt or "").strip(),
                chart_id=str(slide.chart_id or "").strip(),
            )
        )
    return ReviewBriefPptxSpec(
        template_key=prompt_style,
        prompt_style=prompt_style,
        slides=slides,
    )


def render_review_brief_pptx_from_template(deck_path: Path) -> io.BytesIO:
    """Render an editable PPTX from the persisted review-brief template spec."""

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    spec = load_review_brief_pptx_spec(deck_path)
    presentation = _load_template_presentation(Presentation, spec.template_key)
    presentation.slide_width = Inches(_SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(_SLIDE_HEIGHT_IN)
    blank_layout = presentation.slide_layouts[6]
    style = load_notebooklm_style(spec.prompt_style)
    palette = _template_palette(spec.prompt_style)
    use_uniform_layout = resolve_prompt_style_key(spec.prompt_style) == "uniform"
    font_name = style.font_family_primary or style.font_family_fallback
    title_size = float(style.title_size_pt)
    body_size = float(style.body_size_pt)
    line_height = float(style.line_height)

    for index, slide_spec in enumerate(spec.slides, start=1):
        slide = presentation.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(style.bg_color))
        if slide_spec.eyebrow:
            _add_textbox(
                slide,
                text=slide_spec.eyebrow.upper(),
                left=_TITLE_LEFT_IN,
                top=_EYEBROW_TOP_IN,
                width=3.2,
                height=0.24,
                font_name=font_name,
                font_size=10.5,
                color_rgb=palette["muted_text"],
                bold=True,
            )
        _add_textbox(
            slide,
            text=slide_spec.title,
            left=_TITLE_LEFT_IN,
            top=_TITLE_TOP_IN,
            width=_TITLE_WIDTH_IN,
            height=_TITLE_HEIGHT_IN,
            font_name=font_name,
            font_size=title_size,
            color_rgb=_hex_to_rgb(style.text_color),
            bold=True,
        )
        if use_uniform_layout:
            if slide_spec.kind == "summary":
                _render_uniform_summary_slide(
                    slide,
                    slide_spec,
                    font_name=font_name,
                    body_size=body_size,
                    line_height=line_height,
                    text_rgb=_hex_to_rgb(style.text_color),
                    palette=palette,
                )
            else:
                _render_uniform_chart_slide(
                    slide,
                    deck_path=deck_path,
                    slide_spec=slide_spec,
                    font_name=font_name,
                    body_size=body_size,
                    line_height=line_height,
                    text_rgb=_hex_to_rgb(style.text_color),
                    palette=palette,
                )
        else:
            divider = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(_TITLE_LEFT_IN),
                Inches(_DIVIDER_TOP_IN),
                Inches(11.8),
                Inches(0.03),
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = RGBColor(*palette["rule"])
            divider.line.fill.background()
            if slide_spec.kind == "summary":
                _render_summary_slide(
                    slide,
                    slide_spec,
                    font_name=font_name,
                    body_size=body_size,
                    line_height=line_height,
                    text_rgb=_hex_to_rgb(style.text_color),
                    palette=palette,
                )
            else:
                _render_chart_slide(
                    slide,
                    deck_path=deck_path,
                    slide_spec=slide_spec,
                    font_name=font_name,
                    body_size=body_size,
                    line_height=line_height,
                    text_rgb=_hex_to_rgb(style.text_color),
                    palette=palette,
                )
        _add_textbox(
            slide,
            text=str(index),
            left=12.2,
            top=7.06,
            width=0.45,
            height=0.2,
            font_name=font_name,
            font_size=10.5,
            color_rgb=palette["muted_text"],
            alignment=PP_ALIGN.RIGHT,
        )

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


def _render_uniform_summary_slide(
    slide,
    slide_spec: ReviewBriefPptxSlide,
    *,
    font_name: str,
    body_size: float,
    line_height: float,
    text_rgb: tuple[int, int, int],
    palette: Mapping[str, tuple[int, int, int]],
) -> None:
    scope_line = _join_scope_items(slide_spec.scope_items)
    if scope_line:
        _add_textbox(
            slide,
            text=scope_line,
            left=_TITLE_LEFT_IN,
            top=_UNIFORM_META_TOP_IN,
            width=11.55,
            height=0.34,
            font_name=font_name,
            font_size=11.0,
            color_rgb=palette["muted_text"],
        )
    if slide_spec.body:
        _add_textbox(
            slide,
            text=slide_spec.body,
            left=_TITLE_LEFT_IN,
            top=_UNIFORM_SUMMARY_BODY_TOP_IN,
            width=10.8,
            height=0.62,
            font_name=font_name,
            font_size=max(body_size + 1.5, 19.0),
            color_rgb=text_rgb,
        )
    _add_bullet_box(
        slide,
        bullets=slide_spec.bullets,
        left=_TITLE_LEFT_IN,
        top=_UNIFORM_SUMMARY_BULLETS_TOP_IN,
        width=10.2,
        height=3.55,
        font_name=font_name,
        body_size=body_size,
        line_height=line_height,
        text_rgb=text_rgb,
    )


def _render_uniform_chart_slide(
    slide,
    *,
    deck_path: Path,
    slide_spec: ReviewBriefPptxSlide,
    font_name: str,
    body_size: float,
    line_height: float,
    text_rgb: tuple[int, int, int],
    palette: Mapping[str, tuple[int, int, int]],
) -> None:
    from pptx.enum.text import PP_ALIGN

    if slide_spec.body:
        _add_textbox(
            slide,
            text=slide_spec.body,
            left=_UNIFORM_CHART_COPY_LEFT_IN,
            top=1.72,
            width=_UNIFORM_CHART_COPY_WIDTH_IN,
            height=0.74,
            font_name=font_name,
            font_size=max(body_size, 18.0),
            color_rgb=text_rgb,
            bold=True,
        )
    _add_bullet_box(
        slide,
        bullets=slide_spec.bullets,
        left=_UNIFORM_CHART_COPY_LEFT_IN,
        top=2.52,
        width=_UNIFORM_CHART_COPY_WIDTH_IN,
        height=3.42,
        font_name=font_name,
        body_size=max(body_size - 0.5, 15.5),
        line_height=line_height,
        text_rgb=text_rgb,
    )
    chart_path = deck_path / slide_spec.chart_path if slide_spec.chart_path else None
    if chart_path is not None and chart_path.exists():
        left, top, width, height = _fit_image_within_box(
            chart_path,
            left=_UNIFORM_CHART_LEFT_IN,
            top=_UNIFORM_CHART_TOP_IN,
            width=_UNIFORM_CHART_WIDTH_IN,
            height=_UNIFORM_CHART_HEIGHT_IN,
        )
        slide.shapes.add_picture(
            str(chart_path),
            _inches(left),
            _inches(top),
            width=_inches(width),
            height=_inches(height),
        )
    else:
        _add_textbox(
            slide,
            text="Chart preview unavailable",
            left=_UNIFORM_CHART_LEFT_IN + 2.4,
            top=_UNIFORM_CHART_TOP_IN + 1.95,
            width=3.2,
            height=0.25,
            font_name=font_name,
            font_size=12.0,
            color_rgb=palette["muted_text"],
            alignment=PP_ALIGN.CENTER,
        )


def _render_summary_slide(
    slide,
    slide_spec: ReviewBriefPptxSlide,
    *,
    font_name: str,
    body_size: float,
    line_height: float,
    text_rgb: tuple[int, int, int],
    palette: Mapping[str, tuple[int, int, int]],
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    lead_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        _inches(0.78),
        _inches(_CONTENT_TOP_IN),
        _inches(7.15),
        _inches(1.55),
    )
    lead_box.fill.solid()
    lead_box.fill.fore_color.rgb = RGBColor(*palette["panel_fill"])
    lead_box.line.color.rgb = RGBColor(*palette["panel_border"])
    if slide_spec.body:
        _add_textbox(
            slide,
            text=slide_spec.body,
            left=1.02,
            top=1.78,
            width=6.65,
            height=1.0,
            font_name=font_name,
            font_size=max(body_size + 1.5, 19.0),
            color_rgb=text_rgb,
            bold=False,
        )
    _add_bullet_box(
        slide,
        bullets=slide_spec.bullets,
        left=0.82,
        top=3.22,
        width=7.05,
        height=3.18,
        font_name=font_name,
        body_size=body_size,
        line_height=line_height,
        text_rgb=text_rgb,
    )
    scope_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        _inches(8.34),
        _inches(_CONTENT_TOP_IN),
        _inches(4.16),
        _inches(4.95),
    )
    scope_box.fill.solid()
    scope_box.fill.fore_color.rgb = RGBColor(*palette["scope_fill"])
    scope_box.line.color.rgb = RGBColor(*palette["panel_border"])
    _add_textbox(
        slide,
        text="Scope",
        left=8.62,
        top=1.78,
        width=1.2,
        height=0.22,
        font_name=font_name,
        font_size=11.0,
        color_rgb=palette["muted_text"],
        bold=True,
    )
    _add_scope_box(
        slide,
        items=slide_spec.scope_items,
        left=8.64,
        top=2.18,
        width=3.68,
        height=3.78,
        font_name=font_name,
        body_size=max(body_size - 0.5, 15.0),
        text_rgb=text_rgb,
        line_height=line_height,
    )


def _render_chart_slide(
    slide,
    *,
    deck_path: Path,
    slide_spec: ReviewBriefPptxSlide,
    font_name: str,
    body_size: float,
    line_height: float,
    text_rgb: tuple[int, int, int],
    palette: Mapping[str, tuple[int, int, int]],
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN

    narrative_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        _inches(0.78),
        _inches(_CONTENT_TOP_IN),
        _inches(4.15),
        _inches(4.98),
    )
    narrative_box.fill.solid()
    narrative_box.fill.fore_color.rgb = RGBColor(*palette["panel_fill"])
    narrative_box.line.color.rgb = RGBColor(*palette["panel_border"])
    if slide_spec.body:
        _add_textbox(
            slide,
            text=slide_spec.body,
            left=1.02,
            top=1.82,
            width=3.65,
            height=0.72,
            font_name=font_name,
            font_size=max(body_size, 18.0),
            color_rgb=text_rgb,
            bold=True,
        )
    _add_bullet_box(
        slide,
        bullets=slide_spec.bullets,
        left=1.02,
        top=2.56,
        width=3.58,
        height=3.48,
        font_name=font_name,
        body_size=max(body_size - 0.5, 15.5),
        line_height=line_height,
        text_rgb=text_rgb,
    )
    chart_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        _inches(5.15),
        _inches(_CONTENT_TOP_IN),
        _inches(7.44),
        _inches(4.98),
    )
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = RGBColor(*palette["chart_fill"])
    chart_box.line.color.rgb = RGBColor(*palette["panel_border"])
    if slide_spec.chart_caption:
        _add_textbox(
            slide,
            text=slide_spec.chart_caption,
            left=5.45,
            top=1.8,
            width=6.85,
            height=0.22,
            font_name=font_name,
            font_size=10.5,
            color_rgb=palette["muted_text"],
            bold=False,
        )
    chart_path = deck_path / slide_spec.chart_path if slide_spec.chart_path else None
    if chart_path is not None and chart_path.exists():
        left, top, width, height = _fit_image_within_box(
            chart_path,
            left=5.38,
            top=2.08,
            width=6.98,
            height=4.1,
        )
        slide.shapes.add_picture(
            str(chart_path),
            _inches(left),
            _inches(top),
            width=_inches(width),
            height=_inches(height),
        )
    else:
        _add_textbox(
            slide,
            text="Chart preview unavailable",
            left=7.2,
            top=3.78,
            width=3.0,
            height=0.25,
            font_name=font_name,
            font_size=12.0,
            color_rgb=palette["muted_text"],
            alignment=PP_ALIGN.CENTER,
        )


def _add_scope_box(
    slide,
    *,
    items: Sequence[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    body_size: float,
    text_rgb: tuple[int, int, int],
    line_height: float,
) -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    textbox = slide.shapes.add_textbox(
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(8)
        paragraph.line_spacing = line_height
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(body_size)
            run.font.color.rgb = _rgb(text_rgb)


def _add_bullet_box(
    slide,
    *,
    bullets: Sequence[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    body_size: float,
    line_height: float,
    text_rgb: tuple[int, int, int],
) -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    textbox = slide.shapes.add_textbox(
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(10)
        paragraph.line_spacing = line_height
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(body_size)
            run.font.color.rgb = _rgb(text_rgb)


def _add_textbox(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    font_size: float,
    color_rgb: tuple[int, int, int],
    bold: bool = False,
    alignment=None,
) -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    textbox = slide.shapes.add_textbox(
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT if alignment is None else alignment
    paragraph.text = text
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color_rgb)


def _template_palette(prompt_style: str) -> dict[str, tuple[int, int, int]]:
    if resolve_prompt_style_key(prompt_style) == "uniform":
        return {
            "panel_fill": _hex_to_rgb("#FFFFFF"),
            "scope_fill": _hex_to_rgb("#FFFFFF"),
            "chart_fill": _hex_to_rgb("#FFFFFF"),
            "panel_border": _hex_to_rgb("#FFFFFF"),
            "muted_text": _hex_to_rgb("#000000"),
            "rule": _hex_to_rgb("#000000"),
        }
    return {
        "panel_fill": _hex_to_rgb("#F5F7FA"),
        "scope_fill": _hex_to_rgb("#F7F8FB"),
        "chart_fill": _hex_to_rgb("#FBFCFD"),
        "panel_border": _hex_to_rgb("#D6DCE3"),
        "muted_text": _hex_to_rgb("#5F6B7A"),
        "rule": _hex_to_rgb("#D9DEE5"),
    }


def _load_template_presentation(presentation_cls, template_key: str):
    template_path = (
        Path(__file__).resolve().parent / "pptx_templates" / f"{template_key}.pptx"
    )
    if template_path.exists():
        return presentation_cls(str(template_path))
    return presentation_cls()


def _summary_scope_items(payload: Mapping[str, object]) -> list[str]:
    requested_scope = payload.get("requested_scope")
    scope = requested_scope if isinstance(requested_scope, Mapping) else {}
    items: list[str] = []
    retailers = _read_text_list(payload.get("retailers")) or _read_text_list(
        scope.get("retailers")
    )
    if retailers:
        items.append(f"Retailers: {', '.join(retailers)}")
    category = _read_text(payload.get("category")) or _read_text(
        scope.get("category_label")
    )
    if category:
        items.append(f"Category: {category}")
    start_month = _read_text(payload.get("start_month"))
    end_month = _read_text(payload.get("end_month"))
    if start_month and end_month:
        items.append(f"Period: {start_month} to {end_month}")
    return items


def _join_scope_items(items: Sequence[str]) -> str:
    return " | ".join(
        text for text in (str(item or "").strip() for item in items) if text
    )


def _relative_asset_path(image_url: str | None) -> str:
    if not image_url:
        return ""
    path = urlparse(str(image_url)).path
    marker = "/assets/"
    if marker not in path:
        return ""
    return f"assets/{path.split(marker, 1)[1].lstrip('/')}"


def _fit_image_within_box(
    image_path: Path,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
) -> tuple[float, float, float, float]:
    try:
        with Image.open(image_path) as image:
            img_width, img_height = image.size
    except (OSError, ValueError):
        return left, top, width, height
    if img_width <= 0 or img_height <= 0:
        return left, top, width, height
    box_ratio = width / height
    image_ratio = float(img_width) / float(img_height)
    if image_ratio >= box_ratio:
        scaled_width = width
        scaled_height = width / image_ratio
        offset_x = 0.0
        offset_y = (height - scaled_height) / 2.0
    else:
        scaled_height = height
        scaled_width = height * image_ratio
        offset_x = (width - scaled_width) / 2.0
        offset_y = 0.0
    return left + offset_x, top + offset_y, scaled_width, scaled_height


def _read_text_list(value: object) -> list[str]:
    if not isinstance(value, list):
        return []
    return [text for text in (str(item or "").strip() for item in value) if text]


def _read_text(value: object) -> str:
    return str(value or "").strip()


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    normalized = str(value or "").strip().lstrip("#")
    if len(normalized) != 6:
        return (0, 0, 0)
    return tuple(int(normalized[index : index + 2], 16) for index in (0, 2, 4))


def _rgb(value: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(*value)


def _inches(value: float):
    from pptx.util import Inches

    return Inches(value)
