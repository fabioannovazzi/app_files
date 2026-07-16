from __future__ import annotations

import json
import re
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

__all__ = [
    "DECK_PPTX_TEMPLATE_FILENAME",
    "DECK_PPTX_TEMPLATE_MANIFEST_FILENAME",
    "DeckPptxTemplateLayout",
    "DeckPptxTemplateManifest",
    "DeckPptxTemplatePrototype",
    "build_pptx_template_manifest",
    "deck_pptx_template_manifest_path",
    "deck_pptx_template_path",
    "ensure_deck_pptx_template_manifest",
    "load_deck_pptx_template_manifest",
    "write_deck_pptx_template_manifest",
]

DECK_PPTX_TEMPLATE_FILENAME = "pptx_template.pptx"
DECK_PPTX_TEMPLATE_MANIFEST_FILENAME = "pptx_template_manifest.json"

_TITLE_PLACEHOLDER_TYPES = {1, 3}
_TEXT_PLACEHOLDER_TYPES = {2, 7}
_VISUAL_PLACEHOLDER_TYPES = {7, 18}
_IGNORED_PLACEHOLDER_TYPES = {13, 15, 16}
_PROTOTYPE_CLEAR_TEXTS = {
    "client wordmark",
    "confidential | client name",
}
_TEXT_VISUAL_HINTS = {
    "backup exhibit / large table / image",
    "supporting notes",
}
_TEXT_VISUAL_REMOVE_HINTS = {
    "backup exhibit / large table / image",
    "optional mini-exhibit",
}
_DIGITS_RE = re.compile(r"\d+")


@dataclass(frozen=True, slots=True)
class DeckPptxTemplateLayout:
    """Normalized semantic layout mapping extracted from a PPTX template."""

    role: str
    layout_index: int
    layout_name: str
    title_placeholder_idx: int | None = None
    text_placeholder_idx: int | None = None
    visual_placeholder_idx: int | None = None
    visual_placeholder_type: int | None = None


@dataclass(frozen=True, slots=True)
class DeckPptxTemplatePrototype:
    """Reusable sample-slide prototype extracted from a PPTX template deck."""

    role: str
    slide_index: int
    title_shape_index: int | None = None
    subtitle_shape_index: int | None = None
    body_label_shape_index: int | None = None
    body_shape_index: int | None = None
    page_number_shape_index: int | None = None
    clear_shape_indices: list[int] = field(default_factory=list)
    remove_shape_indices: list[int] = field(default_factory=list)
    visual_left: int = 0
    visual_top: int = 0
    visual_width: int = 0
    visual_height: int = 0


@dataclass(frozen=True, slots=True)
class DeckPptxTemplateManifest:
    """Persisted manifest describing usable semantic layouts in a PPTX template."""

    template_name: str
    layouts: list[DeckPptxTemplateLayout]
    prototypes: list[DeckPptxTemplatePrototype] = field(default_factory=list)

    def layout_for_role(self, role: str) -> DeckPptxTemplateLayout | None:
        for layout in self.layouts:
            if layout.role == role:
                return layout
        return None

    def prototype_for_role(self, role: str) -> DeckPptxTemplatePrototype | None:
        for prototype in self.prototypes:
            if prototype.role == role:
                return prototype
        return None


def deck_pptx_template_path(deck_path: Path) -> Path:
    """Return the deck-local uploaded PPTX template path."""

    return deck_path / DECK_PPTX_TEMPLATE_FILENAME


def deck_pptx_template_manifest_path(deck_path: Path) -> Path:
    """Return the persisted deck-local PPTX template manifest path."""

    return deck_path / DECK_PPTX_TEMPLATE_MANIFEST_FILENAME


def write_deck_pptx_template_manifest(
    deck_path: Path, manifest: DeckPptxTemplateManifest
) -> Path:
    """Persist ``manifest`` next to the deck-local PPTX template."""

    output_path = deck_pptx_template_manifest_path(deck_path)
    payload = {
        "templateName": manifest.template_name,
        "layouts": [asdict(layout) for layout in manifest.layouts],
        "prototypes": [asdict(prototype) for prototype in manifest.prototypes],
    }
    output_path.write_text(
        json.dumps(payload, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    return output_path


def load_deck_pptx_template_manifest(deck_path: Path) -> DeckPptxTemplateManifest | None:
    """Load a persisted deck-local PPTX template manifest when available."""

    manifest_path = deck_pptx_template_manifest_path(deck_path)
    if not manifest_path.exists():
        return None
    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise ValueError("PPTX template manifest must be a JSON object.")
    raw_layouts = payload.get("layouts")
    if not isinstance(raw_layouts, list):
        raise ValueError("PPTX template manifest is missing layouts.")
    layouts: list[DeckPptxTemplateLayout] = []
    for item in raw_layouts:
        if not isinstance(item, dict):
            continue
        layouts.append(
            DeckPptxTemplateLayout(
                role=str(item.get("role") or "").strip(),
                layout_index=int(
                    _read_manifest_field(item, "layout_index", "layoutIndex", default=0)
                ),
                layout_name=str(item.get("layout_name") or item.get("layoutName") or "").strip(),
                title_placeholder_idx=_optional_int(
                    _read_manifest_field(
                        item,
                        "title_placeholder_idx",
                        "titlePlaceholderIdx",
                    )
                ),
                text_placeholder_idx=_optional_int(
                    _read_manifest_field(
                        item,
                        "text_placeholder_idx",
                        "textPlaceholderIdx",
                    )
                ),
                visual_placeholder_idx=_optional_int(
                    _read_manifest_field(
                        item,
                        "visual_placeholder_idx",
                        "visualPlaceholderIdx",
                    )
                ),
                visual_placeholder_type=_optional_int(
                    _read_manifest_field(
                        item,
                        "visual_placeholder_type",
                        "visualPlaceholderType",
                    )
                ),
            )
        )
    raw_prototypes = payload.get("prototypes")
    prototypes: list[DeckPptxTemplatePrototype] = []
    if isinstance(raw_prototypes, list):
        for item in raw_prototypes:
            if not isinstance(item, dict):
                continue
            prototypes.append(
                DeckPptxTemplatePrototype(
                    role=str(item.get("role") or "").strip(),
                    slide_index=int(
                        _read_manifest_field(item, "slide_index", "slideIndex", default=0)
                    ),
                    title_shape_index=_optional_int(
                        _read_manifest_field(
                            item,
                            "title_shape_index",
                            "titleShapeIndex",
                        )
                    ),
                    subtitle_shape_index=_optional_int(
                        _read_manifest_field(
                            item,
                            "subtitle_shape_index",
                            "subtitleShapeIndex",
                        )
                    ),
                    body_label_shape_index=_optional_int(
                        _read_manifest_field(
                            item,
                            "body_label_shape_index",
                            "bodyLabelShapeIndex",
                        )
                    ),
                    body_shape_index=_optional_int(
                        _read_manifest_field(
                            item,
                            "body_shape_index",
                            "bodyShapeIndex",
                        )
                    ),
                    page_number_shape_index=_optional_int(
                        _read_manifest_field(
                            item,
                            "page_number_shape_index",
                            "pageNumberShapeIndex",
                        )
                    ),
                    clear_shape_indices=[
                        int(value)
                        for value in (
                            _read_manifest_field(
                                item,
                                "clear_shape_indices",
                                "clearShapeIndices",
                                default=[],
                            )
                            or []
                        )
                    ],
                    remove_shape_indices=[
                        int(value)
                        for value in (
                            _read_manifest_field(
                                item,
                                "remove_shape_indices",
                                "removeShapeIndices",
                                default=[],
                            )
                            or []
                        )
                    ],
                    visual_left=int(
                        _read_manifest_field(item, "visual_left", "visualLeft", default=0)
                    ),
                    visual_top=int(
                        _read_manifest_field(item, "visual_top", "visualTop", default=0)
                    ),
                    visual_width=int(
                        _read_manifest_field(item, "visual_width", "visualWidth", default=0)
                    ),
                    visual_height=int(
                        _read_manifest_field(
                            item,
                            "visual_height",
                            "visualHeight",
                            default=0,
                        )
                    ),
                )
            )
    return DeckPptxTemplateManifest(
        template_name=str(payload.get("templateName") or "").strip(),
        layouts=layouts,
        prototypes=prototypes,
    )


def ensure_deck_pptx_template_manifest(
    deck_path: Path,
) -> DeckPptxTemplateManifest | None:
    """Load or build the deck-local PPTX template manifest when a template exists."""

    template_path = deck_pptx_template_path(deck_path)
    if not template_path.exists():
        return None
    cached_manifest = load_deck_pptx_template_manifest(deck_path)
    if cached_manifest is not None:
        return cached_manifest
    manifest = build_pptx_template_manifest(template_path)
    write_deck_pptx_template_manifest(deck_path, manifest)
    return manifest


def build_pptx_template_manifest(template_path: Path) -> DeckPptxTemplateManifest:
    """Extract a normalized semantic manifest from ``template_path``."""

    from pptx import Presentation

    presentation = Presentation(str(template_path))
    title_only_candidate: tuple[int, DeckPptxTemplateLayout] | None = None
    title_body_candidate: tuple[int, DeckPptxTemplateLayout] | None = None
    text_visual_candidate: tuple[int, DeckPptxTemplateLayout] | None = None

    for layout_index, layout in enumerate(presentation.slide_layouts):
        placeholders = _semantic_placeholders(layout)
        if not placeholders:
            continue
        title_placeholders = [
            entry for entry in placeholders if entry["type"] in _TITLE_PLACEHOLDER_TYPES
        ]
        semantic_content = [
            entry for entry in placeholders if entry["type"] not in _TITLE_PLACEHOLDER_TYPES
        ]
        if title_placeholders and not semantic_content:
            layout_manifest = DeckPptxTemplateLayout(
                role="title_only",
                layout_index=layout_index,
                layout_name=str(layout.name or "").strip(),
                title_placeholder_idx=int(title_placeholders[0]["idx"]),
            )
            score = _title_only_layout_score(str(layout.name or ""), placeholders)
            if title_only_candidate is None or score > title_only_candidate[0]:
                title_only_candidate = (score, layout_manifest)
        if title_placeholders and len(semantic_content) == 1:
            content_entry = semantic_content[0]
            if int(content_entry["type"]) in _TEXT_PLACEHOLDER_TYPES:
                layout_manifest = DeckPptxTemplateLayout(
                    role="title_body",
                    layout_index=layout_index,
                    layout_name=str(layout.name or "").strip(),
                    title_placeholder_idx=int(title_placeholders[0]["idx"]),
                    text_placeholder_idx=int(content_entry["idx"]),
                )
                score = _title_body_layout_score(str(layout.name or ""), semantic_content)
                if title_body_candidate is None or score > title_body_candidate[0]:
                    title_body_candidate = (score, layout_manifest)
        if title_placeholders and len(semantic_content) >= 2:
            text_entry, visual_entry = _pick_text_visual_entries(semantic_content)
            if text_entry is not None and visual_entry is not None:
                layout_manifest = DeckPptxTemplateLayout(
                    role="text_visual",
                    layout_index=layout_index,
                    layout_name=str(layout.name or "").strip(),
                    title_placeholder_idx=int(title_placeholders[0]["idx"]),
                    text_placeholder_idx=int(text_entry["idx"]),
                    visual_placeholder_idx=int(visual_entry["idx"]),
                    visual_placeholder_type=int(visual_entry["type"]),
                )
                score = _text_visual_layout_score(
                    str(layout.name or ""),
                    text_entry=text_entry,
                    visual_entry=visual_entry,
                )
                if text_visual_candidate is None or score > text_visual_candidate[0]:
                    text_visual_candidate = (score, layout_manifest)

    layouts: list[DeckPptxTemplateLayout] = []
    if title_only_candidate is not None:
        layouts.append(title_only_candidate[1])
    if title_body_candidate is not None:
        layouts.append(title_body_candidate[1])
    if text_visual_candidate is not None:
        layouts.append(text_visual_candidate[1])
    prototypes = _build_template_prototypes(presentation)
    return DeckPptxTemplateManifest(
        template_name=template_path.name,
        layouts=layouts,
        prototypes=prototypes,
    )


def _build_template_prototypes(presentation) -> list[DeckPptxTemplatePrototype]:
    prototypes: list[DeckPptxTemplatePrototype] = []
    text_visual = _build_text_visual_prototype(presentation)
    if text_visual is not None:
        prototypes.append(text_visual)
    return prototypes


def _build_text_visual_prototype(presentation) -> DeckPptxTemplatePrototype | None:
    best_candidate: tuple[int, DeckPptxTemplatePrototype] | None = None
    slide_height = int(presentation.slide_height)
    for slide_index, slide in enumerate(presentation.slides):
        entries = _slide_text_entries(slide)
        normalized_texts = {entry["normalized_text"] for entry in entries}
        if not _TEXT_VISUAL_HINTS.issubset(normalized_texts):
            continue
        body_label = _first_entry_with_text(entries, "supporting notes")
        body_entry = _largest_text_entry_below(entries, body_label)
        title_entry, subtitle_entry = _top_title_entries(entries)
        if body_label is None or body_entry is None or title_entry is None:
            continue
        page_number_entry = _page_number_entry(entries)
        clear_shape_indices = [
            int(entry["shape_index"])
            for entry in entries
            if entry["normalized_text"] in _PROTOTYPE_CLEAR_TEXTS
        ]
        remove_shape_indices = [
            int(entry["shape_index"])
            for entry in entries
            if entry["normalized_text"] in _TEXT_VISUAL_REMOVE_HINTS
        ]
        visual_left = min(
            int(entry["left"])
            for entry in entries
            if int(entry["shape_index"]) not in clear_shape_indices
        )
        visual_top = min(
            int(body_label["top"]),
            int(body_entry["top"]),
        )
        visual_right = int(body_label["left"]) - 182880
        visual_bottom = (
            int(page_number_entry["top"]) - 182880
            if page_number_entry is not None
            else slide_height - 182880
        )
        if visual_right <= visual_left or visual_bottom <= visual_top:
            continue
        prototype = DeckPptxTemplatePrototype(
            role="text_visual",
            slide_index=slide_index,
            title_shape_index=int(title_entry["shape_index"]),
            subtitle_shape_index=(
                int(subtitle_entry["shape_index"]) if subtitle_entry is not None else None
            ),
            body_label_shape_index=int(body_label["shape_index"]),
            body_shape_index=int(body_entry["shape_index"]),
            page_number_shape_index=(
                int(page_number_entry["shape_index"]) if page_number_entry is not None else None
            ),
            clear_shape_indices=sorted(set(clear_shape_indices)),
            remove_shape_indices=sorted(set(remove_shape_indices)),
            visual_left=visual_left,
            visual_top=visual_top,
            visual_width=visual_right - visual_left,
            visual_height=visual_bottom - visual_top,
        )
        score = int(body_entry["area"]) + int(prototype.visual_width * prototype.visual_height)
        if best_candidate is None or score > best_candidate[0]:
            best_candidate = (score, prototype)
    return None if best_candidate is None else best_candidate[1]


def _slide_text_entries(slide) -> list[dict[str, Any]]:
    entries: list[dict[str, Any]] = []
    for shape_index, shape in enumerate(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        raw_text = " | ".join(
            paragraph.text.strip()
            for paragraph in shape.text_frame.paragraphs
            if paragraph.text.strip()
        )
        if not raw_text:
            continue
        normalized = _normalize_text(raw_text)
        entries.append(
            {
                "shape_index": shape_index,
                "text": raw_text,
                "normalized_text": normalized,
                "left": int(shape.left),
                "top": int(shape.top),
                "width": int(shape.width),
                "height": int(shape.height),
                "area": int(shape.width) * int(shape.height),
            }
        )
    return entries


def _first_entry_with_text(
    entries: list[dict[str, Any]],
    text: str,
) -> dict[str, Any] | None:
    normalized = _normalize_text(text)
    for entry in entries:
        if entry["normalized_text"] == normalized:
            return entry
    return None


def _largest_text_entry_below(
    entries: list[dict[str, Any]],
    anchor: dict[str, Any] | None,
) -> dict[str, Any] | None:
    if anchor is None:
        return None
    candidates = [
        entry
        for entry in entries
        if int(entry["left"]) >= int(anchor["left"])
        and int(entry["top"]) > int(anchor["top"])
        and int(entry["width"]) >= int(anchor["width"] * 0.8)
    ]
    if not candidates:
        return None
    return max(candidates, key=lambda entry: int(entry["area"]))


def _top_title_entries(
    entries: list[dict[str, Any]],
) -> tuple[dict[str, Any] | None, dict[str, Any] | None]:
    candidates = [
        entry
        for entry in entries
        if entry["normalized_text"] not in _PROTOTYPE_CLEAR_TEXTS
        and not _DIGITS_RE.fullmatch(entry["normalized_text"])
        and int(entry["top"]) < 914400
        and int(entry["width"]) >= 4000000
    ]
    candidates = sorted(candidates, key=lambda entry: (int(entry["top"]), -int(entry["width"])))
    title = candidates[0] if candidates else None
    subtitle = candidates[1] if len(candidates) > 1 else None
    return title, subtitle


def _page_number_entry(entries: list[dict[str, Any]]) -> dict[str, Any] | None:
    candidates = [
        entry
        for entry in entries
        if _DIGITS_RE.fullmatch(entry["normalized_text"])
    ]
    if not candidates:
        return None
    return max(candidates, key=lambda entry: (int(entry["top"]), int(entry["left"])))


def _normalize_text(text: str) -> str:
    return " ".join(str(text or "").strip().lower().split())


def _semantic_placeholders(layout) -> list[dict[str, Any]]:
    entries: list[dict[str, Any]] = []
    for placeholder in layout.placeholders:
        placeholder_type = int(placeholder.placeholder_format.type)
        if placeholder_type in _IGNORED_PLACEHOLDER_TYPES:
            continue
        entries.append(
            {
                "idx": int(placeholder.placeholder_format.idx),
                "type": placeholder_type,
                "name": str(placeholder.name or ""),
                "left": int(placeholder.left),
                "top": int(placeholder.top),
                "width": int(placeholder.width),
                "height": int(placeholder.height),
                "area": int(placeholder.width) * int(placeholder.height),
            }
        )
    return entries


def _pick_text_visual_entries(
    entries: list[dict[str, Any]],
) -> tuple[dict[str, Any] | None, dict[str, Any] | None]:
    picture_entries = [entry for entry in entries if int(entry["type"]) == 18]
    text_entries = [entry for entry in entries if int(entry["type"]) in _TEXT_PLACEHOLDER_TYPES]
    if picture_entries and text_entries:
        visual_entry = max(picture_entries, key=lambda entry: int(entry["area"]))
        remaining_text = [entry for entry in text_entries if int(entry["idx"]) != int(visual_entry["idx"])]
        text_entry = min(
            remaining_text or text_entries,
            key=lambda entry: (int(entry["left"]), int(entry["top"])),
        )
        return text_entry, visual_entry
    if len(text_entries) < 2:
        return None, None
    sorted_entries = sorted(text_entries, key=lambda entry: (int(entry["left"]), int(entry["top"])))
    text_entry = sorted_entries[0]
    visual_entry = max(sorted_entries[1:], key=lambda entry: int(entry["area"]))
    return text_entry, visual_entry


def _title_only_layout_score(name: str, placeholders: list[dict[str, Any]]) -> int:
    lowered = name.lower()
    score = 0
    if "title only" in lowered:
        score += 100
    if "title" in lowered:
        score += 20
    score -= len(placeholders)
    return score


def _title_body_layout_score(name: str, placeholders: list[dict[str, Any]]) -> int:
    lowered = name.lower()
    score = 0
    if "title and content" in lowered:
        score += 100
    elif "title" in lowered and "content" in lowered:
        score += 80
    elif "title" in lowered and "text" in lowered:
        score += 60
    if len(placeholders) == 1:
        score += 10
    return score


def _text_visual_layout_score(
    name: str,
    *,
    text_entry: dict[str, Any],
    visual_entry: dict[str, Any],
) -> int:
    lowered = name.lower()
    score = 0
    if "content with caption" in lowered:
        score += 120
    elif "picture with caption" in lowered:
        score += 100
    elif "two content" in lowered:
        score += 90
    elif "comparison" in lowered:
        score += 80
    if int(visual_entry["left"]) > int(text_entry["left"]):
        score += 20
    score += int(visual_entry["area"] // 1000000)
    return score


def _optional_int(value: object) -> int | None:
    if value in (None, ""):
        return None
    return int(value)


def _read_manifest_field(
    payload: dict[str, Any],
    snake_key: str,
    camel_key: str,
    *,
    default: object = None,
) -> object:
    if snake_key in payload:
        return payload[snake_key]
    if camel_key in payload:
        return payload[camel_key]
    return default
