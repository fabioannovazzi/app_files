from __future__ import annotations

import io
import json
import re
from copy import deepcopy
from dataclasses import asdict, dataclass, field
from math import ceil
from pathlib import Path
from statistics import median
from typing import Any, Mapping, Sequence
from urllib.parse import unquote, urlparse

from bs4 import BeautifulSoup  # type: ignore[import]
from PIL import Image, ImageChops

from src.slides.launch_report_ast import validate_launch_report_payload
from src.slides.layout_semantics import (
    BULLET_BLOCK_TYPES,
    VISUAL_BLOCK_TYPES as _SEMANTIC_VISUAL_BLOCK_TYPES,
    block_sort_key,
    normalize_block_type,
    normalize_list_level,
    normalize_optional_string,
    normalize_render_mode,
)
from src.slides.models import Deck, Slide
from src.slides.notebooklm_style import (
    DEFAULT_PROMPT_STYLE_KEY,
    load_notebooklm_style,
    resolve_prompt_style_key,
)
from src.slides.ocr_cleanup import clean_ocr_text
from src.slides.pptx_template_manifest import (
    DeckPptxTemplateLayout,
    DeckPptxTemplateManifest,
    DeckPptxTemplatePrototype,
    deck_pptx_template_path,
    ensure_deck_pptx_template_manifest,
)

__all__ = [
    "SLIDES_PPTX_SPEC_FILENAME",
    "SlidesPptxSlide",
    "SlidesPptxSpec",
    "build_slides_pptx_spec_from_report_payload",
    "build_slides_pptx_spec",
    "load_slides_pptx_spec",
    "render_slides_pptx_from_template",
    "slides_pptx_spec_path",
    "write_slides_pptx_spec",
]

SLIDES_PPTX_SPEC_FILENAME = "slides_pptx_spec.json"
_PPTX_ASSET_DIRNAME = "pptx_assets"
_SLIDE_WIDTH_IN = 13.333333
_SLIDE_HEIGHT_IN = 7.5
_TITLE_LEFT_IN = 0.58
_TITLE_TOP_IN = 0.42
_TITLE_WIDTH_IN = 12.1
_TITLE_HEIGHT_IN = 0.86
_TITLE_ONLY_LEFT_IN = 1.18
_TITLE_ONLY_TOP_IN = 1.35
_TITLE_ONLY_WIDTH_IN = 10.95
_TITLE_ONLY_HEIGHT_IN = 2.8
_COVER_BODY_LEFT_IN = 0.96
_COVER_BODY_TOP_IN = 2.18
_COVER_BODY_WIDTH_IN = 9.95
_COVER_BODY_HEIGHT_IN = 1.18
_FOOTER_LEFT_IN = 0.82
_FOOTER_TOP_IN = 6.3
_FOOTER_WIDTH_IN = 5.4
_FOOTER_HEIGHT_IN = 0.48
_FULL_TEXT_LEFT_IN = 0.78
_FULL_TEXT_TOP_IN = 1.58
_FULL_TEXT_WIDTH_IN = 11.75
_FULL_TEXT_HEIGHT_IN = 5.2
_FULL_VISUAL_LEFT_IN = 0.76
_FULL_VISUAL_TOP_IN = 1.46
_FULL_VISUAL_WIDTH_IN = 11.82
_FULL_VISUAL_HEIGHT_IN = 5.45
_TABLE_VISUAL_TOP_IN = 2.12
_TABLE_VISUAL_HEIGHT_IN = 4.78
_RIGHT_COPY_LEFT_IN = 0.78
_RIGHT_COPY_TOP_IN = 1.56
_RIGHT_COPY_WIDTH_IN = 4.9
_RIGHT_COPY_HEIGHT_IN = 5.08
_RIGHT_VISUAL_LEFT_IN = 5.94
_RIGHT_VISUAL_TOP_IN = 1.54
_RIGHT_VISUAL_WIDTH_IN = 6.48
_RIGHT_VISUAL_HEIGHT_IN = 5.08
_BOTTOM_BULLETS_TOP_IN = 1.50
_BOTTOM_BULLETS_HEIGHT_IN = 2.05
_BOTTOM_VISUAL_TOP_IN = 3.82
_BOTTOM_VISUAL_HEIGHT_IN = 2.92
_BODY_ABOVE_VISUAL_TOP_IN = 1.45
_BODY_ABOVE_VISUAL_HEIGHT_IN = 0.48
_TEXT_RIGHT_COPY_LEFT_IN = 0.82
_TEXT_RIGHT_COPY_TOP_IN = 1.58
_TEXT_RIGHT_COPY_WIDTH_IN = 4.96
_TEXT_RIGHT_COPY_HEIGHT_IN = 5.02
_TEXT_BOTTOM_TOP_IN = 1.48
_TEXT_BOTTOM_HEIGHT_IN = 1.16
_TEXT_BOTTOM_VISUAL_TOP_IN = 3.02
_TEXT_BOTTOM_VISUAL_HEIGHT_IN = 3.92
_COMPARISON_LEFT_IN = 0.84
_COMPARISON_TOP_IN = 2.2
_COMPARISON_COLUMN_WIDTH_IN = 5.42
_COMPARISON_COLUMN_GAP_IN = 0.48
_COMPARISON_HEADER_HEIGHT_IN = 0.38
_COMPARISON_BULLETS_HEIGHT_IN = 2.2
_COMPARISON_CALLOUT_TOP_IN = 5.22
_COMPARISON_CALLOUT_HEIGHT_IN = 0.82
_BOTTOM_BANNER_LEFT_IN = 0.78
_BOTTOM_BANNER_WIDTH_IN = 11.75
_BOTTOM_BANNER_TOP_IN = 6.12
_BOTTOM_BANNER_HEIGHT_IN = 0.84
_VISUAL_BLOCK_TYPES = {"figure", "image", "chart", "table"}
_NATIVE_TABLE_MIN_CONFIDENCE = 0.72
_SPACE_RE = re.compile(r"\s+")
_NON_WORD_JOIN_RE = re.compile(r"\s+([,.;:%!?])")
_LEADING_BULLET_MARKER_RE = re.compile(r"^\s*[•·▪◦*-]\s*")
_GROUP_AS_IMAGE_RENDER_MODE = "group_as_image"
_LEAD_IN_PREFIXES = ("Governing thought:", "IMPLICATION:")
_IMPLICATION_PREFIX = "IMPLICATION:"
_COMPARISON_GROUP_KINDS = {"comparison", "comparison_panel", "comparison_columns"}
_EXHIBIT_GROUP_KINDS = {"exhibit", "flowchart", "flow_diagram", "diagram"}
_NATIVE_CARD_GROUP_KINDS = _EXHIBIT_GROUP_KINDS | {
    "cards",
    "card_grid",
    "cards_row",
    "process_flow",
}


@dataclass(frozen=True, slots=True)
class SlidesPptxSlide:
    """Semantic slide content rendered into the shared consulting template."""

    slide_id: str
    kind: str
    layout_variant: str
    title: str
    density_hint: str = ""
    body: str = ""
    bullets: list[str] = field(default_factory=list)
    visual_path: str = ""
    visual_type: str = ""
    table_model: dict[str, object] | None = None
    footer_text: str = ""
    implication: str = ""
    table_title: str = ""
    callout_title: str = ""
    callout_body: str = ""
    comparison_columns: list[dict[str, object]] = field(default_factory=list)
    native_visual: dict[str, object] | None = None
    repair_hints: dict[str, object] = field(default_factory=dict)


@dataclass(frozen=True, slots=True)
class SlidesPptxSpec:
    """Persisted semantic PPTX spec for a slides deck."""

    template_key: str
    prompt_style: str
    slides: list[SlidesPptxSlide]


@dataclass(frozen=True, slots=True)
class _SectionHeaderSubsection:
    label: str
    is_current: bool = False


@dataclass(frozen=True, slots=True)
class _SectionHeaderEntry:
    label: str
    is_current: bool = False
    subsections: list[_SectionHeaderSubsection] = field(default_factory=list)


def slides_pptx_spec_path(deck_path: Path) -> Path:
    """Return the persisted semantic PPTX spec path for ``deck_path``."""

    return deck_path / SLIDES_PPTX_SPEC_FILENAME


def write_slides_pptx_spec(deck_path: Path, spec: SlidesPptxSpec) -> Path:
    """Persist ``spec`` next to the deck assets."""

    payload = {
        "templateKey": spec.template_key,
        "promptStyle": spec.prompt_style,
        "slides": [asdict(slide) for slide in spec.slides],
    }
    output_path = slides_pptx_spec_path(deck_path)
    output_path.write_text(
        json.dumps(payload, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    return output_path


def load_slides_pptx_spec(deck_path: Path) -> SlidesPptxSpec:
    """Load the persisted semantic PPTX spec for ``deck_path``."""

    payload = json.loads(slides_pptx_spec_path(deck_path).read_text(encoding="utf-8"))
    if not isinstance(payload, Mapping):
        raise ValueError("Slides PPTX spec must be a JSON object.")
    raw_slides = payload.get("slides")
    if not isinstance(raw_slides, list):
        raise ValueError("Slides PPTX spec is missing slides.")
    slides = []
    for item in raw_slides:
        if not isinstance(item, Mapping):
            continue
        raw_comparison_columns = (
            item.get("comparison_columns")
            if isinstance(item.get("comparison_columns"), list)
            else (
                item.get("comparisonColumns")
                if isinstance(item.get("comparisonColumns"), list)
                else []
            )
        )
        slides.append(
            SlidesPptxSlide(
                slide_id=str(item.get("slide_id") or item.get("slideId") or "").strip(),
                kind=str(item.get("kind") or "").strip(),
                layout_variant=_resolve_layout_variant(
                    str(
                        item.get("layout_variant") or item.get("layoutVariant") or ""
                    ).strip(),
                    kind=str(item.get("kind") or "").strip(),
                    visual_type=str(
                        item.get("visual_type") or item.get("visualType") or ""
                    ).strip(),
                ),
                density_hint=str(
                    item.get("density_hint") or item.get("densityHint") or ""
                ).strip(),
                title=str(item.get("title") or "").strip(),
                body=str(item.get("body") or "").strip(),
                bullets=_read_text_list(item.get("bullets")),
                visual_path=str(
                    item.get("visual_path") or item.get("visualPath") or ""
                ).strip(),
                visual_type=str(
                    item.get("visual_type") or item.get("visualType") or ""
                ).strip(),
                table_model=(
                    item.get("table_model")
                    if isinstance(item.get("table_model"), Mapping)
                    else (
                        item.get("tableModel")
                        if isinstance(item.get("tableModel"), Mapping)
                        else None
                    )
                ),
                footer_text=str(
                    item.get("footer_text") or item.get("footerText") or ""
                ).strip(),
                implication=str(
                    item.get("implication") or item.get("implicationText") or ""
                ).strip(),
                table_title=str(
                    item.get("table_title") or item.get("tableTitle") or ""
                ).strip(),
                callout_title=str(
                    item.get("callout_title") or item.get("calloutTitle") or ""
                ).strip(),
                callout_body=str(
                    item.get("callout_body") or item.get("calloutBody") or ""
                ).strip(),
                comparison_columns=[
                    dict(column)
                    for column in raw_comparison_columns
                    if isinstance(column, Mapping)
                ],
                native_visual=(
                    dict(item.get("native_visual"))
                    if isinstance(item.get("native_visual"), Mapping)
                    else (
                        dict(item.get("nativeVisual"))
                        if isinstance(item.get("nativeVisual"), Mapping)
                        else None
                    )
                ),
                repair_hints=(
                    dict(item.get("repair_hints"))
                    if isinstance(item.get("repair_hints"), Mapping)
                    else (
                        dict(item.get("repairHints"))
                        if isinstance(item.get("repairHints"), Mapping)
                        else {}
                    )
                ),
            )
        )
    prompt_style = resolve_prompt_style_key(
        str(payload.get("promptStyle") or "").strip() or None
    )
    template_key = resolve_prompt_style_key(
        str(payload.get("templateKey") or "").strip() or prompt_style
    )
    return SlidesPptxSpec(
        template_key=template_key,
        prompt_style=prompt_style,
        slides=slides,
    )


def build_slides_pptx_spec(
    deck: Deck,
    deck_path: Path,
    *,
    slide_analysis: Mapping[str, object] | None,
) -> SlidesPptxSpec:
    """Build a deterministic semantic PPTX spec from deck content and analysis."""

    prompt_style = resolve_prompt_style_key(DEFAULT_PROMPT_STYLE_KEY)
    analysis_by_slide_id = _analysis_slide_map(slide_analysis)
    slides: list[SlidesPptxSlide] = []
    asset_dir = deck_path / _PPTX_ASSET_DIRNAME
    asset_dir.mkdir(parents=True, exist_ok=True)
    for stale_path in asset_dir.glob("*.png"):
        stale_path.unlink(missing_ok=True)
    for slide_number, slide in enumerate(deck.slides, start=1):
        if slide.is_section_header:
            slides.append(
                _build_section_header_slide_spec(
                    slide=slide,
                    slide_number=slide_number,
                )
            )
            continue
        analysis_slide = analysis_by_slide_id.get(slide.id)
        if analysis_slide is not None:
            slides.append(
                _build_analysis_slide_spec(
                    slide=slide,
                    slide_number=slide_number,
                    deck_path=deck_path,
                    analysis_slide=analysis_slide,
                    asset_dir=asset_dir,
                )
            )
            continue
        slides.append(
            _build_html_slide_spec(
                slide=slide,
                slide_number=slide_number,
                deck_path=deck_path,
            )
        )
    return SlidesPptxSpec(
        template_key=prompt_style,
        prompt_style=prompt_style,
        slides=slides,
    )


def build_slides_pptx_spec_from_report_payload(
    report_payload: Mapping[str, object],
    *,
    deck_path: Path,
) -> SlidesPptxSpec:
    """Build a semantic PPTX spec directly from a structured report payload."""

    validated_payload = validate_launch_report_payload(report_payload)
    raw_prompt_style = str(
        validated_payload.get("promptStyle")
        or validated_payload.get("prompt_style")
        or ""
    ).strip()
    prompt_style = resolve_prompt_style_key(
        raw_prompt_style or DEFAULT_PROMPT_STYLE_KEY
    )
    raw_template_key = str(
        validated_payload.get("templateKey")
        or validated_payload.get("template_key")
        or ""
    ).strip()
    template_key = resolve_prompt_style_key(raw_template_key or prompt_style)
    raw_slides = validated_payload.get("slides")
    slides: list[SlidesPptxSlide] = []
    if isinstance(raw_slides, list):
        for slide_number, raw_slide in enumerate(raw_slides, start=1):
            if not isinstance(raw_slide, Mapping):
                continue
            slides.append(
                _build_report_payload_slide_spec(
                    slide_payload=raw_slide,
                    slide_number=slide_number,
                    deck_path=deck_path,
                )
            )
    return SlidesPptxSpec(
        template_key=template_key,
        prompt_style=prompt_style,
        slides=slides,
    )


def _build_report_payload_slide_spec(
    *,
    slide_payload: Mapping[str, object],
    slide_number: int,
    deck_path: Path,
) -> SlidesPptxSlide:
    slide_id = (
        _read_report_text(slide_payload.get("slideId") or slide_payload.get("slide_id"))
        or f"slide{slide_number}.html"
    )
    title = _read_report_text(slide_payload.get("title")) or f"Slide {slide_number}"
    body = _read_report_body(
        slide_payload.get("body")
        if slide_payload.get("body") is not None
        else slide_payload.get("bodyLines")
    )
    bullets = _read_text_list(slide_payload.get("bullets"))
    footer_text = _read_report_text(
        slide_payload.get("footerText") or slide_payload.get("footer_text")
    )
    implication = _read_report_text(
        slide_payload.get("implication") or slide_payload.get("implicationText")
    )
    table_title = _read_report_text(
        slide_payload.get("tableTitle") or slide_payload.get("table_title")
    )
    callout_title = _read_report_text(
        slide_payload.get("calloutTitle") or slide_payload.get("callout_title")
    )
    callout_body = _read_report_text(
        slide_payload.get("calloutBody") or slide_payload.get("callout_body")
    )
    comparison_columns = _read_report_mapping_list(
        slide_payload.get("comparisonColumns")
        if slide_payload.get("comparisonColumns") is not None
        else slide_payload.get("comparison_columns")
    )
    table_model = (
        _read_report_mapping(
            slide_payload.get("tableModel")
            if slide_payload.get("tableModel") is not None
            else slide_payload.get("table_model")
        )
        or None
    )
    native_visual = (
        _read_report_mapping(
            slide_payload.get("nativeVisual")
            if slide_payload.get("nativeVisual") is not None
            else slide_payload.get("native_visual")
        )
        or None
    )
    repair_hints = _read_report_mapping(
        slide_payload.get("repairHints")
        if slide_payload.get("repairHints") is not None
        else slide_payload.get("repair_hints")
    )
    visual_payload = _read_report_mapping(slide_payload.get("visual"))
    visual_path = _normalize_report_visual_path(
        _read_report_text(
            slide_payload.get("visualPath")
            or slide_payload.get("visual_path")
            or visual_payload.get("path")
            or visual_payload.get("assetPath")
            or visual_payload.get("asset_path")
        ),
        deck_path=deck_path,
    )
    visual_type = _normalize_report_visual_type(
        _read_report_text(
            slide_payload.get("visualType")
            or slide_payload.get("visual_type")
            or visual_payload.get("kind")
        ),
        has_table_model=table_model is not None,
    )
    visual_layout_hint = _read_report_text(
        slide_payload.get("visualLayoutHint")
        or slide_payload.get("visual_layout_hint")
        or visual_payload.get("layoutHint")
        or visual_payload.get("layout_hint")
    ).lower()
    explicit_kind = _read_report_text(slide_payload.get("kind"))

    has_visual = bool(visual_path or native_visual)
    if comparison_columns:
        kind = "comparison_columns"
    elif explicit_kind:
        kind = explicit_kind
    else:
        kind = _resolve_slide_kind(
            bullets=bullets,
            body=body,
            has_visual=has_visual,
        )
        if footer_text and not has_visual and not bullets and body:
            kind = "cover_with_footer"

    density_text = _read_report_text(
        slide_payload.get("densityHint") or slide_payload.get("density_hint")
    ).lower()
    density_hint = (
        density_text
        if density_text in {"light", "medium", "dense"}
        else _classify_density(
            title=title,
            bullets=bullets,
            body="\n\n".join(
                part
                for part in (
                    body,
                    callout_title,
                    callout_body,
                    implication,
                    footer_text,
                    table_title,
                )
                if part
            ),
        )
    )
    layout_variant = _resolve_layout_variant(
        _read_report_text(
            slide_payload.get("layoutVariant") or slide_payload.get("layout_variant")
        ),
        kind=kind,
        visual_type=visual_type,
    )
    if not _read_report_text(
        slide_payload.get("layoutVariant") or slide_payload.get("layout_variant")
    ):
        layout_variant = _select_layout_variant(
            kind=kind,
            title=title,
            bullets=bullets,
            body=body,
            visual_type=visual_type,
            density_hint=density_hint,
            visual_layout_hint=visual_layout_hint,
        )

    return SlidesPptxSlide(
        slide_id=slide_id,
        kind=kind,
        layout_variant=layout_variant,
        density_hint=density_hint,
        title=title,
        body=body,
        bullets=bullets,
        visual_path=visual_path,
        visual_type=visual_type,
        table_model=table_model,
        footer_text=footer_text,
        implication=implication,
        table_title=table_title,
        callout_title=callout_title,
        callout_body=callout_body,
        comparison_columns=[dict(column) for column in comparison_columns],
        native_visual=dict(native_visual) if native_visual is not None else None,
        repair_hints=dict(repair_hints),
    )


def render_slides_pptx_from_template(deck_path: Path) -> io.BytesIO:
    """Render an editable PPTX from the persisted semantic slides spec."""

    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN

    spec = load_slides_pptx_spec(deck_path)
    custom_template_path = deck_pptx_template_path(deck_path)
    template_manifest = ensure_deck_pptx_template_manifest(deck_path)
    prototype_source_presentation = None
    if custom_template_path.exists():
        presentation = Presentation(str(custom_template_path))
        prototype_source_presentation = Presentation(str(custom_template_path))
    else:
        presentation = _load_template_presentation(Presentation, spec.template_key)
        presentation.slide_width = _inches(_SLIDE_WIDTH_IN)
        presentation.slide_height = _inches(_SLIDE_HEIGHT_IN)
    _clear_presentation_slides(presentation)
    blank_layout = presentation.slide_layouts[6]
    style = load_notebooklm_style(spec.prompt_style)
    font_name = style.font_family_primary or style.font_family_fallback
    title_size = max(float(style.title_size_pt) - 2.0, 22.0)
    body_size = max(float(style.body_size_pt), 15.0)
    line_height = float(style.line_height)
    text_rgb = _hex_to_rgb(style.text_color)
    muted_rgb = _hex_to_rgb("#000000")

    for index, slide_spec in enumerate(spec.slides, start=1):
        if template_manifest is not None:
            rendered_slide = _render_slide_with_custom_template(
                presentation,
                slide_spec=slide_spec,
                deck_path=deck_path,
                manifest=template_manifest,
                prototype_source_presentation=prototype_source_presentation,
                slide_number=index,
            )
            if rendered_slide is not None:
                _populate_slide_number_placeholder(rendered_slide, index)
                continue
        slide = presentation.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(_hex_to_rgb(style.bg_color))
        slide_body_size = _body_font_size_for_slide(
            body_size,
            layout_variant=slide_spec.layout_variant,
            density_hint=slide_spec.density_hint,
        )
        body_scale = _repair_hint_float(
            slide_spec,
            "body_scale",
            minimum=0.82,
            maximum=1.0,
            default=1.0,
        )
        slide_body_size = max(12.0, slide_body_size * body_scale)
        slide_line_height = _line_height_for_slide(
            line_height,
            density_hint=slide_spec.density_hint,
        )
        if body_scale < 1.0:
            slide_line_height = max(
                1.0, slide_line_height - ((1.0 - body_scale) * 0.25)
            )
        slide_paragraph_spacing = _paragraph_spacing_for_slide(
            density_hint=slide_spec.density_hint,
        )
        if body_scale < 1.0:
            slide_paragraph_spacing = max(
                4.0,
                slide_paragraph_spacing - ((1.0 - body_scale) * 8.0),
            )
        text_margin = _text_margin_for_slide(
            layout_variant=slide_spec.layout_variant,
            density_hint=slide_spec.density_hint,
        )
        title_layout = _title_layout_for_slide(title_size, slide_spec)
        banner_box = _implication_banner_box(slide_spec)
        content_bottom = _content_bottom_limit(banner_box)
        if slide_spec.layout_variant == "section_header_agenda":
            _add_section_header_agenda(
                slide,
                deck_path=deck_path,
                slide_spec=slide_spec,
                font_name=font_name,
                title_size=title_size,
                body_size=body_size,
                text_rgb=text_rgb,
            )
        elif slide_spec.layout_variant != "title_only_centered":
            _add_textbox(
                slide,
                text=slide_spec.title,
                left=float(title_layout["left"]),
                top=float(title_layout["top"]),
                width=float(title_layout["width"]),
                height=float(title_layout["height"]),
                font_name=font_name,
                font_size=float(title_layout["font_size"]),
                color_rgb=text_rgb,
                bold=True,
                margin=text_margin,
            )
        if slide_spec.layout_variant == "title_only_centered":
            _add_textbox(
                slide,
                text=slide_spec.title,
                left=_TITLE_ONLY_LEFT_IN,
                top=_TITLE_ONLY_TOP_IN,
                width=_TITLE_ONLY_WIDTH_IN,
                height=_TITLE_ONLY_HEIGHT_IN,
                font_name=font_name,
                font_size=max(title_size + 6.0, 30.0),
                color_rgb=text_rgb,
                bold=True,
                alignment=PP_ALIGN.CENTER,
                margin=text_margin,
            )
        elif slide_spec.layout_variant == "cover_with_footer":
            body_top = _content_top_after_title(
                _COVER_BODY_TOP_IN, title_layout=title_layout
            )
            if slide_spec.body:
                _add_textbox(
                    slide,
                    text=slide_spec.body,
                    left=_COVER_BODY_LEFT_IN,
                    top=body_top,
                    width=_COVER_BODY_WIDTH_IN,
                    height=_COVER_BODY_HEIGHT_IN,
                    font_name=font_name,
                    font_size=max(slide_body_size + 1.0, 16.0),
                    color_rgb=text_rgb,
                    margin=text_margin,
                )
            if slide_spec.footer_text:
                _add_textbox(
                    slide,
                    text=slide_spec.footer_text,
                    left=_FOOTER_LEFT_IN,
                    top=_FOOTER_TOP_IN,
                    width=_FOOTER_WIDTH_IN,
                    height=_FOOTER_HEIGHT_IN,
                    font_name=font_name,
                    font_size=12.0,
                    color_rgb=text_rgb,
                    margin=0.02,
                )
        elif slide_spec.layout_variant == "comparison_columns":
            body_top = _content_top_after_title(1.66, title_layout=title_layout)
            column_top = body_top
            if slide_spec.body:
                body_height = _text_box_height_for_text(
                    slide_spec.body,
                    base_height=0.64,
                    max_height=0.98,
                )
                _add_textbox(
                    slide,
                    text=slide_spec.body,
                    left=_FULL_TEXT_LEFT_IN,
                    top=body_top,
                    width=_FULL_TEXT_WIDTH_IN,
                    height=body_height,
                    font_name=font_name,
                    font_size=max(slide_body_size - 0.35, 13.5),
                    color_rgb=text_rgb,
                    margin=text_margin,
                    lead_prefixes=("Governing thought:",),
                )
                column_top = body_top + body_height + 0.18
            callout_height = 0.0
            if slide_spec.callout_title or slide_spec.callout_body:
                callout_height = min(
                    1.08,
                    max(
                        _COMPARISON_CALLOUT_HEIGHT_IN,
                        0.42
                        + (
                            0.22
                            * _estimate_wrapped_lines(
                                slide_spec.callout_body, line_capacity=88
                            )
                        ),
                    ),
                )
            column_bottom = max(
                column_top + 1.18,
                content_bottom - (callout_height + 0.22 if callout_height else 0.0),
            )
            bullet_top = column_top + _COMPARISON_HEADER_HEIGHT_IN + 0.08
            bullet_height = max(0.92, column_bottom - bullet_top - 0.04)
            for column_index, column in enumerate(slide_spec.comparison_columns[:2]):
                left = _COMPARISON_LEFT_IN + (
                    column_index
                    * (_COMPARISON_COLUMN_WIDTH_IN + _COMPARISON_COLUMN_GAP_IN)
                )
                title_text = _normalize_text(str(column.get("title") or ""))
                bullets = _read_text_list(column.get("bullets"))
                if title_text:
                    _add_textbox(
                        slide,
                        text=title_text,
                        left=left,
                        top=column_top,
                        width=_COMPARISON_COLUMN_WIDTH_IN,
                        height=_COMPARISON_HEADER_HEIGHT_IN,
                        font_name=font_name,
                        font_size=max(slide_body_size + 0.8, 15.0),
                        color_rgb=text_rgb,
                        bold=True,
                        margin=0.01,
                    )
                if bullets:
                    _add_bullet_box(
                        slide,
                        bullets=bullets,
                        left=left,
                        top=bullet_top,
                        width=_COMPARISON_COLUMN_WIDTH_IN,
                        height=bullet_height,
                        font_name=font_name,
                        body_size=max(slide_body_size - 0.8, 12.5),
                        line_height=slide_line_height,
                        paragraph_spacing=max(slide_paragraph_spacing - 1.5, 5.0),
                        text_rgb=text_rgb,
                        margin=0.02,
                    )
            if callout_height:
                callout_top = min(
                    max(_COMPARISON_CALLOUT_TOP_IN, column_bottom + 0.12),
                    content_bottom - callout_height,
                )
                _add_callout_box(
                    slide,
                    title=slide_spec.callout_title,
                    body=slide_spec.callout_body,
                    left=1.08,
                    top=callout_top,
                    width=11.1,
                    height=callout_height,
                    font_name=font_name,
                    body_size=max(slide_body_size - 0.2, 13.0),
                    text_rgb=text_rgb,
                )
        elif slide_spec.layout_variant == "bullets_full_width":
            bullets_top = _content_top_after_title(
                _FULL_TEXT_TOP_IN, title_layout=title_layout
            )
            _add_bullet_box(
                slide,
                bullets=slide_spec.bullets,
                left=_FULL_TEXT_LEFT_IN,
                top=bullets_top,
                width=_FULL_TEXT_WIDTH_IN,
                height=max(1.2, content_bottom - bullets_top),
                font_name=font_name,
                body_size=slide_body_size,
                line_height=slide_line_height,
                paragraph_spacing=slide_paragraph_spacing,
                text_rgb=text_rgb,
                margin=text_margin,
            )
        elif slide_spec.layout_variant == "text_full_width":
            body_top = _content_top_after_title(
                _FULL_TEXT_TOP_IN, title_layout=title_layout
            )
            _add_textbox(
                slide,
                text=slide_spec.body,
                left=_FULL_TEXT_LEFT_IN,
                top=body_top,
                width=_FULL_TEXT_WIDTH_IN,
                height=max(1.2, content_bottom - body_top),
                font_name=font_name,
                font_size=slide_body_size,
                color_rgb=text_rgb,
                margin=text_margin,
                lead_prefixes=_LEAD_IN_PREFIXES,
            )
        elif slide_spec.layout_variant == "visual_full_width":
            visual_top = _content_top_after_title(
                _FULL_VISUAL_TOP_IN, title_layout=title_layout
            )
            fit_left, fit_top, fit_width, fit_height = _apply_visual_box_repairs(
                slide_spec,
                left=_FULL_VISUAL_LEFT_IN,
                top=visual_top,
                width=_FULL_VISUAL_WIDTH_IN,
                height=max(1.5, content_bottom - visual_top),
                content_bottom=content_bottom,
            )
            _add_visual(
                slide,
                deck_path=deck_path,
                relative_path=slide_spec.visual_path,
                native_visual=slide_spec.native_visual,
                left=fit_left,
                top=fit_top,
                width=fit_width,
                height=fit_height,
                font_name=font_name,
                body_size=max(slide_body_size - 0.5, 12.0),
                text_rgb=text_rgb,
            )
        elif slide_spec.layout_variant == "table_focus":
            body_top = _content_top_after_title(
                _BODY_ABOVE_VISUAL_TOP_IN, title_layout=title_layout
            )
            cursor_top = body_top
            if slide_spec.body:
                body_height = _text_box_height_for_text(
                    slide_spec.body,
                    base_height=0.48,
                    max_height=0.92,
                )
                _add_textbox(
                    slide,
                    text=slide_spec.body,
                    left=_FULL_TEXT_LEFT_IN,
                    top=body_top,
                    width=_FULL_TEXT_WIDTH_IN,
                    height=body_height,
                    font_name=font_name,
                    font_size=max(slide_body_size - 0.5, 13.0),
                    color_rgb=text_rgb,
                    margin=text_margin,
                    lead_prefixes=("Governing thought:",),
                )
                cursor_top = body_top + body_height + 0.08
            if slide_spec.table_title:
                _add_textbox(
                    slide,
                    text=slide_spec.table_title,
                    left=_FULL_TEXT_LEFT_IN,
                    top=cursor_top,
                    width=_FULL_TEXT_WIDTH_IN,
                    height=0.34,
                    font_name=font_name,
                    font_size=max(slide_body_size - 0.2, 13.5),
                    color_rgb=text_rgb,
                    bold=True,
                    margin=0.01,
                )
                cursor_top += 0.38
            table_top = max(_TABLE_VISUAL_TOP_IN, cursor_top + 0.06)
            table_height = max(1.6, content_bottom - table_top)
            if not _add_native_table(
                slide,
                table_model=slide_spec.table_model,
                left=_FULL_VISUAL_LEFT_IN,
                top=table_top,
                width=_FULL_VISUAL_WIDTH_IN,
                height=table_height,
                font_name=font_name,
                body_size=max(slide_body_size - 1.0, 11.0),
                text_rgb=text_rgb,
            ):
                fit_left, fit_top, fit_width, fit_height = _apply_visual_box_repairs(
                    slide_spec,
                    left=_FULL_VISUAL_LEFT_IN,
                    top=table_top,
                    width=_FULL_VISUAL_WIDTH_IN,
                    height=table_height,
                    content_bottom=content_bottom,
                )
                _add_visual(
                    slide,
                    deck_path=deck_path,
                    relative_path=slide_spec.visual_path,
                    native_visual=slide_spec.native_visual,
                    left=fit_left,
                    top=fit_top,
                    width=fit_width,
                    height=fit_height,
                    font_name=font_name,
                    body_size=max(slide_body_size - 0.8, 11.5),
                    text_rgb=text_rgb,
                )
        elif slide_spec.layout_variant == "text_visual_right":
            body_top = _content_top_after_title(
                _TEXT_RIGHT_COPY_TOP_IN, title_layout=title_layout
            )
            column_height = max(1.4, content_bottom - body_top)
            if slide_spec.body:
                _add_textbox(
                    slide,
                    text=slide_spec.body,
                    left=_TEXT_RIGHT_COPY_LEFT_IN,
                    top=body_top,
                    width=_TEXT_RIGHT_COPY_WIDTH_IN,
                    height=column_height,
                    font_name=font_name,
                    font_size=slide_body_size,
                    color_rgb=text_rgb,
                    margin=text_margin,
                    lead_prefixes=_LEAD_IN_PREFIXES,
                )
            visual_top = max(_RIGHT_VISUAL_TOP_IN, body_top)
            fit_left, fit_top, fit_width, fit_height = _apply_visual_box_repairs(
                slide_spec,
                left=_RIGHT_VISUAL_LEFT_IN,
                top=visual_top,
                width=_RIGHT_VISUAL_WIDTH_IN,
                height=max(1.4, content_bottom - visual_top),
                content_bottom=content_bottom,
            )
            _add_visual(
                slide,
                deck_path=deck_path,
                relative_path=slide_spec.visual_path,
                native_visual=slide_spec.native_visual,
                left=fit_left,
                top=fit_top,
                width=fit_width,
                height=fit_height,
                font_name=font_name,
                body_size=max(slide_body_size - 0.5, 12.0),
                text_rgb=text_rgb,
            )
        elif slide_spec.layout_variant == "text_visual_bottom":
            body_top = _content_top_after_title(
                _TEXT_BOTTOM_TOP_IN, title_layout=title_layout
            )
            visual_top = _TEXT_BOTTOM_VISUAL_TOP_IN
            if slide_spec.body:
                body_height = _text_box_height_for_text(
                    slide_spec.body,
                    base_height=0.62,
                    max_height=1.02,
                )
                _add_textbox(
                    slide,
                    text=slide_spec.body,
                    left=_FULL_TEXT_LEFT_IN,
                    top=body_top,
                    width=_FULL_TEXT_WIDTH_IN,
                    height=body_height,
                    font_name=font_name,
                    font_size=max(slide_body_size - 0.5, 13.0),
                    color_rgb=text_rgb,
                    margin=text_margin,
                    lead_prefixes=_LEAD_IN_PREFIXES,
                )
                visual_top = max(
                    _TEXT_BOTTOM_VISUAL_TOP_IN, body_top + body_height + 0.12
                )
            fit_left, fit_top, fit_width, fit_height = _apply_visual_box_repairs(
                slide_spec,
                left=_FULL_VISUAL_LEFT_IN,
                top=visual_top,
                width=_FULL_VISUAL_WIDTH_IN,
                height=max(1.5, content_bottom - visual_top),
                content_bottom=content_bottom,
            )
            _add_visual(
                slide,
                deck_path=deck_path,
                relative_path=slide_spec.visual_path,
                native_visual=slide_spec.native_visual,
                left=fit_left,
                top=fit_top,
                width=fit_width,
                height=fit_height,
                font_name=font_name,
                body_size=max(slide_body_size - 0.5, 12.0),
                text_rgb=text_rgb,
            )
        elif slide_spec.layout_variant == "bullets_visual_right":
            bullets_top = _content_top_after_title(
                _RIGHT_COPY_TOP_IN, title_layout=title_layout
            )
            content_height = max(1.4, content_bottom - bullets_top)
            _add_bullet_box(
                slide,
                bullets=slide_spec.bullets,
                left=_RIGHT_COPY_LEFT_IN,
                top=bullets_top,
                width=_RIGHT_COPY_WIDTH_IN,
                height=content_height,
                font_name=font_name,
                body_size=slide_body_size,
                line_height=slide_line_height,
                paragraph_spacing=slide_paragraph_spacing,
                text_rgb=text_rgb,
                margin=text_margin,
            )
            visual_top = max(_RIGHT_VISUAL_TOP_IN, bullets_top)
            fit_left, fit_top, fit_width, fit_height = _apply_visual_box_repairs(
                slide_spec,
                left=_RIGHT_VISUAL_LEFT_IN,
                top=visual_top,
                width=_RIGHT_VISUAL_WIDTH_IN,
                height=max(1.4, content_bottom - visual_top),
                content_bottom=content_bottom,
            )
            _add_visual(
                slide,
                deck_path=deck_path,
                relative_path=slide_spec.visual_path,
                native_visual=slide_spec.native_visual,
                left=fit_left,
                top=fit_top,
                width=fit_width,
                height=fit_height,
                font_name=font_name,
                body_size=max(slide_body_size - 0.5, 12.0),
                text_rgb=text_rgb,
            )
        elif slide_spec.layout_variant == "bullets_visual_bottom":
            bullets_top = _content_top_after_title(
                _BOTTOM_BULLETS_TOP_IN, title_layout=title_layout
            )
            visual_top = _BOTTOM_VISUAL_TOP_IN
            visual_height = max(1.5, content_bottom - visual_top)
            _add_bullet_box(
                slide,
                bullets=slide_spec.bullets,
                left=_FULL_TEXT_LEFT_IN,
                top=bullets_top,
                width=_FULL_TEXT_WIDTH_IN,
                height=max(
                    0.9, min(_BOTTOM_BULLETS_HEIGHT_IN, visual_top - bullets_top - 0.18)
                ),
                font_name=font_name,
                body_size=slide_body_size,
                line_height=slide_line_height,
                paragraph_spacing=slide_paragraph_spacing,
                text_rgb=text_rgb,
                margin=text_margin,
            )
            visual_top = max(
                _BOTTOM_VISUAL_TOP_IN,
                bullets_top + min(_BOTTOM_BULLETS_HEIGHT_IN, 1.95) + 0.12,
            )
            fit_left, fit_top, fit_width, fit_height = _apply_visual_box_repairs(
                slide_spec,
                left=_FULL_VISUAL_LEFT_IN,
                top=visual_top,
                width=_FULL_VISUAL_WIDTH_IN,
                height=max(1.5, content_bottom - visual_top),
                content_bottom=content_bottom,
            )
            _add_visual(
                slide,
                deck_path=deck_path,
                relative_path=slide_spec.visual_path,
                native_visual=slide_spec.native_visual,
                left=fit_left,
                top=fit_top,
                width=fit_width,
                height=fit_height,
                font_name=font_name,
                body_size=max(slide_body_size - 0.5, 12.0),
                text_rgb=text_rgb,
            )
        if banner_box is not None:
            _add_bottom_banner(
                slide,
                text=slide_spec.implication,
                left=float(banner_box["left"]),
                top=float(banner_box["top"]),
                width=float(banner_box["width"]),
                height=float(banner_box["height"]),
                font_name=font_name,
                font_size=max(slide_body_size - 1.0, 11.5),
                text_rgb=text_rgb,
            )
        _add_textbox(
            slide,
            text=str(index),
            left=12.18,
            top=7.02,
            width=0.42,
            height=0.2,
            font_name=font_name,
            font_size=10.0,
            color_rgb=muted_rgb,
            alignment=PP_ALIGN.RIGHT,
            margin=0.0,
        )

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


def _render_slide_with_custom_template(
    presentation,
    *,
    slide_spec: SlidesPptxSlide,
    deck_path: Path,
    manifest: DeckPptxTemplateManifest,
    prototype_source_presentation,
    slide_number: int,
):
    prototype_role = _template_prototype_role_for_slide(slide_spec.kind)
    if prototype_role and prototype_source_presentation is not None:
        prototype_manifest = manifest.prototype_for_role(prototype_role)
        if prototype_manifest is not None:
            slide = _render_slide_with_custom_prototype(
                presentation,
                prototype_source_presentation=prototype_source_presentation,
                prototype_manifest=prototype_manifest,
                slide_spec=slide_spec,
                deck_path=deck_path,
                slide_number=slide_number,
            )
            if slide is not None:
                return slide
    role = _template_role_for_slide(slide_spec.kind)
    if not role:
        return None
    layout_manifest = manifest.layout_for_role(role)
    if layout_manifest is None:
        return None
    layout = presentation.slide_layouts[layout_manifest.layout_index]
    slide = presentation.slides.add_slide(layout)
    _fill_template_title_placeholder(
        slide,
        layout_manifest,
        slide_spec.title,
    )
    _fill_template_text_placeholder(
        slide,
        layout_manifest,
        slide_spec,
    )
    _fill_template_visual_slot(
        slide,
        layout_manifest,
        deck_path=deck_path,
        relative_path=slide_spec.visual_path,
        native_visual=slide_spec.native_visual,
        table_model=slide_spec.table_model,
    )
    return slide


def _template_prototype_role_for_slide(kind: str) -> str | None:
    if kind in {"bullets_visual", "text_visual"}:
        return "text_visual"
    return None


def _template_role_for_slide(kind: str) -> str | None:
    if kind == "title_only":
        return "title_only"
    if kind == "text_only":
        return "title_body"
    if kind in {"bullets_visual", "text_visual"}:
        return "text_visual"
    return None


def _render_slide_with_custom_prototype(
    presentation,
    *,
    prototype_source_presentation,
    prototype_manifest: DeckPptxTemplatePrototype,
    slide_spec: SlidesPptxSlide,
    deck_path: Path,
    slide_number: int,
):
    source_slides = list(prototype_source_presentation.slides)
    if prototype_manifest.slide_index < 0 or prototype_manifest.slide_index >= len(
        source_slides
    ):
        return None
    source_slide = source_slides[prototype_manifest.slide_index]
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _copy_slide_shapes(source_slide, slide)
    for shape_index in prototype_manifest.clear_shape_indices:
        shape = _shape_at_index(slide, shape_index)
        if shape is not None and getattr(shape, "has_text_frame", False):
            shape.text_frame.clear()
    _set_shape_text(
        slide,
        prototype_manifest.title_shape_index,
        slide_spec.title,
    )
    _set_shape_text(
        slide,
        prototype_manifest.subtitle_shape_index,
        "",
    )
    _set_shape_text(
        slide,
        prototype_manifest.body_label_shape_index,
        "Key takeaways" if slide_spec.bullets else "Supporting notes",
    )
    _set_shape_text(
        slide,
        prototype_manifest.body_shape_index,
        _prototype_body_text(slide_spec),
    )
    _set_shape_text(
        slide,
        prototype_manifest.page_number_shape_index,
        str(slide_number),
    )
    for shape_index in sorted(prototype_manifest.remove_shape_indices, reverse=True):
        shape = _shape_at_index(slide, shape_index)
        if shape is not None:
            _remove_shape(shape)
    if not _add_native_table_within_emu_box(
        slide,
        table_model=slide_spec.table_model,
        left=prototype_manifest.visual_left,
        top=prototype_manifest.visual_top,
        width=prototype_manifest.visual_width,
        height=prototype_manifest.visual_height,
    ):
        _add_visual_within_emu_box(
            slide,
            deck_path=deck_path,
            relative_path=slide_spec.visual_path,
            native_visual=slide_spec.native_visual,
            left=prototype_manifest.visual_left,
            top=prototype_manifest.visual_top,
            width=prototype_manifest.visual_width,
            height=prototype_manifest.visual_height,
        )
    return slide


def _fill_template_title_placeholder(
    slide,
    layout_manifest: DeckPptxTemplateLayout,
    text: str,
) -> None:
    placeholder = _find_placeholder(slide, layout_manifest.title_placeholder_idx)
    if placeholder is None:
        return
    text_frame = placeholder.text_frame
    text_frame.clear()
    text_frame.text = text


def _fill_template_text_placeholder(
    slide,
    layout_manifest: DeckPptxTemplateLayout,
    slide_spec: SlidesPptxSlide,
) -> None:
    placeholder = _find_placeholder(slide, layout_manifest.text_placeholder_idx)
    if placeholder is None:
        return
    text_frame = placeholder.text_frame
    text_frame.clear()
    if slide_spec.bullets:
        for index, bullet in enumerate(slide_spec.bullets):
            paragraph = (
                text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            )
            level, paragraph_text = _bullet_level_and_text(bullet)
            paragraph.text = paragraph_text
            _apply_native_bullet_format(paragraph, level=level)
    elif slide_spec.body:
        text_frame.text = slide_spec.body


def _fill_template_visual_slot(
    slide,
    layout_manifest: DeckPptxTemplateLayout,
    *,
    deck_path: Path,
    relative_path: str,
    native_visual: Mapping[str, object] | None = None,
    table_model: Mapping[str, object] | None = None,
) -> None:
    placeholder = _find_placeholder(slide, layout_manifest.visual_placeholder_idx)
    if placeholder is None:
        return
    if _add_native_table_within_emu_box(
        slide,
        table_model=table_model,
        left=int(placeholder.left),
        top=int(placeholder.top),
        width=int(placeholder.width),
        height=int(placeholder.height),
    ):
        _remove_shape(placeholder)
        return
    if isinstance(native_visual, Mapping):
        _add_visual_within_emu_box(
            slide,
            deck_path=deck_path,
            relative_path="",
            native_visual=native_visual,
            left=int(placeholder.left),
            top=int(placeholder.top),
            width=int(placeholder.width),
            height=int(placeholder.height),
        )
        _remove_shape(placeholder)
        return
    if not relative_path:
        return
    image_path = _resolve_local_asset_path(deck_path, relative_path)
    if image_path is None:
        return
    if layout_manifest.visual_placeholder_type == 18:
        try:
            placeholder.insert_picture(str(image_path))
            return
        except (AttributeError, ValueError):
            pass
    fitted = _fit_image_path_within_box(
        image_path,
        left=_emu_to_inches(placeholder.left),
        top=_emu_to_inches(placeholder.top),
        width=_emu_to_inches(placeholder.width),
        height=_emu_to_inches(placeholder.height),
    )
    if fitted is None:
        return
    fit_left, fit_top, fit_width, fit_height = fitted
    slide.shapes.add_picture(
        str(image_path),
        _inches(fit_left),
        _inches(fit_top),
        width=_inches(fit_width),
        height=_inches(fit_height),
    )
    _remove_shape(placeholder)


def _populate_slide_number_placeholder(slide, index: int) -> None:
    placeholder = None
    for candidate in slide.placeholders:
        if int(candidate.placeholder_format.type) == 13:
            placeholder = candidate
            break
    if placeholder is None:
        return
    text_frame = placeholder.text_frame
    text_frame.clear()
    text_frame.text = str(index)


def _shape_at_index(slide, shape_index: int | None):
    if shape_index is None:
        return None
    shapes = list(slide.shapes)
    if shape_index < 0 or shape_index >= len(shapes):
        return None
    return shapes[shape_index]


def _set_shape_text(slide, shape_index: int | None, text: str) -> None:
    shape = _shape_at_index(slide, shape_index)
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.text = text


def _find_placeholder(slide, placeholder_idx: int | None):
    if placeholder_idx is None:
        return None
    for placeholder in slide.placeholders:
        if int(placeholder.placeholder_format.idx) == int(placeholder_idx):
            return placeholder
    return None


def _remove_shape(shape) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def _clear_presentation_slides(presentation) -> None:
    slide_id_list = presentation.slides._sldIdLst  # noqa: SLF001
    for slide_id in list(slide_id_list):
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _copy_slide_shapes(source_slide, destination_slide) -> None:
    destination_tree = destination_slide.shapes._spTree  # noqa: SLF001
    for shape in source_slide.shapes:
        destination_tree.insert_element_before(deepcopy(shape.element), "p:extLst")


def _analysis_slide_map(
    slide_analysis: Mapping[str, object] | None,
) -> dict[str, Mapping[str, object]]:
    if not isinstance(slide_analysis, Mapping):
        return {}
    raw_slides = slide_analysis.get("slides")
    if not isinstance(raw_slides, list):
        return {}
    result: dict[str, Mapping[str, object]] = {}
    for slide in raw_slides:
        if not isinstance(slide, Mapping):
            continue
        slide_id = str(slide.get("slideId") or slide.get("slide_id") or "").strip()
        if slide_id:
            result[slide_id] = slide
    return result


def _build_analysis_slide_spec(
    *,
    slide: Slide,
    slide_number: int,
    deck_path: Path,
    analysis_slide: Mapping[str, object],
    asset_dir: Path,
) -> SlidesPptxSlide:
    title = _normalize_text(
        str(analysis_slide.get("titleText") or analysis_slide.get("title_text") or "")
    ) or _html_text(slide.title_html)
    blocks = _sorted_blocks(analysis_slide.get("blocks"))
    canvas_width, canvas_height = _analysis_canvas_size(analysis_slide, blocks)
    visual_group_ids = _visual_exhibit_group_ids(blocks)
    footer_text, footer_block_ids = _extract_footer_text(
        blocks,
        canvas_width=canvas_width,
        canvas_height=canvas_height,
    )
    implication, implication_block_ids = _extract_implication_text(
        blocks,
        canvas_width=canvas_width,
        canvas_height=canvas_height,
    )
    comparison_columns, callout_title, callout_body, comparison_block_ids = (
        _extract_comparison_structure(
            blocks,
            canvas_width=canvas_width,
            canvas_height=canvas_height,
        )
    )
    table_model = _resolve_table_model(blocks)
    table_title, table_title_block_ids = _extract_table_title(
        blocks,
        table_model=table_model,
    )
    skip_block_ids = (
        footer_block_ids
        | implication_block_ids
        | comparison_block_ids
        | table_title_block_ids
    )
    bullets = _extract_bullets_from_analysis(
        analysis_slide,
        blocks,
        skip_block_ids=skip_block_ids,
        blocked_group_ids=visual_group_ids,
    )
    text_blocks = _extract_text_blocks(
        blocks,
        skip_block_ids=skip_block_ids,
        blocked_group_ids=visual_group_ids,
    )
    excluded_visual_bboxes = _block_bboxes_for_ids(
        blocks,
        footer_block_ids | implication_block_ids,
    )
    native_visual = _extract_cards_row_native_visual(
        deck_path=deck_path,
        analysis_slide=analysis_slide,
        blocks=blocks,
        excluded_block_ids=skip_block_ids,
        excluded_bboxes=excluded_visual_bboxes,
    )
    visual_bbox, visual_type = _resolve_visual_bbox(
        analysis_slide,
        blocks,
        excluded_bboxes=excluded_visual_bboxes,
    )
    visual_path = ""
    if native_visual is not None:
        visual_type = "figure"
    elif visual_bbox is not None:
        visual_path = _write_visual_crop(
            deck_path=deck_path,
            asset_dir=asset_dir,
            slide=slide,
            slide_number=slide_number,
            analysis_slide=analysis_slide,
            bbox=visual_bbox,
            visual_type=visual_type,
        )
    body = ""
    if comparison_columns:
        body = text_blocks[0] if text_blocks else ""
        bullets = []
        kind = "comparison_columns"
    else:
        if bullets:
            body = text_blocks[0] if text_blocks else ""
        elif text_blocks:
            if visual_path or native_visual is not None:
                body = text_blocks[0]
            else:
                body = "\n\n".join(text_blocks)
        kind = _resolve_slide_kind(
            bullets=bullets,
            body=body,
            has_visual=bool(visual_path or native_visual),
        )
        if (
            footer_text
            and not visual_path
            and native_visual is None
            and not bullets
            and body
        ):
            kind = "cover_with_footer"
    if bullets and not body and text_blocks:
        body = text_blocks[0]
    density_hint = _classify_density(
        title=title,
        bullets=bullets,
        body="\n\n".join(
            part
            for part in (
                body,
                callout_title,
                callout_body,
                implication,
                footer_text,
                table_title,
            )
            if part
        ),
    )
    visual_layout_hint = _infer_visual_layout_hint(
        deck_path=deck_path,
        analysis_slide=analysis_slide,
        visual_bbox=visual_bbox,
        visual_type=visual_type,
    )
    layout_variant = _select_layout_variant(
        kind=kind,
        title=title,
        bullets=bullets,
        body=body,
        visual_type=visual_type,
        density_hint=density_hint,
        visual_layout_hint=visual_layout_hint,
    )
    return SlidesPptxSlide(
        slide_id=slide.id,
        kind=kind,
        layout_variant=layout_variant,
        density_hint=density_hint,
        title=title or f"Slide {slide_number}",
        body=body,
        bullets=bullets,
        visual_path=visual_path,
        visual_type=visual_type,
        table_model=table_model,
        footer_text=footer_text,
        implication=implication,
        table_title=table_title,
        callout_title=callout_title,
        callout_body=callout_body,
        comparison_columns=comparison_columns,
        native_visual=native_visual,
    )


def _build_html_slide_spec(
    *,
    slide: Slide,
    slide_number: int,
    deck_path: Path,
) -> SlidesPptxSlide:
    title = _html_text(slide.title_html) or f"Slide {slide_number}"
    soup = BeautifulSoup(slide.body_html or "", "html.parser")
    bullets = _unique_texts(
        [
            _normalize_text(item.get_text(" ", strip=True))
            for item in soup.find_all("li")
        ]
    )
    body_parts = _extract_html_body_parts(soup)
    body = "\n\n".join(_unique_texts(body_parts)) if body_parts else ""
    visual_path = ""
    visual_type = ""
    native_visual = None
    image_tag = soup.find("img")
    if image_tag is not None:
        resolved = _resolve_local_asset_path(deck_path, str(image_tag.get("src") or ""))
        if resolved is not None:
            try:
                visual_path = str(resolved.relative_to(deck_path)).replace("\\", "/")
                visual_type = "figure"
            except ValueError:
                visual_path = ""
    if bullets and body:
        body = ""
    density_hint = _classify_density(
        title=title,
        bullets=bullets,
        body=body,
    )
    kind = _resolve_slide_kind(
        bullets=bullets,
        body=body,
        has_visual=bool(visual_path or native_visual),
    )
    layout_variant = _select_layout_variant(
        kind=kind,
        title=title,
        bullets=bullets,
        body=body,
        visual_type=visual_type,
        density_hint=density_hint,
        visual_layout_hint="",
    )
    return SlidesPptxSlide(
        slide_id=slide.id,
        kind=kind,
        layout_variant=layout_variant,
        density_hint=density_hint,
        title=title,
        body=body,
        bullets=bullets,
        visual_path=visual_path,
        visual_type=visual_type,
        native_visual=native_visual,
        table_model=None,
    )


def _extract_html_body_parts(soup: BeautifulSoup) -> list[str]:
    body_parts: list[str] = []
    for paragraph in soup.find_all(["p", "h2", "h3", "blockquote"]):
        text = _normalize_text(paragraph.get_text(" ", strip=True))
        if text:
            body_parts.append(text)
    if body_parts:
        return _unique_texts(body_parts)
    for node in soup.find_all("div"):
        if node.find(
            [
                "blockquote",
                "div",
                "h1",
                "h2",
                "h3",
                "img",
                "li",
                "ol",
                "p",
                "table",
                "ul",
            ]
        ):
            continue
        text = _normalize_text(node.get_text(" ", strip=True))
        if text:
            body_parts.append(text)
    return _unique_texts(body_parts)


def _build_section_header_slide_spec(
    *,
    slide: Slide,
    slide_number: int,
) -> SlidesPptxSlide:
    entries, placeholder = _parse_section_header_entries(
        slide.body_html or slide.full_html
    )
    current_entry = next((entry for entry in entries if entry.is_current), None)
    title = (
        current_entry.label
        if current_entry is not None
        else placeholder or _html_text(slide.title_html) or f"Slide {slide_number}"
    )
    density_hint = (
        "dense"
        if len(entries) + sum(len(entry.subsections) for entry in entries) >= 12
        else "medium"
    )
    return SlidesPptxSlide(
        slide_id=slide.id,
        kind="section_header",
        layout_variant="section_header_agenda",
        density_hint=density_hint,
        title=title,
        body=placeholder or "",
        bullets=[],
        visual_path="",
        visual_type="",
        table_model=None,
    )


def _parse_section_header_entries(
    raw_html: str,
) -> tuple[list[_SectionHeaderEntry], str]:
    if not raw_html:
        return [], ""
    soup = BeautifulSoup(raw_html, "html.parser")
    placeholder = _normalize_text(
        str(
            (
                soup.select_one(".section-header__placeholder").get_text(
                    " ", strip=True
                )
                if soup.select_one(".section-header__placeholder") is not None
                else ""
            )
        )
    )
    entries: list[_SectionHeaderEntry] = []
    for node in soup.select(".section-header__section"):
        label_node = node.select_one(".section-header__section-label")
        label = _normalize_text(
            label_node.get_text(" ", strip=True) if label_node is not None else ""
        )
        if not label:
            continue
        classes = node.get("class") or []
        subsections: list[_SectionHeaderSubsection] = []
        for subsection_node in node.select(".section-header__subsection"):
            subsection_label = _normalize_text(
                subsection_node.get_text(" ", strip=True)
            )
            if not subsection_label:
                continue
            subsection_classes = subsection_node.get("class") or []
            subsections.append(
                _SectionHeaderSubsection(
                    label=subsection_label,
                    is_current="is-current" in subsection_classes,
                )
            )
        entries.append(
            _SectionHeaderEntry(
                label=label,
                is_current="is-current" in classes,
                subsections=subsections,
            )
        )
    return entries, placeholder


def _resolve_slide_kind(
    *,
    bullets: Sequence[str],
    body: str,
    has_visual: bool,
) -> str:
    has_bullets = bool(bullets)
    has_body = bool(body.strip())
    if has_visual and has_bullets:
        return "bullets_visual"
    if has_visual and has_body:
        return "text_visual"
    if has_visual:
        return "visual_only"
    if has_bullets:
        return "bullets_only"
    if has_body:
        return "text_only"
    return "title_only"


def _resolve_layout_variant(
    layout_variant: str,
    *,
    kind: str,
    visual_type: str,
) -> str:
    normalized = str(layout_variant or "").strip()
    if normalized:
        return normalized
    return _select_layout_variant(
        kind=kind,
        title="",
        bullets=[],
        body="",
        visual_type=visual_type,
        density_hint="",
        visual_layout_hint="",
    )


def _select_layout_variant(
    *,
    kind: str,
    title: str,
    bullets: Sequence[str],
    body: str,
    visual_type: str,
    density_hint: str,
    visual_layout_hint: str,
) -> str:
    del title
    if kind == "section_header":
        return "section_header_agenda"
    if kind == "cover_with_footer":
        return "cover_with_footer"
    if kind == "comparison_columns":
        return "comparison_columns"
    if kind == "title_only":
        return "title_only_centered"
    if kind == "bullets_only":
        return "bullets_full_width"
    if kind == "text_only":
        return "text_full_width"
    if kind == "visual_only":
        return "table_focus" if visual_type == "table" else "visual_full_width"
    if kind == "bullets_visual":
        if visual_type == "table":
            return "table_focus"
        if visual_layout_hint == "right":
            return "bullets_visual_right"
        if visual_layout_hint == "bottom":
            return "bullets_visual_bottom"
        if density_hint in {"light", "medium"}:
            return "bullets_visual_right"
        return "bullets_visual_bottom"
    if kind == "text_visual":
        if visual_type == "table":
            return "table_focus"
        if visual_layout_hint == "right":
            return "text_visual_right"
        if visual_layout_hint == "bottom":
            return "text_visual_bottom"
        if density_hint in {"light", "medium"}:
            return "text_visual_right"
        return "text_visual_bottom"
    return "text_full_width"


def _classify_density(
    *,
    title: str,
    bullets: Sequence[str],
    body: str,
) -> str:
    title_score = len(title)
    bullet_count = len(bullets)
    bullet_chars = sum(len(item) for item in bullets)
    max_bullet = max((len(item) for item in bullets), default=0)
    body_chars = len(body)
    score = title_score + body_chars + bullet_chars + (bullet_count * 35) + max_bullet
    if score >= 520:
        return "dense"
    if score >= 260:
        return "medium"
    return "light"


def _infer_visual_layout_hint(
    *,
    deck_path: Path,
    analysis_slide: Mapping[str, object],
    visual_bbox: tuple[float, float, float, float] | None,
    visual_type: str,
) -> str:
    if visual_bbox is None or visual_type == "table":
        return ""
    asset_path = _resolve_local_asset_path(
        deck_path,
        str(analysis_slide.get("assetPath") or analysis_slide.get("asset_path") or ""),
    )
    if asset_path is None:
        return ""
    try:
        with Image.open(asset_path) as image:
            image_width, image_height = image.size
    except (OSError, ValueError):
        return ""
    if image_width <= 0 or image_height <= 0:
        return ""
    left, top, width, height = visual_bbox
    width_ratio = width / float(image_width)
    height_ratio = height / float(image_height)
    crop_ratio = width / height if height > 0 else 0.0
    center_x = (left + (width / 2.0)) / float(image_width)
    center_y = (top + (height / 2.0)) / float(image_height)

    # Prefer a right-hand placement for square/tall visuals that live on the
    # right side of the source slide. Forcing these into a short bottom box
    # makes them shrink dramatically.
    if center_x >= 0.56 and width_ratio <= 0.62:
        return "right"
    if center_x >= 0.52 and crop_ratio <= 1.35 and width_ratio <= 0.68:
        return "right"

    # Reserve bottom placement for genuinely wide visuals and tables-like
    # compositions where horizontal space matters more than height.
    if width_ratio >= 0.78 and crop_ratio >= 1.25:
        return "bottom"
    if center_y >= 0.62 and crop_ratio >= 1.35:
        return "bottom"
    if height_ratio >= 0.72 and width_ratio >= 0.56 and crop_ratio >= 1.2:
        return "bottom"
    if center_x >= 0.53 and width_ratio <= 0.7:
        return "right"
    return ""


def _body_font_size_for_slide(
    base_size: float,
    *,
    layout_variant: str,
    density_hint: str,
) -> float:
    size = float(base_size)
    if density_hint == "dense":
        if layout_variant in {
            "bullets_full_width",
            "text_full_width",
            "bullets_visual_right",
            "text_visual_right",
        }:
            size -= 2.0
        elif layout_variant in {"bullets_visual_bottom", "text_visual_bottom"}:
            size -= 1.25
    elif density_hint == "medium" and layout_variant in {
        "bullets_full_width",
        "bullets_visual_right",
        "text_visual_right",
    }:
        size -= 0.75
    return max(size, 13.0)


def _line_height_for_slide(base_line_height: float, *, density_hint: str) -> float:
    if density_hint == "dense":
        return max(base_line_height - 0.12, 1.0)
    if density_hint == "medium":
        return max(base_line_height - 0.07, 1.0)
    return base_line_height


def _paragraph_spacing_for_slide(*, density_hint: str) -> float:
    if density_hint == "dense":
        return 6.0
    if density_hint == "medium":
        return 8.0
    return 10.0


def _text_margin_for_slide(*, layout_variant: str, density_hint: str) -> float:
    if layout_variant in {
        "bullets_visual_right",
        "text_visual_right",
        "bullets_full_width",
        "text_full_width",
    }:
        if density_hint == "dense":
            return 0.03
        if density_hint == "medium":
            return 0.04
    return 0.06


def _estimate_wrapped_lines(text: str, *, line_capacity: int = 44) -> int:
    normalized = _normalize_text(text)
    if not normalized:
        return 1
    words = normalized.split()
    if not words:
        return 1
    lines = 1
    current = 0
    for word in words:
        word_length = len(word)
        if current == 0:
            current = word_length
            continue
        if current + 1 + word_length <= line_capacity:
            current += 1 + word_length
            continue
        lines += 1
        current = word_length
    return max(1, min(lines, 5))


def _title_layout_for_slide(
    base_size: float, slide_spec: SlidesPptxSlide
) -> dict[str, float]:
    font_size = float(base_size)
    title_length = len(slide_spec.title)
    if title_length >= 160:
        font_size -= 8.5
    elif title_length >= 125:
        font_size -= 6.5
    elif title_length >= 95:
        font_size -= 4.5
    elif title_length >= 72:
        font_size -= 2.75
    if slide_spec.density_hint == "dense":
        font_size -= 0.75
    if slide_spec.layout_variant in {"comparison_columns", "table_focus"}:
        font_size -= 0.5
    title_scale = _repair_hint_float(
        slide_spec,
        "title_scale",
        minimum=0.72,
        maximum=1.0,
        default=1.0,
    )
    font_size *= title_scale
    font_size = max(font_size, 18.0)
    line_capacity = (
        42 if slide_spec.layout_variant in {"comparison_columns", "table_focus"} else 46
    )
    estimated_lines = _estimate_wrapped_lines(
        slide_spec.title, line_capacity=line_capacity
    )
    height = min(max(_TITLE_HEIGHT_IN, 0.52 + (estimated_lines * 0.24)), 1.42)
    if title_scale < 1.0:
        height = min(1.58, height * (1.0 + ((1.0 - title_scale) * 0.7)))
    return {
        "left": _TITLE_LEFT_IN,
        "top": _TITLE_TOP_IN,
        "width": _TITLE_WIDTH_IN,
        "height": height,
        "font_size": font_size,
    }


def _content_top_after_title(
    default_top: float, *, title_layout: Mapping[str, float]
) -> float:
    return max(
        default_top, float(title_layout["top"]) + float(title_layout["height"]) + 0.12
    )


def _implication_banner_box(slide_spec: SlidesPptxSlide) -> dict[str, float] | None:
    if not slide_spec.implication:
        return None
    extra_height = 0.14 if len(slide_spec.implication) >= 165 else 0.0
    banner_scale = _repair_hint_float(
        slide_spec,
        "banner_scale",
        minimum=0.9,
        maximum=1.3,
        default=1.0,
    )
    height = (_BOTTOM_BANNER_HEIGHT_IN + extra_height) * banner_scale
    return {
        "left": _BOTTOM_BANNER_LEFT_IN,
        "top": _SLIDE_HEIGHT_IN - height - 0.44,
        "width": _BOTTOM_BANNER_WIDTH_IN,
        "height": height,
    }


def _repair_hint_float(
    slide_spec: SlidesPptxSlide,
    key: str,
    *,
    minimum: float,
    maximum: float,
    default: float,
) -> float:
    value = (
        slide_spec.repair_hints.get(key)
        if isinstance(slide_spec.repair_hints, Mapping)
        else None
    )
    if not isinstance(value, (int, float)):
        return default
    return min(maximum, max(minimum, float(value)))


def _repair_hint_choice(
    slide_spec: SlidesPptxSlide,
    key: str,
    *,
    allowed: set[str],
    default: str,
) -> str:
    value = ""
    if isinstance(slide_spec.repair_hints, Mapping):
        value = str(slide_spec.repair_hints.get(key) or "").strip().lower()
    return value if value in allowed else default


def _apply_visual_box_repairs(
    slide_spec: SlidesPptxSlide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    content_bottom: float,
) -> tuple[float, float, float, float]:
    visual_scale = _repair_hint_float(
        slide_spec,
        "visual_scale",
        minimum=0.85,
        maximum=1.2,
        default=1.0,
    )
    visual_anchor = _repair_hint_choice(
        slide_spec,
        "visual_anchor",
        allowed={"top", "center", "bottom"},
        default="center",
    )
    if visual_scale == 1.0:
        return left, top, width, height
    scaled_width = min((_SLIDE_WIDTH_IN - 0.8), width * visual_scale)
    scaled_height = min(max(1.0, content_bottom - 0.3), height * visual_scale)
    width_delta = scaled_width - width
    height_delta = scaled_height - height
    adjusted_left = left - (width_delta / 2.0)
    if visual_anchor == "top":
        adjusted_top = top
    elif visual_anchor == "bottom":
        adjusted_top = top - height_delta
    else:
        adjusted_top = top - (height_delta / 2.0)
    adjusted_left = min(max(0.38, adjusted_left), _SLIDE_WIDTH_IN - 0.38 - scaled_width)
    adjusted_top = min(max(0.92, adjusted_top), content_bottom - scaled_height)
    return adjusted_left, adjusted_top, scaled_width, scaled_height


def _content_bottom_limit(banner_box: Mapping[str, float] | None) -> float:
    if banner_box is None:
        return 6.92
    return float(banner_box["top"]) - 0.12


def _text_box_height_for_text(
    text: str,
    *,
    base_height: float,
    max_height: float,
) -> float:
    estimated_lines = _estimate_wrapped_lines(text, line_capacity=92)
    return min(max_height, max(base_height, 0.24 + (estimated_lines * 0.22)))


def _sorted_blocks(raw_blocks: object) -> list[Mapping[str, object]]:
    if not isinstance(raw_blocks, list):
        return []
    blocks = [block for block in raw_blocks if isinstance(block, Mapping)]
    return sorted(blocks, key=block_sort_key)


def _analysis_block_id(block: Mapping[str, object]) -> str:
    return str(block.get("blockId") or block.get("block_id") or "").strip()


def _analysis_block_text(block: Mapping[str, object]) -> str:
    return _normalize_text(str(block.get("text") or ""))


def _analysis_block_items(block: Mapping[str, object]) -> list[str]:
    raw_items = block.get("items") if isinstance(block.get("items"), list) else []
    if raw_items:
        return _unique_texts([str(item or "") for item in raw_items])
    raw_text = str(block.get("text") or "")
    return _unique_texts(
        segment.strip()
        for segment in raw_text.splitlines()
        if str(segment or "").strip()
    )


def _analysis_block_group_kind(block: Mapping[str, object]) -> str | None:
    return normalize_optional_string(
        block.get("groupKind") if "groupKind" in block else block.get("group_kind")
    )


def _analysis_block_parent_id(block: Mapping[str, object]) -> str | None:
    return normalize_optional_string(
        block.get("parentId") if "parentId" in block else block.get("parent_id")
    )


def _analysis_canvas_size(
    analysis_slide: Mapping[str, object],
    blocks: Sequence[Mapping[str, object]],
) -> tuple[float, float]:
    max_right = 0.0
    max_bottom = 0.0
    for bbox in (_bbox_from_mapping(block.get("bbox")) for block in blocks):
        if bbox is None:
            continue
        max_right = max(max_right, bbox[0] + bbox[2])
        max_bottom = max(max_bottom, bbox[1] + bbox[3])
    raw_figure_regions = analysis_slide.get("figureRegions") or analysis_slide.get(
        "figure_regions"
    )
    if isinstance(raw_figure_regions, list):
        for bbox in (_bbox_from_mapping(region) for region in raw_figure_regions):
            if bbox is None:
                continue
            max_right = max(max_right, bbox[0] + bbox[2])
            max_bottom = max(max_bottom, bbox[1] + bbox[3])
    return max(max_right, 1.0), max(max_bottom, 1.0)


def _group_blocks_by_group_id(
    blocks: Sequence[Mapping[str, object]],
) -> dict[str, list[Mapping[str, object]]]:
    groups: dict[str, list[Mapping[str, object]]] = {}
    for block in blocks:
        group_id = _analysis_block_group_id(block)
        if group_id is None:
            continue
        groups.setdefault(group_id, []).append(block)
    return groups


def _analysis_block_visual_lines(
    block: Mapping[str, object],
) -> list[dict[str, object]]:
    raw_lines = (
        block.get("visualLines")
        if isinstance(block.get("visualLines"), list)
        else (
            block.get("visual_lines")
            if isinstance(block.get("visual_lines"), list)
            else []
        )
    )
    normalized: list[dict[str, object]] = []
    for line in raw_lines:
        if not isinstance(line, Mapping):
            continue
        text = clean_ocr_text(str(line.get("text") or ""))
        bbox = _bbox_from_mapping(line.get("bbox"))
        if not text or bbox is None:
            continue
        normalized.append(
            {
                "text": text,
                "bbox": bbox,
                "confidence": line.get("confidence"),
            }
        )
    return normalized


def _line_center_x(line: Mapping[str, object]) -> float:
    bbox = line.get("bbox")
    assert isinstance(bbox, tuple)
    return float(bbox[0]) + (float(bbox[2]) / 2.0)


def _line_center_y(line: Mapping[str, object]) -> float:
    bbox = line.get("bbox")
    assert isinstance(bbox, tuple)
    return float(bbox[1]) + (float(bbox[3]) / 2.0)


def _collect_visual_line_candidates(
    *,
    blocks: Sequence[Mapping[str, object]],
    figure_block: Mapping[str, object],
    excluded_block_ids: set[str],
    excluded_bboxes: Sequence[tuple[float, float, float, float]],
) -> list[dict[str, object]]:
    figure_bbox = _bbox_from_mapping(figure_block.get("bbox"))
    if figure_bbox is None:
        return []
    candidates = _analysis_block_visual_lines(figure_block)
    if not candidates:
        figure_group_id = _analysis_block_group_id(figure_block)
        for block in blocks:
            if _analysis_block_id(block) in excluded_block_ids:
                continue
            if _analysis_block_id(block) == _analysis_block_id(figure_block):
                continue
            bbox = _bbox_from_mapping(block.get("bbox"))
            text = _analysis_block_text(block)
            if bbox is None or not text:
                continue
            if figure_group_id is not None:
                if _analysis_block_group_id(block) != figure_group_id:
                    continue
            elif _horizontal_overlap(figure_bbox, bbox) <= 0:
                continue
            candidates.append(
                {
                    "text": text,
                    "bbox": bbox,
                    "confidence": block.get("confidence"),
                }
            )
    filtered: list[dict[str, object]] = []
    for candidate in candidates:
        bbox = candidate["bbox"]
        if _horizontal_overlap(figure_bbox, bbox) <= 0:
            continue
        if bbox[1] < figure_bbox[1] or (bbox[1] + bbox[3]) > (
            figure_bbox[1] + figure_bbox[3]
        ):
            continue
        if any(
            _horizontal_overlap(bbox, excluded_bbox) > 0
            and bbox[1] >= excluded_bbox[1] - 8.0
            for excluded_bbox in excluded_bboxes
        ):
            continue
        filtered.append(candidate)
    return sorted(filtered, key=lambda item: (item["bbox"][1], item["bbox"][0]))


def _cluster_card_line_candidates(
    lines: Sequence[Mapping[str, object]],
    *,
    figure_bbox: tuple[float, float, float, float],
) -> list[list[dict[str, object]]]:
    if len(lines) < 4:
        return []
    threshold = max(36.0, figure_bbox[2] * 0.11)
    sorted_lines = sorted(lines, key=_line_center_x)
    clusters: list[list[dict[str, object]]] = []
    centers: list[float] = []
    for line in sorted_lines:
        center_x = _line_center_x(line)
        if not centers:
            centers.append(center_x)
            clusters.append([dict(line)])
            continue
        nearest_index = min(
            range(len(centers)),
            key=lambda index: abs(centers[index] - center_x),
        )
        if abs(centers[nearest_index] - center_x) <= threshold:
            clusters[nearest_index].append(dict(line))
            centers[nearest_index] = (
                centers[nearest_index] * float(len(clusters[nearest_index]) - 1)
                + center_x
            ) / float(len(clusters[nearest_index]))
        else:
            centers.append(center_x)
            clusters.append([dict(line)])
    if len(clusters) < 2 or len(clusters) > 4:
        return []
    clusters = sorted(
        [
            sorted(cluster, key=lambda item: (item["bbox"][1], item["bbox"][0]))
            for cluster in clusters
        ],
        key=lambda cluster: _line_center_x(cluster[0]),
    )
    cluster_bboxes = [
        _union_bboxes([line["bbox"] for line in cluster]) for cluster in clusters
    ]
    figure_height = max(1.0, figure_bbox[3])
    top_values = [bbox[1] for bbox in cluster_bboxes]
    bottom_values = [bbox[1] + bbox[3] for bbox in cluster_bboxes]
    if max(top_values) - min(top_values) > (figure_height * 0.14):
        return []
    if max(bottom_values) - min(bottom_values) > (figure_height * 0.18):
        return []
    if any(bbox[3] > (figure_height * 0.72) for bbox in cluster_bboxes):
        return []
    if any(len(cluster) < 3 for cluster in clusters):
        return []
    return clusters


def _split_card_title_and_body_lines(
    lines: Sequence[Mapping[str, object]],
) -> tuple[list[str], list[str]]:
    if not lines:
        return [], []
    ordered = sorted(lines, key=lambda item: (item["bbox"][1], item["bbox"][0]))
    if len(ordered) == 1:
        return [_normalize_text(str(ordered[0]["text"]))], []
    heights = [float(item["bbox"][3]) for item in ordered]
    baseline_gap = max(8.0, float(median(heights)) * 0.85)
    split_index = 1
    for index in range(min(len(ordered) - 1, 4)):
        current_bbox = ordered[index]["bbox"]
        next_bbox = ordered[index + 1]["bbox"]
        gap = float(next_bbox[1]) - (float(current_bbox[1]) + float(current_bbox[3]))
        if gap >= baseline_gap:
            split_index = index + 1
            break
    else:
        first_two = _normalize_text(" ".join(str(item["text"]) for item in ordered[:2]))
        if len(ordered) >= 4 and len(first_two) <= 64:
            split_index = 2
    title_lines = [_normalize_text(str(item["text"])) for item in ordered[:split_index]]
    body_lines = [_normalize_text(str(item["text"])) for item in ordered[split_index:]]
    return [line for line in title_lines if line], [line for line in body_lines if line]


def _card_body_payload_from_lines(lines: Sequence[str]) -> tuple[str, list[str]]:
    normalized_lines = [
        _normalize_text(line) for line in lines if _normalize_text(line)
    ]
    if not normalized_lines:
        return "", []
    bullet_like_count = sum(
        1 for line in normalized_lines if line.startswith(("-", "•", "·", "▪", "◦"))
    )
    if bullet_like_count >= max(1, len(normalized_lines) // 2):
        items = [_normalize_bullet_text(line) for line in normalized_lines]
        items = [item for item in items if item]
        return "", items
    return _normalize_text(" ".join(normalized_lines)), []


def _sample_card_style(
    image: Image.Image,
    *,
    card_bbox: tuple[float, float, float, float],
) -> tuple[tuple[int, int, int], bool]:
    working = image.convert("RGB")
    left = max(0, int(card_bbox[0]))
    top = max(0, int(card_bbox[1]))
    right = min(working.width, int(card_bbox[0] + card_bbox[2]))
    bottom = min(working.height, int(card_bbox[1] + card_bbox[3]))
    if right - left < 8 or bottom - top < 8:
        return (95, 117, 143), False

    def _collect_nonwhite_pixels(
        region: tuple[int, int, int, int],
    ) -> list[tuple[int, int, int]]:
        cropped = working.crop(region)
        pixels = list(cropped.getdata())
        return [
            pixel
            for pixel in pixels
            if min(pixel) < 235 and max(pixel) - min(pixel) > 18
        ]

    top_band = _collect_nonwhite_pixels(
        (left + 2, top + 1, right - 2, min(bottom, top + 8))
    )
    left_band = _collect_nonwhite_pixels(
        (left + 1, top + 4, min(right, left + 6), bottom - 4)
    )
    right_band = _collect_nonwhite_pixels(
        (max(left, right - 6), top + 4, right - 1, bottom - 4)
    )
    sample = top_band or left_band or right_band
    if not sample:
        return (95, 117, 143), False
    channel_medians = tuple(
        int(median([pixel[channel] for pixel in sample])) for channel in range(3)
    )
    full_border = len(left_band) >= max(8, len(top_band) // 5) and len(
        right_band
    ) >= max(8, len(top_band) // 5)
    return channel_medians, full_border


def _gap_has_connector(
    image: Image.Image,
    *,
    left_card_bbox: tuple[float, float, float, float],
    right_card_bbox: tuple[float, float, float, float],
) -> bool:
    working = image.convert("RGB")
    gap_left = int(left_card_bbox[0] + left_card_bbox[2] + 4)
    gap_right = int(right_card_bbox[0] - 4)
    if gap_right - gap_left < 12:
        return False
    top = int(
        max(left_card_bbox[1], right_card_bbox[1])
        + min(left_card_bbox[3], right_card_bbox[3]) * 0.34
    )
    bottom = int(
        min(
            left_card_bbox[1] + left_card_bbox[3],
            right_card_bbox[1] + right_card_bbox[3],
        )
        - min(left_card_bbox[3], right_card_bbox[3]) * 0.24
    )
    if bottom <= top:
        return False
    cropped = working.crop((gap_left, top, gap_right, bottom))
    pixels = list(cropped.getdata())
    if not pixels:
        return False
    dark_pixels = sum(1 for pixel in pixels if min(pixel) < 200)
    return dark_pixels / float(len(pixels)) >= 0.015


def _extract_cards_row_native_visual(
    *,
    deck_path: Path,
    analysis_slide: Mapping[str, object],
    blocks: Sequence[Mapping[str, object]],
    excluded_block_ids: set[str],
    excluded_bboxes: Sequence[tuple[float, float, float, float]],
) -> dict[str, object] | None:
    asset_path = _resolve_local_asset_path(
        deck_path,
        str(analysis_slide.get("assetPath") or analysis_slide.get("asset_path") or ""),
    )
    if asset_path is None:
        return None
    figure_blocks = [
        block
        for block in blocks
        if _analysis_block_type(block) == "figure"
        and _analysis_block_render_mode(block) != "ignore"
    ]
    if not figure_blocks:
        return None
    try:
        with Image.open(asset_path) as image:
            for figure_block in sorted(
                figure_blocks,
                key=lambda block: (
                    (_bbox_from_mapping(block.get("bbox")) or (0.0, 0.0, 0.0, 0.0))[2]
                    * (_bbox_from_mapping(block.get("bbox")) or (0.0, 0.0, 0.0, 0.0))[3]
                ),
                reverse=True,
            ):
                figure_bbox = _bbox_from_mapping(figure_block.get("bbox"))
                if figure_bbox is None:
                    continue
                group_kind = _analysis_block_group_kind(figure_block)
                if group_kind and group_kind in _COMPARISON_GROUP_KINDS:
                    continue
                line_candidates = _collect_visual_line_candidates(
                    blocks=blocks,
                    figure_block=figure_block,
                    excluded_block_ids=excluded_block_ids,
                    excluded_bboxes=excluded_bboxes,
                )
                clusters = _cluster_card_line_candidates(
                    line_candidates,
                    figure_bbox=figure_bbox,
                )
                if not clusters:
                    continue
                cluster_bboxes = [
                    _union_bboxes([line["bbox"] for line in cluster])
                    for cluster in clusters
                ]
                figure_left, figure_top, figure_width, figure_height = figure_bbox
                row_top = max(
                    figure_top,
                    min(bbox[1] for bbox in cluster_bboxes) - (figure_height * 0.05),
                )
                row_bottom = min(
                    figure_top + figure_height,
                    max(bbox[1] + bbox[3] for bbox in cluster_bboxes)
                    + (figure_height * 0.06),
                )
                if row_bottom - row_top < 40:
                    continue
                cards: list[dict[str, object]] = []
                centers = [bbox[0] + (bbox[2] / 2.0) for bbox in cluster_bboxes]
                for index, (cluster, cluster_bbox) in enumerate(
                    zip(clusters, cluster_bboxes, strict=False)
                ):
                    if index == 0:
                        left_boundary = figure_left + (figure_width * 0.03)
                    else:
                        left_boundary = (centers[index - 1] + centers[index]) / 2.0
                    if index == len(clusters) - 1:
                        right_boundary = (
                            figure_left + figure_width - (figure_width * 0.03)
                        )
                    else:
                        right_boundary = (centers[index] + centers[index + 1]) / 2.0
                    gap = max(8.0, right_boundary - left_boundary)
                    shrink = min(figure_width * 0.012, gap * 0.09)
                    card_left = max(figure_left, left_boundary + shrink)
                    card_right = min(
                        figure_left + figure_width, right_boundary - shrink
                    )
                    card_bbox = (
                        card_left,
                        row_top,
                        max(24.0, card_right - card_left),
                        max(32.0, row_bottom - row_top),
                    )
                    title_lines, body_lines = _split_card_title_and_body_lines(cluster)
                    title = _normalize_text(" ".join(title_lines))
                    body, items = _card_body_payload_from_lines(body_lines)
                    if not title or not (body or items):
                        cards = []
                        break
                    accent_rgb, full_border = _sample_card_style(
                        image,
                        card_bbox=card_bbox,
                    )
                    cards.append(
                        {
                            "left": (card_bbox[0] - figure_left)
                            / max(1.0, figure_width),
                            "top": (card_bbox[1] - figure_top)
                            / max(1.0, figure_height),
                            "width": card_bbox[2] / max(1.0, figure_width),
                            "height": card_bbox[3] / max(1.0, figure_height),
                            "title": title,
                            "body": body,
                            "items": items,
                            "accentRgb": list(accent_rgb),
                            "fullBorder": full_border,
                            "align": "center",
                        }
                    )
                if len(cards) < 2:
                    continue
                connectors = bool(
                    group_kind
                    and group_kind in {"flowchart", "flow_diagram", "process_flow"}
                )
                if not connectors:
                    connectors = all(
                        _gap_has_connector(
                            image,
                            left_card_bbox=(
                                figure_left + (cards[index]["left"] * figure_width),
                                figure_top + (cards[index]["top"] * figure_height),
                                cards[index]["width"] * figure_width,
                                cards[index]["height"] * figure_height,
                            ),
                            right_card_bbox=(
                                figure_left + (cards[index + 1]["left"] * figure_width),
                                figure_top + (cards[index + 1]["top"] * figure_height),
                                cards[index + 1]["width"] * figure_width,
                                cards[index + 1]["height"] * figure_height,
                            ),
                        )
                        for index in range(len(cards) - 1)
                    )
                return {
                    "kind": "cards_row",
                    "cards": cards,
                    "connectors": connectors,
                    "sourceBbox": {
                        "x": figure_left,
                        "y": figure_top,
                        "w": figure_width,
                        "h": figure_height,
                    },
                }
    except (OSError, ValueError):
        return None
    return None


def _visual_exhibit_group_ids(blocks: Sequence[Mapping[str, object]]) -> set[str]:
    visual_group_ids: set[str] = set()
    for group_id, group_blocks in _group_blocks_by_group_id(blocks).items():
        has_visual = any(
            _analysis_block_type(block)
            in (_VISUAL_BLOCK_TYPES | _SEMANTIC_VISUAL_BLOCK_TYPES)
            for block in group_blocks
        )
        has_grouped_visual = any(
            _analysis_block_render_mode(block) == _GROUP_AS_IMAGE_RENDER_MODE
            for block in group_blocks
        )
        group_kinds = {
            str(_analysis_block_group_kind(block) or "").strip().lower()
            for block in group_blocks
            if _analysis_block_group_kind(block)
        }
        if has_visual and has_grouped_visual:
            visual_group_ids.add(group_id)
            continue
        if (
            has_grouped_visual
            and group_kinds
            and not group_kinds & _COMPARISON_GROUP_KINDS
        ):
            visual_group_ids.add(group_id)
    return visual_group_ids


def _looks_like_short_heading(text: str) -> bool:
    normalized = _normalize_text(text)
    if not normalized:
        return False
    words = normalized.split()
    return len(words) <= 8 and len(normalized) <= 72


def _split_callout_parts_from_block(
    block: Mapping[str, object] | None,
) -> tuple[str, str]:
    if block is None:
        return "", ""
    items = _analysis_block_items(block)
    if len(items) >= 2 and _looks_like_short_heading(items[0]):
        return items[0], _normalize_text(" ".join(items[1:]))
    return "", _analysis_block_text(block)


def _extract_footer_text(
    blocks: Sequence[Mapping[str, object]],
    *,
    canvas_width: float,
    canvas_height: float,
) -> tuple[str, set[str]]:
    best_block: Mapping[str, object] | None = None
    best_score = 0.0
    for block in blocks:
        if _analysis_block_render_mode(block) == _GROUP_AS_IMAGE_RENDER_MODE:
            continue
        block_type = _analysis_block_type(block)
        if block_type not in {"footer_meta", "body_text", "group_label"}:
            continue
        text = _analysis_block_text(block)
        bbox = _bbox_from_mapping(block.get("bbox"))
        if not text or bbox is None:
            continue
        bottom_ratio = (bbox[1] + bbox[3]) / max(canvas_height, 1.0)
        width_ratio = bbox[2] / max(canvas_width, 1.0)
        if block_type != "footer_meta" and (bottom_ratio < 0.82 or width_ratio > 0.58):
            continue
        score = bottom_ratio + (0.15 if block_type == "footer_meta" else 0.0)
        if score > best_score:
            best_score = score
            best_block = block
    if best_block is None:
        return "", set()
    return _analysis_block_text(best_block), {_analysis_block_id(best_block)}


def _extract_implication_text(
    blocks: Sequence[Mapping[str, object]],
    *,
    canvas_width: float,
    canvas_height: float,
) -> tuple[str, set[str]]:
    best_block: Mapping[str, object] | None = None
    best_score = 0.0
    for block in blocks:
        if _analysis_block_render_mode(block) == _GROUP_AS_IMAGE_RENDER_MODE:
            continue
        block_type = _analysis_block_type(block)
        text = _analysis_block_text(block)
        bbox = _bbox_from_mapping(block.get("bbox"))
        if not text or bbox is None:
            continue
        bottom_ratio = (bbox[1] + bbox[3]) / max(canvas_height, 1.0)
        width_ratio = bbox[2] / max(canvas_width, 1.0)
        explicit_implication = text.upper().startswith(_IMPLICATION_PREFIX)
        geometry_implication = (
            block_type in {"body_text", "group_label", "implication_banner"}
            and bottom_ratio >= 0.78
            and width_ratio >= 0.68
            and len(_normalize_text(text)) >= 28
        )
        if not (
            block_type == "implication_banner"
            or explicit_implication
            or geometry_implication
        ):
            continue
        score = bottom_ratio
        if block_type == "implication_banner":
            score += 0.2
        elif explicit_implication:
            score += 0.1
        elif geometry_implication:
            score += 0.05
        if score > best_score:
            best_score = score
            best_block = block
    if best_block is None:
        return "", set()
    return _analysis_block_text(best_block), {_analysis_block_id(best_block)}


def _extract_table_title(
    blocks: Sequence[Mapping[str, object]],
    *,
    table_model: Mapping[str, object] | None,
) -> tuple[str, set[str]]:
    table_block = next(
        (
            block
            for block in blocks
            if _analysis_block_type(block) == "table"
            and _analysis_block_render_mode(block) != _GROUP_AS_IMAGE_RENDER_MODE
        ),
        None,
    )
    if table_block is None:
        return "", set()
    explicit_title = next(
        (
            block
            for block in blocks
            if _analysis_block_type(block) == "table_title"
            and _analysis_block_render_mode(block) != _GROUP_AS_IMAGE_RENDER_MODE
            and _analysis_block_text(block)
        ),
        None,
    )
    if explicit_title is not None:
        return _analysis_block_text(explicit_title), {
            _analysis_block_id(explicit_title)
        }
    table_bbox = _bbox_from_mapping(table_block.get("bbox"))
    if table_bbox is not None:
        candidate_blocks = [
            block
            for block in blocks
            if _analysis_block_render_mode(block) != _GROUP_AS_IMAGE_RENDER_MODE
            and _analysis_block_id(block) != _analysis_block_id(table_block)
            and _analysis_block_type(block) in {"group_label", "body_text"}
        ]
        best_candidate: Mapping[str, object] | None = None
        best_gap = float("inf")
        for block in candidate_blocks:
            bbox = _bbox_from_mapping(block.get("bbox"))
            text = _analysis_block_text(block)
            if bbox is None or not text or not _looks_like_short_heading(text):
                continue
            gap = table_bbox[1] - (bbox[1] + bbox[3])
            horizontal_overlap = min(
                table_bbox[0] + table_bbox[2],
                bbox[0] + bbox[2],
            ) - max(table_bbox[0], bbox[0])
            if gap < -6.0 or gap > 120.0 or horizontal_overlap <= 0:
                continue
            if gap < best_gap:
                best_gap = gap
                best_candidate = block
        if best_candidate is not None:
            return _analysis_block_text(best_candidate), {
                _analysis_block_id(best_candidate)
            }
    table_text = _analysis_block_text(table_block)
    raw_table_text = str(table_block.get("text") or "")
    if raw_table_text:
        first_line = clean_ocr_text(
            str(raw_table_text.splitlines()[0] if raw_table_text.splitlines() else "")
        )
        if first_line and _looks_like_short_heading(first_line):
            return first_line, set()
    if table_text:
        first_line = clean_ocr_text(table_text)
        if first_line and _looks_like_short_heading(first_line):
            return first_line, set()
    return "", set()


def _extract_comparison_structure(
    blocks: Sequence[Mapping[str, object]],
    *,
    canvas_width: float,
    canvas_height: float,
) -> tuple[list[dict[str, object]], str, str, set[str]]:
    comparison_candidates: list[Mapping[str, object]] = []
    callout_title_block: Mapping[str, object] | None = None
    callout_body_block: Mapping[str, object] | None = None

    for block in blocks:
        group_kind = str(_analysis_block_group_kind(block) or "").strip().lower()
        is_comparison_group = group_kind in _COMPARISON_GROUP_KINDS
        if (
            _analysis_block_render_mode(block) == _GROUP_AS_IMAGE_RENDER_MODE
            and not is_comparison_group
        ):
            continue
        block_type = _analysis_block_type(block)
        text = _analysis_block_text(block)
        bbox = _bbox_from_mapping(block.get("bbox"))
        if not text or bbox is None:
            continue
        center_x = bbox[0] + (bbox[2] / 2.0)
        if group_kind == "callout" and _looks_like_short_heading(text):
            callout_title_block = block
            continue
        if block_type == "callout_banner" or group_kind == "callout":
            callout_body_block = block
            continue
        if not is_comparison_group:
            continue
        if (
            block_type == "group_label"
            and center_x >= canvas_width * 0.32
            and center_x <= canvas_width * 0.68
            and bbox[1] >= canvas_height * 0.6
            and _looks_like_short_heading(text)
        ):
            callout_title_block = block
            continue
        comparison_candidates.append(block)

    comparison_bounds = [
        _bbox_from_mapping(block.get("bbox"))
        for block in comparison_candidates
        if _bbox_from_mapping(block.get("bbox")) is not None
    ]
    if comparison_bounds:
        min_y = min(bbox[1] for bbox in comparison_bounds)
        max_y = max(bbox[1] + bbox[3] for bbox in comparison_bounds)
        callout_body_top = (
            (
                _bbox_from_mapping(callout_body_block.get("bbox"))
                or (0.0, max_y + 120.0, 0.0, 0.0)
            )[1]
            if callout_body_block is not None
            else max_y + 120.0
        )
        if callout_title_block is None:
            best_candidate: Mapping[str, object] | None = None
            best_gap = float("inf")
            callout_body_bbox = (
                _bbox_from_mapping(callout_body_block.get("bbox"))
                if callout_body_block is not None
                else None
            )
            for block in blocks:
                if _analysis_block_render_mode(block) == _GROUP_AS_IMAGE_RENDER_MODE:
                    continue
                if block in comparison_candidates or block is callout_body_block:
                    continue
                if _analysis_block_type(block) not in {
                    "group_label",
                    "body_text",
                    "table_title",
                }:
                    continue
                text = _analysis_block_text(block)
                bbox = _bbox_from_mapping(block.get("bbox"))
                if not text or bbox is None or not _looks_like_short_heading(text):
                    continue
                center_x = bbox[0] + (bbox[2] / 2.0)
                horizontal_overlap = (
                    _horizontal_overlap(callout_body_bbox, bbox)
                    if callout_body_bbox is not None
                    else bbox[2]
                )
                gap = callout_body_top - (bbox[1] + bbox[3])
                if center_x < canvas_width * 0.32 or center_x > canvas_width * 0.68:
                    continue
                if bbox[1] < (max_y - 30.0) or gap < -10.0 or gap > 140.0:
                    continue
                if callout_body_block is not None and horizontal_overlap <= 0:
                    continue
                if gap < best_gap:
                    best_gap = gap
                    best_candidate = block
            callout_title_block = best_candidate
        callout_title_y = (
            (
                _bbox_from_mapping(callout_title_block.get("bbox"))
                or (0.0, max_y, 0.0, 0.0)
            )[1]
            if callout_title_block is not None
            else callout_body_top
        )
        for block in blocks:
            if block in comparison_candidates:
                continue
            group_kind = str(_analysis_block_group_kind(block) or "").strip().lower()
            if (
                _analysis_block_render_mode(block) == _GROUP_AS_IMAGE_RENDER_MODE
                and group_kind not in _COMPARISON_GROUP_KINDS
            ):
                continue
            if _analysis_block_type(block) not in {
                "body_text",
                "bullet_item",
                "group_label",
            }:
                continue
            text = _analysis_block_text(block)
            bbox = _bbox_from_mapping(block.get("bbox"))
            if not text or bbox is None:
                continue
            if not str(text).lstrip().startswith(("-", "•", "·", "▪", "◦")):
                continue
            if bbox[1] < (min_y - 40.0) or (bbox[1] + bbox[3]) > (
                callout_title_y - 20.0
            ):
                continue
            comparison_candidates.append(block)

    if len(comparison_candidates) < 4:
        return [], "", "", set()

    columns = [
        [
            block
            for block in comparison_candidates
            if (_bbox_from_mapping(block.get("bbox")) or (0.0, 0.0, 0.0, 0.0))[0]
            + ((_bbox_from_mapping(block.get("bbox")) or (0.0, 0.0, 0.0, 0.0))[2] / 2.0)
            < (canvas_width / 2.0)
        ],
        [
            block
            for block in comparison_candidates
            if (_bbox_from_mapping(block.get("bbox")) or (0.0, 0.0, 0.0, 0.0))[0]
            + ((_bbox_from_mapping(block.get("bbox")) or (0.0, 0.0, 0.0, 0.0))[2] / 2.0)
            >= (canvas_width / 2.0)
        ],
    ]
    if not columns[0] or not columns[1]:
        return [], "", "", set()

    used_block_ids: set[str] = set()
    serialized_columns: list[dict[str, object]] = []
    for column_blocks in columns:
        sorted_column = sorted(column_blocks, key=block_sort_key)
        header_text = ""
        bullet_texts: list[str] = []
        if sorted_column:
            first_text = _analysis_block_text(sorted_column[0])
            if _looks_like_short_heading(first_text):
                header_text = first_text
                used_block_ids.add(_analysis_block_id(sorted_column[0]))
                remaining_blocks = sorted_column[1:]
            else:
                remaining_blocks = sorted_column
        else:
            remaining_blocks = []
        for block in remaining_blocks:
            text = _normalize_bullet_text(_analysis_block_text(block))
            if not text:
                continue
            bullet_texts.append(text)
            used_block_ids.add(_analysis_block_id(block))
        if not header_text and bullet_texts:
            header_text = bullet_texts.pop(0)
        serialized_columns.append(
            {
                "title": header_text,
                "bullets": _unique_bullets(bullet_texts),
            }
        )

    if not all(
        column.get("title") or column.get("bullets") for column in serialized_columns
    ):
        return [], "", "", set()
    callout_title = (
        _analysis_block_text(callout_title_block)
        if callout_title_block is not None
        else ""
    )
    implicit_callout_title, callout_body = _split_callout_parts_from_block(
        callout_body_block
    )
    if not callout_title:
        callout_title = implicit_callout_title
    if callout_title_block is not None and callout_title:
        used_block_ids.add(_analysis_block_id(callout_title_block))
    if callout_body_block is not None and (callout_title or callout_body):
        used_block_ids.add(_analysis_block_id(callout_body_block))
    return serialized_columns, callout_title, callout_body, used_block_ids


def _extract_bullets_from_analysis(
    analysis_slide: Mapping[str, object],
    blocks: Sequence[Mapping[str, object]],
    *,
    skip_block_ids: set[str] | None = None,
    blocked_group_ids: set[str] | None = None,
) -> list[str]:
    skip_block_ids = skip_block_ids or set()
    blocked_group_ids = blocked_group_ids or set()
    block_bullets: list[str] = []
    grouped_visual_texts = {
        _normalize_bullet_text(_analysis_block_text(block)).lower()
        for block in blocks
        if _analysis_block_group_id(block) in blocked_group_ids
        and _analysis_block_text(block)
    }
    for block in blocks:
        block_id = _analysis_block_id(block)
        if block_id in skip_block_ids:
            continue
        if _analysis_block_group_id(block) in blocked_group_ids:
            continue
        if _analysis_block_render_mode(block) == _GROUP_AS_IMAGE_RENDER_MODE:
            continue
        if _analysis_block_type(block) not in BULLET_BLOCK_TYPES:
            continue
        block_bullets.extend(_analysis_block_bullet_items(block))
    normalized_blocks = _unique_bullets(block_bullets)
    raw_bullets = analysis_slide.get("bulletTexts") or analysis_slide.get(
        "bullet_texts"
    )
    if isinstance(raw_bullets, list):
        normalized = _unique_texts(
            [_normalize_bullet_text(str(item or "")) for item in raw_bullets]
        )
        if normalized_blocks:
            seen = {_normalize_bullet_text(item).lower() for item in normalized_blocks}
            extras = [
                item
                for item in normalized
                if item.lower() not in seen and item.lower() not in grouped_visual_texts
            ]
            return normalized_blocks + extras
        filtered_normalized = [
            item for item in normalized if item.lower() not in grouped_visual_texts
        ]
        if filtered_normalized:
            return filtered_normalized
    return normalized_blocks


def _extract_text_blocks(
    blocks: Sequence[Mapping[str, object]],
    *,
    skip_block_ids: set[str] | None = None,
    blocked_group_ids: set[str] | None = None,
) -> list[str]:
    skip_block_ids = skip_block_ids or set()
    blocked_group_ids = blocked_group_ids or set()
    texts = [
        _analysis_block_text(block)
        for block in blocks
        if _analysis_block_id(block) not in skip_block_ids
        and _analysis_block_group_id(block) not in blocked_group_ids
        if _analysis_block_render_mode(block) != _GROUP_AS_IMAGE_RENDER_MODE
        and _analysis_block_type(block)
        in {
            "text",
            "body_text",
            "metric",
            "exhibit_label",
            "group_label",
            "footer_meta",
            "implication_banner",
            "callout_banner",
            "table_title",
        }
    ]
    return _unique_texts(texts)


def _block_bboxes_for_ids(
    blocks: Sequence[Mapping[str, object]],
    block_ids: set[str],
) -> list[tuple[float, float, float, float]]:
    if not block_ids:
        return []
    return [
        bbox
        for block in blocks
        if _analysis_block_id(block) in block_ids
        for bbox in [_bbox_from_mapping(block.get("bbox"))]
        if bbox is not None
    ]


def _trim_visual_bbox_for_excluded_bands(
    visual_bbox: tuple[float, float, float, float] | None,
    *,
    excluded_bboxes: Sequence[tuple[float, float, float, float]] | None = None,
) -> tuple[float, float, float, float] | None:
    if visual_bbox is None:
        return None
    bands = list(excluded_bboxes or [])
    if not bands:
        return visual_bbox
    left, top, width, height = visual_bbox
    right = left + width
    bottom = top + height
    trim_candidates: list[float] = []
    for band_left, band_top, band_width, band_height in bands:
        band_right = band_left + band_width
        band_bottom = band_top + band_height
        overlap_width = max(0.0, min(right, band_right) - max(left, band_left))
        if overlap_width <= 0.0:
            continue
        overlap_ratio = overlap_width / max(1.0, min(width, band_width))
        if overlap_ratio < 0.38 and (band_width / max(1.0, width)) < 0.55:
            continue
        if band_top <= top + (height * 0.28):
            continue
        if band_top >= bottom:
            continue
        trim_candidates.append(band_top)
        if band_bottom > bottom and band_top > top:
            trim_candidates.append(band_top)
    if not trim_candidates:
        return visual_bbox
    gap = 6.0
    trimmed_bottom = min(trim_candidates) - gap
    if trimmed_bottom <= top + 32.0:
        return visual_bbox
    return left, top, width, max(1.0, trimmed_bottom - top)


def _resolve_table_model(
    blocks: Sequence[Mapping[str, object]],
) -> dict[str, object] | None:
    table_blocks = [
        block
        for block in blocks
        if _analysis_block_type(block) == "table"
        and _analysis_block_render_mode(block) != _GROUP_AS_IMAGE_RENDER_MODE
    ]
    if not table_blocks:
        return None
    best_block = max(
        table_blocks,
        key=lambda block: (
            (
                _bbox_from_mapping(block.get("bbox"))[2]
                * _bbox_from_mapping(block.get("bbox"))[3]
            )
            if _bbox_from_mapping(block.get("bbox")) is not None
            else 0.0
        ),
    )
    raw_model = best_block.get("tableModel")
    if not isinstance(raw_model, Mapping):
        raw_model = best_block.get("table_model")
    return dict(raw_model) if isinstance(raw_model, Mapping) else None


def _resolve_visual_bbox(
    analysis_slide: Mapping[str, object],
    blocks: Sequence[Mapping[str, object]],
    *,
    excluded_bboxes: Sequence[tuple[float, float, float, float]] | None = None,
) -> tuple[tuple[float, float, float, float] | None, str]:
    grouped_visual_bboxes: list[tuple[float, float, float, float]] = []
    groups: dict[str, list[Mapping[str, object]]] = {}
    for block in blocks:
        group_id = _analysis_block_group_id(block)
        if group_id is None:
            continue
        groups.setdefault(group_id, []).append(block)
    visual_group_ids = _visual_exhibit_group_ids(blocks)
    for group_id, group_blocks in groups.items():
        if group_id not in visual_group_ids:
            continue
        bboxes = [
            bbox
            for bbox in (
                _bbox_from_mapping(block.get("bbox")) for block in group_blocks
            )
            if bbox is not None
        ]
        if not bboxes:
            continue
        grouped_visual_bboxes.append(_union_bboxes(bboxes))
    if grouped_visual_bboxes:
        visual_bbox = _union_bboxes(grouped_visual_bboxes)
        return (
            _trim_visual_bbox_for_excluded_bands(
                visual_bbox,
                excluded_bboxes=excluded_bboxes,
            ),
            "figure",
        )
    visual_blocks = [
        block
        for block in blocks
        if _analysis_block_type(block)
        in (_VISUAL_BLOCK_TYPES | _SEMANTIC_VISUAL_BLOCK_TYPES)
    ]
    if visual_blocks:
        visual_type = (
            "table"
            if any(_analysis_block_type(block) == "table" for block in visual_blocks)
            else "figure"
        )
        bboxes = [
            bbox
            for bbox in (
                _bbox_from_mapping(block.get("bbox")) for block in visual_blocks
            )
            if bbox is not None
        ]
        if bboxes:
            visual_bbox = _union_bboxes(bboxes)
            return (
                _trim_visual_bbox_for_excluded_bands(
                    visual_bbox,
                    excluded_bboxes=excluded_bboxes,
                ),
                visual_type,
            )
    raw_figure_regions = analysis_slide.get("figureRegions") or analysis_slide.get(
        "figure_regions"
    )
    if isinstance(raw_figure_regions, list):
        bboxes = [
            bbox
            for bbox in (_bbox_from_mapping(region) for region in raw_figure_regions)
            if bbox is not None
        ]
        if bboxes:
            visual_bbox = _union_bboxes(bboxes)
            return (
                _trim_visual_bbox_for_excluded_bands(
                    visual_bbox,
                    excluded_bboxes=excluded_bboxes,
                ),
                "figure",
            )
    return None, ""


def _write_visual_crop(
    *,
    deck_path: Path,
    asset_dir: Path,
    slide: Slide,
    slide_number: int,
    analysis_slide: Mapping[str, object],
    bbox: tuple[float, float, float, float],
    visual_type: str,
) -> str:
    asset_path = _resolve_local_asset_path(
        deck_path,
        str(analysis_slide.get("assetPath") or analysis_slide.get("asset_path") or ""),
    )
    if asset_path is None or not asset_path.exists():
        return ""
    try:
        with Image.open(asset_path) as image:
            working = image.convert("RGB")
            left, top, width, height = bbox
            pad = 12
            crop_left = max(0, int(left) - pad)
            crop_top = max(0, int(top) - pad)
            crop_right = min(working.width, int(left + width) + pad)
            crop_bottom = min(working.height, int(top + height) + pad)
            if crop_right <= crop_left or crop_bottom <= crop_top:
                return ""
            cropped = working.crop((crop_left, crop_top, crop_right, crop_bottom))
            trimmed = _trim_white_image(cropped)
            output_name = f"{slide_number:03d}-{_safe_stem(slide.id)}-{visual_type or 'visual'}.png"
            output_path = asset_dir / output_name
            trimmed.save(output_path, format="PNG")
    except (OSError, ValueError):
        return ""
    return str(output_path.relative_to(deck_path)).replace("\\", "/")


def _bbox_from_mapping(value: object) -> tuple[float, float, float, float] | None:
    if not isinstance(value, Mapping):
        return None
    try:
        x = float(value["x"])
        y = float(value["y"])
        w = float(value["w"])
        h = float(value["h"])
    except (KeyError, TypeError, ValueError):
        return None
    if w <= 0 or h <= 0:
        return None
    return x, y, w, h


def _union_bboxes(
    bboxes: Sequence[tuple[float, float, float, float]],
) -> tuple[float, float, float, float]:
    left = min(bbox[0] for bbox in bboxes)
    top = min(bbox[1] for bbox in bboxes)
    right = max(bbox[0] + bbox[2] for bbox in bboxes)
    bottom = max(bbox[1] + bbox[3] for bbox in bboxes)
    return left, top, max(1.0, right - left), max(1.0, bottom - top)


def _horizontal_overlap(
    left_bbox: tuple[float, float, float, float],
    right_bbox: tuple[float, float, float, float],
) -> float:
    return max(
        0.0,
        min(left_bbox[0] + left_bbox[2], right_bbox[0] + right_bbox[2])
        - max(left_bbox[0], right_bbox[0]),
    )


def _safe_stem(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9]+", "-", str(value or "").strip().lower()).strip("-")
    return cleaned or "slide"


def _normalize_text(text: str) -> str:
    normalized = _SPACE_RE.sub(" ", str(text or "").replace("\n", " ")).strip()
    normalized = _NON_WORD_JOIN_RE.sub(r"\1", normalized)
    return normalized


def _normalize_bullet_text(text: str) -> str:
    normalized = _normalize_text(text)
    if not normalized:
        return ""
    return _LEADING_BULLET_MARKER_RE.sub("", normalized).strip()


def _analysis_block_type(block: Mapping[str, object]) -> str:
    return normalize_block_type(block.get("type"))


def _analysis_block_render_mode(block: Mapping[str, object]) -> str:
    render_mode = normalize_render_mode(
        block.get("renderMode") if "renderMode" in block else block.get("render_mode")
    )
    if render_mode is not None:
        return render_mode
    return "native"


def _analysis_block_group_id(block: Mapping[str, object]) -> str | None:
    return normalize_optional_string(
        block.get("groupId") if "groupId" in block else block.get("group_id")
    )


def _analysis_block_list_level(block: Mapping[str, object]) -> int:
    level = normalize_list_level(
        block.get("listLevel") if "listLevel" in block else block.get("list_level")
    )
    return level if level is not None else 0


def _prefixed_bullet_text(*, text: str, level: int) -> str:
    if level <= 0:
        return text
    return ("\t" * level) + text


def _analysis_block_bullet_items(block: Mapping[str, object]) -> list[str]:
    raw_items = block.get("items") if isinstance(block.get("items"), list) else []
    normalized_items = _unique_texts(
        [_normalize_bullet_text(str(item or "")) for item in raw_items]
    )
    level = _analysis_block_list_level(block)
    if normalized_items:
        return [
            _prefixed_bullet_text(text=item, level=level) for item in normalized_items
        ]
    fallback = _normalize_bullet_text(str(block.get("text") or ""))
    return [_prefixed_bullet_text(text=fallback, level=level)] if fallback else []


def _unique_texts(values: Sequence[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for value in values:
        normalized = _normalize_text(value)
        if not normalized:
            continue
        if normalized in seen:
            continue
        seen.add(normalized)
        result.append(normalized)
    return result


def _unique_bullets(values: Sequence[str]) -> list[str]:
    seen: set[tuple[int, str]] = set()
    result: list[str] = []
    for value in values:
        level, paragraph_text = _bullet_level_and_text(value)
        normalized = _normalize_bullet_text(str(value).lstrip("\t"))
        if not normalized:
            continue
        key = (level, normalized.lower())
        if key in seen:
            continue
        seen.add(key)
        result.append(("\t" * level) + normalized)
    return result


def _html_text(value: str) -> str:
    if not value:
        return ""
    return _normalize_text(
        BeautifulSoup(value, "html.parser").get_text(" ", strip=True)
    )


def _resolve_local_asset_path(deck_path: Path, raw_path: str) -> Path | None:
    normalized = str(raw_path or "").strip()
    if not normalized:
        return None
    parsed = urlparse(normalized)
    if parsed.scheme or parsed.netloc:
        return None
    path_text = unquote(parsed.path or "")
    if not path_text:
        return None
    relative_candidates = _local_asset_relative_candidates(path_text)
    deck_root = deck_path.resolve()
    for relative in relative_candidates:
        candidate = (deck_root / relative).resolve()
        try:
            candidate.relative_to(deck_root)
        except ValueError:
            continue
        if candidate.exists():
            return candidate
    absolute_path = Path(path_text)
    if absolute_path.is_absolute():
        candidate = absolute_path.resolve()
        try:
            candidate.relative_to(deck_root)
        except ValueError:
            return None
        return candidate if candidate.exists() else None
    return None


def _local_asset_relative_candidates(path_text: str) -> list[Path]:
    path = Path(path_text.lstrip("/"))
    parts = path.parts
    if not parts:
        return []
    candidates: list[Path] = []
    if (
        len(parts) >= 5
        and parts[0] == "slides"
        and parts[1] == "deck"
        and parts[3] == "assets"
    ):
        candidates.append(Path("assets", *parts[4:]))
    if parts[0] == "assets":
        candidates.append(Path(*parts))
    candidates.append(Path(*parts))
    unique_candidates: list[Path] = []
    seen: set[str] = set()
    for candidate in candidates:
        key = candidate.as_posix()
        if key in seen:
            continue
        seen.add(key)
        unique_candidates.append(candidate)
    return unique_candidates


def _trim_white_image(image: Image.Image, *, padding_px: int = 12) -> Image.Image:
    working = image.convert("RGB")
    background = Image.new("RGB", working.size, "white")
    diff = ImageChops.difference(working, background)
    bbox = diff.getbbox()
    if bbox is None:
        return working
    left, top, right, bottom = bbox
    left = max(0, left - padding_px)
    top = max(0, top - padding_px)
    right = min(working.width, right + padding_px)
    bottom = min(working.height, bottom + padding_px)
    return working.crop((left, top, right, bottom))


def _read_report_text(value: object) -> str:
    return str(value or "").strip()


def _read_report_body(value: object) -> str:
    if isinstance(value, list):
        return "\n\n".join(_read_text_list(value))
    return _read_report_text(value)


def _read_report_mapping(value: object) -> Mapping[str, object]:
    return value if isinstance(value, Mapping) else {}


def _read_report_mapping_list(value: object) -> list[Mapping[str, object]]:
    if not isinstance(value, list):
        return []
    return [item for item in value if isinstance(item, Mapping)]


def _normalize_report_visual_path(raw_path: str, *, deck_path: Path) -> str:
    normalized = str(raw_path or "").strip()
    if not normalized:
        return ""
    parsed = urlparse(normalized)
    if parsed.scheme or parsed.netloc:
        return ""
    path_text = parsed.path.lstrip("/")
    if not path_text:
        return ""
    path_obj = Path(path_text)
    if path_obj.is_absolute():
        try:
            return str(path_obj.resolve().relative_to(deck_path.resolve())).replace(
                "\\", "/"
            )
        except ValueError:
            return ""
    return path_text.replace("\\", "/")


def _normalize_report_visual_type(raw_type: str, *, has_table_model: bool) -> str:
    normalized = str(raw_type or "").strip().lower()
    if has_table_model:
        return "table"
    if normalized in {"chart", "figure", "image", "photo"}:
        return "figure"
    if normalized in {"table", "matrix"}:
        return "table"
    if normalized in {"product_collage", "hero_image", "moodboard"}:
        return "figure"
    return normalized


def _read_text_list(value: object) -> list[str]:
    if not isinstance(value, list):
        return []
    return [text for text in (str(item or "").strip() for item in value) if text]


def _load_template_presentation(presentation_cls, template_key: str):
    template_path = (
        Path(__file__).resolve().parent / "pptx_templates" / f"{template_key}.pptx"
    )
    if template_path.exists():
        return presentation_cls(str(template_path))
    return presentation_cls()


def _add_visual(
    slide,
    *,
    deck_path: Path,
    relative_path: str,
    native_visual: Mapping[str, object] | None = None,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str | None = None,
    body_size: float = 12.0,
    text_rgb: tuple[int, int, int] = (17, 24, 39),
) -> None:
    if isinstance(native_visual, Mapping):
        if _add_native_visual(
            slide,
            native_visual=native_visual,
            left=left,
            top=top,
            width=width,
            height=height,
            font_name=font_name or "Aptos",
            body_size=body_size,
            text_rgb=text_rgb,
        ):
            return
    image_path = _resolve_local_asset_path(deck_path, relative_path)
    if image_path is None:
        return
    fitted = _fit_image_path_within_box(
        image_path,
        left=left,
        top=top,
        width=width,
        height=height,
    )
    if fitted is None:
        return
    fit_left, fit_top, fit_width, fit_height = fitted
    slide.shapes.add_picture(
        str(image_path),
        _inches(fit_left),
        _inches(fit_top),
        width=_inches(fit_width),
        height=_inches(fit_height),
    )


def _add_visual_within_emu_box(
    slide,
    *,
    deck_path: Path,
    relative_path: str,
    native_visual: Mapping[str, object] | None = None,
    left: int,
    top: int,
    width: int,
    height: int,
) -> None:
    if isinstance(native_visual, Mapping) and width > 0 and height > 0:
        if _add_native_visual(
            slide,
            native_visual=native_visual,
            left=_emu_to_inches(left),
            top=_emu_to_inches(top),
            width=_emu_to_inches(width),
            height=_emu_to_inches(height),
            font_name="Aptos",
            body_size=12.0,
            text_rgb=(17, 24, 39),
        ):
            return
    image_path = _resolve_local_asset_path(deck_path, relative_path)
    if image_path is None or width <= 0 or height <= 0:
        return
    fitted = _fit_image_path_within_box(
        image_path,
        left=_emu_to_inches(left),
        top=_emu_to_inches(top),
        width=_emu_to_inches(width),
        height=_emu_to_inches(height),
    )
    if fitted is None:
        return
    fit_left, fit_top, fit_width, fit_height = fitted
    slide.shapes.add_picture(
        str(image_path),
        _inches(fit_left),
        _inches(fit_top),
        width=_inches(fit_width),
        height=_inches(fit_height),
    )


def _fit_image_path_within_box(
    image_path: Path,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
) -> tuple[float, float, float, float] | None:
    try:
        with Image.open(image_path) as image:
            img_width, img_height = image.size
    except (OSError, ValueError):
        return None
    if img_width <= 0 or img_height <= 0:
        return None
    box_ratio = width / height
    image_ratio = float(img_width) / float(img_height)
    if image_ratio >= box_ratio:
        scaled_width = width
        scaled_height = width / image_ratio
        offset_x = 0.0
        offset_y = (height - scaled_height) / 2.0
    else:
        scaled_height = height
        scaled_width = height * image_ratio
        offset_x = (width - scaled_width) / 2.0
        offset_y = 0.0
    return left + offset_x, top + offset_y, scaled_width, scaled_height


def _add_bordered_card(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
    full_border: bool,
) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb((255, 255, 255))
    shape.line.color.rgb = _rgb(accent_rgb if full_border else (198, 205, 214))
    shape.line.width = _inches(0.014 if full_border else 0.01)
    if not full_border:
        _add_filled_rectangle(
            slide,
            left=left,
            top=top,
            width=width,
            height=min(0.04, height * 0.03),
            fill_rgb=accent_rgb,
        )


def _add_connector_arrow(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_rgb: tuple[int, int, int],
) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    if width <= 0.08 or height <= 0.08:
        return
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_rgb)
    shape.line.fill.background()


def _add_native_visual(
    slide,
    *,
    native_visual: Mapping[str, object],
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    body_size: float,
    text_rgb: tuple[int, int, int],
) -> bool:
    visual_kind = str(native_visual.get("kind") or "").strip()
    if visual_kind == "launch_product_tiles":
        return _add_launch_product_tiles(
            slide,
            native_visual=native_visual,
            left=left,
            top=top,
            width=width,
            height=height,
            font_name=font_name,
            body_size=body_size,
            text_rgb=text_rgb,
        )
    if visual_kind != "cards_row":
        return False
    raw_cards = native_visual.get("cards")
    if not isinstance(raw_cards, list):
        return False
    cards = [card for card in raw_cards if isinstance(card, Mapping)]
    if len(cards) < 2:
        return False
    for card in cards:
        try:
            card_left = left + (float(card.get("left") or 0.0) * width)
            card_top = top + (float(card.get("top") or 0.0) * height)
            card_width = max(0.4, float(card.get("width") or 0.0) * width)
            card_height = max(0.6, float(card.get("height") or 0.0) * height)
        except (TypeError, ValueError):
            return False
        accent_rgb_raw = card.get("accentRgb")
        accent_rgb = (
            tuple(int(value) for value in accent_rgb_raw[:3])
            if isinstance(accent_rgb_raw, list) and len(accent_rgb_raw) >= 3
            else (95, 117, 143)
        )
        full_border = bool(card.get("fullBorder"))
        _add_bordered_card(
            slide,
            left=card_left,
            top=card_top,
            width=card_width,
            height=card_height,
            accent_rgb=accent_rgb,
            full_border=full_border,
        )
        title = _normalize_text(str(card.get("title") or ""))
        body = _normalize_text(str(card.get("body") or ""))
        items = _read_text_list(card.get("items"))
        title_height = min(
            0.72,
            max(0.34, 0.24 + (0.18 * _estimate_wrapped_lines(title, line_capacity=22))),
        )
        _add_textbox(
            slide,
            text=title,
            left=card_left + 0.08,
            top=card_top + 0.11,
            width=max(0.2, card_width - 0.16),
            height=title_height,
            font_name=font_name,
            font_size=max(body_size + 0.2, 12.0),
            color_rgb=text_rgb,
            bold=True,
            alignment=None,
            margin=0.02,
        )
        content_top = card_top + title_height + 0.18
        content_height = max(0.28, (card_top + card_height) - content_top - 0.08)
        if items:
            _add_bullet_box(
                slide,
                bullets=items,
                left=card_left + 0.08,
                top=content_top,
                width=max(0.2, card_width - 0.16),
                height=content_height,
                font_name=font_name,
                body_size=max(body_size - 0.9, 10.5),
                line_height=1.02,
                paragraph_spacing=4.0,
                text_rgb=text_rgb,
                margin=0.01,
            )
        elif body:
            _add_textbox(
                slide,
                text=body,
                left=card_left + 0.1,
                top=content_top,
                width=max(0.2, card_width - 0.2),
                height=content_height,
                font_name=font_name,
                font_size=max(body_size - 0.55, 10.8),
                color_rgb=text_rgb,
                margin=0.01,
            )
    if bool(native_visual.get("connectors")) and len(cards) >= 2:
        for index in range(len(cards) - 1):
            try:
                left_card = cards[index]
                right_card = cards[index + 1]
                left_card_right = left + (
                    (
                        float(left_card.get("left") or 0.0)
                        + float(left_card.get("width") or 0.0)
                    )
                    * width
                )
                right_card_left = left + (float(right_card.get("left") or 0.0) * width)
                left_card_top = top + (float(left_card.get("top") or 0.0) * height)
                left_card_height = float(left_card.get("height") or 0.0) * height
                right_card_top = top + (float(right_card.get("top") or 0.0) * height)
                right_card_height = float(right_card.get("height") or 0.0) * height
            except (TypeError, ValueError):
                continue
            gap_width = right_card_left - left_card_right
            arrow_left = left_card_right + 0.05
            arrow_width = max(0.14, gap_width - 0.1)
            arrow_mid = max(
                left_card_top + (left_card_height / 2.0),
                right_card_top + (right_card_height / 2.0),
            )
            _add_connector_arrow(
                slide,
                left=arrow_left,
                top=arrow_mid - 0.1,
                width=arrow_width,
                height=0.2,
                fill_rgb=(145, 153, 166),
            )
    return True


def _add_launch_product_tiles(
    slide,
    *,
    native_visual: Mapping[str, object],
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    body_size: float,
    text_rgb: tuple[int, int, int],
) -> bool:
    raw_tiles = native_visual.get("tiles")
    if not isinstance(raw_tiles, list):
        return False
    tiles = [tile for tile in raw_tiles if isinstance(tile, Mapping)]
    if len(tiles) < 2:
        return False

    accent_rgb = (44, 78, 116)
    muted_text_rgb = (92, 101, 112)
    neutral_border_rgb = (214, 220, 227)
    neutral_fill_rgb = (255, 255, 255)
    neutral_rule_rgb = (230, 234, 239)
    neutral_chip_fill_rgb = (243, 245, 247)

    tile_gap = 0.18 if len(tiles) <= 3 else 0.14
    if len(tiles) == 4:
        columns = 2
        rows = 2
    else:
        columns = min(3, len(tiles))
        rows = 1
    tile_width = max(0.8, (width - (tile_gap * (columns - 1))) / columns)
    tile_height = max(0.9, (height - (tile_gap * (rows - 1))) / rows)

    for index, tile in enumerate(tiles):
        row = index // columns
        column = index % columns
        tile_left = left + (column * (tile_width + tile_gap))
        tile_top = top + (row * (tile_height + tile_gap))
        badge = _normalize_text(str(tile.get("badge") or ""))
        _add_tile_card(
            slide,
            left=tile_left,
            top=tile_top,
            width=tile_width,
            height=tile_height,
            fill_rgb=neutral_fill_rgb,
            border_rgb=neutral_border_rgb,
        )
        _add_filled_rectangle(
            slide,
            left=tile_left,
            top=tile_top,
            width=tile_width,
            height=min(0.11, tile_height * 0.08),
            fill_rgb=accent_rgb if badge else neutral_rule_rgb,
        )

        brand = _normalize_text(str(tile.get("brand") or tile.get("eyebrow") or ""))
        product = _normalize_text(str(tile.get("product") or tile.get("title") or ""))
        body = _normalize_text(str(tile.get("body") or tile.get("note") or ""))
        tags = _read_text_list(tile.get("tags"))

        current_top = tile_top + 0.18
        if brand:
            _add_textbox(
                slide,
                text=brand.upper(),
                left=tile_left + 0.12,
                top=current_top,
                width=max(0.2, tile_width - 0.24),
                height=0.22,
                font_name=font_name,
                font_size=max(body_size - 1.9, 9.0),
                color_rgb=muted_text_rgb,
                bold=True,
                margin=0.01,
            )
            current_top += 0.24
        if badge:
            _add_label_chip(
                slide,
                text=badge.upper(),
                left=max(tile_left + 0.12, tile_left + tile_width - 1.08),
                top=tile_top + 0.17,
                width=0.9,
                height=0.22,
                fill_rgb=accent_rgb,
                text_rgb=(255, 255, 255),
                font_name=font_name,
                font_size=max(body_size - 2.0, 8.8),
            )
        title_height = min(
            0.66,
            max(
                0.34, 0.22 + (0.15 * _estimate_wrapped_lines(product, line_capacity=18))
            ),
        )
        _add_textbox(
            slide,
            text=product,
            left=tile_left + 0.12,
            top=current_top,
            width=max(0.2, tile_width - 0.24),
            height=title_height,
            font_name=font_name,
            font_size=max(body_size + 0.2, 12.0),
            color_rgb=text_rgb,
            bold=True,
            margin=0.01,
        )
        current_top += title_height + 0.06
        body_height = max(0.42, tile_height - 0.9)
        _add_textbox(
            slide,
            text=body,
            left=tile_left + 0.12,
            top=current_top,
            width=max(0.2, tile_width - 0.24),
            height=min(body_height, 0.62),
            font_name=font_name,
            font_size=max(body_size - 0.75, 10.5),
            color_rgb=text_rgb,
            margin=0.01,
        )
        if tags:
            chip_top = tile_top + tile_height - 0.34
            chip_left = tile_left + 0.12
            max_chip_width = tile_width - 0.24
            for tag in tags[:3]:
                chip_width = min(
                    0.95,
                    max(0.42, 0.18 + (0.055 * len(str(tag)))),
                )
                if chip_left + chip_width > tile_left + max_chip_width:
                    break
                _add_label_chip(
                    slide,
                    text=tag,
                    left=chip_left,
                    top=chip_top,
                    width=chip_width,
                    height=0.22,
                    fill_rgb=neutral_chip_fill_rgb,
                    text_rgb=muted_text_rgb,
                    font_name=font_name,
                    font_size=max(body_size - 2.0, 8.8),
                )
                chip_left += chip_width + 0.06
    return True


def _add_section_header_agenda(
    slide,
    *,
    deck_path: Path,
    slide_spec: SlidesPptxSlide,
    font_name: str,
    title_size: float,
    body_size: float,
    text_rgb: tuple[int, int, int],
) -> None:
    slide_path = deck_path / slide_spec.slide_id
    raw_html = ""
    if slide_path.exists():
        raw_html = slide_path.read_text(encoding="utf-8")
    entries, placeholder = _parse_section_header_entries(raw_html)
    fallback_text = placeholder or slide_spec.body or slide_spec.title
    if not entries:
        _add_textbox(
            slide,
            text=fallback_text,
            left=_FULL_TEXT_LEFT_IN,
            top=_FULL_TEXT_TOP_IN,
            width=_FULL_TEXT_WIDTH_IN,
            height=_FULL_TEXT_HEIGHT_IN,
            font_name=font_name,
            font_size=max(body_size, 16.0),
            color_rgb=text_rgb,
            margin=0.04,
        )
        return

    accent_rgb = (0, 83, 186)
    total_lines = len(entries) + sum(
        len(entry.subsections) for entry in entries if entry.is_current
    )
    if total_lines >= 14:
        current_size = max(title_size - 10.0, 22.0)
        section_size = max(body_size - 4.0, 18.0)
        subsection_size = max(body_size - 9.0, 12.0)
    elif total_lines >= 10:
        current_size = max(title_size - 8.0, 24.0)
        section_size = max(body_size - 3.0, 19.0)
        subsection_size = max(body_size - 8.0, 13.0)
    else:
        current_size = max(title_size - 6.0, 26.0)
        section_size = max(body_size - 2.0, 20.0)
        subsection_size = max(body_size - 7.0, 14.0)

    left = 0.34
    label_left = 0.58
    subsection_left = 0.82
    content_width = 12.2
    y = 0.18
    section_gap = 0.08
    subsection_gap = 0.03

    for entry in entries:
        group_top = y
        label_height = 0.46 if entry.is_current else 0.36
        _add_textbox(
            slide,
            text=entry.label,
            left=label_left,
            top=y,
            width=content_width,
            height=label_height,
            font_name=font_name,
            font_size=current_size if entry.is_current else section_size,
            color_rgb=accent_rgb if entry.is_current else text_rgb,
            bold=True,
            margin=0.01,
        )
        y += label_height
        if entry.is_current:
            for subsection in entry.subsections:
                subsection_height = 0.24
                _add_textbox(
                    slide,
                    text=subsection.label,
                    left=subsection_left,
                    top=y,
                    width=content_width - (subsection_left - label_left),
                    height=subsection_height,
                    font_name=font_name,
                    font_size=subsection_size,
                    color_rgb=accent_rgb if subsection.is_current else text_rgb,
                    bold=subsection.is_current,
                    margin=0.0,
                )
                y += subsection_height + subsection_gap
            _add_filled_rectangle(
                slide,
                left=left,
                top=group_top + 0.03,
                width=0.04,
                height=max(0.28, y - group_top - subsection_gap),
                fill_rgb=accent_rgb,
            )
        y += section_gap


def _add_filled_rectangle(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_rgb: tuple[int, int, int],
) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_rgb)
    shape.line.fill.background()


def _add_tile_card(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_rgb: tuple[int, int, int],
    border_rgb: tuple[int, int, int],
) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_rgb)
    shape.line.color.rgb = _rgb(border_rgb)
    shape.line.width = _inches(0.012)


def _add_label_chip(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_rgb: tuple[int, int, int],
    text_rgb: tuple[int, int, int],
    font_name: str,
    font_size: float,
) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_rgb)
    shape.line.fill.background()
    _add_textbox(
        slide,
        text=text,
        left=left,
        top=top + 0.005,
        width=width,
        height=height - 0.01,
        font_name=font_name,
        font_size=font_size,
        color_rgb=text_rgb,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        margin=0.0,
    )


def _add_callout_box(
    slide,
    *,
    title: str,
    body: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    body_size: float,
    text_rgb: tuple[int, int, int],
) -> None:
    _add_filled_rectangle(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        fill_rgb=(245, 247, 250),
    )
    cursor_top = top + 0.06
    if title:
        _add_textbox(
            slide,
            text=title,
            left=left + 0.14,
            top=cursor_top,
            width=width - 0.28,
            height=0.26,
            font_name=font_name,
            font_size=max(body_size - 0.3, 12.0),
            color_rgb=text_rgb,
            bold=True,
            margin=0.01,
        )
        cursor_top += 0.28
    if body:
        _add_textbox(
            slide,
            text=body,
            left=left + 0.14,
            top=cursor_top,
            width=width - 0.28,
            height=max(0.3, (top + height) - cursor_top - 0.08),
            font_name=font_name,
            font_size=max(body_size - 0.7, 11.5),
            color_rgb=text_rgb,
            margin=0.01,
        )


def _add_bottom_banner(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    font_size: float,
    text_rgb: tuple[int, int, int],
) -> None:
    _add_filled_rectangle(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        fill_rgb=(240, 243, 247),
    )
    _add_textbox(
        slide,
        text=text,
        left=left + 0.14,
        top=top + 0.06,
        width=width - 0.28,
        height=max(0.3, height - 0.12),
        font_name=font_name,
        font_size=font_size,
        color_rgb=text_rgb,
        margin=0.01,
        lead_prefixes=(_IMPLICATION_PREFIX,),
    )


def _coerce_table_int(value: object) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return 0


def _table_model_is_native_ready(table_model: Mapping[str, object] | None) -> bool:
    if not isinstance(table_model, Mapping):
        return False
    if bool(table_model.get("has_merged_cells") or table_model.get("hasMergedCells")):
        return False
    row_count = _coerce_table_int(
        table_model.get("row_count") or table_model.get("rowCount")
    )
    column_count = _coerce_table_int(
        table_model.get("column_count") or table_model.get("columnCount")
    )
    confidence = table_model.get("confidence")
    try:
        confidence_value = float(confidence)
    except (TypeError, ValueError):
        confidence_value = 0.0
    return (
        row_count >= 1
        and column_count >= 1
        and confidence_value >= _NATIVE_TABLE_MIN_CONFIDENCE
        and not _table_model_is_structurally_suspicious(table_model)
    )


def _table_model_rows(
    table_model: Mapping[str, object],
) -> list[list[Mapping[str, object]]]:
    raw_rows = table_model.get("rows")
    if not isinstance(raw_rows, list):
        return []
    rows: list[list[Mapping[str, object]]] = []
    for raw_row in raw_rows:
        if not isinstance(raw_row, Mapping):
            continue
        raw_cells = raw_row.get("cells")
        if not isinstance(raw_cells, list):
            continue
        cells = [cell for cell in raw_cells if isinstance(cell, Mapping)]
        if cells:
            rows.append(cells)
    return rows


def _table_model_is_structurally_suspicious(
    table_model: Mapping[str, object] | None,
) -> bool:
    if not isinstance(table_model, Mapping):
        return True
    rows = _table_model_rows(table_model)
    if not rows:
        return True
    column_count = max((len(row) for row in rows), default=0)
    if column_count < 1:
        return True

    text_rows = [
        [clean_ocr_text(str(cell.get("text") or "")) for cell in row] for row in rows
    ]
    nonempty_counts = [
        sum(1 for text in row if clean_ocr_text(text)) for row in text_rows
    ]

    if text_rows and column_count >= 4:
        first_nonempty = nonempty_counts[0]
        first_row_texts = [text for text in text_rows[0] if clean_ocr_text(text)]
        if (
            len(text_rows) >= 2
            and first_nonempty <= max(1, column_count // 3)
            and nonempty_counts[1] >= max(2, column_count - 1)
            and any(len(text) >= 24 for text in first_row_texts)
        ):
            return True

    try:
        header_rows = int(
            table_model.get("header_rows") or table_model.get("headerRows") or 0
        )
    except (TypeError, ValueError):
        header_rows = 0
    if header_rows >= 1 and text_rows:
        header_fragments = 0
        for text in text_rows[0]:
            cleaned = clean_ocr_text(text)
            if not cleaned:
                continue
            if cleaned.endswith(",") or cleaned.startswith(("(", ")", "≤", "≥", "%")):
                header_fragments += 1
        if column_count >= 5 and header_fragments >= 2:
            return True

    return False


def _table_model_column_widths(
    table_model: Mapping[str, object],
    *,
    column_count: int,
) -> list[float]:
    raw_widths = table_model.get("column_widths")
    if not isinstance(raw_widths, list):
        raw_widths = table_model.get("columnWidths")
    widths = []
    for item in raw_widths if isinstance(raw_widths, list) else []:
        try:
            value = float(item)
        except (TypeError, ValueError):
            continue
        if value > 0:
            widths.append(value)
    if len(widths) != column_count:
        return [1.0 / float(column_count)] * column_count
    total = sum(widths)
    if total <= 0:
        return [1.0 / float(column_count)] * column_count
    return [value / total for value in widths]


def _table_alignment(value: object):
    from pptx.enum.text import PP_ALIGN

    normalized = str(value or "").strip().lower()
    if normalized == "right":
        return PP_ALIGN.RIGHT
    if normalized == "center":
        return PP_ALIGN.CENTER
    return PP_ALIGN.LEFT


def _add_native_table(
    slide,
    *,
    table_model: Mapping[str, object] | None,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str | None = None,
    body_size: float = 12.0,
    text_rgb: tuple[int, int, int] | None = None,
) -> bool:
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Pt

    if not _table_model_is_native_ready(table_model):
        return False
    assert table_model is not None
    row_count = _coerce_table_int(
        table_model.get("row_count") or table_model.get("rowCount")
    )
    column_count = _coerce_table_int(
        table_model.get("column_count") or table_model.get("columnCount")
    )
    header_rows = _coerce_table_int(
        table_model.get("header_rows") or table_model.get("headerRows")
    )
    rows = _table_model_rows(table_model)
    column_widths = _table_model_column_widths(table_model, column_count=column_count)
    table_shape = slide.shapes.add_table(
        row_count,
        column_count,
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    table = table_shape.table
    table.first_row = header_rows > 0
    table.horz_banding = row_count > max(1, header_rows + 1)
    default_text_rgb = text_rgb if text_rgb is not None else (17, 24, 39)
    header_fill_rgb = (236, 239, 243)
    band_fill_rgb = (248, 250, 252)
    row_height = _inches(height / float(max(1, row_count)))
    for row_index, row in enumerate(table.rows):
        row.height = row_height
        for column_index in range(column_count):
            if column_index < len(column_widths):
                table.columns[column_index].width = _inches(
                    width * column_widths[column_index]
                )
            cell = row.cells[column_index]
            cell.margin_left = _inches(0.04)
            cell.margin_right = _inches(0.04)
            cell.margin_top = _inches(0.02)
            cell.margin_bottom = _inches(0.02)
            if row_index < header_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(header_fill_rgb)
            elif row_index % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(band_fill_rgb)
            cell_payload = (
                rows[row_index][column_index]
                if row_index < len(rows) and column_index < len(rows[row_index])
                else None
            )
            text = (
                str(cell_payload.get("text") or "").strip()
                if isinstance(cell_payload, Mapping)
                else ""
            )
            text_frame = cell.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = _table_alignment(
                cell_payload.get("align") if isinstance(cell_payload, Mapping) else ""
            )
            paragraph.text = text
            for run in paragraph.runs:
                if font_name:
                    run.font.name = font_name
                run.font.size = Pt(body_size)
                run.font.bold = row_index < header_rows
                run.font.color.rgb = _rgb(default_text_rgb)
    return True


def _add_native_table_within_emu_box(
    slide,
    *,
    table_model: Mapping[str, object] | None,
    left: int,
    top: int,
    width: int,
    height: int,
) -> bool:
    return _add_native_table(
        slide,
        table_model=table_model,
        left=_emu_to_inches(left),
        top=_emu_to_inches(top),
        width=_emu_to_inches(width),
        height=_emu_to_inches(height),
    )


def _prototype_body_text(slide_spec: SlidesPptxSlide) -> str:
    if slide_spec.comparison_columns:
        parts: list[str] = []
        if slide_spec.body:
            parts.append(slide_spec.body)
        for column in slide_spec.comparison_columns:
            title = _normalize_text(str(column.get("title") or ""))
            bullets = _read_text_list(column.get("bullets"))
            if title:
                parts.append(title)
            parts.extend(_bullet_level_and_text(bullet)[1] for bullet in bullets)
        if slide_spec.callout_title:
            parts.append(slide_spec.callout_title)
        if slide_spec.callout_body:
            parts.append(slide_spec.callout_body)
        if slide_spec.implication:
            parts.append(slide_spec.implication)
        return "\n".join(part for part in parts if part)
    if slide_spec.bullets:
        parts = [_bullet_level_and_text(bullet)[1] for bullet in slide_spec.bullets]
        if slide_spec.implication:
            parts.append(slide_spec.implication)
        return "\n".join(parts)
    parts = [slide_spec.body]
    if slide_spec.footer_text:
        parts.append(slide_spec.footer_text)
    if slide_spec.implication:
        parts.append(slide_spec.implication)
    return "\n".join(part for part in parts if part)


def _bullet_level_and_text(text: str) -> tuple[int, str]:
    raw = str(text or "")
    level = len(raw) - len(raw.lstrip("\t"))
    normalized = _normalize_bullet_text(raw.lstrip("\t"))
    return level, normalized


def _bullet_paragraph_text(text: str) -> str:
    return _bullet_level_and_text(text)[1]


def _apply_native_bullet_format(paragraph, *, level: int) -> None:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Pt

    paragraph.level = level
    p_pr = paragraph._pPr
    base_margin = int(Pt(20))
    level_step = int(Pt(18))
    hanging_indent = int(Pt(14))
    p_pr.set("marL", str(base_margin + (level * level_step)))
    p_pr.set("indent", str(-hanging_indent))

    bullet_tags = {
        qn("a:buNone"),
        qn("a:buAutoNum"),
        qn("a:buBlip"),
        qn("a:buChar"),
    }
    for child in list(p_pr):
        if child.tag in bullet_tags:
            p_pr.remove(child)

    bullet_char = OxmlElement("a:buChar")
    bullet_char.set("char", "•")
    p_pr.insert_element_before(bullet_char, "a:tabLst", "a:defRPr", "a:extLst")


def _add_bullet_box(
    slide,
    *,
    bullets: Sequence[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    body_size: float,
    line_height: float,
    paragraph_spacing: float,
    text_rgb: tuple[int, int, int],
    margin: float = 0.06,
) -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    textbox = slide.shapes.add_textbox(
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    margin_emu = _inches(margin)
    frame.margin_left = margin_emu
    frame.margin_right = margin_emu
    frame.margin_top = margin_emu
    frame.margin_bottom = margin_emu
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        level, paragraph_text = _bullet_level_and_text(bullet)
        paragraph.text = paragraph_text
        _apply_native_bullet_format(paragraph, level=level)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(paragraph_spacing)
        paragraph.line_spacing = line_height
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(body_size)
            run.font.color.rgb = _rgb(text_rgb)


def _add_textbox(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    font_size: float,
    color_rgb: tuple[int, int, int],
    bold: bool = False,
    alignment=None,
    margin: float = 0.06,
    lead_prefixes: Sequence[str] = (),
) -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    textbox = slide.shapes.add_textbox(
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    margin_emu = _inches(margin)
    frame.margin_left = margin_emu
    frame.margin_right = margin_emu
    frame.margin_top = margin_emu
    frame.margin_bottom = margin_emu
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT if alignment is None else alignment
    _set_paragraph_runs(
        paragraph,
        text=text,
        font_name=font_name,
        font_size=font_size,
        color_rgb=color_rgb,
        bold=bold,
        lead_prefixes=lead_prefixes,
    )


def _set_paragraph_runs(
    paragraph,
    *,
    text: str,
    font_name: str,
    font_size: float,
    color_rgb: tuple[int, int, int],
    bold: bool = False,
    lead_prefixes: Sequence[str] = (),
) -> None:
    from pptx.util import Pt

    normalized_text = str(text or "")
    prefix = ""
    for candidate in lead_prefixes:
        if normalized_text.startswith(candidate):
            prefix = candidate
            break
    if prefix:
        lead_run = paragraph.add_run()
        lead_run.text = prefix
        lead_run.font.name = font_name
        lead_run.font.size = Pt(font_size)
        lead_run.font.bold = True
        lead_run.font.color.rgb = _rgb(color_rgb)
        remainder = normalized_text[len(prefix) :]
        if remainder:
            remainder_run = paragraph.add_run()
            remainder_run.text = remainder
            remainder_run.font.name = font_name
            remainder_run.font.size = Pt(font_size)
            remainder_run.font.bold = bold
            remainder_run.font.color.rgb = _rgb(color_rgb)
        return
    run = paragraph.add_run()
    run.text = normalized_text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_rgb)


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    normalized = str(value or "").strip().lstrip("#")
    if len(normalized) != 6:
        return (0, 0, 0)
    return tuple(int(normalized[index : index + 2], 16) for index in (0, 2, 4))


def _blend_rgb(
    base_rgb: tuple[int, int, int],
    other_rgb: tuple[int, int, int],
    ratio: float,
) -> tuple[int, int, int]:
    clamped_ratio = max(0.0, min(1.0, float(ratio)))
    return tuple(
        int(round((base_value * (1.0 - clamped_ratio)) + (other_value * clamped_ratio)))
        for base_value, other_value in zip(base_rgb, other_rgb, strict=False)
    )


def _rgb(value: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(*value)


def _inches(value: float):
    from pptx.util import Inches

    return Inches(value)


def _emu_to_inches(value: int) -> float:
    return float(value) / 914400.0
