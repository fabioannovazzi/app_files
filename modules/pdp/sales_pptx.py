from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Mapping, Sequence

from PIL import Image, ImageChops

from src.slides.notebooklm_style import load_notebooklm_style, resolve_prompt_style_key

__all__ = [
    "render_sales_authored_deck_plan_pptx",
]

_SLIDE_WIDTH_IN = 13.333333
_SLIDE_HEIGHT_IN = 7.5
_TITLE_LEFT_IN = 0.58
_TITLE_TOP_IN = 0.42
_TITLE_WIDTH_IN = 12.1
_TITLE_HEIGHT_IN = 0.8
_SUMMARY_BULLETS_LEFT_IN = 0.76
_SUMMARY_BULLETS_TOP_IN = 1.6
_SUMMARY_BULLETS_WIDTH_IN = 11.7
_SUMMARY_BULLETS_HEIGHT_IN = 4.8
_SUBTITLE_LEFT_IN = 0.6
_SUBTITLE_TOP_IN = 1.15
_SUBTITLE_WIDTH_IN = 3.0
_SUBTITLE_HEIGHT_IN = 0.24
_INSIGHT_BULLETS_LEFT_IN = 0.74
_INSIGHT_BULLETS_TOP_IN = 1.55
_INSIGHT_BULLETS_WIDTH_IN = 2.9
_INSIGHT_BULLETS_HEIGHT_IN = 4.95
_CHART_LEFT_IN = 4.02
_CHART_TOP_IN = 1.48
_CHART_WIDTH_IN = 8.7
_CHART_HEIGHT_IN = 5.15


def render_sales_authored_deck_plan_pptx(
    authored_deck_plan: Mapping[str, Any],
    *,
    chart_images: Mapping[str, bytes],
    template_key: str = "uniform",
) -> io.BytesIO:
    """Render an editable PPTX from the authored sales deck plan."""

    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN

    style_key = resolve_prompt_style_key(template_key or "uniform")
    presentation = _load_template_presentation(Presentation, style_key)
    presentation.slide_width = _inches(_SLIDE_WIDTH_IN)
    presentation.slide_height = _inches(_SLIDE_HEIGHT_IN)
    blank_layout = presentation.slide_layouts[6]
    style = load_notebooklm_style(style_key)
    font_name = style.font_family_primary or style.font_family_fallback
    title_size = max(float(style.title_size_pt) - 6.0, 22.0)
    body_size = max(float(style.body_size_pt) - 1.0, 15.0)
    line_height = float(style.line_height)
    text_rgb = _hex_to_rgb(style.text_color)
    muted_rgb = _hex_to_rgb("#000000")

    slides = _read_slide_list(authored_deck_plan.get("slides"))
    for index, slide_payload in enumerate(slides, start=1):
        slide = presentation.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(_hex_to_rgb(style.bg_color))

        kind = _read_text(slide_payload.get("kind"))
        if kind == "summary":
            _add_textbox(
                slide,
                text=_read_text(slide_payload.get("title")),
                left=_TITLE_LEFT_IN,
                top=_TITLE_TOP_IN,
                width=_TITLE_WIDTH_IN,
                height=_TITLE_HEIGHT_IN,
                font_name=font_name,
                font_size=title_size,
                color_rgb=text_rgb,
                bold=True,
            )
            _add_bullet_box(
                slide,
                bullets=_read_text_list(slide_payload.get("bullets")),
                left=_SUMMARY_BULLETS_LEFT_IN,
                top=_SUMMARY_BULLETS_TOP_IN,
                width=_SUMMARY_BULLETS_WIDTH_IN,
                height=_SUMMARY_BULLETS_HEIGHT_IN,
                font_name=font_name,
                body_size=max(body_size + 1.0, 17.0),
                line_height=line_height,
                text_rgb=text_rgb,
            )
        else:
            chart_id = _read_text(slide_payload.get("chart_id"))
            image_bytes = chart_images.get(chart_id) if chart_id else None
            if image_bytes:
                image_bytes = _trim_white_image_bytes(image_bytes)
                left, top, width, height = _fit_image_bytes_within_box(
                    image_bytes,
                    left=_CHART_LEFT_IN,
                    top=_CHART_TOP_IN,
                    width=_CHART_WIDTH_IN,
                    height=_CHART_HEIGHT_IN,
                )
                slide.shapes.add_picture(
                    io.BytesIO(image_bytes),
                    _inches(left),
                    _inches(top),
                    width=_inches(width),
                    height=_inches(height),
                )
            else:
                _add_textbox(
                    slide,
                    text="Chart preview unavailable",
                    left=_CHART_LEFT_IN + 2.5,
                    top=_CHART_TOP_IN + 2.2,
                    width=3.2,
                    height=0.3,
                    font_name=font_name,
                    font_size=11.0,
                    color_rgb=muted_rgb,
                    alignment=PP_ALIGN.CENTER,
                )
            _add_textbox(
                slide,
                text=_read_text(slide_payload.get("title")),
                left=_TITLE_LEFT_IN,
                top=_TITLE_TOP_IN,
                width=_TITLE_WIDTH_IN,
                height=_TITLE_HEIGHT_IN,
                font_name=font_name,
                font_size=title_size,
                color_rgb=text_rgb,
                bold=True,
            )
            subtitle = _read_text(slide_payload.get("subtitle"))
            if subtitle:
                _add_textbox(
                    slide,
                    text=subtitle.upper(),
                    left=_SUBTITLE_LEFT_IN,
                    top=_SUBTITLE_TOP_IN,
                    width=_SUBTITLE_WIDTH_IN,
                    height=_SUBTITLE_HEIGHT_IN,
                    font_name=font_name,
                    font_size=10.5,
                    color_rgb=muted_rgb,
                    bold=True,
                )
            _add_bullet_box(
                slide,
                bullets=_read_text_list(slide_payload.get("bullets")),
                left=_INSIGHT_BULLETS_LEFT_IN,
                top=_INSIGHT_BULLETS_TOP_IN,
                width=_INSIGHT_BULLETS_WIDTH_IN,
                height=_INSIGHT_BULLETS_HEIGHT_IN,
                font_name=font_name,
                body_size=body_size,
                line_height=line_height,
                text_rgb=text_rgb,
            )

        _add_textbox(
            slide,
            text=str(index),
            left=12.2,
            top=7.03,
            width=0.45,
            height=0.2,
            font_name=font_name,
            font_size=10.0,
            color_rgb=muted_rgb,
            alignment=PP_ALIGN.RIGHT,
        )

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


def _read_slide_list(value: object) -> list[Mapping[str, Any]]:
    if not isinstance(value, list):
        return []
    return [slide for slide in value if isinstance(slide, Mapping)]


def _read_text(value: object) -> str:
    return str(value or "").strip()


def _read_text_list(value: object) -> list[str]:
    if not isinstance(value, list):
        return []
    return [text for text in (str(item or "").strip() for item in value) if text]


def _load_template_presentation(presentation_cls, template_key: str):
    template_path = (
        Path(__file__).resolve().parents[2]
        / "src"
        / "review_brief"
        / "pptx_templates"
        / f"{template_key}.pptx"
    )
    if template_path.exists():
        return presentation_cls(str(template_path))
    return presentation_cls()


def _fit_image_bytes_within_box(
    image_bytes: bytes,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
) -> tuple[float, float, float, float]:
    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            img_width, img_height = image.size
    except (OSError, ValueError):
        return left, top, width, height
    if img_width <= 0 or img_height <= 0:
        return left, top, width, height
    box_ratio = width / height
    image_ratio = float(img_width) / float(img_height)
    if image_ratio >= box_ratio:
        scaled_width = width
        scaled_height = width / image_ratio
        offset_x = 0.0
        offset_y = (height - scaled_height) / 2.0
    else:
        scaled_height = height
        scaled_width = height * image_ratio
        offset_x = (width - scaled_width) / 2.0
        offset_y = 0.0
    return left + offset_x, top + offset_y, scaled_width, scaled_height


def _trim_white_image_bytes(image_bytes: bytes, *, padding_px: int = 12) -> bytes:
    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            working = image.convert("RGBA")
            flattened = Image.new("RGBA", working.size, (255, 255, 255, 255))
            flattened.alpha_composite(working)
            working_rgb = flattened.convert("RGB")
            background = Image.new("RGB", working.size, "white")
            diff = ImageChops.difference(working_rgb, background)
            bbox = diff.getbbox()
            if bbox is None:
                return image_bytes
            left, top, right, bottom = bbox
            left = max(0, left - padding_px)
            top = max(0, top - padding_px)
            right = min(working_rgb.width, right + padding_px)
            bottom = min(working_rgb.height, bottom + padding_px)
            cropped = working_rgb.crop((left, top, right, bottom))
            buffer = io.BytesIO()
            cropped.save(buffer, format="PNG")
            return buffer.getvalue()
    except (OSError, ValueError):
        return image_bytes


def _add_bullet_box(
    slide,
    *,
    bullets: Sequence[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    body_size: float,
    line_height: float,
    text_rgb: tuple[int, int, int],
) -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    textbox = slide.shapes.add_textbox(
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(10)
        paragraph.line_spacing = line_height
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(body_size)
            run.font.color.rgb = _rgb(text_rgb)


def _add_textbox(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str,
    font_size: float,
    color_rgb: tuple[int, int, int],
    bold: bool = False,
    alignment=None,
) -> None:
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    textbox = slide.shapes.add_textbox(
        _inches(left),
        _inches(top),
        _inches(width),
        _inches(height),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT if alignment is None else alignment
    paragraph.text = text
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color_rgb)


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    normalized = str(value or "").strip().lstrip("#")
    if len(normalized) != 6:
        return (0, 0, 0)
    return tuple(int(normalized[index : index + 2], 16) for index in (0, 2, 4))


def _rgb(value: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(*value)


def _inches(value: float):
    from pptx.util import Inches

    return Inches(value)
