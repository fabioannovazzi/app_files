from __future__ import annotations

import base64
import hashlib
import io
import json
import logging
import math
import mimetypes
import os
import re
import secrets
import shutil
import sys
import threading
import time
import uuid
from concurrent.futures import ThreadPoolExecutor
from copy import deepcopy
from dataclasses import replace
from datetime import UTC, datetime, timedelta
from functools import lru_cache
from html import escape
from importlib.util import find_spec
from pathlib import Path
from typing import Any, Callable, Iterable, List, Literal, NamedTuple, Optional, cast
from urllib.parse import ParseResult, urlparse

import fitz  # type: ignore[import-not-found]
from bs4 import BeautifulSoup  # type: ignore[import]
from fastapi import (
    APIRouter,
    Depends,
    File,
    Form,
    HTTPException,
    Request,
    UploadFile,
    status,
)
from fastapi.responses import FileResponse, JSONResponse, Response, StreamingResponse
from fastapi.templating import Jinja2Templates
from PIL import Image, ImageDraw, UnidentifiedImageError
from pydantic import BaseModel, ConfigDict, Field

from modules.auth.dependencies import (
    require_authenticated_user,
    require_authenticated_user_for_site,
)
from modules.auth.session import AuthenticatedUser
from modules.llm.llm_call_wrapper import LLMCallWrapper
from modules.notifications.notifier import notify_failed, notify_finished
from modules.pdp.language import (
    get_navigation_label,
    get_page_copy,
    resolve_language,
)
from modules.slides.ocr import (
    SlideOcrEngineUnavailableError,
    extract_layout_summary_from_raw_layout,
    extract_raw_layout_from_data_url,
    extract_raw_ocr_from_data_url,
    extract_text_from_raw_ocr_result,
)
from modules.slides.pptx_jobs import (
    PptxJob,
)
from modules.slides.pptx_jobs import cleanup_expired_jobs as cleanup_expired_pptx_jobs
from modules.slides.pptx_jobs import create_job as create_pptx_job
from modules.slides.pptx_jobs import get_job as get_pptx_job_record
from modules.slides.pptx_jobs import set_output_path as set_pptx_output_path
from modules.slides.pptx_jobs import update_job_status as update_pptx_job_status
from modules.slides.print_jobs import (
    PrintJob,
    cleanup_expired_jobs,
    create_job,
    get_job,
    set_output_path,
    update_job_status,
)
from modules.utilities.config import get_naming_params
from src.slides import deck_from_payload, generate_slide_filename
from src.slides.chart_type_classifier import classify_deck_chart_regions
from src.slides.errors import DeckNotFoundError, InvalidDeckError, SlideNotFoundError
from src.slides.html_normalizer import update_slide_document
from src.slides.layout_service import build_deck_layout_payload
from src.slides.loader import INDEX_HTML, INDEX_JSON, find_index_file
from src.slides.models import Deck, Section, Slide, Subsection
from src.slides.notebooklm_style import (
    DEFAULT_PROMPT_STYLE_KEY,
    build_notebooklm_css_variables,
    load_notebooklm_style,
    resolve_prompt_style_key,
)
from src.slides.ocr_payload import (
    DeckOcrAnalysisPayload,
    DeckOcrPayload,
    DeckOcrSlide,
    build_analysis_payload,
    find_slide_payload,
    normalize_ocr_payload,
)
from src.slides.ocr_service import (
    OCR_STRATEGY_LAYOUT_GUIDED,
    build_deck_ocr_payload,
    ensure_deck_ocr_payload,
)
from src.slides.pdf_import import build_image_only_slide_content
from src.slides.pptx_post_render import apply_post_render_compare_loop
from src.slides.pptx_template_manifest import (
    DECK_PPTX_TEMPLATE_FILENAME,
    DECK_PPTX_TEMPLATE_MANIFEST_FILENAME,
)
from src.slides.pptx_template_store import (
    PptxTemplateRecord,
    apply_saved_pptx_template_to_deck,
    list_saved_pptx_templates,
    save_uploaded_pptx_template,
    set_default_pptx_template,
)
from src.slides.semantic_pptx import (
    build_slides_pptx_spec,
    render_slides_pptx_from_template,
    write_slides_pptx_spec,
)
from src.slides.storage import DeckStorage

try:  # pragma: no cover - optional dependency for uploads
    import multipart  # type: ignore  # noqa: F401

    HAS_MULTIPART = True
except Exception:  # noqa: BLE001
    HAS_MULTIPART = False

PlaywrightError = RuntimeError
try:  # pragma: no cover - optional dependency for PDF exports
    from playwright.sync_api import sync_playwright

    HAS_PLAYWRIGHT = True
    try:
        from playwright.sync_api import Error as PlaywrightError
    except Exception:  # noqa: BLE001
        PlaywrightError = RuntimeError
except Exception:  # noqa: BLE001
    HAS_PLAYWRIGHT = False

HAS_REPORTLAB = find_spec("reportlab") is not None
HAS_PPTX = find_spec("pptx") is not None

LOGGER = logging.getLogger(__name__)

_SLIDES_EXPORT_PAGE_WIDTH_PT = 1280
_SLIDES_EXPORT_PAGE_HEIGHT_PT = 720
_SLIDES_RENDERED_PPTX_HEIGHT_IN = 7.5
_SLIDES_PDF_ALLOWED_DIMENSIONS_PT: tuple[tuple[float, float], ...] = ((1376.0, 768.0),)
_SLIDES_PDF_REQUIRED_DIMENSION_TOLERANCE_PT = 0.5
_SLIDES_EXPORT_ENABLE_VECTOR_PDF = str(
    os.getenv("SLIDES_EXPORT_ENABLE_VECTOR_PDF", "0")
).strip().lower() in {"1", "true", "yes", "on"}
try:
    _SLIDES_EXPORT_SCREENSHOT_SCALE = max(
        1.0, float(os.getenv("SLIDES_EXPORT_SCREENSHOT_SCALE", "2"))
    )
except ValueError:
    _SLIDES_EXPORT_SCREENSHOT_SCALE = 2.0

try:
    _SLIDES_PDF_IMPORT_RASTER_SCALE = max(
        1.0, float(os.getenv("SLIDES_PDF_IMPORT_RASTER_SCALE", "2"))
    )
except ValueError:
    _SLIDES_PDF_IMPORT_RASTER_SCALE = 2.0

_SLIDES_PPTX_POST_RENDER_COMPARE_ENABLED = str(
    os.getenv("SLIDES_PPTX_POST_RENDER_COMPARE", "1")
).strip().lower() in {"1", "true", "yes", "on"}

templates = Jinja2Templates(directory="templates")
router = APIRouter(
    prefix="/slides",
    tags=["slides"],
    dependencies=[Depends(require_authenticated_user)],
)
site_router = APIRouter(
    prefix="/slides", dependencies=[Depends(require_authenticated_user_for_site)]
)

_DECK_ROOT = Path("slide_decks")
_DECK_ROOT.mkdir(parents=True, exist_ok=True)
_storage = DeckStorage(_DECK_ROOT)

_OCR_LANG_BY_LOCALE = {
    "en": "eng",
    "it": "ita",
    "fr": "fra",
    "de": "deu",
    "es": "spa",
}

_NOTIFY_LANG_BY_OCR_LANG = {
    "eng": "en",
    "ita": "it",
    "fra": "fr",
    "deu": "de",
    "spa": "es",
}


def _resolve_ocr_lang(locale: str | None, *, fallback: str = "eng") -> str:
    if not locale:
        return fallback
    return _OCR_LANG_BY_LOCALE.get(locale.lower(), fallback)


def _resolve_notify_lang_from_ocr_lang(ocr_lang: str | None) -> str:
    normalized = str(ocr_lang or "").strip().lower()
    return _NOTIFY_LANG_BY_OCR_LANG.get(normalized, "en")


def _build_deck_processing_notify_context(
    user: AuthenticatedUser | None,
    *,
    ocr_lang: str | None,
) -> dict[str, str] | None:
    if user is None or not user.email:
        return None
    return {
        "notify_email": user.email,
        "notify_lang": _resolve_notify_lang_from_ocr_lang(ocr_lang),
    }


def _get_pptx_import_error() -> str | None:
    try:
        import pptx  # type: ignore  # noqa: F401
    except Exception as exc:  # noqa: BLE001
        return _truncate_dependency_error(f"{type(exc).__name__}: {exc}")
    return None


def _ensure_pptx_available() -> None:
    err = _get_pptx_import_error()
    if err is None:
        return
    raise HTTPException(
        status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
        detail=(
            "PPTX export requires `python-pptx` in the API runtime. "
            f"Python executable: `{sys.executable}`. Error: {err}"
        ),
    )


def _ocr_engine_unavailable_detail(
    *,
    lang: str | None = None,
    exc: BaseException | None = None,
) -> str:
    detail = (
        "PaddleOCR is unavailable or misconfigured on this server. "
        "Install `paddleocr` and its runtime dependencies in the API environment."
    )
    if lang and lang.lower() != "eng":
        detail = f"{detail} OCR model data for '{lang}' may be missing."
    if exc:
        reason = _truncate_dependency_error(str(exc))
        if reason:
            detail = f"{detail} Details: {reason}"
    return detail


_ENV_PREVIEW_STYLES = "SLIDES_PREVIEW_STYLE_BUNDLES"
_ENV_PREVIEW_SCRIPTS = "SLIDES_PREVIEW_SCRIPT_BUNDLES"
_ENV_PREVIEW_ALLOWLIST = "SLIDES_PREVIEW_ORIGIN_ALLOWLIST"


_PRINT_JOB_DIR = Path("tmp/pdf_jobs").resolve()
_PRINT_JOB_DIR.mkdir(parents=True, exist_ok=True)
_PRINT_JOB_TTL = timedelta(hours=12)
_PRINT_EXECUTOR = ThreadPoolExecutor(max_workers=2)

_PPTX_JOB_DIR = Path("tmp/pptx_jobs").resolve()
_PPTX_JOB_DIR.mkdir(parents=True, exist_ok=True)
_PPTX_JOB_TTL = timedelta(hours=12)
_PPTX_EXECUTOR = ThreadPoolExecutor(max_workers=2)

_OCR_EXECUTOR = ThreadPoolExecutor(max_workers=1)
_LAYOUT_EXECUTOR = ThreadPoolExecutor(max_workers=1)

_OCR_PROGRESS: dict[str, dict[str, object]] = {}
_OCR_PROGRESS_LOCK = threading.Lock()
_LAYOUT_PROGRESS: dict[str, dict[str, object]] = {}
_LAYOUT_PROGRESS_LOCK = threading.Lock()


class PreviewBundleConfig(NamedTuple):
    styles: tuple[str, ...]
    scripts: tuple[str, ...]
    allowlist: tuple[str, ...]


def _utc_now_iso() -> str:
    return datetime.now(UTC).isoformat()


def _set_progress_record(
    progress_store: dict[str, dict[str, object]],
    progress_lock: threading.Lock,
    deck_id: str,
    *,
    status: Literal["idle", "running", "completed", "failed", "skipped"],
    built_pages: int = 0,
    total_pages: int = 0,
    message: str | None = None,
    error: str | None = None,
    lang: str | None = None,
    step: str | None = None,
    last_completed_step: str | None = None,
) -> None:
    now = _utc_now_iso()
    normalized_step = str(step or "").strip() or None
    normalized_last_completed_step = str(last_completed_step or "").strip() or None
    with progress_lock:
        previous = progress_store.get(deck_id, {})
        previous_status = str(previous.get("status") or "").strip().lower()
        previous_started_at = previous.get("startedAt")
        previous_lang = (
            str(previous.get("lang")).strip()
            if isinstance(previous.get("lang"), str)
            and str(previous.get("lang")).strip()
            else None
        )
        started_at = (
            str(previous_started_at).strip()
            if isinstance(previous_started_at, str) and previous_started_at.strip()
            else None
        )
        resolved_lang = (
            str(lang).strip()
            if isinstance(lang, str) and str(lang).strip()
            else previous_lang
        )
        if status == "running":
            if previous_status != "running" or started_at is None:
                started_at = now
        elif status != "idle" and started_at is None:
            started_at = now
        progress_store[deck_id] = {
            "status": status,
            "builtPages": built_pages,
            "totalPages": total_pages,
            "message": message,
            "error": error,
            "lang": resolved_lang,
            "step": normalized_step,
            "startedAt": started_at,
            "updatedAt": now,
            "lastCompletedStep": normalized_last_completed_step,
        }


def _set_ocr_progress(
    deck_id: str,
    *,
    status: Literal["idle", "running", "completed", "failed", "skipped"],
    built_pages: int = 0,
    total_pages: int = 0,
    message: str | None = None,
    error: str | None = None,
    lang: str | None = None,
    step: str | None = None,
    last_completed_step: str | None = None,
) -> None:
    _set_progress_record(
        _OCR_PROGRESS,
        _OCR_PROGRESS_LOCK,
        deck_id,
        status=status,
        built_pages=built_pages,
        total_pages=total_pages,
        message=message,
        error=error,
        lang=lang,
        step=step,
        last_completed_step=last_completed_step,
    )


def _get_ocr_progress(deck_id: str) -> dict[str, object] | None:
    with _OCR_PROGRESS_LOCK:
        return _OCR_PROGRESS.get(deck_id)


def is_any_ocr_running() -> bool:
    with _OCR_PROGRESS_LOCK:
        return any(
            str(progress.get("status") or "").strip().lower() == "running"
            for progress in _OCR_PROGRESS.values()
            if isinstance(progress, dict)
        )


def _set_layout_progress(
    deck_id: str,
    *,
    status: Literal["idle", "running", "completed", "failed", "skipped"],
    built_pages: int = 0,
    total_pages: int = 0,
    message: str | None = None,
    error: str | None = None,
    step: str | None = None,
    last_completed_step: str | None = None,
) -> None:
    _set_progress_record(
        _LAYOUT_PROGRESS,
        _LAYOUT_PROGRESS_LOCK,
        deck_id,
        status=status,
        built_pages=built_pages,
        total_pages=total_pages,
        message=message,
        error=error,
        step=step,
        last_completed_step=last_completed_step,
    )


def _get_layout_progress(deck_id: str) -> dict[str, object] | None:
    with _LAYOUT_PROGRESS_LOCK:
        return _LAYOUT_PROGRESS.get(deck_id)


def _build_layout_processing_message(built_pages: int, total_pages: int) -> str:
    total = max(0, int(total_pages))
    built = min(max(0, int(built_pages)), total) if total else max(0, int(built_pages))
    if total > 0:
        return f"Analyzing layout. Slides done: {built} of {total}."
    return "Analyzing layout..."


class SlidePayload(BaseModel):
    id: str
    titleHtml: str = ""
    bodyHtml: str = ""
    notesHtml: str = ""
    sourceHtml: str = ""
    fullHtml: str = ""
    kind: Literal["normal", "sectionHeader"] = "normal"
    sectionId: Optional[str] = None
    subsectionId: Optional[str] = None

    model_config = ConfigDict(populate_by_name=True)

    def to_domain(self) -> Slide:
        return Slide(
            id=self.id,
            title_html=self.titleHtml,
            body_html=self.bodyHtml,
            notes_html=self.notesHtml,
            source_html=self.sourceHtml,
            full_html=self.fullHtml,
            kind=self.kind,
            section_id=self.sectionId,
            subsection_id=self.subsectionId,
        )

    @classmethod
    def from_domain(cls, slide: Slide) -> "SlidePayload":
        return cls(
            id=slide.id,
            titleHtml=slide.title_html,
            bodyHtml=slide.body_html,
            notesHtml=slide.notes_html,
            sourceHtml=slide.source_html,
            fullHtml=slide.full_html,
            kind=slide.kind,
            sectionId=slide.section_id,
            subsectionId=slide.subsection_id,
        )


class SubsectionPayload(BaseModel):
    id: str
    title: str
    startSlide: str

    @classmethod
    def from_domain(cls, subsection: Subsection) -> "SubsectionPayload":
        return cls(
            id=subsection.id, title=subsection.title, startSlide=subsection.start_slide
        )


class SectionPayload(BaseModel):
    id: str
    title: str
    startSlide: str
    subsections: List[SubsectionPayload] = Field(default_factory=list)

    @classmethod
    def from_domain(cls, section: Section) -> "SectionPayload":
        return cls(
            id=section.id,
            title=section.title,
            startSlide=section.start_slide,
            subsections=[
                SubsectionPayload.from_domain(sub) for sub in section.subsections
            ],
        )


class DeckResponse(BaseModel):
    deckId: str
    promptStyle: str = Field(default=DEFAULT_PROMPT_STYLE_KEY, alias="promptStyle")
    ownerEmail: str | None = Field(default=None, alias="ownerEmail")
    sharedWith: List[str] = Field(default_factory=list, alias="sharedWith")
    slides: List[SlidePayload]
    sections: List[SectionPayload] = Field(default_factory=list)
    thumbnails: dict[str, str] = Field(default_factory=dict, alias="thumbnails")
    hasPdf: bool = Field(default=False, alias="hasPdf")
    hasImages: bool = Field(default=False, alias="hasImages")
    hasLayout: bool = Field(default=False, alias="hasLayout")
    archivedSourceDecks: List[str] = Field(
        default_factory=list, alias="archivedSourceDecks"
    )

    model_config = ConfigDict(populate_by_name=True)


class DeckSaveRequest(BaseModel):
    slides: List[SlidePayload]
    sections: List[SectionPayload] = Field(default_factory=list)
    promptStyle: str | None = Field(default=None, alias="promptStyle")


class DeckShareRequest(BaseModel):
    sharedWith: List[str] = Field(default_factory=list, alias="sharedWith")

    model_config = ConfigDict(populate_by_name=True)


class SlidesOcrRequest(BaseModel):
    imageDataUrl: str = Field(..., alias="imageDataUrl")
    lang: str = "eng"
    slideId: str | None = Field(default=None, alias="slideId")
    slideNumber: int | None = Field(default=None, alias="slideNumber", ge=1)

    model_config = ConfigDict(populate_by_name=True)


class SlidesOcrResponse(BaseModel):
    ocrText: str = Field(..., alias="ocrText")
    rawOcr: object = Field(..., alias="rawOcr")

    model_config = ConfigDict(populate_by_name=True)


class SlidesLayoutRequest(BaseModel):
    imageDataUrl: str = Field(..., alias="imageDataUrl")
    lang: str = "eng"
    slideId: str | None = Field(default=None, alias="slideId")
    slideNumber: int | None = Field(default=None, alias="slideNumber", ge=1)

    model_config = ConfigDict(populate_by_name=True)


class SlidesLayoutBlock(BaseModel):
    blockId: str | None = Field(default=None, alias="blockId")
    type: str = "unknown"
    detectedType: str | None = Field(default=None, alias="detectedType")
    text: str = ""
    items: List[str] = Field(default_factory=list)
    groupId: str | None = Field(default=None, alias="groupId")
    groupKind: str | None = Field(default=None, alias="groupKind")
    parentId: str | None = Field(default=None, alias="parentId")
    listLevel: int | None = Field(default=None, alias="listLevel")
    readingOrder: int | None = Field(default=None, alias="readingOrder")
    renderMode: str | None = Field(default=None, alias="renderMode")
    bbox: dict[str, float] | None = None
    confidence: float | None = None
    tableModel: dict[str, object] | None = Field(default=None, alias="tableModel")
    auditStatus: str | None = Field(default=None, alias="auditStatus")
    auditReason: str | None = Field(default=None, alias="auditReason")
    auditSuggestedText: str | None = Field(default=None, alias="auditSuggestedText")
    visualStatus: str | None = Field(default=None, alias="visualStatus")
    visualReason: str | None = Field(default=None, alias="visualReason")
    visualSuggestedText: str | None = Field(default=None, alias="visualSuggestedText")
    visualConfidence: float | None = Field(default=None, alias="visualConfidence")
    visualText: str | None = Field(default=None, alias="visualText")
    visualItems: List[str] = Field(default_factory=list, alias="visualItems")
    visualLines: List[dict[str, object]] = Field(
        default_factory=list, alias="visualLines"
    )

    model_config = ConfigDict(populate_by_name=True)


class SlidesLayoutFigureRegion(BaseModel):
    x: float
    y: float
    w: float
    h: float


class SlidesLayoutResponse(BaseModel):
    rawLayout: object = Field(..., alias="rawLayout")
    blocks: List[SlidesLayoutBlock] = Field(default_factory=list)
    titleText: str = Field(default="", alias="titleText")
    bulletTexts: List[str] = Field(default_factory=list, alias="bulletTexts")
    figureRegions: List[SlidesLayoutFigureRegion] = Field(
        default_factory=list, alias="figureRegions"
    )

    model_config = ConfigDict(populate_by_name=True)


class DeckLayoutRequest(BaseModel):
    lang: str = "eng"
    force: bool = False


class DeckLayoutSlide(BaseModel):
    slideId: str = Field(..., alias="slideId")
    slideNumber: int = Field(..., alias="slideNumber")
    pageNumber: int = Field(..., alias="pageNumber")
    assetPath: str = Field(default="", alias="assetPath")
    rawLayout: object | None = Field(default=None, alias="rawLayout")
    blocks: List[SlidesLayoutBlock] = Field(default_factory=list)
    titleText: str = Field(default="", alias="titleText")
    bulletTexts: List[str] = Field(default_factory=list, alias="bulletTexts")
    figureRegions: List[SlidesLayoutFigureRegion] = Field(
        default_factory=list, alias="figureRegions"
    )

    model_config = ConfigDict(populate_by_name=True)


class DeckLayoutPayload(BaseModel):
    deckId: str = Field(..., alias="deckId")
    lang: str = "eng"
    generatedAt: str = Field(..., alias="generatedAt")
    slides: List[DeckLayoutSlide] = Field(default_factory=list)

    model_config = ConfigDict(populate_by_name=True)


class DeckLayoutResponse(BaseModel):
    payload: DeckLayoutPayload
    cached: bool = False


class DeckLayoutStatusResponse(BaseModel):
    status: Literal["idle", "running", "completed", "failed", "skipped"]
    builtPages: int = Field(default=0, alias="builtPages")
    totalPages: int = Field(default=0, alias="totalPages")
    message: str | None = None
    error: str | None = None
    step: str | None = None
    startedAt: str | None = Field(default=None, alias="startedAt")
    updatedAt: str | None = Field(default=None, alias="updatedAt")
    lastCompletedStep: str | None = Field(default=None, alias="lastCompletedStep")

    model_config = ConfigDict(populate_by_name=True)


class DeckOcrRequest(BaseModel):
    lang: str = "eng"
    includeBboxes: bool = Field(default=True, alias="includeBboxes")
    pdfPath: str | None = Field(default=None, alias="pdfPath")
    force: bool = False

    model_config = ConfigDict(populate_by_name=True)


class DeckOcrResponse(BaseModel):
    payload: DeckOcrPayload
    cached: bool


class DeckOcrStatusResponse(BaseModel):
    status: Literal["idle", "running", "completed", "failed", "skipped"]
    builtPages: int = Field(default=0, alias="builtPages")
    totalPages: int = Field(default=0, alias="totalPages")
    message: str | None = None
    error: str | None = None
    lang: str | None = None
    step: str | None = None
    startedAt: str | None = Field(default=None, alias="startedAt")
    updatedAt: str | None = Field(default=None, alias="updatedAt")
    lastCompletedStep: str | None = Field(default=None, alias="lastCompletedStep")

    model_config = ConfigDict(populate_by_name=True)


class ConcatenateDecksRequest(BaseModel):
    newDeckId: str = Field(..., alias="newDeckId")
    sourceDeckIds: List[str] = Field(..., alias="sourceDeckIds", min_length=1)
    interleaveByIndex: bool = Field(default=False, alias="interleaveByIndex")
    interleaveGuide: Literal["first", "longest", "shortest"] = Field(
        default="shortest", alias="interleaveGuide"
    )
    deleteSourceDecks: bool = Field(default=False, alias="deleteSourceDecks")
    deleteMode: Literal["archive"] = Field(default="archive", alias="deleteMode")

    model_config = ConfigDict(populate_by_name=True)


class DeckArchiveResponse(BaseModel):
    deckId: str
    status: Literal["archived"]


class ImportRequest(BaseModel):
    sourceDeckId: str
    sourceSlideId: str
    afterSlideId: Optional[str] = None
    currentOrder: List[str] = Field(default_factory=list)

    model_config = ConfigDict(populate_by_name=True)


class ImportResponse(BaseModel):
    slide: SlidePayload
    order: List[str]


class DeckSummary(BaseModel):
    deckId: str
    promptStyle: str = Field(default=DEFAULT_PROMPT_STYLE_KEY, alias="promptStyle")
    ownerEmail: str | None = Field(default=None, alias="ownerEmail")
    sharedWith: List[str] = Field(default_factory=list, alias="sharedWith")
    slides: List[SlidePayload]
    sections: List[SectionPayload] = Field(default_factory=list)

    model_config = ConfigDict(populate_by_name=True)


class DeckListResponse(BaseModel):
    decks: List[DeckSummary]


class PptxTemplateSummary(BaseModel):
    templateId: str = Field(..., alias="templateId")
    name: str
    originalFilename: str = Field(default="", alias="originalFilename")
    uploadedAt: str = Field(default="", alias="uploadedAt")
    isDefault: bool = Field(default=False, alias="isDefault")

    model_config = ConfigDict(populate_by_name=True)

    @classmethod
    def from_record(cls, record: PptxTemplateRecord) -> "PptxTemplateSummary":
        return cls(
            templateId=record.template_id,
            name=record.name,
            originalFilename=record.original_filename,
            uploadedAt=record.uploaded_at,
            isDefault=record.is_default,
        )


class PptxTemplateListResponse(BaseModel):
    templates: List[PptxTemplateSummary] = Field(default_factory=list)
    defaultTemplateId: str | None = Field(default=None, alias="defaultTemplateId")

    model_config = ConfigDict(populate_by_name=True)


class PptxTemplateDefaultRequest(BaseModel):
    templateId: str = Field(..., alias="templateId")

    model_config = ConfigDict(populate_by_name=True)


def get_storage() -> DeckStorage:
    return _storage


def set_storage(storage: DeckStorage) -> None:
    global _storage
    _storage = storage


def _load_existing_prompt_style(deck_id: str, storage: DeckStorage) -> str | None:
    deck_path = storage.root / deck_id
    if not deck_path.exists():
        return None
    index_path = deck_path / INDEX_JSON
    if not index_path.exists():
        return resolve_prompt_style_key(None)
    try:
        payload = json.loads(index_path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError) as exc:
        LOGGER.warning("Failed to read promptStyle for deck %s: %s", deck_id, exc)
        return resolve_prompt_style_key(None)
    prompt_style_value = (
        payload.get("promptStyle") if isinstance(payload, dict) else None
    )
    return resolve_prompt_style_key(
        prompt_style_value if isinstance(prompt_style_value, str) else None
    )


def _resolve_prompt_style_strict(prompt_style: str) -> str:
    try:
        return resolve_prompt_style_key(prompt_style, strict=True)
    except ValueError as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)
        ) from exc


def _resolve_prompt_style_for_save(
    *,
    existing_prompt_style: str | None,
    requested_prompt_style: str | None,
) -> str:
    base_style = existing_prompt_style or resolve_prompt_style_key(None)
    if existing_prompt_style:
        return base_style
    if requested_prompt_style is None:
        return base_style
    resolved_requested = _resolve_prompt_style_strict(requested_prompt_style)
    return resolved_requested


def _validate_shared_prompt_style(
    deck_ids: list[str], storage: DeckStorage, *, context: str
) -> str:
    styles: set[str] = set()
    for deck_id in deck_ids:
        deck = storage.load_deck(deck_id)
        styles.add(deck.prompt_style)
    if len(styles) > 1:
        keys = sorted(styles)
        detail = f"Cannot {context} decks with different promptStyle values: {', '.join(keys)}."
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=detail)
    if styles:
        return next(iter(styles))
    return resolve_prompt_style_key(None)


@lru_cache(maxsize=1)
def get_preview_bundle_config() -> PreviewBundleConfig:
    """Return configured preview bundle URLs and allowed origins."""

    return PreviewBundleConfig(
        styles=_parse_preview_entries(os.environ.get(_ENV_PREVIEW_STYLES)),
        scripts=_parse_preview_entries(os.environ.get(_ENV_PREVIEW_SCRIPTS)),
        allowlist=_parse_preview_entries(os.environ.get(_ENV_PREVIEW_ALLOWLIST)),
    )


def _parse_preview_entries(raw_value: str | None) -> tuple[str, ...]:
    if not raw_value:
        return ()
    entries: list[str] = []
    for part in re.split(r"[,\s]+", raw_value):
        item = part.strip()
        if not item or item in entries:
            continue
        entries.append(item)
    return tuple(entries)


def _static_asset_version(path: Path) -> str:
    try:
        return str(int(path.stat().st_mtime))
    except OSError:
        return str(int(time.time()))


def _normalize_shared_with(entries: Iterable[str]) -> list[str]:
    normalized: list[str] = []
    seen: set[str] = set()
    for entry in entries:
        cleaned = entry.strip()
        if not cleaned or cleaned in seen:
            continue
        normalized.append(cleaned)
        seen.add(cleaned)
    return normalized


def _can_access_deck(
    deck: Deck,
    user: AuthenticatedUser | None,
    *,
    allow_ownerless: bool = True,
) -> bool:
    if user is None:
        return True
    if deck.owner_email is None:
        return allow_ownerless
    if deck.owner_email == user.email:
        return True
    return user.email in deck.shared_with


def _assert_deck_access(
    deck: Deck,
    user: AuthenticatedUser | None,
    *,
    allow_ownerless: bool = True,
) -> None:
    if _can_access_deck(deck, user, allow_ownerless=allow_ownerless):
        return
    raise HTTPException(
        status_code=status.HTTP_403_FORBIDDEN,
        detail="You do not have access to this deck.",
    )


def _resolve_owner_email(
    existing_owner: str | None, user: AuthenticatedUser | None
) -> str | None:
    if existing_owner is not None:
        return existing_owner
    if user is None:
        return None
    return user.email


@site_router.get("/page", include_in_schema=False)
def slides_page(request: Request) -> Any:
    lang = resolve_language(request)
    page_label = get_navigation_label(lang, "/slides/page")
    preview_config = get_preview_bundle_config()
    asset_versions = {
        "css": _static_asset_version(Path("static/css/slides_editor.css")),
        "js": _static_asset_version(Path("static/js/slides-react.js")),
    }
    return templates.TemplateResponse(
        request,
        "slides_editor_react.html",
        {
            "lang": lang,
            "page_label": page_label,
            "copy": get_page_copy("slides_editor", lang),
            "preview_styles": list(preview_config.styles),
            "preview_scripts": list(preview_config.scripts),
            "preview_allowlist": list(preview_config.allowlist),
            "asset_versions": asset_versions,
        },
    )


@router.get("/decks", response_model=DeckListResponse)
def list_decks(
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckListResponse:
    deck_ids = storage.list_decks()
    summaries: List[DeckSummary] = []
    for deck_id in deck_ids:
        try:
            deck = storage.load_deck(deck_id)
        except Exception:
            LOGGER.exception("Failed to load deck %s", deck_id)
            continue
        if not _can_access_deck(deck, user):
            continue
        summaries.append(
            DeckSummary(
                deckId=deck.deck_id,
                promptStyle=deck.prompt_style,
                ownerEmail=deck.owner_email,
                sharedWith=list(deck.shared_with),
                slides=[SlidePayload.from_domain(slide) for slide in deck.slides],
                sections=[
                    SectionPayload.from_domain(section) for section in deck.sections
                ],
            )
        )
    return DeckListResponse(decks=summaries)


@router.get("/pptx-templates", response_model=PptxTemplateListResponse)
def list_pptx_templates(
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> PptxTemplateListResponse:
    owner_email = _resolve_owner_email(None, user)
    templates = list_saved_pptx_templates(storage.root, owner_email)
    default_template_id = next(
        (record.template_id for record in templates if record.is_default),
        None,
    )
    return PptxTemplateListResponse(
        templates=[PptxTemplateSummary.from_record(record) for record in templates],
        defaultTemplateId=default_template_id,
    )


if HAS_MULTIPART:

    @router.post(
        "/pptx-templates/upload",
        response_model=PptxTemplateSummary,
        status_code=status.HTTP_201_CREATED,
    )
    async def upload_pptx_template(
        file: UploadFile = File(..., alias="file"),
        set_default: bool = Form(True, alias="setDefault"),
        user: AuthenticatedUser | None = Depends(require_authenticated_user),
        storage: DeckStorage = Depends(get_storage),
    ) -> PptxTemplateSummary:
        if file is None:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Upload a .pptx or .potx template.",
            )
        file_bytes = await file.read()
        await file.close()
        if not file_bytes:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Uploaded PPTX template is empty.",
            )
        owner_email = _resolve_owner_email(None, user)
        try:
            record = save_uploaded_pptx_template(
                storage.root,
                owner_email,
                filename=file.filename or "template.pptx",
                file_bytes=file_bytes,
                set_default=bool(set_default),
            )
        except ValueError as exc:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=str(exc),
            ) from exc
        return PptxTemplateSummary.from_record(record)


@router.post("/pptx-templates/default", response_model=PptxTemplateListResponse)
def update_default_pptx_template(
    payload: PptxTemplateDefaultRequest,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> PptxTemplateListResponse:
    owner_email = _resolve_owner_email(None, user)
    try:
        set_default_pptx_template(storage.root, owner_email, payload.templateId)
    except FileNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=str(exc),
        ) from exc
    templates = list_saved_pptx_templates(storage.root, owner_email)
    default_template_id = next(
        (record.template_id for record in templates if record.is_default),
        None,
    )
    return PptxTemplateListResponse(
        templates=[PptxTemplateSummary.from_record(record) for record in templates],
        defaultTemplateId=default_template_id,
    )


@router.get("/deck/{deck_id}", response_model=DeckResponse)
def get_deck(
    deck_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckResponse:
    try:
        deck = storage.load_deck(deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    return _build_deck_response(deck, storage)


@router.get("/deck/{deck_id}/assets/{asset_path:path}", include_in_schema=False)
def get_deck_asset(
    deck_id: str,
    asset_path: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> FileResponse:
    normalized = _normalize_deck_id(deck_id)
    normalized_asset = _normalize_asset_path(asset_path)
    try:
        deck = storage.load_deck(normalized)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    deck_path = (storage.root / normalized).resolve()
    candidate_roots = (deck_path, deck_path / "assets")
    for root in candidate_roots:
        target = (root / normalized_asset).resolve()
        try:
            target.relative_to(root)
        except ValueError:  # pragma: no cover - defensive guard
            continue
        if target.exists() and not target.is_dir():
            return FileResponse(target)
    raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found.")


@router.post("/deck/{deck_id}/archive", response_model=DeckArchiveResponse)
def archive_deck_endpoint(
    deck_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckArchiveResponse:
    normalized = _normalize_deck_id(deck_id)
    try:
        deck = storage.load_deck(normalized)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    try:
        archived = storage.archive_deck(normalized)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT, detail=str(exc)
        ) from exc
    if not archived:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=f"Deck {normalized} not found.",
        )
    return DeckArchiveResponse(deckId=normalized, status="archived")


@router.post("/deck/{deck_id}/save", response_model=DeckResponse)
def save_deck(
    deck_id: str,
    payload: DeckSaveRequest,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckResponse:
    normalized_deck_id = _normalize_deck_id(deck_id)
    existing_prompt_style = _load_existing_prompt_style(normalized_deck_id, storage)
    resolved_prompt_style = _resolve_prompt_style_for_save(
        existing_prompt_style=existing_prompt_style,
        requested_prompt_style=payload.promptStyle,
    )
    existing_deck: Deck | None = None
    try:
        existing_deck = storage.load_deck(normalized_deck_id)
    except DeckNotFoundError:
        existing_deck = None
    if existing_deck is not None:
        _assert_deck_access(existing_deck, user)
    owner_email = _resolve_owner_email(
        existing_deck.owner_email if existing_deck else None, user
    )
    shared_with = list(existing_deck.shared_with) if existing_deck is not None else []
    try:
        deck = deck_from_payload(
            normalized_deck_id,
            [slide.model_dump() for slide in payload.slides],
            sections_data=[section.model_dump() for section in payload.sections],
            prompt_style=resolved_prompt_style,
            owner_email=owner_email,
            shared_with=shared_with,
        )
    except InvalidDeckError as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)
        ) from exc
    storage.save_deck(deck)
    _reconcile_deck_processing_artifacts(
        normalized_deck_id,
        storage,
        slide_ids=deck.slide_ids(),
    )
    deck = storage.load_deck(normalized_deck_id)
    return _build_deck_response(deck, storage)


@router.post("/deck/{deck_id}/share", response_model=DeckResponse)
def update_deck_sharing(
    deck_id: str,
    payload: DeckShareRequest,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckResponse:
    normalized_deck_id = _normalize_deck_id(deck_id)
    try:
        deck = storage.load_deck(normalized_deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    if user is not None:
        if deck.owner_email is None:
            deck.owner_email = user.email
        elif deck.owner_email != user.email:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="Only the deck owner can update sharing.",
            )
    deck.shared_with = _normalize_shared_with(payload.sharedWith)
    storage.save_deck(deck)
    deck = storage.load_deck(normalized_deck_id)
    return _build_deck_response(deck, storage)


@router.post(
    "/deck/concatenate",
    response_model=DeckResponse,
    status_code=status.HTTP_201_CREATED,
)
def concatenate_decks(
    payload: ConcatenateDecksRequest,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckResponse:
    new_deck_id = _normalize_deck_id(payload.newDeckId)
    source_ids = [
        _normalize_deck_id(identifier) for identifier in payload.sourceDeckIds
    ]
    if not source_ids:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Select at least one source deck.",
        )
    target_path = storage.root / new_deck_id
    if target_path.exists():
        raise HTTPException(
            status_code=status.HTTP_409_CONFLICT,
            detail=f"Deck {new_deck_id} already exists.",
        )
    try:
        prompt_style = _validate_shared_prompt_style(
            source_ids, storage, context="combine"
        )
        combined_deck, slide_mappings = _concatenate_deck_payload(
            new_deck_id,
            source_ids,
            storage,
            prompt_style=prompt_style,
            interleave_by_index=payload.interleaveByIndex,
            interleave_guide=payload.interleaveGuide,
        )
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    for source_id in source_ids:
        source_deck = storage.load_deck(source_id)
        _assert_deck_access(source_deck, user)
    combined_deck.owner_email = _resolve_owner_email(None, user)
    combined_deck.shared_with = []
    storage.save_deck(combined_deck)
    _copy_concatenated_assets(new_deck_id, source_ids, storage)
    _copy_shared_concatenated_pptx_template(new_deck_id, source_ids, storage)
    _merge_concatenated_layout_payload(new_deck_id, slide_mappings, storage)
    _merge_concatenated_ocr_payload(new_deck_id, slide_mappings, storage)
    _reconcile_deck_processing_artifacts(
        new_deck_id,
        storage,
        slide_ids=combined_deck.slide_ids(),
    )
    deck = storage.load_deck(new_deck_id)
    archived_sources: list[str] = []
    if payload.deleteSourceDecks:
        for deck_id in source_ids:
            if deck_id == new_deck_id:
                continue
            if payload.deleteMode != "archive":
                LOGGER.warning(
                    "Delete mode %s not supported for deck %s. Skipping.",
                    payload.deleteMode,
                    deck_id,
                )
                continue
            try:
                if storage.archive_deck(deck_id):
                    archived_sources.append(deck_id)
            except DeckNotFoundError as exc:
                LOGGER.warning("Failed to archive deck %s: %s", deck_id, exc)
    return _build_deck_response(deck, storage, archived_source_decks=archived_sources)


@router.post("/deck/{deck_id}/print")
def enqueue_print_deck(
    deck_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> JSONResponse:
    if not HAS_PLAYWRIGHT:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="PDF export requires Playwright to be installed on the server.",
        )
    normalized = _normalize_deck_id(deck_id)
    try:
        deck = storage.load_deck(normalized)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    if not deck.slides:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail="Deck is empty."
        )
    job_id = _start_print_job(normalized)
    return JSONResponse(
        {"jobId": job_id, "status": "pending"}, status_code=status.HTTP_202_ACCEPTED
    )


@router.get("/deck/print/{job_id}")
def get_print_job(
    job_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> dict[str, Any]:
    job = _get_print_job_or_404(job_id)
    try:
        deck = storage.load_deck(job.deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    output_ready = bool(_resolve_print_job_output(job))
    payload: dict[str, Any] = {
        "jobId": job.job_id,
        "status": (
            job.status if job.status != "succeeded" or output_ready else "running"
        ),
    }
    if job.detail:
        payload["detail"] = job.detail
    if job.status == "succeeded" and output_ready:
        payload["downloadUrl"] = f"/slides/deck/print/{job.job_id}/download"
    return payload


@router.get("/deck/print/{job_id}/download")
def download_print_job(
    job_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> Response:
    job = _get_print_job_or_404(job_id)
    try:
        deck = storage.load_deck(job.deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    output_path = _resolve_print_job_output(job)
    if job.status != "succeeded" or output_path is None:
        return JSONResponse({"status": "running"}, status_code=status.HTTP_202_ACCEPTED)
    headers = {"Content-Disposition": f'attachment; filename="{job.deck_id}.pdf"'}
    return StreamingResponse(
        output_path.open("rb"), media_type="application/pdf", headers=headers
    )


@router.post("/deck/{deck_id}/export-pptx")
def enqueue_export_pptx(
    deck_id: str,
    source: Literal["rendered", "template"] = "rendered",
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> JSONResponse:
    _ensure_pptx_available()
    normalized = _normalize_deck_id(deck_id)
    try:
        deck = storage.load_deck(normalized)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    if not deck.slides:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail="Deck is empty."
        )
    job_id = _start_pptx_job(normalized, source=source)
    return JSONResponse(
        {"jobId": job_id, "status": "pending", "source": source},
        status_code=status.HTTP_202_ACCEPTED,
    )


@router.get("/deck/export-pptx/{job_id}")
def get_pptx_job(
    job_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> dict[str, Any]:
    job = _get_pptx_job_or_404(job_id)
    try:
        deck = storage.load_deck(job.deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    output_ready = bool(_resolve_pptx_job_output(job))
    payload: dict[str, Any] = {
        "jobId": job.job_id,
        "status": (
            job.status if job.status != "succeeded" or output_ready else "running"
        ),
        "source": job.source,
    }
    if job.detail:
        payload["detail"] = job.detail
    if job.status == "succeeded" and output_ready:
        payload["downloadUrl"] = f"/slides/deck/export-pptx/{job.job_id}/download"
    return payload


@router.get("/deck/export-pptx/{job_id}/download")
def download_pptx_job(
    job_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> Response:
    job = _get_pptx_job_or_404(job_id)
    try:
        deck = storage.load_deck(job.deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    output_path = _resolve_pptx_job_output(job)
    if job.status != "succeeded" or output_path is None:
        return JSONResponse({"status": "running"}, status_code=status.HTTP_202_ACCEPTED)
    filename = f"{job.deck_id}.pptx"
    headers = {"Content-Disposition": f'attachment; filename="{filename}"'}
    return StreamingResponse(
        output_path.open("rb"),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers,
    )


if HAS_MULTIPART:

    @router.post(
        "/deck/upload", response_model=DeckResponse, status_code=status.HTTP_410_GONE
    )
    async def upload_deck_unavailable() -> DeckResponse:
        raise HTTPException(
            status_code=status.HTTP_410_GONE,
            detail="HTML deck uploads are no longer supported. Use the PDF/image upload instead.",
        )

    @router.post(
        "/deck/upload-pdf",
        response_model=DeckResponse,
        status_code=status.HTTP_201_CREATED,
    )
    async def upload_pdf_deck(
        deck_id: str = Form(..., alias="deckId"),
        prompt_style: str = Form(DEFAULT_PROMPT_STYLE_KEY, alias="promptStyle"),
        lang: str = Form("eng", alias="lang"),
        run_ocr: bool = Form(True, alias="runOcr"),
        pptx_template_id: str | None = Form(default=None, alias="pptxTemplateId"),
        use_uniform_template: bool = Form(False, alias="useUniformTemplate"),
        file: UploadFile = File(..., alias="file"),
        user: AuthenticatedUser | None = Depends(require_authenticated_user),
        storage: DeckStorage = Depends(get_storage),
    ) -> DeckResponse:
        deck_id = _normalize_deck_id(deck_id)
        resolved_prompt_style = _resolve_prompt_style_strict(prompt_style)
        resolved_lang = _resolve_ocr_lang(lang, fallback="eng")
        if file is None:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Upload a PDF or PNG/JPG file.",
            )
        deck_path = storage.root / deck_id
        if deck_path.exists():
            raise HTTPException(
                status_code=status.HTTP_409_CONFLICT,
                detail=f"Deck {deck_id} already exists.",
            )
        deck_path.mkdir(parents=True, exist_ok=False)
        try:
            file_bytes = await file.read()
            await file.close()
            if not file_bytes:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Uploaded PDF or image is empty.",
                )
            is_pdf = _is_pdf_upload(file)
            is_image = _is_supported_image_upload(file)
            owner_email = _resolve_owner_email(None, user)
            shared_with: list[str] = []
            if is_pdf:
                _render_pdf_deck(
                    deck_id,
                    deck_path,
                    file_bytes,
                    storage,
                    prompt_style=resolved_prompt_style,
                    owner_email=owner_email,
                    shared_with=shared_with,
                )
            elif is_image:
                _render_image_deck(
                    deck_id,
                    deck_path,
                    file_bytes,
                    storage,
                    file.filename,
                    prompt_style=resolved_prompt_style,
                    owner_email=owner_email,
                    shared_with=shared_with,
                )
            else:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail="Upload a PDF or PNG/JPG file.",
                )
            apply_saved_pptx_template_to_deck(
                storage.root,
                owner_email,
                deck_path=deck_path,
                template_id=pptx_template_id,
                use_uniform_template=use_uniform_template,
            )
            deck = storage.load_deck(deck_id)
            if run_ocr:
                _enqueue_deck_ocr_job(deck, lang=resolved_lang)
        except (DeckNotFoundError, InvalidDeckError) as exc:
            shutil.rmtree(deck_path, ignore_errors=True)
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)
            ) from exc
        except SlideOcrEngineUnavailableError as exc:
            shutil.rmtree(deck_path, ignore_errors=True)
            detail = _ocr_engine_unavailable_detail(lang=resolved_lang, exc=exc)
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail=detail,
            ) from exc
        except (ValueError, UnidentifiedImageError) as exc:
            shutil.rmtree(deck_path, ignore_errors=True)
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=str(exc),
            ) from exc
        except FileNotFoundError as exc:
            shutil.rmtree(deck_path, ignore_errors=True)
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=str(exc),
            ) from exc
        except HTTPException:
            shutil.rmtree(deck_path, ignore_errors=True)
            raise
        except Exception as exc:  # pragma: no cover - defensive
            shutil.rmtree(deck_path, ignore_errors=True)
            LOGGER.exception("Unexpected error while uploading deck")
            detail = _format_upload_failure_detail(exc)
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=detail,
            ) from exc
        return _build_deck_response(deck, storage)

else:  # pragma: no cover - optional dependency fallback

    @router.post("/deck/upload", response_model=DeckResponse)
    async def upload_deck_unavailable() -> DeckResponse:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="File uploads require the 'python-multipart' package to be installed.",
        )

    @router.post("/deck/upload-pdf", response_model=DeckResponse)
    async def upload_pdf_deck_unavailable() -> DeckResponse:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="File uploads require the 'python-multipart' package to be installed.",
        )


def _get_or_create_ocr_payload(
    deck: Deck,
    storage: DeckStorage,
    *,
    lang: str = "eng",
    include_bboxes: bool = True,
    force_rebuild: bool = False,
    progress_callback: Callable[[int, int], None] | None = None,
    event_callback: Callable[[str, dict[str, object] | None], None] | None = None,
) -> dict[str, object] | None:
    def _safe_report_progress(built: int, total: int) -> None:
        if progress_callback is None:
            return
        try:
            built_int = max(0, int(built))
            total_int = max(0, int(total))
            if total_int < built_int:
                total_int = built_int
            progress_callback(built_int, total_int)
        except (TypeError, ValueError, OSError) as exc:
            LOGGER.warning("Ignoring OCR progress callback error: %s", exc)

    def _report_done(payload: dict[str, object] | None) -> None:
        if progress_callback is None or payload is None:
            return
        slides = (
            payload.get("slides") if isinstance(payload.get("slides"), list) else []
        )
        slide_count = len(slides)
        _safe_report_progress(slide_count, slide_count)

    def _safe_report_event(
        event: str,
        details: dict[str, object],
    ) -> None:
        if event_callback is None:
            return
        event_name = str(event or "").strip()
        if not event_name:
            return
        payload_details: dict[str, object] | None = (
            details if isinstance(details, dict) else None
        )
        try:
            event_callback(event_name, payload_details)
        except (TypeError, ValueError, OSError) as exc:
            LOGGER.warning("Ignoring OCR event callback error: %s", exc)

    deck_path = storage.root / deck.deck_id
    cached_layout = _load_raw_layout_payload(deck.deck_id, storage, lang=lang)
    if force_rebuild:
        ensure_kwargs: dict[str, object] = {
            "lang": lang,
            "include_bboxes": include_bboxes,
            "layout_payload": cached_layout,
        }
        if progress_callback is not None:
            ensure_kwargs["progress_callback"] = _safe_report_progress
        if event_callback is not None:
            ensure_kwargs["event_callback"] = _safe_report_event
        payload = ensure_deck_ocr_payload(
            deck,
            deck_path,
            **ensure_kwargs,
        )
        if payload:
            storage.save_ocr_payload(deck.deck_id, payload)
        _report_done(payload)
        return payload

    cached_payload = storage.load_ocr_payload(deck.deck_id)
    if cached_payload and _ocr_payload_uses_layout_guided_strategy(cached_payload):
        ensure_kwargs: dict[str, object] = {
            "lang": lang,
            "include_bboxes": include_bboxes,
            "cached_payload": cached_payload,
            "layout_payload": cached_layout,
        }
        if progress_callback is not None:
            ensure_kwargs["progress_callback"] = _safe_report_progress
        if event_callback is not None:
            ensure_kwargs["event_callback"] = _safe_report_event
        try:
            normalized = ensure_deck_ocr_payload(
                deck,
                deck_path,
                **ensure_kwargs,
            )
        except (
            SlideOcrEngineUnavailableError,
            RuntimeError,
            OSError,
            ValueError,
        ) as exc:
            LOGGER.warning(
                "OCR payload completion failed for deck %s; falling back to cached OCR payload: %s",
                deck.deck_id,
                exc,
            )
            normalized = normalize_ocr_payload(
                cached_payload,
                deck_id=deck.deck_id,
                lang=lang,
            )
        if normalized and normalized != cached_payload:
            storage.save_ocr_payload(deck.deck_id, normalized)
        _report_done(normalized)
        return normalized

    has_pdf = (deck_path / "source.pdf").exists()
    has_images = any("<img" in (slide.body_html or "") for slide in deck.slides)
    if not has_pdf and not has_images:
        LOGGER.info(
            "OCR skipped for deck %s: no PDF or slide images available.", deck.deck_id
        )
        return None
    ensure_kwargs = {
        "lang": lang,
        "include_bboxes": include_bboxes,
        "layout_payload": cached_layout,
    }
    if progress_callback is not None:
        ensure_kwargs["progress_callback"] = _safe_report_progress
    if event_callback is not None:
        ensure_kwargs["event_callback"] = _safe_report_event
    payload = ensure_deck_ocr_payload(
        deck,
        deck_path,
        **ensure_kwargs,
    )
    if payload:
        storage.save_ocr_payload(deck.deck_id, payload)
    _report_done(payload)
    return payload


def _ocr_payload_uses_layout_guided_strategy(payload: dict[str, object] | None) -> bool:
    if not isinstance(payload, dict):
        return False
    strategy = str(
        payload.get("ocr_strategy") or payload.get("ocrStrategy") or ""
    ).strip()
    if not strategy:
        return False
    normalized_strategy = strategy.lower()
    normalized_expected = OCR_STRATEGY_LAYOUT_GUIDED.strip().lower()
    return normalized_strategy == normalized_expected or normalized_strategy.startswith(
        "layout_guided_"
    )


def _load_cached_ocr_payload(
    deck_id: str,
    storage: DeckStorage,
    *,
    lang: str | None = None,
) -> dict[str, object] | None:
    payload = storage.load_ocr_payload(deck_id)
    if payload is None:
        return None
    resolved_lang = str(payload.get("lang") or lang or "eng")
    normalized = normalize_ocr_payload(payload, deck_id=deck_id, lang=resolved_lang)
    if normalized != payload:
        storage.save_ocr_payload(deck_id, normalized)
    return normalized


def _normalize_layout_slide_payload(
    slide_payload: dict[str, object],
) -> dict[str, object] | None:
    slide_id = str(
        slide_payload.get("slide_id") or slide_payload.get("slideId") or ""
    ).strip()
    if not slide_id:
        return None
    raw_blocks = (
        slide_payload.get("blocks")
        if isinstance(slide_payload.get("blocks"), list)
        else []
    )
    normalized_blocks: list[dict[str, object]] = []
    for raw_block in raw_blocks:
        if not isinstance(raw_block, dict):
            continue
        normalized_block: dict[str, object] = {
            "blockId": str(
                raw_block.get("block_id")
                or raw_block.get("blockId")
                or raw_block.get("id")
                or ""
            ),
            "type": str(raw_block.get("type") or "unknown"),
            "detectedType": str(
                raw_block.get("detected_type") or raw_block.get("detectedType") or ""
            ).strip(),
            "text": str(raw_block.get("text") or ""),
            "items": [
                str(item).strip()
                for item in (
                    raw_block.get("items")
                    if isinstance(raw_block.get("items"), list)
                    else []
                )
                if str(item).strip()
            ],
        }
        group_id = str(
            raw_block.get("group_id") or raw_block.get("groupId") or ""
        ).strip()
        if group_id:
            normalized_block["groupId"] = group_id
        group_kind = str(
            raw_block.get("group_kind") or raw_block.get("groupKind") or ""
        ).strip()
        if group_kind:
            normalized_block["groupKind"] = group_kind
        parent_id = str(
            raw_block.get("parent_id") or raw_block.get("parentId") or ""
        ).strip()
        if parent_id:
            normalized_block["parentId"] = parent_id
        list_level = raw_block.get("list_level")
        if list_level is None:
            list_level = raw_block.get("listLevel")
        if isinstance(list_level, int) and list_level >= 0:
            normalized_block["listLevel"] = list_level
        reading_order = raw_block.get("reading_order")
        if reading_order is None:
            reading_order = raw_block.get("readingOrder")
        if isinstance(reading_order, int) and reading_order >= 0:
            normalized_block["readingOrder"] = reading_order
        render_mode = str(
            raw_block.get("render_mode") or raw_block.get("renderMode") or ""
        ).strip()
        if render_mode:
            normalized_block["renderMode"] = render_mode
        visual_text = str(
            raw_block.get("visual_text") or raw_block.get("visualText") or ""
        ).strip()
        if visual_text:
            normalized_block["visualText"] = visual_text
        visual_items = [
            str(item).strip()
            for item in (
                raw_block.get("visual_items")
                if isinstance(raw_block.get("visual_items"), list)
                else (
                    raw_block.get("visualItems")
                    if isinstance(raw_block.get("visualItems"), list)
                    else []
                )
            )
            if str(item).strip()
        ]
        if visual_items:
            normalized_block["visualItems"] = visual_items
        visual_lines = [
            line
            for line in (
                raw_block.get("visual_lines")
                if isinstance(raw_block.get("visual_lines"), list)
                else (
                    raw_block.get("visualLines")
                    if isinstance(raw_block.get("visualLines"), list)
                    else []
                )
            )
            if isinstance(line, dict) and str(line.get("text") or "").strip()
        ]
        if visual_lines:
            normalized_block["visualLines"] = visual_lines
        bbox = raw_block.get("bbox")
        if isinstance(bbox, dict):
            normalized_block["bbox"] = bbox
        confidence = raw_block.get("confidence")
        if isinstance(confidence, (int, float)):
            normalized_block["confidence"] = float(confidence)
        table_model = raw_block.get("table_model") or raw_block.get("tableModel")
        if isinstance(table_model, dict):
            normalized_block["tableModel"] = table_model
        normalized_blocks.append(normalized_block)
    raw_figure_regions = (
        slide_payload.get("figure_regions")
        if isinstance(slide_payload.get("figure_regions"), list)
        else (
            slide_payload.get("figureRegions")
            if isinstance(slide_payload.get("figureRegions"), list)
            else []
        )
    )
    figure_regions = [
        region
        for region in raw_figure_regions
        if isinstance(region, dict)
        and all(
            isinstance(region.get(key), (int, float)) for key in ("x", "y", "w", "h")
        )
    ]
    return {
        "slideId": slide_id,
        "slideNumber": int(
            slide_payload.get("slide_number") or slide_payload.get("slideNumber") or 0
        ),
        "pageNumber": int(
            slide_payload.get("page_number") or slide_payload.get("pageNumber") or 0
        ),
        "assetPath": str(
            slide_payload.get("asset_path") or slide_payload.get("assetPath") or ""
        ),
        "blocks": normalized_blocks,
        "titleText": str(
            slide_payload.get("title_text") or slide_payload.get("titleText") or ""
        ),
        "bulletTexts": [
            str(item).strip()
            for item in (
                slide_payload.get("bullet_texts")
                if isinstance(slide_payload.get("bullet_texts"), list)
                else (
                    slide_payload.get("bulletTexts")
                    if isinstance(slide_payload.get("bulletTexts"), list)
                    else []
                )
            )
            if str(item).strip()
        ],
        "figureRegions": figure_regions,
    }


def _normalize_layout_payload(
    payload: dict[str, object],
    *,
    deck_id: str,
    lang: str,
) -> dict[str, object]:
    raw_slides = (
        payload.get("slides") if isinstance(payload.get("slides"), list) else []
    )
    normalized_slides = [
        normalized
        for raw_slide in raw_slides
        if isinstance(raw_slide, dict)
        for normalized in [_normalize_layout_slide_payload(raw_slide)]
        if normalized is not None
    ]
    return {
        "deckId": deck_id,
        "lang": str(payload.get("lang") or lang or "eng"),
        "generatedAt": str(
            payload.get("generated_at")
            or payload.get("generatedAt")
            or datetime.now(UTC).isoformat()
        ),
        "slides": normalized_slides,
    }


def _merge_layout_slide_with_ocr(
    layout_slide: dict[str, object],
    ocr_slide: dict[str, object] | None,
) -> dict[str, object]:
    if not isinstance(ocr_slide, dict):
        return layout_slide
    ocr_blocks = (
        ocr_slide.get("blocks") if isinstance(ocr_slide.get("blocks"), list) else []
    )
    ocr_blocks_by_id: dict[str, dict[str, object]] = {}
    for raw_block in ocr_blocks:
        if not isinstance(raw_block, dict):
            continue
        block_id = str(
            raw_block.get("block_id")
            or raw_block.get("blockId")
            or raw_block.get("id")
            or ""
        ).strip()
        if block_id:
            ocr_blocks_by_id[block_id] = raw_block

    merged_blocks: list[dict[str, object]] = []
    for raw_block in (
        layout_slide.get("blocks")
        if isinstance(layout_slide.get("blocks"), list)
        else []
    ):
        if not isinstance(raw_block, dict):
            continue
        merged_block = dict(raw_block)
        block_id = str(
            raw_block.get("blockId")
            or raw_block.get("block_id")
            or raw_block.get("id")
            or ""
        ).strip()
        matching_block = ocr_blocks_by_id.get(block_id)
        if isinstance(matching_block, dict):
            matching_type = str(
                matching_block.get("type") or matching_block.get("blockType") or ""
            ).strip()
            if matching_type:
                merged_block["type"] = matching_type
            matching_detected_type = str(
                matching_block.get("detected_type")
                or matching_block.get("detectedType")
                or ""
            ).strip()
            if matching_detected_type:
                merged_block["detectedType"] = matching_detected_type
            text = str(matching_block.get("text") or "").strip()
            if text:
                merged_block["text"] = text
            items = [
                str(item).strip()
                for item in (
                    matching_block.get("items")
                    if isinstance(matching_block.get("items"), list)
                    else []
                )
                if str(item).strip()
            ]
            if items:
                merged_block["items"] = items
            confidence = matching_block.get("confidence")
            if isinstance(confidence, (int, float)):
                merged_block["confidence"] = float(confidence)
            table_model = matching_block.get("table_model")
            if not isinstance(table_model, dict):
                table_model = matching_block.get("tableModel")
            if isinstance(table_model, dict):
                merged_block["tableModel"] = table_model
            audit_status = str(
                matching_block.get("audit_status")
                or matching_block.get("auditStatus")
                or ""
            ).strip()
            if audit_status:
                merged_block["auditStatus"] = audit_status
            audit_reason = str(
                matching_block.get("audit_reason")
                or matching_block.get("auditReason")
                or ""
            ).strip()
            if audit_reason:
                merged_block["auditReason"] = audit_reason
            audit_suggested_text = str(
                matching_block.get("audit_suggested_text")
                or matching_block.get("auditSuggestedText")
                or ""
            ).strip()
            if audit_suggested_text:
                merged_block["auditSuggestedText"] = audit_suggested_text
            visual_status = str(
                matching_block.get("visual_status")
                or matching_block.get("visualStatus")
                or ""
            ).strip()
            if visual_status:
                merged_block["visualStatus"] = visual_status
            visual_reason = str(
                matching_block.get("visual_reason")
                or matching_block.get("visualReason")
                or ""
            ).strip()
            if visual_reason:
                merged_block["visualReason"] = visual_reason
            visual_suggested_text = str(
                matching_block.get("visual_suggested_text")
                or matching_block.get("visualSuggestedText")
                or ""
            ).strip()
            if visual_suggested_text:
                merged_block["visualSuggestedText"] = visual_suggested_text
            visual_confidence = matching_block.get("visual_confidence")
            if not isinstance(visual_confidence, (int, float)):
                visual_confidence = matching_block.get("visualConfidence")
            if isinstance(visual_confidence, (int, float)):
                merged_block["visualConfidence"] = float(visual_confidence)
            visual_text = str(
                matching_block.get("visual_text")
                or matching_block.get("visualText")
                or ""
            ).strip()
            if visual_text:
                merged_block["visualText"] = visual_text
            visual_items = [
                str(item).strip()
                for item in (
                    matching_block.get("visual_items")
                    if isinstance(matching_block.get("visual_items"), list)
                    else (
                        matching_block.get("visualItems")
                        if isinstance(matching_block.get("visualItems"), list)
                        else []
                    )
                )
                if str(item).strip()
            ]
            if visual_items:
                merged_block["visualItems"] = visual_items
            visual_lines = [
                line
                for line in (
                    matching_block.get("visual_lines")
                    if isinstance(matching_block.get("visual_lines"), list)
                    else (
                        matching_block.get("visualLines")
                        if isinstance(matching_block.get("visualLines"), list)
                        else []
                    )
                )
                if isinstance(line, dict) and str(line.get("text") or "").strip()
            ]
            if visual_lines:
                merged_block["visualLines"] = visual_lines
            group_id = str(
                matching_block.get("group_id") or matching_block.get("groupId") or ""
            ).strip()
            if group_id:
                merged_block["groupId"] = group_id
            group_kind = str(
                matching_block.get("group_kind")
                or matching_block.get("groupKind")
                or ""
            ).strip()
            if group_kind:
                merged_block["groupKind"] = group_kind
            parent_id = str(
                matching_block.get("parent_id") or matching_block.get("parentId") or ""
            ).strip()
            if parent_id:
                merged_block["parentId"] = parent_id
            list_level = matching_block.get("list_level")
            if list_level is None:
                list_level = matching_block.get("listLevel")
            if isinstance(list_level, int) and list_level >= 0:
                merged_block["listLevel"] = list_level
            reading_order = matching_block.get("reading_order")
            if reading_order is None:
                reading_order = matching_block.get("readingOrder")
            if isinstance(reading_order, int) and reading_order >= 0:
                merged_block["readingOrder"] = reading_order
            render_mode = str(
                matching_block.get("render_mode")
                or matching_block.get("renderMode")
                or ""
            ).strip()
            if render_mode:
                merged_block["renderMode"] = render_mode
        merged_blocks.append(merged_block)

    merged_slide = dict(layout_slide)
    merged_slide["blocks"] = merged_blocks
    title_text = str(
        ocr_slide.get("title_text") or ocr_slide.get("titleText") or ""
    ).strip()
    if title_text:
        merged_slide["titleText"] = title_text
    bullet_texts = [
        str(item).strip()
        for item in (
            ocr_slide.get("bullet_texts")
            if isinstance(ocr_slide.get("bullet_texts"), list)
            else (
                ocr_slide.get("bulletTexts")
                if isinstance(ocr_slide.get("bulletTexts"), list)
                else []
            )
        )
        if str(item).strip()
    ]
    if bullet_texts:
        merged_slide["bulletTexts"] = bullet_texts
    return merged_slide


def _merge_layout_payload_with_ocr(
    layout_payload: dict[str, object],
    ocr_payload: dict[str, object] | None,
    *,
    deck_id: str,
    lang: str,
) -> dict[str, object]:
    normalized_layout = _normalize_layout_payload(
        layout_payload, deck_id=deck_id, lang=lang
    )
    if not isinstance(ocr_payload, dict):
        return normalized_layout
    normalized_ocr = normalize_ocr_payload(ocr_payload, deck_id=deck_id, lang=lang)
    ocr_slides = (
        normalized_ocr.get("slides")
        if isinstance(normalized_ocr.get("slides"), list)
        else []
    )
    ocr_by_slide_id = {
        str(slide.get("slide_id") or slide.get("slideId") or "").strip(): slide
        for slide in ocr_slides
        if isinstance(slide, dict)
    }
    merged_slides = [
        _merge_layout_slide_with_ocr(
            slide,
            ocr_by_slide_id.get(
                str(slide.get("slideId") or slide.get("slide_id") or "").strip()
            ),
        )
        for slide in normalized_layout.get("slides", [])
        if isinstance(slide, dict)
    ]
    merged_payload = dict(normalized_layout)
    merged_payload["slides"] = merged_slides
    return merged_payload


def _load_raw_layout_payload(
    deck_id: str,
    storage: DeckStorage,
    *,
    lang: str | None = None,
) -> dict[str, object] | None:
    payload = storage.load_layout_payload(deck_id)
    if payload is None:
        return None
    resolved_lang = str(payload.get("lang") or lang or "eng")
    normalized = _normalize_layout_payload(
        payload,
        deck_id=deck_id,
        lang=resolved_lang,
    )
    if normalized != payload:
        storage.save_layout_payload(deck_id, normalized)
    return normalized


def _build_slide_analysis_payload(
    layout_payload: dict[str, object] | None,
    ocr_payload: dict[str, object] | None,
    *,
    deck_id: str,
    lang: str,
) -> dict[str, object] | None:
    if not isinstance(layout_payload, dict):
        return None
    return _merge_layout_payload_with_ocr(
        layout_payload,
        ocr_payload,
        deck_id=deck_id,
        lang=lang,
    )


def _copy_layout_slide_for_target(
    layout_slide: dict[str, object],
    *,
    target_slide_id: str,
    slide_number: int,
    asset_subdir: str | None = None,
) -> dict[str, object]:
    copied = deepcopy(layout_slide)
    copied["slideId"] = target_slide_id
    copied["slideNumber"] = slide_number
    copied["pageNumber"] = slide_number
    asset_path = str(
        layout_slide.get("assetPath") or layout_slide.get("asset_path") or ""
    ).strip()
    if asset_path:
        copied["assetPath"] = _rewrite_concatenated_asset_path(
            asset_path,
            asset_subdir=asset_subdir,
        )
    return copied


def _copy_ocr_slide_for_target(
    ocr_slide: dict[str, object],
    *,
    target_slide_id: str,
    slide_number: int,
) -> dict[str, object]:
    copied = deepcopy(ocr_slide)
    copied["slide_id"] = target_slide_id
    copied["slide_number"] = slide_number
    copied["page_number"] = slide_number
    return copied


def _reindex_layout_payload_for_slide_order(
    payload: dict[str, object],
    *,
    deck_id: str,
    slide_ids: list[str],
    lang: str,
) -> dict[str, object]:
    normalized = _normalize_layout_payload(payload, deck_id=deck_id, lang=lang)
    slide_map = {
        str(slide.get("slideId") or "").strip(): slide
        for slide in normalized.get("slides", [])
        if isinstance(slide, dict) and str(slide.get("slideId") or "").strip()
    }
    reordered_slides = [
        _copy_layout_slide_for_target(
            slide_map[slide_id], target_slide_id=slide_id, slide_number=index + 1
        )
        for index, slide_id in enumerate(slide_ids)
        if slide_id in slide_map
    ]
    return {
        "deckId": deck_id,
        "lang": str(normalized.get("lang") or lang or "eng"),
        "generatedAt": str(
            normalized.get("generatedAt") or datetime.now(UTC).isoformat()
        ),
        "slides": reordered_slides,
    }


def _reindex_ocr_payload_for_slide_order(
    payload: dict[str, object],
    *,
    deck_id: str,
    slide_ids: list[str],
    lang: str,
) -> dict[str, object]:
    normalized = normalize_ocr_payload(payload, deck_id=deck_id, lang=lang)
    slide_map = {
        str(slide.get("slide_id") or "").strip(): slide
        for slide in normalized.get("slides", [])
        if isinstance(slide, dict) and str(slide.get("slide_id") or "").strip()
    }
    reordered_slides = [
        _copy_ocr_slide_for_target(
            slide_map[slide_id], target_slide_id=slide_id, slide_number=index + 1
        )
        for index, slide_id in enumerate(slide_ids)
        if slide_id in slide_map
    ]
    return normalize_ocr_payload(
        {
            "deck_id": deck_id,
            "lang": str(normalized.get("lang") or lang or "eng"),
            "ocr_strategy": normalized.get("ocr_strategy"),
            "prompt_style": normalized.get("prompt_style"),
            "style_hint": (
                deepcopy(normalized.get("style_hint"))
                if isinstance(normalized.get("style_hint"), dict)
                else None
            ),
            "generated_at": str(
                normalized.get("generated_at") or datetime.now(UTC).isoformat()
            ),
            "slides": reordered_slides,
        },
        deck_id=deck_id,
        lang=lang,
    )


def _reconcile_deck_processing_artifacts(
    deck_id: str,
    storage: DeckStorage,
    *,
    slide_ids: list[str],
    lang: str | None = None,
) -> None:
    resolved_lang = str(lang or "eng")
    raw_layout = storage.load_layout_payload(deck_id)
    normalized_layout: dict[str, object] | None = None
    if isinstance(raw_layout, dict):
        normalized_layout = _reindex_layout_payload_for_slide_order(
            raw_layout,
            deck_id=deck_id,
            slide_ids=slide_ids,
            lang=resolved_lang,
        )
        resolved_lang = str(normalized_layout.get("lang") or resolved_lang)
        storage.save_layout_payload(deck_id, normalized_layout)

    raw_ocr = storage.load_ocr_payload(deck_id)
    normalized_ocr: dict[str, object] | None = None
    if isinstance(raw_ocr, dict):
        normalized_ocr = _reindex_ocr_payload_for_slide_order(
            raw_ocr,
            deck_id=deck_id,
            slide_ids=slide_ids,
            lang=resolved_lang,
        )
        resolved_lang = str(normalized_ocr.get("lang") or resolved_lang)
        storage.save_ocr_payload(deck_id, normalized_ocr)

    if normalized_layout is None:
        return
    merged_analysis = _build_slide_analysis_payload(
        normalized_layout,
        normalized_ocr,
        deck_id=deck_id,
        lang=resolved_lang,
    )
    if merged_analysis is not None:
        storage.save_slide_analysis_payload(deck_id, merged_analysis)


def _concatenated_asset_subdir(source_deck_id: str) -> str:
    return str(source_deck_id or "").strip()


def _rewrite_concatenated_asset_path(
    asset_path: str,
    *,
    asset_subdir: str | None,
) -> str:
    normalized_relative = _normalize_asset_path(asset_path)
    if not normalized_relative:
        return ""
    if asset_subdir:
        return f"assets/{asset_subdir}/{normalized_relative}"
    return f"assets/{normalized_relative}"


def _build_and_persist_deck_processing_artifacts(
    deck: Deck,
    storage: DeckStorage,
    *,
    lang: str,
) -> None:
    deck_path = storage.root / deck.deck_id
    built_layout = build_deck_layout_payload(deck, deck_path, lang=lang)
    normalized_layout = _normalize_layout_payload(
        built_layout,
        deck_id=deck.deck_id,
        lang=lang,
    )
    storage.save_layout_payload(deck.deck_id, normalized_layout)
    ocr_payload = build_deck_ocr_payload(
        deck,
        deck_path,
        lang=lang,
        include_bboxes=True,
        layout_payload=normalized_layout,
        pdf_path=_resolve_ocr_pdf_path(None, deck_path),
    )
    normalized_ocr = normalize_ocr_payload(ocr_payload, deck_id=deck.deck_id, lang=lang)
    storage.save_ocr_payload(deck.deck_id, normalized_ocr)
    merged_analysis = _build_slide_analysis_payload(
        normalized_layout,
        normalized_ocr,
        deck_id=deck.deck_id,
        lang=lang,
    )
    if merged_analysis is not None:
        storage.save_slide_analysis_payload(deck.deck_id, merged_analysis)


def _copy_imported_deck_processing_artifacts(
    *,
    source_deck_id: str,
    source_slide_id: str,
    target_deck_id: str,
    target_slide_id: str,
    storage: DeckStorage,
) -> None:
    source_layout = _load_raw_layout_payload(source_deck_id, storage)
    target_layout = _load_raw_layout_payload(target_deck_id, storage)
    if isinstance(source_layout, dict):
        source_slide_map = {
            str(slide.get("slideId") or "").strip(): slide
            for slide in source_layout.get("slides", [])
            if isinstance(slide, dict) and str(slide.get("slideId") or "").strip()
        }
        source_layout_slide = source_slide_map.get(source_slide_id)
        if isinstance(source_layout_slide, dict):
            resolved_lang = str(
                (target_layout or {}).get("lang") or source_layout.get("lang") or "eng"
            )
            target_slides = [
                slide
                for slide in (
                    target_layout.get("slides")
                    if isinstance(target_layout, dict)
                    and isinstance(target_layout.get("slides"), list)
                    else []
                )
                if isinstance(slide, dict)
                and str(slide.get("slideId") or "").strip() != target_slide_id
            ]
            target_slides.append(
                _copy_layout_slide_for_target(
                    source_layout_slide,
                    target_slide_id=target_slide_id,
                    slide_number=len(target_slides) + 1,
                )
            )
            storage.save_layout_payload(
                target_deck_id,
                _normalize_layout_payload(
                    {
                        "deckId": target_deck_id,
                        "lang": resolved_lang,
                        "generatedAt": datetime.now(UTC).isoformat(),
                        "slides": target_slides,
                    },
                    deck_id=target_deck_id,
                    lang=resolved_lang,
                ),
            )

    source_ocr = _load_cached_ocr_payload(source_deck_id, storage)
    target_ocr = _load_cached_ocr_payload(target_deck_id, storage)
    if isinstance(source_ocr, dict):
        source_slide_map = {
            str(slide.get("slide_id") or "").strip(): slide
            for slide in source_ocr.get("slides", [])
            if isinstance(slide, dict) and str(slide.get("slide_id") or "").strip()
        }
        source_ocr_slide = source_slide_map.get(source_slide_id)
        if isinstance(source_ocr_slide, dict):
            resolved_lang = str(
                (target_ocr or {}).get("lang") or source_ocr.get("lang") or "eng"
            )
            target_slides = [
                slide
                for slide in (
                    target_ocr.get("slides")
                    if isinstance(target_ocr, dict)
                    and isinstance(target_ocr.get("slides"), list)
                    else []
                )
                if isinstance(slide, dict)
                and str(slide.get("slide_id") or "").strip() != target_slide_id
            ]
            target_slides.append(
                _copy_ocr_slide_for_target(
                    source_ocr_slide,
                    target_slide_id=target_slide_id,
                    slide_number=len(target_slides) + 1,
                )
            )
            storage.save_ocr_payload(
                target_deck_id,
                normalize_ocr_payload(
                    {
                        "deck_id": target_deck_id,
                        "lang": resolved_lang,
                        "ocr_strategy": source_ocr.get("ocr_strategy"),
                        "prompt_style": source_ocr.get("prompt_style"),
                        "style_hint": (
                            deepcopy(source_ocr.get("style_hint"))
                            if isinstance(source_ocr.get("style_hint"), dict)
                            else None
                        ),
                        "generated_at": datetime.now(UTC).isoformat(),
                        "slides": target_slides,
                    },
                    deck_id=target_deck_id,
                    lang=resolved_lang,
                ),
            )


def _merge_concatenated_layout_payload(
    new_deck_id: str,
    slide_mappings: list[ConcatenatedSlideMapping],
    storage: DeckStorage,
) -> None:
    if not slide_mappings:
        return
    source_slide_payloads: dict[str, dict[str, dict[str, object]]] = {}
    resolved_lang: str | None = None
    for source_deck_id in {mapping.source_deck_id for mapping in slide_mappings}:
        payload = _load_raw_layout_payload(source_deck_id, storage)
        if payload is None:
            continue
        if resolved_lang is None:
            resolved_lang = str(payload.get("lang") or "eng")
        slide_map = {
            str(slide.get("slideId") or "").strip(): slide
            for slide in payload.get("slides", [])
            if isinstance(slide, dict) and str(slide.get("slideId") or "").strip()
        }
        source_slide_payloads[source_deck_id] = slide_map
    if not source_slide_payloads:
        return
    merged_slides: list[dict[str, object]] = []
    for index, mapping in enumerate(slide_mappings):
        slide_payloads = source_slide_payloads.get(mapping.source_deck_id)
        if not slide_payloads:
            continue
        source_slide = slide_payloads.get(mapping.source_slide_id)
        if not isinstance(source_slide, dict):
            continue
        merged_slides.append(
            _copy_layout_slide_for_target(
                source_slide,
                target_slide_id=mapping.target_slide_id,
                slide_number=index + 1,
                asset_subdir=_concatenated_asset_subdir(mapping.source_deck_id),
            )
        )
    if not merged_slides:
        return
    storage.save_layout_payload(
        new_deck_id,
        _normalize_layout_payload(
            {
                "deckId": new_deck_id,
                "lang": resolved_lang or "eng",
                "generatedAt": datetime.now(UTC).isoformat(),
                "slides": merged_slides,
            },
            deck_id=new_deck_id,
            lang=resolved_lang or "eng",
        ),
    )


def _load_cached_slide_analysis_payload(
    deck_id: str,
    storage: DeckStorage,
    *,
    lang: str | None = None,
) -> dict[str, object] | None:
    cached_analysis = storage.load_slide_analysis_payload(deck_id)
    raw_layout = _load_raw_layout_payload(deck_id, storage, lang=lang)
    if raw_layout is None:
        if cached_analysis is None:
            return None
        resolved_lang = str(cached_analysis.get("lang") or lang or "eng")
        normalized_analysis = _normalize_layout_payload(
            cached_analysis,
            deck_id=deck_id,
            lang=resolved_lang,
        )
        if normalized_analysis != cached_analysis:
            storage.save_slide_analysis_payload(deck_id, normalized_analysis)
        return normalized_analysis

    resolved_lang = str(raw_layout.get("lang") or lang or "eng")
    ocr_payload = _load_cached_ocr_payload(deck_id, storage, lang=resolved_lang)
    merged_analysis = _build_slide_analysis_payload(
        raw_layout,
        ocr_payload,
        deck_id=deck_id,
        lang=resolved_lang,
    )
    if merged_analysis is None:
        return None
    if merged_analysis != cached_analysis:
        storage.save_slide_analysis_payload(deck_id, merged_analysis)
    return merged_analysis


def _run_deck_layout_job(deck_id: str, *, lang: str) -> None:
    try:
        deck = _storage.load_deck(deck_id)
        deck_path = _storage.root / deck_id
        total_pages = sum(
            1 for slide in deck.slides if "<img" in (slide.body_html or "")
        )
        _set_layout_progress(
            deck_id,
            status="running",
            built_pages=0,
            total_pages=total_pages,
            message=_build_layout_processing_message(0, total_pages),
            step="layout",
        )

        def _progress_callback(built: int, total: int) -> None:
            _set_layout_progress(
                deck_id,
                status="running",
                built_pages=built,
                total_pages=total,
                message=_build_layout_processing_message(built, total),
                step="layout",
            )

        built_payload = build_deck_layout_payload(
            deck,
            deck_path,
            lang=lang,
            progress_callback=_progress_callback,
        )
        normalized = _normalize_layout_payload(
            built_payload,
            deck_id=deck_id,
            lang=lang,
        )
        _storage.save_layout_payload(deck_id, normalized)
        merged_analysis = _build_slide_analysis_payload(
            normalized,
            _load_cached_ocr_payload(deck_id, _storage, lang=lang),
            deck_id=deck_id,
            lang=lang,
        )
        if merged_analysis is not None:
            _storage.save_slide_analysis_payload(deck_id, merged_analysis)
        slide_count = len(normalized.get("slides", []))
        _set_layout_progress(
            deck_id,
            status="completed",
            built_pages=slide_count,
            total_pages=slide_count,
            message="Layout payload saved.",
            step="complete",
            last_completed_step="layout",
        )
    except SlideOcrEngineUnavailableError as exc:
        detail = _ocr_engine_unavailable_detail(lang=lang, exc=exc)
        LOGGER.error(
            "Deck layout detection failed because PaddleOCR is unavailable.",
            exc_info=exc,
        )
        _set_layout_progress(deck_id, status="failed", error=detail, step="layout")
    except Exception as exc:  # noqa: BLE001
        detail = str(exc) or "Deck layout analysis failed unexpectedly."
        LOGGER.exception(
            "Deck layout analysis failed unexpectedly for deck %s", deck_id
        )
        _set_layout_progress(deck_id, status="failed", error=detail, step="layout")


def _run_deck_ocr_job(deck_id: str, *, lang: str) -> None:
    try:
        deck_ocr_endpoint(
            deck_id,
            DeckOcrRequest(lang=lang),
            user=None,
            storage=get_storage(),
        )
    except HTTPException as exc:
        LOGGER.warning(
            "Background deck processing finished with HTTP %s for deck %s: %s",
            exc.status_code,
            deck_id,
            exc.detail,
        )
    except Exception as exc:  # noqa: BLE001
        detail = str(exc) or "Deck OCR failed unexpectedly."
        LOGGER.exception("Background deck OCR failed unexpectedly for deck %s", deck_id)
        _set_ocr_progress(
            deck_id,
            status="failed",
            error=detail,
            lang=lang,
            step="ocr",
        )


@router.post("/deck/{deck_id}/ocr", response_model=DeckOcrResponse)
def deck_ocr_endpoint(
    deck_id: str,
    payload: DeckOcrRequest,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckOcrResponse:
    """Run OCR across all slides in a deck, reusing cached results when available."""
    notify_context = _build_deck_processing_notify_context(user, ocr_lang=payload.lang)
    started_at = time.perf_counter()
    try:
        deck = storage.load_deck(deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)

    cached_payload = storage.load_ocr_payload(deck_id)
    cached_prompt_style = (
        str(
            cached_payload.get("prompt_style")
            or cached_payload.get("promptStyle")
            or ""
        )
        .strip()
        .lower()
        if isinstance(cached_payload, dict)
        else ""
    )
    cached_style_hint = (
        cached_payload.get("style_hint") or cached_payload.get("styleHint")
        if isinstance(cached_payload, dict)
        else None
    )
    style_cache_matches = cached_prompt_style == str(
        deck.prompt_style
    ).strip().lower() and isinstance(cached_style_hint, dict)
    cached_layout = _load_raw_layout_payload(deck_id, storage, lang=payload.lang)
    if (
        cached_payload
        and not payload.force
        and cached_payload.get("lang") == payload.lang
        and style_cache_matches
        and cached_layout is not None
        and _ocr_payload_uses_layout_guided_strategy(cached_payload)
    ):
        normalized = normalize_ocr_payload(
            cached_payload, deck_id=deck_id, lang=payload.lang
        )
        if normalized != cached_payload:
            storage.save_ocr_payload(deck_id, normalized)
        _set_ocr_progress(
            deck_id,
            status="completed",
            built_pages=len(normalized.get("slides", [])),
            total_pages=len(normalized.get("slides", [])),
            message="Deck processing already available.",
            lang=payload.lang,
            step="complete",
            last_completed_step="ocr",
        )
        return DeckOcrResponse(
            payload=DeckOcrPayload.model_validate(normalized), cached=True
        )

    deck_path = storage.root / deck_id
    has_pdf = (deck_path / "source.pdf").exists()
    has_images = any("<img" in (slide.body_html or "") for slide in deck.slides)
    if not has_pdf and not has_images:
        message = "OCR skipped: no PDF or slide images available."
        _set_ocr_progress(
            deck_id,
            status="skipped",
            message=message,
            lang=payload.lang,
            step="complete",
        )
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=message)
    pdf_path = _resolve_ocr_pdf_path(payload.pdfPath, deck_path)
    normalized_layout = cached_layout
    did_work = False

    if cached_layout is None or payload.force:
        _set_ocr_progress(
            deck_id,
            status="running",
            built_pages=0,
            total_pages=0,
            message=_build_layout_processing_message(0, 0),
            lang=payload.lang,
            step="layout",
        )

        def _layout_progress_callback(built: int, total: int) -> None:
            _set_ocr_progress(
                deck_id,
                status="running",
                built_pages=built,
                total_pages=total,
                message=_build_layout_processing_message(built, total),
                step="layout",
            )

        try:
            built_layout = build_deck_layout_payload(
                deck,
                deck_path,
                lang=payload.lang,
                progress_callback=_layout_progress_callback,
            )
        except SlideOcrEngineUnavailableError as exc:
            LOGGER.error(
                "Deck layout detection failed because PaddleOCR is unavailable.",
                exc_info=exc,
            )
            detail = _ocr_engine_unavailable_detail(lang=payload.lang, exc=exc)
            _set_ocr_progress(
                deck_id,
                status="failed",
                error=detail,
                step="layout",
            )
            if notify_context is not None:
                notify_failed("deck_processing", notify_context)
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail=detail,
            ) from exc
        except (ValueError, UnidentifiedImageError) as exc:
            _set_ocr_progress(
                deck_id,
                status="failed",
                error=str(exc),
                step="layout",
            )
            if notify_context is not None:
                notify_failed("deck_processing", notify_context)
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)
            ) from exc
        except Exception as exc:  # noqa: BLE001
            LOGGER.exception(
                "Deck layout analysis failed unexpectedly for deck %s", deck_id
            )
            detail = str(exc) or "Deck layout analysis failed unexpectedly."
            _set_ocr_progress(
                deck_id,
                status="failed",
                error=detail,
                step="layout",
            )
            if notify_context is not None:
                notify_failed("deck_processing", notify_context)
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=detail,
            ) from exc

        normalized_layout = _normalize_layout_payload(
            built_layout,
            deck_id=deck_id,
            lang=payload.lang,
        )
        try:
            storage.save_layout_payload(deck_id, normalized_layout)
        except Exception as exc:  # noqa: BLE001
            LOGGER.exception("Deck layout payload save failed for deck %s", deck_id)
            detail = str(exc) or "Unable to persist layout payload."
            _set_ocr_progress(
                deck_id,
                status="failed",
                error=detail,
                step="layout",
            )
            if notify_context is not None:
                notify_failed("deck_processing", notify_context)
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=detail,
            ) from exc
        did_work = True
    slide_count = len(
        normalized_layout.get("slides", [])
        if isinstance(normalized_layout, dict)
        and isinstance(normalized_layout.get("slides"), list)
        else []
    )

    def _ocr_progress_callback(built: int, total: int) -> None:
        _set_ocr_progress(
            deck_id,
            status="running",
            built_pages=built,
            total_pages=total,
            step="ocr",
            last_completed_step="layout",
        )

    normalized_ocr_payload: dict[str, object]
    if (
        cached_payload
        and not payload.force
        and cached_payload.get("lang") == payload.lang
        and style_cache_matches
        and _ocr_payload_uses_layout_guided_strategy(cached_payload)
    ):
        normalized_ocr_payload = normalize_ocr_payload(
            cached_payload, deck_id=deck_id, lang=payload.lang
        )
        if normalized_ocr_payload != cached_payload:
            storage.save_ocr_payload(deck_id, normalized_ocr_payload)
    else:
        _set_ocr_progress(
            deck_id,
            status="running",
            built_pages=0,
            total_pages=slide_count,
            lang=payload.lang,
            step="ocr",
            last_completed_step="layout",
        )
        try:
            ocr_payload = build_deck_ocr_payload(
                deck,
                deck_path,
                lang=payload.lang,
                include_bboxes=payload.includeBboxes,
                layout_payload=normalized_layout,
                pdf_path=pdf_path,
                progress_callback=_ocr_progress_callback,
            )
        except (ValueError, UnidentifiedImageError) as exc:
            _set_ocr_progress(
                deck_id,
                status="failed",
                error=str(exc),
                step="ocr",
                last_completed_step="layout",
            )
            if notify_context is not None:
                notify_failed("deck_processing", notify_context)
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)
            ) from exc
        except SlideOcrEngineUnavailableError as exc:
            LOGGER.error(
                "Deck OCR failed because PaddleOCR is unavailable.", exc_info=exc
            )
            detail = _ocr_engine_unavailable_detail(lang=payload.lang, exc=exc)
            _set_ocr_progress(
                deck_id,
                status="failed",
                error=detail,
                step="ocr",
                last_completed_step="layout",
            )
            if notify_context is not None:
                notify_failed("deck_processing", notify_context)
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail=detail,
            ) from exc
        except Exception as exc:  # noqa: BLE001 - ensure running state is finalized
            LOGGER.exception("Deck OCR failed unexpectedly for deck %s", deck_id)
            detail = str(exc) or "Deck OCR failed unexpectedly."
            _set_ocr_progress(
                deck_id,
                status="failed",
                error=detail,
                step="ocr",
                last_completed_step="layout",
            )
            if notify_context is not None:
                notify_failed("deck_processing", notify_context)
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=detail,
            ) from exc

        normalized_ocr_payload = normalize_ocr_payload(
            ocr_payload, deck_id=deck_id, lang=payload.lang
        )
        try:
            storage.save_ocr_payload(deck_id, normalized_ocr_payload)
        except (
            Exception
        ) as exc:  # noqa: BLE001 - persistence must flip status to failed
            LOGGER.exception("Deck OCR payload save failed for deck %s", deck_id)
            detail = str(exc) or "Unable to persist OCR payload."
            _set_ocr_progress(
                deck_id,
                status="failed",
                error=detail,
                step="ocr",
                last_completed_step="layout",
            )
            if notify_context is not None:
                notify_failed("deck_processing", notify_context)
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=detail,
            ) from exc
        did_work = True

    if normalized_layout is not None:
        merged_analysis = _build_slide_analysis_payload(
            normalized_layout,
            normalized_ocr_payload,
            deck_id=deck_id,
            lang=payload.lang,
        )
        if merged_analysis is not None:
            try:
                storage.save_slide_analysis_payload(deck_id, merged_analysis)
            except Exception as exc:  # noqa: BLE001
                LOGGER.exception(
                    "Merged slide analysis payload save failed for deck %s", deck_id
                )
                detail = str(exc) or "Unable to persist merged slide analysis payload."
                _set_ocr_progress(
                    deck_id,
                    status="failed",
                    error=detail,
                    step="ocr",
                    last_completed_step="layout",
                )
                if notify_context is not None:
                    notify_failed("deck_processing", notify_context)
                raise HTTPException(
                    status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                    detail=detail,
                ) from exc

    _set_ocr_progress(
        deck_id,
        status="completed",
        built_pages=slide_count,
        total_pages=slide_count,
        message="Deck processing complete.",
        step="complete",
        last_completed_step="ocr",
    )
    if notify_context is not None and did_work:
        notify_finished(
            time.perf_counter() - started_at,
            "deck_processing",
            notify_context,
        )
    return DeckOcrResponse(
        payload=DeckOcrPayload.model_validate(normalized_ocr_payload),
        cached=not did_work,
    )


@router.post(
    "/deck/{deck_id}/layout",
    response_model=DeckLayoutResponse,
    response_model_exclude_none=True,
)
def deck_layout_endpoint(
    deck_id: str,
    payload: DeckLayoutRequest,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckLayoutResponse:
    """Build or reuse a persisted Paddle layout payload for a deck."""

    try:
        deck = storage.load_deck(deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)

    cached_payload = _load_cached_slide_analysis_payload(
        deck_id, storage, lang=payload.lang
    )
    if cached_payload is not None and not payload.force:
        return DeckLayoutResponse(
            payload=DeckLayoutPayload.model_validate(cached_payload),
            cached=True,
        )

    deck_path = storage.root / deck_id
    has_images = any("<img" in (slide.body_html or "") for slide in deck.slides)
    if not has_images:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Layout analysis skipped: no slide images available.",
        )
    try:
        built_payload = build_deck_layout_payload(deck, deck_path, lang=payload.lang)
    except SlideOcrEngineUnavailableError as exc:
        LOGGER.error(
            "Deck layout detection failed because PaddleOCR is unavailable.",
            exc_info=exc,
        )
        detail = _ocr_engine_unavailable_detail(lang=payload.lang, exc=exc)
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=detail,
        ) from exc
    except (ValueError, UnidentifiedImageError) as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)
        ) from exc
    except Exception as exc:  # noqa: BLE001 - persist real failure cause
        LOGGER.exception(
            "Deck layout analysis failed unexpectedly for deck %s", deck_id
        )
        detail = str(exc) or "Deck layout analysis failed unexpectedly."
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=detail,
        ) from exc

    normalized = _normalize_layout_payload(
        built_payload,
        deck_id=deck_id,
        lang=payload.lang,
    )
    try:
        storage.save_layout_payload(deck_id, normalized)
        merged_analysis = _build_slide_analysis_payload(
            normalized,
            _load_cached_ocr_payload(deck_id, storage, lang=payload.lang),
            deck_id=deck_id,
            lang=payload.lang,
        )
        if merged_analysis is not None:
            storage.save_slide_analysis_payload(deck_id, merged_analysis)
    except Exception as exc:  # noqa: BLE001
        LOGGER.exception("Deck layout payload save failed for deck %s", deck_id)
        detail = str(exc) or "Unable to persist layout payload."
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=detail,
        ) from exc
    return DeckLayoutResponse(
        payload=DeckLayoutPayload.model_validate(
            merged_analysis if merged_analysis is not None else normalized
        ),
        cached=False,
    )


@router.post(
    "/deck/{deck_id}/layout/start",
    response_model=DeckLayoutStatusResponse,
    status_code=status.HTTP_202_ACCEPTED,
)
def start_deck_layout_endpoint(
    deck_id: str,
    payload: DeckLayoutRequest,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckLayoutStatusResponse:
    """Start persisted layout analysis in the background."""

    try:
        deck = storage.load_deck(deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)

    cached_payload = _load_cached_slide_analysis_payload(
        deck_id, storage, lang=payload.lang
    )
    if cached_payload is not None and not payload.force:
        total_pages = len(cached_payload.get("slides", []))
        _set_layout_progress(
            deck_id,
            status="completed",
            built_pages=total_pages,
            total_pages=total_pages,
            message="Layout payload already available.",
            step="complete",
            last_completed_step="layout",
        )
        return DeckLayoutStatusResponse(
            status="completed",
            builtPages=total_pages,
            totalPages=total_pages,
            message="Layout payload already available.",
            step="complete",
            updatedAt=_utc_now_iso(),
            lastCompletedStep="layout",
        )

    has_images = any("<img" in (slide.body_html or "") for slide in deck.slides)
    if not has_images:
        message = "Layout analysis skipped: no slide images available."
        _set_layout_progress(
            deck_id, status="skipped", message=message, step="complete"
        )
        return DeckLayoutStatusResponse(
            status="skipped",
            message=message,
            step="complete",
            updatedAt=_utc_now_iso(),
        )

    existing = _get_layout_progress(deck_id)
    if (
        existing is not None
        and str(existing.get("status") or "").strip().lower() == "running"
    ):
        return DeckLayoutStatusResponse(**existing)

    total_pages = sum(1 for slide in deck.slides if "<img" in (slide.body_html or ""))
    _set_layout_progress(
        deck_id,
        status="running",
        built_pages=0,
        total_pages=total_pages,
        message=_build_layout_processing_message(0, total_pages),
        step="layout",
    )
    _LAYOUT_EXECUTOR.submit(_run_deck_layout_job, deck_id, lang=payload.lang)
    return DeckLayoutStatusResponse(
        status="running",
        builtPages=0,
        totalPages=total_pages,
        message=_build_layout_processing_message(0, total_pages),
        step="layout",
        startedAt=_utc_now_iso(),
        updatedAt=_utc_now_iso(),
    )


@router.get("/deck/{deck_id}/layout/status", response_model=DeckLayoutStatusResponse)
def deck_layout_status(
    deck_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckLayoutStatusResponse:
    try:
        deck = storage.load_deck(deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)

    progress = _get_layout_progress(deck_id)
    cached_payload = _load_cached_slide_analysis_payload(deck_id, storage)
    if progress is not None:
        status_value = str(progress.get("status") or "").strip().lower()
        if status_value in {"idle", "running"} and cached_payload:
            total_pages = len(cached_payload.get("slides", []))
            completed_message = "Layout payload already available."
            _set_layout_progress(
                deck_id,
                status="completed",
                built_pages=total_pages,
                total_pages=total_pages,
                message=completed_message,
                step="complete",
                last_completed_step="layout",
            )
            return DeckLayoutStatusResponse(
                status="completed",
                builtPages=total_pages,
                totalPages=total_pages,
                message=completed_message,
                step="complete",
                updatedAt=_utc_now_iso(),
                lastCompletedStep="layout",
            )
        return DeckLayoutStatusResponse(**progress)

    if cached_payload:
        total_pages = len(cached_payload.get("slides", []))
        return DeckLayoutStatusResponse(
            status="completed",
            builtPages=total_pages,
            totalPages=total_pages,
            message="Layout payload already available.",
            step="complete",
            updatedAt=_utc_now_iso(),
            lastCompletedStep="layout",
        )

    has_images = any("<img" in (slide.body_html or "") for slide in deck.slides)
    if not has_images:
        return DeckLayoutStatusResponse(
            status="skipped",
            message="Layout analysis skipped: no slide images available.",
            step="complete",
            updatedAt=_utc_now_iso(),
        )
    return DeckLayoutStatusResponse(
        status="idle",
        message="Layout analysis not started.",
        updatedAt=_utc_now_iso(),
    )


@router.get("/deck/{deck_id}/ocr/status", response_model=DeckOcrStatusResponse)
def deck_ocr_status(
    deck_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckOcrStatusResponse:
    try:
        deck = storage.load_deck(deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)

    progress = _get_ocr_progress(deck_id)
    cached_payload = storage.load_ocr_payload(deck_id)
    cached_layout = _load_cached_slide_analysis_payload(deck_id, storage)
    resolved_lang = (
        str((progress or {}).get("lang") or "").strip()
        or str((cached_payload or {}).get("lang") or "").strip()
        or str((cached_layout or {}).get("lang") or "").strip()
        or None
    )
    if progress is not None:
        status_value = str(progress.get("status") or "").strip().lower()
        progress_step = str(progress.get("step") or "").strip().lower()
        if (
            status_value in {"idle", "running"}
            and cached_payload
            and progress_step in {"", "ocr"}
        ):
            total_pages = len(cached_payload.get("slides", []))
            completed_message = "OCR payload already available."
            _set_ocr_progress(
                deck_id,
                status="completed",
                built_pages=total_pages,
                total_pages=total_pages,
                message=completed_message,
                lang=resolved_lang,
                step="complete",
                last_completed_step="ocr",
            )
            return DeckOcrStatusResponse(
                status="completed",
                builtPages=total_pages,
                totalPages=total_pages,
                message=completed_message,
                lang=resolved_lang,
                step="complete",
                updatedAt=_utc_now_iso(),
                lastCompletedStep="ocr",
            )
        return DeckOcrStatusResponse(**{**progress, "lang": resolved_lang})

    if cached_payload:
        total_pages = len(cached_payload.get("slides", []))
        return DeckOcrStatusResponse(
            status="completed",
            builtPages=total_pages,
            totalPages=total_pages,
            message="OCR payload already available.",
            lang=resolved_lang,
            step="complete",
            updatedAt=_utc_now_iso(),
            lastCompletedStep="ocr",
        )

    if cached_layout:
        total_pages = len(cached_layout.get("slides", []))
        return DeckOcrStatusResponse(
            status="idle",
            builtPages=total_pages,
            totalPages=total_pages,
            message="Layout payload already available. OCR not started.",
            lang=resolved_lang,
            step="layout",
            updatedAt=_utc_now_iso(),
            lastCompletedStep="layout",
        )

    deck_path = storage.root / deck_id
    has_pdf = (deck_path / "source.pdf").exists()
    has_images = any("<img" in (slide.body_html or "") for slide in deck.slides)
    if not has_pdf and not has_images:
        return DeckOcrStatusResponse(
            status="skipped",
            message="OCR skipped: no PDF or slide images available.",
            lang=resolved_lang,
            step="complete",
            updatedAt=_utc_now_iso(),
        )
    return DeckOcrStatusResponse(
        status="idle",
        message="OCR not started.",
        lang=resolved_lang,
        updatedAt=_utc_now_iso(),
    )


@router.get("/deck/{deck_id}/ocr/audit", response_model=DeckOcrPayload)
def deck_ocr_audit_payload(
    deck_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckOcrPayload:
    try:
        deck = storage.load_deck(deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    payload = _load_cached_ocr_payload(deck_id, storage)
    if payload is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="OCR payload not available."
        )
    return DeckOcrPayload.model_validate(payload)


@router.get("/deck/{deck_id}/ocr/slides/{slide_id}", response_model=DeckOcrSlide)
def deck_ocr_slide_payload(
    deck_id: str,
    slide_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckOcrSlide:
    try:
        deck = storage.load_deck(deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    payload = _load_cached_ocr_payload(deck_id, storage)
    if payload is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="OCR payload not available."
        )
    slide_payload = find_slide_payload(payload, slide_id)
    if slide_payload is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="Slide OCR not available."
        )
    return DeckOcrSlide.model_validate(slide_payload)


@router.get("/deck/{deck_id}/ocr/analysis", response_model=DeckOcrAnalysisPayload)
def deck_ocr_analysis_payload(
    deck_id: str,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> DeckOcrAnalysisPayload:
    try:
        deck = storage.load_deck(deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(deck, user)
    payload = _load_cached_ocr_payload(deck_id, storage)
    if payload is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="OCR payload not available."
        )
    normalized_payload = normalize_ocr_payload(payload, deck_id=deck_id)
    llm_wrapper = LLMCallWrapper(mode="live")
    figure_classifications = classify_deck_chart_regions(
        llm_wrapper,
        deck=deck,
        deck_path=storage.root / deck_id,
        ocr_payload=normalized_payload,
    )
    analysis_payload = build_analysis_payload(
        normalized_payload,
        figure_classifications=figure_classifications,
    )
    return DeckOcrAnalysisPayload.model_validate(analysis_payload)


@router.post("/ocr", response_model=SlidesOcrResponse)
def slides_ocr(payload: SlidesOcrRequest) -> SlidesOcrResponse:
    """Run PaddleOCR for slide image selection and return the raw OCR payload."""
    try:
        raw_ocr = extract_raw_ocr_from_data_url(
            payload.imageDataUrl,
            payload.lang,
        )
    except (ValueError, UnidentifiedImageError) as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)
        ) from exc
    except SlideOcrEngineUnavailableError as exc:
        LOGGER.error("Slide OCR failed because PaddleOCR is unavailable.", exc_info=exc)
        detail = _ocr_engine_unavailable_detail(lang=payload.lang, exc=exc)
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=detail,
        ) from exc
    return SlidesOcrResponse(
        ocrText=extract_text_from_raw_ocr_result(raw_ocr),
        rawOcr=raw_ocr,
    )


@router.post("/layout", response_model=SlidesLayoutResponse)
def slides_layout(payload: SlidesLayoutRequest) -> SlidesLayoutResponse:
    """Run Paddle layout detection for slide image selection and return inspectable blocks."""
    try:
        raw_layout = extract_raw_layout_from_data_url(
            payload.imageDataUrl,
            payload.lang,
        )
        layout_summary = extract_layout_summary_from_raw_layout(
            raw_layout,
            slide_id=payload.slideId,
            slide_number=payload.slideNumber,
        )
    except (ValueError, UnidentifiedImageError) as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)
        ) from exc
    except SlideOcrEngineUnavailableError as exc:
        LOGGER.error(
            "Slide layout detection failed because PaddleOCR is unavailable.",
            exc_info=exc,
        )
        detail = _ocr_engine_unavailable_detail(lang=payload.lang, exc=exc)
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=detail,
        ) from exc

    raw_blocks = (
        layout_summary.get("blocks")
        if isinstance(layout_summary.get("blocks"), list)
        else []
    )
    raw_regions = (
        layout_summary.get("figure_regions")
        if isinstance(layout_summary.get("figure_regions"), list)
        else []
    )
    return SlidesLayoutResponse(
        rawLayout=raw_layout,
        blocks=[
            SlidesLayoutBlock(
                blockId=str(block.get("block_id") or block.get("blockId") or "").strip()
                or None,
                type=str(block.get("type") or "unknown"),
                text=str(block.get("text") or ""),
                items=[
                    str(item).strip()
                    for item in (
                        block.get("items")
                        if isinstance(block.get("items"), list)
                        else []
                    )
                    if str(item).strip()
                ],
                bbox=(
                    block.get("bbox") if isinstance(block.get("bbox"), dict) else None
                ),
                confidence=(
                    float(block.get("confidence"))
                    if isinstance(block.get("confidence"), (int, float))
                    else None
                ),
                tableModel=(
                    block.get("table_model")
                    if isinstance(block.get("table_model"), dict)
                    else (
                        block.get("tableModel")
                        if isinstance(block.get("tableModel"), dict)
                        else None
                    )
                ),
            )
            for block in raw_blocks
            if isinstance(block, dict)
        ],
        titleText=str(layout_summary.get("title_text") or ""),
        bulletTexts=[
            str(item).strip()
            for item in (
                layout_summary.get("bullet_texts")
                if isinstance(layout_summary.get("bullet_texts"), list)
                else []
            )
            if str(item).strip()
        ],
        figureRegions=[
            SlidesLayoutFigureRegion.model_validate(region)
            for region in raw_regions
            if isinstance(region, dict)
        ],
    )


@router.post("/deck/{deck_id}/import", response_model=ImportResponse)
def import_slide(
    deck_id: str,
    payload: ImportRequest,
    user: AuthenticatedUser | None = Depends(require_authenticated_user),
    storage: DeckStorage = Depends(get_storage),
) -> ImportResponse:
    normalized_deck_id = _normalize_deck_id(deck_id)
    source_deck_id = _normalize_deck_id(payload.sourceDeckId)
    try:
        target_deck = storage.load_deck(normalized_deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(target_deck, user)
    try:
        source_deck = storage.load_deck(source_deck_id)
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    _assert_deck_access(source_deck, user)
    if source_deck.prompt_style != target_deck.prompt_style:
        detail = (
            "Cannot import slides between decks with different promptStyle values: "
            f"{target_deck.prompt_style}, {source_deck.prompt_style}."
        )
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=detail)

    existing_ids = set(storage.list_slide_ids(normalized_deck_id)) | set(
        payload.currentOrder
    )
    new_slide_id = generate_slide_filename(existing_ids)
    try:
        slide = storage.import_slide(
            source_deck_id=source_deck_id,
            source_slide_id=payload.sourceSlideId,
            target_deck_id=normalized_deck_id,
            new_slide_id=new_slide_id,
        )
    except DeckNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc
    except SlideNotFoundError as exc:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail=str(exc)
        ) from exc

    order = list(payload.currentOrder)
    insert_after = payload.afterSlideId
    if insert_after:
        if insert_after not in order:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Slide {insert_after} not present in current order.",
            )
        index = order.index(insert_after) + 1
    else:
        index = 0
    order.insert(index, slide.id)
    _copy_imported_deck_processing_artifacts(
        source_deck_id=source_deck_id,
        source_slide_id=payload.sourceSlideId,
        target_deck_id=normalized_deck_id,
        target_slide_id=slide.id,
        storage=storage,
    )
    _reconcile_deck_processing_artifacts(
        normalized_deck_id,
        storage,
        slide_ids=order,
    )
    return ImportResponse(slide=SlidePayload.from_domain(slide), order=order)


__all__ = [
    "router",
    "site_router",
    "get_storage",
    "set_storage",
    "is_any_ocr_running",
]


_DECK_ID_PATTERN = re.compile(r"^[A-Za-z0-9._-]+$")


def _normalize_deck_id(deck_id: str) -> str:
    normalized = deck_id.strip()
    if not normalized:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail="Deck ID is required."
        )
    if not _DECK_ID_PATTERN.fullmatch(normalized):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Deck ID may only contain letters, numbers, dots, underscores, and hyphens.",
        )
    return normalized


def _normalize_asset_path(asset_path: str) -> str:
    normalized = asset_path.lstrip("/")
    for _ in range(2):
        if normalized.startswith("slides/"):
            normalized = normalized[len("slides/") :]
            continue
        if normalized.startswith("assets/"):
            normalized = normalized[len("assets/") :]
            continue
        break
    return normalized


def _resolve_ocr_pdf_path(pdf_path: str | None, deck_path: Path) -> Path | None:
    if not pdf_path:
        candidate = deck_path / "source.pdf"
        return candidate if candidate.exists() else None
    candidate = Path(pdf_path)
    if not candidate.is_absolute():
        candidate = deck_path / candidate
    candidate = candidate.resolve()
    deck_root = deck_path.resolve()
    if candidate == deck_root or deck_root not in candidate.parents:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="PDF path must be within the deck directory.",
        )
    if not candidate.exists():
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="PDF not found.",
        )
    if not candidate.is_file():
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="PDF path must point to a file.",
        )
    return candidate


def _enqueue_deck_ocr_job(deck: Deck, *, lang: str) -> None:
    total_pages = sum(1 for slide in deck.slides if "<img" in (slide.body_html or ""))
    _set_ocr_progress(
        deck.deck_id,
        status="running",
        built_pages=0,
        total_pages=total_pages,
        message=_build_layout_processing_message(0, total_pages),
        lang=lang,
        step="layout",
    )
    _OCR_EXECUTOR.submit(_run_deck_ocr_job, deck.deck_id, lang=lang)


def _is_pdf_upload(file: UploadFile) -> bool:
    content_type = (file.content_type or "").lower()
    if content_type == "application/pdf":
        return True
    if file.filename and file.filename.lower().endswith(".pdf"):
        return True
    return False


def _is_supported_image_upload(file: UploadFile) -> bool:
    content_type = (file.content_type or "").lower()
    if content_type in {"image/png", "image/jpeg", "image/jpg"}:
        return True
    if file.filename and file.filename.lower().endswith((".png", ".jpg", ".jpeg")):
        return True
    return False


def _promote_single_nested_directory(deck_path: Path) -> None:
    """Move files up when uploads include a redundant top-level folder."""

    if find_index_file(deck_path) is not None:
        return

    candidates = [
        entry
        for entry in deck_path.iterdir()
        if entry.is_dir() and find_index_file(entry) is not None
    ]
    if len(candidates) != 1:
        return

    nested_root = candidates[0]
    for child in nested_root.iterdir():
        target = deck_path / child.name
        shutil.move(str(child), target)
    nested_root.rmdir()


def _cover_notebooklm_logo(image: Image.Image) -> None:
    """Fill the cover area with a median, filtering near-black unless all pixels are near-black."""
    if not image.size:
        return
    width, height = image.size
    if width <= 0 or height <= 0:
        return
    cover_width = max(200, int(math.floor(width * 0.10)))
    cover_height = max(30, int(math.floor(height * 0.03)))
    x0 = max(0, width - cover_width)
    y0 = max(0, height - cover_height)
    x1 = width
    y1 = height
    if x1 <= x0 or y1 <= y0:
        return
    sampling_image = image if image.mode == "RGB" else image.convert("RGB")

    samples: list[tuple[int, int, int]] = []
    filtered_samples: list[tuple[int, int, int]] = []
    for py in range(y0, y1):
        for px in range(x0, x1):
            pixel = sampling_image.getpixel((px, py))
            samples.append(pixel)
            if max(pixel) > 20:
                filtered_samples.append(pixel)

    fill_color = (255, 255, 255)
    sample_set = filtered_samples if filtered_samples else samples
    if sample_set:
        channel_values = [
            sorted(pixel[idx] for pixel in sample_set) for idx in range(3)
        ]
        mid = len(sample_set) // 2
        fill_color = tuple(values[mid] for values in channel_values)
    draw = ImageDraw.Draw(image)
    draw.rectangle([x0, y0, x1, y1], fill=fill_color)


def _cover_notebooklm_logo_in_image_bytes(image_bytes: bytes) -> bytes:
    """Return PNG bytes with the NotebookLM logo area covered."""

    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            has_alpha = image.mode in {"RGBA", "LA"} or (
                image.mode == "P" and "transparency" in image.info
            )
            if has_alpha:
                rgba = image.convert("RGBA")
                rendered = Image.new("RGB", rgba.size, (255, 255, 255))
                rendered.paste(rgba, mask=rgba.split()[-1])
            else:
                rendered = image.convert("RGB")
            _cover_notebooklm_logo(rendered)
            buffer = io.BytesIO()
            rendered.save(buffer, format="PNG")
            return buffer.getvalue()
    except (UnidentifiedImageError, OSError, ValueError):
        return image_bytes


def _render_pdf_deck(
    deck_id: str,
    deck_path: Path,
    pdf_bytes: bytes,
    storage: DeckStorage,
    *,
    prompt_style: str,
    owner_email: str | None,
    shared_with: list[str],
) -> None:
    assets_path = deck_path / "assets"
    assets_path.mkdir(parents=True, exist_ok=True)
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    except (fitz.FileDataError, RuntimeError, ValueError) as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Uploaded file is not a valid PDF.",
        ) from exc
    pdf_path = deck_path / "source.pdf"
    pdf_path.write_bytes(pdf_bytes)
    with doc:
        if doc.page_count == 0:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Uploaded PDF has no pages.",
            )
        selected_dimension: tuple[float, float] | None = None
        for page_number in range(doc.page_count):
            page = doc.load_page(page_number)
            crop_box = page.cropbox
            width_pt = float(crop_box.width)
            height_pt = float(crop_box.height)
            matched_dimension: tuple[float, float] | None = next(
                (
                    (allowed_w, allowed_h)
                    for allowed_w, allowed_h in _SLIDES_PDF_ALLOWED_DIMENSIONS_PT
                    if (
                        abs(width_pt - allowed_w)
                        <= _SLIDES_PDF_REQUIRED_DIMENSION_TOLERANCE_PT
                        and abs(height_pt - allowed_h)
                        <= _SLIDES_PDF_REQUIRED_DIMENSION_TOLERANCE_PT
                    )
                ),
                None,
            )
            if matched_dimension is None:
                allowed_sizes_text = " or ".join(
                    f"{int(width)}x{int(height)}"
                    for width, height in _SLIDES_PDF_ALLOWED_DIMENSIONS_PT
                )
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=(
                        "Uploaded PDF must use NotebookLM slide pages "
                        f"({allowed_sizes_text}). "
                        f"Page {page_number + 1} is "
                        f"{width_pt:.1f}x{height_pt:.1f}."
                    ),
                )
            if selected_dimension is None:
                selected_dimension = matched_dimension
            elif matched_dimension != selected_dimension:
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=(
                        "Uploaded PDF pages must all have the same dimensions. "
                        f"Page 1 matched {int(selected_dimension[0])}x"
                        f"{int(selected_dimension[1])}, but page {page_number + 1} is "
                        f"{width_pt:.1f}x{height_pt:.1f}."
                    ),
                )
        slides: list[Slide] = []
        raster_matrix = fitz.Matrix(
            _SLIDES_PDF_IMPORT_RASTER_SCALE, _SLIDES_PDF_IMPORT_RASTER_SCALE
        )
        for page_number in range(doc.page_count):
            page = doc.load_page(page_number)
            pix = page.get_pixmap(matrix=raster_matrix, alpha=True)
            crop_box = page.cropbox
            image = Image.frombytes("RGBA", (pix.width, pix.height), pix.samples)
            flattened = Image.new("RGB", image.size, (255, 255, 255))
            flattened.paste(image, mask=image.split()[-1])
            _cover_notebooklm_logo(flattened)
            buffer = io.BytesIO()
            flattened.save(buffer, format="PNG")
            image_bytes = buffer.getvalue()
            image_hash = hashlib.sha256(image_bytes).hexdigest()
            image_name = f"{image_hash}.png"
            image_path = assets_path / image_name
            image_path.write_bytes(image_bytes)
            slide_id = f"slide-{page_number + 1:03d}.html"
            title_html, body_html, full_html = build_image_only_slide_content(
                deck_id,
                Path(image_name),
                page_index=page_number,
                crop_w_pt=crop_box.width,
                crop_h_pt=crop_box.height,
                crop_x0_pt=crop_box.x0,
                crop_y0_pt=crop_box.y0,
                rotation_deg=page.rotation,
            )
            slides.append(
                Slide(
                    id=slide_id,
                    title_html=title_html,
                    body_html=body_html,
                    full_html=full_html,
                )
            )
    deck = Deck(
        deck_id=deck_id,
        prompt_style=prompt_style,
        owner_email=owner_email,
        shared_with=shared_with,
        slides=slides,
    )
    storage.save_deck(deck)


def _resolve_image_extension(
    image_format: Optional[str], filename: Optional[str]
) -> str:
    """Return a file extension for a supported image format."""

    format_map = {
        "PNG": ".png",
        "JPEG": ".jpg",
        "JPG": ".jpg",
    }
    if image_format and image_format.upper() in format_map:
        return format_map[image_format.upper()]
    if filename:
        suffix = Path(filename).suffix.lower()
        if suffix in {".png", ".jpg", ".jpeg"}:
            return ".jpg" if suffix == ".jpeg" else suffix
    return ".png"


def _render_image_deck(
    deck_id: str,
    deck_path: Path,
    image_bytes: bytes,
    storage: DeckStorage,
    filename: Optional[str],
    *,
    prompt_style: str,
    owner_email: str | None,
    shared_with: list[str],
) -> None:
    assets_path = deck_path / "assets"
    assets_path.mkdir(parents=True, exist_ok=True)
    try:
        with Image.open(io.BytesIO(image_bytes)) as image:
            image_format = (image.format or "").upper()
            width, height = image.size
            has_alpha = image.mode in {"RGBA", "LA"} or (
                image.mode == "P" and "transparency" in image.info
            )
            if has_alpha:
                rgba = image.convert("RGBA")
                rendered = Image.new("RGB", rgba.size, (255, 255, 255))
                rendered.paste(rgba, mask=rgba.split()[-1])
            else:
                rendered = image.convert("RGB")
    except (UnidentifiedImageError, OSError, ValueError) as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Uploaded file is not a valid image.",
        ) from exc
    _cover_notebooklm_logo(rendered)
    buffer = io.BytesIO()
    if image_format in {"JPEG", "JPG"}:
        rendered.save(buffer, format="JPEG", quality=95, optimize=True, subsampling=0)
        image_extension = ".jpg"
    else:
        rendered.save(buffer, format="PNG")
        image_extension = ".png"
    patched_bytes = buffer.getvalue()
    image_hash = hashlib.sha256(patched_bytes).hexdigest()
    image_name = f"{image_hash}{image_extension}"
    image_path = assets_path / image_name
    image_path.write_bytes(patched_bytes)
    slide_id = "slide-001.html"
    title_html, body_html, full_html = build_image_only_slide_content(
        deck_id,
        Path(image_name),
        page_index=0,
        crop_w_pt=width,
        crop_h_pt=height,
        crop_x0_pt=0,
        crop_y0_pt=0,
        rotation_deg=0,
    )
    deck = Deck(
        deck_id=deck_id,
        prompt_style=prompt_style,
        slides=[
            Slide(
                id=slide_id,
                title_html=title_html,
                body_html=body_html,
                full_html=full_html,
            )
        ],
    )
    storage.save_deck(deck)


def _format_upload_failure_detail(exc: Exception) -> str:
    message = str(exc).strip()
    if not message:
        message = exc.__class__.__name__
    return f"Failed to upload deck: {message}"


class ConcatenatedSlideMapping(NamedTuple):
    source_deck_id: str
    source_slide_id: str
    target_slide_id: str


class _ConcatenateDeckContext(NamedTuple):
    deck: Deck
    section_id_map: dict[str, str]
    subsection_id_map: dict[str, str]
    slide_id_map: dict[str, str]
    ocr_tokens_by_index: list[frozenset[str]]


_OCR_INTERLEAVE_WINDOW = 2
_OCR_INTERLEAVE_MIN_TOKENS = 3
_OCR_INTERLEAVE_SIMILARITY_THRESHOLD = 0.4


def _normalize_ocr_tokens(text: str) -> frozenset[str]:
    """Return normalized OCR tokens for similarity comparisons."""

    tokens = {
        token for token in re.findall(r"[a-z0-9]+", text.lower()) if len(token) > 1
    }
    return frozenset(tokens)


def _build_ocr_tokens_by_slide_id(
    deck_id: str, storage: DeckStorage
) -> dict[str, frozenset[str]]:
    """Build OCR token fingerprints keyed by slide id for ``deck_id``."""

    payload = storage.load_ocr_payload(deck_id)
    if payload is None:
        return {}
    normalized = normalize_ocr_payload(payload, deck_id=deck_id)
    tokens_by_slide_id: dict[str, frozenset[str]] = {}
    for slide_payload in normalized.get("slides", []):
        if not isinstance(slide_payload, dict):
            continue
        slide_id = str(slide_payload.get("slide_id") or "")
        if not slide_id:
            continue
        ocr_text = str(slide_payload.get("ocr_text") or "")
        tokens = _normalize_ocr_tokens(ocr_text)
        if tokens:
            tokens_by_slide_id[slide_id] = tokens
    return tokens_by_slide_id


def _tokens_are_usable(tokens: frozenset[str]) -> bool:
    return len(tokens) >= _OCR_INTERLEAVE_MIN_TOKENS


def _jaccard_similarity(tokens_a: frozenset[str], tokens_b: frozenset[str]) -> float:
    if not tokens_a or not tokens_b:
        return 0.0
    intersection = len(tokens_a & tokens_b)
    if intersection == 0:
        return 0.0
    union = len(tokens_a | tokens_b)
    return intersection / union


def _first_unused_index(
    slide_count: int, used_indices: set[int], *, start_at: int
) -> int | None:
    """Return the first unused slide index starting at ``start_at``."""

    for candidate_index in range(start_at, slide_count):
        if candidate_index not in used_indices:
            return candidate_index
    for candidate_index in range(0, min(start_at, slide_count)):
        if candidate_index not in used_indices:
            return candidate_index
    return None


def _find_windowed_ocr_match(
    tokens_by_index: list[frozenset[str]],
    *,
    index: int,
    guide_tokens: frozenset[str],
    used_indices: set[int],
) -> int | None:
    """Return the best OCR match within a small window around ``index``."""

    if not _tokens_are_usable(guide_tokens):
        return None
    slide_count = len(tokens_by_index)
    if slide_count == 0:
        return None
    start_index = max(0, index - _OCR_INTERLEAVE_WINDOW)
    end_index = min(slide_count - 1, index + _OCR_INTERLEAVE_WINDOW)
    best_index: int | None = None
    best_score = 0.0
    best_distance = math.inf
    epsilon = 1e-9
    for candidate_index in range(start_index, end_index + 1):
        if candidate_index in used_indices:
            continue
        candidate_tokens = tokens_by_index[candidate_index]
        if not _tokens_are_usable(candidate_tokens):
            continue
        score = _jaccard_similarity(guide_tokens, candidate_tokens)
        if score < _OCR_INTERLEAVE_SIMILARITY_THRESHOLD:
            continue
        distance = abs(candidate_index - index)
        if score > best_score + epsilon:
            best_index = candidate_index
            best_score = score
            best_distance = distance
            continue
        if abs(score - best_score) <= epsilon and distance < best_distance:
            best_index = candidate_index
            best_distance = distance
            continue
        if (
            abs(score - best_score) <= epsilon
            and distance == best_distance
            and best_index is not None
            and candidate_index < best_index
        ):
            best_index = candidate_index
    return best_index


def _select_interleave_guide(
    contexts: list[_ConcatenateDeckContext],
    interleave_guide: Literal["first", "longest", "shortest"],
) -> _ConcatenateDeckContext:
    """Select the guide deck context used to anchor interleaving."""

    if interleave_guide == "longest":
        return max(contexts, key=lambda context: len(context.deck.slides))
    if interleave_guide == "shortest":
        return min(contexts, key=lambda context: len(context.deck.slides))
    return contexts[0]


def _concatenate_deck_payload(
    new_deck_id: str,
    source_ids: List[str],
    storage: DeckStorage,
    *,
    prompt_style: str,
    interleave_by_index: bool = False,
    interleave_guide: Literal["first", "longest", "shortest"] = "shortest",
) -> tuple[Deck, list[ConcatenatedSlideMapping]]:
    combined_slides: list[Slide] = []
    combined_sections: list[Section] = []
    existing_ids: set[str] = set()
    seen_signatures: set[str] = set()
    slide_mappings: list[ConcatenatedSlideMapping] = []
    dedupe_enabled = not interleave_by_index
    contexts: list[_ConcatenateDeckContext] = []
    for source_id in source_ids:
        deck = storage.load_deck(source_id)
        tokens_by_slide_id = _build_ocr_tokens_by_slide_id(source_id, storage)
        section_id_map: dict[str, str] = {}
        subsection_id_map: dict[str, str] = {}
        for section in deck.sections:
            new_section_id = f"{deck.deck_id}__{section.id}"
            section_id_map[section.id] = new_section_id
            for subsection in section.subsections:
                new_subsection_id = f"{deck.deck_id}__{subsection.id}"
                subsection_id_map[subsection.id] = new_subsection_id
        slide_id_map: dict[str, str] = {}
        contexts.append(
            _ConcatenateDeckContext(
                deck=deck,
                section_id_map=section_id_map,
                subsection_id_map=subsection_id_map,
                slide_id_map=slide_id_map,
                ocr_tokens_by_index=[
                    tokens_by_slide_id.get(slide.id, frozenset())
                    for slide in deck.slides
                ],
            )
        )

    guide_length = 0
    guide_context: _ConcatenateDeckContext | None = None
    if contexts:
        deck_lengths = [len(context.deck.slides) for context in contexts]
        if interleave_by_index:
            guide_context = _select_interleave_guide(contexts, interleave_guide)
            guide_length = len(guide_context.deck.slides)
        else:
            guide_length = deck_lengths[0]

    used_indices_by_deck: dict[str, set[int]] = {
        context.deck.deck_id: set() for context in contexts
    }

    def iter_slides():
        if not interleave_by_index:
            for context in contexts:
                for slide in context.deck.slides:
                    yield context, slide
            return
        if guide_context is None:
            return
        guide_has_usable_tokens = any(
            _tokens_are_usable(tokens) for tokens in guide_context.ocr_tokens_by_index
        )
        for index in range(guide_length):
            guide_tokens = guide_context.ocr_tokens_by_index[index]
            use_ocr_similarity = guide_has_usable_tokens and _tokens_are_usable(
                guide_tokens
            )
            for context in contexts:
                slide_count = len(context.deck.slides)
                if slide_count == 0:
                    continue
                used_indices = used_indices_by_deck[context.deck.deck_id]
                selected_index: int | None = None
                if use_ocr_similarity and context is not guide_context:
                    selected_index = _find_windowed_ocr_match(
                        context.ocr_tokens_by_index,
                        index=index,
                        guide_tokens=guide_tokens,
                        used_indices=used_indices,
                    )
                if (
                    selected_index is None
                    and index < slide_count
                    and index not in used_indices
                ):
                    selected_index = index
                if selected_index is None:
                    selected_index = _first_unused_index(
                        slide_count, used_indices, start_at=index
                    )
                if selected_index is None:
                    continue
                used_indices.add(selected_index)
                yield context, context.deck.slides[selected_index]
        for context in contexts:
            used_indices = used_indices_by_deck[context.deck.deck_id]
            for slide_index, slide in enumerate(context.deck.slides):
                if slide_index in used_indices:
                    continue
                yield context, slide

    for context, slide in iter_slides():
        source_id = context.deck.deck_id
        signature = _slide_signature(slide)
        if (
            dedupe_enabled
            and slide.kind == "normal"
            and slide.section_id is None
            and slide.subsection_id is None
            and signature in seen_signatures
        ):
            continue
        new_id = generate_slide_filename(existing_ids)
        existing_ids.add(new_id)
        context.slide_id_map[slide.id] = new_id
        if (
            dedupe_enabled
            and slide.kind == "normal"
            and slide.section_id is None
            and slide.subsection_id is None
        ):
            seen_signatures.add(signature)
        rewritten = _rewrite_slide_assets_for_concatenation(
            slide,
            source_deck_id=source_id,
            target_deck_id=new_deck_id,
            asset_subdir=_concatenated_asset_subdir(source_id),
        )
        slide_mappings.append(
            ConcatenatedSlideMapping(
                source_deck_id=source_id,
                source_slide_id=slide.id,
                target_slide_id=new_id,
            )
        )
        combined_slides.append(
            replace(
                rewritten,
                id=new_id,
                section_id=context.section_id_map.get(slide.section_id),
                subsection_id=context.subsection_id_map.get(slide.subsection_id),
            )
        )

    for context in contexts:
        deck = context.deck
        for section in deck.sections:
            combined_sections.append(
                Section(
                    id=context.section_id_map.get(section.id, section.id),
                    title=section.title,
                    start_slide=context.slide_id_map.get(
                        section.start_slide, section.start_slide
                    ),
                    subsections=[
                        Subsection(
                            id=context.subsection_id_map.get(
                                subsection.id, subsection.id
                            ),
                            title=subsection.title,
                            start_slide=context.slide_id_map.get(
                                subsection.start_slide, subsection.start_slide
                            ),
                        )
                        for subsection in section.subsections
                    ],
                )
            )
    return (
        Deck(
            deck_id=new_deck_id,
            prompt_style=prompt_style,
            slides=combined_slides,
            sections=combined_sections,
        ),
        slide_mappings,
    )


def _merge_concatenated_ocr_payload(
    new_deck_id: str,
    slide_mappings: list[ConcatenatedSlideMapping],
    storage: DeckStorage,
) -> None:
    if not slide_mappings:
        return
    source_slide_payloads: dict[str, dict[str, dict[str, object]]] = {}
    resolved_lang: str | None = None
    for source_deck_id in {mapping.source_deck_id for mapping in slide_mappings}:
        payload = storage.load_ocr_payload(source_deck_id)
        if payload is None:
            continue
        normalized = normalize_ocr_payload(payload, deck_id=source_deck_id)
        if resolved_lang is None:
            resolved_lang = str(normalized.get("lang") or "eng")
        slide_map: dict[str, dict[str, object]] = {}
        for slide_payload in normalized.get("slides", []):
            if not isinstance(slide_payload, dict):
                continue
            slide_id = str(slide_payload.get("slide_id") or "")
            if slide_id:
                slide_map[slide_id] = slide_payload
        source_slide_payloads[source_deck_id] = slide_map
    if not source_slide_payloads:
        return
    merged_slides: list[dict[str, object]] = []
    resolved_ocr_strategy: str | None = None
    resolved_prompt_style: str | None = None
    resolved_style_hint: dict[str, object] | None = None
    for index, mapping in enumerate(slide_mappings):
        slide_payloads = source_slide_payloads.get(mapping.source_deck_id)
        if not slide_payloads:
            continue
        source_slide = slide_payloads.get(mapping.source_slide_id)
        if not source_slide:
            continue
        if resolved_ocr_strategy is None:
            source_payload = storage.load_ocr_payload(mapping.source_deck_id)
            if isinstance(source_payload, dict):
                normalized_source_payload = normalize_ocr_payload(
                    source_payload, deck_id=mapping.source_deck_id
                )
                resolved_ocr_strategy = cast(
                    str | None, normalized_source_payload.get("ocr_strategy")
                )
                resolved_prompt_style = cast(
                    str | None, normalized_source_payload.get("prompt_style")
                )
                if isinstance(normalized_source_payload.get("style_hint"), dict):
                    resolved_style_hint = deepcopy(
                        cast(dict[str, object], normalized_source_payload["style_hint"])
                    )
        merged_slides.append(
            _copy_ocr_slide_for_target(
                source_slide,
                target_slide_id=mapping.target_slide_id,
                slide_number=index + 1,
            )
        )
    if not merged_slides:
        return
    merged_payload = normalize_ocr_payload(
        {
            "deck_id": new_deck_id,
            "lang": resolved_lang or "eng",
            "ocr_strategy": resolved_ocr_strategy,
            "prompt_style": resolved_prompt_style,
            "style_hint": resolved_style_hint,
            "generated_at": datetime.now(UTC).isoformat(),
            "slides": merged_slides,
        },
        deck_id=new_deck_id,
        lang=resolved_lang or "eng",
    )
    storage.save_ocr_payload(new_deck_id, merged_payload)


def _copy_shared_concatenated_pptx_template(
    new_deck_id: str,
    source_ids: list[str],
    storage: DeckStorage,
) -> None:
    candidate_paths: list[tuple[Path, Path]] = []
    digests: set[str] = set()
    for source_id in source_ids:
        source_deck_path = storage.root / source_id
        template_path = source_deck_path / DECK_PPTX_TEMPLATE_FILENAME
        manifest_path = source_deck_path / DECK_PPTX_TEMPLATE_MANIFEST_FILENAME
        if not template_path.exists() and not manifest_path.exists():
            continue
        if not template_path.exists() or not manifest_path.exists():
            LOGGER.warning(
                "Skipping PPTX template propagation for combined deck %s: deck %s has an incomplete template bundle.",
                new_deck_id,
                source_id,
            )
            return
        digests.add(hashlib.sha256(template_path.read_bytes()).hexdigest())
        candidate_paths.append((template_path, manifest_path))
    if not candidate_paths:
        return
    if len(candidate_paths) != len(source_ids):
        LOGGER.info(
            "Combined deck %s will not inherit a PPTX template because not all source decks have one.",
            new_deck_id,
        )
        return
    if len(digests) != 1:
        LOGGER.info(
            "Combined deck %s will not inherit a PPTX template because source decks use different templates.",
            new_deck_id,
        )
        return
    target_deck_path = storage.root / new_deck_id
    target_deck_path.mkdir(parents=True, exist_ok=True)
    template_path, manifest_path = candidate_paths[0]
    (target_deck_path / DECK_PPTX_TEMPLATE_FILENAME).write_bytes(
        template_path.read_bytes()
    )
    (target_deck_path / DECK_PPTX_TEMPLATE_MANIFEST_FILENAME).write_text(
        manifest_path.read_text(encoding="utf-8"),
        encoding="utf-8",
    )


def _copy_concatenated_assets(
    new_deck_id: str, source_ids: List[str], storage: DeckStorage
) -> None:
    target_assets_root = storage.root / new_deck_id / "assets"
    target_assets_root.mkdir(parents=True, exist_ok=True)
    for source_id in source_ids:
        source_assets = storage.root / source_id / "assets"
        if not source_assets.exists():
            continue
        target_assets = target_assets_root / _concatenated_asset_subdir(source_id)
        shutil.copytree(source_assets, target_assets, dirs_exist_ok=True)


def _rewrite_slide_assets_for_concatenation(
    slide: Slide,
    *,
    source_deck_id: str,
    target_deck_id: str,
    asset_subdir: str | None,
) -> Slide:
    def _rewrite(value: str) -> str:
        return _rewrite_asset_sources_for_concatenation(
            value,
            source_deck_id=source_deck_id,
            target_deck_id=target_deck_id,
            asset_subdir=asset_subdir,
        )

    return replace(
        slide,
        title_html=_rewrite(slide.title_html),
        body_html=_rewrite(slide.body_html),
        notes_html=_rewrite(slide.notes_html),
        source_html=_rewrite(slide.source_html),
        full_html=_rewrite(slide.full_html),
    )


def _rewrite_asset_sources_for_concatenation(
    html: str,
    *,
    source_deck_id: str,
    target_deck_id: str,
    asset_subdir: str | None,
) -> str:
    if not html:
        return html
    source_asset_prefix = f"/slides/deck/{source_deck_id}/assets/"
    target_asset_prefix = f"/slides/deck/{target_deck_id}/assets/"
    if asset_subdir:
        target_asset_prefix = f"{target_asset_prefix}{asset_subdir}/"

    def _build_target_url(raw_url: str) -> str | None:
        parsed = urlparse(raw_url)
        relative_path = _extract_asset_relative_path(raw_url, source_asset_prefix)
        if not relative_path:
            relative_path = _extract_relative_asset_path(parsed)
        if not relative_path:
            return None
        suffix = ""
        if parsed.query:
            suffix = f"?{parsed.query}"
        if parsed.fragment:
            suffix = f"{suffix}#{parsed.fragment}"
        return f"{target_asset_prefix}{relative_path}{suffix}"

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup.find_all(True):
        for attr in ("src", "data-src", "href", "poster"):
            value = tag.get(attr)
            if not value:
                continue
            replacement = _build_target_url(value)
            if replacement:
                tag[attr] = replacement
        style_value = tag.get("style")
        if style_value:
            tag["style"] = _rewrite_css_urls(style_value, _build_target_url)
    for style_tag in soup.find_all("style"):
        if not style_tag.string:
            continue
        style_tag.string = _rewrite_css_urls(style_tag.string, _build_target_url)
    return soup.decode()


def _extract_relative_asset_path(parsed_url: ParseResult) -> str | None:
    if parsed_url.scheme or parsed_url.netloc:
        return None
    path = parsed_url.path
    if path.startswith("./"):
        path = path[2:]
    if path.startswith("/"):
        path = path[1:]
    if not path.startswith("assets/"):
        return None
    return path[len("assets/") :]


def _rewrite_css_urls(style_value: str, replacer: Callable[[str], str | None]) -> str:
    def _replace(match: re.Match[str]) -> str:
        raw_url = match.group(1).strip().strip("'\"")
        replacement = replacer(raw_url)
        if not replacement:
            return match.group(0)
        return f'url("{replacement}")'

    return re.sub(r"url\(([^)]+)\)", _replace, style_value)


def _slide_signature(slide: Slide) -> str:
    content = slide.full_html or ""
    return hashlib.sha256(content.encode("utf-8")).hexdigest()


def _extract_slide_page_size_from_html(
    html: str | None,
) -> tuple[float, float] | None:
    if not html:
        return None
    soup = BeautifulSoup(html, "html.parser")
    image = soup.find("img")
    if image is None:
        return None
    width_raw = image.get("data-pdf-crop-w-pt")
    height_raw = image.get("data-pdf-crop-h-pt")
    try:
        width = float(width_raw) if width_raw is not None else 0.0
        height = float(height_raw) if height_raw is not None else 0.0
    except (TypeError, ValueError):
        return None
    if width <= 0 or height <= 0:
        return None
    return width, height


def _resolve_deck_export_page_size(
    deck: Deck, *, deck_path: Path | None = None
) -> tuple[float, float]:
    fallback = (
        float(_SLIDES_EXPORT_PAGE_WIDTH_PT),
        float(_SLIDES_EXPORT_PAGE_HEIGHT_PT),
    )
    resolved: tuple[float, float] | None = None
    for slide in deck.slides:
        size = _extract_slide_page_size_from_html(
            slide.body_html
        ) or _extract_slide_page_size_from_html(slide.full_html)
        if size is None:
            continue
        if resolved is None:
            resolved = size
            continue
        if (
            abs(size[0] - resolved[0]) > _SLIDES_PDF_REQUIRED_DIMENSION_TOLERANCE_PT
            or abs(size[1] - resolved[1]) > _SLIDES_PDF_REQUIRED_DIMENSION_TOLERANCE_PT
        ):
            LOGGER.warning(
                "Deck %s has mixed page sizes in slide metadata; using first detected %.1fx%.1f for export.",
                deck.deck_id,
                resolved[0],
                resolved[1],
            )
            break
    if resolved is not None:
        return resolved
    if deck_path is not None:
        source_pdf_path = deck_path / "source.pdf"
        if source_pdf_path.exists():
            try:
                with fitz.open(source_pdf_path) as source_doc:
                    if source_doc.page_count > 0:
                        first_crop = source_doc.load_page(0).cropbox
                        width = float(first_crop.width)
                        height = float(first_crop.height)
                        if width > 0 and height > 0:
                            for page_index in range(1, source_doc.page_count):
                                crop_box = source_doc.load_page(page_index).cropbox
                                if (
                                    abs(float(crop_box.width) - width)
                                    > _SLIDES_PDF_REQUIRED_DIMENSION_TOLERANCE_PT
                                    or abs(float(crop_box.height) - height)
                                    > _SLIDES_PDF_REQUIRED_DIMENSION_TOLERANCE_PT
                                ):
                                    LOGGER.warning(
                                        "Deck %s source PDF has mixed page sizes; using page 1 size %.1fx%.1f for export.",
                                        deck.deck_id,
                                        width,
                                        height,
                                    )
                                    break
                            return (width, height)
            except (fitz.FileDataError, RuntimeError, ValueError, OSError):
                LOGGER.warning(
                    "Failed to read source PDF size for deck %s; using fallback %.1fx%.1f.",
                    deck.deck_id,
                    fallback[0],
                    fallback[1],
                )
    return fallback


def _render_deck_pdf_with_reportlab(
    slide_images: list[bytes],
    *,
    prompt_style: str,
    page_size_pt: tuple[float, float] | None = None,
) -> io.BytesIO:
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    buffer = io.BytesIO()
    if page_size_pt is None:
        width = float(_SLIDES_EXPORT_PAGE_WIDTH_PT)
        height = float(_SLIDES_EXPORT_PAGE_HEIGHT_PT)
    else:
        width = float(page_size_pt[0])
        height = float(page_size_pt[1])
    pdf = canvas.Canvas(buffer, pagesize=(width, height))
    style = load_notebooklm_style(prompt_style)
    font_size = max(9.0, style.body_size_pt * 0.6)
    box_height = 0.22 * 96
    text_offset = 3
    for slide_index, image_bytes in enumerate(slide_images, start=1):
        image = ImageReader(io.BytesIO(image_bytes))
        pdf.drawImage(image, 0, 0, width=width, height=height)
        pdf.setFont("Helvetica", font_size)
        pdf.drawCentredString(
            width / 2, min(box_height - 1, text_offset), str(slide_index)
        )
        pdf.showPage()
    pdf.save()
    buffer.seek(0)
    return buffer


def _render_deck_pdf_with_fitz(
    slide_images: list[bytes],
    *,
    prompt_style: str,
    page_size_pt: tuple[float, float] | None = None,
) -> io.BytesIO:
    doc = fitz.open()
    if page_size_pt is None:
        width = float(_SLIDES_EXPORT_PAGE_WIDTH_PT)
        height = float(_SLIDES_EXPORT_PAGE_HEIGHT_PT)
    else:
        width = float(page_size_pt[0])
        height = float(page_size_pt[1])
    style = load_notebooklm_style(prompt_style)
    font_size = max(9.0, style.body_size_pt * 0.6)
    box_width = 0.9 * 96
    box_height = 0.22 * 96
    text_offset = 3
    for slide_index, image_bytes in enumerate(slide_images, start=1):
        page = doc.new_page(width=width, height=height)
        page.insert_image(fitz.Rect(0, 0, width, height), stream=image_bytes)
        box_left = (width - box_width) / 2
        box_rect = fitz.Rect(
            box_left, height - box_height, box_left + box_width, height
        )
        page.insert_textbox(
            box_rect,
            str(slide_index),
            fontsize=font_size,
            align=fitz.TEXT_ALIGN_CENTER,
        )
    buffer = io.BytesIO()
    doc.save(buffer)
    doc.close()
    buffer.seek(0)
    return buffer


def _render_deck_pptx_from_slide_images(
    slide_images: list[bytes],
    *,
    page_size_pt: tuple[float, float] | None = None,
) -> io.BytesIO:
    """Render browser screenshots into one full-slide image per PPTX slide."""

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    if not slide_images:
        raise ValueError("No slides could be rendered to PPTX.")
    if page_size_pt is None:
        page_width = float(_SLIDES_EXPORT_PAGE_WIDTH_PT)
        page_height = float(_SLIDES_EXPORT_PAGE_HEIGHT_PT)
    else:
        page_width = max(1.0, float(page_size_pt[0]))
        page_height = max(1.0, float(page_size_pt[1]))
    slide_height_in = float(_SLIDES_RENDERED_PPTX_HEIGHT_IN)
    slide_width_in = slide_height_in * (page_width / page_height)

    presentation = Presentation()
    presentation.slide_width = Inches(slide_width_in)
    presentation.slide_height = Inches(slide_height_in)
    blank_layout = presentation.slide_layouts[6]
    for image_bytes in slide_images:
        slide = presentation.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        slide.shapes.add_picture(
            io.BytesIO(image_bytes),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


def _render_deck_pdf_with_playwright(
    deck: Deck,
    deck_path: Path,
    *,
    page_size_pt: tuple[float, float] | None = None,
) -> io.BytesIO:
    if not HAS_PLAYWRIGHT:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="Playwright is not available on this server.",
        )

    doc = fitz.open()
    if page_size_pt is None:
        width = float(_SLIDES_EXPORT_PAGE_WIDTH_PT)
        height = float(_SLIDES_EXPORT_PAGE_HEIGHT_PT)
    else:
        width = float(page_size_pt[0])
        height = float(page_size_pt[1])
    viewport_width = max(1, int(round(width)))
    viewport_height = max(1, int(round(height)))
    style = load_notebooklm_style(deck.prompt_style)
    font_size = max(9.0, style.body_size_pt * 0.6)
    box_width = 0.9 * 96
    box_height = 0.22 * 96
    margin_spec = {"top": "0", "right": "0", "bottom": "0", "left": "0"}
    with sync_playwright() as playwright:  # pragma: no cover - requires Playwright
        browser = _launch_playwright_chromium(playwright)
        page = browser.new_page(
            viewport={
                "width": viewport_width,
                "height": viewport_height,
            },
        )
        page.emulate_media(media="screen")
        for slide in deck.slides:
            slide_file = deck_path / slide.id
            if slide.kind != "sectionHeader" and not slide_file.exists():
                LOGGER.warning(
                    "Slide %s missing at %s; skipping PDF render", slide.id, slide_file
                )
                continue
            try:
                html = (
                    slide.full_html
                    if slide.kind == "sectionHeader" and slide.full_html
                    else ""
                )
                if not html:
                    html = slide_file.read_text(encoding="utf-8")
                html = _rewrite_asset_sources_for_print(
                    html,
                    deck_id=deck.deck_id,
                    deck_path=deck_path,
                    prompt_style=deck.prompt_style,
                    page_size_pt=(width, height),
                )
                page.set_content(html, wait_until="load")
                page.wait_for_timeout(2000)
                slide_pdf = page.pdf(
                    width=f"{width:g}px",
                    height=f"{height:g}px",
                    print_background=True,
                    prefer_css_page_size=False,
                    margin=margin_spec,
                )
            except Exception:  # pragma: no cover - defensive
                LOGGER.exception("Failed to render slide %s as vector PDF", slide.id)
                continue
            try:
                source_doc = fitz.open(stream=slide_pdf, filetype="pdf")
            except Exception:  # pragma: no cover - defensive
                LOGGER.exception(
                    "Failed to parse Playwright PDF output for slide %s", slide.id
                )
                continue
            try:
                if source_doc.page_count > 0:
                    doc.insert_pdf(source_doc, from_page=0, to_page=0)
            finally:
                source_doc.close()
        browser.close()

    if doc.page_count <= 0:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No slides could be rendered to PDF.",
        )

    for slide_index in range(doc.page_count):
        page = doc[slide_index]
        _cover_notebooklm_logo_on_pdf_page(page)
        box_left = (width - box_width) / 2
        box_rect = fitz.Rect(
            box_left,
            height - box_height,
            box_left + box_width,
            height,
        )
        page.insert_textbox(
            box_rect,
            str(slide_index + 1),
            fontsize=font_size,
            align=fitz.TEXT_ALIGN_CENTER,
        )

    buffer = io.BytesIO()
    doc.save(buffer)
    doc.close()
    buffer.seek(0)
    return buffer


def _cover_notebooklm_logo_on_pdf_page(page: fitz.Page) -> None:
    """Draw a white rectangle over the NotebookLM logo area on a PDF page."""

    page_width = float(page.cropbox.width)
    page_height = float(page.cropbox.height)
    if page_width <= 0 or page_height <= 0:
        return
    # Match the image-based cover ratio used during upload.
    cover_width = max(100.0, page_width * 0.10)
    cover_height = max(15.0, page_height * 0.03)
    x0 = max(0.0, page_width - cover_width)
    y0 = max(0.0, page_height - cover_height)
    if page_width <= x0 or page_height <= y0:
        return
    page.draw_rect(
        fitz.Rect(x0, y0, page_width, page_height),
        color=None,
        fill=(1, 1, 1),
        width=0,
        overlay=True,
    )


def _asset_path_to_data_uri(asset_path: Path) -> str | None:
    try:
        payload = asset_path.read_bytes()
    except OSError as exc:
        LOGGER.warning("Failed to read slide asset for print export: %s", exc)
        return None
    mime_type, _encoding = mimetypes.guess_type(asset_path.name)
    if not mime_type:
        mime_type = "application/octet-stream"
    encoded = base64.b64encode(payload).decode("ascii")
    return f"data:{mime_type};base64,{encoded}"


def _extract_asset_relative_path(value: str, asset_prefix: str) -> str | None:
    if value.startswith(asset_prefix):
        return value[len(asset_prefix) :]
    parsed = urlparse(value)
    if parsed.path.startswith(asset_prefix):
        return parsed.path[len(asset_prefix) :]
    return None


def _rewrite_style_urls(style_value: str, *, deck_id: str, deck_path: Path) -> str:
    asset_prefix = f"/slides/deck/{deck_id}/assets/"
    if asset_prefix not in style_value:
        return style_value

    def _replace(match: re.Match[str]) -> str:
        raw_url = match.group(1).strip().strip("'\"")
        relative_path = _extract_asset_relative_path(raw_url, asset_prefix)
        if not relative_path:
            return match.group(0)
        local_path = (deck_path / "assets" / relative_path).resolve()
        data_uri = _asset_path_to_data_uri(local_path)
        if not data_uri:
            return match.group(0)
        return f'url("{data_uri}")'

    return re.sub(r"url\(([^)]+)\)", _replace, style_value)


def _rewrite_asset_sources_for_print(
    html: str,
    *,
    deck_id: str,
    deck_path: Path,
    prompt_style: str,
    page_size_pt: tuple[float, float] | None = None,
) -> str:
    asset_prefix = f"/slides/deck/{deck_id}/assets/"
    if not html:
        return html
    should_process = (
        asset_prefix in html
        or "section_header.css" in html
        or "slide-title" in html
        or "slide-body" in html
    )
    if not should_process:
        return html
    soup = BeautifulSoup(html, "html.parser")
    if page_size_pt is None:
        page_width = float(_SLIDES_EXPORT_PAGE_WIDTH_PT)
        page_height = float(_SLIDES_EXPORT_PAGE_HEIGHT_PT)
    else:
        page_width = float(page_size_pt[0])
        page_height = float(page_size_pt[1])
    page_width = max(1.0, page_width)
    page_height = max(1.0, page_height)

    def _ensure_head() -> Any:
        head = soup.head
        if head is not None:
            return head
        head = soup.new_tag("head")
        if soup.html is not None:
            soup.html.insert(0, head)
        else:  # pragma: no cover - defensive
            soup.insert(0, head)
        return head

    title_element = soup.select_one(".slide-title")
    body_element = soup.select_one(".slide-body")
    title_text = (
        _clean_text(title_element.get_text(" ", strip=True))
        if title_element is not None
        else ""
    )
    body_text = (
        _clean_text(body_element.get_text(" ", strip=True))
        if body_element is not None
        else ""
    )
    body_has_image = body_element is not None and body_element.find("img") is not None
    is_image_only_slide = bool(body_has_image and not body_text and not title_text)

    head = _ensure_head()
    if (
        head.find("style", attrs={"data-inline-stylesheet": "slides_export_base.css"})
        is None
    ):
        base_css = (
            "html, body {\n"
            "  margin: 0;\n"
            "  padding: 0;\n"
            f"  width: {page_width:g}px;\n"
            f"  height: {page_height:g}px;\n"
            "  overflow: hidden;\n"
            "}\n"
            "*, *::before, *::after {\n"
            "  box-sizing: border-box;\n"
            "}\n"
            ".slide-container {\n"
            "  position: relative;\n"
            f"  width: {page_width:g}px;\n"
            f"  height: {page_height:g}px;\n"
            "  margin: 0;\n"
            "  padding: 0;\n"
            "  overflow: hidden;\n"
            "}\n"
            ".slide-title,\n"
            ".slide-body {\n"
            "  margin: 0;\n"
            "  padding: 0;\n"
            "}\n"
        )
        style_tag = soup.new_tag("style")
        style_tag["data-inline-stylesheet"] = "slides_export_base.css"
        style_tag.string = base_css
        head.append(style_tag)

    if is_image_only_slide:
        container = soup.select_one(".slide-container")
        if container is not None:
            classes = container.get("class") or []
            if "slide-image-only-export" not in classes:
                container["class"] = [*classes, "slide-image-only-export"]
        if (
            head.find(
                "style", attrs={"data-inline-stylesheet": "image_slide_export.css"}
            )
            is None
        ):
            image_css = (
                ".slide-image-only-export .slide-title {\n"
                "  display: none !important;\n"
                "  margin: 0 !important;\n"
                "  padding: 0 !important;\n"
                "  font-size: 0 !important;\n"
                "  line-height: 0 !important;\n"
                "}\n"
                ".slide-image-only-export .slide-body {\n"
                "  position: absolute;\n"
                "  inset: 0;\n"
                "  margin: 0 !important;\n"
                "  padding: 0 !important;\n"
                "}\n"
                ".slide-image-only-export .slide-body > div {\n"
                "  position: absolute !important;\n"
                "  inset: 0 !important;\n"
                "  margin: 0 !important;\n"
                "  padding: 0 !important;\n"
                "  width: 100% !important;\n"
                "  height: 100% !important;\n"
                "}\n"
                ".slide-image-only-export .slide-body img {\n"
                "  display: block !important;\n"
                "  margin: 0 !important;\n"
                "  width: 100% !important;\n"
                "  height: 100% !important;\n"
                "  max-width: none !important;\n"
                "  max-height: none !important;\n"
                "  object-fit: contain !important;\n"
                "  object-position: center !important;\n"
                "}\n"
            )
            style_tag = soup.new_tag("style")
            style_tag["data-inline-stylesheet"] = "image_slide_export.css"
            style_tag.string = image_css
            head.append(style_tag)

    section_header_links: list[Any] = []
    for link in soup.find_all("link"):
        rel = link.get("rel") or []
        rel_values = [rel] if isinstance(rel, str) else rel
        if not any(str(item).lower() == "stylesheet" for item in rel_values):
            continue
        href = link.get("href")
        if not href:
            continue
        parsed = urlparse(href)
        if parsed.scheme or parsed.netloc:
            continue
        if Path(parsed.path).name != "section_header.css":
            continue
        section_header_links.append(link)
    if section_header_links:
        css_text = ""
        css_candidates = (
            (deck_path / "section_header.css").resolve(),
            (deck_path / "assets" / "section_header.css").resolve(),
        )
        for candidate in css_candidates:
            try:
                css_text = candidate.read_text(encoding="utf-8")
            except FileNotFoundError:
                continue
            except OSError as exc:
                LOGGER.warning("Failed to read %s for PDF export: %s", candidate, exc)
                continue
            if css_text.strip():
                break
        if not css_text.strip():
            source = Path("static/css/section_header_viewer.css")
            if source.exists():
                try:
                    style = load_notebooklm_style(prompt_style)
                    css_vars = build_notebooklm_css_variables(style)
                    css_text = f"{css_vars}\n{source.read_text(encoding='utf-8')}"
                except (OSError, ValueError, KeyError) as exc:
                    LOGGER.warning(
                        "Failed to build fallback section header CSS for PDF export: %s",
                        exc,
                    )
                    css_text = ""
        if css_text.strip():
            for link in section_header_links:
                link.decompose()
            head = _ensure_head()
            if (
                head.find(
                    "style", attrs={"data-inline-stylesheet": "section_header.css"}
                )
                is None
            ):
                style_tag = soup.new_tag("style")
                style_tag["data-inline-stylesheet"] = "section_header.css"
                style_tag.string = css_text
                head.append(style_tag)
    if (
        soup.select_one(".slide-title")
        and soup.select_one(".slide-body")
        and not is_image_only_slide
    ):
        head = _ensure_head()
        if (
            head.find("style", attrs={"data-inline-stylesheet": "intro_slide.css"})
            is None
        ):
            intro_css = ""
            try:
                style = load_notebooklm_style(prompt_style)
                css_vars = build_notebooklm_css_variables(style)
                intro_css = (
                    f"{css_vars}\n"
                    ".slide-title,\n"
                    ".slide-title * {\n"
                    '  font-family: var(--notebooklm-font-stack, "Inter", "Roboto", sans-serif);\n'
                    "}\n"
                    ".slide-body,\n"
                    ".slide-body * {\n"
                    '  font-family: var(--notebooklm-font-stack, "Inter", "Roboto", sans-serif);\n'
                    "}\n"
                    ".slide-title {\n"
                    "  font-size: var(--notebooklm-title-size-px, 42.67px);\n"
                    "  line-height: var(--notebooklm-line-height, 1.25);\n"
                    "}\n"
                    ".slide-body {\n"
                    "  font-size: var(--notebooklm-body-size-px, 24px);\n"
                    "  line-height: var(--notebooklm-line-height, 1.25);\n"
                    "}\n"
                )
            except (OSError, ValueError, KeyError) as exc:
                LOGGER.warning(
                    "Failed to build intro slide CSS for PDF export: %s", exc
                )
                intro_css = ""
            if intro_css.strip():
                style_tag = soup.new_tag("style")
                style_tag["data-inline-stylesheet"] = "intro_slide.css"
                style_tag.string = intro_css
                head.append(style_tag)
    for tag in soup.find_all(["img", "source"]):
        for attr in ("src", "data-src"):
            value = tag.get(attr)
            if not value:
                continue
            relative_path = _extract_asset_relative_path(value, asset_prefix)
            if not relative_path:
                continue
            local_path = (deck_path / "assets" / relative_path).resolve()
            data_uri = _asset_path_to_data_uri(local_path)
            if data_uri:
                tag[attr] = data_uri
            else:
                tag[attr] = local_path.as_uri()
    for tag in soup.find_all(style=True):
        style_value = tag.get("style")
        if not style_value:
            continue
        tag["style"] = _rewrite_style_urls(
            style_value, deck_id=deck_id, deck_path=deck_path
        )
    for style_tag in soup.find_all("style"):
        if not style_tag.string:
            continue
        style_tag.string = _rewrite_style_urls(
            style_tag.string, deck_id=deck_id, deck_path=deck_path
        )
    return soup.decode()


def _render_deck_slide_images(
    deck: Deck,
    deck_path: Path,
    *,
    page_size_pt: tuple[float, float] | None = None,
) -> list[bytes]:
    if not HAS_PLAYWRIGHT:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="Playwright is not available on this server.",
        )
    slide_images: list[bytes] = []
    if page_size_pt is None:
        width = float(_SLIDES_EXPORT_PAGE_WIDTH_PT)
        height = float(_SLIDES_EXPORT_PAGE_HEIGHT_PT)
    else:
        width = float(page_size_pt[0])
        height = float(page_size_pt[1])
    viewport_width = max(1, int(round(width)))
    viewport_height = max(1, int(round(height)))
    with sync_playwright() as playwright:  # pragma: no cover - requires Playwright
        browser = _launch_playwright_chromium(playwright)
        page = browser.new_page(
            viewport={
                "width": viewport_width,
                "height": viewport_height,
            },
            device_scale_factor=_SLIDES_EXPORT_SCREENSHOT_SCALE,
        )
        for slide in deck.slides:
            slide_file = deck_path / slide.id
            if slide.kind != "sectionHeader" and not slide_file.exists():
                LOGGER.warning(
                    "Slide %s missing at %s; skipping PDF render", slide.id, slide_file
                )
                continue
            try:
                html = (
                    slide.full_html
                    if slide.kind == "sectionHeader" and slide.full_html
                    else ""
                )
                if not html:
                    html = slide_file.read_text(encoding="utf-8")
                html = _rewrite_asset_sources_for_print(
                    html,
                    deck_id=deck.deck_id,
                    deck_path=deck_path,
                    prompt_style=deck.prompt_style,
                    page_size_pt=(width, height),
                )
                page.set_content(html, wait_until="load")
                page.wait_for_timeout(2000)
            except Exception:  # pragma: no cover - defensive
                LOGGER.exception("Failed to load slide %s for PDF export", slide.id)
                continue
            screenshot = page.screenshot(
                full_page=False,
                type="png",
                clip={
                    "x": 0,
                    "y": 0,
                    "width": viewport_width,
                    "height": viewport_height,
                },
            )
            slide_images.append(_cover_notebooklm_logo_in_image_bytes(screenshot))
        browser.close()
    return slide_images


def _is_missing_playwright_browser_error(exc: BaseException) -> bool:
    message = str(exc).lower()
    return "executable doesn't exist" in message and "playwright install" in message


def _truncate_dependency_error(detail: str, *, max_chars: int = 600) -> str:
    normalized = " ".join(str(detail).split())
    if len(normalized) <= max_chars:
        return normalized
    return f"{normalized[:max_chars].rstrip()}..."


def _discover_playwright_browser_executables() -> list[Path]:
    override_root = os.getenv("SLIDES_PLAYWRIGHT_BROWSERS_PATH", "").strip()
    if override_root:
        roots = [Path(override_root).expanduser()]
    else:
        roots = [_runtime_user_home_dir() / ".cache" / "ms-playwright"]

    candidates: list[Path] = []
    seen: set[Path] = set()
    for root in roots:
        if not root.exists():
            continue
        for pattern in ("chromium-*", "chromium_headless_shell-*"):
            for folder in sorted(root.glob(pattern), reverse=True):
                for binary_name in ("chrome", "headless_shell"):
                    binary = folder / "chrome-linux" / binary_name
                    if binary.exists():
                        resolved = binary.resolve()
                        if resolved in seen:
                            continue
                        seen.add(resolved)
                        candidates.append(resolved)
    return candidates


def _runtime_user_home_dir() -> Path:
    try:
        import pwd  # local import for compatibility across platforms

        return Path(pwd.getpwuid(os.getuid()).pw_dir)
    except Exception:  # pragma: no cover - platform-specific fallback
        return Path.home()


def _resolve_slides_executable_path_override() -> Path | None:
    raw_path = os.getenv("SLIDES_PLAYWRIGHT_EXECUTABLE_PATH", "").strip()
    if not raw_path:
        return None
    candidate = Path(raw_path).expanduser()
    if candidate.exists():
        return candidate.resolve()
    LOGGER.warning("SLIDES_PLAYWRIGHT_EXECUTABLE_PATH does not exist: %s", candidate)
    return None


def _launch_playwright_chromium(playwright: Any) -> Any:
    missing_exc: BaseException | None = None
    try:
        return playwright.chromium.launch()
    except PlaywrightError as exc:
        if not _is_missing_playwright_browser_error(exc):
            raise
        missing_exc = exc

    executable_override = _resolve_slides_executable_path_override()
    if executable_override is not None:
        try:
            return playwright.chromium.launch(executable_path=str(executable_override))
        except PlaywrightError:
            LOGGER.warning(
                "Configured slide browser executable failed to launch: %s",
                executable_override,
            )

    # Playwright Python package is present, but its expected browser payload is missing.
    # First try system browser channels, then existing cached browser binaries.
    for channel in ("chrome", "msedge"):
        try:
            return playwright.chromium.launch(channel=channel)
        except PlaywrightError:
            continue

    for executable in _discover_playwright_browser_executables():
        try:
            return playwright.chromium.launch(executable_path=str(executable))
        except PlaywrightError:
            continue

    raise RuntimeError(
        "Playwright browser executable is unavailable for this runtime user. "
        "Set `SLIDES_PLAYWRIGHT_EXECUTABLE_PATH` or `SLIDES_PLAYWRIGHT_BROWSERS_PATH` "
        "for the API service user."
    ) from missing_exc


def _render_deck_pdf(deck: Deck, deck_path: Path) -> io.BytesIO:
    page_size_pt = _resolve_deck_export_page_size(deck, deck_path=deck_path)
    if _SLIDES_EXPORT_ENABLE_VECTOR_PDF:
        try:
            return _render_deck_pdf_with_playwright(
                deck,
                deck_path,
                page_size_pt=page_size_pt,
            )
        except (
            PlaywrightError,
            RuntimeError,
            OSError,
            ValueError,
            HTTPException,
        ) as exc:
            LOGGER.exception(
                "Vector PDF export failed for deck %s; falling back to raster export",
                deck.deck_id,
                exc_info=exc,
            )
    slide_images = _render_deck_slide_images(
        deck,
        deck_path,
        page_size_pt=page_size_pt,
    )
    if not slide_images:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No slides could be rendered to PDF.",
        )
    if HAS_REPORTLAB:
        return _render_deck_pdf_with_reportlab(
            slide_images,
            prompt_style=deck.prompt_style,
            page_size_pt=page_size_pt,
        )
    return _render_deck_pdf_with_fitz(
        slide_images,
        prompt_style=deck.prompt_style,
        page_size_pt=page_size_pt,
    )


def _render_deck_pptx(deck: Deck, deck_path: Path) -> io.BytesIO:
    page_size_pt = _resolve_deck_export_page_size(deck, deck_path=deck_path)
    slide_images = _render_deck_slide_images(
        deck,
        deck_path,
        page_size_pt=page_size_pt,
    )
    return _render_deck_pptx_from_slide_images(
        slide_images,
        page_size_pt=page_size_pt,
    )


def _extract_bbox(candidate: object) -> tuple[float, float, float, float] | None:
    if not isinstance(candidate, dict):
        return None
    try:
        x = float(candidate.get("x", 0.0))
        y = float(candidate.get("y", 0.0))
        w = float(candidate.get("w", 0.0))
        h = float(candidate.get("h", 0.0))
    except (TypeError, ValueError):
        return None
    if w <= 0 or h <= 0:
        return None
    return x, y, w, h


def _parse_hex_color(value: object) -> tuple[int, int, int] | None:
    if not isinstance(value, str):
        return None
    normalized = value.strip().lstrip("#")
    if len(normalized) != 6:
        return None
    try:
        return tuple(int(normalized[index : index + 2], 16) for index in (0, 2, 4))
    except ValueError:
        return None


def _normalize_joined_ocr_text(text: str) -> str:
    normalized = " ".join(str(text or "").split())
    if not normalized:
        return ""
    normalized = re.sub(r"\s+([,.;:%!?])", r"\1", normalized)
    normalized = re.sub(r"([\(\[\{])\s+", r"\1", normalized)
    normalized = re.sub(r"\s+([\)\]\}])", r"\1", normalized)
    normalized = re.sub(r"\s+([’'])\s+", r"\1", normalized)
    return normalized.strip()


def _bbox_overlap_ratio(
    a: tuple[float, float, float, float], b: tuple[float, float, float, float]
) -> float:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    if aw <= 0 or ah <= 0 or bw <= 0 or bh <= 0:
        return 0.0
    ax1 = ax + aw
    ay1 = ay + ah
    bx1 = bx + bw
    by1 = by + bh
    overlap_w = max(0.0, min(ax1, bx1) - max(ax, bx))
    overlap_h = max(0.0, min(ay1, by1) - max(ay, by))
    overlap_area = overlap_w * overlap_h
    return overlap_area / max(aw * ah, 1.0)


def _segment_ocr_lines(
    *,
    line_tokens: list[dict[str, object]],
    image_width: int,
    split_gap_ratio: float = 0.06,
    split_gap_min: float = 90.0,
) -> list[dict[str, object]]:
    if not line_tokens:
        return []
    sorted_tokens = sorted(
        line_tokens,
        key=lambda item: (
            float(item["y"]) + (float(item["h"]) / 2.0),
            float(item["x"]),
        ),
    )
    rows: list[dict[str, object]] = []
    for token in sorted_tokens:
        cy = float(token["y"]) + (float(token["h"]) / 2.0)
        token_h = float(token["h"])
        target_row: dict[str, object] | None = None
        for row in rows:
            row_cy = float(row["cy"])
            row_h = float(row["h"])
            if abs(cy - row_cy) <= max(token_h, row_h) * 0.65:
                target_row = row
                break
        if target_row is None:
            rows.append({"cy": cy, "h": token_h, "items": [token]})
            continue
        target_items = target_row["items"]
        if isinstance(target_items, list):
            target_items.append(token)
            target_row["cy"] = (
                (float(target_row["cy"]) * (len(target_items) - 1)) + cy
            ) / max(len(target_items), 1)
            target_row["h"] = max(float(target_row["h"]), token_h)

    rows = sorted(rows, key=lambda row: float(row["cy"]))
    segmented_lines: list[dict[str, object]] = []
    split_gap = max(float(split_gap_min), float(image_width) * float(split_gap_ratio))
    for row in rows:
        row_items = row.get("items")
        if not isinstance(row_items, list):
            continue
        ordered = sorted(
            (item for item in row_items if isinstance(item, dict)),
            key=lambda item: float(item["x"]),
        )
        if not ordered:
            continue
        current_segment: list[dict[str, object]] = [ordered[0]]
        for token in ordered[1:]:
            previous = current_segment[-1]
            prev_right = float(previous["x"]) + float(previous["w"])
            gap = float(token["x"]) - prev_right
            if gap > split_gap:
                segmented_lines.append({"items": list(current_segment)})
                current_segment = [token]
            else:
                current_segment.append(token)
        segmented_lines.append({"items": list(current_segment)})

    merged: list[dict[str, object]] = []
    for segment in segmented_lines:
        items = segment.get("items")
        if not isinstance(items, list):
            continue
        cleaned_items = [item for item in items if isinstance(item, dict)]
        if not cleaned_items:
            continue
        texts = [str(item.get("text") or "").strip() for item in cleaned_items]
        joined = _normalize_joined_ocr_text(" ".join(texts))
        if not joined:
            continue
        min_x = min(float(item["x"]) for item in cleaned_items)
        min_y = min(float(item["y"]) for item in cleaned_items)
        max_x = max(float(item["x"]) + float(item["w"]) for item in cleaned_items)
        max_y = max(float(item["y"]) + float(item["h"]) for item in cleaned_items)
        merged.append(
            {
                "text": joined,
                "x": min_x,
                "y": min_y,
                "w": max_x - min_x,
                "h": max_y - min_y,
            }
        )
    return merged


def _merge_adjacent_title_blocks(
    *,
    blocks: list[dict[str, object]],
    image_size: tuple[int, int],
) -> list[dict[str, object]]:
    if not blocks:
        return []
    image_w = max(1, int(image_size[0]))
    image_h = max(1, int(image_size[1]))
    ordered = [
        dict(block)
        for block in sorted(
            blocks, key=lambda item: (float(item["y"]), float(item["x"]))
        )
    ]
    title_present = any(bool(item.get("is_title")) for item in ordered)
    if not title_present:
        promotable_indices = [
            index
            for index, block in enumerate(ordered)
            if float(block.get("y", 0.0)) <= float(image_h) * 0.30
            and float(block.get("w", 0.0)) >= float(image_w) * 0.22
            and not re.match(
                r"^\s*(?:[-*•·▪◦]|\d+[.)]|[A-Za-z][.)])\s+",
                str(block.get("text") or "").strip(),
            )
        ]
        if promotable_indices:
            first_index = promotable_indices[0]
            ordered[first_index]["is_title"] = True
            if first_index + 1 < len(ordered):
                first = ordered[first_index]
                second = ordered[first_index + 1]
                if float(second.get("y", 0.0)) <= float(
                    image_h
                ) * 0.34 and not re.match(
                    r"^\s*(?:[-*•·▪◦]|\d+[.)]|[A-Za-z][.)])\s+",
                    str(second.get("text") or "").strip(),
                ):
                    first_x = float(first.get("x", 0.0))
                    first_w = float(first.get("w", 0.0))
                    second_x = float(second.get("x", 0.0))
                    second_w = float(second.get("w", 0.0))
                    overlap_w = max(
                        0.0,
                        min(first_x + first_w, second_x + second_w)
                        - max(first_x, second_x),
                    )
                    overlap_ratio = overlap_w / max(second_w, 1.0)
                    gap = float(second.get("y", 0.0)) - (
                        float(first.get("y", 0.0)) + float(first.get("h", 0.0))
                    )
                    if overlap_ratio >= 0.35 and gap <= max(
                        42.0, float(second.get("h", 0.0)) * 1.8
                    ):
                        ordered[first_index + 1]["is_title"] = True

    merged: list[dict[str, object]] = []
    top_title_limit = float(image_h) * 0.36
    for block in ordered:
        current = dict(block)
        current["is_title"] = bool(current.get("is_title"))
        if not merged:
            merged.append(current)
            continue
        previous = merged[-1]
        if not (bool(previous.get("is_title")) and bool(current.get("is_title"))):
            merged.append(current)
            continue
        prev_y = float(previous.get("y", 0.0))
        curr_y = float(current.get("y", 0.0))
        if prev_y > top_title_limit or curr_y > top_title_limit:
            merged.append(current)
            continue
        prev_x = float(previous.get("x", 0.0))
        prev_w = float(previous.get("w", 0.0))
        prev_h = float(previous.get("h", 0.0))
        curr_x = float(current.get("x", 0.0))
        curr_w = float(current.get("w", 0.0))
        curr_h = float(current.get("h", 0.0))
        overlap_w = max(
            0.0, min(prev_x + prev_w, curr_x + curr_w) - max(prev_x, curr_x)
        )
        overlap_ratio = overlap_w / max(curr_w, 1.0)
        same_column = overlap_ratio >= 0.32 or abs(curr_x - prev_x) <= max(
            42.0, prev_w * 0.20
        )
        gap = curr_y - (prev_y + prev_h)
        if not same_column or gap > max(46.0, curr_h * 1.9, prev_h * 1.9):
            merged.append(current)
            continue
        previous_text = str(previous.get("text") or "").strip()
        current_text = str(current.get("text") or "").strip()
        if current_text and current_text.lower() != previous_text.lower():
            previous["text"] = (
                f"{previous_text}\n{current_text}" if previous_text else current_text
            )
        merged_left = min(prev_x, curr_x)
        merged_top = min(prev_y, curr_y)
        merged_right = max(prev_x + prev_w, curr_x + curr_w)
        merged_bottom = max(prev_y + prev_h, curr_y + curr_h)
        previous["x"] = merged_left
        previous["y"] = merged_top
        previous["w"] = max(1.0, merged_right - merged_left)
        previous["h"] = max(1.0, merged_bottom - merged_top)
        previous["is_title"] = True
    return merged


def _build_export_text_blocks_from_ocr_slide(
    *,
    ocr_slide: dict[str, object] | None,
    image_size: tuple[int, int],
    chart_bboxes: list[tuple[float, float, float, float]],
    split_gap_ratio: float = 0.06,
    split_gap_min: float = 90.0,
) -> list[dict[str, object]]:
    if ocr_slide is None:
        return []
    lines = ocr_slide.get("lines")
    if not isinstance(lines, list):
        return []
    source_w = max(1, int(image_size[0]))
    source_h = max(1, int(image_size[1]))
    ignored_values = {
        "text",
        "chart",
        "figure_title",
        "paragraph_title",
        "region",
        "doc_title",
    }
    line_tokens: list[dict[str, object]] = []
    for raw_line in lines:
        if not isinstance(raw_line, dict):
            continue
        text = str(raw_line.get("text") or "").strip()
        if not text:
            continue
        if text.lower() in ignored_values:
            continue
        bbox = _extract_bbox(raw_line.get("bbox"))
        if bbox is None:
            continue
        if any(
            _bbox_overlap_ratio(bbox, chart_bbox) > 0.18 for chart_bbox in chart_bboxes
        ):
            continue
        x, y, w, h = bbox
        line_tokens.append({"text": text, "x": x, "y": y, "w": w, "h": h})
    if not line_tokens:
        return []

    segmented_lines = _segment_ocr_lines(
        line_tokens=line_tokens,
        image_width=source_w,
        split_gap_ratio=split_gap_ratio,
        split_gap_min=split_gap_min,
    )
    if not segmented_lines:
        return []

    segmented_lines = sorted(
        segmented_lines,
        key=lambda item: (float(item["y"]), float(item["x"])),
    )

    blocks: list[dict[str, object]] = []
    x_tolerance = float(source_w) * 0.05
    right_tolerance = float(source_w) * 0.10
    for line in segmented_lines:
        x = float(line["x"])
        y = float(line["y"])
        w = float(line["w"])
        h = float(line["h"])
        text = str(line["text"] or "").strip()
        if not text:
            continue
        if not blocks:
            blocks.append(
                {
                    "x": x,
                    "y": y,
                    "w": w,
                    "h": h,
                    "text_lines": [text],
                }
            )
            continue
        current = blocks[-1]
        current_x = float(current["x"])
        current_w = float(current["w"])
        current_bottom = float(current["y"]) + float(current["h"])
        current_right = current_x + current_w
        line_right = x + w
        vertical_gap = y - current_bottom
        same_column = (
            abs(x - current_x) <= x_tolerance
            and abs(line_right - current_right) <= right_tolerance
        )
        line_height_ref = max(float(current["h"]), h)
        should_merge = same_column and vertical_gap <= (line_height_ref * 1.35 + 6.0)
        if should_merge:
            current["x"] = min(current_x, x)
            current["y"] = min(float(current["y"]), y)
            merged_right = max(current_right, line_right)
            merged_bottom = max(current_bottom, y + h)
            current["w"] = max(1.0, merged_right - float(current["x"]))
            current["h"] = max(1.0, merged_bottom - float(current["y"]))
            if isinstance(current["text_lines"], list):
                current["text_lines"].append(text)
            continue
        blocks.append({"x": x, "y": y, "w": w, "h": h, "text_lines": [text]})

    if not blocks:
        return []

    title_candidates = [
        block
        for block in blocks
        if float(block["y"]) <= (float(source_h) * 0.24)
        and float(block["w"]) >= (float(source_w) * 0.30)
    ]
    title_block = (
        min(title_candidates, key=lambda block: float(block["y"]))
        if title_candidates
        else None
    )

    finalized: list[dict[str, object]] = []
    for block in blocks:
        text_lines = (
            block.get("text_lines") if isinstance(block.get("text_lines"), list) else []
        )
        cleaned_lines: list[str] = []
        for item in text_lines:
            normalized = _normalize_joined_ocr_text(str(item or ""))
            if normalized:
                cleaned_lines.append(normalized)
        if not cleaned_lines:
            continue
        text = "\n".join(cleaned_lines)
        is_title = title_block is block
        finalized.append(
            {
                "text": text,
                "x": float(block["x"]),
                "y": float(block["y"]),
                "w": float(block["w"]),
                "h": float(block["h"]),
                "is_title": is_title,
            }
        )

    if chart_bboxes and finalized:
        narrative = [item for item in finalized if not bool(item["is_title"])]
        if len(narrative) >= 2:
            min_x = min(float(item["x"]) for item in narrative)
            max_x = max(float(item["x"] + item["w"]) for item in narrative)
            if (max_x - min_x) <= float(source_w) * 0.40:
                top = min(float(item["y"]) for item in narrative)
                bottom = max(float(item["y"] + item["h"]) for item in narrative)
                merged_lines: list[str] = []
                for item in sorted(narrative, key=lambda entry: float(entry["y"])):
                    line_text = str(item["text"] or "").strip()
                    if not line_text:
                        continue
                    if not re.match(r"^\s*(?:[-*•])\s+", line_text):
                        line_text = f"• {line_text}"
                    merged_lines.append(line_text)
                if merged_lines:
                    merged_item = {
                        "text": "\n".join(merged_lines),
                        "x": min_x,
                        "y": top,
                        "w": max(1.0, max_x - min_x),
                        "h": max(1.0, bottom - top),
                        "is_title": False,
                    }
                    finalized = [item for item in finalized if bool(item["is_title"])]
                    finalized.append(merged_item)

    finalized = _merge_adjacent_title_blocks(blocks=finalized, image_size=image_size)
    return sorted(finalized, key=lambda item: (float(item["y"]), float(item["x"])))


def _build_deck_response(
    deck: Deck,
    storage: DeckStorage,
    *,
    archived_source_decks: list[str] | None = None,
) -> DeckResponse:
    deck_path = storage.root / deck.deck_id
    has_pdf = (deck_path / "source.pdf").exists()
    has_images = any("<img" in (slide.body_html or "") for slide in deck.slides)
    has_layout = storage.load_layout_payload(deck.deck_id) is not None
    return DeckResponse(
        deckId=deck.deck_id,
        promptStyle=deck.prompt_style,
        ownerEmail=deck.owner_email,
        sharedWith=list(deck.shared_with),
        slides=[SlidePayload.from_domain(slide) for slide in deck.slides],
        sections=[SectionPayload.from_domain(section) for section in deck.sections],
        thumbnails=_build_thumbnails(deck),
        hasPdf=has_pdf,
        hasImages=has_images,
        hasLayout=has_layout,
        archivedSourceDecks=archived_source_decks or [],
    )


def _build_thumbnails(deck: Deck) -> dict[str, str]:
    thumbnails: dict[str, str] = {}
    for slide in deck.slides:
        thumbnails[slide.id] = _extract_slide_thumbnail(slide)
    return thumbnails


def _extract_slide_thumbnail(slide: Slide) -> str:
    if slide.kind == "sectionHeader":
        label = _clean_text(slide.title_html or "") or "Section"
        return f"<strong>{label}</strong>"
    if slide.body_html:
        soup = BeautifulSoup(slide.body_html, "html.parser")
        image_tag = soup.find("img")
        if image_tag:
            image_src = image_tag.get("src") or image_tag.get("data-src")
            if image_src:
                image_alt = image_tag.get("alt", "")
                return (
                    f'<img src="{escape(image_src, quote=True)}" '
                    f'alt="{escape(image_alt, quote=True)}" />'
                )
    body_text = _clean_text(slide.body_html or "")
    if not body_text:
        body_text = _clean_text(slide.title_html or slide.id)
    snippet = body_text[:160]
    return f"<p>{escape(snippet)}</p>"


_CLEAN_PATTERN = re.compile(r"<[^>]+>")


def _clean_text(value: str) -> str:
    if not value:
        return ""
    stripped = _CLEAN_PATTERN.sub("", value)
    return stripped.replace("\n", " ").strip()


def _start_print_job(deck_id: str) -> str:
    _cleanup_print_jobs()
    job_id = secrets.token_hex(16)
    create_job(job_id, deck_id)
    _PRINT_EXECUTOR.submit(_run_print_job, job_id)
    return job_id


def _resolve_print_job_output(job: PrintJob) -> Path | None:
    output_path = job.output_path
    if output_path and output_path.exists():
        return output_path
    fallback_path = _PRINT_JOB_DIR / f"{job.job_id}.pdf"
    if fallback_path.exists():
        set_output_path(job.job_id, fallback_path)
        return fallback_path
    return None


def _run_print_job(job_id: str) -> None:
    storage = get_storage()
    job = get_job(job_id)
    if job is None:
        return
    update_job_status(job_id, "running")
    output_path = (_PRINT_JOB_DIR / f"{job_id}.pdf").resolve()
    temp_output_path = output_path.with_suffix(".pdf.tmp")
    try:
        deck = storage.load_deck(job.deck_id)
        if not deck.slides:
            raise ValueError("Deck is empty.")
        deck_path = storage.root / job.deck_id
        if not deck_path.exists():
            raise ValueError(f"Deck path {deck_path} missing.")
        pdf_buffer = _render_deck_pdf(deck, deck_path)
        with temp_output_path.open("wb") as handle:
            handle.write(pdf_buffer.getvalue())
        temp_output_path.replace(output_path)
        update_job_status(job_id, "succeeded")
        set_output_path(job_id, output_path)
    except Exception as exc:  # pragma: no cover - background task
        LOGGER.exception("Print job %s failed", job_id)
        temp_output_path.unlink(missing_ok=True)
        output_path.unlink(missing_ok=True)
        update_job_status(job_id, "failed", detail=str(exc))


def _get_print_job_or_404(job_id: str) -> PrintJob:
    _cleanup_print_jobs()
    job = get_job(job_id)
    if job is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="Print job not found."
        )
    return job


def _cleanup_print_jobs() -> None:
    cleanup_expired_jobs(_PRINT_JOB_TTL)


def _start_pptx_job(
    deck_id: str,
    *,
    source: Literal["rendered", "template"] = "rendered",
) -> str:
    _cleanup_pptx_jobs()
    job_id = secrets.token_hex(16)
    create_pptx_job(job_id, deck_id, source=source)
    _PPTX_EXECUTOR.submit(_run_pptx_job, job_id)
    return job_id


def _resolve_pptx_job_output(job: PptxJob) -> Path | None:
    output_path = job.output_path
    if output_path and output_path.exists():
        return output_path
    fallback_path = _PPTX_JOB_DIR / f"{job.job_id}.pptx"
    if fallback_path.exists():
        set_pptx_output_path(job.job_id, fallback_path)
        return fallback_path
    return None


def _run_pptx_job(job_id: str) -> None:
    storage = get_storage()
    job = get_pptx_job_record(job_id)
    if job is None:
        return
    update_pptx_job_status(job_id, "running")
    output_path = (_PPTX_JOB_DIR / f"{job_id}.pptx").resolve()
    temp_output_path = output_path.with_suffix(".pptx.tmp")
    try:
        deck = storage.load_deck(job.deck_id)
        if not deck.slides:
            raise ValueError("Deck is empty.")
        deck_path = storage.root / job.deck_id
        if not deck_path.exists():
            raise ValueError(f"Deck path {deck_path} missing.")
        if job.source == "template":
            slide_analysis = _load_cached_slide_analysis_payload(job.deck_id, storage)
            spec = build_slides_pptx_spec(
                deck,
                deck_path,
                slide_analysis=slide_analysis,
            )
            write_slides_pptx_spec(deck_path, spec)
            pptx_buffer = render_slides_pptx_from_template(deck_path)
        else:
            spec = None
            pptx_buffer = _render_deck_pptx(deck, deck_path)
        with temp_output_path.open("wb") as handle:
            handle.write(pptx_buffer.getvalue())
        if spec is not None and _SLIDES_PPTX_POST_RENDER_COMPARE_ENABLED:
            try:
                repaired_spec, _report = apply_post_render_compare_loop(
                    deck=deck,
                    deck_path=deck_path,
                    spec=spec,
                    pptx_path=temp_output_path,
                    job_id=job_id,
                )
                spec = repaired_spec
            except Exception as exc:  # pragma: no cover - best-effort validation
                LOGGER.warning(
                    "Skipping PPTX post-render compare loop for job %s after failure: %s",
                    job_id,
                    exc,
                )
        temp_output_path.replace(output_path)
        update_pptx_job_status(job_id, "succeeded")
        set_pptx_output_path(job_id, output_path)
    except Exception as exc:  # pragma: no cover - background task
        LOGGER.exception("PPTX job %s failed", job_id)
        temp_output_path.unlink(missing_ok=True)
        output_path.unlink(missing_ok=True)
        update_pptx_job_status(job_id, "failed", detail=str(exc))


def _get_pptx_job_or_404(job_id: str) -> PptxJob:
    _cleanup_pptx_jobs()
    job = get_pptx_job_record(job_id)
    if job is None:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND, detail="PPTX job not found."
        )
    return job


def _cleanup_pptx_jobs() -> None:
    cleanup_expired_pptx_jobs(_PPTX_JOB_TTL)


def _resolve_slide_image_path(
    deck: Deck,
    deck_path: Path,
    slide: Slide,
) -> Path | None:
    html_sources = [slide.body_html or "", slide.full_html or ""]
    for source_html in html_sources:
        if not source_html:
            continue
        soup = BeautifulSoup(source_html, "html.parser")
        image = soup.find("img")
        if image is None:
            continue
        src = str(image.get("src") or "").strip()
        if not src or src.startswith("data:"):
            continue
        prefix = f"/slides/deck/{deck.deck_id}/assets/"
        if src.startswith(prefix):
            relative = src[len(prefix) :]
        else:
            relative = src
        relative = relative.split("?", 1)[0].split("#", 1)[0]
        relative = relative.lstrip("/")
        if relative.startswith("assets/"):
            relative = relative[len("assets/") :]
        if not relative:
            continue
        candidate = deck_path / "assets" / Path(relative)
        if candidate.exists():
            return candidate
    return None


def _resolve_local_image_dimensions(image_path: Path) -> tuple[int, int] | None:
    try:
        with Image.open(image_path) as image:
            width, height = image.size
    except (UnidentifiedImageError, OSError, ValueError):
        return None
    if width <= 0 or height <= 0:
        return None
    return width, height


def _resolve_slide_image_dimensions(
    deck: Deck, deck_path: Path, slide: Slide
) -> tuple[int, int] | None:
    image_path = _resolve_slide_image_path(deck, deck_path, slide)
    if image_path is None:
        return None
    return _resolve_local_image_dimensions(image_path)


def _encode_local_image_to_data_uri(image_path: Path) -> str:
    try:
        with Image.open(image_path) as image:
            has_alpha = image.mode in {"RGBA", "LA"} or (
                image.mode == "P" and "transparency" in image.info
            )
            if has_alpha:
                rgba = image.convert("RGBA")
                flattened = Image.new("RGB", rgba.size, (255, 255, 255))
                flattened.paste(rgba, mask=rgba.split()[-1])
                export = flattened
            else:
                export = image.convert("RGB")
            _cover_notebooklm_logo(export)
            buffer = io.BytesIO()
            export.save(buffer, format="PNG")
            payload = buffer.getvalue()
    except (UnidentifiedImageError, OSError, ValueError):
        payload = image_path.read_bytes()
    encoded = base64.b64encode(payload).decode("ascii")
    return f"data:image/png;base64,{encoded}"
