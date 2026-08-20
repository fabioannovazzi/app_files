#!/usr/bin/env python3
"""Prepare and mechanically audit Clara advisory-deliverable validation runs.

The helper is deterministic because supported-format extraction, hashing,
declared-schema validation, cross-field consistency, and packaging are
mechanically verifiable. It does not select material claims or judge support,
reasoning, recommendations, corrections, or professional readiness.
"""

from __future__ import annotations

import argparse
import hashlib
import html
import importlib.util
import json
import logging
import re
import shutil
import sys
from html.parser import HTMLParser
from pathlib import Path
from typing import Any
from zipfile import BadZipFile

from jsonschema import Draft202012Validator

__all__ = [
    "AdvisoryValidationError",
    "prepare_validation",
    "package_validation",
    "read_supported_deliverable",
    "validate_advisory_contract",
    "validate_review_record",
]

LOGGER = logging.getLogger(__name__)
CLARA_ROOT = Path(__file__).resolve().parents[3]
REPOSITORY_ROOT = Path(__file__).resolve().parents[5]
CONTRACT_SCHEMA_PATH = CLARA_ROOT / "contracts" / "advisory_contract.v1.schema.json"
REVIEW_SCHEMA_PATH = (
    Path(__file__).resolve().parents[1]
    / "references"
    / "advisory_validation_review.schema.json"
)
LINEAGE_SCRIPT_PATH = CLARA_ROOT / "scripts" / "advisory_evidence_lineage.py"
EVIDENCE_REGISTER_FILENAME = "advisory_evidence_register.json"
CLAIM_REGISTER_FILENAME = "advisory_claim_register.json"
SUPPORTED_PRIMARY_SUFFIXES = {
    ".docx",
    ".htm",
    ".html",
    ".markdown",
    ".md",
    ".pdf",
    ".pptx",
    ".txt",
}
SUPPORTED_SOURCE_SUFFIXES = SUPPORTED_PRIMARY_SUFFIXES | {".csv", ".parquet", ".xlsx"}
LANGUAGES = {"it", "en", "fr", "de", "es"}
REVIEW_DIMENSIONS = (
    "contract_conformance",
    "factual_source_support",
    "calculations_data_provenance",
    "reasoning_assumptions",
    "contradictions_missing_evidence",
    "recommendation_evidence_decision_fit",
    "professional_judgement_boundaries",
    "correction_needs",
    "residual_uncertainty",
    "delivery_readiness",
)
REQUIRED_CONTRACT_FIELDS = {
    "schema_version",
    "decision",
    "purpose",
    "audience",
    "deliverable_type",
    "output_language",
    "scope_included",
    "scope_excluded",
    "available_inputs",
    "evidence_requirements",
    "analysis_plan",
    "assumptions",
    "unresolved_questions",
    "success_criteria",
    "selected_clara_workflow",
    "validation_profile",
    "validation_scope",
    "correction_policy",
    "professional_judgement_policy",
}
DIMENSION_STATUSES = {
    "conforms",
    "partially_conforms",
    "does_not_conform",
    "contradicted",
    "uncertain",
    "judgment_required",
    "not_applicable",
}
CORRECTION_STATUSES = {
    "not_needed",
    "proposed",
    "completed",
    "blocked",
    "professional_review_required",
}
FORMAT_CHECK_STATUSES = {
    "passed",
    "issues_found",
    "blocked",
    "not_run",
    "not_applicable",
}
READINESS_STATUSES = {
    "ready",
    "ready_with_residual_uncertainty",
    "not_ready",
    "blocked",
}
APPROVAL_STATUSES = {"approved", "pending", "not_required"}
PROVENANCE_MODES = {"generation_time", "matched_support"}
CHAIN_SUPPORT_STATUSES = {
    "adequate",
    "partial",
    "unsupported",
    "contradicted",
    "uncertain",
    "not_applicable",
}
CHAIN_REASONING_STATUSES = {
    "sound",
    "gap",
    "contradicted",
    "uncertain",
    "not_applicable",
}
RECHECK_KINDS = {"none", "web", "calculation", "file", "transcript", "other"}
RECHECK_STATUSES = {"not_required", "completed", "pending", "blocked"}
CHAIN_RESOLUTION_STATUSES = {
    "no_change",
    "corrected",
    "removed",
    "qualified",
    "pending",
    "professional_review_required",
}
URL_RE = re.compile(r"https?://[^\s)\]>\"']+", re.IGNORECASE)
CITATION_RE = re.compile(r"\[(?:\^?[A-Za-z0-9_-]+|\d+(?:,\s*\d+)*)\]")
FOOTNOTE_RE = re.compile(r"^\[\^?([A-Za-z0-9_-]+)\]:\s*(.+)$", re.MULTILINE)
NUMBER_RE = re.compile(
    r"(?<![\w.])(?:(?:EUR|USD|GBP|[$€£])\s*)?-?\d[\d.,]*(?:\s*%|\s*(?:bps|EUR|USD|GBP))?(?!\w)",
    re.IGNORECASE,
)
FORMULA_RE = re.compile(
    r"(?:[$€£]?\s*-?\d[\d.,]*\s*)(?:[+*/=]|\s-\s)(?:\s*[$€£]?\s*-?\d[\d.,%]*)"
)


class AdvisoryValidationError(ValueError):
    """Raised when a mechanically verifiable validation contract is invalid."""


class _HTMLTextExtractor(HTMLParser):
    _NON_VISIBLE_ELEMENTS = {"script", "style", "template"}

    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._non_visible_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in self._NON_VISIBLE_ELEMENTS:
            self._non_visible_depth += 1
            return
        if self._non_visible_depth:
            return
        if tag == "a":
            for key, value in attrs:
                if key.casefold() == "href" and value:
                    self.parts.append(f" {value} ")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._NON_VISIBLE_ELEMENTS and self._non_visible_depth:
            self._non_visible_depth -= 1
            return
        if self._non_visible_depth:
            return
        if tag in {"p", "div", "br", "li", "section", "article", "h1", "h2", "h3"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._non_visible_depth:
            self.parts.append(html.unescape(data))

    def text(self) -> str:
        return re.sub(r"\n{3,}", "\n\n", "".join(self.parts)).strip()


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _write_json(path: Path, payload: Any) -> None:
    path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def _load_json(path: Path) -> Any:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        raise AdvisoryValidationError(f"invalid JSON in {path}: {exc}") from exc


def _lineage_module() -> Any:
    """Load the shared Clara lineage helper without changing import paths."""

    module_name = "clara_advisory_evidence_lineage"
    existing = sys.modules.get(module_name)
    if existing is not None:
        return existing
    spec = importlib.util.spec_from_file_location(module_name, LINEAGE_SCRIPT_PATH)
    if spec is None or spec.loader is None:
        raise AdvisoryValidationError(
            f"cannot load advisory evidence lineage helper: {LINEAGE_SCRIPT_PATH}"
        )
    module = importlib.util.module_from_spec(spec)
    sys.modules[module_name] = module
    spec.loader.exec_module(module)
    return module


def _is_non_empty_string(value: Any) -> bool:
    return isinstance(value, str) and bool(value.strip())


def _is_string_list(value: Any, *, allow_empty: bool) -> bool:
    return (
        isinstance(value, list)
        and (allow_empty or bool(value))
        and all(_is_non_empty_string(item) for item in value)
    )


def _within(path: Path, parent: Path) -> bool:
    try:
        path.resolve().relative_to(parent.resolve())
    except ValueError:
        return False
    return True


def _ensure_output_location(output_dir: Path) -> None:
    if _within(output_dir, REPOSITORY_ROOT):
        raise AdvisoryValidationError(
            "validation outputs must be outside the Git workspace"
        )


def _reject_output_collisions(
    outputs: dict[str, Path], protected_inputs: dict[str, Path]
) -> None:
    """Reject exact path aliases because source preservation is mechanical."""

    resolved_outputs = {name: path.resolve() for name, path in outputs.items()}
    for input_name, input_path in protected_inputs.items():
        resolved_input = input_path.resolve()
        for output_name, output_path in resolved_outputs.items():
            if resolved_input == output_path:
                raise AdvisoryValidationError(
                    f"refusing to overwrite {input_name} with {output_name}: "
                    f"{output_path}"
                )


def validate_advisory_contract(payload: Any) -> list[str]:
    """Return mechanical shape errors for advisory_contract.json."""

    if not isinstance(payload, dict):
        return ["advisory contract must be an object"]
    missing = sorted(REQUIRED_CONTRACT_FIELDS - payload.keys())
    if missing:
        return [f"missing contract fields: {', '.join(missing)}"]
    try:
        schema = json.loads(CONTRACT_SCHEMA_PATH.read_text(encoding="utf-8"))
        Draft202012Validator.check_schema(schema)
    except (OSError, json.JSONDecodeError) as exc:
        return [f"cannot load advisory contract schema: {exc}"]
    validator = Draft202012Validator(schema)
    errors = [
        "contract schema "
        + (".".join(str(part) for part in error.absolute_path) or "$")
        + f": {error.message}"
        for error in sorted(
            validator.iter_errors(payload),
            key=lambda item: tuple(str(part) for part in item.absolute_path),
        )
    ]
    profile = payload.get("validation_profile")
    if isinstance(profile, dict):
        format_checks = profile.get("format_checks")
        if isinstance(format_checks, list):
            workflows = [
                item.get("workflow")
                for item in format_checks
                if isinstance(item, dict) and _is_non_empty_string(item.get("workflow"))
            ]
            if len(workflows) != len(set(workflows)):
                errors.append(
                    "validation_profile.format_checks workflows must be unique"
                )
    return errors


def _read_pdf(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError as exc:
        raise AdvisoryValidationError(
            "PyMuPDF is required for PDF extraction; run Clara's dependency check"
        ) from exc
    try:
        with fitz.open(path) as document:
            pages = [page.get_text("text") for page in document]
    except (fitz.FileDataError, RuntimeError) as exc:
        raise AdvisoryValidationError(f"PDF is unreadable or damaged: {path}") from exc
    text = "\n\n".join(pages).strip()
    if not text:
        raise AdvisoryValidationError(
            "PDF has no readable text layer; run Clara's input-aware OCR preflight"
        )
    return text, {"parser": "pymupdf_text", "page_count": len(pages)}


def _read_docx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from docx import Document  # type: ignore[import-not-found]
        from docx.opc.exceptions import PackageNotFoundError
    except ImportError as exc:
        raise AdvisoryValidationError(
            "python-docx is required for DOCX extraction; run Clara's dependency check"
        ) from exc
    try:
        document = Document(path)
    except (BadZipFile, PackageNotFoundError, KeyError) as exc:
        raise AdvisoryValidationError(f"DOCX is unreadable or damaged: {path}") from exc
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    table_cell_count = 0
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            table_cell_count += len(values)
            parts.append(" | ".join(values))
    return "\n\n".join(parts).strip(), {
        "parser": "python_docx",
        "paragraph_count": len(document.paragraphs),
        "table_count": len(document.tables),
        "table_cell_count": table_cell_count,
    }


def _read_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
        from pptx.exc import PackageNotFoundError
    except ImportError as exc:
        raise AdvisoryValidationError(
            "python-pptx is required for PPTX extraction; run Clara's dependency check"
        ) from exc
    try:
        presentation = Presentation(path)
    except (BadZipFile, PackageNotFoundError, KeyError) as exc:
        raise AdvisoryValidationError(f"PPTX is unreadable or damaged: {path}") from exc
    parts: list[str] = []
    text_shape_count = 0
    table_cell_count = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts.append(f"# Slide {slide_number}")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                parts.append(shape.text.strip())
                text_shape_count += 1
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    table_cell_count += len(values)
                    parts.append(" | ".join(values))
    return "\n\n".join(parts).strip(), {
        "parser": "python_pptx_visible_text",
        "slide_count": len(presentation.slides),
        "text_shape_count": text_shape_count,
        "table_cell_count": table_cell_count,
        "limitations": [
            "Visible text and table cells were extracted; visual meaning, charts, images, and speaker notes require separate review."
        ],
    }


def read_supported_deliverable(path: Path) -> tuple[str, dict[str, Any]]:
    """Extract text and mechanical parser metadata from one supported artifact."""

    suffix = path.suffix.casefold()
    if suffix not in SUPPORTED_PRIMARY_SUFFIXES:
        supported = ", ".join(sorted(SUPPORTED_PRIMARY_SUFFIXES))
        raise AdvisoryValidationError(
            f"unsupported primary deliverable format {suffix or '<none>'}; supported: {supported}"
        )
    if suffix in {".md", ".markdown", ".txt"}:
        return path.read_text(encoding="utf-8", errors="replace").strip(), {
            "parser": "plain_text"
        }
    if suffix in {".html", ".htm"}:
        extractor = _HTMLTextExtractor()
        extractor.feed(path.read_text(encoding="utf-8", errors="replace"))
        return extractor.text(), {"parser": "html_text"}
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    return _read_pptx(path)


def _ordered_unique(
    items: list[str], *, limit: int, strip_terminal_punctuation: bool = True
) -> list[str]:
    seen: set[str] = set()
    output: list[str] = []
    for item in items:
        cleaned = item.strip()
        if strip_terminal_punctuation:
            cleaned = cleaned.rstrip(".,;:!?)\\]}>'\"")
        if not cleaned or cleaned.casefold() in seen:
            continue
        seen.add(cleaned.casefold())
        output.append(cleaned)
        if len(output) >= limit:
            break
    return output


def _citation_inventory(text: str) -> dict[str, Any]:
    return {
        "schema_version": "1.0",
        "purpose": "Mechanical navigation inventory; not semantic support assessment.",
        "urls": _ordered_unique(URL_RE.findall(text), limit=500),
        "citation_markers": _ordered_unique(
            CITATION_RE.findall(text),
            limit=1000,
            strip_terminal_punctuation=False,
        ),
        "footnotes": [
            {"id": match.group(1), "text": match.group(2).strip()}
            for match in FOOTNOTE_RE.finditer(text)
        ][:500],
    }


def _calculation_inventory(text: str) -> dict[str, Any]:
    return {
        "schema_version": "1.0",
        "purpose": "Mechanical numeric-token inventory; not calculation verification or semantic materiality selection.",
        "numeric_tokens": _ordered_unique(NUMBER_RE.findall(text), limit=2000),
        "formula_like_fragments": _ordered_unique(FORMULA_RE.findall(text), limit=500),
    }


def _coverage_inventory(
    text: str, *, target_characters: int = 12_000
) -> dict[str, Any]:
    """Split extracted text into bounded navigation units without classifying it."""

    units: list[dict[str, Any]] = []
    start = 0
    unit_number = 1
    text_length = len(text)
    while start < text_length:
        candidate_end = min(text_length, start + target_characters)
        end = candidate_end
        if candidate_end < text_length:
            paragraph_break = text.rfind("\n\n", start, candidate_end)
            if paragraph_break > start + target_characters // 2:
                end = paragraph_break
        if end <= start:
            end = candidate_end
        unit_text = text[start:end]
        unit_id = f"unit-{unit_number:04d}"
        units.append(
            {
                "id": unit_id,
                "character_start": start,
                "character_end": end,
                "character_count": len(unit_text),
                "sha256": hashlib.sha256(unit_text.encode("utf-8")).hexdigest(),
            }
        )
        start = end
        while start < text_length and text[start].isspace():
            start += 1
        unit_number += 1
    return {
        "schema_version": "1.0",
        "purpose": (
            "Mechanical coverage units for scalable review; units do not identify "
            "claims, importance, support, or meaning."
        ),
        "extracted_text_sha256": hashlib.sha256(text.encode("utf-8")).hexdigest(),
        "target_characters_per_unit": target_characters,
        "units": units,
    }


def _copy_declared_file(source: Path, target: Path, *, label: str) -> Path:
    if source.resolve() == target.resolve():
        return target
    if target.exists() and target.read_bytes() != source.read_bytes():
        raise AdvisoryValidationError(
            f"refusing to overwrite a different {label}: {target}"
        )
    shutil.copyfile(source, target)
    return target


def _prepare_lineage(
    output_dir: Path,
    evidence_register: Path | None,
    claim_register: Path | None,
) -> tuple[dict[str, Any], dict[str, Path]]:
    """Validate and copy a complete generation-time lineage pair when supplied."""

    if (evidence_register is None) != (claim_register is None):
        raise AdvisoryValidationError(
            "evidence and claim registers must be supplied together"
        )
    paths: dict[str, Path] = {}
    if evidence_register is None or claim_register is None:
        return (
            {
                "schema_version": "1.0",
                "provenance_mode": "matched_support",
                "purpose": (
                    "No generation-time lineage was supplied. Clara must match "
                    "material deliverable claims to available support without "
                    "calling the result original provenance."
                ),
                "evidence_register": None,
                "claim_register": None,
                "counts": {"evidence": 0, "claims": 0, "active_claims": 0},
            },
            paths,
        )
    if not evidence_register.is_file():
        raise AdvisoryValidationError(
            f"evidence register does not exist: {evidence_register}"
        )
    if not claim_register.is_file():
        raise AdvisoryValidationError(
            f"claim register does not exist: {claim_register}"
        )
    if evidence_register.parent.resolve() != claim_register.parent.resolve():
        raise AdvisoryValidationError(
            "evidence and claim registers must come from the same case directory"
        )
    lineage = _lineage_module()
    audit = lineage.validate_lineage(evidence_register.parent)
    if not audit["valid"]:
        raise AdvisoryValidationError(
            "invalid advisory evidence lineage: " + "; ".join(audit["errors"])
        )
    evidence_copy = _copy_declared_file(
        evidence_register,
        output_dir / EVIDENCE_REGISTER_FILENAME,
        label="advisory evidence register",
    )
    claim_copy = _copy_declared_file(
        claim_register,
        output_dir / CLAIM_REGISTER_FILENAME,
        label="advisory claim register",
    )
    paths = {
        "evidence_register": evidence_copy,
        "claim_register": claim_copy,
    }
    has_generation_claims = audit["counts"]["claims"] > 0
    return (
        {
            "schema_version": "1.0",
            "provenance_mode": (
                "generation_time" if has_generation_claims else "matched_support"
            ),
            "purpose": (
                "Generation-time claim and evidence lineage copied from the case "
                "workspace after mechanical validation."
                if has_generation_claims
                else "Empty lineage registers were supplied. They do not establish generation-time claim provenance, so review must use matched support."
            ),
            "empty_generation_registers_supplied": not has_generation_claims,
            "evidence_register": {
                "source_path": str(evidence_register.resolve()),
                "copied_path": str(evidence_copy.resolve()),
                "sha256": _sha256(evidence_register),
            },
            "claim_register": {
                "source_path": str(claim_register.resolve()),
                "copied_path": str(claim_copy.resolve()),
                "sha256": _sha256(claim_register),
            },
            "counts": audit["counts"],
        },
        paths,
    )


def _copy_contract(source: Path, output_dir: Path) -> Path:
    target = output_dir / "advisory_contract.json"
    if source.resolve() == target.resolve():
        return target
    if target.exists() and target.read_bytes() != source.read_bytes():
        raise AdvisoryValidationError(
            f"refusing to overwrite a different advisory contract: {target}"
        )
    shutil.copyfile(source, target)
    return target


def _source_inventory(paths: list[Path]) -> dict[str, Any]:
    items: list[dict[str, Any]] = []
    seen: set[Path] = set()
    for source_number, source in enumerate(paths, start=1):
        resolved = source.resolve()
        if resolved in seen:
            continue
        seen.add(resolved)
        if not source.is_file():
            raise AdvisoryValidationError(f"source file does not exist: {source}")
        items.append(
            {
                "id": f"source-{source_number:04d}",
                "path": str(resolved),
                "name": source.name,
                "suffix": source.suffix.casefold(),
                "supported_source_type": source.suffix.casefold()
                in SUPPORTED_SOURCE_SUFFIXES,
                "byte_count": source.stat().st_size,
                "sha256": _sha256(source),
            }
        )
    return {
        "schema_version": "1.0",
        "purpose": "Selected-source identity inventory; contents remain subject to model-led review.",
        "sources": items,
    }


def _resolve_artifact_ref(reference: str, base_dir: Path) -> Path:
    path = Path(reference).expanduser()
    if not path.is_absolute():
        path = base_dir / path
    return path.resolve()


def _authoritative_format_result(
    workflow: str, payload: Any
) -> tuple[str, bool, str | None] | None:
    """Read only workflow-owned mechanical result shapes.

    This is deterministic because it verifies explicit status fields emitted by
    the authoritative workflow. It does not reinterpret the semantic review.
    """

    if not isinstance(payload, dict):
        return None
    if workflow == "clara:reporting-engine":
        if (
            payload.get("owner") != "clara.reporting-engine"
            or payload.get("schema_version") != "0.2"
        ):
            return None
        runner = payload.get("runner")
        render_proof = payload.get("render_proof")
        if not isinstance(runner, dict) or not isinstance(render_proof, dict):
            return "reporting_engine_render", False, None
        return (
            "reporting_engine_render",
            runner.get("returncode") == 0
            and render_proof.get("status") in {"rendered", "not_required_data_only"},
            None,
        )
    if workflow == "clara:html-deck":
        input_record = payload.get("input")
        input_sha256 = (
            str(input_record.get("sha256"))
            if isinstance(input_record, dict) and input_record.get("sha256")
            else None
        )
        if "browser" in payload and "viewports" in payload:
            return "html_browser_qa", payload.get("result") == "pass", input_sha256
        if "deck" in payload and "checks" in payload and "summary" in payload:
            return (
                "html_static_validation",
                payload.get("result") == "pass",
                input_sha256,
            )
        return None
    if workflow == "clara:deck-correction":
        if payload.get("source") == "clara_deck_revision_output_review_completion":
            summary = payload.get("summary")
            return (
                "deck_correction_completion",
                isinstance(summary, dict)
                and summary.get("status") == "complete"
                and summary.get("final_delivery_allowed") is True,
                None,
            )
        return None
    if workflow == "clara:claim-basis-map":
        if payload.get("source") != "clara_claim_basis_map_audit":
            return None
        return "claim_basis_map_audit", payload.get("result") == "pass", None
    return None


def _audit_format_check_artifacts(
    review: dict[str, Any],
    contract: dict[str, Any],
    base_dir: Path,
    *,
    deliverable_sha256: str,
) -> tuple[list[dict[str, Any]], list[str]]:
    """Verify local artifact identity without judging a check's semantics."""

    metadata: list[dict[str, Any]] = []
    errors: list[str] = []
    checks = review.get("format_specific_checks")
    if not isinstance(checks, list):
        return metadata, errors
    declared_checks = contract.get("validation_profile", {}).get("format_checks", [])
    declared_by_workflow = {
        item.get("workflow"): item
        for item in declared_checks
        if isinstance(item, dict) and _is_non_empty_string(item.get("workflow"))
    }
    for check in checks:
        if not isinstance(check, dict) or check.get("status") != "passed":
            continue
        workflow = check.get("workflow")
        references = check.get("artifact_refs")
        if not _is_non_empty_string(workflow) or not isinstance(references, list):
            continue
        actual_paths: set[Path] = set()
        authoritative_results: dict[str, bool] = {}
        for reference in references:
            if not _is_non_empty_string(reference):
                continue
            artifact_path = _resolve_artifact_ref(reference, base_dir)
            actual_paths.add(artifact_path)
            if not artifact_path.is_file():
                errors.append(
                    f"passed format check artifact does not exist: {workflow}: "
                    f"{artifact_path}"
                )
                continue
            metadata.append(
                {
                    "workflow": workflow,
                    "reference": reference,
                    "path": str(artifact_path),
                    "sha256": _sha256(artifact_path),
                    "byte_count": artifact_path.stat().st_size,
                }
            )
            try:
                artifact_payload = _load_json(artifact_path)
            except (AdvisoryValidationError, OSError, UnicodeDecodeError):
                artifact_payload = None
            result = _authoritative_format_result(str(workflow), artifact_payload)
            if result is not None:
                result_kind, passed, input_sha256 = result
                authoritative_results[result_kind] = passed
                metadata[-1]["authoritative_result"] = {
                    "kind": result_kind,
                    "passed": passed,
                }
                if input_sha256 is not None:
                    metadata[-1]["authoritative_result"]["input_sha256"] = input_sha256
                if (
                    result_kind in {"html_static_validation", "html_browser_qa"}
                    and input_sha256 != deliverable_sha256
                ):
                    errors.append(
                        "passed HTML Deck result is not bound to the prepared "
                        f"deliverable: {result_kind}"
                    )
        declared = declared_by_workflow.get(workflow)
        if check.get("status") == "passed":
            required_result_kinds = {
                "clara:claim-basis-map": {"claim_basis_map_audit"},
                "clara:html-deck": {
                    "html_static_validation",
                    "html_browser_qa",
                },
                "clara:reporting-engine": {"reporting_engine_render"},
                "clara:deck-correction": {"deck_correction_completion"},
            }.get(str(workflow))
            if required_result_kinds is None:
                errors.append(
                    f"passed format check has no authoritative result adapter: {workflow}"
                )
            else:
                missing_results = sorted(
                    required_result_kinds - set(authoritative_results)
                )
                if missing_results:
                    errors.append(
                        f"passed format check lacks authoritative result artifacts: {workflow}: "
                        + ", ".join(missing_results)
                    )
                failed_results = sorted(
                    kind for kind, passed in authoritative_results.items() if not passed
                )
                if failed_results:
                    errors.append(
                        f"passed format check contradicts authoritative failed result: {workflow}: "
                        + ", ".join(failed_results)
                    )
        if isinstance(declared, dict) and declared.get("requirement") == "required":
            for declared_reference in declared.get("artifact_refs", []):
                declared_path = _resolve_artifact_ref(declared_reference, base_dir)
                if declared_path not in actual_paths:
                    errors.append(
                        f"required format check omitted declared artifact: {workflow}: "
                        f"{declared_path}"
                    )
    return metadata, errors


def prepare_validation(
    deliverable: Path,
    advisory_contract: Path,
    output_dir: Path,
    *,
    source_files: list[Path] | None = None,
    evidence_register: Path | None = None,
    claim_register: Path | None = None,
) -> dict[str, Path]:
    """Write deterministic preparation artifacts for a validator run."""

    _ensure_output_location(output_dir)
    if not deliverable.is_file():
        raise AdvisoryValidationError(f"deliverable does not exist: {deliverable}")
    if not advisory_contract.is_file():
        raise AdvisoryValidationError(
            "advisory_contract.json is required; Clara must create it from explicit or user-confirmed context"
        )
    contract_payload = _load_json(advisory_contract)
    contract_errors = validate_advisory_contract(contract_payload)
    if contract_errors:
        raise AdvisoryValidationError("; ".join(contract_errors))
    text, parser_metadata = read_supported_deliverable(deliverable)
    if not text:
        raise AdvisoryValidationError(
            "deliverable extraction produced no readable text"
        )

    paths = {
        "advisory_contract": output_dir / "advisory_contract.json",
        "deliverable_inventory": output_dir / "deliverable_inventory.json",
        "extracted_deliverable": output_dir / "extracted_deliverable.md",
        "citation_inventory": output_dir / "citation_inventory.json",
        "calculation_inventory": output_dir / "calculation_inventory.json",
        "source_inventory": output_dir / "source_inventory.json",
        "coverage_inventory": output_dir / "coverage_inventory.json",
        "lineage_inventory": output_dir / "lineage_inventory.json",
    }
    if evidence_register is not None:
        paths["evidence_register"] = output_dir / EVIDENCE_REGISTER_FILENAME
    if claim_register is not None:
        paths["claim_register"] = output_dir / CLAIM_REGISTER_FILENAME
    protected_inputs = {"primary deliverable": deliverable}
    protected_inputs.update(
        {
            f"selected source {index}": source
            for index, source in enumerate(source_files or [], start=1)
        }
    )
    if evidence_register is not None:
        protected_inputs["advisory evidence register"] = evidence_register
    if claim_register is not None:
        protected_inputs["advisory claim register"] = claim_register
    _reject_output_collisions(paths, protected_inputs)
    contract_protected_outputs = dict(paths)
    if advisory_contract.resolve() == paths["advisory_contract"].resolve():
        contract_protected_outputs.pop("advisory_contract")
    _reject_output_collisions(
        contract_protected_outputs,
        {"advisory contract": advisory_contract},
    )

    output_dir.mkdir(parents=True, exist_ok=True)
    contract_copy = _copy_contract(advisory_contract, output_dir)
    source_inventory = _source_inventory(source_files or [])
    lineage_inventory, lineage_paths = _prepare_lineage(
        output_dir,
        evidence_register,
        claim_register,
    )
    paths.update(lineage_paths)
    coverage_inventory = _coverage_inventory(text)
    _write_json(paths["coverage_inventory"], coverage_inventory)
    _write_json(paths["lineage_inventory"], lineage_inventory)
    deliverable_hash = _sha256(deliverable)
    inventory = {
        "schema_version": "1.0",
        "source_path": str(deliverable.resolve()),
        "source_name": deliverable.name,
        "source_suffix": deliverable.suffix.casefold(),
        "source_sha256": deliverable_hash,
        "source_byte_count": deliverable.stat().st_size,
        "advisory_contract_path": str(contract_copy.resolve()),
        "advisory_contract_sha256": _sha256(contract_copy),
        "extraction": {
            **parser_metadata,
            "character_count": len(text),
            "word_count": len(re.findall(r"\S+", text)),
        },
        "boundary": {
            "semantic_selection": "model_led",
            "semantic_assessment": "model_led",
            "hidden_model_api_calls": False,
            "original_preserved": True,
        },
        "lineage": {
            "provenance_mode": lineage_inventory["provenance_mode"],
            "lineage_inventory_path": str(paths["lineage_inventory"].resolve()),
            "lineage_inventory_sha256": _sha256(paths["lineage_inventory"]),
        },
        "coverage_inventory_path": str(paths["coverage_inventory"].resolve()),
        "coverage_inventory_sha256": _sha256(paths["coverage_inventory"]),
        "source_inventory_path": str(paths["source_inventory"].resolve()),
    }
    extracted_path = paths["extracted_deliverable"]
    extracted_path.write_text(text.rstrip() + "\n", encoding="utf-8")
    paths["advisory_contract"] = contract_copy
    _write_json(paths["deliverable_inventory"], inventory)
    _write_json(paths["citation_inventory"], _citation_inventory(text))
    _write_json(paths["calculation_inventory"], _calculation_inventory(text))
    _write_json(paths["source_inventory"], source_inventory)
    inventory["source_inventory_sha256"] = _sha256(paths["source_inventory"])
    _write_json(paths["deliverable_inventory"], inventory)
    return paths


def _validate_dimension_review(value: Any, *, subject: str) -> list[str]:
    if not isinstance(value, dict):
        return [f"{subject} must be an object"]
    errors: list[str] = []
    required = {
        "status",
        "analysis",
        "evidence_refs",
        "issues",
        "correction_status",
        "professional_review_required",
    }
    missing = sorted(required - value.keys())
    if missing:
        errors.append(f"{subject} missing fields: {', '.join(missing)}")
        return errors
    if value["status"] not in DIMENSION_STATUSES:
        errors.append(f"{subject}.status is invalid")
    if not _is_non_empty_string(value["analysis"]):
        errors.append(f"{subject}.analysis is required")
    for field in ("evidence_refs", "issues"):
        if not _is_string_list(value[field], allow_empty=True):
            errors.append(f"{subject}.{field} must be a string array")
    if value["correction_status"] not in CORRECTION_STATUSES:
        errors.append(f"{subject}.correction_status is invalid")
    if not isinstance(value["professional_review_required"], bool):
        errors.append(f"{subject}.professional_review_required must be boolean")
    return errors


def _validate_approval_record(value: Any, *, subject: str) -> list[str]:
    if not isinstance(value, dict):
        return [f"{subject} must be an object"]
    errors: list[str] = []
    required = {"status", "approved_by", "evidence_refs"}
    missing = sorted(required - value.keys())
    if missing:
        return [f"{subject} missing fields: {', '.join(missing)}"]
    if value["status"] not in APPROVAL_STATUSES:
        errors.append(f"{subject}.status is invalid")
    if not isinstance(value["approved_by"], str):
        errors.append(f"{subject}.approved_by must be a string")
    if not _is_string_list(value["evidence_refs"], allow_empty=True):
        errors.append(f"{subject}.evidence_refs must be a string array")
    if value["status"] == "approved":
        if not _is_non_empty_string(value["approved_by"]):
            errors.append(f"{subject}.approved_by is required when approved")
        if not value["evidence_refs"]:
            errors.append(f"{subject}.evidence_refs is required when approved")
    elif value["approved_by"] or value["evidence_refs"]:
        errors.append(f"{subject} approver and evidence require approved status")
    return errors


def _claim_dependency_closure(
    selected_ids: set[str],
    claim_by_id: dict[str, dict[str, Any]],
) -> set[str]:
    closure = set(selected_ids)
    pending = list(selected_ids)
    while pending:
        claim_id = pending.pop()
        claim = claim_by_id.get(claim_id, {})
        dependency = claim.get("dependency")
        if not isinstance(dependency, dict):
            continue
        for dependency_id in dependency.get("claim_ids", []):
            dependency_id = str(dependency_id)
            if dependency_id not in closure:
                closure.add(dependency_id)
                pending.append(dependency_id)
    return closure


def _evidence_reference_closure(
    selected_ids: set[str],
    evidence_by_id: dict[str, dict[str, Any]],
) -> set[str]:
    """Return every receipt in the declared history of selected evidence."""

    closure = set(selected_ids)
    pending = list(selected_ids)
    while pending:
        evidence_id = pending.pop()
        receipt = evidence_by_id.get(evidence_id, {})
        for field in ("rechecks_evidence_id", "supersedes_evidence_id"):
            referenced_id = str(receipt.get(field, ""))
            if referenced_id and referenced_id not in closure:
                closure.add(referenced_id)
                pending.append(referenced_id)
    return closure


def _validate_recheck(value: Any, *, subject: str) -> list[str]:
    if not isinstance(value, dict):
        return [f"{subject} must be an object"]
    errors: list[str] = []
    required = {"required", "kind", "status", "evidence_ids", "analysis"}
    missing = sorted(required - value.keys())
    if missing:
        return [f"{subject} missing fields: {', '.join(missing)}"]
    if not isinstance(value["required"], bool):
        errors.append(f"{subject}.required must be boolean")
    if value["kind"] not in RECHECK_KINDS:
        errors.append(f"{subject}.kind is invalid")
    if value["status"] not in RECHECK_STATUSES:
        errors.append(f"{subject}.status is invalid")
    if not _is_string_list(value["evidence_ids"], allow_empty=True):
        errors.append(f"{subject}.evidence_ids must be a string array")
    if not _is_non_empty_string(value["analysis"]):
        errors.append(f"{subject}.analysis is required")
    if value["required"]:
        if value["kind"] == "none" or value["status"] == "not_required":
            errors.append(f"{subject} required recheck needs a kind and active status")
        if value["status"] == "completed" and not value["evidence_ids"]:
            errors.append(f"{subject} completed recheck requires evidence_ids")
    elif (
        value["kind"] != "none"
        or value["status"] != "not_required"
        or value["evidence_ids"]
    ):
        errors.append(f"{subject} non-required recheck must use none/not_required")
    return errors


def _validate_resolution(value: Any, *, subject: str) -> list[str]:
    if not isinstance(value, dict):
        return [f"{subject} must be an object"]
    if set(value) != {"status", "explanation"}:
        return [f"{subject} must contain status and explanation"]
    errors: list[str] = []
    if value["status"] not in CHAIN_RESOLUTION_STATUSES:
        errors.append(f"{subject}.status is invalid")
    if not _is_non_empty_string(value["explanation"]):
        errors.append(f"{subject}.explanation is required")
    return errors


def _validate_chain_item(
    value: Any,
    *,
    subject: str,
    id_field: str,
    known_evidence_ids: set[str],
    require_known_evidence_ids: bool = True,
) -> list[str]:
    if not isinstance(value, dict):
        return [f"{subject} must be an object"]
    required = {
        id_field,
        "statement",
        "deliverable_locations",
        "evidence_ids",
        "dependency_claim_ids",
        "support_status",
        "reasoning_status",
        "contradiction_resolution",
        "analysis",
        "recheck",
        "resolution",
    }
    missing = sorted(required - value.keys())
    if missing:
        return [f"{subject} missing fields: {', '.join(missing)}"]
    errors: list[str] = []
    if not _is_non_empty_string(value[id_field]):
        errors.append(f"{subject}.{id_field} is required")
    if not _is_non_empty_string(value["statement"]):
        errors.append(f"{subject}.statement is required")
    for field in ("deliverable_locations", "evidence_ids", "dependency_claim_ids"):
        if not _is_string_list(value[field], allow_empty=True):
            errors.append(f"{subject}.{field} must be a string array")
    if value["support_status"] not in CHAIN_SUPPORT_STATUSES:
        errors.append(f"{subject}.support_status is invalid")
    if value["reasoning_status"] not in CHAIN_REASONING_STATUSES:
        errors.append(f"{subject}.reasoning_status is invalid")
    if not isinstance(value["contradiction_resolution"], str):
        errors.append(f"{subject}.contradiction_resolution must be a string")
    if not _is_non_empty_string(value["analysis"]):
        errors.append(f"{subject}.analysis is required")
    if (
        value["support_status"] == "adequate"
        and not value["evidence_ids"]
        and not value["dependency_claim_ids"]
    ):
        errors.append(
            f"{subject}: adequate support requires declared evidence_ids or claim dependencies"
        )
    errors.extend(_validate_recheck(value["recheck"], subject=f"{subject}.recheck"))
    errors.extend(
        _validate_resolution(value["resolution"], subject=f"{subject}.resolution")
    )
    for evidence_id in value.get("evidence_ids", []):
        if require_known_evidence_ids and evidence_id not in known_evidence_ids:
            errors.append(f"{subject}: unknown evidence id {evidence_id}")
    recheck = value.get("recheck")
    if isinstance(recheck, dict):
        for evidence_id in recheck.get("evidence_ids", []):
            if require_known_evidence_ids and evidence_id not in known_evidence_ids:
                errors.append(f"{subject}.recheck: unknown evidence id {evidence_id}")
    return errors


def _validate_lineage_review(
    value: Any,
    *,
    provenance_mode: str | None,
    claim_register: dict[str, Any] | None,
    evidence_register: dict[str, Any] | None,
    deliverable_sha256: str | None = None,
    matched_support_source_ids: set[str] | None = None,
) -> list[str]:
    if not isinstance(value, dict):
        return ["lineage_review must be an object"]
    required = {
        "provenance_mode",
        "selection_method",
        "reviewed_claim_ids",
        "chain_assessments",
        "untracked_material_claims",
        "limitations",
        "analysis",
    }
    missing = sorted(required - value.keys())
    if missing:
        return [f"lineage_review missing fields: {', '.join(missing)}"]
    errors: list[str] = []
    if value["provenance_mode"] not in PROVENANCE_MODES:
        errors.append("lineage_review.provenance_mode is invalid")
    if provenance_mode and value["provenance_mode"] != provenance_mode:
        errors.append("lineage_review.provenance_mode must match preparation")
    if value["selection_method"] != "model_led_claim_chain_review":
        errors.append(
            "lineage_review.selection_method must be model_led_claim_chain_review"
        )
    for field in ("reviewed_claim_ids", "limitations"):
        if not _is_string_list(value[field], allow_empty=True):
            errors.append(f"lineage_review.{field} must be a string array")
    if not isinstance(value["chain_assessments"], list):
        errors.append("lineage_review.chain_assessments must be an array")
    if not isinstance(value["untracked_material_claims"], list):
        errors.append("lineage_review.untracked_material_claims must be an array")
    if not _is_non_empty_string(value["analysis"]):
        errors.append("lineage_review.analysis is required")

    claims = (
        claim_register.get("claims", []) if isinstance(claim_register, dict) else []
    )
    evidence = (
        evidence_register.get("evidence", [])
        if isinstance(evidence_register, dict)
        else []
    )
    claim_by_id = {
        str(item.get("id")): item
        for item in claims
        if isinstance(item, dict) and item.get("id")
    }
    evidence_ids = {
        str(item.get("id"))
        for item in evidence
        if isinstance(item, dict) and item.get("id")
    }
    allowed_basis_ids = (
        set(matched_support_source_ids or set())
        if value.get("provenance_mode") == "matched_support"
        else evidence_ids
    )
    evidence_by_id = {
        str(item.get("id")): item
        for item in evidence
        if isinstance(item, dict) and item.get("id")
    }
    reviewed_ids = {
        str(item) for item in value.get("reviewed_claim_ids", []) if str(item)
    }
    if len(reviewed_ids) != len(value.get("reviewed_claim_ids", [])):
        errors.append("lineage_review.reviewed_claim_ids must be unique")
    if value.get("provenance_mode") == "matched_support" and reviewed_ids:
        errors.append("matched-support review cannot claim upstream reviewed_claim_ids")
    if value.get("provenance_mode") == "matched_support" and value.get(
        "chain_assessments"
    ):
        errors.append(
            "matched-support review must put reconstructed claims in untracked_material_claims"
        )
    for claim_id in sorted(reviewed_ids):
        if claim_id not in claim_by_id:
            errors.append(f"lineage_review: unknown reviewed claim id {claim_id}")

    assessment_ids: list[str] = []
    for index, assessment in enumerate(value.get("chain_assessments", [])):
        subject = f"lineage_review.chain_assessments[{index}]"
        errors.extend(
            _validate_chain_item(
                assessment,
                subject=subject,
                id_field="claim_id",
                known_evidence_ids=allowed_basis_ids,
            )
        )
        if not isinstance(assessment, dict):
            continue
        claim_id = str(assessment.get("claim_id", ""))
        assessment_ids.append(claim_id)
        upstream_claim = claim_by_id.get(claim_id)
        if upstream_claim is None and claim_by_id:
            errors.append(f"{subject}: unknown claim id {claim_id}")
            continue
        if upstream_claim is not None:
            if assessment.get("statement") != upstream_claim.get("statement"):
                errors.append(f"{subject}.statement must match the upstream claim")
            dependency = upstream_claim.get("dependency", {})
            expected_dependencies = (
                dependency.get("claim_ids", []) if isinstance(dependency, dict) else []
            )
            if set(assessment.get("dependency_claim_ids", [])) != set(
                expected_dependencies
            ):
                errors.append(f"{subject}.dependency_claim_ids must match lineage")
            directly_linked_evidence = {
                str(link.get("evidence_id"))
                for link in upstream_claim.get("evidence_links", [])
                if isinstance(link, dict) and link.get("evidence_id")
            }
            if isinstance(dependency, dict) and dependency.get(
                "calculation_evidence_id"
            ):
                directly_linked_evidence.add(str(dependency["calculation_evidence_id"]))
            expected_evidence = _evidence_reference_closure(
                directly_linked_evidence,
                evidence_by_id,
            )
            if set(assessment.get("evidence_ids", [])) != expected_evidence:
                errors.append(f"{subject}.evidence_ids must match lineage")
            relationships = {
                str(link.get("relationship"))
                for link in upstream_claim.get("evidence_links", [])
                if isinstance(link, dict)
            }
            if (
                relationships & {"weakens", "contradicts"}
                and assessment.get("support_status") == "adequate"
                and not str(assessment.get("contradiction_resolution", "")).strip()
            ):
                errors.append(
                    f"{subject}: adequate support requires contradiction_resolution when lineage weakens or contradicts the claim"
                )
            if claim_id in reviewed_ids:
                current_appearances = [
                    appearance
                    for appearance in upstream_claim.get("appearances", [])
                    if isinstance(appearance, dict)
                    and appearance.get("artifact_sha256") == deliverable_sha256
                ]
                if not current_appearances:
                    errors.append(
                        f"{subject}: reviewed generation-time claim has no hash-bound appearance in the prepared deliverable"
                    )
                appearance_locators = {
                    str(appearance.get("locator"))
                    for appearance in current_appearances
                    if appearance.get("locator")
                }
                reviewed_locations = set(assessment.get("deliverable_locations", []))
                if not reviewed_locations:
                    errors.append(
                        f"{subject}: selected material claim requires a deliverable location"
                    )
                elif not reviewed_locations <= appearance_locators:
                    errors.append(
                        f"{subject}.deliverable_locations must match hash-bound claim appearances"
                    )
            recheck = assessment.get("recheck")
            if isinstance(recheck, dict) and recheck.get("status") == "completed":
                for recheck_id in recheck.get("evidence_ids", []):
                    receipt = evidence_by_id.get(str(recheck_id))
                    if not isinstance(receipt, dict):
                        continue
                    prior_id = str(receipt.get("rechecks_evidence_id", ""))
                    if not prior_id or prior_id not in expected_evidence:
                        errors.append(
                            f"{subject}.recheck: completed receipt {recheck_id} must recheck an evidence receipt in the reviewed chain"
                        )
                    verification = receipt.get("verification")
                    status = (
                        verification.get("status")
                        if isinstance(verification, dict)
                        else None
                    )
                    if status not in {"rechecked_unchanged", "rechecked_changed"}:
                        errors.append(
                            f"{subject}.recheck: receipt {recheck_id} lacks completed recheck verification"
                        )
                    expected_types = {
                        "web": "web_capture",
                        "calculation": "calculation_run",
                        "transcript": "interview_transcript",
                    }
                    expected_type = expected_types.get(str(recheck.get("kind")))
                    if expected_type and receipt.get("evidence_type") != expected_type:
                        errors.append(
                            f"{subject}.recheck: {recheck_id} must be {expected_type} evidence"
                        )
    if len(assessment_ids) != len(set(assessment_ids)):
        errors.append("lineage_review chain assessment claim ids must be unique")
    if claim_by_id:
        required_chain = _claim_dependency_closure(reviewed_ids, claim_by_id)
        missing_chain = sorted(required_chain - set(assessment_ids))
        if missing_chain:
            errors.append(
                "lineage_review omitted dependency chain claims: "
                + ", ".join(missing_chain)
            )
        appeared_direct_claims = {
            claim_id
            for claim_id, claim in claim_by_id.items()
            if claim.get("state") == "active"
            and claim.get("decision_use") == "direct"
            and any(
                isinstance(appearance, dict)
                and appearance.get("artifact_sha256") == deliverable_sha256
                for appearance in claim.get("appearances", [])
            )
        }
        omitted_direct_claims = sorted(appeared_direct_claims - reviewed_ids)
        if omitted_direct_claims:
            errors.append(
                "lineage_review omitted active direct claims appearing in the deliverable: "
                + ", ".join(omitted_direct_claims)
            )

    untracked_ids: list[str] = []
    for index, assessment in enumerate(value.get("untracked_material_claims", [])):
        subject = f"lineage_review.untracked_material_claims[{index}]"
        errors.extend(
            _validate_chain_item(
                assessment,
                subject=subject,
                id_field="id",
                known_evidence_ids=allowed_basis_ids,
            )
        )
        if isinstance(assessment, dict):
            untracked_ids.append(str(assessment.get("id", "")))
            if not assessment.get("deliverable_locations"):
                errors.append(
                    f"{subject}: untracked material claim requires a deliverable location"
                )
    if len(untracked_ids) != len(set(untracked_ids)):
        errors.append("lineage_review untracked claim ids must be unique")
    return errors


def validate_review_record(
    payload: Any,
    contract: dict[str, Any],
    *,
    provenance_mode: str | None = None,
    claim_register: dict[str, Any] | None = None,
    evidence_register: dict[str, Any] | None = None,
    coverage_inventory: dict[str, Any] | None = None,
    matched_support_source_ids: set[str] | None = None,
) -> list[str]:
    """Return mechanical shape and internal-consistency errors for a review."""

    if not isinstance(payload, dict):
        return ["review record must be an object"]
    try:
        review_schema = _load_json(REVIEW_SCHEMA_PATH)
        Draft202012Validator.check_schema(review_schema)
        errors = [
            "review schema "
            + (".".join(str(part) for part in error.absolute_path) or "$")
            + f": {error.message}"
            for error in sorted(
                Draft202012Validator(review_schema).iter_errors(payload),
                key=lambda item: tuple(str(part) for part in item.absolute_path),
            )
        ]
    except (OSError, AdvisoryValidationError) as exc:
        errors = [f"cannot load advisory review schema: {exc}"]
    required = {
        "schema_version",
        "language",
        "advisory_contract_sha256",
        "deliverable_sha256",
        "coverage_inventory_sha256",
        "lineage_inventory_sha256",
        "coverage_review",
        "lineage_review",
        "dimension_reviews",
        "findings",
        "format_specific_checks",
        "correction",
        "approvals",
        "overall_assessment",
        "delivery_readiness",
    }
    missing = sorted(required - payload.keys())
    if missing:
        return [f"review record missing fields: {', '.join(missing)}"]
    if payload["schema_version"] != "1.3":
        errors.append('review schema_version must be "1.3"')
    if payload["language"] != contract["output_language"]:
        errors.append("review language must match the advisory contract")
    for field in (
        "advisory_contract_sha256",
        "deliverable_sha256",
        "coverage_inventory_sha256",
        "lineage_inventory_sha256",
    ):
        value = payload[field]
        if not isinstance(value, str) or re.fullmatch(r"[0-9a-f]{64}", value) is None:
            errors.append(f"{field} must be a lowercase SHA-256")

    coverage = payload["coverage_review"]
    if not isinstance(coverage, dict):
        errors.append("coverage_review must be an object")
    else:
        if coverage.get("selection_method") != "model_led_materiality_review":
            errors.append(
                "coverage_review.selection_method must be model_led_materiality_review"
            )
        if coverage.get("scope") != contract["validation_scope"]["coverage"]:
            errors.append(
                "coverage_review.scope must match advisory_contract.validation_scope"
            )
        if not _is_string_list(coverage.get("reviewed_sections"), allow_empty=False):
            errors.append("coverage_review.reviewed_sections must be non-empty")
        for field in ("omitted_sections", "limitations"):
            if not _is_string_list(coverage.get(field), allow_empty=True):
                errors.append(f"coverage_review.{field} must be a string array")
        if not _is_non_empty_string(coverage.get("analysis")):
            errors.append("coverage_review.analysis is required")
        considered_unit_ids = coverage.get("considered_unit_ids")
        omitted_unit_ids = coverage.get("omitted_unit_ids")
        if not _is_string_list(considered_unit_ids, allow_empty=False):
            errors.append("coverage_review.considered_unit_ids must be non-empty")
        if not _is_string_list(omitted_unit_ids, allow_empty=True):
            errors.append("coverage_review.omitted_unit_ids must be a string array")
        unit_assessments = coverage.get("unit_assessments")
        assessed_unit_ids: list[str] = []
        selected_claim_ids: set[str] = set()
        selected_untracked_ids: set[str] = set()
        if not isinstance(unit_assessments, list) or not unit_assessments:
            errors.append("coverage_review.unit_assessments must be non-empty")
            unit_assessments = []
        for index, assessment in enumerate(unit_assessments):
            subject = f"coverage_review.unit_assessments[{index}]"
            if not isinstance(assessment, dict):
                errors.append(f"{subject} must be an object")
                continue
            unit_id = str(assessment.get("unit_id", ""))
            assessed_unit_ids.append(unit_id)
            status = assessment.get("status")
            material_ids = assessment.get("material_claim_ids")
            untracked_ids = assessment.get("untracked_claim_ids")
            if status not in {
                "reviewed_material_claims",
                "reviewed_no_material_claims",
                "omitted",
            }:
                errors.append(f"{subject}.status is invalid")
            if not _is_string_list(material_ids, allow_empty=True):
                errors.append(f"{subject}.material_claim_ids must be a string array")
                material_ids = []
            if not _is_string_list(untracked_ids, allow_empty=True):
                errors.append(f"{subject}.untracked_claim_ids must be a string array")
                untracked_ids = []
            if not _is_non_empty_string(assessment.get("analysis")):
                errors.append(f"{subject}.analysis is required")
            if status == "reviewed_material_claims" and not (
                material_ids or untracked_ids
            ):
                errors.append(
                    f"{subject}: reviewed_material_claims requires a tracked or untracked claim id"
                )
            if status in {"reviewed_no_material_claims", "omitted"} and (
                material_ids or untracked_ids
            ):
                errors.append(f"{subject}: {status} cannot carry claim ids")
            selected_claim_ids.update(str(value) for value in material_ids)
            selected_untracked_ids.update(str(value) for value in untracked_ids)
        if len(assessed_unit_ids) != len(set(assessed_unit_ids)):
            errors.append("coverage review unit assessment ids must be unique")
        considered = set(considered_unit_ids or [])
        omitted = set(omitted_unit_ids or [])
        for assessment in unit_assessments:
            if not isinstance(assessment, dict):
                continue
            unit_id = str(assessment.get("unit_id", ""))
            status = assessment.get("status")
            if status == "omitted" and unit_id not in omitted:
                errors.append(
                    f"coverage unit {unit_id} is assessed omitted but is not in omitted_unit_ids"
                )
            if status != "omitted" and unit_id not in considered:
                errors.append(
                    f"coverage unit {unit_id} is reviewed but is not in considered_unit_ids"
                )
        lineage_review_for_coverage = payload.get("lineage_review")
        if isinstance(lineage_review_for_coverage, dict):
            declared_selected = set(
                str(value)
                for value in lineage_review_for_coverage.get("reviewed_claim_ids", [])
            )
            declared_untracked = {
                str(item.get("id"))
                for item in lineage_review_for_coverage.get(
                    "untracked_material_claims", []
                )
                if isinstance(item, dict) and item.get("id")
            }
            if selected_claim_ids != declared_selected:
                errors.append(
                    "coverage unit material_claim_ids must exactly match lineage_review.reviewed_claim_ids"
                )
            if selected_untracked_ids != declared_untracked:
                errors.append(
                    "coverage unit untracked_claim_ids must exactly match lineage_review.untracked_material_claims"
                )
        if isinstance(coverage_inventory, dict):
            known_units = {
                str(unit.get("id"))
                for unit in coverage_inventory.get("units", [])
                if isinstance(unit, dict) and unit.get("id")
            }
            considered = set(considered_unit_ids or [])
            omitted = set(omitted_unit_ids or [])
            if considered & omitted:
                errors.append("coverage units cannot be both considered and omitted")
            unknown_units = sorted((considered | omitted) - known_units)
            if unknown_units:
                errors.append(
                    "coverage review references unknown units: "
                    + ", ".join(unknown_units)
                )
            missing_units = sorted(known_units - considered - omitted)
            if missing_units:
                errors.append(
                    "coverage review does not account for units: "
                    + ", ".join(missing_units)
                )
            if set(assessed_unit_ids) != known_units:
                errors.append(
                    "coverage_review.unit_assessments must contain exactly one assessment for every coverage unit"
                )

    errors.extend(
        _validate_lineage_review(
            payload["lineage_review"],
            provenance_mode=provenance_mode,
            claim_register=claim_register,
            evidence_register=evidence_register,
            deliverable_sha256=str(payload.get("deliverable_sha256", "")),
            matched_support_source_ids=matched_support_source_ids,
        )
    )

    dimensions = payload["dimension_reviews"]
    if not isinstance(dimensions, dict) or set(dimensions) != set(REVIEW_DIMENSIONS):
        errors.append(
            "dimension_reviews must contain exactly the ten required dimensions"
        )
    else:
        for dimension in REVIEW_DIMENSIONS:
            errors.extend(
                _validate_dimension_review(
                    dimensions[dimension], subject=f"dimension_reviews.{dimension}"
                )
            )

    findings = payload["findings"]
    finding_ids: list[str] = []
    if not isinstance(findings, list):
        errors.append("findings must be an array")
    else:
        for index, finding in enumerate(findings):
            subject = f"findings[{index}]"
            if not isinstance(finding, dict):
                errors.append(f"{subject} must be an object")
                continue
            if not _is_non_empty_string(finding.get("id")):
                errors.append(f"{subject}.id is required")
            else:
                finding_ids.append(finding["id"])
            if finding.get("dimension") not in REVIEW_DIMENSIONS:
                errors.append(f"{subject}.dimension is invalid")
            if not _is_non_empty_string(finding.get("finding")):
                errors.append(f"{subject}.finding is required")
            if finding.get("status") not in DIMENSION_STATUSES:
                errors.append(f"{subject}.status is invalid")
            if not _is_string_list(finding.get("evidence_refs"), allow_empty=True):
                errors.append(f"{subject}.evidence_refs must be a string array")
            if not isinstance(finding.get("correction_action"), str):
                errors.append(f"{subject}.correction_action must be a string")
            if finding.get("correction_status") not in CORRECTION_STATUSES:
                errors.append(f"{subject}.correction_status is invalid")
            if not isinstance(finding.get("professional_review_required"), bool):
                errors.append(f"{subject}.professional_review_required must be boolean")
        if len(finding_ids) != len(set(finding_ids)):
            errors.append("finding ids must be unique")

    checks = payload["format_specific_checks"]
    check_by_workflow: dict[str, dict[str, Any]] = {}
    if not isinstance(checks, list):
        errors.append("format_specific_checks must be an array")
    else:
        for index, check in enumerate(checks):
            subject = f"format_specific_checks[{index}]"
            if not isinstance(check, dict):
                errors.append(f"{subject} must be an object")
                continue
            workflow = check.get("workflow")
            if not _is_non_empty_string(workflow):
                errors.append(f"{subject}.workflow is required")
            elif workflow in check_by_workflow:
                errors.append(f"duplicate format-specific check: {workflow}")
            else:
                check_by_workflow[workflow] = check
            if check.get("status") not in FORMAT_CHECK_STATUSES:
                errors.append(f"{subject}.status is invalid")
            if not _is_string_list(check.get("artifact_refs"), allow_empty=True):
                errors.append(f"{subject}.artifact_refs must be a string array")
            if not _is_non_empty_string(check.get("analysis")):
                errors.append(f"{subject}.analysis is required")
            if check.get("status") == "passed" and not check.get("artifact_refs"):
                errors.append(f"{subject} passed without an artifact reference")

    for declared in contract["validation_profile"]["format_checks"]:
        if declared["requirement"] != "required":
            continue
        workflow = declared["workflow"]
        actual = check_by_workflow.get(workflow)
        if actual is None:
            errors.append(f"required format check is missing: {workflow}")
        elif actual.get("status") != "passed":
            errors.append(f"required format check did not pass: {workflow}")

    correction = payload["correction"]
    if not isinstance(correction, dict):
        errors.append("correction must be an object")
    else:
        if correction.get("status") not in {
            "not_required",
            "required",
            "completed",
            "blocked",
        }:
            errors.append("correction.status is invalid")
        if not _is_non_empty_string(correction.get("summary")):
            errors.append("correction.summary is required")
        if not isinstance(correction.get("corrected_artifact"), str):
            errors.append("correction.corrected_artifact must be a string")
        corrected_hash = correction.get("corrected_artifact_sha256")
        if not isinstance(corrected_hash, str):
            errors.append("correction.corrected_artifact_sha256 must be a string")
        corrected_inventory_hash = correction.get("corrected_inventory_sha256")
        corrected_review_hash = correction.get("corrected_review_sha256")
        if not isinstance(corrected_inventory_hash, str):
            errors.append("correction.corrected_inventory_sha256 must be a string")
        if not isinstance(corrected_review_hash, str):
            errors.append("correction.corrected_review_sha256 must be a string")
        if not _is_string_list(correction.get("unresolved_changes"), allow_empty=True):
            errors.append("correction.unresolved_changes must be a string array")
        if (
            correction.get("status") == "completed"
            and not contract["correction_policy"]["allowed"]
        ):
            errors.append(
                "correction completed even though the contract disallows correction"
            )
        if correction.get("status") == "completed":
            if not _is_non_empty_string(correction.get("corrected_artifact")):
                errors.append(
                    "correction.corrected_artifact is required when completed"
                )
            if (
                not isinstance(corrected_hash, str)
                or re.fullmatch(r"[0-9a-f]{64}", corrected_hash) is None
            ):
                errors.append(
                    "correction.corrected_artifact_sha256 must be a lowercase SHA-256 when completed"
                )
            for field, value in (
                ("corrected_inventory_sha256", corrected_inventory_hash),
                ("corrected_review_sha256", corrected_review_hash),
            ):
                if (
                    not isinstance(value, str)
                    or re.fullmatch(r"[0-9a-f]{64}", value) is None
                ):
                    errors.append(
                        f"correction.{field} must be a lowercase SHA-256 when completed"
                    )
        elif (
            correction.get("corrected_artifact")
            or corrected_hash
            or corrected_inventory_hash
            or corrected_review_hash
        ):
            errors.append(
                "correction artifact path and review hashes require completed status"
            )

    approvals = payload["approvals"]
    if not isinstance(approvals, dict):
        errors.append("approvals must be an object")
    else:
        approval_names = {"professional_judgement", "correction"}
        if set(approvals) != approval_names:
            errors.append(
                "approvals must contain professional_judgement and correction"
            )
        else:
            for approval_name in sorted(approval_names):
                errors.extend(
                    _validate_approval_record(
                        approvals[approval_name],
                        subject=f"approvals.{approval_name}",
                    )
                )

    overall = payload["overall_assessment"]
    readiness = payload["delivery_readiness"]
    lineage_review = payload["lineage_review"]
    if not isinstance(overall, dict):
        errors.append("overall_assessment must be an object")
    if not isinstance(readiness, dict):
        errors.append("delivery_readiness must be an object")
    if isinstance(overall, dict) and isinstance(readiness, dict):
        outcome = overall.get("outcome")
        status = readiness.get("status")
        if outcome not in READINESS_STATUSES:
            errors.append("overall_assessment.outcome is invalid")
        if status not in READINESS_STATUSES:
            errors.append("delivery_readiness.status is invalid")
        if outcome != status:
            errors.append("overall assessment and delivery readiness must match")
        if not _is_non_empty_string(overall.get("analysis")):
            errors.append("overall_assessment.analysis is required")
        for field in ("residual_uncertainties", "professional_review_items"):
            if not _is_string_list(overall.get(field), allow_empty=True):
                errors.append(f"overall_assessment.{field} must be a string array")
        if not _is_string_list(readiness.get("conditions"), allow_empty=True):
            errors.append("delivery_readiness.conditions must be a string array")

        attention_statuses = {
            "partially_conforms",
            "does_not_conform",
            "contradicted",
            "uncertain",
            "judgment_required",
        }
        unresolved_attention = isinstance(dimensions, dict) and any(
            isinstance(dimensions.get(dimension), dict)
            and (
                dimensions[dimension].get("status") in attention_statuses
                or dimensions[dimension].get("correction_status")
                in {"proposed", "blocked", "professional_review_required"}
                or dimensions[dimension].get("professional_review_required") is True
            )
            for dimension in REVIEW_DIMENSIONS
        )
        finding_review_required = isinstance(findings, list) and any(
            isinstance(finding, dict)
            and finding.get("professional_review_required") is True
            for finding in findings
        )
        dimension_review_required = isinstance(dimensions, dict) and any(
            isinstance(dimensions.get(dimension), dict)
            and (
                dimensions[dimension].get("professional_review_required") is True
                or dimensions[dimension].get("correction_status")
                == "professional_review_required"
            )
            for dimension in REVIEW_DIMENSIONS
        )
        delivery_ready = status in {"ready", "ready_with_residual_uncertainty"}
        chain_items: list[dict[str, Any]] = []
        if isinstance(lineage_review, dict):
            for field in ("chain_assessments", "untracked_material_claims"):
                values = lineage_review.get(field, [])
                if isinstance(values, list):
                    chain_items.extend(
                        item for item in values if isinstance(item, dict)
                    )
        unresolved_chain = any(
            item.get("support_status") in {"unsupported", "contradicted", "uncertain"}
            or item.get("reasoning_status") in {"gap", "contradicted", "uncertain"}
            or (
                isinstance(item.get("recheck"), dict)
                and item["recheck"].get("status") in {"pending", "blocked"}
            )
            or (
                isinstance(item.get("resolution"), dict)
                and item["resolution"].get("status")
                in {"pending", "professional_review_required"}
            )
            for item in chain_items
        )
        if delivery_ready and not chain_items:
            errors.append(
                "delivery-ready status requires at least one model-reviewed material claim"
            )
        if delivery_ready and unresolved_chain:
            errors.append(
                "delivery-ready status cannot coexist with an unresolved claim chain"
            )
        if status == "ready" and unresolved_attention:
            errors.append(
                "ready status cannot coexist with unresolved review attention"
            )
        if status == "ready" and (
            overall.get("residual_uncertainties")
            or overall.get("professional_review_items")
        ):
            errors.append(
                "ready status cannot carry residual uncertainty or professional-review items"
            )
        if (
            isinstance(correction, dict)
            and correction.get("status")
            in {
                "required",
                "blocked",
            }
            and status in {"ready", "ready_with_residual_uncertainty"}
        ):
            errors.append("unresolved correction cannot be delivery-ready")
        if delivery_ready and finding_review_required:
            errors.append(
                "delivery-ready status cannot coexist with a finding requiring professional review"
            )
        if delivery_ready and dimension_review_required:
            errors.append(
                "delivery-ready status cannot coexist with a dimension requiring professional review"
            )
        if delivery_ready and isinstance(approvals, dict):
            professional_approval = approvals.get("professional_judgement", {})
            correction_approval = approvals.get("correction", {})
            if not isinstance(professional_approval, dict):
                professional_approval = {}
            if not isinstance(correction_approval, dict):
                correction_approval = {}
            if (
                contract["professional_judgement_policy"][
                    "approval_required_before_delivery"
                ]
                and professional_approval.get("status") != "approved"
            ):
                errors.append(
                    "professional-judgement approval is required before delivery"
                )
            if (
                isinstance(correction, dict)
                and correction.get("status") == "completed"
                and contract["correction_policy"]["approval_required_before_delivery"]
                and correction_approval.get("status") != "approved"
            ):
                errors.append("correction approval is required before delivery")
    return errors


def _render_package(
    inventory: dict[str, Any],
    review: dict[str, Any],
    audit: dict[str, Any],
) -> str:
    dimensions = review.get("dimension_reviews", {})
    if not isinstance(dimensions, dict):
        dimensions = {}
    lines = [
        "# Advisory deliverable validation",
        "",
        f"Original: {inventory.get('source_name', 'unknown')}",
        f"Original SHA-256: {inventory.get('source_sha256', 'unknown')}",
        f"Contract SHA-256: {inventory.get('advisory_contract_sha256', 'unknown')}",
        f"Record complete: {'yes' if audit['record_complete'] else 'no'}",
        f"Delivery readiness: {audit['effective_delivery_readiness']}",
        "",
        "## Review dimensions",
        "",
    ]
    for dimension in REVIEW_DIMENSIONS:
        record = dimensions.get(dimension, {})
        if not isinstance(record, dict):
            record = {}
        lines.append(
            f"- **{dimension.replace('_', ' ').title()}** — {record.get('status', 'missing')}: {record.get('analysis', 'No analysis recorded.')}"
        )
    lines.extend(["", "## Findings", ""])
    findings = review.get("findings", [])
    if not isinstance(findings, list):
        findings = []
    if findings:
        for finding in findings:
            if not isinstance(finding, dict):
                continue
            lines.append(
                f"- **{finding.get('id', 'finding')}** ({finding.get('dimension', 'unknown')}): {finding.get('finding', '')}"
            )
    else:
        lines.append("- No individual findings recorded.")
    lines.extend(["", "## Claim and evidence chains", ""])
    lineage_review = review.get("lineage_review", {})
    if not isinstance(lineage_review, dict):
        lineage_review = {}
    lines.append(f"Provenance mode: {lineage_review.get('provenance_mode', 'missing')}")
    chain_items: list[tuple[str, dict[str, Any]]] = []
    for field, label in (
        ("chain_assessments", "Tracked"),
        ("untracked_material_claims", "Untracked"),
    ):
        values = lineage_review.get(field, [])
        if isinstance(values, list):
            chain_items.extend(
                (label, item) for item in values if isinstance(item, dict)
            )
    if not chain_items:
        lines.append("- No claim-chain assessment recorded.")
    for label, item in chain_items:
        item_id = item.get("claim_id") or item.get("id") or "claim"
        recheck = item.get("recheck")
        if not isinstance(recheck, dict):
            recheck = {}
        resolution = item.get("resolution")
        if not isinstance(resolution, dict):
            resolution = {}
        lines.append(
            f"- **{label} {item_id}** — support: {item.get('support_status', 'missing')}; "
            f"reasoning: {item.get('reasoning_status', 'missing')}; "
            f"recheck: {recheck.get('status', 'missing')}; "
            f"resolution: {resolution.get('status', 'missing')}. "
            f"{item.get('statement', '')}"
        )
    lines.extend(["", "## Format-specific checks", ""])
    checks = review.get("format_specific_checks", [])
    if not isinstance(checks, list):
        checks = []
    if checks:
        for check in checks:
            if not isinstance(check, dict):
                continue
            lines.append(
                f"- **{check.get('workflow', 'unknown')}** — {check.get('status', 'missing')}: {check.get('analysis', '')}"
            )
    else:
        lines.append("- No format-specific checks recorded.")
    correction = review.get("correction", {})
    if not isinstance(correction, dict):
        correction = {}
    lines.extend(
        [
            "",
            "## Correction",
            "",
            f"Status: {correction.get('status', 'missing')}",
            "",
            str(correction.get("summary", "No correction summary recorded.")),
            "",
            "## Residual uncertainty and professional review",
            "",
        ]
    )
    overall = review.get("overall_assessment", {})
    if not isinstance(overall, dict):
        overall = {}
    residual_uncertainties = overall.get("residual_uncertainties", [])
    if not isinstance(residual_uncertainties, list):
        residual_uncertainties = []
    professional_review_items = overall.get("professional_review_items", [])
    if not isinstance(professional_review_items, list):
        professional_review_items = []
    for item in residual_uncertainties:
        lines.append(f"- Residual uncertainty: {item}")
    for item in professional_review_items:
        lines.append(f"- Professional review: {item}")
    if not residual_uncertainties and not professional_review_items:
        lines.append("- None recorded.")
    if audit["errors"]:
        lines.extend(["", "## Mechanical audit errors", ""])
        lines.extend(f"- {error}" for error in audit["errors"])
    return "\n".join(lines).rstrip() + "\n"


def _load_preparation_context(
    inventory: dict[str, Any],
) -> tuple[
    dict[str, Any],
    dict[str, Any],
    dict[str, Any] | None,
    dict[str, Any] | None,
    list[str],
]:
    """Load hash-bound coverage and lineage artifacts recorded at preparation."""

    errors: list[str] = []
    coverage_inventory: dict[str, Any] = {}
    lineage_inventory: dict[str, Any] = {}
    claim_register: dict[str, Any] | None = None
    evidence_register: dict[str, Any] | None = None

    coverage_path = Path(str(inventory.get("coverage_inventory_path", "")))
    if not coverage_path.is_file():
        errors.append("coverage inventory is missing")
    else:
        if _sha256(coverage_path) != inventory.get("coverage_inventory_sha256"):
            errors.append("coverage inventory changed since preparation")
        loaded = _load_json(coverage_path)
        if isinstance(loaded, dict):
            coverage_inventory = loaded
        else:
            errors.append("coverage inventory must be an object")

    lineage_record = inventory.get("lineage")
    if not isinstance(lineage_record, dict):
        errors.append("deliverable inventory has no lineage record")
        return (
            coverage_inventory,
            lineage_inventory,
            claim_register,
            evidence_register,
            errors,
        )
    lineage_path = Path(str(lineage_record.get("lineage_inventory_path", "")))
    if not lineage_path.is_file():
        errors.append("lineage inventory is missing")
    else:
        if _sha256(lineage_path) != lineage_record.get("lineage_inventory_sha256"):
            errors.append("lineage inventory changed since preparation")
        loaded = _load_json(lineage_path)
        if isinstance(loaded, dict):
            lineage_inventory = loaded
        else:
            errors.append("lineage inventory must be an object")

    if lineage_record.get("provenance_mode") == "generation_time" and lineage_inventory:
        for label, filename in (
            ("evidence_register", EVIDENCE_REGISTER_FILENAME),
            ("claim_register", CLAIM_REGISTER_FILENAME),
        ):
            record = lineage_inventory.get(label)
            if not isinstance(record, dict):
                errors.append(f"lineage inventory has no {label}")
                continue
            copied_path = Path(str(record.get("copied_path", "")))
            source_path = Path(str(record.get("source_path", "")))
            expected_hash = record.get("sha256")
            if copied_path.name != filename or not copied_path.is_file():
                errors.append(f"prepared {label} is missing")
                continue
            if _sha256(copied_path) != expected_hash:
                errors.append(f"prepared {label} changed since preparation")
            if not source_path.is_file() or _sha256(source_path) != expected_hash:
                errors.append(f"source {label} changed since preparation")
            loaded = _load_json(copied_path)
            if not isinstance(loaded, dict):
                errors.append(f"prepared {label} must be an object")
            elif label == "claim_register":
                claim_register = loaded
            else:
                evidence_register = loaded
    return (
        coverage_inventory,
        lineage_inventory,
        claim_register,
        evidence_register,
        errors,
    )


def _load_source_inventory(
    inventory: dict[str, Any],
) -> tuple[dict[str, Any], list[str]]:
    errors: list[str] = []
    path = Path(str(inventory.get("source_inventory_path", "")))
    expected_hash = inventory.get("source_inventory_sha256")
    if not path.is_file():
        return {}, ["source inventory is missing"]
    if _sha256(path) != expected_hash:
        errors.append("source inventory changed since preparation")
    payload = _load_json(path)
    if not isinstance(payload, dict):
        errors.append("source inventory must be an object")
        return {}, errors
    return payload, errors


def _recheck_tasks(review: dict[str, Any]) -> dict[str, Any]:
    """Package model-selected rechecks without selecting or executing them."""

    tasks: list[dict[str, Any]] = []
    lineage_review = review.get("lineage_review")
    if isinstance(lineage_review, dict):
        for field, id_field in (
            ("chain_assessments", "claim_id"),
            ("untracked_material_claims", "id"),
        ):
            values = lineage_review.get(field, [])
            if not isinstance(values, list):
                continue
            for item in values:
                if not isinstance(item, dict):
                    continue
                recheck = item.get("recheck")
                if not isinstance(recheck, dict) or recheck.get("required") is not True:
                    continue
                tasks.append(
                    {
                        "claim_review_id": str(item.get(id_field, "")),
                        "statement": str(item.get("statement", "")),
                        "kind": recheck.get("kind"),
                        "status": recheck.get("status"),
                        "evidence_ids": recheck.get("evidence_ids", []),
                        "analysis": recheck.get("analysis", ""),
                    }
                )
    return {
        "schema_version": "1.0",
        "purpose": (
            "Model-selected final rechecks. This file does not choose sources, "
            "browse, rerun calculations, or decide whether evidence supports a claim."
        ),
        "tasks": tasks,
    }


def _audit_corrected_validation(
    *,
    corrected_deliverable: Path,
    corrected_inventory_path: Path | None,
    corrected_review_path: Path | None,
    contract: dict[str, Any],
    contract_path: Path,
    correction_record: dict[str, Any],
) -> tuple[dict[str, Any] | None, list[str]]:
    """Require a complete second review bound to the corrected bytes."""

    errors: list[str] = []
    if corrected_inventory_path is None or not corrected_inventory_path.is_file():
        errors.append(
            "completed correction requires --corrected-deliverable-inventory from a second prepare run"
        )
    if corrected_review_path is None or not corrected_review_path.is_file():
        errors.append(
            "completed correction requires --corrected-review from a second model-led review"
        )
    if errors:
        return None, errors
    assert corrected_inventory_path is not None
    assert corrected_review_path is not None
    if correction_record.get("corrected_inventory_sha256") != _sha256(
        corrected_inventory_path
    ):
        errors.append("corrected inventory does not match the hash bound in correction")
    if correction_record.get("corrected_review_sha256") != _sha256(
        corrected_review_path
    ):
        errors.append("corrected review does not match the hash bound in correction")
    inventory = _load_json(corrected_inventory_path)
    review = _load_json(corrected_review_path)
    if not isinstance(inventory, dict):
        return None, [*errors, "corrected deliverable inventory must be an object"]
    if not isinstance(review, dict):
        return None, [*errors, "corrected review must be an object"]
    corrected_hash = _sha256(corrected_deliverable)
    if (
        Path(str(inventory.get("source_path", ""))).resolve()
        != corrected_deliverable.resolve()
    ):
        errors.append(
            "corrected inventory is not for the supplied corrected deliverable"
        )
    if inventory.get("source_sha256") != corrected_hash:
        errors.append("corrected inventory is not bound to the corrected bytes")
    if inventory.get("advisory_contract_sha256") != _sha256(contract_path):
        errors.append("corrected inventory is not bound to the current contract")
    (
        coverage_inventory,
        lineage_inventory,
        claim_register,
        evidence_register,
        preparation_errors,
    ) = _load_preparation_context(inventory)
    source_inventory, source_inventory_errors = _load_source_inventory(inventory)
    errors.extend(f"corrected review: {error}" for error in preparation_errors)
    errors.extend(f"corrected review: {error}" for error in source_inventory_errors)
    matched_support_source_ids = {
        str(item.get("id"))
        for item in source_inventory.get("sources", [])
        if isinstance(item, dict) and item.get("id")
    }
    provenance_mode = (
        str(lineage_inventory.get("provenance_mode"))
        if isinstance(lineage_inventory, dict)
        else None
    )
    errors.extend(
        "corrected review: " + error
        for error in validate_review_record(
            review,
            contract,
            provenance_mode=provenance_mode,
            claim_register=claim_register,
            evidence_register=evidence_register,
            coverage_inventory=coverage_inventory,
            matched_support_source_ids=matched_support_source_ids,
        )
    )
    if review.get("deliverable_sha256") != corrected_hash:
        errors.append("corrected review is not bound to the corrected bytes")
    contract_hash = _sha256(contract_path)
    if review.get("advisory_contract_sha256") != contract_hash:
        errors.append("corrected review is not bound to the current contract")
    if review.get("coverage_inventory_sha256") != inventory.get(
        "coverage_inventory_sha256"
    ):
        errors.append("corrected review is not bound to its coverage inventory")
    inventory_lineage = inventory.get("lineage")
    expected_lineage_hash = (
        inventory_lineage.get("lineage_inventory_sha256")
        if isinstance(inventory_lineage, dict)
        else None
    )
    if review.get("lineage_inventory_sha256") != expected_lineage_hash:
        errors.append("corrected review is not bound to its lineage inventory")
    _, corrected_format_errors = _audit_format_check_artifacts(
        review,
        contract,
        contract_path.parent,
        deliverable_sha256=corrected_hash,
    )
    errors.extend(f"corrected review: {error}" for error in corrected_format_errors)
    corrected_correction = review.get("correction")
    if (
        not isinstance(corrected_correction, dict)
        or corrected_correction.get("status") != "not_required"
    ):
        errors.append("corrected review must conclude correction.status not_required")
    corrected_readiness = review.get("delivery_readiness")
    if not isinstance(corrected_readiness, dict) or corrected_readiness.get(
        "status"
    ) not in {"ready", "ready_with_residual_uncertainty"}:
        errors.append("corrected review is not delivery-ready")
    metadata = {
        "inventory_path": str(corrected_inventory_path.resolve()),
        "inventory_sha256": _sha256(corrected_inventory_path),
        "review_path": str(corrected_review_path.resolve()),
        "review_sha256": _sha256(corrected_review_path),
        "delivery_readiness": (
            corrected_readiness.get("status")
            if isinstance(corrected_readiness, dict)
            else "blocked"
        ),
    }
    return metadata, errors


def package_validation(
    inventory_path: Path,
    review_draft_path: Path,
    advisory_contract_path: Path,
    output_dir: Path,
    *,
    corrected_deliverable: Path | None = None,
    corrected_deliverable_inventory: Path | None = None,
    corrected_review: Path | None = None,
) -> tuple[dict[str, Path], dict[str, Any]]:
    """Audit a model-authored review and write the validation package."""

    _ensure_output_location(output_dir)
    inventory = _load_json(inventory_path)
    review = _load_json(review_draft_path)
    contract = _load_json(advisory_contract_path)
    contract_errors = validate_advisory_contract(contract)
    errors = list(contract_errors)
    if not isinstance(inventory, dict):
        errors.append("deliverable inventory must be an object")
        inventory = {}
    (
        coverage_inventory,
        lineage_inventory,
        claim_register,
        evidence_register,
        preparation_errors,
    ) = _load_preparation_context(inventory)
    errors.extend(preparation_errors)
    source_inventory, source_inventory_errors = _load_source_inventory(inventory)
    errors.extend(source_inventory_errors)
    matched_support_source_ids = {
        str(item.get("id"))
        for item in source_inventory.get("sources", [])
        if isinstance(item, dict) and item.get("id")
    }
    provenance_mode = None
    if isinstance(lineage_inventory, dict):
        candidate_mode = lineage_inventory.get("provenance_mode")
        if isinstance(candidate_mode, str):
            provenance_mode = candidate_mode
    if isinstance(review, dict) and isinstance(contract, dict) and not contract_errors:
        errors.extend(
            validate_review_record(
                review,
                contract,
                provenance_mode=provenance_mode,
                claim_register=claim_register,
                evidence_register=evidence_register,
                coverage_inventory=coverage_inventory,
                matched_support_source_ids=matched_support_source_ids,
            )
        )
    elif not isinstance(review, dict):
        errors.append("review record must be an object")
        review = {}

    original_path = Path(str(inventory.get("source_path", "")))
    expected_original_hash = inventory.get("source_sha256")
    original_unchanged = (
        original_path.is_file() and _sha256(original_path) == expected_original_hash
    )
    if not original_unchanged:
        errors.append("original deliverable is missing or changed since preparation")
    contract_hash = _sha256(advisory_contract_path)
    if inventory.get("advisory_contract_sha256") != contract_hash:
        errors.append("advisory contract no longer matches the preparation inventory")
    if review.get("advisory_contract_sha256") != contract_hash:
        errors.append("review is not bound to the current advisory contract")
    if review.get("deliverable_sha256") != expected_original_hash:
        errors.append("review is not bound to the prepared deliverable")
    if review.get("coverage_inventory_sha256") != inventory.get(
        "coverage_inventory_sha256"
    ):
        errors.append("review is not bound to the prepared coverage inventory")
    inventory_lineage = inventory.get("lineage")
    expected_lineage_hash = (
        inventory_lineage.get("lineage_inventory_sha256")
        if isinstance(inventory_lineage, dict)
        else None
    )
    if review.get("lineage_inventory_sha256") != expected_lineage_hash:
        errors.append("review is not bound to the prepared lineage inventory")

    format_check_artifacts: list[dict[str, Any]] = []
    artifact_errors: list[str] = []
    if isinstance(review, dict) and isinstance(contract, dict) and not contract_errors:
        format_check_artifacts, artifact_errors = _audit_format_check_artifacts(
            review,
            contract,
            advisory_contract_path.parent,
            deliverable_sha256=str(expected_original_hash),
        )
        errors.extend(artifact_errors)

    corrected_metadata: dict[str, Any] | None = None
    corrected_validation_metadata: dict[str, Any] | None = None
    correction_record = review.get("correction", {})
    if not isinstance(correction_record, dict):
        correction_record = {}
    correction_status = correction_record.get("status")
    if correction_status == "completed":
        if corrected_deliverable is None or not corrected_deliverable.is_file():
            errors.append("completed correction requires --corrected-deliverable")
        elif corrected_deliverable.resolve() == original_path.resolve():
            errors.append("corrected deliverable must use a separate path")
        else:
            corrected_hash = _sha256(corrected_deliverable)
            declared_corrected_path = _resolve_artifact_ref(
                str(correction_record.get("corrected_artifact", "")),
                advisory_contract_path.parent,
            )
            if corrected_hash == expected_original_hash:
                errors.append("completed correction has the same bytes as the original")
            if declared_corrected_path != corrected_deliverable.resolve():
                errors.append(
                    "corrected deliverable does not match the path bound in the review"
                )
            if correction_record.get("corrected_artifact_sha256") != corrected_hash:
                errors.append(
                    "corrected deliverable does not match the SHA-256 bound in the review"
                )
            corrected_metadata = {
                "path": str(corrected_deliverable.resolve()),
                "name": corrected_deliverable.name,
                "sha256": corrected_hash,
                "byte_count": corrected_deliverable.stat().st_size,
            }
            (
                candidate_corrected_validation,
                corrected_validation_errors,
            ) = _audit_corrected_validation(
                corrected_deliverable=corrected_deliverable,
                corrected_inventory_path=corrected_deliverable_inventory,
                corrected_review_path=corrected_review,
                contract=contract,
                contract_path=advisory_contract_path,
                correction_record=correction_record,
            )
            errors.extend(corrected_validation_errors)
            if not corrected_validation_errors:
                corrected_validation_metadata = candidate_corrected_validation
    elif corrected_deliverable is not None:
        errors.append(
            "a corrected deliverable was supplied but correction.status is not completed"
        )
    elif corrected_deliverable_inventory is not None or corrected_review is not None:
        errors.append(
            "corrected inventory/review were supplied but correction.status is not completed"
        )

    paths = {
        "review": output_dir / "advisory_validation_review.json",
        "audit": output_dir / "validation_audit.json",
        "package": output_dir / "advisory_validation_package.md",
        "recheck_tasks": output_dir / "recheck_tasks.json",
    }
    protected_inputs = {
        "deliverable inventory": inventory_path,
        "review draft": review_draft_path,
        "advisory contract": advisory_contract_path,
    }
    if original_path.is_file():
        protected_inputs["original deliverable"] = original_path
    if corrected_deliverable is not None:
        protected_inputs["corrected deliverable"] = corrected_deliverable
    if corrected_deliverable_inventory is not None:
        protected_inputs["corrected deliverable inventory"] = (
            corrected_deliverable_inventory
        )
    if corrected_review is not None:
        protected_inputs["corrected review"] = corrected_review
    for index, artifact in enumerate(format_check_artifacts, start=1):
        protected_inputs[f"format-check artifact {index}"] = Path(artifact["path"])
    _reject_output_collisions(paths, protected_inputs)

    delivery_readiness = review.get("delivery_readiness", {})
    if not isinstance(delivery_readiness, dict):
        delivery_readiness = {}
    declared_readiness = delivery_readiness.get("status", "blocked")
    lineage_review = review.get("lineage_review", {})
    if not isinstance(lineage_review, dict):
        lineage_review = {}
    evidence_lineage = lineage_inventory.get("evidence_register")
    claim_lineage = lineage_inventory.get("claim_register")
    audit = {
        "schema_version": "1.1",
        "record_complete": not errors,
        "errors": errors,
        "deliverable": {
            "path": str(original_path.resolve()),
            "sha256": expected_original_hash,
            "byte_count": (
                original_path.stat().st_size if original_path.is_file() else 0
            ),
        },
        "lineage": {
            "provenance_mode": provenance_mode,
            "evidence_register": (
                {
                    "source_path": evidence_lineage.get("source_path"),
                    "sha256": evidence_lineage.get("sha256"),
                }
                if isinstance(evidence_lineage, dict)
                else None
            ),
            "claim_register": (
                {
                    "source_path": claim_lineage.get("source_path"),
                    "sha256": claim_lineage.get("sha256"),
                }
                if isinstance(claim_lineage, dict)
                else None
            ),
            "reviewed_claim_ids": lineage_review.get("reviewed_claim_ids", []),
            "assessed_claim_ids": [
                str(item.get("claim_id"))
                for item in lineage_review.get("chain_assessments", [])
                if isinstance(item, dict) and item.get("claim_id")
            ],
        },
        "checks": {
            "advisory_contract_shape": not validate_advisory_contract(contract),
            "contract_hash_bound": review.get("advisory_contract_sha256")
            == contract_hash,
            "deliverable_hash_bound": review.get("deliverable_sha256")
            == expected_original_hash,
            "coverage_inventory_hash_bound": review.get("coverage_inventory_sha256")
            == inventory.get("coverage_inventory_sha256"),
            "lineage_inventory_hash_bound": review.get("lineage_inventory_sha256")
            == (
                inventory.get("lineage", {}).get("lineage_inventory_sha256")
                if isinstance(inventory.get("lineage"), dict)
                else None
            ),
            "preparation_artifacts_unchanged": not preparation_errors,
            "source_inventory_unchanged": not source_inventory_errors,
            "original_unchanged": original_unchanged,
            "separate_corrected_artifact": correction_status != "completed"
            or corrected_metadata is not None,
            "corrected_artifact_path_bound": correction_status != "completed"
            or (
                corrected_metadata is not None
                and _resolve_artifact_ref(
                    str(correction_record.get("corrected_artifact", "")),
                    advisory_contract_path.parent,
                )
                == Path(corrected_metadata["path"])
            ),
            "corrected_artifact_hash_bound": correction_status != "completed"
            or (
                corrected_metadata is not None
                and correction_record.get("corrected_artifact_sha256")
                == corrected_metadata["sha256"]
            ),
            "corrected_artifact_re_reviewed": correction_status != "completed"
            or corrected_validation_metadata is not None,
            "format_check_artifacts_exist": not artifact_errors,
            "hidden_model_api_calls": False,
        },
        "declared_delivery_readiness": declared_readiness,
        "effective_delivery_readiness": declared_readiness if not errors else "blocked",
        "corrected_artifact": corrected_metadata,
        "corrected_validation": corrected_validation_metadata,
        "format_check_artifacts": format_check_artifacts,
        "interpretation": "Record completeness proves mechanical consistency only; semantic support, reasoning, recommendation quality, and professional judgement remain model-led and professionally reviewed.",
    }
    output_dir.mkdir(parents=True, exist_ok=True)
    _write_json(paths["review"], review)
    _write_json(paths["audit"], audit)
    _write_json(paths["recheck_tasks"], _recheck_tasks(review))
    paths["package"].write_text(
        _render_package(inventory, review, audit), encoding="utf-8"
    )
    return paths, audit


def _prepare_command(args: argparse.Namespace) -> int:
    paths = prepare_validation(
        args.deliverable,
        args.advisory_contract,
        args.output_dir,
        source_files=args.source_file,
        evidence_register=args.evidence_register,
        claim_register=args.claim_register,
    )
    LOGGER.info("Prepared advisory validation: %s", paths["deliverable_inventory"])
    return 0


def _package_command(args: argparse.Namespace) -> int:
    paths, audit = package_validation(
        args.deliverable_inventory,
        args.review_draft,
        args.advisory_contract,
        args.output_dir,
        corrected_deliverable=args.corrected_deliverable,
        corrected_deliverable_inventory=args.corrected_deliverable_inventory,
        corrected_review=args.corrected_review,
    )
    LOGGER.info("Wrote advisory validation package: %s", paths["package"])
    return 0 if audit["record_complete"] else 1


def _parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    subparsers = parser.add_subparsers(dest="command", required=True)
    prepare = subparsers.add_parser(
        "prepare", help="Extract and inventory a deliverable."
    )
    prepare.add_argument("deliverable", type=Path)
    prepare.add_argument("--advisory-contract", type=Path, required=True)
    prepare.add_argument("--output-dir", type=Path, required=True)
    prepare.add_argument("--source-file", type=Path, action="append", default=[])
    prepare.add_argument("--evidence-register", type=Path)
    prepare.add_argument("--claim-register", type=Path)
    prepare.set_defaults(handler=_prepare_command)

    package = subparsers.add_parser(
        "package", help="Audit and package a model-led review."
    )
    package.add_argument("deliverable_inventory", type=Path)
    package.add_argument("review_draft", type=Path)
    package.add_argument("--advisory-contract", type=Path, required=True)
    package.add_argument("--output-dir", type=Path, required=True)
    package.add_argument("--corrected-deliverable", type=Path)
    package.add_argument("--corrected-deliverable-inventory", type=Path)
    package.add_argument("--corrected-review", type=Path)
    package.set_defaults(handler=_package_command)
    return parser


def main() -> int:
    logging.basicConfig(level=logging.INFO, format="%(levelname)s %(message)s")
    parser = _parser()
    args = parser.parse_args()
    try:
        return int(args.handler(args))
    except AdvisoryValidationError as exc:
        parser.error(str(exc))
    return 2


if __name__ == "__main__":
    raise SystemExit(main())
